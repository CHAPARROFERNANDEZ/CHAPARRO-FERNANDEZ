import calendar
import json
import os
import re
import smtplib
import sys
import ssl
import zipfile
from datetime import datetime
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from io import BytesIO
from typing import Optional

import time
import numpy as np
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components

try:
    import plotly.express as px
except Exception:
    px = None

try:
    import yfinance as yf
except Exception:
    yf = None

import requests

from openpyxl import load_workbook
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.utils import get_column_letter
import bcrypt

if __name__ == "__main__":  # page config: solo se ejecuta con `streamlit run`, no al importar
    st.set_page_config(
        page_title="Chaparro Fernández Wealth",
        page_icon="assets/favicon_cf.png",
        layout="wide",
        initial_sidebar_state="collapsed",
    )

ARCHIVO = "inversiones.xlsx"
GDRIVE_FILE_ID = "1CImiIbg7kSLrYNpWgzPHEBCmI3KRVlBX"

# --- Migración a Postgres: interruptor de origen de datos -------------------------------------
# "drive" (default, sin tocar la variable en Railway) = comportamiento IDÉNTICO al actual.
# "postgres" = intenta leer de Postgres primero; si falla por cualquier motivo, cae a Drive
# automáticamente y avisa por log — la app nunca se queda sin datos por esto.
DATA_SOURCE = os.environ.get("DATA_SOURCE", "drive").strip().lower()
HOJA_INVERSIONES = "INVERSIONES"
HOJA_CALENDARIO = "CALENDARIO_NOTAS"
HOJA_CONTROL = "CONTROL_NOTAS"
HOJA_MOTOCLICK = "MOVIMIENTOS_MOTOCLICK"
HOJA_MOVIMIENTOS_BANCO = "MOVIMIENTOS_BANCO"

TASA_ANUAL_FUTBOL = 0.15
TASA_ANUAL_MOTOCLICK = 0.25
TASA_ANUAL_PARAGUAY = 0.15
TASA_ANUAL_BOLIVIA = 0.15
TASA_ANUAL_BITCOIN = 0.20

# Posición real del fondo en el ETF de Bitcoin, según confirmación de compra StoneX/Aragon Capital
BITCOIN_ETF_TICKER = "IBIT"
BITCOIN_ETF_NOMBRE = "iShares Bitcoin Trust ETF"
BITCOIN_ETF_UNIDADES = 1030
BITCOIN_ETF_PRECIO_COMPRA = 38.2899
BITCOIN_ETF_FECHA_COMPRA = pd.Timestamp("2026-03-30")
BITCOIN_ETF_CAPITAL_INVERTIDO = 40000.0   # capital aportado por los inversores (Jordi Especial + Chaparro Fernández)
BITCOIN_ETF_COSTE_REAL = 39443.55         # total pagado en la compra, incluida comisión

MESES_ES_EMAIL = {
    1: "enero", 2: "febrero", 3: "marzo", 4: "abril",
    5: "mayo", 6: "junio", 7: "julio", 8: "agosto",
    9: "septiembre", 10: "octubre", 11: "noviembre", 12: "diciembre",
}


# =========================
# MÓDULO EMAIL EXTRACTOS
# =========================

def _parse_lista_emails(cadena: str) -> list:
    """Convierte 'a@x.com, b@y.com; c@z.com' en ['a@x.com', 'b@y.com', 'c@z.com'].
    Admite coma, punto y coma o espacio como separador, y quita duplicados conservando el orden."""
    import re as _re_mail
    if not cadena:
        return []
    partes = [p.strip() for p in _re_mail.split(r"[,;\s]+", cadena) if p.strip()]
    vistos = set()
    limpio = []
    for p in partes:
        pl = p.lower()
        if pl not in vistos and "@" in p:
            vistos.add(pl)
            limpio.append(p)
    return limpio


def _leer_emails_inversores(df_inv: pd.DataFrame) -> dict:
    """
    Devuelve {inversor_upper: email_o_emails} leyendo la columna 'email' del Excel.
    Admite varios correos para el mismo inversor separados por coma en la misma celda
    (ej. 'yuri@x.com, contable@x.com') — el extracto se envía a todos ellos a la vez.
    """
    if "email" not in df_inv.columns:
        return {}
    mapa = {}
    for _, row in df_inv.iterrows():
        inv = str(row.get("inversor", "")).strip()
        mail = str(row.get("email", "")).strip()
        if inv and mail and mail.lower() not in ("", "nan", "none"):
            key = inv.upper()
            if key not in mapa:
                mapa[key] = mail
    return mapa


def _construir_cuerpo_html_email(inversor: str, mes: int, anio: int, total_intereses: float) -> str:
    """Genera el cuerpo HTML del email con diseño premium CF."""
    mes_str = MESES_ES_EMAIL.get(mes, str(mes)).capitalize()
    fecha_str = f"{mes_str} {anio}"
    total_fmt = f"${total_intereses:,.2f}"
    return f"""<!DOCTYPE html>
<html lang="es">
<head><meta charset="UTF-8"><meta name="viewport" content="width=device-width, initial-scale=1.0"></head>
<body style="margin:0;padding:0;background:#f4f6f9;font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
  <table width="100%" cellpadding="0" cellspacing="0" style="background:#f4f6f9;padding:40px 20px;">
    <tr><td align="center">
      <table width="600" cellpadding="0" cellspacing="0" style="background:#ffffff;border-radius:20px;overflow:hidden;box-shadow:0 8px 40px rgba(7,20,37,0.14);">
        <tr>
          <td style="background:linear-gradient(135deg,#0e2338 0%,#173b5c 60%,#bf9a5f 100%);padding:36px 40px;text-align:center;">
            <table cellpadding="0" cellspacing="0" style="margin:0 auto 16px auto;">
              <tr><td style="width:72px;height:72px;background:rgba(255,255,255,0.95);border-radius:16px;font-size:26px;font-weight:800;color:#0e2338;text-align:center;vertical-align:middle;border:1px solid rgba(191,154,95,0.5);">CF</td></tr>
            </table>
            <div style="color:#ffffff;font-size:22px;font-weight:800;letter-spacing:-0.5px;margin-bottom:6px;">Chaparro Fernández Wealth</div>
            <div style="color:rgba(255,255,255,0.75);font-size:13px;font-weight:500;">Extracto mensual de inversiones — {fecha_str}</div>
          </td>
        </tr>
        <tr>
          <td style="padding:36px 40px;">
            <p style="color:#334155;font-size:15px;margin:0 0 24px;">Estimado/a <strong style="color:#0e2338;">{inversor}</strong>,</p>
            <p style="color:#475569;font-size:14px;line-height:1.7;margin:0 0 28px;">Adjunto encontrará su extracto de inversiones correspondiente al mes de <strong>{fecha_str}</strong>. El documento detalla todas sus posiciones activas, los intereses devengados y el acumulado histórico desde el inicio de su relación con nosotros.</p>
            <table width="100%" cellpadding="0" cellspacing="0" style="background:linear-gradient(145deg,#ffffff,#f5f0e8);border:1px solid rgba(191,154,95,0.30);border-radius:14px;margin-bottom:28px;">
              <tr><td style="padding:24px;text-align:center;">
                <div style="font-size:12px;font-weight:600;color:#64748b;text-transform:uppercase;letter-spacing:1px;margin-bottom:8px;">Total intereses {fecha_str}</div>
                <div style="font-size:32px;font-weight:800;color:#0e2338;letter-spacing:-1px;">{total_fmt}</div>
              </td></tr>
            </table>
            <p style="color:#475569;font-size:14px;line-height:1.7;margin:0 0 28px;">El extracto completo en formato Excel está adjunto a este correo. Incluye el detalle operación por operación, el cierre mensual, el cierre anual y el resumen acumulado de toda su cartera.</p>
            <table width="100%" cellpadding="0" cellspacing="0" style="background:#f0f4ff;border-left:4px solid #0e2338;border-radius:4px;margin-bottom:28px;">
              <tr><td style="padding:14px 18px;"><p style="color:#1e293b;font-size:13px;margin:0;line-height:1.6;">🔒 <strong>Confidencial.</strong> Este correo y el documento adjunto contienen información financiera privada. Por favor no lo reenvíe a terceros sin autorización expresa.</p></td></tr>
            </table>
            <p style="color:#475569;font-size:14px;line-height:1.7;margin:0;">Ante cualquier consulta no dude en ponerse en contacto con nosotros. Gracias por su confianza.</p>
          </td>
        </tr>
        <tr>
          <td style="background:#071425;padding:24px 40px;text-align:center;">
            <p style="color:rgba(255,255,255,0.5);font-size:12px;margin:0;line-height:1.8;">Chaparro Fernández Wealth · Sistema privado de inversiones<br>Este mensaje ha sido generado automáticamente — por favor no responda directamente.</p>
          </td>
        </tr>
      </table>
    </td></tr>
  </table>
</body>
</html>"""


def _excel_a_pdf(extracto_bytes: bytes, inversor: str = "", mes: int = 0, anio: int = 0,
                  total_intereses: float = 0.0) -> bytes:
    """
    Convierte el extracto Excel a PDF usando reportlab.
    Lee las hojas PORTADA y DETALLE del Excel generado
    y reproduce fielmente colores, tipografía y estructura. Sin dependencias externas.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors as rl_colors
        from reportlab.lib.units import mm
        from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                         Paragraph, Spacer, PageBreak)
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

        def _hex_color(hex_str):
            if not hex_str or hex_str in ('00000000', '000000', 'FF000000'):
                return None
            h = hex_str.lstrip('#')
            if len(h) == 8:
                h = h[2:]
            if len(h) == 6:
                try:
                    return rl_colors.Color(int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255)
                except Exception:
                    return None
            return None

        wb = load_workbook(BytesIO(extracto_bytes), data_only=True)
        output = BytesIO()
        doc = SimpleDocTemplate(output, pagesize=A4,
                                 leftMargin=15*mm, rightMargin=15*mm,
                                 topMargin=15*mm, bottomMargin=15*mm)
        story = []

        # Paleta corporativa
        C_AZUL_OSC = rl_colors.Color(13/255,  33/255,  55/255)
        C_AZUL_MED = rl_colors.Color(26/255,  63/255,  92/255)
        C_AZUL_CL  = rl_colors.Color(214/255, 233/255, 248/255)
        C_VERDE    = rl_colors.Color(217/255, 234/255, 211/255)
        C_DORADO   = rl_colors.Color(255/255, 242/255, 204/255)
        C_NARANJA  = rl_colors.Color(252/255, 229/255, 205/255)
        C_GRIS     = rl_colors.Color(0.97, 0.98, 0.99)
        C_BLANCO   = rl_colors.white
        C_BORDE    = rl_colors.Color(0.80, 0.80, 0.80)
        C_BORDE_L  = rl_colors.Color(0.88, 0.88, 0.88)

        def _tbl_style_base():
            return [
                ('TOPPADDING',    (0,0), (-1,-1), 3),
                ('BOTTOMPADDING', (0,0), (-1,-1), 3),
                ('BOX',           (0,0), (-1,-1), 0.5, C_BORDE),
                ('INNERGRID',     (0,0), (-1,-1), 0.3, C_BORDE_L),
            ]

        # ── PORTADA ──────────────────────────────────────────────────────────
        ws_p = wb['PORTADA'] if 'PORTADA' in wb.sheetnames else None
        if ws_p:
            rows_p = list(ws_p.iter_rows(min_row=1, max_row=ws_p.max_row, values_only=True))
            inv_nombre = ''; fecha_corte = ''; kpi_labels = []; kpi_values = []
            mes_rows = []; tot_row = None; nota_desc = ''; footer_txt = ''

            for i, row in enumerate(rows_p, 1):
                vals = [v for v in row if v is not None]
                if i == 8  and vals: inv_nombre  = str(vals[0])
                elif i == 10 and vals: fecha_corte = str(vals[0])
                elif i == 14: kpi_labels = [str(v) for v in row if v is not None]
                elif i == 15: kpi_values = [v     for v in row if v is not None]
                elif i == 17 and vals: nota_desc  = str(vals[0])
                elif 21 <= i <= 29:
                    r = [v for v in row if v is not None]
                    if r: mes_rows.append(r)
                elif i == 30:
                    r = [v for v in row if v is not None]
                    if r: tot_row = r
                elif i == 32 and vals: footer_txt = str(vals[0])

            # Header
            hdr = Table([[
                Paragraph('<b>CF</b>', ParagraphStyle('cf', fontName='Helvetica-Bold', fontSize=20, textColor=C_BLANCO, alignment=TA_CENTER)),
                Paragraph('Chaparro Fernández Wealth  ·  Extracto de inversiones',
                          ParagraphStyle('ht', fontName='Helvetica', fontSize=11, textColor=C_BLANCO))
            ]], colWidths=[22*mm, 158*mm])
            hdr.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1), C_AZUL_OSC),
                ('VALIGN',   (0,0),(-1,-1), 'MIDDLE'),
                ('TOPPADDING',(0,0),(-1,-1), 10), ('BOTTOMPADDING',(0,0),(-1,-1), 10),
                ('LEFTPADDING',(0,0),(-1,-1), 8),
            ]))
            story.append(hdr)
            story.append(Spacer(1, 5*mm))

            # Nombre inversor + fecha
            story.append(Paragraph(inv_nombre, ParagraphStyle('inv', fontName='Helvetica-Bold', fontSize=18, textColor=C_AZUL_OSC, spaceAfter=2)))
            story.append(Paragraph(fecha_corte, ParagraphStyle('fc', fontName='Helvetica', fontSize=9, textColor=rl_colors.Color(0.33,0.33,0.33), spaceAfter=5)))

            # KPIs
            if kpi_labels and kpi_values:
                bg_map = [C_AZUL_CL, C_VERDE, C_DORADO]
                tc_map = [C_AZUL_OSC, rl_colors.Color(30/255,70/255,32/255), rl_colors.Color(74/255,48/255,0)]
                cw = 180*mm / max(len(kpi_labels[:3]), 1)
                lbl_row = [Paragraph(f'<b>{lbl}</b>', ParagraphStyle(f'kl{i}', fontName='Helvetica-Bold', fontSize=7, textColor=tc_map[i], alignment=TA_CENTER)) for i, lbl in enumerate(kpi_labels[:3])]
                val_row = [Paragraph(f'<b>${float(v):,.2f}</b>', ParagraphStyle(f'kv{i}', fontName='Helvetica-Bold', fontSize=14, textColor=tc_map[i], alignment=TA_CENTER)) for i, v in enumerate(kpi_values[:3])]
                n = len(kpi_labels[:3])
                kpi_lbl_t = Table([lbl_row], colWidths=[cw]*n)
                kpi_val_t = Table([val_row], colWidths=[cw]*n)
                for kt, pad_b in [(kpi_lbl_t, 3), (kpi_val_t, 8)]:
                    style_kpi = _tbl_style_base()
                    for ki in range(n):
                        style_kpi.append(('BACKGROUND', (ki,0), (ki,0), bg_map[ki]))
                    style_kpi.append(('BOTTOMPADDING', (0,0), (-1,-1), pad_b))
                    kt.setStyle(TableStyle(style_kpi))
                story.append(kpi_lbl_t)
                story.append(kpi_val_t)
                story.append(Spacer(1, 4*mm))

            # Nota descriptiva
            if nota_desc:
                story.append(Paragraph(f'<i>{nota_desc}</i>',
                    ParagraphStyle('nd', fontName='Helvetica-Oblique', fontSize=8,
                                   textColor=rl_colors.Color(0.33,0.33,0.33), spaceAfter=5)))

            # Resumen mensual en portada
            if mes_rows:
                res_hdr_t = Table([[Paragraph('<b>RESUMEN MENSUAL DE INTERESES</b>',
                    ParagraphStyle('rh', fontName='Helvetica-Bold', fontSize=9, textColor=C_BLANCO, alignment=TA_CENTER))]], colWidths=[180*mm])
                res_hdr_t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1), C_AZUL_MED),
                    ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5)]))
                story.append(res_hdr_t)

                n_cols_r = len(mes_rows[0]) if mes_rows else 3
                if n_cols_r == 4:
                    hdrs_r = ['MES','GENERADO ($)','PAGADO ($)','SALDO ($)']
                    cw_r = [45*mm,45*mm,45*mm,45*mm]
                else:
                    hdrs_r = ['MES','INTERESES ($)','ACUMULADO ($)']
                    cw_r = [60*mm,60*mm,60*mm]

                col_hdr_t = Table([hdrs_r], colWidths=cw_r)
                col_hdr_t.setStyle(TableStyle([
                    ('BACKGROUND',(0,0),(-1,-1), C_AZUL_CL),
                    ('TEXTCOLOR', (0,0),(-1,-1), C_AZUL_OSC),
                    ('FONTNAME',  (0,0),(-1,-1), 'Helvetica-Bold'),
                    ('FONTSIZE',  (0,0),(-1,-1), 8),
                    ('ALIGN',     (0,0),(-1,-1), 'CENTER'),
                    ('TOPPADDING',(0,0),(-1,-1),4),('BOTTOMPADDING',(0,0),(-1,-1),4),
                    ('BOX',(0,0),(-1,-1),0.5,C_BORDE),('INNERGRID',(0,0),(-1,-1),0.3,C_BORDE_L),
                ]))
                story.append(col_hdr_t)

                data_r = []
                for mr in mes_rows:
                    row_r = [str(mr[0])] + [f'${float(v):,.2f}' for v in mr[1:n_cols_r]]
                    data_r.append(row_r)
                if data_r:
                    t_data_r = Table(data_r, colWidths=cw_r)
                    s = _tbl_style_base()
                    s += [('ROWBACKGROUNDS',(0,0),(-1,-1),[C_BLANCO, C_GRIS]),
                          ('FONTNAME',(0,0),(-1,-1),'Helvetica'),('FONTSIZE',(0,0),(-1,-1),8),
                          ('ALIGN',(0,0),(0,-1),'CENTER'),('ALIGN',(1,0),(-1,-1),'RIGHT')]
                    t_data_r.setStyle(TableStyle(s))
                    story.append(t_data_r)

                if tot_row:
                    row_t = [str(tot_row[0])] + [f'${float(v):,.2f}' for v in tot_row[1:n_cols_r]]
                    t_tot_r = Table([row_t], colWidths=cw_r)
                    t_tot_r.setStyle(TableStyle([
                        ('BACKGROUND',(0,0),(-1,-1), C_DORADO),
                        ('TEXTCOLOR', (0,0),(-1,-1), rl_colors.Color(74/255,48/255,0)),
                        ('FONTNAME',  (0,0),(-1,-1), 'Helvetica-Bold'),
                        ('FONTSIZE',  (0,0),(-1,-1), 9),
                        ('ALIGN',(0,0),(0,-1),'CENTER'),('ALIGN',(1,0),(-1,-1),'RIGHT'),
                        ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
                        ('BOX',(0,0),(-1,-1),0.5,C_BORDE),('INNERGRID',(0,0),(-1,-1),0.3,C_BORDE_L),
                    ]))
                    story.append(t_tot_r)

            story.append(Spacer(1, 5*mm))
            if footer_txt:
                ft = Table([[Paragraph(footer_txt, ParagraphStyle('ft', fontName='Helvetica', fontSize=7,
                    textColor=rl_colors.Color(0.67,0.67,0.67), alignment=TA_CENTER))]], colWidths=[180*mm])
                ft.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1), C_AZUL_OSC),
                    ('TOPPADDING',(0,0),(-1,-1),6),('BOTTOMPADDING',(0,0),(-1,-1),6)]))
                story.append(ft)

        # ── PÁGINA 2: DETALLE ─────────────────────────────────────────────────
        story.append(PageBreak())
        ws_d = wb['DETALLE'] if 'DETALLE' in wb.sheetnames else None
        if ws_d:
            det_rows   = list(ws_d.iter_rows(min_row=1, max_row=ws_d.max_row, values_only=True))
            det_styles = []
            for row in ws_d.iter_rows(min_row=1, max_row=ws_d.max_row):
                first = row[0]
                bg = fc = None
                bold = False
                if first.fill and first.fill.fgColor and first.fill.fgColor.type == 'rgb':
                    bg = _hex_color(first.fill.fgColor.rgb)
                if first.font:
                    if first.font.color and first.font.color.type == 'rgb':
                        fc = _hex_color(first.font.color.rgb)
                    bold = bool(first.font.bold)
                det_styles.append((bg, fc, bold))

            titulo_det    = str(det_rows[0][0])  if det_rows and det_rows[0][0]  else ''
            subtitulo_det = str(det_rows[1][0])  if len(det_rows) > 1 and det_rows[1][0] else ''

            t_hdr_d = Table([[Paragraph(f'<b>{titulo_det}</b>',
                ParagraphStyle('thd', fontName='Helvetica-Bold', fontSize=11, textColor=C_BLANCO))]], colWidths=[180*mm])
            t_hdr_d.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1), C_AZUL_OSC),
                ('TOPPADDING',(0,0),(-1,-1),8),('BOTTOMPADDING',(0,0),(-1,-1),8),('LEFTPADDING',(0,0),(-1,-1),8)]))
            story.append(t_hdr_d)
            story.append(Paragraph(subtitulo_det, ParagraphStyle('sd', fontName='Helvetica', fontSize=8,
                textColor=rl_colors.Color(0.27,0.27,0.27), spaceAfter=4, spaceBefore=3)))

            # Columnas visibles en el PDF: Activo(0), Mes(1), Fecha inversión(2), Capital(3), Interés mes(4)
            # Ya no hay columnas ocultas — el Excel contiene exactamente lo que se muestra
            COLS_VISIBLES = [0, 1, 2, 3, 4]
            if len(det_rows) > 3:
                full_hdr = [str(v) if v else '' for v in det_rows[3]]
                col_hdr = [full_hdr[ci] if ci < len(full_hdr) else '' for ci in COLS_VISIBLES]
                cw_d = [38*mm, 22*mm, 28*mm, 46*mm, 46*mm]

                hdr_d_t = Table([col_hdr], colWidths=cw_d)
                hdr_d_t.setStyle(TableStyle([
                    ('BACKGROUND',(0,0),(-1,-1), C_AZUL_MED),
                    ('TEXTCOLOR', (0,0),(-1,-1), C_BLANCO),
                    ('FONTNAME',  (0,0),(-1,-1), 'Helvetica-Bold'),
                    ('FONTSIZE',  (0,0),(-1,-1), 9),
                    ('ALIGN',     (0,0),(-1,-1), 'CENTER'),
                    ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
                    ('BOX',(0,0),(-1,-1),0.5,C_BORDE),('INNERGRID',(0,0),(-1,-1),0.3,C_BORDE_L),
                ]))
                story.append(hdr_d_t)

                for ri, row in enumerate(det_rows[4:], 4):
                    bg_c, fc_c, is_bold = det_styles[ri]
                    if bg_c is None: bg_c = C_BLANCO
                    if fc_c is None: fc_c = rl_colors.black
                    fn = 'Helvetica-Bold' if is_bold else 'Helvetica'

                    def fmt_cell(ci, v):
                        if v is None: return ''
                        if ci in (3, 4) and isinstance(v, (int, float)): return f'${float(v):,.2f}'
                        return str(v)

                    is_cierre = str(row[0] or '').startswith('CIERRE') or str(row[0] or '').startswith('   ')
                    if is_cierre:
                        label   = str(row[0] or '')
                        val_cap = fmt_cell(3, row[3] if len(row) > 3 else None)
                        val_int = fmt_cell(4, row[4] if len(row) > 4 else None)
                        t_row = Table([[label, '', val_cap, val_int, '']], colWidths=cw_d)
                        t_row.setStyle(TableStyle([
                            ('SPAN',       (0,0),(1,0)),
                            ('SPAN',       (3,0),(4,0)),
                            ('BACKGROUND', (0,0),(-1,-1), bg_c),
                            ('TEXTCOLOR',  (0,0),(-1,-1), fc_c),
                            ('FONTNAME',   (0,0),(-1,-1), fn),
                            ('FONTSIZE',   (0,0),(-1,-1), 8.5),
                            ('ALIGN',      (2,0),(-1,-1), 'RIGHT'),
                            ('ALIGN',      (0,0),(1,0),   'LEFT'),
                            ('LEFTPADDING',(0,0),(0,-1), 6),
                            ('TOPPADDING', (0,0),(-1,-1), 3),('BOTTOMPADDING',(0,0),(-1,-1),3),
                            ('BOX',(0,0),(-1,-1),0.3,C_BORDE),('INNERGRID',(0,0),(-1,-1),0.3,C_BORDE_L),
                        ]))
                    else:
                        row_data = [fmt_cell(ci, row[ci] if ci < len(row) else None) for ci in COLS_VISIBLES]
                        t_row = Table([row_data], colWidths=cw_d)
                        t_row.setStyle(TableStyle([
                            ('BACKGROUND',(0,0),(-1,-1), bg_c),
                            ('TEXTCOLOR', (0,0),(-1,-1), fc_c),
                            ('FONTNAME',  (0,0),(-1,-1), fn),
                            ('FONTSIZE',  (0,0),(-1,-1), 9),
                            ('ALIGN',     (0,0),(-1,-1), 'CENTER'),
                            ('ALIGN',     (2,0),(-1,-1), 'RIGHT'),
                            ('TOPPADDING',(0,0),(-1,-1),3),('BOTTOMPADDING',(0,0),(-1,-1),3),
                            ('BOX',(0,0),(-1,-1),0.3,C_BORDE),('INNERGRID',(0,0),(-1,-1),0.3,C_BORDE_L),
                        ]))
                    story.append(t_row)

        # NOTA: se ha quitado la antigua "PÁGINA 3: RESUMEN MENSUAL" porque duplicaba
        # exactamente la misma tabla "RESUMEN MENSUAL DE INTERESES" que ya aparece en
        # la PORTADA (página 1) — el PDF ahora tiene solo 2 páginas: PORTADA y DETALLE.
        doc.build(story)
        return output.getvalue()

    except Exception as e:
        return b""



def enviar_extracto_email(destinatario, inversor: str, mes: int, anio: int,
                           extracto_bytes: bytes, nombre_archivo: str,
                           total_intereses: float, smtp_sender: str,
                           smtp_password: str,
                           display_name: str = "Chaparro Fernández Wealth") -> tuple:
    """
    Envía el extracto como PDF adjunto por email usando Gmail SMTP.
    'destinatario' puede ser un string con un solo email, un string con varios
    separados por coma ('a@x.com, b@y.com'), o directamente una lista de emails —
    en todos los casos se manda EL MISMO extracto a todos los destinatarios de una vez.
    """
    mes_str = MESES_ES_EMAIL.get(mes, str(mes)).capitalize()
    if isinstance(destinatario, str):
        destinatarios = _parse_lista_emails(destinatario)
    else:
        destinatarios = list(destinatario)
    if not destinatarios:
        return False, "No hay ninguna dirección de email válida."
    try:
        msg = MIMEMultipart("mixed")
        msg["Subject"] = f"Extracto de inversiones — {mes_str} {anio} | Chaparro Fernández Wealth"
        msg["From"]    = f"{display_name} <{smtp_sender}>"
        msg["To"]      = ", ".join(destinatarios)
        cuerpo = _construir_cuerpo_html_email(inversor, mes, anio, total_intereses)
        msg.attach(MIMEText(cuerpo, "html", "utf-8"))

        if extracto_bytes:
            # Convertir Excel → PDF
            pdf_bytes = _excel_a_pdf(extracto_bytes, inversor, mes, anio, total_intereses)
            nombre_pdf = nombre_archivo.replace(".xlsx", ".pdf")
            if pdf_bytes:
                adjunto = MIMEBase("application", "pdf")
                adjunto.set_payload(pdf_bytes)
                encoders.encode_base64(adjunto)
                adjunto.add_header("Content-Disposition", "attachment", filename=nombre_pdf)
                msg.attach(adjunto)
            else:
                # Fallback: si falla la conversión adjuntar el Excel
                adjunto = MIMEBase("application", "vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                adjunto.set_payload(extracto_bytes)
                encoders.encode_base64(adjunto)
                adjunto.add_header("Content-Disposition", "attachment", filename=nombre_archivo)
                msg.attach(adjunto)

        contexto = ssl.create_default_context()
        with smtplib.SMTP("smtp.gmail.com", 587) as servidor:
            servidor.ehlo()
            servidor.starttls(context=contexto)
            servidor.login(smtp_sender, smtp_password)
            servidor.sendmail(smtp_sender, destinatarios, msg.as_bytes())
        return True, ""
    except smtplib.SMTPAuthenticationError:
        return False, "Error de autenticación Gmail. Usa una contraseña de aplicación (no tu contraseña normal)."
    except smtplib.SMTPRecipientsRefused:
        return False, f"Alguna(s) de estas direcciones fue(ron) rechazada(s) por el servidor: {', '.join(destinatarios)}."
    except Exception as e:
        return False, str(e)


def seccion_envio_extractos_email(df_inv: pd.DataFrame, generar_extractos_fn,
                                   anio_override=None, mes_override=None,
                                   inversor_override=None, modo_override=None):
    """Panel de envío por email. Recibe año/mes/inversor ya seleccionados desde seccion_extractos."""

    mapa_emails = _leer_emails_inversores(df_inv)

    if not mapa_emails:
        st.error("⚠️ No hay emails registrados en el Excel.")
        with st.expander("¿Cómo añadir emails?"):
            st.markdown("""
Abre `inversiones.xlsx`, en la hoja **INVERSIONES** añade una columna llamada `email`
y rellena el correo de cada inversor en al menos una de sus filas.

**¿Quieres que el mismo extracto llegue a más de una persona** (por ejemplo, al inversor
y a su asesor o pareja)? Pon los dos correos en la misma celda separados por coma:
`inversor@gmail.com, asesor@gmail.com` — se enviará el mismo extracto a ambos a la vez,
en un único correo.

Luego recarga el Excel desde el menú Gestión de Excel.
""")
        return

    # ── Credenciales Gmail ────────────────────────────────────────────────────
    secrets_ok = False
    try:
        smtp_sender   = st.secrets["email"]["sender"]
        smtp_password = st.secrets["email"]["password"]
        display_name  = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
        st.success(f"✅ Gmail configurado: **{smtp_sender}**")
        secrets_ok = True
    except Exception:
        pass

    if not secrets_ok:
        with st.expander("⚙️ Configurar cuenta Gmail", expanded=True):
            st.caption("Solo necesitas hacer esto una vez. Usa una contraseña de aplicación, no tu contraseña normal de Gmail. Puedes crearla en [myaccount.google.com/apppasswords](https://myaccount.google.com/apppasswords).")
            col1, col2 = st.columns(2)
            smtp_sender   = col1.text_input("Cuenta Gmail emisora", placeholder="tucuenta@gmail.com", key="smtp_sender_input")
            smtp_password = col2.text_input("Contraseña de aplicación", type="password", placeholder="xxxx xxxx xxxx xxxx", key="smtp_password_input")
            display_name  = st.text_input("Nombre remitente", value="Chaparro Fernández Wealth", key="smtp_display_name")

    ocultar_activo_email = st.checkbox(
        "🔒 Ocultar en qué está invertido (igual que en el Portal de inversor)",
        value=True, key="email_ocultar_activo",
    )

    st.divider()

    # ── Destinatarios: usar el inversor ya seleccionado arriba ────────────────
    anio_email = anio_override or datetime.today().year
    mes_email  = mes_override  or (datetime.today().month - 1 or 12)
    mes_str_email = MESES_ES_EMAIL.get(mes_email, str(mes_email)).capitalize()
    inversores_con_email = sorted(mapa_emails.keys())

    # Si viene un inversor concreto de la pestaña Descargar, preseleccionarlo
    if inversor_override and modo_override == "Un inversor":
        inv_upper = str(inversor_override).upper()
        if inv_upper in mapa_emails:
            inversores_enviar = [inv_upper]
            st.info(f"Se enviará el extracto de **{inv_upper}** a **{mapa_emails[inv_upper]}**")
        else:
            st.warning(f"**{inv_upper}** no tiene email registrado en el Excel.")
            return
    else:
        todos_str = [f"{inv}  →  {mapa_emails[inv]}" for inv in inversores_con_email]
        seleccionados_str = st.multiselect(
            "Inversores a incluir",
            options=todos_str, default=todos_str, key="email_seleccion"
        )
        inversores_enviar = [s.split("  →  ")[0].strip() for s in seleccionados_str]
        if not inversores_enviar:
            st.warning("Selecciona al menos un inversor.")
            return

    # ── Test de conexión ──────────────────────────────────────────────────────
    with st.expander("🔌 Probar conexión Gmail", expanded=False):
        email_prueba = st.text_input("Email de destino para la prueba", placeholder="tuemail@gmail.com", key="email_prueba_dest")
        if st.button("Enviar email de prueba", key="btn_email_prueba"):
            if not smtp_sender or not smtp_password:
                st.error("Configura las credenciales Gmail primero.")
            elif not email_prueba:
                st.error("Introduce un email de destino.")
            else:
                ok, err = enviar_extracto_email(
                    destinatario=email_prueba, inversor="INVERSOR TEST",
                    mes=mes_email, anio=anio_email, extracto_bytes=b"",
                    nombre_archivo="test.xlsx", total_intereses=12345.67,
                    smtp_sender=smtp_sender, smtp_password=smtp_password, display_name=display_name,
                )
                st.success(f"✅ Email de prueba enviado a {email_prueba}.") if ok else st.error(f"❌ {err}")

    st.divider()

    # ── Botón de envío ────────────────────────────────────────────────────────
    if st.button(
        f"📨 Enviar extracto{'s' if len(inversores_enviar) > 1 else ''} de {mes_str_email} {anio_email}  ({len(inversores_enviar)} inversor{'es' if len(inversores_enviar) > 1 else ''})",
        type="primary", use_container_width=True, key="btn_enviar_extractos"
    ):
        if not smtp_sender or not smtp_password:
            st.error("Configura las credenciales Gmail antes de enviar.")
            return

        progreso   = st.progress(0, text="Preparando envíos...")
        resultados = []

        for i, inversor in enumerate(inversores_enviar):
            progreso.progress(i / len(inversores_enviar), text=f"Generando extracto de {inversor}... ({i+1}/{len(inversores_enviar)})")
            try:
                archivos = generar_extractos_fn(df_inv, "Un inversor", inversor, anio_email, mes_email)
            except Exception as e:
                resultados.append({"Inversor": inversor, "Email": mapa_emails.get(inversor, "—"), "Estado": "❌ Error al generar", "Detalle": str(e)})
                continue

            if not archivos:
                resultados.append({"Inversor": inversor, "Email": mapa_emails.get(inversor, "—"), "Estado": "⚠️ Sin datos", "Detalle": "Sin inversiones en este periodo."})
                continue

            t = archivos[0]
            nombre_archivo   = t[0]
            extracto_bytes   = t[1]  # t[1] = Excel formateado (PORTADA/DETALLE) — para adjuntar/enviar
            excel_crudo      = t[2] if len(t) > 2 else None  # t[2] = Excel con TOTALES_MES/DETALLE — para leer el total
            email_dest       = mapa_emails.get(inversor, "")

            total_intereses_email = 0.0
            try:
                from openpyxl import load_workbook as _lw
                wb_tmp = _lw(BytesIO(excel_crudo if excel_crudo else extracto_bytes))
                if "TOTALES_MES" in wb_tmp.sheetnames:
                    for row in wb_tmp["TOTALES_MES"].iter_rows(min_row=2, values_only=True):
                        if row and row[0] and f"{mes_email:02d}/{anio_email}" in str(row[0]):
                            try:
                                total_intereses_email += float(row[1] or 0)
                            except Exception:
                                pass
            except Exception:
                pass

            extracto_bytes_envio = preparar_extracto_privado_inversor(extracto_bytes) if ocultar_activo_email else extracto_bytes

            progreso.progress((i + 0.5) / len(inversores_enviar), text=f"Enviando a {email_dest}...")
            ok, err = enviar_extracto_email(
                destinatario=email_dest, inversor=inversor,
                mes=mes_email, anio=anio_email,
                extracto_bytes=extracto_bytes_envio, nombre_archivo=nombre_archivo,
                total_intereses=total_intereses_email,
                smtp_sender=smtp_sender, smtp_password=smtp_password, display_name=display_name,
            )
            resultados.append({
                "Inversor": inversor, "Email": email_dest,
                "Estado": "✅ Enviado" if ok else "❌ Error",
                "Detalle": nombre_archivo if ok else err,
            })

        progreso.progress(1.0, text="¡Proceso completado!")
        df_res   = pd.DataFrame(resultados)
        enviados = (df_res["Estado"] == "✅ Enviado").sum()
        errores  = df_res["Estado"].str.startswith("❌").sum()

        col1, col2, col3 = st.columns(3)
        col1.metric("✅ Enviados", enviados)
        col2.metric("⚠️ Sin datos", (df_res["Estado"] == "⚠️ Sin datos").sum())
        col3.metric("❌ Errores", errores)
        st.dataframe(df_res, use_container_width=True, hide_index=True)

        if errores == 0 and enviados > 0:
            st.balloons()
            st.success(f"🎉 Todos los extractos de {mes_str_email} {anio_email} enviados correctamente.")
        elif errores > 0:
            st.warning(f"Se enviaron {enviados} pero hubo {errores} error(es). Revisa la columna Detalle.")


# =========================
# ESTILO PROFESIONAL
# =========================
def aplicar_estilo_profesional():
    st.markdown(
        """
        <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');
        html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

        /* ── Paleta ───────────────────────────────────────────────────────
           Fondo claro y uniforme (nada de diagonales oscuras). El azul
           marino se usa solo como color de acento puntual (títulos,
           botones, iconos) y el dorado como detalle muy sutil. */
        .stApp {
            background: #eef1f5;
            background-attachment: fixed;
        }
        .block-container {
            max-width: 1240px;
            padding-top: 2rem;
            padding-bottom: 3rem;
            background: #ffffff;
            border: 1px solid #e4e8ee;
            border-radius: 20px;
            box-shadow: 0 10px 34px rgba(15, 35, 60, 0.06);
            margin-top: 1.2rem;
            margin-bottom: 1.2rem;
        }

        /* ── Sidebar clara: mismo tono que el resto de la app, sin
           necesidad de forzar texto blanco sobre fondo oscuro. En móvil
           dejamos que Streamlit la abra/cierre de forma normal (antes se
           forzaba siempre visible, y eso es justo lo que la dejaba como
           una tira estrecha con el texto partido letra a letra en el
           borde izquierdo cuando estaba "cerrada"). */
        section[data-testid="stSidebar"] {
            background: #f7f8fb;
            border-right: 1px solid #e4e8ee;
        }
        section[data-testid="stSidebar"] label,
        section[data-testid="stSidebar"] p,
        section[data-testid="stSidebar"] span,
        section[data-testid="stSidebar"] div { color: #1f2937 !important; }
        section[data-testid="stSidebar"] input { color: #1f2937 !important; }

        h1, h2, h3 { color: #16324f; letter-spacing: -0.02em; }

        div[data-testid="stMetric"] {
            background: #ffffff;
            border: 1px solid #e4e8ee;
            border-radius: 16px;
            padding: 18px 20px;
            box-shadow: 0 6px 18px rgba(15, 35, 60, 0.05);
        }
        div[data-testid="stMetricValue"] { color: #16324f; font-weight: 800; }

        .stButton > button, .stDownloadButton > button, button[kind="primary"],
        [data-testid="stAppViewContainer"] .stFormSubmitButton button,
        section[data-testid="stSidebar"] .stFormSubmitButton button {
            border-radius: 10px !important;
            background: #1e3a5f !important;
            color: white !important;
            border: 0 !important;
            font-weight: 600 !important;
            padding: 0.5rem 1.1rem !important;
            box-shadow: none;
            transition: background 0.15s ease;
        }
        .stButton > button:hover, .stDownloadButton > button:hover, button[kind="primary"]:hover,
        [data-testid="stAppViewContainer"] .stFormSubmitButton button:hover,
        section[data-testid="stSidebar"] .stFormSubmitButton button:hover {
            background: #2a4d78 !important;
        }

        div[data-testid="stDataFrame"] {
            border-radius: 14px;
            overflow: hidden;
            border: 1px solid #e4e8ee;
        }

        /* ── Cabecera de marca: tarjeta clara con un acento de color,
           en vez del degradado oscuro que costaba leer. */
        .brand-hero {
            display: flex; align-items: center; justify-content: space-between; gap: 24px;
            padding: 22px 28px; margin-bottom: 26px; border-radius: 18px;
            background: #ffffff; border: 1px solid #e4e8ee; border-left: 5px solid #1e3a5f;
            box-shadow: 0 6px 18px rgba(15, 35, 60, 0.05); color: #16324f;
        }
        .brand-left { display: flex; align-items: center; gap: 18px; }
        .brand-logo {
            width: 60px; height: 60px; border-radius: 14px; background: #1e3a5f; color: #ffffff;
            display: flex; align-items: center; justify-content: center; font-size: 24px; font-weight: 800;
            letter-spacing: -0.06em;
        }
        .brand-title { font-size: 24px; line-height: 1.15; font-weight: 800; letter-spacing: -0.02em; color: #16324f; }
        .brand-subtitle { margin-top: 6px; color: #64748b; font-size: 14px; font-weight: 500; }
        .brand-tag {
            padding: 8px 14px; border-radius: 999px; background: #f0f4f9;
            border: 1px solid #d8e0ea; font-size: 13px; font-weight: 700; color: #1e3a5f; white-space: nowrap;
        }

        .login-card {
            max-width: 440px; margin: 6vh auto 0 auto; padding: 34px 34px 30px 34px; border-radius: 20px;
            background: #ffffff; border: 1px solid #e4e8ee;
            box-shadow: 0 14px 40px rgba(15, 35, 60, 0.10); text-align: center;
        }
        .login-logo {
            width: 76px; height: 76px; border-radius: 18px; margin: 0 auto 18px auto;
            background: #1e3a5f; color: #ffffff;
            display: flex; align-items: center; justify-content: center; font-size: 30px; font-weight: 800;
            letter-spacing: -0.06em;
        }
        .login-title { font-size: 24px; font-weight: 800; color: #16324f; letter-spacing: -0.02em; margin-bottom: 6px; }
        .login-subtitle { font-size: 14px; color: #64748b; margin-bottom: 20px; }
        #MainMenu, footer {visibility: hidden;}
        header {visibility: visible;}

        /* ── Texto siempre legible, sin depender del modo claro/oscuro
           del dispositivo (evita el bug de texto invisible en móvil). */
        [data-testid="stAppViewContainer"] p,
        [data-testid="stAppViewContainer"] span,
        [data-testid="stAppViewContainer"] li,
        [data-testid="stAppViewContainer"] label,
        [data-testid="stAppViewContainer"] div:not([class*="brand-"]):not(.login-card):not(.login-logo),
        [data-testid="stMarkdownContainer"],
        [data-testid="stMarkdownContainer"] *,
        [data-testid="stChatMessage"],
        [data-testid="stChatMessage"] *,
        [data-testid="stChatMessageContent"],
        [data-testid="stChatMessageContent"] * {
            color: #1f2937 !important;
        }
        [data-testid="stChatMessage"] {
            background: #f7f8fb !important;
            border-radius: 14px;
            border: 1px solid #e4e8ee;
            padding: 6px 10px;
        }

        /* ── Botón para abrir el menú lateral, grande y fácil de tocar */
        [data-testid="collapsedControl"] {
            background: #1e3a5f !important;
            border-radius: 10px !important;
            padding: 8px !important;
            box-shadow: 0 4px 12px rgba(15, 35, 60, 0.18);
            top: 0.6rem !important;
            left: 0.6rem !important;
        }
        [data-testid="collapsedControl"] svg {
            width: 28px !important;
            height: 28px !important;
            color: #ffffff !important;
            fill: #ffffff !important;
        }

        /* ── El texto de dentro de los botones (Descargar Excel, Cerrar
           sesión, Entrar, Actualizar contraseña... TODOS los botones con
           fondo azul marino) no debe quedar afectado por las reglas de
           arriba que oscurecen el texto general: aquí se fuerza blanco
           con máxima prioridad, tanto en la sidebar como en el resto de
           la app, cubriendo también los botones dentro de formularios. */
        [data-testid="stAppViewContainer"] .stButton button,
        [data-testid="stAppViewContainer"] .stButton button *,
        [data-testid="stAppViewContainer"] .stDownloadButton button,
        [data-testid="stAppViewContainer"] .stDownloadButton button *,
        [data-testid="stAppViewContainer"] .stFormSubmitButton button,
        [data-testid="stAppViewContainer"] .stFormSubmitButton button *,
        section[data-testid="stSidebar"] .stButton button,
        section[data-testid="stSidebar"] .stButton button *,
        section[data-testid="stSidebar"] .stDownloadButton button,
        section[data-testid="stSidebar"] .stDownloadButton button *,
        section[data-testid="stSidebar"] .stFormSubmitButton button,
        section[data-testid="stSidebar"] .stFormSubmitButton button *,
        button[kind="primary"], button[kind="primary"] *,
        button[kind="secondary"][class*="stFormSubmitButton"], button[kind="secondary"][class*="stFormSubmitButton"] * {
            color: #ffffff !important;
        }

        /* ── Cuadro de escribir al Asistente IA (chat_input): en móvil,
           al seguir el modo claro/oscuro del sistema, si el teléfono está
           en oscuro este cuadro se queda con fondo negro y letra negra —
           ilegible. Se fuerza aquí un fondo claro y letra oscura siempre,
           tanto en el propio cuadro como en el texto que se escribe. */
        [data-testid="stChatInput"],
        [data-testid="stChatInputContainer"],
        [data-testid="stBottomBlockContainer"],
        [data-testid="stChatFloatingInputContainer"] {
            background: #ffffff !important;
        }
        [data-testid="stChatInput"] textarea,
        [data-testid="stChatInputContainer"] textarea {
            background: #ffffff !important;
            color: #1f2937 !important;
            caret-color: #1f2937 !important;
        }
        [data-testid="stChatInput"] textarea::placeholder,
        [data-testid="stChatInputContainer"] textarea::placeholder {
            color: #94a3b8 !important;
        }
        [data-testid="stChatInput"] button svg,
        [data-testid="stChatInputContainer"] button svg {
            color: #1e3a5f !important;
            fill: #1e3a5f !important;
        }

        /* ── Permitir copiar/pegar el texto (incluidas las respuestas
           del asistente de IA), por si algún navegador móvil lo bloquea */
        [data-testid="stAppViewContainer"], [data-testid="stAppViewContainer"] * {
            -webkit-user-select: text !important;
            user-select: text !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def _forzar_scroll_arriba_del_todo():
    """Fuerza que la página vuelva al principio (arriba del todo) justo después de un login.
    Sin esto, en móvil la app a veces se queda anclada abajo del todo (por el cuadro fijo del
    asistente de IA) y el inversor ve el chat antes que su bienvenida y sus datos."""
    components.html(
        """
        <script>
        (function() {
            function irArriba() {
                try {
                    var doc = window.parent.document;
                    doc.documentElement.scrollTop = 0;
                    doc.body.scrollTop = 0;
                    var contenedor = doc.querySelector('[data-testid="stAppViewContainer"]');
                    if (contenedor) { contenedor.scrollTop = 0; }
                    window.parent.scrollTo(0, 0);
                } catch (e) {}
            }
            // Se repite varias veces porque en móvil el layout (teclado, barras del navegador)
            // puede seguir moviéndose un instante después de la primera carga.
            irArriba();
            setTimeout(irArriba, 100);
            setTimeout(irArriba, 350);
            setTimeout(irArriba, 700);
        })();
        </script>
        """,
        height=0,
    )


def _forzar_scroll_a_ultimo_mensaje_chat():
    """Lleva la vista hasta el último mensaje del chat del asistente (la pregunta que se acaba
    de hacer y su respuesta), para que el inversor no tenga que subir manualmente a verla."""
    components.html(
        """
        <script>
        (function() {
            function irAlUltimoMensaje() {
                try {
                    var doc = window.parent.document;
                    var mensajes = doc.querySelectorAll('[data-testid="stChatMessage"]');
                    if (mensajes.length > 0) {
                        mensajes[mensajes.length - 1].scrollIntoView({behavior: "smooth", block: "start"});
                    }
                } catch (e) {}
            }
            setTimeout(irAlUltimoMensaje, 120);
            setTimeout(irAlUltimoMensaje, 400);
        })();
        </script>
        """,
        height=0,
    )


def mostrar_hero(usuario=None):
    tag = f"Sesión: {usuario}" if usuario else "Private Wealth Dashboard"
    st.markdown(
        f"""
        <div class="brand-hero">
            <div class="brand-left">
                <div class="brand-logo">CF</div>
                <div>
                    <div class="brand-title">Chaparro Fernández Wealth</div>
                </div>
            </div>
            <div class="brand-tag">{tag}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


# Excepciones puntuales: inversores que, desde su propio portal, también pueden consultar la
# posición de otro inversor concreto (a petición expresa de Yuri). El titular sigue viendo
# SIEMPRE primero su propia información — esto solo añade acceso extra, nunca lo quita.
INVERSORES_ADICIONALES_VISIBLES = {
    "ROBERTO VISCAFE": ["CROWE BOLIVIA", "JR REAL ESTATE"],
    "JORDI CHAPARRO": ["PAM", "JEP", "JR REAL ESTATE", "2012 JACC GROUP", "JORDI ESPECIAL"],
    "JAPAN JORDI": ["JEP"],
    "EVA CHAPARRO": ["JEP"],
    "PAOLA CHAPARRO": ["JEP"],
    "PEDRO MAGAÑA": ["2012 JACC GROUP", "PAM"],
}

# Excepción para el equipo interno (admin): además de su panel de administración, algunos
# también quieren poder ver un perfil de inversor concreto sin cerrar sesión y volver a entrar
# por el portal de inversor — cambiando de vista con un simple selector en la barra lateral.
# Clave = usuario admin (tal cual está en USUARIOS, sin distinguir mayúsculas/minúsculas),
# valor = lista de nombres de inversor (tal cual en INVERSIONES) que puede ver desde su sesión.
# EN PRUEBAS: solo con Yuri por ahora, viendo el perfil de PAM. Si funciona bien, se añade Jordi.
ADMIN_VISTAS_INVERSOR_ADICIONALES = {
    "YURI": ["PAM"],
}

# Inversores piloto con verificación en dos pasos por email disponible en su portal.
INVERSORES_CON_2FA = {"EVA CHAPARRO", "JORDI CHAPARRO", "PEDRO MAGAÑA"}

HOJA_USUARIOS = "USUARIOS"
HOJA_LOG_IA = "LOG_IA_USO"
HOJA_GASTOS_PLATAFORMA = "GASTOS_PLATAFORMA"
# Precio por millón de tokens (USD) de cada modelo de IA usado en la app, para estimar el coste
# de cada llamada al guardar el log de uso. Si cambia la tarifa de Anthropic o el modelo usado,
# actualiza este diccionario — es la ÚNICA fuente del cálculo de coste en toda la app.
PRECIOS_MODELOS_IA = {
    "claude-sonnet-4-5": (3.0, 15.0),  # (precio input, precio output) por millón de tokens
}
PRECIO_MODELO_IA_POR_DEFECTO = (3.0, 15.0)
# NOTA DE SEGURIDAD: las contraseñas se guardan hasheadas con bcrypt en la hoja USUARIOS del
# Excel del fondo (ver _hash_password / _verificar_password más abajo). Las filas antiguas que
# aún estuvieran en texto plano se migran a hash de forma silenciosa y transparente en el
# primer login correcto (ver _verificar_credencial). Nadie, ni con acceso al Excel, puede leer
# la contraseña real de un usuario ya migrado — solo comprobar si una contraseña dada coincide.


def _descargar_excel_para_credenciales():
    """Se asegura de que el Excel esté en disco para poder leer/actualizar la hoja USUARIOS,
    sin depender de funciones definidas más abajo en el archivo (el login se ejecuta antes)."""
    import os as _os_cred
    if _os_cred.path.exists(ARCHIVO):
        return
    try:
        url = f"https://docs.google.com/spreadsheets/d/{GDRIVE_FILE_ID}/export?format=xlsx"
        r = requests.get(url, timeout=30)
        if r.status_code == 200:
            with open(ARCHIVO, "wb") as f:
                f.write(r.content)
    except Exception:
        pass


def _texto_seguro_excel(valor) -> str:
    """Convierte un valor leído de Excel a texto de forma segura: si Excel lo interpretó como
    número entero (float 654321.0 en vez de texto '654321'), quita el '.0' sobrante."""
    if pd.isna(valor):
        return ""
    if isinstance(valor, float) and valor.is_integer():
        return str(int(valor))
    return str(valor)


def _leer_hoja_usuarios() -> pd.DataFrame:
    """Lee la hoja USUARIOS (usuario, tipo_usuario, password, debe_cambiar_password, totp_secret,
    totp_activo, email) del Excel del fondo. Si la hoja aún no existe (primera vez), devuelve un
    DataFrame vacío con las columnas correctas y el sistema cae automáticamente en las
    contraseñas iniciales del código."""
    _descargar_excel_para_credenciales()
    try:
        df = pd.read_excel(ARCHIVO, sheet_name=HOJA_USUARIOS)
        df.columns = [str(c).strip().lower() for c in df.columns]
        for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"]:
            if col in df.columns:
                df[col] = df[col].apply(_texto_seguro_excel)
        if "debe_cambiar_password" not in df.columns:
            df["debe_cambiar_password"] = "NO"
        if "totp_secret" not in df.columns:
            df["totp_secret"] = ""
        if "totp_activo" not in df.columns:
            df["totp_activo"] = "NO"
        if "email" not in df.columns:
            # Columna nueva (email de contacto para "olvidé mi contraseña", separada de
            # totp_secret que solo se rellena si el usuario activó la verificación en dos
            # pasos) — para filas ya existentes, se rellena sola con el email de 2FA si lo
            # tenían, así no se pierde nada al añadir la columna.
            df["email"] = df["totp_secret"] if "totp_secret" in df.columns else ""
        return df
    except Exception:
        return pd.DataFrame(columns=["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"])


def _guardar_hoja_usuarios(df_usuarios: pd.DataFrame) -> tuple[bool, str]:
    """Guarda la hoja USUARIOS actualizada, preservando intactas todas las demás hojas del
    Excel, y la sube a Google Drive si hay credenciales de servicio configuradas."""
    import os as _os_cred
    try:
        if _os_cred.path.exists(ARCHIVO):
            hojas = pd.read_excel(ARCHIVO, sheet_name=None)
        else:
            hojas = {}
    except Exception:
        hojas = {}
    hojas[HOJA_USUARIOS] = df_usuarios

    salida = BytesIO()
    with pd.ExcelWriter(salida, engine="openpyxl") as writer:
        for nombre_hoja, df in hojas.items():
            nombre_limpio = str(nombre_hoja)[:31] if str(nombre_hoja).strip() else "Hoja"
            (df if df is not None else pd.DataFrame()).to_excel(writer, sheet_name=nombre_limpio, index=False)
    contenido = salida.getvalue()
    with open(ARCHIVO, "wb") as f:
        f.write(contenido)

    try:
        from postgres_writer import sincronizar_usuarios_postgres
        sincronizar_usuarios_postgres(df_usuarios)
    except Exception:
        pass

    if "gcp_service_account" not in st.secrets:
        return False, "No hay credenciales de Google configuradas (falta [gcp_service_account] en Secrets) — el cambio no se sincronizó con Drive y se perderá si el servidor se reinicia."
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseUpload

        credenciales = service_account.Credentials.from_service_account_info(
            dict(st.secrets["gcp_service_account"]),
            scopes=["https://www.googleapis.com/auth/drive"],
        )
        servicio = build("drive", "v3", credentials=credenciales)
        media = MediaIoBaseUpload(
            BytesIO(contenido),
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            resumable=False,
        )
        servicio.files().update(fileId=GDRIVE_FILE_ID, media_body=media, fields="id").execute()
        return True, "Guardada y sincronizada con Google Drive."
    except Exception as e:
        return False, f"Se guardó localmente pero no se pudo sincronizar con Drive: {e}"


# =========================
# PROTECCIÓN CONTRA FUERZA BRUTA EN EL LOGIN
# =========================
# Estado compartido entre todas las sesiones del mismo proceso (cacheado con st.cache_resource
# para que sobreviva a los reruns de Streamlit, que si no reinician cualquier variable de módulo
# en cada interacción). No es persistente entre redeploys de Railway — es una protección básica
# en memoria, suficiente para frenar intentos automatizados de adivinar contraseñas.
MAX_INTENTOS_LOGIN = 5
BLOQUEO_LOGIN_SEGUNDOS = 15 * 60  # 15 minutos de bloqueo tras agotar los intentos


@st.cache_resource
def _estado_intentos_login() -> dict:
    return {}


def _clave_intentos(tipo: str, usuario: str) -> str:
    return f"{tipo}:{usuario.strip().lower()}"


def _login_bloqueado(tipo: str, usuario: str) -> tuple[bool, int]:
    """Devuelve (bloqueado, minutos_restantes) para este usuario+tipo."""
    estado = _estado_intentos_login()
    clave = _clave_intentos(tipo, usuario)
    intentos, bloqueado_hasta = estado.get(clave, (0, 0.0))
    ahora = time.time()
    if bloqueado_hasta and ahora < bloqueado_hasta:
        return True, int((bloqueado_hasta - ahora) // 60) + 1
    if bloqueado_hasta and ahora >= bloqueado_hasta:
        estado[clave] = (0, 0.0)  # el bloqueo ya expiró: reseteamos contador
    return False, 0


def _registrar_intento_fallido(tipo: str, usuario: str):
    estado = _estado_intentos_login()
    clave = _clave_intentos(tipo, usuario)
    intentos, _bloqueado_hasta = estado.get(clave, (0, 0.0))
    intentos += 1
    bloqueado_hasta = time.time() + BLOQUEO_LOGIN_SEGUNDOS if intentos >= MAX_INTENTOS_LOGIN else 0.0
    estado[clave] = (intentos, bloqueado_hasta)


def _resetear_intentos_login(tipo: str, usuario: str):
    estado = _estado_intentos_login()
    estado.pop(_clave_intentos(tipo, usuario), None)


def _hash_password(password_plano: str) -> str:
    """Genera un hash bcrypt de una contraseña en texto plano, listo para guardar."""
    return bcrypt.hashpw(password_plano.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")


def _generar_password_temporal(longitud: int = 14) -> str:
    """Genera una contraseña temporal aleatoria y criptográficamente segura (módulo `secrets`,
    no `random`). Se usa una sola vez para enviarla por email; nunca se guarda en texto plano
    en ningún sitio — solo su hash bcrypt."""
    import secrets as _secrets
    import string as _string
    alfabeto = _string.ascii_letters + _string.digits + "!@#$%&*-_"
    while True:
        pw = "".join(_secrets.choice(alfabeto) for _ in range(longitud))
        # Nos aseguramos de que tenga al menos una mayúscula, una minúscula, un dígito y un símbolo.
        if (any(c.islower() for c in pw) and any(c.isupper() for c in pw)
                and any(c.isdigit() for c in pw) and any(c in "!@#$%&*-_" for c in pw)):
            return pw


def _debe_cambiar_password(usuario: str, tipo: str) -> bool:
    """Comprueba si este usuario tiene pendiente el cambio obligatorio de contraseña (recién
    creado con una contraseña temporal enviada por email, aún sin personalizar)."""
    df_u = _leer_hoja_usuarios()
    if df_u.empty or not {"usuario", "tipo_usuario", "debe_cambiar_password"}.issubset(df_u.columns):
        return False
    fila = df_u[
        (df_u["usuario"].astype(str).str.strip().str.lower() == usuario.strip().lower())
        & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
    ]
    if fila.empty:
        return False
    return str(fila.iloc[0].get("debe_cambiar_password", "NO")).strip().upper() == "SI"


def _debe_cambiar_password_login(usuario: str, tipo: str) -> bool:
    """Como _debe_cambiar_password, pero además obliga el cambio en el PRIMER login de una
    cuenta de equipo interno (admin) que todavía no tiene fila propia en USUARIOS — es decir,
    que sigue entrando con la contraseña temporal definida en el código. Los inversores no se
    ven afectados por esta regla adicional: solo se obliga el cambio si su fila lo marca."""
    fila = _fila_usuario(_leer_hoja_usuarios(), usuario, tipo)
    if fila is None:
        return tipo == "admin"
    return str(fila.get("debe_cambiar_password", "NO")).strip().upper() == "SI"


# =========================
# VERIFICACIÓN EN DOS PASOS (TOTP — Google Authenticator / Authy)
# =========================
def _fila_usuario(df_u: pd.DataFrame, usuario: str, tipo: str):
    """Devuelve la fila (Series) de un usuario+tipo en la hoja USUARIOS, o None si no existe."""
    if df_u.empty or not {"usuario", "tipo_usuario"}.issubset(df_u.columns):
        return None
    fila = df_u[
        (df_u["usuario"].astype(str).str.strip().str.lower() == usuario.strip().lower())
        & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
    ]
    return fila.iloc[0] if not fila.empty else None


def _2fa_activo(usuario: str, tipo: str) -> bool:
    """Comprueba si este usuario tiene la verificación en dos pasos (código por email) activada."""
    fila = _fila_usuario(_leer_hoja_usuarios(), usuario, tipo)
    if fila is None:
        return False
    return str(fila.get("totp_activo", "NO")).strip().upper() == "SI"


def _2fa_email_de(usuario: str, tipo: str) -> str:
    """Devuelve el email al que se envía el código de verificación (reutiliza la columna
    'totp_secret' de la hoja USUARIOS, que ahora guarda el email en vez de un secreto TOTP)."""
    fila = _fila_usuario(_leer_hoja_usuarios(), usuario, tipo)
    if fila is None:
        return ""
    return str(fila.get("totp_secret", "") or "")


def _email_contacto_de(usuario: str, tipo: str) -> str:
    """Email 'oficial' de contacto de un usuario, usado para autoservicio (olvidé mi
    contraseña). Prioridad: columna 'email' de USUARIOS y, si está vacía, el email de la
    verificación en dos pasos (columna 'totp_secret') — por si el usuario tiene 2FA activado
    pero, por lo que sea, la columna 'email' específica no se llegó a rellenar."""
    fila = _fila_usuario(_leer_hoja_usuarios(), usuario, tipo)
    if fila is None:
        return ""
    email = str(fila.get("email", "") or "").strip()
    if email:
        return email
    return str(fila.get("totp_secret", "") or "").strip()


@st.cache_resource
def _estado_codigos_reset_password() -> dict:
    """Códigos de un solo uso para 'olvidé mi contraseña', en memoria (no en Excel — son
    efímeros, caducan solos a los 15 minutos). Compartido entre reruns con st.cache_resource,
    igual que el estado de los códigos de 2FA (son dos flujos independientes a propósito: pedir
    un código de reseteo no debe interferir con un login en curso con 2FA pendiente, y viceversa)."""
    return {}


def _clave_reset_password(usuario: str, tipo: str) -> str:
    return f"{tipo}:{usuario.strip().lower()}"


def _generar_y_enviar_codigo_reset_password(usuario: str, tipo: str, email_destino: str,
                                             smtp_sender: str, smtp_password: str,
                                             display_name: str = "Chaparro Fernández Wealth") -> tuple:
    """Genera un código de 6 dígitos para restablecer la contraseña, lo guarda en memoria con
    caducidad de 15 minutos, y lo envía por email a la dirección de contacto de este usuario."""
    import secrets as _secrets
    codigo = f"{_secrets.randbelow(1_000_000):06d}"
    estado = _estado_codigos_reset_password()
    estado[_clave_reset_password(usuario, tipo)] = (codigo, time.time() + 15 * 60)
    try:
        destinatarios = _parse_lista_emails(email_destino)
        if not destinatarios:
            return False, "El email registrado no es válido."
        msg = MIMEMultipart("mixed")
        msg["Subject"] = f"Restablecer tu contraseña — {display_name}"
        msg["From"] = f"{display_name} <{smtp_sender}>"
        msg["To"] = ", ".join(destinatarios)
        cuerpo = f"""<!DOCTYPE html><html lang="es"><head><meta charset="UTF-8"></head>
<body style="margin:0;padding:0;background:#f4f6f9;font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
  <table width="100%" cellpadding="0" cellspacing="0" style="background:#f4f6f9;padding:40px 20px;">
    <tr><td align="center">
      <table width="460" cellpadding="0" cellspacing="0" style="background:#ffffff;border-radius:20px;overflow:hidden;box-shadow:0 8px 40px rgba(7,20,37,0.14);">
        <tr><td style="background:linear-gradient(135deg,#0e2338 0%,#173b5c 60%,#bf9a5f 100%);padding:28px 40px;text-align:center;">
          <div style="color:#ffffff;font-size:18px;font-weight:800;">{display_name}</div>
        </td></tr>
        <tr><td style="padding:32px 40px;text-align:center;">
          <p style="color:#334155;font-size:14px;margin:0 0 8px;">Alguien (esperamos que tú, {usuario}) ha pedido restablecer la contraseña de esta cuenta.</p>
          <p style="color:#334155;font-size:14px;margin:0 0 18px;">Tu código para elegir una nueva contraseña es:</p>
          <div style="font-size:36px;font-weight:800;letter-spacing:8px;color:#0e2338;font-family:monospace;margin-bottom:18px;">{codigo}</div>
          <p style="color:#94a3b8;font-size:12px;">Caduca en 15 minutos. Si no has sido tú, ignora este correo — tu contraseña actual sigue siendo válida y no se ha cambiado nada.</p>
        </td></tr>
      </table>
    </td></tr>
  </table>
</body></html>"""
        msg.attach(MIMEText(cuerpo, "html", "utf-8"))
        contexto = ssl.create_default_context()
        with smtplib.SMTP("smtp.gmail.com", 587) as servidor:
            servidor.ehlo()
            servidor.starttls(context=contexto)
            servidor.login(smtp_sender, smtp_password)
            servidor.sendmail(smtp_sender, destinatarios, msg.as_bytes())
        return True, ""
    except smtplib.SMTPAuthenticationError:
        return False, "Error de autenticación Gmail. Usa una contraseña de aplicación (no tu contraseña normal)."
    except smtplib.SMTPRecipientsRefused:
        return False, "El email registrado fue rechazado por el servidor de correo."
    except Exception as e:
        return False, str(e)


def _verificar_codigo_reset_password(usuario: str, tipo: str, codigo_input: str) -> bool:
    """Comprueba el código sin consumirlo todavía (se consume aparte, tras guardar la nueva
    contraseña con éxito — así un fallo al guardar en Drive no obliga a pedir un código nuevo)."""
    estado = _estado_codigos_reset_password()
    clave = _clave_reset_password(usuario, tipo)
    par = estado.get(clave)
    if not par:
        return False
    codigo_guardado, expira = par
    if time.time() > expira:
        estado.pop(clave, None)
        return False
    return (codigo_input or "").strip() == codigo_guardado


def _consumir_codigo_reset_password(usuario: str, tipo: str):
    _estado_codigos_reset_password().pop(_clave_reset_password(usuario, tipo), None)


def _fijar_password_autoservicio(usuario: str, tipo: str, nueva_password: str) -> tuple[bool, str]:
    """Guarda la nueva contraseña elegida por el propio usuario tras verificar el código de
    reseteo — crea su fila en USUARIOS si todavía no existía (usuario que seguía con la
    contraseña inicial del código), o la actualiza si ya existía, preservando su email, y su
    2FA (secreto/estado) intactos igual que hace formulario_cambiar_password. Quita el flag de
    cambio obligatorio: el usuario acaba de elegir la contraseña él mismo, no hace falta
    pedírsela otra vez en el próximo login."""
    df_u = _leer_hoja_usuarios()
    if df_u.empty or not {"usuario", "password", "tipo_usuario"}.issubset(df_u.columns):
        df_u = pd.DataFrame(columns=["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"])
    for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"]:
        if col not in df_u.columns:
            df_u[col] = "NO" if col in ("debe_cambiar_password", "totp_activo") else ""
        df_u[col] = df_u[col].astype(object)
    _fila_previa = _fila_usuario(df_u, usuario, tipo)
    _totp_secret_prev = str(_fila_previa.get("totp_secret", "") or "") if _fila_previa is not None else ""
    _totp_activo_prev = str(_fila_previa.get("totp_activo", "NO") or "NO") if _fila_previa is not None else "NO"
    _email_prev = str(_fila_previa.get("email", "") or "") if _fila_previa is not None else ""
    mascara = (
        (df_u["usuario"].astype(str).str.strip().str.lower() == usuario.strip().lower())
        & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
    )
    df_u = df_u[~mascara]
    fila_nueva = pd.DataFrame([{
        "usuario": usuario, "tipo_usuario": tipo,
        "password": _hash_password(nueva_password), "debe_cambiar_password": "NO",
        "totp_secret": _totp_secret_prev, "totp_activo": _totp_activo_prev, "email": _email_prev,
    }])
    df_u = pd.concat([df_u, fila_nueva], ignore_index=True)
    return _guardar_hoja_usuarios(df_u)


@st.cache_resource
def _estado_codigos_2fa_email() -> dict:
    """Códigos de un solo uso pendientes de verificar, en memoria (no en Excel — son efímeros,
    caducan solos a los 10 minutos). Compartido entre reruns gracias a st.cache_resource."""
    return {}


def _clave_2fa(usuario: str, tipo: str) -> str:
    return f"{tipo}:{usuario.strip().lower()}"


def _generar_y_enviar_codigo_2fa(usuario: str, tipo: str, smtp_sender: str, smtp_password: str,
                                  display_name: str = "Chaparro Fernández Wealth") -> tuple:
    """Genera un código de 6 dígitos, lo guarda en memoria con caducidad de 10 minutos, y lo
    envía por email a la dirección configurada para este usuario."""
    import secrets as _secrets
    email_destino = _2fa_email_de(usuario, tipo)
    if not email_destino:
        return False, "No hay ningún email configurado para la verificación en dos pasos."
    codigo = f"{_secrets.randbelow(1_000_000):06d}"
    estado = _estado_codigos_2fa_email()
    estado[_clave_2fa(usuario, tipo)] = (codigo, time.time() + 10 * 60)
    try:
        destinatarios = _parse_lista_emails(email_destino)
        if not destinatarios:
            return False, "El email configurado no es válido."
        msg = MIMEMultipart("mixed")
        msg["Subject"] = f"Tu código de acceso — {display_name}"
        msg["From"] = f"{display_name} <{smtp_sender}>"
        msg["To"] = ", ".join(destinatarios)
        cuerpo = f"""<!DOCTYPE html><html lang="es"><head><meta charset="UTF-8"></head>
<body style="margin:0;padding:0;background:#f4f6f9;font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
  <table width="100%" cellpadding="0" cellspacing="0" style="background:#f4f6f9;padding:40px 20px;">
    <tr><td align="center">
      <table width="460" cellpadding="0" cellspacing="0" style="background:#ffffff;border-radius:20px;overflow:hidden;box-shadow:0 8px 40px rgba(7,20,37,0.14);">
        <tr><td style="background:linear-gradient(135deg,#0e2338 0%,#173b5c 60%,#bf9a5f 100%);padding:28px 40px;text-align:center;">
          <div style="color:#ffffff;font-size:18px;font-weight:800;">{display_name}</div>
        </td></tr>
        <tr><td style="padding:32px 40px;text-align:center;">
          <p style="color:#334155;font-size:14px;margin:0 0 18px;">Tu código de acceso es:</p>
          <div style="font-size:36px;font-weight:800;letter-spacing:8px;color:#0e2338;font-family:monospace;margin-bottom:18px;">{codigo}</div>
          <p style="color:#94a3b8;font-size:12px;">Caduca en 10 minutos. Si no has sido tú quien ha intentado entrar, ignora este correo.</p>
        </td></tr>
      </table>
    </td></tr>
  </table>
</body></html>"""
        msg.attach(MIMEText(cuerpo, "html", "utf-8"))
        contexto = ssl.create_default_context()
        with smtplib.SMTP("smtp.gmail.com", 587) as servidor:
            servidor.ehlo()
            servidor.starttls(context=contexto)
            servidor.login(smtp_sender, smtp_password)
            servidor.sendmail(smtp_sender, destinatarios, msg.as_bytes())
        return True, email_destino
    except Exception as e:
        return False, str(e)


def _verificar_codigo_2fa(usuario: str, tipo: str, codigo_input: str) -> bool:
    estado = _estado_codigos_2fa_email()
    clave = _clave_2fa(usuario, tipo)
    par = estado.get(clave)
    if not par:
        return False
    codigo_guardado, expira = par
    if time.time() > expira:
        estado.pop(clave, None)
        return False
    if (codigo_input or "").strip() == codigo_guardado:
        estado.pop(clave, None)
        return True
    return False


def _guardar_2fa(usuario: str, tipo: str, email: str, activo: bool):
    """Activa/desactiva la verificación en dos pasos de un usuario, preservando el resto de sus
    datos (password, debe_cambiar_password) intactos — nunca se tocan al activar/desactivar."""
    df_u = _leer_hoja_usuarios()
    for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo"]:
        if col not in df_u.columns:
            df_u[col] = "" if col in ("totp_secret",) else ("NO" if col in ("debe_cambiar_password", "totp_activo") else "")
        df_u[col] = df_u[col].astype(object)
    idx = df_u[
        (df_u["usuario"].astype(str).str.strip().str.lower() == usuario.strip().lower())
        & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
    ].index
    df_u.loc[idx, "totp_secret"] = email
    df_u.loc[idx, "totp_activo"] = "SI" if activo else "NO"
    return _guardar_hoja_usuarios(df_u)


def _completar_login(usuario_match: str, tipo: str):
    """Centraliza qué pasa cuando un login se da por bueno (con o sin segundo factor de por
    medio): arranca la sesión, calcula los avisos pendientes y limpia el estado de TOTP."""
    st.session_state.autenticado = True
    st.session_state.usuario = usuario_match
    st.session_state.tipo_usuario = tipo
    st.session_state.login_timestamp = time.time()
    st.session_state.ultima_actividad = time.time()
    if tipo == "admin":
        _df_u_check = _leer_hoja_usuarios()
        _ya_migrado = _fila_usuario(_df_u_check, usuario_match, "admin") is not None
        st.session_state.mostrar_aviso_pw_temporal = not _ya_migrado
    else:
        st.session_state.mostrar_aviso_pw_temporal = False
    st.session_state.forzar_cambio_password = _debe_cambiar_password_login(usuario_match, tipo)
    st.session_state.totp_pendiente = None
    # Al entrar recién logueado, la vista debe arrancar arriba del todo (bienvenida, capital,
    # etc.) y no abajo donde vive el cuadro fijo del asistente de IA.
    st.session_state.forzar_scroll_arriba_pendiente = True


def enviar_email_credenciales_nuevas(destinatario: str, usuario: str, password_temporal: str,
                                      tipo: str, smtp_sender: str, smtp_password: str,
                                      display_name: str = "Chaparro Fernández Wealth") -> tuple:
    """Envía por email la contraseña temporal de un acceso recién creado. Esta es la ÚNICA vez
    que la contraseña en texto plano existe en algún sitio (en este correo) — ni se muestra en
    pantalla al admin que la crea, ni se guarda así en ningún fichero; solo se persiste su hash
    bcrypt en la hoja USUARIOS."""
    destinatarios = _parse_lista_emails(destinatario)
    if not destinatarios:
        return False, "No hay ninguna dirección de email válida."
    portal_str = "equipo interno" if tipo == "admin" else "portal de inversor"
    try:
        msg = MIMEMultipart("mixed")
        msg["Subject"] = f"Tu acceso a {display_name}"
        msg["From"] = f"{display_name} <{smtp_sender}>"
        msg["To"] = ", ".join(destinatarios)
        cuerpo = f"""<!DOCTYPE html>
<html lang="es"><head><meta charset="UTF-8"></head>
<body style="margin:0;padding:0;background:#f4f6f9;font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
  <table width="100%" cellpadding="0" cellspacing="0" style="background:#f4f6f9;padding:40px 20px;">
    <tr><td align="center">
      <table width="560" cellpadding="0" cellspacing="0" style="background:#ffffff;border-radius:20px;overflow:hidden;box-shadow:0 8px 40px rgba(7,20,37,0.14);">
        <tr><td style="background:linear-gradient(135deg,#0e2338 0%,#173b5c 60%,#bf9a5f 100%);padding:32px 40px;text-align:center;">
          <div style="color:#ffffff;font-size:20px;font-weight:800;">{display_name}</div>
          <div style="color:rgba(255,255,255,0.75);font-size:13px;margin-top:6px;">Acceso al {portal_str}</div>
        </td></tr>
        <tr><td style="padding:32px 40px;">
          <p style="color:#334155;font-size:14px;line-height:1.6;">Se ha creado un acceso para ti. Estos son tus datos de entrada:</p>
          <table width="100%" cellpadding="10" cellspacing="0" style="background:#f5f0e8;border-radius:10px;margin:18px 0;">
            <tr><td style="font-size:13px;color:#64748b;">Usuario</td><td style="font-size:14px;color:#0e2338;font-weight:700;">{usuario}</td></tr>
            <tr><td style="font-size:13px;color:#64748b;">Contraseña temporal</td><td style="font-size:14px;color:#0e2338;font-weight:700;font-family:monospace;">{password_temporal}</td></tr>
          </table>
          <p style="color:#475569;font-size:13px;line-height:1.6;">Por seguridad, <strong>se te pedirá cambiarla</strong> obligatoriamente en cuanto inicies sesión por primera vez. Esta contraseña temporal solo existe en este correo — nadie más, ni siquiera el equipo de administración, puede volver a consultarla.</p>
          <p style="color:#94a3b8;font-size:12px;margin-top:22px;">🔒 Correo confidencial. Si no esperabas este acceso, ignóralo y contacta con nosotros.</p>
        </td></tr>
      </table>
    </td></tr>
  </table>
</body></html>"""
        msg.attach(MIMEText(cuerpo, "html", "utf-8"))
        contexto = ssl.create_default_context()
        with smtplib.SMTP("smtp.gmail.com", 587) as servidor:
            servidor.ehlo()
            servidor.starttls(context=contexto)
            servidor.login(smtp_sender, smtp_password)
            servidor.sendmail(smtp_sender, destinatarios, msg.as_bytes())
        return True, ""
    except smtplib.SMTPAuthenticationError:
        return False, "Error de autenticación Gmail. Usa una contraseña de aplicación (no tu contraseña normal)."
    except smtplib.SMTPRecipientsRefused:
        return False, f"Alguna(s) de estas direcciones fue(ron) rechazada(s): {', '.join(destinatarios)}."
    except Exception as e:
        return False, str(e)


def _es_hash_bcrypt(valor: str) -> bool:
    """Los hashes bcrypt siempre empiezan por $2a$, $2b$ o $2y$ — así distinguimos una
    contraseña ya hasheada de una todavía en texto plano (fila antigua sin migrar)."""
    return isinstance(valor, str) and valor.startswith(("$2a$", "$2b$", "$2y$"))


def _verificar_password(password_input: str, valor_guardado: str) -> bool:
    """Compara una contraseña introducida contra lo guardado, sea hash bcrypt o texto plano
    (para las filas todavía no migradas)."""
    valor_guardado = str(valor_guardado)
    if _es_hash_bcrypt(valor_guardado):
        try:
            return bcrypt.checkpw(password_input.encode("utf-8"), valor_guardado.encode("utf-8"))
        except ValueError:
            return False
    return password_input == valor_guardado


def _verificar_credencial(usuario_input: str, password_input: str, tipo: str, usuarios_codigo: dict) -> Optional[str]:
    """Comprueba usuario+contraseña para tipo 'admin' o 'inversor'.
    Prioridad: hoja USUARIOS en Drive (contraseñas que el propio usuario ha cambiado).
    Si el usuario no está migrado ahí todavía, cae a las contraseñas iniciales del código.
    Devuelve el nombre canónico del usuario si es correcto, o None.

    Las contraseñas de la hoja USUARIOS se guardan hasheadas con bcrypt (ya no en texto
    plano). Para no romper el acceso de nadie que tuviera la contraseña antigua en texto
    plano, si encontramos una fila todavía sin migrar y la contraseña es correcta, la
    volvemos a guardar ya hasheada en el mismo momento (migración silenciosa y transparente)."""
    df_u = _leer_hoja_usuarios()
    if not df_u.empty and {"usuario", "password", "tipo_usuario"}.issubset(df_u.columns):
        fila = df_u[
            (df_u["usuario"].astype(str).str.strip().str.lower() == usuario_input.strip().lower())
            & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
        ]
        if not fila.empty:
            fila = fila.iloc[0]
            valor_guardado = str(fila["password"])
            if _verificar_password(password_input, valor_guardado):
                if not _es_hash_bcrypt(valor_guardado):
                    try:
                        df_u.loc[
                            (df_u["usuario"].astype(str).str.strip().str.lower() == usuario_input.strip().lower())
                            & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo),
                            "password",
                        ] = _hash_password(password_input)
                        _guardar_hoja_usuarios(df_u)
                    except Exception:
                        pass  # si la migración silenciosa falla, el login ya validó igualmente
                return str(fila["usuario"])
            return None  # ya migrado a Drive: la contraseña del código queda obsoleta para este usuario

    match = next((u for u in usuarios_codigo if u.strip().lower() == usuario_input.strip().lower()), None)
    if match and _verificar_password(password_input, usuarios_codigo.get(match, "")):
        return match
    return None


def formulario_cambiar_password(usuario_actual: str, tipo: str, usuarios_codigo: dict):
    """Formulario de autoservicio para que cualquier usuario (equipo interno o inversor)
    cambie su propia contraseña, sin que haga falta tocar el código."""
    with st.sidebar.expander("🔑 Cambiar mi contraseña"):
        with st.form(f"form_cambiar_pw_{tipo}_{usuario_actual}"):
            pw_actual = st.text_input("Contraseña actual", type="password", key=f"pw_actual_{tipo}")
            pw_nueva = st.text_input("Nueva contraseña", type="password", key=f"pw_nueva_{tipo}")
            pw_nueva2 = st.text_input("Repite la nueva contraseña", type="password", key=f"pw_nueva2_{tipo}")
            enviar = st.form_submit_button("Actualizar contraseña")
        if enviar:
            if _verificar_credencial(usuario_actual, pw_actual, tipo, usuarios_codigo) is None:
                st.error("La contraseña actual no es correcta.")
            elif len(pw_nueva) < 6:
                st.error("La nueva contraseña debe tener al menos 6 caracteres.")
            elif pw_nueva != pw_nueva2:
                st.error("Las dos contraseñas nuevas no coinciden.")
            else:
                df_u = _leer_hoja_usuarios()
                if df_u.empty or not {"usuario", "password", "tipo_usuario"}.issubset(df_u.columns):
                    df_u = pd.DataFrame(columns=["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"])
                # Todas las columnas como texto: si la hoja se creó vacía, pandas puede haberlas
                # inferido como float64 (todo NaN), y asignar un string ahí con .loc revienta con
                # TypeError en pandas 3.x (ya no hace upcast silencioso). Forzamos texto primero.
                for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"]:
                    if col not in df_u.columns:
                        df_u[col] = "NO" if col in ("debe_cambiar_password", "totp_activo") else ""
                    df_u[col] = df_u[col].astype(object)
                # Conservamos el secreto TOTP, si está activo, y el email de contacto — cambiar
                # la contraseña NUNCA debe desactivar la verificación en dos pasos ni borrar el
                # email registrado de nadie (lo necesita 'olvidé mi contraseña' más adelante).
                _fila_previa = _fila_usuario(df_u, usuario_actual, tipo)
                _totp_secret_prev = str(_fila_previa.get("totp_secret", "") or "") if _fila_previa is not None else ""
                _totp_activo_prev = str(_fila_previa.get("totp_activo", "NO") or "NO") if _fila_previa is not None else "NO"
                _email_prev = str(_fila_previa.get("email", "") or "") if _fila_previa is not None else ""
                mascara = (
                    (df_u["usuario"].astype(str).str.strip().str.lower() == usuario_actual.strip().lower())
                    & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
                )
                # En vez de asignar in-place con .loc (fuente del TypeError), quitamos la fila
                # antigua si existía y añadimos la fila nueva — evita por completo el problema
                # de tipos, igual de robusto y más simple.
                df_u = df_u[~mascara]
                fila_nueva = pd.DataFrame([{
                    "usuario": usuario_actual, "tipo_usuario": tipo,
                    "password": _hash_password(pw_nueva), "debe_cambiar_password": "NO",
                    "totp_secret": _totp_secret_prev, "totp_activo": _totp_activo_prev, "email": _email_prev,
                }])
                df_u = pd.concat([df_u, fila_nueva], ignore_index=True)
                exito, mensaje = _guardar_hoja_usuarios(df_u)
                if exito:
                    st.session_state.forzar_cambio_password = False
                    st.success(f"✅ Contraseña actualizada. {mensaje}")
                else:
                    st.warning(f"⚠️ Contraseña actualizada localmente. {mensaje}")


def formulario_email_recuperacion(usuario_actual: str, tipo: str):
    """Autoservicio para que cualquier usuario registre o actualice su propio email de
    recuperación — el que se usa para 'olvidé mi contraseña' en la pantalla de login. No exige
    la contraseña actual (a diferencia de cambiar la contraseña): quien ya inició sesión ya se
    autenticó, y pedir la contraseña otra vez aquí no añade seguridad real, solo fricción."""
    email_actual = _email_contacto_de(usuario_actual, tipo)
    etiqueta = f"📧 Mi email de recuperación ({'registrado' if email_actual else 'sin registrar'})"
    with st.sidebar.expander(etiqueta):
        if email_actual:
            _oculto = (email_actual[:2] + "***@" + email_actual.split("@")[-1]) if "@" in email_actual else email_actual
            st.caption(f"Email actual: **{_oculto}**")
        else:
            st.caption("Todavía no tenés un email registrado — sin él, no vas a poder usar '¿Olvidaste tu contraseña?' en el login.")
        with st.form(f"form_email_recuperacion_{tipo}_{usuario_actual}"):
            nuevo_email = st.text_input("Email", key=f"input_email_recuperacion_{tipo}")
            guardar_email = st.form_submit_button("Guardar email")
        if guardar_email:
            nuevo_email_limpio = (nuevo_email or "").strip()
            if not nuevo_email_limpio or "@" not in nuevo_email_limpio:
                st.error("Introduce un email válido.")
            else:
                df_u = _leer_hoja_usuarios()
                for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"]:
                    if col not in df_u.columns:
                        df_u[col] = "NO" if col in ("debe_cambiar_password", "totp_activo") else ""
                    df_u[col] = df_u[col].astype(object)
                if _fila_usuario(df_u, usuario_actual, tipo) is None:
                    st.error("Todavía no tenés una fila propia en el sistema — primero cambiá tu contraseña una vez en '🔑 Cambiar mi contraseña' (aunque sea por la misma), y después volvé a registrar tu email aquí.")
                else:
                    idx = df_u[
                        (df_u["usuario"].astype(str).str.strip().str.lower() == usuario_actual.strip().lower())
                        & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
                    ].index
                    df_u.loc[idx, "email"] = nuevo_email_limpio
                    exito, mensaje = _guardar_hoja_usuarios(df_u)
                    if exito:
                        st.success("✅ Email de recuperación guardado.")
                    else:
                        st.warning(f"⚠️ Guardado localmente. {mensaje}")


def formulario_cambio_obligatorio_password(usuario_actual: str, tipo: str):
    """Pantalla BLOQUEANTE (no se puede saltar) para usuarios recién creados con una contraseña
    temporal enviada por email. No pide la contraseña actual (ya la usaron para entrar) — solo
    obliga a fijar una nueva antes de dejar pasar al resto de la aplicación."""
    st.markdown(
        """
        <div class="login-card">
            <div class="login-logo">CF</div>
            <div class="login-title">Elige tu contraseña</div>
            <div class="login-subtitle">Por seguridad, debes personalizar tu contraseña temporal antes de continuar</div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    with st.form(f"form_cambio_obligatorio_{tipo}_{usuario_actual}"):
        pw_nueva = st.text_input("Nueva contraseña", type="password", key=f"pw_obl_nueva_{tipo}")
        pw_nueva2 = st.text_input("Repite la nueva contraseña", type="password", key=f"pw_obl_nueva2_{tipo}")
        enviar = st.form_submit_button("Guardar y continuar")
    if enviar:
        if len(pw_nueva) < 8:
            st.error("La nueva contraseña debe tener al menos 8 caracteres.")
        elif pw_nueva != pw_nueva2:
            st.error("Las dos contraseñas no coinciden.")
        else:
            df_u = _leer_hoja_usuarios()
            if df_u.empty or not {"usuario", "password", "tipo_usuario"}.issubset(df_u.columns):
                df_u = pd.DataFrame(columns=["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"])
            for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "totp_secret", "totp_activo", "email"]:
                if col not in df_u.columns:
                    df_u[col] = "NO" if col in ("debe_cambiar_password", "totp_activo") else ""
                df_u[col] = df_u[col].astype(object)
            _fila_previa = _fila_usuario(df_u, usuario_actual, tipo)
            _totp_secret_prev = str(_fila_previa.get("totp_secret", "") or "") if _fila_previa is not None else ""
            _totp_activo_prev = str(_fila_previa.get("totp_activo", "NO") or "NO") if _fila_previa is not None else "NO"
            _email_prev = str(_fila_previa.get("email", "") or "") if _fila_previa is not None else ""
            mascara = (
                (df_u["usuario"].astype(str).str.strip().str.lower() == usuario_actual.strip().lower())
                & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo)
            )
            df_u = df_u[~mascara]
            fila_nueva = pd.DataFrame([{
                "usuario": usuario_actual, "tipo_usuario": tipo,
                "password": _hash_password(pw_nueva), "debe_cambiar_password": "NO",
                "totp_secret": _totp_secret_prev, "totp_activo": _totp_activo_prev, "email": _email_prev,
            }])
            df_u = pd.concat([df_u, fila_nueva], ignore_index=True)
            exito, mensaje = _guardar_hoja_usuarios(df_u)
            if exito:
                st.session_state.forzar_cambio_password = False
                st.session_state.mostrar_aviso_pw_temporal = False
                st.success("✅ Contraseña guardada. Entrando...")
                st.rerun()
            else:
                st.warning(f"⚠️ Contraseña guardada localmente pero no sincronizada con Drive. {mensaje} Vuelve a intentarlo antes de continuar.")
    st.stop()


def seccion_configurar_totp(usuario_actual: str, tipo: str):
    """Panel de autoservicio para activar/desactivar la verificación en dos pasos por EMAIL.
    Disponible para todo el equipo interno (admin: Yuri, Jordi, Alan). Mucho más simple que
    TOTP con app: no requiere instalar nada ni sincronizar relojes, solo reutiliza el email."""
    activo = _2fa_activo(usuario_actual, tipo)
    with st.sidebar.expander(f"🔐 Verificación en dos pasos ({'activada' if activo else 'desactivada'})"):
        # Blindaje: si el usuario todavía no tiene fila propia en la hoja USUARIOS (nunca ha
        # cambiado su contraseña desde el portal), no dejamos activar el 2FA — guardarlo no
        # tendría dónde persistir de forma fiable. Se lo decimos claro en vez de fallar silencioso.
        if not activo and _fila_usuario(_leer_hoja_usuarios(), usuario_actual, tipo) is None:
            st.info("Antes de activar esto, cambia tu contraseña arriba en '🔑 Cambiar mi contraseña' (aunque sea por la misma). Así se crea tu perfil y el 2FA queda guardado de forma fiable.")
            return
        if activo:
            email_actual = _2fa_email_de(usuario_actual, tipo)
            st.success(f"Activada. Se te enviará un código a **{email_actual}** al iniciar sesión.")
            with st.form(f"form_desactivar_totp_{tipo}_{usuario_actual}"):
                st.caption("Para desactivarla, confirma con tu contraseña actual.")
                pw_confirmar = st.text_input("Contraseña actual", type="password", key=f"pw_desactivar_totp_{tipo}")
                desactivar = st.form_submit_button("Desactivar verificación en dos pasos")
            if desactivar:
                usuarios_codigo_local = USUARIOS if tipo == "admin" else USUARIOS_INVERSORES
                if _verificar_credencial(usuario_actual, pw_confirmar, tipo, usuarios_codigo_local) is None:
                    st.error("Contraseña incorrecta.")
                else:
                    exito, mensaje = _guardar_2fa(usuario_actual, tipo, "", False)
                    if exito:
                        st.success(f"Verificación en dos pasos desactivada. {mensaje}")
                    else:
                        st.warning(f"Desactivada localmente. {mensaje}")
            return

        try:
            _smtp_sender_2fa = st.secrets["email"]["sender"]
            _smtp_password_2fa = st.secrets["email"]["password"]
            _display_name_2fa = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
            _email_ok_2fa = True
        except Exception:
            _smtp_sender_2fa = _smtp_password_2fa = _display_name_2fa = ""
            _email_ok_2fa = False

        if not _email_ok_2fa:
            st.warning("Configura el email de envío (Secrets → [email]) para poder activar esto.")
            return

        clave_pendiente = f"email_2fa_pendiente_{usuario_actual}"
        st.caption("Recibirás un código de 6 dígitos por email cada vez que inicies sesión. Sin apps, sin QR.")
        email_2fa = st.text_input(
            "Email al que enviar el código", key=f"input_email_2fa_{tipo}_{usuario_actual}",
            value=st.session_state.get(clave_pendiente, ""),
        )
        if st.button("Enviar código de prueba", key=f"btn_probar_2fa_{tipo}_{usuario_actual}"):
            if not email_2fa or "@" not in email_2fa:
                st.error("Introduce un email válido.")
            else:
                # Guardamos temporalmente el email en la hoja para poder reutilizar
                # _generar_y_enviar_codigo_2fa tal cual (lee el email desde ahí).
                _guardar_2fa(usuario_actual, tipo, email_2fa, False)
                st.session_state[clave_pendiente] = email_2fa
                env_ok, env_msg = _generar_y_enviar_codigo_2fa(
                    usuario_actual, tipo, _smtp_sender_2fa, _smtp_password_2fa, _display_name_2fa,
                )
                if env_ok:
                    st.success(f"Código enviado a {email_2fa}. Revisa tu bandeja e introdúcelo abajo.")
                else:
                    st.error(f"No se pudo enviar el código: {env_msg}")

        if st.session_state.get(clave_pendiente):
            with st.form(f"form_activar_2fa_{tipo}_{usuario_actual}"):
                codigo_confirmar = st.text_input("Código de 6 dígitos recibido por email", key=f"codigo_activar_2fa_{tipo}")
                activar = st.form_submit_button("Activar verificación en dos pasos")
            if activar:
                if _verificar_codigo_2fa(usuario_actual, tipo, codigo_confirmar):
                    exito, mensaje = _guardar_2fa(usuario_actual, tipo, st.session_state[clave_pendiente], True)
                    del st.session_state[clave_pendiente]
                    if exito:
                        st.success(f"✅ Verificación en dos pasos activada. {mensaje}")
                        st.rerun()
                    else:
                        st.warning(f"Activada localmente. {mensaje}")
                else:
                    st.error("Código incorrecto o caducado (10 min). Pulsa 'Enviar código de prueba' otra vez si hace falta.")


if __name__ == "__main__":  # login y sidebar: solo se ejecuta con `streamlit run`, no al importar
    print("[DIAG] 1. Arranca bloque login/sidebar", file=sys.stderr, flush=True)
    aplicar_estilo_profesional()
    print("[DIAG] 2. Estilo aplicado", file=sys.stderr, flush=True)



    # =========================
    # LOGIN
    # =========================
    # Contraseñas TEMPORALES del equipo interno — guardadas aquí como HASH bcrypt, no en texto
    # plano: ni siquiera leyendo este código se puede recuperar la contraseña real. Solo se usan
    # si el usuario todavía no ha cambiado su contraseña desde el portal (formulario "Cambiar mi
    # contraseña" del sidebar, que la re-hashea y la guarda en la hoja USUARIOS de Drive). En
    # cuanto Yuri, Alan y Jordi la cambien una vez cada uno, estos hashes dejan de tener efecto.
    USUARIOS = {
        "Yuri": "$2b$12$zMHrZeMDoZHdO6uCtIIFEOhsbnM5PvqjkaSyQ/58s.QqwKwfihdcu",
        "Jordi": "$2b$12$YetePZJ1Zh0mvAJ8hMai9eIUJaGhh30bo4/R18zUfp2KjxyYrnQR.",
        "Alan": "$2b$12$I3PbU6.DklNw4bLulPfdVePmVDxZGARbTxTJ/pCxlU5G..TdHOOum",
    }
    # Portal de inversores: acceso limitado, solo ven su propia posición.
    # El "usuario" debe coincidir exactamente (en mayúsculas) con el valor de la columna 'inversor' en INVERSIONES.
    # Guardadas como HASH bcrypt (no texto plano) — mismo motivo que USUARIOS arriba: ni siquiera
    # leyendo el código se puede recuperar la contraseña real de ningún inversor.
    USUARIOS_INVERSORES = {
        "PAM": "$2b$12$JtXQhdVOtAUIaswhBIY/S.kcevThLgAPe6S6dJ0sf0x2ss2WaAdqi",
        "LEO": "$2b$12$3RLi8kGMlGzLj94BPk0WQOLZFky1RJbngeUazi8sUEiQPnn74i/da",
        "JORDI CHAPARRO": "$2b$12$lWEmQAL5XI01O7tJjNkmwOoy0hYHuJ6oyb63rr2KVEI.mgI0VZNpG",
        "ROBERTO VISCAFE": "$2b$12$47ZgRROusfgRqn4NGaNofOeSJmMtJr97nN1wznNSnolyBDEIulUnO",
        "CROWE BOLIVIA": "$2b$12$2AUnwWCBSqGEJCqR.68wrOD0EMmbpuOqfV6WzF0G6aaWxC6TDB2CW",
        "JR REAL ESTATE": "$2b$12$2Wnhe4CpwAe2YDUDYWYzPOwb3YQbEGjmQCho0VRLAmrAz.vfoNk7a",
        "2012 JACC GROUP": "$2b$12$qzESlhy/kWWzyMnHtp4o2e0YjQaX8roQ148zj1.djF4CexFL6FlS.",
        "PEDRO MAGAÑA": "$2b$12$e7YD2CCtAg.jK1gnMtG6rulTJuB5FyKBkuG.1wvrfr8u.TuhwfbMS",
        "GOLDEN BRICKS": "$2b$12$mN2hcF9UHn8YfLglwFbl7uHn7S7jeuYvrUeJuwsF11JSrRjZG5NYG",
        "TERESA": "$2b$12$zQoopWqN.5g/VL1Q6E/i5eBOXezLOXPACn8G/P/Ga6hzPx6bzr03u",
        "JEP": "$2b$12$yvHDFYNuEVpVnPdYaRW1fO8tNaO0DvJMu4OKuDOXY7ZCR9OEbOFjS",
        "JORDI ESPECIAL": "$2b$12$1WgxEK3N4BF2a8wfl/SxX.6bW5i4ZtBWy8k/t.9eSewdcn00A3XoS",
        "EVA CHAPARRO": "$2b$12$X3Fmv6uEXN4m/xuBJbvfDO8Nm9.xAz/4ks7kLVUGPCaYF85ad9dCm",
        "PAOLA CHAPARRO": "$2b$12$C7sfiEt5NDO6SniNixFTfOKVGXuV.k1xs5HU.Wuio/dImv7ZFN/dC",
        "JAPAN JORDI": "$2b$12$sZ1ry3OZmJFoDC5.isu.AOqMWgEKJs7Wu9eU7JSSsChXH0Qn9BmMu",
        # Usuario DEMO: portal con datos 100% ficticios para presentar a inversores potenciales.
        # No corresponde a ningún inversor real ni toca el Excel del fondo (ver _construir_datos_demo_inversor).
        "DEMO": "$2b$12$2xF3csT1htw4f1N20nZgb.F8/LN7MtfWAp5eynYlYdAZyMlYn9LLa",
    }

    if "autenticado" not in st.session_state:
        st.session_state.autenticado = False
    if "usuario" not in st.session_state:
        st.session_state.usuario = None
    if "tipo_usuario" not in st.session_state:
        st.session_state.tipo_usuario = None  # "admin" o "inversor"
    if "totp_pendiente" not in st.session_state:
        st.session_state.totp_pendiente = None
    if "reset_pw_pendiente" not in st.session_state:
        st.session_state.reset_pw_pendiente = None  # {"usuario", "tipo"} tras pedir código de reseteo

    if not st.session_state.autenticado:
        st.markdown(
            """
            <div class="login-card">
                <div class="login-logo">CF</div>
                <div class="login-title">Chaparro Fernández Wealth</div>
                <div class="login-subtitle">Acceso privado al sistema financiero interno</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        # ── Paso 2: código recibido por email (solo si el usuario ya superó el paso 1) ──
        if st.session_state.totp_pendiente:
            _pend = st.session_state.totp_pendiente
            _email_2fa_pend = _2fa_email_de(_pend["usuario"], _pend["tipo"])
            _email_oculto = (_email_2fa_pend[:2] + "***@" + _email_2fa_pend.split("@")[-1]) if "@" in _email_2fa_pend else _email_2fa_pend
            st.info(f"🔐 Hola {_pend['usuario']}. Te hemos enviado un código a **{_email_oculto}**. Revisa tu correo (y la carpeta de spam).")
            with st.form("form_totp_login"):
                codigo_totp = st.text_input("Código de 6 dígitos", key="input_codigo_totp_login")
                confirmar_totp = st.form_submit_button("Verificar", use_container_width=True, type="primary")
                c1, c2 = st.columns(2)
                reenviar_totp = c1.form_submit_button("Reenviar código", use_container_width=True)
                cancelar_totp = c2.form_submit_button("Cancelar", use_container_width=True)
            if cancelar_totp:
                st.session_state.totp_pendiente = None
                st.rerun()
            if reenviar_totp:
                try:
                    _smtp_sender_l = st.secrets["email"]["sender"]
                    _smtp_password_l = st.secrets["email"]["password"]
                    _display_name_l = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
                    env_ok, env_msg = _generar_y_enviar_codigo_2fa(_pend["usuario"], _pend["tipo"], _smtp_sender_l, _smtp_password_l, _display_name_l)
                    if env_ok:
                        st.success("Código reenviado.")
                    else:
                        st.error(f"No se pudo reenviar: {env_msg}")
                except Exception:
                    st.error("El email de envío no está configurado (Secrets → [email]).")
            if confirmar_totp:
                _bloqueado, _min_restantes = _login_bloqueado("totp", _pend["usuario"])
                if _bloqueado:
                    st.error(f"🔒 Demasiados intentos fallidos. Inténtalo de nuevo en {_min_restantes} minuto(s).")
                else:
                    if _verificar_codigo_2fa(_pend["usuario"], _pend["tipo"], codigo_totp):
                        _resetear_intentos_login("totp", _pend["usuario"])
                        _completar_login(_pend["usuario"], _pend["tipo"])
                        st.rerun()
                    else:
                        _registrar_intento_fallido("totp", _pend["usuario"])
                        st.error("Código incorrecto o caducado. Puedes pedir uno nuevo con 'Reenviar código'.")
            st.stop()

        # ── "Olvidé mi contraseña" — paso 2: código recibido + nueva contraseña ──
        if st.session_state.reset_pw_pendiente:
            _pend_r = st.session_state.reset_pw_pendiente
            _email_oculto_r = (_pend_r["email"][:2] + "***@" + _pend_r["email"].split("@")[-1]) if "@" in _pend_r["email"] else _pend_r["email"]
            st.info(f"📧 Hola {_pend_r['usuario']}. Te hemos enviado un código a **{_email_oculto_r}** para restablecer tu contraseña. Revisa tu correo (y la carpeta de spam).")
            with st.form("form_reset_pw_confirmar"):
                codigo_reset = st.text_input("Código de 6 dígitos", key="input_codigo_reset_pw")
                pw_reset_nueva = st.text_input("Nueva contraseña", type="password", key="input_pw_reset_nueva")
                pw_reset_nueva2 = st.text_input("Repite la nueva contraseña", type="password", key="input_pw_reset_nueva2")
                confirmar_reset = st.form_submit_button("Restablecer contraseña", use_container_width=True, type="primary")
                c1r, c2r = st.columns(2)
                reenviar_reset = c1r.form_submit_button("Reenviar código", use_container_width=True)
                cancelar_reset = c2r.form_submit_button("Cancelar", use_container_width=True)
            if cancelar_reset:
                _consumir_codigo_reset_password(_pend_r["usuario"], _pend_r["tipo"])
                st.session_state.reset_pw_pendiente = None
                st.rerun()
            if reenviar_reset:
                try:
                    _smtp_sender_r = st.secrets["email"]["sender"]
                    _smtp_password_r = st.secrets["email"]["password"]
                    _display_name_r = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
                    env_ok_r, env_msg_r = _generar_y_enviar_codigo_reset_password(
                        _pend_r["usuario"], _pend_r["tipo"], _pend_r["email"], _smtp_sender_r, _smtp_password_r, _display_name_r,
                    )
                    st.success("Código reenviado.") if env_ok_r else st.error(f"No se pudo reenviar: {env_msg_r}")
                except Exception:
                    st.error("El email de envío no está configurado (Secrets → [email]).")
            if confirmar_reset:
                _bloqueado_r, _min_restantes_r = _login_bloqueado("reset_pw", _pend_r["usuario"])
                if _bloqueado_r:
                    st.error(f"🔒 Demasiados intentos fallidos. Inténtalo de nuevo en {_min_restantes_r} minuto(s).")
                elif not _verificar_codigo_reset_password(_pend_r["usuario"], _pend_r["tipo"], codigo_reset):
                    _registrar_intento_fallido("reset_pw", _pend_r["usuario"])
                    st.error("Código incorrecto o caducado. Puedes pedir uno nuevo con 'Reenviar código'.")
                elif len(pw_reset_nueva) < 8:
                    st.error("La nueva contraseña debe tener al menos 8 caracteres.")
                elif pw_reset_nueva != pw_reset_nueva2:
                    st.error("Las dos contraseñas no coinciden.")
                else:
                    exito_r, mensaje_r = _fijar_password_autoservicio(_pend_r["usuario"], _pend_r["tipo"], pw_reset_nueva)
                    if exito_r:
                        _resetear_intentos_login("reset_pw", _pend_r["usuario"])
                        _consumir_codigo_reset_password(_pend_r["usuario"], _pend_r["tipo"])
                        st.session_state.reset_pw_pendiente = None
                        st.success("✅ Contraseña actualizada. Ya puedes iniciar sesión con ella.")
                    else:
                        st.warning(f"⚠️ La contraseña se guardó localmente pero no se sincronizó con Drive. {mensaje_r} Prueba de nuevo antes de cerrar esta pestaña.")
            st.stop()

        # Login único: no se distingue en pantalla entre "equipo interno" y "portal de
        # inversor" — nadie que vea la pantalla de entrada puede saber que el equipo interno
        # accede desde el mismo sitio. Se prueban las credenciales primero contra admin y,
        # si no coinciden, contra inversor; el tipo se determina solo, en silencio.
        with st.form("login_form_unico"):
            usuario_txt = st.text_input("Usuario")
            password = st.text_input("Contraseña", type="password")
            entrar = st.form_submit_button("Entrar")
        if entrar:
            _bloq_admin, _min_admin = _login_bloqueado("admin", usuario_txt or "")
            _bloq_inv, _min_inv = _login_bloqueado("inversor", usuario_txt or "")
            if _bloq_admin and _bloq_inv:
                st.error(f"🔒 Demasiados intentos fallidos. Inténtalo de nuevo en {max(_min_admin, _min_inv)} minuto(s).")
            else:
                usuario_match, tipo_match = None, None
                if not _bloq_admin:
                    usuario_match = _verificar_credencial(usuario_txt, password, "admin", USUARIOS)
                    if usuario_match:
                        tipo_match = "admin"
                if not usuario_match and not _bloq_inv:
                    usuario_match = _verificar_credencial(usuario_txt, password, "inversor", USUARIOS_INVERSORES)
                    if usuario_match:
                        tipo_match = "inversor"
                if usuario_match:
                    _resetear_intentos_login(tipo_match, usuario_txt)
                    if _2fa_activo(usuario_match, tipo_match):
                        try:
                            _smtp_sender_l = st.secrets["email"]["sender"]
                            _smtp_password_l = st.secrets["email"]["password"]
                            _display_name_l = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
                            _generar_y_enviar_codigo_2fa(usuario_match, tipo_match, _smtp_sender_l, _smtp_password_l, _display_name_l)
                        except Exception:
                            pass
                        st.session_state.totp_pendiente = {"usuario": usuario_match, "tipo": tipo_match}
                    else:
                        _completar_login(usuario_match, tipo_match)
                    st.rerun()
                else:
                    if not _bloq_admin:
                        _registrar_intento_fallido("admin", usuario_txt or "")
                    if not _bloq_inv:
                        _registrar_intento_fallido("inversor", usuario_txt or "")
                    st.error("Usuario o contraseña incorrectos")

        # ── "Olvidé mi contraseña" — paso 1: pedir usuario y enviar código por email ──
        with st.expander("¿Olvidaste tu contraseña?"):
            st.caption("Te enviaremos un código de 6 dígitos al email que tengamos registrado para tu usuario.")
            with st.form("form_reset_pw_solicitar"):
                usuario_reset_txt = st.text_input("Tu usuario", key="input_usuario_reset_solicitar")
                solicitar_reset = st.form_submit_button("Enviarme el código")
            if solicitar_reset:
                usuario_reset_limpio = (usuario_reset_txt or "").strip()
                if not usuario_reset_limpio:
                    st.error("Escribe tu usuario.")
                else:
                    # Igual que en el login: se prueba primero como admin y, si no coincide,
                    # como inversor — el tipo se determina solo, sin pedirlo en pantalla.
                    usuario_canon, tipo_canon = None, None
                    fila_admin = _fila_usuario(_leer_hoja_usuarios(), usuario_reset_limpio, "admin")
                    if fila_admin is not None:
                        usuario_canon, tipo_canon = str(fila_admin["usuario"]), "admin"
                    else:
                        match_admin = next((u for u in USUARIOS if u.strip().lower() == usuario_reset_limpio.lower()), None)
                        if match_admin:
                            usuario_canon, tipo_canon = match_admin, "admin"
                    if usuario_canon is None:
                        fila_inv = _fila_usuario(_leer_hoja_usuarios(), usuario_reset_limpio, "inversor")
                        if fila_inv is not None:
                            usuario_canon, tipo_canon = str(fila_inv["usuario"]), "inversor"
                        else:
                            match_inv = next((u for u in USUARIOS_INVERSORES if u.strip().lower() == usuario_reset_limpio.lower()), None)
                            if match_inv:
                                usuario_canon, tipo_canon = match_inv, "inversor"
                    if usuario_canon is None:
                        # Mensaje deliberadamente genérico (no confirma si el usuario existe o
                        # no) — evita que alguien use este formulario para descubrir usuarios
                        # válidos por prueba y error.
                        st.info("Si el usuario existe y tiene un email registrado, te llegará un código en unos segundos.")
                    else:
                        email_contacto = _email_contacto_de(usuario_canon, tipo_canon)
                        if not email_contacto:
                            st.error(
                                "No hay ningún email de recuperación registrado para este usuario todavía. "
                                "Pídele al administrador que te resetee la contraseña desde el panel — de paso, "
                                "quedará tu email guardado para poder usar esta opción la próxima vez."
                            )
                        else:
                            try:
                                _smtp_sender_r2 = st.secrets["email"]["sender"]
                                _smtp_password_r2 = st.secrets["email"]["password"]
                                _display_name_r2 = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
                            except Exception:
                                st.error("El email de envío no está configurado (Secrets → [email]). Pídele al administrador que te resetee la contraseña manualmente.")
                                _smtp_sender_r2 = None
                            if _smtp_sender_r2:
                                env_ok_r2, env_msg_r2 = _generar_y_enviar_codigo_reset_password(
                                    usuario_canon, tipo_canon, email_contacto, _smtp_sender_r2, _smtp_password_r2, _display_name_r2,
                                )
                                if env_ok_r2:
                                    st.session_state.reset_pw_pendiente = {"usuario": usuario_canon, "tipo": tipo_canon, "email": email_contacto}
                                    st.rerun()
                                else:
                                    st.error(f"No se pudo enviar el código: {env_msg_r2}")
        st.stop()

    # ── Timeout de sesión: expira por inactividad o por duración máxima absoluta ──
    SESION_INACTIVIDAD_SEGUNDOS = 30 * 60   # 30 minutos sin ninguna interacción
    SESION_MAXIMA_SEGUNDOS = 12 * 3600      # 12 horas desde el login, aunque haya actividad
    _ahora = time.time()
    _login_ts = st.session_state.get("login_timestamp", _ahora)
    _ultima_act = st.session_state.get("ultima_actividad", _ahora)
    if (_ahora - _ultima_act > SESION_INACTIVIDAD_SEGUNDOS) or (_ahora - _login_ts > SESION_MAXIMA_SEGUNDOS):
        st.session_state.autenticado = False
        st.session_state.usuario = None
        st.session_state.tipo_usuario = None
        st.warning("⏱️ Tu sesión ha expirado por seguridad. Vuelve a iniciar sesión.")
        st.stop()
    st.session_state.ultima_actividad = _ahora
    print("[DIAG] 3. Sesión validada, empieza sidebar autenticado", file=sys.stderr, flush=True)

    # ── Cambio obligatorio de contraseña: bloquea el resto de la app hasta que se complete ──
    if st.session_state.get("forzar_cambio_password"):
        formulario_cambio_obligatorio_password(st.session_state.usuario, st.session_state.tipo_usuario)

    st.sidebar.markdown(f"**Usuario conectado:** {st.session_state.usuario}")

    # ── Cambio de vista admin ↔ inversor: solo para los admins con perfil(es) de inversor
    # adicional autorizados (ver ADMIN_VISTAS_INVERSOR_ADICIONALES). No afecta a nadie más. ──
    if st.session_state.tipo_usuario == "admin":
        _vistas_extra_admin = ADMIN_VISTAS_INVERSOR_ADICIONALES.get(str(st.session_state.usuario).strip().upper(), [])
        if _vistas_extra_admin:
            _opciones_vista = ["🛠️ Panel de administración"] + [f"👤 {v}" for v in _vistas_extra_admin]
            _vista_elegida = st.sidebar.selectbox(
                "Vista", _opciones_vista, key="admin_vista_selector",
                label_visibility="collapsed",
            )
            st.session_state.vista_admin_como_inversor = (
                None if _vista_elegida == "🛠️ Panel de administración" else _vista_elegida[2:]
            )
            st.sidebar.divider()
        else:
            st.session_state.vista_admin_como_inversor = None

    # ── Atajo fijo para el asistente de IA: siempre visible en la barra lateral para el
    # equipo interno, sin tener que buscarlo en el desplegable del menú principal. Solo se
    # muestra cuando el admin está en su panel (no mientras ve un perfil de inversor, que ya
    # tiene su propio asistente incrustado en esa vista). ──
    if st.session_state.tipo_usuario == "admin" and not st.session_state.get("vista_admin_como_inversor"):
        if st.sidebar.button("💬 Preguntar al asistente", use_container_width=True, key="btn_atajo_asistente_ia"):
            st.session_state["menu_principal_selector"] = "✨ Asistente IA"
            st.rerun()
        st.sidebar.divider()

    if st.session_state.get("mostrar_aviso_pw_temporal"):
        st.sidebar.warning("⚠️ Estás usando la contraseña temporal. Cámbiala ahora abajo, en '🔑 Cambiar mi contraseña'.")
    _usuarios_codigo_actual = USUARIOS if st.session_state.tipo_usuario == "admin" else USUARIOS_INVERSORES
    if str(st.session_state.usuario).strip().upper() != "DEMO":
        formulario_cambiar_password(st.session_state.usuario, st.session_state.tipo_usuario, _usuarios_codigo_actual)
        formulario_email_recuperacion(st.session_state.usuario, st.session_state.tipo_usuario)
    # Verificación en dos pasos por email: equipo interno (admin) + inversores piloto autorizados.
    if st.session_state.tipo_usuario == "admin" or str(st.session_state.usuario).strip().upper() in INVERSORES_CON_2FA:
        seccion_configurar_totp(st.session_state.usuario, st.session_state.tipo_usuario)
    if st.sidebar.button("Cerrar sesión"):
        st.session_state.autenticado = False
        st.session_state.usuario = None
        st.session_state.tipo_usuario = None
        st.rerun()


# =========================
# UTILIDADES Y CARGA
# =========================
def fmt(x):
    try:
        return f"${float(x):,.2f}"
    except Exception:
        return "$0.00"


def fmt_pct(x):
    try:
        if pd.isna(x):
            return "0.00%"
        return f"{float(x) * 100:.2f}%"
    except Exception:
        return "0.00%"


def nombre_mes_es(mes: int) -> str:
    meses = {1: "enero", 2: "febrero", 3: "marzo", 4: "abril", 5: "mayo", 6: "junio", 7: "julio", 8: "agosto", 9: "septiembre", 10: "octubre", 11: "noviembre", 12: "diciembre"}
    return meses.get(int(mes), str(mes))


def ultimo_dia_mes(anio: int, mes: int) -> int:
    return calendar.monthrange(int(anio), int(mes))[1]


def limpiar_texto(x) -> str:
    if pd.isna(x):
        return ""
    return str(x).strip().lower()


def es_chaparro_fernandez_row(row) -> bool:
    """Detecta SOLO las inversiones internas de la sociedad Chaparro Fernández.

    IMPORTANTE:
    No se debe marcar como Chaparro Fernández a personas que se apellidan
    Chaparro o Fernández, como JORDI CHAPARRO, YURI FERNANDEZ, EVA CHAPARRO
    o PAOLA CHAPARRO. Por eso la detección se hace únicamente sobre la columna
    inversor y exige el nombre completo de la sociedad.
    """
    inversor = limpiar_texto(row.get("inversor", ""))
    inversor = (
        inversor.replace("á", "a")
        .replace("é", "e")
        .replace("í", "i")
        .replace("ó", "o")
        .replace("ú", "u")
        .replace("ñ", "n")
        .replace("-", " ")
    )
    inversor = " ".join(inversor.split())

    nombres_sociedad = {
        "chaparro fernandez",
        "chaparro fernandez sl",
        "chaparro fernandez s.l.",
        "chaparro fernandez sociedad",
        "chaparro fernandez wealth",
    }
    return inversor in nombres_sociedad


def aplicar_filtro_chaparro_fernandez(df: pd.DataFrame, incluir_chaparro: bool) -> pd.DataFrame:
    """Incluye o excluye las inversiones internas de Chaparro Fernández.

    Si incluir_chaparro=True, no toca el dataframe.
    Si incluir_chaparro=False, elimina filas detectadas como Chaparro Fernández.
    Además añade una columna auxiliar es_chaparro_fernandez para poder auditarlo.
    """
    if df is None or df.empty:
        return df
    out = df.copy()
    out["es_chaparro_fernandez"] = out.apply(es_chaparro_fernandez_row, axis=1)
    if incluir_chaparro:
        return out
    return out[~out["es_chaparro_fernandez"]].copy()


def descargar_excel_desde_drive():
    """Descarga el Excel desde Google Sheets y lo guarda localmente.

    Intenta primero por la API autenticada de Drive (con la cuenta de servicio) — no depende del
    enlace público de exportación, que puede ralentizarse o limitarse tras muchas descargas
    seguidas en poco tiempo. Si no hay credenciales configuradas o falla, cae al enlace público
    como respaldo, con un timeout corto para no dejar la app colgada esperando indefinidamente.
    """
    if "gcp_service_account" in st.secrets:
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            from googleapiclient.http import MediaIoBaseDownload
            import io as _io_dl

            credenciales = service_account.Credentials.from_service_account_info(
                dict(st.secrets["gcp_service_account"]),
                scopes=["https://www.googleapis.com/auth/drive.readonly"],
            )
            servicio = build("drive", "v3", credentials=credenciales)
            request = servicio.files().export_media(
                fileId=GDRIVE_FILE_ID,
                mimeType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
            buffer = _io_dl.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _status, done = downloader.next_chunk(num_retries=1)
            with open(ARCHIVO, "wb") as f:
                f.write(buffer.getvalue())
            return True
        except Exception:
            pass  # sin credenciales válidas o fallo puntual: caemos al enlace público de abajo

    try:
        url = f"https://docs.google.com/spreadsheets/d/{GDRIVE_FILE_ID}/export?format=xlsx"
        response = requests.get(url, timeout=12)
        if response.status_code == 200:
            with open(ARCHIVO, "wb") as f:
                f.write(response.content)
            return True
        else:
            st.warning(f"No se pudo descargar el Excel desde Google Drive (status {response.status_code}).")
            return False
    except Exception as e:
        st.warning(f"No se pudo descargar el Excel desde Google Drive: {e}")
        return False


@st.cache_data(show_spinner=False, ttl=300)
def parsear_fecha_robusta(serie: pd.Series) -> pd.Series:
    """
    Parsea una columna de fechas admitiendo tanto DD/MM/AAAA (formato español, el habitual en
    INVERSIONES) como AAAA-MM-DD (ISO, el que a veces se escribe sin querer al editar el Sheet a
    mano o al pegar datos). A diferencia de un pd.to_datetime(..., dayfirst=True) aplicado de golpe
    a toda la columna —que puede "aprender" un formato equivocado de las primeras filas y destrozar
    en silencio filas posteriores con otro formato—, aquí se prueba fecha a fecha, probando primero
    los formatos explícitos conocidos y solo cayendo a la detección automática como último recurso.
    Así da igual si la fecha se escribió desde la app (siempre en DD/MM/AAAA) o a mano en Google
    Sheets (a veces en ISO): las dos se leen bien.
    """
    def _parsear_una(valor):
        if pd.isna(valor):
            return pd.NaT
        if isinstance(valor, (pd.Timestamp, datetime)):
            return pd.Timestamp(valor)
        texto = str(valor).strip()
        if not texto:
            return pd.NaT
        for fmt in ("%d/%m/%Y", "%Y-%m-%d", "%d-%m-%Y", "%Y/%m/%d"):
            try:
                return pd.to_datetime(texto, format=fmt)
            except (ValueError, TypeError):
                continue
        return pd.to_datetime(texto, errors="coerce", dayfirst=True)
    return serie.apply(_parsear_una)


def cargar_excel_completo():
    if DATA_SOURCE == "postgres":
        try:
            from postgres_reader import cargar_excel_completo_postgres
            return cargar_excel_completo_postgres()
        except Exception as e:
            print(f"[DATA_SOURCE=postgres] Fallo leyendo de Postgres, cayendo a Drive: {e}", file=sys.stderr)
    return _cargar_excel_completo_desde_drive()


def _cargar_excel_completo_desde_drive():
    descargar_excel_desde_drive()
    inv = pd.read_excel(ARCHIVO, sheet_name=HOJA_INVERSIONES)
    cal = pd.read_excel(ARCHIVO, sheet_name=HOJA_CALENDARIO)
    try:
        control = pd.read_excel(ARCHIVO, sheet_name=HOJA_CONTROL)
    except Exception:
        control = pd.DataFrame()

    inv.columns = [str(c).strip().lower() for c in inv.columns]
    cal.columns = [str(c).strip().lower() for c in cal.columns]
    control.columns = [str(c).strip().lower() for c in control.columns]

    if "unnamed: 6" in inv.columns and "cuenta_cobro" not in inv.columns:
        inv = inv.rename(columns={"unnamed: 6": "cuenta_cobro"})

    for col in ["id_inversion", "inversor", "tipo_inversion", "subtipo_inversion", "nombre_activo", "metodo_calculo", "activo_generador_interes", "tipo_operacion", "capital_nuevo_real", "cuenta_cobro", "motivo"]:
        if col in inv.columns:
            inv[col] = inv[col].fillna("").astype(str).str.strip()

    for col in ["fecha_inversion", "fecha_final_inversion"]:
        if col in inv.columns:
            inv[col] = parsear_fecha_robusta(inv[col])

    for col in ["capital_invertido", "interes_inversor_anual", "interes_nota_anual"]:
        if col in inv.columns:
            inv[col] = pd.to_numeric(inv[col], errors="coerce").fillna(0)
        else:
            inv[col] = 0

    # periodicidad_meses: cuántos meses cubre cada pago (1=mensual, 3=trimestral, etc.)
    if "periodicidad_meses" in inv.columns:
        inv["periodicidad_meses"] = pd.to_numeric(inv["periodicidad_meses"], errors="coerce").fillna(1).astype(int)
    else:
        inv["periodicidad_meses"] = 1
    if "nota" in cal.columns:
        cal["nota"] = pd.to_numeric(cal["nota"], errors="coerce").astype("Int64")
    if "tipo_evento" in cal.columns:
        cal["tipo_evento"] = cal["tipo_evento"].fillna("").astype(str).str.strip().str.upper()
    if "fecha" in cal.columns:
        # BUG CRÍTICO (encontrado 18/07/2026): CALENDARIO_NOTAS guarda las fechas en texto ISO
        # (AAAA-MM-DD). Parsear la columna ENTERA de golpe con dayfirst=True (pensado para el
        # formato español DD/MM/AAAA de otras hojas como INVERSIONES) hace que pandas descarte
        # como NaT ~2 de cada 3 fechas ISO válidas y sin ambigüedad — cada fecha, una por una,
        # se parsea perfectamente; es solo el parseo en bloque con dayfirst=True el que falla.
        # Probamos primero el formato ISO explícito (cubre el caso real de esta hoja al 100%),
        # y solo usamos dayfirst=True como red de seguridad para lo que quede sin resolver (por
        # si alguna fila se escribió a mano en formato español DD/MM/AAAA).
        _fecha_iso = pd.to_datetime(cal["fecha"], format="%Y-%m-%d", errors="coerce")
        _pendientes = _fecha_iso.isna() & cal["fecha"].notna()
        if _pendientes.any():
            _fecha_iso.loc[_pendientes] = pd.to_datetime(cal.loc[_pendientes, "fecha"], errors="coerce", dayfirst=True)
        cal["fecha"] = _fecha_iso.dt.normalize()

    if not control.empty:
        if "nota" in control.columns:
            control["nota"] = pd.to_numeric(control["nota"], errors="coerce").astype("Int64")
        if "ticker" in control.columns:
            control["ticker"] = control["ticker"].fillna("").astype(str).str.strip().str.upper()
        for col in ["precio_compra", "barrera_cupon", "contingency", "barrera_capital"]:
            if col in control.columns:
                control[col] = pd.to_numeric(control[col], errors="coerce")
    return inv, cal, control


def _md_seguro(texto: str) -> str:
    """
    Escapa los signos '$' antes de mostrar texto con st.markdown().
    Streamlit interpreta '$...$' como fórmulas LaTeX — sin este escape, un texto generado
    por IA que menciona precios como '$450 a $640' se renderiza roto (ej. '450a640' en
    cursiva de fórmula matemática) en vez de mostrarse como texto normal.
    """
    if not texto:
        return texto
    return str(texto).replace("$", "\\$")


def leer_hoja_excel(nombre_hoja: str) -> pd.DataFrame:
    if DATA_SOURCE == "postgres":
        try:
            from postgres_reader import leer_hoja_excel_postgres
            return leer_hoja_excel_postgres(nombre_hoja)
        except Exception as e:
            print(f"[DATA_SOURCE=postgres] Fallo leyendo '{nombre_hoja}' de Postgres, cayendo a Drive: {e}", file=sys.stderr)
    try:
        descargar_excel_desde_drive()
        df = pd.read_excel(ARCHIVO, sheet_name=nombre_hoja)
        df.columns = [str(c).strip().lower() for c in df.columns]
        return df
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=60, show_spinner=False)
def _leer_calendario_calls_cached() -> pd.DataFrame:
    """Versión cacheada de leer_hoja_excel('CALENDARIO_CALLS') — se usa en varias secciones y,
    sin caché, se releía en cada rerun de TODA la app (Streamlit reejecuta todas las pestañas en
    cada clic, no solo la activa), contribuyendo a la lentitud al cambiar de pestaña o de nota."""
    return leer_hoja_excel("CALENDARIO_CALLS")


@st.cache_data(ttl=60, show_spinner=False)
def _leer_auditoria_notas_cached() -> pd.DataFrame:
    """
    Versión cacheada de leer_hoja_excel('AUDITORIA_NOTAS') para el visor de solo lectura de
    'Auditar nota existente'. Sin caché, esta lectura (openpyxl + posible descarga de Drive) se
    repetía en CADA interacción de CUALQUIER parte de la app (Streamlit reejecuta todas las pestañas
    en cada rerun, no solo la que estás mirando), lo que provocaba lentitud/bloqueos notables al
    cambiar de pestaña o de número de nota. 60 segundos de caché es un buen equilibrio: la auditoría
    no cambia cada segundo, así que no hace falta leerla en tiempo real en cada clic.
    """
    return leer_hoja_excel("AUDITORIA_NOTAS")


@st.cache_data(ttl=60, show_spinner=False)
def _leer_movimientos_banco_cached() -> pd.DataFrame:
    """Versión cacheada de leer_hoja_excel('MOVIMIENTOS_BANCO') — histórico de movimientos ya
    importados desde extractos del broker (custodio de las notas estructuradas). Vacío si
    todavía no se ha importado ningún extracto (la hoja ni siquiera existe en el Excel)."""
    return leer_hoja_excel(HOJA_MOVIMIENTOS_BANCO)


ETIQUETAS_CATEGORIA_BANCO = {
    "RUIDO_SWEEP": "Barrido de caja (FDIC Sweep)",
    "INTERES_COBRADO": "Interés cobrado",
    "ALTA_INVERSION": "Alta de inversión",
    "CANCELACION_INVERSION": "Cancelación / vencimiento",
    "APORTACION_CAPITAL": "Aportación de capital",
    "TRASPASO_INTERNO": "Traspaso interno",
    "COMISION_GASTO": "Comisión / gasto",
    "SIN_CLASIFICAR": "Sin clasificar",
}


def clasificar_movimiento_extracto_banco(fila) -> str:
    """Clasifica una fila del extracto del broker (custodio de las notas estructuradas) en las
    categorías contables del fondo. Ver ETIQUETAS_CATEGORIA_BANCO para el texto de cada una."""
    tipo = str(fila.get("Tipo de Transacción", "")).strip()
    accion = str(fila.get("Acción", "")).strip()
    desc_trans = str(fila.get("Descripción de Transacción", "")).strip().upper()

    if tipo == "FDIC Sweep":
        # Barrido automático de efectivo: compra/vende el fondo de barrido cada vez que hay
        # actividad de caja en la cuenta. Se anula solo con el tiempo — no es un ingreso/gasto
        # real del fondo, pero SÍ forma parte del saldo real de la cuenta, así que se incluye
        # en el saldo aunque se excluya de los totales por categoría.
        return "RUIDO_SWEEP"
    if tipo == "Dividends and Interest":
        return "INTERES_COBRADO"
    if tipo == "Trade Activity":
        return "ALTA_INVERSION" if accion == "Buy" else "CANCELACION_INVERSION"
    if tipo == "Withdrawal":
        # La mayoría de las salidas vistas hasta ahora son la comisión de Aragon Capital, pero
        # cualquier "Withdrawal" cae aquí — revisar en la tabla de importación antes de guardar
        # si el concepto no es una comisión reconocible.
        return "COMISION_GASTO"
    if tipo == "Deposit":
        return "APORTACION_CAPITAL"
    if tipo == "Journal":
        return "TRASPASO_INTERNO"
    return "SIN_CLASIFICAR"


def parsear_extracto_banco_bytes(contenido: bytes) -> pd.DataFrame:
    """Lee el Excel de movimientos exportado del broker y devuelve un DataFrame clasificado en
    las categorías contables del fondo, con un identificador único por fila (id_movimiento) para
    poder detectar duplicados si se reimporta un extracto que se solape con uno ya guardado."""
    df = pd.read_excel(BytesIO(contenido))
    df.columns = [str(c).strip() for c in df.columns]

    columnas_esperadas = [
        "Fecha", "Monto", "Tipo de Transacción", "Descripción de Activo",
        "Descripción de Transacción", "Acción", "Símbolo / ID",
    ]
    faltan = [c for c in columnas_esperadas if c not in df.columns]
    if faltan:
        raise ValueError(
            "El archivo no tiene el formato esperado del extracto del broker — "
            f"faltan columnas: {', '.join(faltan)}"
        )

    df["Fecha"] = pd.to_datetime(df["Fecha"], errors="coerce")
    df["Monto"] = pd.to_numeric(df["Monto"], errors="coerce").fillna(0)
    df["categoria"] = df.apply(clasificar_movimiento_extracto_banco, axis=1)
    df["cusip"] = df["Símbolo / ID"].astype(str).str.strip().str.upper()
    df.loc[df["cusip"].isin(["—", "NAN", "NONE", ""]), "cusip"] = ""

    def _id_movimiento(fila):
        partes = [
            fila["Fecha"].strftime("%Y-%m-%d") if pd.notna(fila["Fecha"]) else "SINFECHA",
            f"{float(fila['Monto']):.2f}",
            str(fila.get("Símbolo / ID", "")).strip(),
            str(fila.get("Descripción de Transacción", "")).strip()[:40],
        ]
        return "|".join(partes)

    df["id_movimiento"] = df.apply(_id_movimiento, axis=1)
    return df


def matchear_movimientos_con_notas(df_mov: pd.DataFrame, df_control: pd.DataFrame) -> pd.DataFrame:
    """Añade la columna nota_asociada a df_mov cruzando la columna cusip contra el CUSIP guardado
    en CONTROL_NOTAS. Si CONTROL_NOTAS todavía no tiene columna cusip (o está vacía), no casa nada
    en vez de fallar."""
    df_mov = df_mov.copy()
    if df_control is None or df_control.empty or "cusip" not in df_control.columns or "nota" not in df_control.columns:
        df_mov["nota_asociada"] = pd.NA
        return df_mov
    mapa_cusip = (
        df_control[df_control["cusip"].astype(str).str.strip() != ""]
        .assign(cusip_norm=lambda d: d["cusip"].astype(str).str.strip().str.upper())
        .drop_duplicates("cusip_norm")
        .set_index("cusip_norm")["nota"]
    )
    df_mov["nota_asociada"] = df_mov["cusip"].map(mapa_cusip)
    return df_mov


def preparar_tabla_monetaria(df: pd.DataFrame, columnas_monetarias) -> pd.DataFrame:
    if df is None or df.empty:
        return df
    out = df.copy()
    for col in out.columns:
        if "fecha" in col:
            out[col] = pd.to_datetime(out[col], errors="coerce").dt.strftime("%d/%m/%Y")
    for col in columnas_monetarias:
        if col in out.columns:
            out[col] = out[col].map(fmt)
    return out


def mostrar_metricas(titulo, valores):
    st.subheader(titulo)
    cols = st.columns(len(valores))
    for col, (label, value) in zip(cols, valores):
        col.metric(label, value)


# =========================
# CÁLCULOS ACTIVOS
# =========================
def filtrar_activo(df_base: pd.DataFrame, activo: str) -> pd.DataFrame:
    activo_l = activo.lower()
    subtipo = df_base.get("subtipo_inversion", pd.Series(index=df_base.index, dtype=str)).astype(str).str.lower()
    nombre = df_base.get("nombre_activo", pd.Series(index=df_base.index, dtype=str)).astype(str).str.lower()
    if activo_l == "futbol":
        return df_base[subtipo.isin(["futbol", "fútbol"]) | nombre.isin(["futbol", "fútbol"])].copy()
    if activo_l == "prestamo":
        return df_base[
            subtipo.str.contains("prestamo", na=False) | subtipo.str.contains("préstamo", na=False) |
            nombre.str.contains("prestamo", na=False) | nombre.str.contains("préstamo", na=False)
        ].copy()
    return df_base[subtipo.eq(activo_l) | nombre.eq(activo_l)].copy()


def dias_activos_en_mes(fecha_inicio, fecha_fin, anio: int, mes: int) -> int:
    inicio_mes = pd.Timestamp(anio, mes, 1)
    fin_mes = pd.Timestamp(anio, mes, ultimo_dia_mes(anio, mes))
    if pd.isna(fecha_inicio) or fecha_inicio > fin_mes:
        return 0
    if pd.notna(fecha_fin) and fecha_fin < inicio_mes:
        return 0
    inicio_real = max(fecha_inicio, inicio_mes)
    fin_real = fin_mes if pd.isna(fecha_fin) else min(fecha_fin, fin_mes)
    if inicio_real > fin_real:
        return 0
    return (fin_real - inicio_real).days + 1



def cargar_movimientos_motoclick() -> pd.DataFrame:
    """Lee la hoja MOVIMIENTOS_MOTOCLICK del Excel.
    Columnas esperadas: fecha, tipo (DEVOLUCION/REINVERSION), importe, descripcion.
    """
    try:
        df = pd.read_excel(ARCHIVO, sheet_name=HOJA_MOTOCLICK)
        df.columns = [str(c).strip().lower() for c in df.columns]
        df["fecha"] = pd.to_datetime(df["fecha"], dayfirst=True, errors="coerce")
        df["importe"] = pd.to_numeric(df["importe"], errors="coerce").fillna(0)
        df["tipo"] = df["tipo"].astype(str).str.strip().str.upper()
        return df.dropna(subset=["fecha"])
    except Exception:
        return pd.DataFrame(columns=["fecha", "tipo", "importe", "descripcion"])


def ingreso_bruto_motoclick_mes(df_inv: pd.DataFrame, df_mov: pd.DataFrame, anio: int, mes: int) -> float:
    """
    Calcula el ingreso bruto de MotoClick para la compañía en el mes,
    usando el capital real desplegado día a día.

    Lógica:
    - Capital base cada día = suma del capital de todas las inversiones activas en MotoClick ese día
    - Ajuste = devoluciones acumuladas (restan) y reinversiones acumuladas (suman) hasta ese día
    - Capital real día = capital base + ajuste
    - Ingreso = promedio(capital_real_día) × 25% / 12
    """
    dias_mes = ultimo_dia_mes(anio, mes)
    inicio_mes = pd.Timestamp(anio, mes, 1)
    fin_mes = pd.Timestamp(anio, mes, dias_mes)

    mc = filtrar_activo(df_inv, "motoclick")
    mc = mc[mc["tipo_operacion"].astype(str).str.lower() != "cancelada"].copy()
    mc["fecha_inversion"] = parsear_fecha_robusta(mc["fecha_inversion"])
    mc["fecha_final_inversion"] = parsear_fecha_robusta(mc["fecha_final_inversion"])

    # Movimientos del mes y anteriores (acumulados)
    movs = df_mov.copy() if not df_mov.empty else pd.DataFrame(columns=["fecha","tipo","importe"])

    suma_capital_diaria = 0.0
    for dia in range(1, dias_mes + 1):
        fecha_dia = pd.Timestamp(anio, mes, dia)

        # Capital base: inversiones activas ese día
        cap_base = 0.0
        for _, r in mc.iterrows():
            fi = r["fecha_inversion"]
            ff = r["fecha_final_inversion"]
            if pd.isna(fi) or fi > fecha_dia:
                continue
            if pd.notna(ff) and ff < fecha_dia:
                continue
            cap_base += float(r.get("capital_invertido", 0))

        # Ajuste acumulado por movimientos hasta ese día
        ajuste = 0.0
        if not movs.empty:
            movs_hasta = movs[movs["fecha"] <= fecha_dia]
            for _, m in movs_hasta.iterrows():
                if m["tipo"] == "DEVOLUCION":
                    ajuste -= float(m["importe"])
                elif m["tipo"] == "REINVERSION":
                    ajuste += float(m["importe"])

        capital_real = max(cap_base + ajuste, 0.0)
        suma_capital_diaria += capital_real

    capital_promedio = suma_capital_diaria / dias_mes
    return capital_promedio * TASA_ANUAL_MOTOCLICK / 12


def ajustar_ingreso_motoclick(d_fijos: pd.DataFrame, df_inv: pd.DataFrame, anio: int, mes: int) -> pd.DataFrame:
    """
    Reemplaza el ingreso_bruto y beneficio_empresa_mes de las filas de MotoClick
    en d_fijos por el ingreso real calculado con movimientos (capital promedio diario).
    El pago_inversor_mes NO se toca.
    """
    if d_fijos.empty or "activo" not in d_fijos.columns:
        return d_fijos
    df_mov = cargar_movimientos_motoclick()
    ingreso_real = ingreso_bruto_motoclick_mes(df_inv, df_mov, anio, mes)
    # Ingreso actual de MotoClick en d_fijos (sin ajuste)
    mask_mc = d_fijos["activo"].str.lower() == "motoclick"
    ingreso_actual = d_fijos.loc[mask_mc, "ingreso_bruto"].sum()
    pago_actual = d_fijos.loc[mask_mc, "pago_inversor_mes"].sum()
    if ingreso_actual == 0:
        return d_fijos
    # Distribuir el ingreso_real proporcionalmente entre las filas de MotoClick
    factor = ingreso_real / ingreso_actual
    d_fijos = d_fijos.copy()
    d_fijos.loc[mask_mc, "ingreso_bruto"] = d_fijos.loc[mask_mc, "ingreso_bruto"] * factor
    d_fijos.loc[mask_mc, "beneficio_empresa_mes"] = d_fijos.loc[mask_mc, "ingreso_bruto"] - d_fijos.loc[mask_mc, "pago_inversor_mes"]
    return d_fijos


def detalle_activo_mes(df_base: pd.DataFrame, activo: str, tasa_anual: float, anio: int, mes: int) -> pd.DataFrame:
    df_activo = filtrar_activo(df_base, activo)
    dias_mes = ultimo_dia_mes(anio, mes)
    filas = []
    for _, fila in df_activo.iterrows():
        dias = dias_activos_en_mes(fila.get("fecha_inversion"), fila.get("fecha_final_inversion"), anio, mes)
        if dias == 0:
            continue
        proporcion = dias / dias_mes
        capital = float(fila.get("capital_invertido", 0))
        ingreso_bruto = capital * tasa_anual / 12 * proporcion
        pago_inversor = capital * float(fila.get("interes_inversor_anual", 0)) / 12 * proporcion
        filas.append({
            "id_inversion": fila.get("id_inversion", ""), "inversor": fila.get("inversor", ""),
            "capital_invertido": capital, "fecha_inversion": fila.get("fecha_inversion"),
            "fecha_final_inversion": fila.get("fecha_final_inversion"), "dias_activos": dias,
            "dias_mes": dias_mes, "ingreso_bruto": ingreso_bruto,
            "pago_inversor_mes": pago_inversor, "beneficio_empresa_mes": ingreso_bruto - pago_inversor,
        })
    return pd.DataFrame(filas)


def detalle_prestamos_mes(df_base: pd.DataFrame, anio: int, mes: int) -> pd.DataFrame:
    """
    Igual que detalle_activo_mes, pero para préstamos internos hechos a Yuri. A diferencia de
    Paraguay/MotoClick/etc. — que usan una tasa fija global (TASA_ANUAL_PARAGUAY y similares) —
    aquí cada fila ya trae su propia interes_nota_anual (lo que gana la empresa) e
    interes_inversor_anual (lo que se le paga al inversor), porque cada préstamo puede tener una
    tasa distinta según lo que haya negociado cada inversor. Por eso se usan tal cual, fila a fila,
    en vez de una constante compartida.
    Se detecta igual que cualquier otro activo fijo: por subtipo_inversion o nombre_activo
    conteniendo 'prestamo'/'préstamo' (ver detectar_activo/filtrar_activo) — no por tipo_inversion,
    que normalmente es 'empresa_privada' como el resto de inversiones manuales (Paraguay, MotoClick...).
    """
    if df_base is None or df_base.empty:
        return pd.DataFrame()
    df_activo = filtrar_activo(df_base, "prestamo")
    if df_activo.empty:
        return pd.DataFrame()
    dias_mes = ultimo_dia_mes(anio, mes)
    filas = []
    for _, fila in df_activo.iterrows():
        dias = dias_activos_en_mes(fila.get("fecha_inversion"), fila.get("fecha_final_inversion"), anio, mes)
        if dias == 0:
            continue
        proporcion = dias / dias_mes
        capital = float(fila.get("capital_invertido", 0))
        ingreso_bruto = capital * float(fila.get("interes_nota_anual", 0)) / 12 * proporcion
        pago_inversor = capital * float(fila.get("interes_inversor_anual", 0)) / 12 * proporcion
        filas.append({
            "id_inversion": fila.get("id_inversion", ""), "inversor": fila.get("inversor", ""),
            "capital_invertido": capital, "fecha_inversion": fila.get("fecha_inversion"),
            "fecha_final_inversion": fila.get("fecha_final_inversion"), "dias_activos": dias,
            "dias_mes": dias_mes, "ingreso_bruto": ingreso_bruto,
            "pago_inversor_mes": pago_inversor, "beneficio_empresa_mes": ingreso_bruto - pago_inversor,
        })
    return pd.DataFrame(filas)




def capital_activo_en_fecha(df_base: pd.DataFrame, fecha_consulta, activo: Optional[str] = None, solo_real: bool = False) -> float:
    """Capital activo en una fecha aplicando la misma lógica que inversiones_activas_global.

    - NOTAS: todas las activas cuentan (nueva, reinversion, call).
    - NO-NOTAS: se excluyen las canceladas.
    - solo_real: si True, aplica además el filtro capital_nuevo_real=si (solo para notas).
    """
    fecha_consulta = pd.Timestamp(fecha_consulta).normalize()
    trabajo = df_base.copy()
    if activo:
        trabajo = filtrar_activo(trabajo, activo)

    activas = trabajo[
        (trabajo["fecha_inversion"].notna()) &
        (trabajo["fecha_inversion"] <= fecha_consulta) &
        (trabajo["fecha_final_inversion"].isna() | (trabajo["fecha_final_inversion"] >= fecha_consulta))
    ].copy()

    if activas.empty:
        return 0.0

    es_nota = activas["tipo_inversion"].apply(limpiar_texto) == "nota"
    es_cancelada = activas["tipo_operacion"].apply(limpiar_texto) == "cancelada"
    filtrado = activas[es_nota | (~es_nota & ~es_cancelada)].copy()

    if solo_real and "capital_nuevo_real" in filtrado.columns:
        filtrado = filtrado[filtrado["capital_nuevo_real"].astype(str).str.lower() == "si"].copy()

    return float(filtrado["capital_invertido"].sum()) if not filtrado.empty else 0.0


def calcular_intereses_acumulados_inversor(df_inv: pd.DataFrame, inversor: str, fecha_fin=None) -> dict:
    """
    Total de intereses generados por un inversor desde su primera posición hasta fecha_fin (hoy por defecto),
    usando la misma lógica de extractos (solo NUEVA/CANCELADA, ignorando fecha_final_inversion en filas NUEVA).
    Devuelve total acumulado, capital activo hoy y tasa media ponderada.
    """
    fecha_fin = pd.Timestamp(fecha_fin).normalize() if fecha_fin is not None else pd.Timestamp.today().normalize()
    df = df_inv.copy()
    df["inversor"] = df["inversor"].astype(str).str.strip().str.upper()
    df["tipo_operacion"] = df["tipo_operacion"].astype(str).str.strip().str.lower()
    df = df[df["inversor"] == inversor.strip().upper()]
    df = df[df["tipo_operacion"].isin(["nueva", "cancelada"])].copy()
    if df.empty:
        return {"total_intereses": 0.0, "capital_activo": 0.0, "tasa_media": 0.0, "fecha_inicio": None}

    df["fecha_inversion"] = parsear_fecha_robusta(df["fecha_inversion"])
    df["fecha_final_inversion"] = parsear_fecha_robusta(df["fecha_final_inversion"])
    df["capital_invertido"] = pd.to_numeric(df["capital_invertido"], errors="coerce").fillna(0)
    df["interes_inversor_anual"] = pd.to_numeric(df["interes_inversor_anual"], errors="coerce").fillna(0)

    fecha_inicio_global = df["fecha_inversion"].min()
    if pd.isna(fecha_inicio_global):
        return {"total_intereses": 0.0, "capital_activo": 0.0, "tasa_media": 0.0, "fecha_inicio": None}

    total_intereses = 0.0
    f_iter = pd.Timestamp(fecha_inicio_global.year, fecha_inicio_global.month, 1)
    while f_iter <= fecha_fin:
        dias_m = ultimo_dia_mes(f_iter.year, f_iter.month)
        inicio_mes = datetime(f_iter.year, f_iter.month, 1)
        fin_mes = datetime(f_iter.year, f_iter.month, dias_m)
        for _, row in df.iterrows():
            fi = row["fecha_inversion"]
            if pd.isna(fi):
                continue
            fi_dt = fi.to_pydatetime()
            ff = row["fecha_final_inversion"]
            if row["tipo_operacion"] == "cancelada" and pd.notna(ff):
                fecha_fin_row = min(ff.to_pydatetime(), fin_mes)
            else:
                fecha_fin_row = fin_mes
            inicio_calc = max(fi_dt, inicio_mes)
            fin_calc = min(fecha_fin_row, fin_mes, fecha_fin.to_pydatetime())
            if inicio_calc > fin_calc:
                continue
            dias = (fin_calc - inicio_calc).days + 1
            total_intereses += row["capital_invertido"] * row["interes_inversor_anual"] / 12 * dias / dias_m
        f_iter += pd.DateOffset(months=1)

    capital_activo = capital_activo_en_fecha(df_inv[df_inv["inversor"].astype(str).str.strip().str.upper() == inversor.strip().upper()], fecha_fin)
    activas_hoy = df[(df["fecha_inversion"] <= fecha_fin) & ((df["tipo_operacion"] != "cancelada") | (df["fecha_final_inversion"] >= fecha_fin))]
    tasa_media = float((activas_hoy["capital_invertido"] * activas_hoy["interes_inversor_anual"]).sum() / activas_hoy["capital_invertido"].sum()) if not activas_hoy.empty and activas_hoy["capital_invertido"].sum() > 0 else 0.0

    return {
        "total_intereses": round(total_intereses, 2),
        "capital_activo": capital_activo,
        "tasa_media": tasa_media,
        "fecha_inicio": fecha_inicio_global,
    }


def preparar_extracto_privado_inversor(contenido_bytes: bytes) -> bytes:
    """
    Quita del extracto, antes de dárselo a un inversor, cualquier dato que revele EN QUÉ
    está invertido (columnas 'Activo' y 'Estado intereses' de la hoja DETALLE), dejando
    solo mes, fecha de inversión, capital e interés del mes.

    IMPORTANTE: las filas de "CIERRE MENSUAL/ANUAL/FINAL" son celdas COMBINADAS que empiezan
    justo en la columna 'Activo' — si se borrara esa columna entera, se destruiría también
    la combinación y desaparecerían los totales de cierre. Por eso aquí NO se borran columnas:
    se vacía únicamente el contenido de las filas de DETALLE individuales (una por posición),
    dejando intactas las filas de cierre con sus totales de capital e interés.
    PORTADA ya es solo totales, no revela activos, así que se deja igual.
    """
    try:
        wb = load_workbook(BytesIO(contenido_bytes))
        if "DETALLE" in wb.sheetnames:
            ws = wb["DETALLE"]
            fila_cabecera = 4
            col_activo = col_estado = None
            for col_idx in range(1, ws.max_column + 1):
                valor = ws.cell(row=fila_cabecera, column=col_idx).value
                if valor and str(valor).strip().lower() == "activo":
                    col_activo = col_idx
                elif valor and str(valor).strip().lower() == "estado intereses":
                    col_estado = col_idx

            for fila_idx in range(fila_cabecera + 1, ws.max_row + 1):
                valor_col1 = ws.cell(row=fila_idx, column=1).value
                es_fila_cierre = valor_col1 and str(valor_col1).strip().upper().startswith("CIERRE")
                if es_fila_cierre:
                    continue  # no tocar las filas de resumen (mensual/anual/final)
                if col_activo:
                    ws.cell(row=fila_idx, column=col_activo).value = None
                if col_estado:
                    ws.cell(row=fila_idx, column=col_estado).value = None
        salida = BytesIO()
        wb.save(salida)
        return salida.getvalue()
    except Exception:
        return contenido_bytes  # si algo falla, mejor devolver el original que romper la descarga


def construir_historial_mensual_inversor(df_inv: pd.DataFrame, inversor: str, fecha_fin=None) -> list:
    """Desglose mes a mes (capital activo ese mes e interés generado ese mes) para UN inversor,
    usando la misma lógica de extractos (solo NUEVA/CANCELADA). No revela en qué activo está
    invertido el capital — solo capital e interés, que es toda la información a la que un
    inversor tiene derecho a acceder sobre sí mismo.
    """
    fecha_fin = pd.Timestamp(fecha_fin).normalize() if fecha_fin is not None else pd.Timestamp.today().normalize()
    df = df_inv.copy()
    df["inversor"] = df["inversor"].astype(str).str.strip().str.upper()
    df["tipo_operacion"] = df["tipo_operacion"].astype(str).str.strip().str.lower()
    df = df[df["inversor"] == inversor.strip().upper()]
    df = df[df["tipo_operacion"].isin(["nueva", "cancelada"])].copy()
    if df.empty:
        return []

    df["fecha_inversion"] = parsear_fecha_robusta(df["fecha_inversion"])
    df["fecha_final_inversion"] = parsear_fecha_robusta(df["fecha_final_inversion"])
    df["capital_invertido"] = pd.to_numeric(df["capital_invertido"], errors="coerce").fillna(0)
    df["interes_inversor_anual"] = pd.to_numeric(df["interes_inversor_anual"], errors="coerce").fillna(0)

    fecha_inicio_global = df["fecha_inversion"].min()
    if pd.isna(fecha_inicio_global):
        return []

    historial = []
    f_iter = pd.Timestamp(fecha_inicio_global.year, fecha_inicio_global.month, 1)
    while f_iter <= fecha_fin:
        dias_m = ultimo_dia_mes(f_iter.year, f_iter.month)
        inicio_mes = datetime(f_iter.year, f_iter.month, 1)
        fin_mes = datetime(f_iter.year, f_iter.month, dias_m)
        capital_mes = 0.0
        interes_mes = 0.0
        for _, row in df.iterrows():
            fi = row["fecha_inversion"]
            if pd.isna(fi):
                continue
            fi_dt = fi.to_pydatetime()
            ff = row["fecha_final_inversion"]
            if row["tipo_operacion"] == "cancelada" and pd.notna(ff):
                fecha_fin_row = min(ff.to_pydatetime(), fin_mes)
            else:
                fecha_fin_row = fin_mes
            inicio_calc = max(fi_dt, inicio_mes)
            fin_calc = min(fecha_fin_row, fin_mes, fecha_fin.to_pydatetime())
            if inicio_calc > fin_calc:
                continue
            dias = (fin_calc - inicio_calc).days + 1
            interes_mes += row["capital_invertido"] * row["interes_inversor_anual"] / 12 * dias / dias_m
            if inicio_mes <= fi_dt <= fin_mes or fi_dt <= inicio_mes:
                if row["tipo_operacion"] != "cancelada" or pd.isna(ff) or ff.to_pydatetime() >= inicio_mes:
                    capital_mes += row["capital_invertido"]
        if capital_mes > 0 or interes_mes > 0:
            historial.append({
                "anio": f_iter.year, "mes": f_iter.month,
                "capital": round(capital_mes, 2), "interes": round(interes_mes, 2),
            })
        f_iter += pd.DateOffset(months=1)
    return historial


def construir_contexto_ia_inversor(nombre_inversor, df_inv: pd.DataFrame, fecha_limite=None) -> str:
    """Construye el contexto EXCLUSIVO de uno o varios inversores para el asistente de IA del
    portal. Contiene ÚNICAMENTE el capital, la tasa e intereses de los inversores indicados —
    nunca datos del fondo, de otros inversores, ni de en qué activos está invertido el capital.

    'nombre_inversor' puede ser un string (un solo inversor, comportamiento de siempre) o una
    lista de strings (varios inversores autorizados — ver INVERSORES_ADICIONALES_VISIBLES),
    en cuyo caso se genera un bloque "=== POSICIÓN DE <nombre> ===" por cada uno.
    """
    hoy = pd.Timestamp.today().normalize()
    fecha_fin = pd.Timestamp(fecha_limite).normalize() if fecha_limite else hoy
    lista_inversores = [nombre_inversor] if isinstance(nombre_inversor, str) else list(nombre_inversor)

    meses_nombre = ["", "Enero","Febrero","Marzo","Abril","Mayo","Junio","Julio","Agosto","Septiembre","Octubre","Noviembre","Diciembre"]

    lineas = [f"Fecha de hoy: {hoy.strftime('%d/%m/%Y')}"]
    if len(lista_inversores) > 1:
        lineas.append(f"Inversores cuya posición puedes consultar en esta sesión: {', '.join(lista_inversores)}")
    else:
        lineas.append(f"Nombre del inversor (titular de esta sesión): {lista_inversores[0]}")

    for nombre in lista_inversores:
        datos = calcular_intereses_acumulados_inversor(df_inv, nombre, fecha_fin)
        historial = construir_historial_mensual_inversor(df_inv, nombre, fecha_fin)
        lineas.append("")
        lineas.append(f"=== POSICIÓN DE {nombre} ===")
        if datos["fecha_inicio"] is None:
            lineas.append("No se encontraron posiciones registradas a este nombre.")
            continue
        lineas.append(f"Capital activo hoy: ${datos['capital_activo']:,.2f}")
        lineas.append(f"Fecha de inicio de la inversión: {datos['fecha_inicio'].strftime('%d/%m/%Y')}")
        lineas.append(f"Tasa de interés anual media contratada: {datos['tasa_media']*100:.2f}%")
        lineas.append(f"Intereses totales acumulados desde el inicio hasta hoy: ${datos['total_intereses']:,.2f}")
        lineas.append("Historial mensual (capital activo e interés generado cada mes):")
        for h in historial:
            lineas.append(f"  {meses_nombre[h['mes']]} {h['anio']}: capital activo ${h['capital']:,.2f} — interés generado ${h['interes']:,.2f}")

    return "\n".join(lineas)


def preguntar_asistente_ia_inversor(pregunta: str, nombre_inversor, df_inv,
                                     historial_previo=None) -> str:
    """Asistente de IA EXCLUSIVO del portal del inversor. Responde solo sobre la(s) posición(es)
    de 'nombre_inversor' (string o lista — ver construir_contexto_ia_inversor). Nunca sobre el
    fondo, otros inversores no autorizados, ni en qué activos está invertido el capital.
    """
    import re as _re_inv, requests as _req_inv
    try:
        lista_inversores = [nombre_inversor] if isinstance(nombre_inversor, str) else list(nombre_inversor)
        nombre_sesion = lista_inversores[0]  # titular de la sesión, siempre el primero de la lista
        es_multi = len(lista_inversores) > 1
        etiqueta_inversores = " y ".join(lista_inversores)

        fecha_limite = None
        m_fecha = _re_inv.search(r"hasta\s+(?:el\s+)?(\d{1,2})[/\-](\d{1,2})(?:[/\-](\d{4}))?", pregunta.lower())
        if m_fecha:
            d, mo = int(m_fecha.group(1)), int(m_fecha.group(2))
            yr = int(m_fecha.group(3)) if m_fecha.group(3) else pd.Timestamp.today().year
            fecha_limite = f"{yr}-{mo:02d}-{d:02d}"

        ctx = construir_contexto_ia_inversor(lista_inversores, df_inv, fecha_limite=fecha_limite)

        historial = []
        mensajes_prev = historial_previo or []
        for m in mensajes_prev[-4:]:
            if isinstance(m["content"], str):
                historial.append({"role": m["role"], "content": m["content"]})
        historial.append({"role": "user", "content": f"DATOS AUTORIZADOS PARA ESTA SESIÓN:\n\n{ctx}\n\n---\nPREGUNTA: {pregunta}"})

        try:
            api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
        except Exception:
            api_key = ""
        if not api_key:
            import os as _os_key
            api_key = _os_key.environ.get("ANTHROPIC_API_KEY", "")

        if es_multi:
            bloque_alcance = f"""== LÍMITE ABSOLUTO E INQUEBRANTABLE DE LO QUE PUEDES RESPONDER ==
Esta sesión está autorizada EXCEPCIONALMENTE a consultar la posición de {etiqueta_inversores} — ningún otro inversor. Puedes hablar del capital activo, tasa de interés contratada, intereses generados mes a mes e histórico de CUALQUIERA de estos {len(lista_inversores)} inversores (los datos de ambos están en "DATOS AUTORIZADOS PARA ESTA SESIÓN"), y ayudar a interpretar sus extractos. Si la pregunta no deja claro de cuál de los dos habla, pregunta primero a cuál se refiere en vez de asumirlo o mezclar ambos."""
        else:
            bloque_alcance = f"""== LÍMITE ABSOLUTO E INQUEBRANTABLE DE LO QUE PUEDES RESPONDER ==
SOLO puedes hablar de la posición personal de {nombre_sesion}: su capital activo, su tasa de interés contratada, sus intereses generados mes a mes, el histórico de su posición, y ayuda para interpretar o encontrar información en sus propios extractos."""

        system_prompt_inversor = f"""Eres el asistente personal del portal de inversores de Chaparro Fernández Wealth Management. Estás hablando EXCLUSIVAMENTE con {nombre_sesion}, un inversor del fondo, a través de su portal privado.

{bloque_alcance}

NUNCA, bajo ninguna circunstancia, aunque te lo pidan de forma insistente, indirecta, hipotética, "solo para entender mejor", disfrazada de otra pregunta, o alegando cualquier motivo:
- Reveles en qué activos, notas, empresas o instrumentos está invertido el capital del fondo (Paraguay, Bolivia, MotoClick, Fútbol, Bitcoin, notas estructuradas, tickers, etc.). No confirmes ni desmientas si el capital de {etiqueta_inversores if es_multi else nombre_sesion} está en tal o cual activo.
- Reveles información del fondo en general: capital total gestionado, número de inversores, beneficio de la empresa, tasas de otros activos, estrategia de inversión, estructura societaria, ni cómo funciona internamente el negocio.
- Reveles absolutamente nada sobre inversores que NO estén en la lista autorizada de esta sesión ({etiqueta_inversores}): sus nombres, si existen, su capital, sus tasas, ni ninguna comparación.
- Reveles información de socios, gestión interna, cuentas, Google Drive, Excel, código de la aplicación, ni cómo se calculan internamente los datos.
- Des consejo de inversión, opiniones sobre mercados, ni recomendaciones financieras.

Si te preguntan por cualquiera de estos temas, responde amablemente que esa información no está disponible en el portal del inversor y que, si lo necesita, contacte directamente con Chaparro Fernández Wealth Management. No des pistas ni respuestas parciales que insinúen la respuesta.

== LO QUE SÍ PUEDES HACER ==
Con los datos que se te entregan en "DATOS AUTORIZADOS PARA ESTA SESIÓN" (exclusivamente los de {etiqueta_inversores}), para cada uno de ellos puedes:
- Decir su capital activo actual.
- Decir su tasa de interés contratada.
- Decir cuánto interés ha generado en total o en un mes/periodo concreto.
- Explicar su historial mensual.
- Ayudarle a entender su extracto.

== FORMATO ==
Responde SIEMPRE en español, de forma cercana y clara. Sé conciso. Importes con $ y 2 decimales. Fechas DD/MM/YYYY."""

        resp = _req_inv.post(
            "https://api.anthropic.com/v1/messages",
            headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
            json={
                "model": "claude-sonnet-4-5", "max_tokens": 1200,
                "system": system_prompt_inversor,
                "messages": historial,
            },
            timeout=60,
        )
        data = resp.json()
        _uso = data.get("usage", {}) or {}
        log_uso_ia(nombre_sesion, "inversor",
                   _uso.get("input_tokens", 0), _uso.get("output_tokens", 0))
        if "content" in data:
            textos = [b.get("text", "") for b in data["content"] if b.get("type") == "text"]
            return "\n".join(t for t in textos if t).strip() or "No he podido generar una respuesta."
        return f"Error del asistente: {data.get('error', {}).get('message', 'respuesta inesperada')}"
    except Exception as e:
        return f"Error: {e}"


def seccion_asistente_ia_inversor(nombre_inversor: str, df_inv, inversores_visibles=None):
    """Chat del asistente de IA personal dentro del portal del inversor.
    'inversores_visibles' es la lista completa de inversores que esta sesión puede consultar
    (normalmente solo [nombre_inversor]; para casos autorizados, ver INVERSORES_ADICIONALES_VISIBLES).

    OJO: se usa st.form + st.text_input en vez de st.chat_input a propósito. El chat_input de
    Streamlit queda SIEMPRE fijo/flotando en la parte inferior de la pantalla, fuera del sitio
    donde lo dibujamos en el código — eso era la causa real de que la app pareciera abrirse
    desplazada hacia abajo en móvil, y de que el cuadro de escribir pareciera pertenecer a la
    sección de Noticias (que va justo después) en vez de a esta. Con un formulario normal, todo
    el bloque (título, preguntas de ejemplo y cuadro de escribir) queda dentro de una tarjeta
    con borde propio, en el flujo normal de la página — consecutivo y visualmente separado de
    Noticias."""
    lista_inversores = inversores_visibles if inversores_visibles else [nombre_inversor]
    es_multi = len(lista_inversores) > 1

    key_chat = f"chat_ia_inversor_{nombre_inversor}"
    if key_chat not in st.session_state:
        st.session_state[key_chat] = []

    tarjeta_asistente = st.container(border=True)
    with tarjeta_asistente:
        st.markdown("### 💬 Tu asistente personal")
        if es_multi:
            st.caption(f"Pregunta sobre el capital, el interés o los extractos de {' o '.join(lista_inversores)}. Este asistente solo conoce estas posiciones — no tiene acceso a información del fondo ni de ningún otro inversor.")
        else:
            st.caption("Pregunta sobre tu capital, tu interés o tus extractos. Este asistente solo conoce tu propia posición — no tiene acceso a información del fondo ni de otros inversores.")

        for msg in st.session_state[key_chat]:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        if not st.session_state[key_chat]:
            st.markdown("**Ejemplos:**")
            sugs = ["¿Cuánto capital tengo activo?", "¿Cuánto interés he ganado este año?", "¿Cuál es mi tasa contratada?"]
            cols = st.columns(3)
            for i, s in enumerate(sugs):
                if cols[i].button(s, key=f"sug_ia_inv_{i}", use_container_width=True):
                    st.session_state[key_chat].append({"role": "user", "content": s})
                    st.rerun()

        # Cuadro de escribir, justo debajo de las preguntas de ejemplo, dentro del flujo
        # normal de la tarjeta (nunca flotante — ver nota al inicio de la función).
        with st.form(f"form_chat_ia_inv_{nombre_inversor}", clear_on_submit=True):
            col_txt, col_btn = st.columns([5, 1])
            with col_txt:
                texto_pregunta = st.text_input(
                    "Pregúntale a tu asistente personal",
                    placeholder="Escribe aquí tu pregunta a tu asistente personal...",
                    label_visibility="collapsed",
                    key=f"txt_chat_ia_inv_{nombre_inversor}",
                )
            with col_btn:
                enviar = st.form_submit_button("Enviar ➤", use_container_width=True)

        if enviar and texto_pregunta and texto_pregunta.strip():
            st.session_state[key_chat].append({"role": "user", "content": texto_pregunta.strip()})
            st.rerun()

        if st.session_state[key_chat] and st.session_state[key_chat][-1]["role"] == "user":
            ultima = st.session_state[key_chat][-1]["content"]
            with st.chat_message("assistant"):
                with st.spinner("Consultando tu posición..."):
                    mensajes_prev = st.session_state[key_chat][:-1]
                    respuesta = preguntar_asistente_ia_inversor(ultima, lista_inversores, df_inv, historial_previo=mensajes_prev)
                    st.markdown(_md_seguro(respuesta))
            st.session_state[key_chat].append({"role": "assistant", "content": respuesta})
            # Lleva la vista hasta la pregunta+respuesta que se acaban de generar, para que no
            # haya que subir manualmente a verla.
            _forzar_scroll_a_ultimo_mensaje_chat()

        if st.session_state[key_chat]:
            if st.button("🗑️ Limpiar conversación", key=f"btn_limpiar_ia_inv_{nombre_inversor}"):
                st.session_state[key_chat] = []
                st.rerun()


def _construir_datos_demo_inversor():
    """
    Genera un set de datos 100% ficticio con la misma forma exacta que produce
    cargar_excel_completo(), para alimentar el portal de inversor en modo DEMO.

    Objetivo: dejar hacer una demo comercial a un inversor potencial usando la interfaz
    real del portal (mismos KPIs, extracto descargable, asistente IA...), sin tocar
    nunca el Excel real ni exponer datos de ningún inversor existente. Las fechas se
    calculan siempre relativas a "hoy" para que la demo no quede desfasada con el tiempo.
    """
    hoy = pd.Timestamp.today().normalize()
    filas = [
        {
            "id_inversion": "DEMO-001", "inversor": "DEMO", "tipo_operacion": "cancelada",
            "tipo_inversion": "nota", "subtipo_inversion": "ESTRUCTURADA",
            "nombre_activo": "Nota estructurada — cesta bancaria EU",
            "fecha_inversion": hoy - pd.DateOffset(months=14),
            "fecha_final_inversion": hoy - pd.DateOffset(months=8),
            "capital_invertido": 50000.0, "interes_inversor_anual": 0.10, "interes_nota_anual": 0.10,
        },
        {
            "id_inversion": "DEMO-002", "inversor": "DEMO", "tipo_operacion": "cancelada",
            "tipo_inversion": "nota", "subtipo_inversion": "ESTRUCTURADA",
            "nombre_activo": "Nota estructurada — cesta tecnológica US",
            "fecha_inversion": hoy - pd.DateOffset(months=12),
            "fecha_final_inversion": hoy - pd.DateOffset(months=3),
            "capital_invertido": 75000.0, "interes_inversor_anual": 0.10, "interes_nota_anual": 0.10,
        },
        {
            "id_inversion": "DEMO-003", "inversor": "DEMO", "tipo_operacion": "nueva",
            "tipo_inversion": "nota", "subtipo_inversion": "ESTRUCTURADA",
            "nombre_activo": "Nota estructurada — cesta bancaria global",
            "fecha_inversion": hoy - pd.DateOffset(months=5),
            "fecha_final_inversion": pd.NaT,
            "capital_invertido": 100000.0, "interes_inversor_anual": 0.10, "interes_nota_anual": 0.10,
        },
        {
            "id_inversion": "DEMO-004", "inversor": "DEMO", "tipo_operacion": "nueva",
            "tipo_inversion": "fijo", "subtipo_inversion": "ACTIVO FIJO",
            "nombre_activo": "Activo fijo — desarrollo inmobiliario",
            "fecha_inversion": hoy - pd.DateOffset(months=10),
            "fecha_final_inversion": pd.NaT,
            "capital_invertido": 30000.0, "interes_inversor_anual": 0.10, "interes_nota_anual": 0.0,
        },
        {
            "id_inversion": "DEMO-005", "inversor": "DEMO", "tipo_operacion": "nueva",
            "tipo_inversion": "fijo", "subtipo_inversion": "ACTIVO FIJO",
            "nombre_activo": "Activo fijo — cripto",
            "fecha_inversion": hoy - pd.DateOffset(months=4),
            "fecha_final_inversion": pd.NaT,
            "capital_invertido": 20000.0, "interes_inversor_anual": 0.10, "interes_nota_anual": 0.0,
        },
    ]
    df_inv = pd.DataFrame(filas)
    for col in ["id_inversion", "inversor", "tipo_operacion", "tipo_inversion", "subtipo_inversion",
                "nombre_activo", "metodo_calculo", "activo_generador_interes", "capital_nuevo_real",
                "cuenta_cobro", "motivo"]:
        if col not in df_inv.columns:
            df_inv[col] = ""
        df_inv[col] = df_inv[col].fillna("").astype(str).str.strip()
    df_inv["fecha_inversion"] = pd.to_datetime(df_inv["fecha_inversion"])
    df_inv["fecha_final_inversion"] = pd.to_datetime(df_inv["fecha_final_inversion"])
    for col in ["capital_invertido", "interes_inversor_anual", "interes_nota_anual"]:
        df_inv[col] = pd.to_numeric(df_inv[col], errors="coerce").fillna(0)
    df_inv["periodicidad_meses"] = 1
    df_inv["pago_intereses"] = "REINVIERTE"

    df_cal = pd.DataFrame(columns=["nota", "fecha", "tipo_evento"])
    df_control = pd.DataFrame()
    return df_inv, df_cal, df_control


def _inversores_visibles_para(nombre_inversor: str) -> list:
    """Devuelve la lista de inversores cuya posición puede consultar este login: siempre el
    propio titular primero, más cualquier inversor adicional autorizado explícitamente."""
    extra = INVERSORES_ADICIONALES_VISIBLES.get(str(nombre_inversor).strip().upper(), [])
    vistos = []
    for n in [nombre_inversor] + list(extra):
        if n not in vistos:
            vistos.append(n)
    return vistos


def _mostrar_aviso_seguridad_pendiente_inversor(nombre_inversor: str):
    """Aviso grande y persistente arriba del portal para los inversores piloto con 2FA
    disponible: no bloquea nada (siguen viendo todo su portal con normalidad), pero deja bien
    claro qué les falta hacer hasta que completen los dos pasos (cambiar contraseña + activar
    verificación en dos pasos)."""
    nombre_up = str(nombre_inversor).strip().upper()
    if nombre_up not in INVERSORES_CON_2FA:
        return

    ya_cambio_password = _fila_usuario(_leer_hoja_usuarios(), nombre_inversor, "inversor") is not None
    tiene_2fa_activo = _2fa_activo(nombre_inversor, "inversor")

    if ya_cambio_password and tiene_2fa_activo:
        return  # ya hizo los dos pasos: no se muestra nada más

    pasos_pendientes = []
    if not ya_cambio_password:
        pasos_pendientes.append("**1. Cambia tu contraseña** — despliega '🔑 Cambiar mi contraseña' en el menú de la izquierda.")
    if not tiene_2fa_activo:
        numero = "2" if pasos_pendientes else "1"
        pasos_pendientes.append(f"**{numero}. Activa la verificación en dos pasos** — despliega '🔐 Verificación en dos pasos' en el menú de la izquierda y sigue los pasos con tu email.")

    st.markdown(
        f"""
        <div style="background:linear-gradient(135deg,#fff4e0,#ffe4d6);border:2px solid #e8862c;
                    border-radius:16px;padding:20px 24px;margin-bottom:24px;">
            <div style="font-size:17px;font-weight:800;color:#7a3a00;margin-bottom:10px;">
                ⚠️ Antes de nada — completa la configuración de seguridad de tu cuenta
            </div>
            <div style="font-size:14px;color:#7a3a00;line-height:1.7;">
                {"<br>".join(pasos_pendientes)}
            </div>
            <div style="font-size:13px;color:#9a5a1a;margin-top:10px;">
                El resto de tu portal ya funciona con normalidad — esto es solo un recordatorio hasta que lo completes.
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def seccion_portal_inversor(nombre_inversor: str):
    """Portal de acceso limitado para un inversor: solo ve su propia posición y, en casos
    excepcionales autorizados explícitamente (ver INVERSORES_ADICIONALES_VISIBLES), la de otro
    inversor concreto adicional — nunca la de nadie más."""
    if st.session_state.get("forzar_scroll_arriba_pendiente"):
        _forzar_scroll_arriba_del_todo()
        st.session_state.forzar_scroll_arriba_pendiente = False

    es_demo = str(nombre_inversor).strip().upper() == "DEMO"
    if es_demo:
        df_inv, df_cal, df_control = _construir_datos_demo_inversor()
    else:
        df_inv, df_cal, df_control = cargar_excel_completo()
    hoy = pd.Timestamp.today().normalize()

    inversores_visibles = _inversores_visibles_para(nombre_inversor)

    st.header(f"👋 Bienvenido/a, {nombre_inversor}")
    if not es_demo:
        _mostrar_aviso_seguridad_pendiente_inversor(nombre_inversor)
    if len(inversores_visibles) > 1:
        st.caption("Este es tu portal personal. Tienes acceso adicional autorizado a otra posición concreta — nadie más puede ver esta información.")
        inversor_activo = st.selectbox(
            "Viendo la información de:", inversores_visibles, key="portal_inversor_activo_selector",
        )
    else:
        st.caption("Este es tu portal personal — solo tú puedes ver esta información. Ningún otro inversor tiene acceso a tu posición.")
        inversor_activo = nombre_inversor

    datos = calcular_intereses_acumulados_inversor(df_inv, inversor_activo, hoy)

    if datos["fecha_inicio"] is None:
        st.warning(f"No se encontraron posiciones registradas a nombre de {inversor_activo}.")
        return

    c1, c2, c3 = st.columns(3)
    with c1:
        tarjeta_kpi("💰 Capital activo hoy", fmt(datos["capital_activo"]), f"Invirtiendo desde {datos['fecha_inicio'].strftime('%d/%m/%Y')}", "normal")
    with c2:
        tarjeta_kpi("📈 Intereses totales ganados", fmt(datos["total_intereses"]), "Acumulado desde tu primera inversión", "positivo")
    with c3:
        tarjeta_kpi("🎯 Tu rentabilidad contratada", fmt_pct(datos["tasa_media"]), "Tasa anual media ponderada", "positivo")

    st.markdown("---")
    st.markdown(f"### 📄 Extracto de {inversor_activo}" if len(inversores_visibles) > 1 else "### 📄 Tu extracto")

    # Solo se pueden consultar extractos de meses YA CERRADOS: el mes en curso queda
    # excluido (aún no ha terminado). El último mes disponible es siempre el mes anterior
    # al actual — el día 1 de cada mes, ese mes anterior ya está disponible.
    if hoy.month == 1:
        anio_max, mes_max = hoy.year - 1, 12
    else:
        anio_max, mes_max = hoy.year, hoy.month - 1

    col_mes, col_anio = st.columns(2)
    with col_anio:
        anio_extracto = st.selectbox("Año", list(range(2024, anio_max + 1)), index=anio_max - 2024)
    with col_mes:
        meses_disponibles = list(range(1, mes_max + 1)) if anio_extracto == anio_max else list(range(1, 13))
        mes_extracto = st.selectbox(
            "Mes", meses_disponibles, index=len(meses_disponibles) - 1,
            format_func=lambda m: ["Enero","Febrero","Marzo","Abril","Mayo","Junio","Julio","Agosto","Septiembre","Octubre","Noviembre","Diciembre"][m-1],
        )

    if st.button("🔎 Ver extracto", type="primary"):
        archivos = generar_extractos(df_inv, "Un inversor", inversor_activo, anio_extracto, mes_extracto)
        if not archivos:
            st.warning("No se encontró extracto para ese mes — puede que no hubiera capital activo entonces.")
        else:
            nombre_archivo, contenido, _ = archivos[0]
            contenido_privado = preparar_extracto_privado_inversor(contenido)
            st.session_state["portal_extracto_actual"] = (nombre_archivo, contenido_privado)
            st.success("Extracto listo.")

    extracto_actual = st.session_state.get("portal_extracto_actual")
    if extracto_actual:
        nombre_archivo, contenido = extracto_actual
        st.download_button("⬇️ Descargar extracto (Excel)", contenido, file_name=nombre_archivo,
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            type="primary")

    st.markdown("---")
    seccion_asistente_ia_inversor(nombre_inversor, df_inv, inversores_visibles=inversores_visibles)

    st.markdown("---")
    st.markdown("### 📰 Noticias")
    _widget_busqueda_libre_noticias(key_prefix=f"inversor_{nombre_inversor}")



def total_pagado_activo_desde_inicio(df_base: pd.DataFrame, activo: str, tasa_anual: float) -> float:
    df_activo = filtrar_activo(df_base, activo)
    if df_activo.empty:
        return 0.0
    fecha_min = df_activo["fecha_inversion"].dropna().min()
    if pd.isna(fecha_min):
        return 0.0
    hoy = pd.Timestamp.today().normalize()
    total = 0.0
    anio, mes = fecha_min.year, fecha_min.month
    while (anio < hoy.year) or (anio == hoy.year and mes <= hoy.month):
        detalle = detalle_activo_mes(df_base, activo, tasa_anual, anio, mes)
        if not detalle.empty:
            total += detalle["pago_inversor_mes"].sum()
        mes += 1
        if mes == 13:
            mes = 1
            anio += 1
    return float(total)


def total_ingresado_activo_desde_inicio(df_base: pd.DataFrame, activo: str, tasa_anual: float) -> float:
    df_activo = filtrar_activo(df_base, activo)
    if df_activo.empty:
        return 0.0
    fecha_min = df_activo["fecha_inversion"].dropna().min()
    if pd.isna(fecha_min):
        return 0.0
    hoy = pd.Timestamp.today().normalize()
    total = 0.0
    anio, mes = fecha_min.year, fecha_min.month
    while (anio < hoy.year) or (anio == hoy.year and mes <= hoy.month):
        detalle = detalle_activo_mes(df_base, activo, tasa_anual, anio, mes)
        if not detalle.empty:
            total += detalle["ingreso_bruto"].sum()
        mes += 1
        if mes == 13:
            mes = 1
            anio += 1
    return float(total)


def totales_activo_desde_inicio(df_base: pd.DataFrame, activo: str, tasa_anual: float) -> dict:
    """Recorre mes a mes, desde la primera inversión del activo hasta el mes actual incluido,
    y devuelve el ingreso acumulado de la compañía y el pago acumulado a inversores, con
    desglose por año.

    Usa exactamente la misma fuente y la misma lógica que total_ingresado_activo_desde_inicio()
    y total_pagado_activo_desde_inicio() (detalle_activo_mes), pero en UNA sola pasada en vez de
    dos, y devolviendo además el detalle por año. Sirve para alimentar el contexto del asistente
    IA sin duplicar el coste de cálculo."""
    vacio = {"ingresado": 0.0, "pagado": 0.0, "beneficio": 0.0, "por_anio": {}, "desde": None}
    df_activo = filtrar_activo(df_base, activo)
    if df_activo.empty:
        return vacio
    fecha_min = df_activo["fecha_inversion"].dropna().min()
    if pd.isna(fecha_min):
        return vacio
    hoy = pd.Timestamp.today().normalize()
    ing_total = 0.0
    pag_total = 0.0
    por_anio = {}
    anio, mes = int(fecha_min.year), int(fecha_min.month)
    while (anio < hoy.year) or (anio == hoy.year and mes <= hoy.month):
        detalle = detalle_activo_mes(df_base, activo, tasa_anual, anio, mes)
        if not detalle.empty:
            ing_mes = float(detalle["ingreso_bruto"].sum())
            pag_mes = float(detalle["pago_inversor_mes"].sum())
            ing_total += ing_mes
            pag_total += pag_mes
            acum = por_anio.setdefault(anio, {"ingresado": 0.0, "pagado": 0.0})
            acum["ingresado"] += ing_mes
            acum["pagado"] += pag_mes
        mes += 1
        if mes == 13:
            mes = 1
            anio += 1
    return {
        "ingresado": ing_total,
        "pagado": pag_total,
        "beneficio": ing_total - pag_total,
        "por_anio": por_anio,
        "desde": pd.Timestamp(fecha_min),
    }


# =========================
# NOTAS
# =========================
def normalizar_cuenta(valor):
    texto = str(valor).strip().lower()
    if texto in ["jordi", "cuenta jordi"]:
        return "JORDI"
    if texto in ["compañia", "compania", "empresa", "sociedad"]:
        return "COMPAÑÍA"
    return "SIN CLASIFICAR"


def extraer_numero_nota(nombre_activo):
    if pd.isna(nombre_activo):
        return pd.NA
    match = re.search(r"NOTA[_\s]?(\d+)", str(nombre_activo).strip().upper())
    return int(match.group(1)) if match else pd.NA


def filtrar_notas(df_base: pd.DataFrame) -> pd.DataFrame:
    trabajo = df_base.copy()
    if "tipo_inversion" in trabajo.columns:
        trabajo = trabajo[trabajo["tipo_inversion"].astype(str).str.lower() == "nota"].copy()
    if "nombre_activo" not in trabajo.columns:
        trabajo["nombre_activo"] = ""
    trabajo["nota_num"] = trabajo["nombre_activo"].apply(extraer_numero_nota)
    trabajo["nota_num"] = pd.to_numeric(trabajo["nota_num"], errors="coerce").astype("Int64")
    if "activo_generador_interes" in trabajo.columns:
        trabajo = trabajo[trabajo["activo_generador_interes"].astype(str).str.upper() == "SI"].copy()
    if "cuenta_cobro" not in trabajo.columns:
        trabajo["cuenta_cobro"] = "SIN CLASIFICAR"
    trabajo["cuenta_cobro"] = trabajo["cuenta_cobro"].apply(normalizar_cuenta)
    return trabajo


def inversiones_activas_para_nota(df_base: pd.DataFrame, nota: int, fecha_pago) -> pd.DataFrame:
    fecha_pago = pd.Timestamp(fecha_pago).normalize()
    trabajo = filtrar_notas(df_base)
    return trabajo[(trabajo["nota_num"] == nota) & (trabajo["fecha_inversion"].notna()) & (trabajo["fecha_inversion"] <= fecha_pago) & (trabajo["fecha_final_inversion"].isna() | (trabajo["fecha_final_inversion"] >= fecha_pago))].copy()


def pagos_notas_mes(df_cal: pd.DataFrame, anio: int, mes: int) -> pd.DataFrame:
    if df_cal.empty:
        return pd.DataFrame()
    return df_cal[(df_cal["tipo_evento"] == "PAGO") & (df_cal["fecha"].notna()) & (df_cal["fecha"].dt.year == anio) & (df_cal["fecha"].dt.month == mes)].copy().sort_values(["fecha", "nota"])


def pagos_notas_hasta_hoy(df_cal: pd.DataFrame) -> pd.DataFrame:
    hoy = pd.Timestamp.today().normalize()
    return df_cal[(df_cal["tipo_evento"] == "PAGO") & (df_cal["fecha"].notna()) & (df_cal["fecha"] <= hoy)].copy().sort_values(["fecha", "nota"])


def obtener_observacion_previa_nota(df_cal: pd.DataFrame, nota: int, fecha_pago):
    if df_cal is None or df_cal.empty:
        return None
    fecha_pago = pd.Timestamp(fecha_pago).normalize()
    obs = df_cal[(df_cal["nota"] == nota) & (df_cal["tipo_evento"] == "OBSERVACION") & (df_cal["fecha"].notna()) & (df_cal["fecha"] <= fecha_pago)].copy().sort_values("fecha")
    return None if obs.empty else obs.iloc[-1]["fecha"]


def normalizar_barrera(valor):
    if pd.isna(valor):
        return None
    try:
        valor = float(valor)
        return valor / 100 if valor > 1 else valor
    except Exception:
        return None


@st.cache_data(show_spinner=False, ttl=3600)
def obtener_cierre_ticker_fecha(ticker: str, fecha_objetivo):
    if yf is None:
        return None
    try:
        fecha_objetivo = pd.Timestamp(fecha_objetivo).normalize()
        inicio = fecha_objetivo - pd.Timedelta(days=10)
        fin = fecha_objetivo + pd.Timedelta(days=2)
        data = yf.download(str(ticker).strip().upper(), start=inicio.strftime("%Y-%m-%d"), end=fin.strftime("%Y-%m-%d"), progress=False, auto_adjust=False)
        if data.empty:
            return None
        if isinstance(data.columns, pd.MultiIndex):
            if "Close" in data.columns.get_level_values(0):
                cierres = data["Close"].iloc[:, 0].dropna()
            else:
                return None
        else:
            if "Close" not in data.columns:
                return None
            cierres = data["Close"].dropna()
        cierres.index = pd.to_datetime(cierres.index).normalize()
        cierres = cierres[cierres.index <= fecha_objetivo]
        if cierres.empty:
            return None
        return float(cierres.iloc[-1])
    except Exception:
        return None


def columna_barrera_control(df_control: pd.DataFrame, preferida="contingency"):
    for col in [preferida, "contingency", "barrera_cupon", "barrera_capital"]:
        if col in df_control.columns:
            return col
    return None


# ══════════════════════════════════════════════════════════════════════
# ANÁLISIS DE COMPAÑÍA + SIMULACIÓN MONTE CARLO PARA NOTAS ESTRUCTURADAS
# ══════════════════════════════════════════════════════════════════════

@st.cache_data(show_spinner=False, ttl=900)
def obtener_datos_fundamentales(ticker: str) -> dict:
    """
    Precio actual, objetivo de precio de consenso de analistas, volatilidad histórica
    anualizada, próxima fecha de resultados, salto histórico medio en earnings y
    datos generales de la compañía. Todo dato real de mercado (yfinance) — nada
    inventado ni "previsto" por IA.

    NOTA DE FIABILIDAD: el endpoint de "info" de Yahoo Finance (de donde sale el precio
    objetivo de analistas) es conocido por ser inestable, especialmente desde servidores
    en la nube — a veces Yahoo limita o bloquea temporalmente estas peticiones. Por eso
    se reintenta y se distingue entre "sin cobertura de analistas" (dato real: no hay
    analistas cubriendo esta acción) y "no se pudo consultar" (fallo temporal de conexión).
    """
    resultado = {
        "ticker": ticker, "precio_actual": None, "sector": None, "industria": None, "nombre": None,
        "market_cap": None, "target_medio": None, "target_alto": None, "target_bajo": None,
        "n_analistas": None, "recomendacion": None, "volatilidad_anual_pct": None,
        "proxima_fecha_resultados": None, "proxima_fecha_resultados_dt": None,
        "salto_medio_earnings_pct": None, "n_earnings_medidos": 0,
        "retorno_implicito_analistas_pct": None,
        "variacion_1m_pct": None, "variacion_ytd_pct": None,
        "error": None, "aviso_analistas": None,
        # --- Añadido para la ficha de subyacente rediseñada (rango 52 sem., PER, gráfico) ---
        "precio_max_52sem": None, "precio_min_52sem": None,
        "posicion_en_rango_52sem_pct": None, "pct_desde_max_52sem": None,
        "trailing_pe": None, "forward_pe": None,
        "variacion_dia_pct": None,
        "historico_precios": None,  # lista [{"fecha": "YYYY-MM-DD", "close": float}, ...] ~1 año, para el gráfico
    }
    if yf is None:
        resultado["error"] = "yfinance no disponible"
        return resultado
    try:
        t = yf.Ticker(ticker)

        # El endpoint "info" es el más inestable de yfinance/Yahoo — reintenta con espera creciente.
        info = {}
        for intento in range(4):
            try:
                info = t.info or {}
                if info and len(info) > 5:
                    break
            except Exception:
                pass
            if intento < 3:
                time.sleep(1.5 * (intento + 1))  # 1.5s, 3s, 4.5s — da tiempo real a que se libere el límite

        if not info or len(info) <= 5:
            resultado["aviso_analistas"] = (
                "Yahoo Finance no devolvió datos de analistas/fundamentales en este momento "
                "(suele ser un límite temporal de peticiones, no que falten datos de verdad). Prueba de nuevo en un rato."
            )
        else:
            resultado["nombre"] = info.get("longName") or info.get("shortName")
            resultado["sector"] = info.get("sector")
            resultado["industria"] = info.get("industry")
            resultado["market_cap"] = info.get("marketCap")
            resultado["target_medio"] = info.get("targetMeanPrice")
            resultado["target_alto"] = info.get("targetHighPrice")
            resultado["target_bajo"] = info.get("targetLowPrice")
            resultado["n_analistas"] = info.get("numberOfAnalystOpinions")
            resultado["recomendacion"] = info.get("recommendationKey")
            resultado["trailing_pe"] = info.get("trailingPE")
            resultado["forward_pe"] = info.get("forwardPE")
            if info.get("targetMeanPrice") is None:
                resultado["aviso_analistas"] = "Esta compañía no tiene cobertura de analistas en Yahoo Finance (dato real, no es un fallo)."

        # Si el precio objetivo sigue sin salir (endpoint "info" caído/limitado), prueba una fuente
        # alternativa de yfinance que consulta un endpoint distinto de Yahoo — a veces uno está
        # limitado y el otro no, así que merece la pena intentarlo antes de rendirse.
        if resultado["target_medio"] is None:
            try:
                targets = t.analyst_price_targets
                if targets:
                    resultado["target_medio"] = targets.get("mean")
                    resultado["target_alto"] = targets.get("high")
                    resultado["target_bajo"] = targets.get("low")
                    if resultado["target_medio"] is not None:
                        resultado["aviso_analistas"] = None
            except Exception:
                pass

        hist = None
        for intento in range(2):
            try:
                hist = t.history(period="2y")  # 2 años: 1 para volatilidad, más margen para earnings pasados
                if hist is not None and not hist.empty:
                    break
            except Exception:
                pass
            if intento == 0:
                time.sleep(1.0)

        if hist is not None and not hist.empty:
            cierres_todo = hist["Close"].dropna()
            cierres = cierres_todo.tail(253)  # ~1 año para la volatilidad "normal" (sin contar los saltos de earnings aparte)
            resultado["precio_actual"] = float(cierres.iloc[-1])
            retornos = np.log(cierres / cierres.shift(1)).dropna()
            resultado["volatilidad_anual_pct"] = float(retornos.std() * np.sqrt(252) * 100)
            if len(cierres) > 21:
                resultado["variacion_1m_pct"] = float((cierres.iloc[-1] / cierres.iloc[-22] - 1) * 100)
            primer_dia_anio = cierres[cierres.index.year == cierres.index[-1].year]
            if len(primer_dia_anio) > 1:
                resultado["variacion_ytd_pct"] = float((cierres.iloc[-1] / primer_dia_anio.iloc[0] - 1) * 100)
            if len(cierres) > 1:
                resultado["variacion_dia_pct"] = float((cierres.iloc[-1] / cierres.iloc[-2] - 1) * 100)

            # Rango de 52 semanas (mismo tramo de 253 sesiones ~1 año usado para volatilidad, así
            # que la barra visual y el % de recorrido son consistentes con el resto de la ficha).
            resultado["precio_max_52sem"] = float(cierres.max())
            resultado["precio_min_52sem"] = float(cierres.min())
            if resultado["precio_max_52sem"] > resultado["precio_min_52sem"]:
                resultado["posicion_en_rango_52sem_pct"] = float(
                    (resultado["precio_actual"] - resultado["precio_min_52sem"])
                    / (resultado["precio_max_52sem"] - resultado["precio_min_52sem"]) * 100
                )
            resultado["pct_desde_max_52sem"] = float(
                (resultado["precio_actual"] / resultado["precio_max_52sem"] - 1) * 100
            )

            # Histórico de cierres para el gráfico de la ficha (últimos ~12 meses, serializado a
            # lista de dicts para que sobreviva al cacheo de st.cache_data sin problemas).
            resultado["historico_precios"] = [
                {"fecha": fecha.strftime("%Y-%m-%d"), "close": float(valor)}
                for fecha, valor in cierres.items()
            ]

            # Salto histórico medio el día después de resultados — así el riesgo de earnings se mide
            # con el movimiento REAL de esta acción en sus últimos informes, no con la vol. media anual.
            try:
                earn = t.get_earnings_dates(limit=8)
                if earn is not None and not earn.empty:
                    idx_naive = cierres_todo.index.tz_localize(None) if cierres_todo.index.tz is not None else cierres_todo.index
                    idx_naive = idx_naive.normalize()
                    saltos = []
                    for fecha_earn in earn.index:
                        fe = pd.Timestamp(fecha_earn)
                        if fe.tz is not None:
                            fe = fe.tz_localize(None)
                        fe = fe.normalize()
                        mask_despues = idx_naive > fe
                        mask_antes = idx_naive <= fe
                        if not mask_despues.any() or not mask_antes.any():
                            continue
                        precio_antes = float(cierres_todo.iloc[np.where(mask_antes)[0][-1]])
                        precio_despues = float(cierres_todo.iloc[np.where(mask_despues)[0][0]])
                        saltos.append(abs(precio_despues / precio_antes - 1))
                    if saltos:
                        resultado["salto_medio_earnings_pct"] = float(np.mean(saltos) * 100)
                        resultado["n_earnings_medidos"] = len(saltos)
            except Exception:
                pass  # si yfinance no da earnings_dates para este ticker, seguimos sin ese dato (no rompe el resto)
        else:
            resultado["error"] = "No se pudo obtener el histórico de precios (fallo de conexión con Yahoo Finance, prueba de nuevo)."

        # Retorno implícito del consenso de analistas, asumiendo el horizonte habitual de ~12 meses
        if resultado.get("target_medio") and resultado.get("precio_actual"):
            resultado["retorno_implicito_analistas_pct"] = float(
                (resultado["target_medio"] / resultado["precio_actual"] - 1) * 100
            )

        try:
            cal = t.calendar
            fecha_dt = None
            if isinstance(cal, dict) and cal.get("Earnings Date"):
                fechas = cal["Earnings Date"]
                fecha_dt = fechas[0] if isinstance(fechas, list) else fechas
            elif hasattr(cal, "empty") and not cal.empty and "Earnings Date" in cal.index:
                fecha_dt = cal.loc["Earnings Date"].iloc[0]
            if fecha_dt is not None:
                resultado["proxima_fecha_resultados"] = str(fecha_dt)
                resultado["proxima_fecha_resultados_dt"] = pd.Timestamp(fecha_dt)
        except Exception:
            pass
    except Exception as e:
        resultado["error"] = str(e)
    return resultado


def calcular_matriz_correlacion_tickers(tickers_datos: list) -> np.ndarray:
    """Matriz de correlación histórica (retornos diarios log, ~1 año) entre los tickers de una
    nota worst-of, calculada a partir del histórico de precios que YA se descargó para cada ticker
    (historico_precios, de obtener_datos_fundamentales) — sin llamadas de red adicionales.
    Si algún ticker no tiene histórico suficiente, esa fila/columna se queda en 0 de correlación
    con el resto (supuesto conservador: nunca se inventa una correlación que no se puede medir)."""
    n = len(tickers_datos)
    if n <= 1:
        return np.eye(max(n, 1))

    series = {}
    for td in tickers_datos:
        hist = td.get("historico_precios")
        if not hist:
            continue
        df_h = pd.DataFrame(hist)
        df_h["fecha"] = pd.to_datetime(df_h["fecha"])
        df_h = df_h.drop_duplicates(subset="fecha").set_index("fecha").sort_index()
        retornos = np.log(df_h["close"] / df_h["close"].shift(1)).dropna()
        if len(retornos) > 20:
            series[td["ticker"]] = retornos

    corr = np.eye(n)
    tickers_ok = list(series.keys())
    if len(tickers_ok) >= 2:
        df_retornos = pd.DataFrame(series).dropna()
        if len(df_retornos) > 20:
            corr_calc = df_retornos.corr()
            for i, tdi in enumerate(tickers_datos):
                for j, tdj in enumerate(tickers_datos):
                    ti, tj = tdi["ticker"], tdj["ticker"]
                    if ti in corr_calc.index and tj in corr_calc.columns:
                        corr[i, j] = corr_calc.loc[ti, tj]
    return corr


def simular_montecarlo_nota(tickers_datos: list, dias_hasta_eventos: list, n_simulaciones: int = 5000,
                             tipo_proteccion: str = "barrera", tiene_memoria: bool = False,
                             matriz_correlacion: np.ndarray = None) -> dict:
    """
    Simulación Monte Carlo (movimiento browniano geométrico con deriva configurable + shocks de
    earnings) para una nota worst-of con uno o varios tickers subyacentes.

    tickers_datos: lista de dicts con, por ticker:
        - precio_actual, precio_inicial, volatilidad_anual_pct, barrera_cupon_pct, call_level_pct, barrera_capital_pct
        - drift_anual_pct (opcional, default 0): deriva anual asumida en la simulación. 0 = paseo
          aleatorio puro (más conservador, no asume ni sube ni baja). Puede venir del consenso de
          analistas o de una convicción manual — NUNCA se inventa sola, la decide quien usa el comparador.
        - dias_earnings (opcional, lista de int): días desde hoy en los que cae un informe de
          resultados dentro del horizonte de la nota (el próximo conocido + proyecciones trimestrales).
        - salto_earnings_pct (opcional, float): magnitud media histórica del movimiento del día
          después de resultados de ESTE ticker (dato real, de obtener_datos_fundamentales). Si se
          da, esos días concretos llevan una sacudida extra de esa magnitud, en vez de repartir el
          riesgo de earnings uniformemente a lo largo de todo el año (que es lo que hace un GBM puro
          y por eso subestima el riesgo real alrededor de esas fechas concretas).
    dias_hasta_eventos: lista de dicts {dias, tipo} ordenada cronológicamente, tipo="cupon" o "call" o "vencimiento"
    tipo_proteccion: "barrera" (europea) o "buffer".
        - "barrera": si el peor activo rompe su nivel de protección al vencimiento, el fondo recibe
          EXACTAMENTE la performance del peor activo (pérdida total desde el nivel inicial, sin suelo).
        - "buffer": el fondo está protegido hasta el tamaño del buffer (1 - barrera_capital_pct); si
          la caída del peor activo supera ese buffer, solo se pierde el EXCESO sobre el buffer.
          Para la misma barrera nominal, un buffer siempre pierde menos (o igual) que una barrera.
    tiene_memoria: si True, un periodo en el que NO se cumple la barrera de cupón no se pierde sin
        más — se acumula, y en el primer periodo futuro en el que SÍ se cumpla la barrera, se cobran
        de golpe TODOS los cupones acumulados desde el último pago (efecto retroactivo). Solo se
        pierden definitivamente los periodos acumulados que quedan sin recuperar hasta el vencimiento
        (la "cola" final si la nota nunca vuelve a cumplir la barrera antes de vencer).
    matriz_correlacion: matriz n_tickers×n_tickers de correlación entre los subyacentes (típicamente
        de calcular_matriz_correlacion_tickers, con la correlación histórica real). Si es None, se
        asume independencia total entre tickers (matriz identidad) — comportamiento previo, más
        conservador/simplificado en notas de un único ticker no cambia nada.

    LIMITACIÓN HONESTA: usa volatilidad histórica (no implícita de opciones, más precisa pero no
    disponible gratis) y correlación histórica de precios (no correlación implícita de mercado de
    opciones, que tampoco está disponible gratis). Es una estimación con supuestos simplificados
    apoyados en datos reales, no una certeza.
    """
    np.random.seed(42)
    n_tickers = len(tickers_datos)
    max_dias = max(e["dias"] for e in dias_hasta_eventos) if dias_hasta_eventos else 1
    pasos = max(max_dias, 1)
    dt = 1 / 252

    if matriz_correlacion is None:
        matriz_correlacion = np.eye(n_tickers)

    # Cholesky para generar shocks diarios CORRELACIONADOS entre tickers (si hay más de uno) —
    # antes cada ticker se simulaba de forma completamente independiente, lo cual subestima el
    # riesgo real de una nota worst-of: si los subyacentes suelen moverse juntos (correlación
    # positiva, lo habitual entre acciones del mismo sector), la probabilidad de que TODOS caigan
    # a la vez es mayor que si se tratan como independientes.
    try:
        L = np.linalg.cholesky(matriz_correlacion)
    except np.linalg.LinAlgError:
        # La matriz no es semidefinida positiva (puede pasar con correlaciones estimadas de pocos
        # datos, o si se pasa una matriz manual inconsistente) — se reconstruye la matriz de
        # correlación válida más cercana recortando autovalores negativos a un mínimo positivo,
        # en vez de romper la simulación.
        valores, vectores = np.linalg.eigh(matriz_correlacion)
        valores_clip = np.clip(valores, 1e-6, None)
        matriz_psd = vectores @ np.diag(valores_clip) @ vectores.T
        diag_sqrt = np.sqrt(np.diag(matriz_psd))
        matriz_correlacion = matriz_psd / np.outer(diag_sqrt, diag_sqrt)  # renormaliza a diagonal 1
        L = np.linalg.cholesky(matriz_correlacion)

    z_indep = np.random.normal(0, 1, size=(n_simulaciones, pasos, n_tickers))
    z_corr = z_indep @ L.T  # correlaciona entre tickers en cada paso; los pasos siguen siendo independientes entre sí

    precios_simulados = {}
    for k, td in enumerate(tickers_datos):
        s0 = td["precio_actual"]
        sigma = td["volatilidad_anual_pct"] / 100
        mu = td.get("drift_anual_pct", 0.0) / 100
        incrementos = (mu - 0.5 * sigma**2) * dt + sigma * np.sqrt(dt) * z_corr[:, :, k]
        # Shock de earnings: en vez de repartir ese riesgo uniformemente en el año (lo que hace un
        # GBM puro), lo concentramos en los días concretos donde cae un informe de resultados,
        # con la magnitud REAL que esta acción se ha movido en sus últimos informes.
        dias_earnings = td.get("dias_earnings") or []
        salto_pct = td.get("salto_earnings_pct")
        if salto_pct and dias_earnings:
            for de in dias_earnings:
                idx = min(max(int(de), 1), pasos) - 1
                incrementos[:, idx] += np.random.normal(0, salto_pct / 100, size=n_simulaciones)
        log_precios = np.log(s0) + np.cumsum(incrementos, axis=1)
        precios_simulados[td["ticker"]] = np.exp(log_precios)  # shape (n_sim, pasos)

    resultados_eventos = []
    ya_llamada = np.zeros(n_simulaciones, dtype=bool)
    cupones_acumulados = np.zeros(n_simulaciones)  # solo se usa si tiene_memoria=True
    for evento in sorted(dias_hasta_eventos, key=lambda e: e["dias"]):
        idx_dia = min(evento["dias"], max_dias) - 1
        cumple_todas_cupon = np.ones(n_simulaciones, dtype=bool)
        cumple_todas_call = np.ones(n_simulaciones, dtype=bool)
        cumple_todas_capital = np.ones(n_simulaciones, dtype=bool)
        for td in tickers_datos:
            precio_dia = precios_simulados[td["ticker"]][:, idx_dia]
            precio_inicial = td["precio_inicial"]
            cumple_todas_cupon &= (precio_dia >= precio_inicial * td["barrera_cupon_pct"])
            cumple_todas_call &= (precio_dia >= precio_inicial * td["call_level_pct"])
            cumple_todas_capital &= (precio_dia >= precio_inicial * td["barrera_capital_pct"])

        if evento["tipo"] == "cupon":
            activo = ~ya_llamada
            cumple = cumple_todas_cupon & activo
            prob = float(np.mean(cumple))  # prob. de cobrar ALGO este periodo concreto (con o sin memoria)
            if tiene_memoria:
                # Quien cumple cobra el cupón de este periodo MÁS todo lo acumulado desde el último pago.
                # Quien no cumple (y sigue activo) acumula un periodo más para el futuro.
                pagos_este_evento = np.where(cumple, 1.0 + cupones_acumulados, 0.0)
                cupones_acumulados = np.where(cumple, 0.0, np.where(activo, cupones_acumulados + 1.0, cupones_acumulados))
            else:
                pagos_este_evento = np.where(cumple, 1.0, 0.0)
            cupones_pagados_esperado = float(np.mean(pagos_este_evento))
            resultados_eventos.append({
                "tipo": "cupon", "dias": evento["dias"], "probabilidad": prob,
                "cupones_pagados_esperado": cupones_pagados_esperado,
            })
        elif evento["tipo"] == "call":
            call_ahora = cumple_todas_call & ~ya_llamada
            prob = float(np.mean(call_ahora))
            resultados_eventos.append({"tipo": "call", "dias": evento["dias"], "probabilidad": prob})
            ya_llamada = ya_llamada | call_ahora
        elif evento["tipo"] == "vencimiento":
            no_llamadas = ~ya_llamada
            incumple = ~cumple_todas_capital & no_llamadas
            prob_perdida = float(np.mean(incumple))
            # Pérdida media (%) en los escenarios que sí incumplen, usando el peor ticker (worst-of).
            # Con barrera: pérdida = caída completa del peor activo. Con buffer: solo el exceso sobre
            # el tamaño del buffer (1 - barrera_capital_pct) del PEOR activo en cada escenario.
            perdida_pct_promedio = 0.0
            if incumple.sum() > 0:
                peor_performance = np.ones(n_simulaciones)
                peor_buffer_size = np.zeros(n_simulaciones)  # buffer del ticker que resulta ser el peor en cada escenario
                for td in tickers_datos:
                    precio_final = precios_simulados[td["ticker"]][:, idx_dia]
                    performance = precio_final / td["precio_inicial"]
                    es_nuevo_peor = performance < peor_performance
                    peor_buffer_size = np.where(es_nuevo_peor, 1 - td["barrera_capital_pct"], peor_buffer_size)
                    peor_performance = np.minimum(peor_performance, performance)
                caida = 1 - peor_performance
                if tipo_proteccion == "buffer":
                    perdida_fraccion = np.maximum(0.0, caida - peor_buffer_size)
                else:  # "barrera": pérdida total desde el nivel inicial, sin suelo
                    perdida_fraccion = caida
                perdida_pct_promedio = float(np.mean(perdida_fraccion[incumple]) * 100)
            resultados_eventos.append({"tipo": "vencimiento_perdida_capital", "dias": evento["dias"], "probabilidad": prob_perdida, "perdida_pct_promedio": perdida_pct_promedio})

    prob_call_total = sum(e["probabilidad"] for e in resultados_eventos if e["tipo"] == "call")
    evento_venc = next((e for e in resultados_eventos if e["tipo"] == "vencimiento_perdida_capital"), None)
    prob_perdida_capital = evento_venc["probabilidad"] if evento_venc else 0.0
    perdida_pct_promedio = evento_venc.get("perdida_pct_promedio", 0.0) if evento_venc else 0.0
    eventos_cupon_calc = [e for e in resultados_eventos if e["tipo"] == "cupon"]
    cupones_totales_esperados = float(sum(e["cupones_pagados_esperado"] for e in eventos_cupon_calc))
    cupones_perdidos_definitivo_esperado = float(np.mean(cupones_acumulados)) if tiene_memoria else None

    return {
        "eventos": resultados_eventos,
        "probabilidad_call_total": prob_call_total,
        "probabilidad_perdida_capital": prob_perdida_capital,
        "perdida_pct_promedio_si_incumple": perdida_pct_promedio,
        "n_simulaciones": n_simulaciones,
        "tipo_proteccion": tipo_proteccion,
        "tiene_memoria": tiene_memoria,
        "cupones_totales_esperados": cupones_totales_esperados,
        "cupones_perdidos_definitivo_esperado": cupones_perdidos_definitivo_esperado,
        "matriz_correlacion": matriz_correlacion,
    }




def evaluar_nota_en_fecha(df_control: pd.DataFrame, nota: int, fecha_obs, preferida="contingency") -> tuple[str, pd.DataFrame]:
    """
    Evalúa una nota en una fecha de observación.

    Lógica definitiva:
    - Si la fecha de observación es futura: PENDIENTE y cuenta como cobro previsto.
    - Si ya llegó la observación y no hay dato de precio: SIN DATO, pero NO se trata como negativa.
    - Solo es NEGATIVA si ya llegó la observación y existe precio real por debajo de la barrera.
    """
    hoy = pd.Timestamp.today().normalize()

    if fecha_obs is None or pd.isna(fecha_obs):
        return "SIN_OBSERVACION", pd.DataFrame()

    fecha_obs = pd.Timestamp(fecha_obs).normalize()

    if df_control is None or df_control.empty:
        return "SIN_CONTROL", pd.DataFrame()

    sub = df_control[df_control.get("nota") == nota].copy()
    if sub.empty:
        return "SIN_CONTROL", pd.DataFrame()

    barrera_col = columna_barrera_control(sub, preferida=preferida)
    if barrera_col is None or "precio_compra" not in sub.columns or "ticker" not in sub.columns:
        return "SIN_COLUMNAS", pd.DataFrame()

    filas = []

    # Si la observación aún no ha llegado, no se descarga precio histórico.
    # Se marca como pendiente y se mantiene el cobro previsto.
    if fecha_obs > hoy:
        for _, row in sub.iterrows():
            ticker = row.get("ticker", "")
            compra = pd.to_numeric(row.get("precio_compra"), errors="coerce")
            barrera_pct = normalizar_barrera(row.get(barrera_col))
            precio_barrera = float(compra) * float(barrera_pct) if pd.notna(compra) and barrera_pct is not None else None
            filas.append({
                "ticker": ticker,
                "precio_compra": float(compra) if pd.notna(compra) else None,
                "barrera_%": barrera_pct,
                "precio_barrera": precio_barrera,
                "cierre_usado": None,
                "estado": "PENDIENTE",
            })
        return "PENDIENTE", pd.DataFrame(filas)

    hay_negativa_real = False
    hay_sin_dato = False

    for _, row in sub.iterrows():
        ticker = row.get("ticker", "")
        compra = pd.to_numeric(row.get("precio_compra"), errors="coerce")
        barrera_pct = normalizar_barrera(row.get(barrera_col))

        if pd.isna(compra) or barrera_pct is None:
            filas.append({
                "ticker": ticker,
                "precio_compra": float(compra) if pd.notna(compra) else None,
                "barrera_%": barrera_pct,
                "precio_barrera": None,
                "cierre_usado": None,
                "estado": "FALTAN DATOS",
            })
            hay_sin_dato = True
            continue

        precio_barrera = float(compra) * float(barrera_pct)
        cierre = obtener_cierre_ticker_fecha(ticker, fecha_obs)

        if cierre is None:
            filas.append({
                "ticker": ticker,
                "precio_compra": float(compra),
                "barrera_%": barrera_pct,
                "precio_barrera": precio_barrera,
                "cierre_usado": None,
                "estado": "SIN DATO",
            })
            hay_sin_dato = True
            continue

        estado = "OK" if cierre >= precio_barrera else "NO OK"
        if estado == "NO OK":
            hay_negativa_real = True

        filas.append({
            "ticker": ticker,
            "precio_compra": float(compra),
            "barrera_%": barrera_pct,
            "precio_barrera": precio_barrera,
            "cierre_usado": cierre,
            "estado": estado,
        })

    detalle = pd.DataFrame(filas)

    if hay_negativa_real:
        return "NEGATIVA", detalle
    if hay_sin_dato:
        return "SIN DATO", detalle
    return "POSITIVA", detalle


def resumen_detalle_observacion(detalle_obs: pd.DataFrame) -> str:
    if detalle_obs is None or detalle_obs.empty:
        return ""
    partes = []
    for _, row in detalle_obs.iterrows():
        ticker = row.get("ticker", "")
        estado = row.get("estado", "")
        cierre = row.get("cierre_usado", None)
        barrera = row.get("precio_barrera", None)
        if pd.notna(cierre) and pd.notna(barrera):
            partes.append(f"{ticker}: {estado} cierre {float(cierre):.2f} / barrera {float(barrera):.2f}")
        else:
            partes.append(f"{ticker}: {estado}")
    return " | ".join(partes)


def detectar_periodicidad_nota(df_cal: pd.DataFrame, nota: int, fecha_pago) -> int:
    """
    Detecta cada cuántos meses paga una nota (mensual=1, trimestral=3...)
    mirando el espaciado REAL entre eventos PAGO en CALENDARIO_NOTAS,
    en vez de depender de una columna periodicidad_meses que no existe en el Excel.
    Por defecto devuelve 1 (mensual) si no hay suficiente historial para deducirlo.
    """
    if df_cal is None or df_cal.empty:
        return 1
    fecha_pago = pd.Timestamp(fecha_pago).normalize()
    pagos_nota = df_cal[
        (df_cal["nota"] == nota) &
        (df_cal["tipo_evento"] == "PAGO") &
        (df_cal["fecha"].notna())
    ].sort_values("fecha")
    if len(pagos_nota) < 2:
        return 1
    anteriores = pagos_nota[pagos_nota["fecha"] < fecha_pago]
    if not anteriores.empty:
        dias = (fecha_pago - anteriores.iloc[-1]["fecha"]).days
    else:
        posteriores = pagos_nota[pagos_nota["fecha"] > fecha_pago]
        if posteriores.empty:
            return 1
        dias = (posteriores.iloc[0]["fecha"] - fecha_pago).days
    meses = round(dias / 30.44)
    return max(meses, 1)


def preparar_detalle_notas(df_inv: pd.DataFrame, df_pagos: pd.DataFrame, df_cal: pd.DataFrame | None = None, df_control: pd.DataFrame | None = None) -> pd.DataFrame:
    filas = []
    cache_observaciones = {}

    for _, evento in df_pagos.iterrows():
        nota = evento.get("nota")
        fecha_pago = evento.get("fecha")

        if pd.isna(nota) or pd.isna(fecha_pago):
            continue

        nota_int = int(nota)
        fecha_pago = pd.Timestamp(fecha_pago).normalize()
        fecha_obs = obtener_observacion_previa_nota(df_cal, nota_int, fecha_pago) if df_cal is not None else None

        resultado_obs = "NO_EVALUADA"
        detalle_obs = pd.DataFrame()

        # Por defecto el cobro cuenta como previsto.
        # Solo se elimina si la observación es NEGATIVA real.
        ingreso_habilitado = True

        if fecha_obs is not None and df_control is not None and not df_control.empty:
            clave = (nota_int, pd.Timestamp(fecha_obs).normalize())
            if clave not in cache_observaciones:
                cache_observaciones[clave] = evaluar_nota_en_fecha(df_control, nota_int, fecha_obs, preferida="contingency")
            resultado_obs, detalle_obs = cache_observaciones[clave]

            if resultado_obs == "NEGATIVA":
                ingreso_habilitado = False

        activas = inversiones_activas_para_nota(df_inv, nota_int, fecha_pago)

        for _, fila in activas.iterrows():
            capital = float(fila.get("capital_invertido", 0))
            periodicidad = detectar_periodicidad_nota(df_cal, nota_int, fecha_pago)
            cobro_teorico = capital * float(fila.get("interes_nota_anual", 0)) / 12 * periodicidad
            cobro_compania = cobro_teorico if ingreso_habilitado else 0.0

            # Tratamiento especial Chaparro Fernández:
            # - CF es la sociedad gestora, no un inversor externo.
            # - Su interes_inversor_anual es siempre 0% porque no se pagan a sí mismos.
            # - Todo el cobro de la nota es beneficio íntegro de la empresa.
            # - Si se excluye Chaparro, la fila ya no llega aquí porque se filtra antes.
            es_chaparro = bool(fila.get("es_chaparro_fernandez", False)) or es_chaparro_fernandez_row(fila)
            if es_chaparro:
                pago_inversor = 0.0
                beneficio_empresa = cobro_compania
                tratamiento_chaparro = "INTERNO: pago = 0, beneficio = cobro nota"
            else:
                pago_inversor = capital * float(fila.get("interes_inversor_anual", 0)) / 12 * periodicidad
                beneficio_empresa = cobro_compania - pago_inversor
                tratamiento_chaparro = "NO"

            filas.append({
                "fecha_pago": fecha_pago,
                "nota": nota_int,
                "fecha_observacion_usada": fecha_obs,
                "resultado_observacion": resultado_obs,
                "detalle_observacion": resumen_detalle_observacion(detalle_obs),
                "ingreso_habilitado": "SI" if ingreso_habilitado else "NO",
                "id_inversion": fila.get("id_inversion", ""),
                "inversor": fila.get("inversor", ""),
                "cuenta_cobro": fila.get("cuenta_cobro", "SIN CLASIFICAR"),
                "es_chaparro_fernandez": es_chaparro,
                "tratamiento_chaparro": tratamiento_chaparro,
                "capital_invertido": capital,
                "interes_nota_anual": fila.get("interes_nota_anual", 0),
                "interes_inversor_anual": fila.get("interes_inversor_anual", 0),
                "cobro_teorico_compania": cobro_teorico,
                "cobro_compania": cobro_compania,
                "pago_inversor": pago_inversor,
                "beneficio_empresa": beneficio_empresa,
            })

    return pd.DataFrame(filas)


def _filtrar_notas_activas_en_mes(df_inv: pd.DataFrame, inicio_mes, fin_mes) -> pd.DataFrame:
    """Filas de NOTAS (NUEVA/CANCELADA) con posicion activa en algun momento del mes.
    Fuente unica reutilizada por pago_inversores_notas_mes y por el prorrateo de ingresos."""
    return df_inv[
        (df_inv["tipo_inversion"].apply(limpiar_texto) == "nota") &
        (df_inv["tipo_operacion"].apply(limpiar_texto).str.upper().isin(["NUEVA", "CANCELADA"])) &
        (df_inv["fecha_inversion"].notna()) &
        (df_inv["fecha_inversion"] <= fin_mes) &
        (df_inv["fecha_final_inversion"].isna() | (df_inv["fecha_final_inversion"] >= inicio_mes))
    ].copy()


def _pago_inversor_nota_devengo(row, inicio_mes, fin_mes, dias_mes) -> float:
    """Pago devengado a UN inversor de UNA fila de nota en el mes (pro-rata de dias +
    tramo de tipos Biscafe/Crowe Bolivia/JR Real Estate). Fuente unica: la usan tanto
    pago_inversores_notas_mes (total agregado) como el detalle prorrateado de ingresos,
    para que el pago al inversor sea siempre identico se mire por donde se mire."""
    INVERSORES_TRAMO = {"ROBERTO VISCAFE", "CROWE BOLIVIA", "JR REAL ESTATE"}
    FIN_T1 = pd.Timestamp("2026-01-31")
    INI_T2 = pd.Timestamp("2026-02-01")
    FIN_T2 = pd.Timestamp("2026-06-30")
    INI_T3 = pd.Timestamp("2026-07-01")

    if es_chaparro_fernandez_row(row):
        return 0.0
    capital = float(row.get("capital_invertido", 0) or 0)
    inicio_calc = max(row["fecha_inversion"], inicio_mes)
    fin_calc = fin_mes if pd.isna(row["fecha_final_inversion"]) else min(row["fecha_final_inversion"], fin_mes)
    if inicio_calc > fin_calc:
        return 0.0
    dias = (fin_calc - inicio_calc).days + 1
    inv_upper = str(row.get("inversor", "")).strip().upper()
    if inv_upper in INVERSORES_TRAMO:
        pago = 0.0
        # Tramo 1: hasta 31/01/2026 -> 5%
        if inicio_calc <= FIN_T1:
            d1 = (min(fin_calc, FIN_T1) - inicio_calc).days + 1
            pago += (capital * 0.05 / 12) * d1 / dias_mes
        # Tramo 2: 01/02/2026 - 30/06/2026 -> 7.5%
        ini_t2 = max(inicio_calc, INI_T2)
        fin_t2 = min(fin_calc, FIN_T2)
        if ini_t2 <= fin_t2:
            d2 = (fin_t2 - ini_t2).days + 1
            pago += (capital * 0.075 / 12) * d2 / dias_mes
        # Tramo 3: desde 01/07/2026 -> 10%
        if fin_calc >= INI_T3:
            ini_t3 = max(inicio_calc, INI_T3)
            d3 = (fin_calc - ini_t3).days + 1
            pago += (capital * 0.10 / 12) * d3 / dias_mes
        return pago
    tasa = float(row.get("interes_inversor_anual", 0) or 0)
    return (capital * tasa / 12) * dias / dias_mes


def pago_inversores_notas_mes(df_inv: pd.DataFrame, anio: int, mes: int) -> float:
    """Calcula el pago a inversores de notas por DEVENGO MENSUAL.

    Logica alineada con extractos:
    - Solo cuenta tipo_operacion NUEVA y CANCELADA (igual que generar_extractos).
    - Las REINVERSIONES no generan pago propio: el interes ya esta en la operacion NUEVA origen.
    - Excluye Chaparro Fernandez (tasa 0%, no reciben pago).
    - Aplica pro-rata de dias para el primer y ultimo mes de cada posicion.
    - Aplica tramo de tipos para Biscafe/Crowe Bolivia (5% hasta ene2026, 7.5% desde feb2026).
    """
    import calendar as _cal
    dias_mes = _cal.monthrange(anio, mes)[1]
    inicio_mes = pd.Timestamp(anio, mes, 1)
    fin_mes = pd.Timestamp(anio, mes, dias_mes)

    df_notas = _filtrar_notas_activas_en_mes(df_inv, inicio_mes, fin_mes)

    total_pago = 0.0
    for _, row in df_notas.iterrows():
        total_pago += _pago_inversor_nota_devengo(row, inicio_mes, fin_mes, dias_mes)

    return total_pago


def detalle_notas_mes_prorrateado(df_inv: pd.DataFrame, anio: int, mes: int, df_cal: pd.DataFrame | None = None, df_control: pd.DataFrame | None = None) -> pd.DataFrame:
    """Detalle de INGRESOS de notas por DEVENGO MENSUAL (prorrateo), en vez de por evento de
    cobro real del calendario.

    Por que existe: algunas notas cobran trimestral (o semestralmente) en vez de mensual. Con
    la logica normal (calendario), todo el cobro del trimestre aparece de golpe en el mes del
    pago y los otros meses aparecen "vacios", haciendo bailar la rentabilidad mensual. Aqui se
    reconoce cada mes su parte proporcional de ese cobro (capital x interes_nota_anual / 12,
    pro-rateado por dias activos en el mes), independientemente de si ese mes coincide con un
    pago real en el calendario.

    Universo de capital (INGRESOS): igual que el cobro real (inversiones_activas_para_nota /
    filtrar_notas) -- incluye posiciones NUEVA, CANCELADA y REINVERSION, porque el capital
    reinvertido sigue generando interes real de la nota para la empresa. Si aqui solo se
    contaran NUEVA/CANCELADA (como en el pago a inversores), el prorrateo infravaloraria el
    ingreso de cualquier nota con reinversiones. Tambien aplica el mismo filtro
    activo_generador_interes == "SI" que usa filtrar_notas.

    Barrera / cupon condicional (fix): estas notas son "Contingent Income Notes" -- si la
    observacion mas reciente hasta fin de mes salio NEGATIVA (barrera rota), NO se devenga
    ingreso ese mes para esa nota, igual que ya pasa en el calculo por calendario. Sin este
    chequeo, el prorrateo devengaba ingreso teorico incluso en notas con la barrera ya rota
    (p.ej. notas con alerta "Nota negativa real"), inflando el total.

    Alcance (a peticion de Yuri): SOLO afecta a los INGRESOS de la nota (cobro_compania). El pago
    al inversor NO se prorratea de forma distinta: se reutiliza el mismo devengo mensual de
    siempre (_pago_inversor_nota_devengo), y las filas REINVERSION siguen sin generar pago propio
    (el interes ya esta en la operacion NUEVA origen), igual que en pago_inversores_notas_mes.
    Esto es solo para visualizacion en el Dashboard general y en la vista de Notas: no toca
    extractos, comparador de notas ni activos fijos (futbol, motoclick, paraguay, bolivia,
    bitcoin).
    """
    import calendar as _cal
    dias_mes = _cal.monthrange(anio, mes)[1]
    inicio_mes = pd.Timestamp(anio, mes, 1)
    fin_mes = pd.Timestamp(anio, mes, dias_mes)

    notas_todas = filtrar_notas(df_inv)
    activas_ingreso = notas_todas[
        notas_todas["fecha_inversion"].notna()
        & (notas_todas["fecha_inversion"] <= fin_mes)
        & (notas_todas["fecha_final_inversion"].isna() | (notas_todas["fecha_final_inversion"] >= inicio_mes))
    ].copy()

    cache_barrera = {}

    def _ingreso_habilitado_nota(nota_int):
        """Ultima observacion conocida hasta fin de mes: si salio NEGATIVA, sin ingreso este mes."""
        if nota_int in cache_barrera:
            return cache_barrera[nota_int]
        habilitado = True
        if df_cal is not None and df_control is not None and not df_control.empty:
            fecha_obs = obtener_observacion_previa_nota(df_cal, nota_int, fin_mes)
            if fecha_obs is not None:
                resultado_obs, _ = evaluar_nota_en_fecha(df_control, nota_int, fecha_obs, preferida="contingency")
                if resultado_obs == "NEGATIVA":
                    habilitado = False
        cache_barrera[nota_int] = habilitado
        return habilitado

    filas = []
    for _, row in activas_ingreso.iterrows():
        capital = float(row.get("capital_invertido", 0) or 0)
        inicio_calc = max(row["fecha_inversion"], inicio_mes)
        fin_calc = fin_mes if pd.isna(row["fecha_final_inversion"]) else min(row["fecha_final_inversion"], fin_mes)
        if inicio_calc > fin_calc:
            continue
        dias = (fin_calc - inicio_calc).days + 1
        tasa_nota = float(row.get("interes_nota_anual", 0) or 0)

        try:
            nota_int = int(row.get("nota_num"))
        except (TypeError, ValueError):
            nota_int = None
        ingreso_habilitado = _ingreso_habilitado_nota(nota_int) if nota_int is not None else True
        cobro_compania = (capital * tasa_nota / 12) * dias / dias_mes if ingreso_habilitado else 0.0

        es_reinversion = str(row.get("tipo_operacion", "")).strip().upper() == "REINVERSION"
        es_chaparro = es_chaparro_fernandez_row(row)
        if es_reinversion:
            # La reinversion no genera pago propio: el interes ya esta en la operacion NUEVA origen.
            pago_inversor = 0.0
        else:
            pago_inversor = _pago_inversor_nota_devengo(row, inicio_mes, fin_mes, dias_mes)
        beneficio_empresa = cobro_compania if es_chaparro else (cobro_compania - pago_inversor)

        filas.append({
            "fecha_pago": pd.NaT,
            "nota": row.get("nota_num", ""),
            "fecha_observacion_usada": pd.NaT,
            "detalle_observacion": "",
            "id_inversion": row.get("id_inversion", ""),
            "inversor": row.get("inversor", ""),
            "cuenta_cobro": row.get("cuenta_cobro", "SIN CLASIFICAR"),
            "es_chaparro_fernandez": es_chaparro,
            "tratamiento_chaparro": "INTERNO: pago = 0, beneficio = cobro nota" if es_chaparro else "NO",
            "capital_invertido": capital,
            "interes_nota_anual": tasa_nota,
            "interes_inversor_anual": row.get("interes_inversor_anual", 0),
            "cobro_teorico_compania": (capital * tasa_nota / 12) * dias / dias_mes,
            "cobro_compania": cobro_compania,
            "pago_inversor": pago_inversor,
            "beneficio_empresa": beneficio_empresa,
            "resultado_observacion": "PRORRATEADO (devengo mensual)" if ingreso_habilitado else "NEGATIVA (barrera rota, sin ingreso este mes)",
            "ingreso_habilitado": "SI" if ingreso_habilitado else "NO",
        })

    return pd.DataFrame(filas)


def interes_devengado_no_cobrado_notas(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, anio: int | None = None, mes: int | None = None) -> pd.DataFrame:
    """Interes corrido por nota: lo que cada nota ya ha generado como INGRESO de la empresa
    (cobro_compania / interes_nota_anual) desde el ULTIMO PAGO REAL cobrado (evento PAGO en
    CALENDARIO_NOTAS con fecha <= fin del mes seleccionado) hasta fin de ese mes, y que
    todavia no se ha cobrado porque el proximo PAGO cae mas adelante (tipico de notas
    trimestrales/semestrales).

    anio/mes: mes del dashboard a analizar (igual que anio_dashboard/mes_dashboard). Si se
    omiten, se usa el mes en curso. Esto permite navegar mes a mes en vez de acumular siempre
    hasta hoy: para agosto se ve el corrido hasta fin de agosto, para julio hasta fin de julio,
    cada uno con su propio "ultimo pago" (que puede ser distinto si hubo un PAGO real entre medias).

    Formula (interes corrido ACT/365 -- respuesta literal a lo que pidio el jefe de Yuri:
    "cuanto llevo devengado desde el ultimo pago hasta fin de mes"):
        dias_devengo = (fin_de_mes - ultimo_pago_cobrado).days
        interes_devengado = capital_invertido x interes_nota_anual / 365 x dias_devengo

    Reglas (confirmadas con Yuri 10/08/2026):
    - Todas las notas activas en el mes seleccionado, incluyendo REINVERSION (el capital
      reinvertido sigue generando interes real de la nota para la empresa, igual que en el
      prorrateo mensual).
    - Si la nota nunca tuvo un PAGO real todavia (hasta fin del mes seleccionado), el devengo
      arranca en su fecha_inversion.
    - Respeta la barrera de cupon: si la ultima observacion conocida hasta fin de mes fue
      NEGATIVA, esa nota no devenga interes (mismo criterio que detalle_notas_mes_prorrateado).
    - Es el INGRESO de la empresa (interes_nota_anual), NO el pago a inversores
      (interes_inversor_anual) -- el jefe pregunto por lo que "tengo" devengado, no por lo
      que se debe pagar a los inversores.
    """
    hoy = pd.Timestamp.today().normalize()
    anio = int(anio) if anio else hoy.year
    mes = int(mes) if mes else hoy.month
    import calendar as _cal
    dias_mes_sel = _cal.monthrange(anio, mes)[1]
    inicio_mes = pd.Timestamp(anio, mes, 1)
    fin_mes = pd.Timestamp(anio, mes, dias_mes_sel)

    cols_vacias = ["nota", "inversor", "capital_invertido", "interes_nota_anual",
                   "ultimo_pago_cobrado", "dias_devengo", "interes_devengado_no_cobrado",
                   "ingreso_habilitado"]

    notas_todas = filtrar_notas(df_inv)
    activas = notas_todas[
        notas_todas["fecha_inversion"].notna()
        & (notas_todas["fecha_inversion"] <= fin_mes)
        & (notas_todas["fecha_final_inversion"].isna() | (notas_todas["fecha_final_inversion"] >= inicio_mes))
    ].copy()
    if activas.empty:
        return pd.DataFrame(columns=cols_vacias)

    pagos_hasta_fin_mes = df_cal[
        (df_cal["tipo_evento"] == "PAGO") & (df_cal["fecha"].notna()) & (df_cal["fecha"] <= fin_mes)
    ].copy().sort_values(["fecha", "nota"]) if df_cal is not None and not df_cal.empty else pd.DataFrame()
    cache_barrera = {}

    def _ultimo_pago(nota_int):
        if pagos_hasta_fin_mes.empty:
            return None
        pagos_nota = pagos_hasta_fin_mes[pagos_hasta_fin_mes["nota"] == nota_int]
        return None if pagos_nota.empty else pagos_nota.iloc[-1]["fecha"]

    def _ingreso_habilitado(nota_int):
        if nota_int in cache_barrera:
            return cache_barrera[nota_int]
        habilitado = True
        if df_cal is not None and df_control is not None and not df_control.empty:
            fecha_obs = obtener_observacion_previa_nota(df_cal, nota_int, fin_mes)
            if fecha_obs is not None:
                resultado_obs, _ = evaluar_nota_en_fecha(df_control, nota_int, fecha_obs, preferida="contingency")
                if resultado_obs == "NEGATIVA":
                    habilitado = False
        cache_barrera[nota_int] = habilitado
        return habilitado

    filas = []
    for _, row in activas.iterrows():
        try:
            nota_int = int(row.get("nota_num"))
        except (TypeError, ValueError):
            continue

        capital = float(row.get("capital_invertido", 0) or 0)
        tasa_nota = float(row.get("interes_nota_anual", 0) or 0)

        ultimo_pago = _ultimo_pago(nota_int)
        inicio_devengo = pd.Timestamp(ultimo_pago) if ultimo_pago is not None else pd.Timestamp(row["fecha_inversion"])
        inicio_devengo = max(inicio_devengo, pd.Timestamp(row["fecha_inversion"]))

        dias_devengo = max((fin_mes - inicio_devengo).days, 0)
        habilitado = _ingreso_habilitado(nota_int)
        interes = (capital * tasa_nota / 365.0) * dias_devengo if habilitado else 0.0

        filas.append({
            "nota": nota_int,
            "inversor": row.get("inversor", ""),
            "capital_invertido": capital,
            "interes_nota_anual": tasa_nota,
            "ultimo_pago_cobrado": ultimo_pago if ultimo_pago is not None else row["fecha_inversion"],
            "dias_devengo": dias_devengo,
            "interes_devengado_no_cobrado": interes,
            "ingreso_habilitado": "SI" if habilitado else "NO (barrera rota)",
        })

    return pd.DataFrame(filas) if filas else pd.DataFrame(columns=cols_vacias)


def resumen_notas_mes(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, anio: int, mes: int, prorratear: bool = False):
    """Devuelve cobro compañía, pago inversores (devengo mensual) y beneficio.

    - cobro_compania: por defecto (prorratear=False) se calcula desde CALENDARIO_NOTAS
      (eventos PAGO del mes) — es el cobro REAL, el mismo que usan extractos y comparador.
      Si prorratear=True, se calcula por devengo mensual (capital × interes_nota_anual / 12,
      pro-rateado por días activos del mes) en vez de por evento de calendario. Esto evita que
      una nota trimestral/semestral "infle" el mes de cobro real y "vacíe" los demás — solo
      pensado para suavizar la vista del Dashboard general y de Notas.
    - pago_inversores: siempre por devengo mensual (capital × tasa / 12), igual que fijos e
      igual con o sin prorrateo, porque el inversor cobra lo mismo cada mes exista o no cobro
      real ese mes.
    """
    if prorratear:
        detalle = detalle_notas_mes_prorrateado(df_inv, anio, mes, df_cal=df_cal, df_control=df_control)
        cobro_compania = float(detalle["cobro_compania"].sum()) if not detalle.empty else 0.0
        pagos = pagos_notas_mes(df_cal, anio, mes)
    else:
        pagos = pagos_notas_mes(df_cal, anio, mes)
        detalle = preparar_detalle_notas(df_inv, pagos, df_cal=df_cal, df_control=df_control)
        cobro_compania = float(detalle["cobro_compania"].sum()) if not detalle.empty else 0.0
    # Pago al inversor por devengo mensual (lógica correcta, igual con o sin prorrateo)
    pago_inversores = pago_inversores_notas_mes(df_inv, anio, mes)
    beneficio = cobro_compania - pago_inversores
    return cobro_compania, pago_inversores, beneficio, detalle, pagos


def resumen_por_cuenta_cobro(detalle: pd.DataFrame) -> pd.DataFrame:
    if detalle.empty:
        return pd.DataFrame(columns=["cuenta_cobro", "cobro_compania"])
    return detalle.groupby("cuenta_cobro", as_index=False)["cobro_compania"].sum().sort_values("cobro_compania", ascending=False)


def resumen_capital_por_inversor_notas(df_inv: pd.DataFrame, solo_activo: bool = False) -> pd.DataFrame:
    trabajo = filtrar_notas(df_inv)
    hoy = pd.Timestamp.today().normalize()
    if solo_activo:
        trabajo = trabajo[(trabajo["fecha_inversion"].notna()) & (trabajo["fecha_inversion"] <= hoy) & (trabajo["fecha_final_inversion"].isna() | (trabajo["fecha_final_inversion"] >= hoy))]
    if trabajo.empty:
        return pd.DataFrame(columns=["inversor", "capital"])
    return trabajo.groupby("inversor", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False)


def proximo_evento_nota(df_cal: pd.DataFrame, nota: int, tipo: str):
    hoy = pd.Timestamp.today().normalize()
    eventos = df_cal[(df_cal["tipo_evento"] == tipo) & (df_cal["nota"] == nota) & (df_cal["fecha"].notna()) & (df_cal["fecha"] >= hoy)].sort_values("fecha")
    return None if eventos.empty else eventos.iloc[0]["fecha"]


# =========================
# GLOBAL Y DASHBOARD
# =========================
def detectar_activo(row):
    tipo = limpiar_texto(row.get("tipo_inversion", ""))
    subtipo = limpiar_texto(row.get("subtipo_inversion", ""))
    nombre = limpiar_texto(row.get("nombre_activo", ""))
    if tipo == "nota" or nombre.startswith("nota"):
        return "notas"
    if tipo == "prestamo":
        return "prestamo"
    for activo in ["paraguay", "bolivia", "motoclick", "futbol", "fútbol", "bitcoin", "prestamo", "préstamo"]:
        if activo in subtipo or activo in nombre:
            return "futbol" if activo == "fútbol" else ("prestamo" if activo == "préstamo" else activo)
    return "otros"


def inversiones_activas_global(df_inv: pd.DataFrame, fecha=None) -> pd.DataFrame:
    """Devuelve las inversiones que cuentan como capital activo en una fecha.

    Reglas:
    - NOTAS (tipo_inversion=nota): todas las activas en la fecha cuentan,
      sean nueva, reinversion o call. Solo se excluyen las cerradas (fecha_final < fecha).
    - TODO LO DEMÁS: solo se excluyen las que tienen tipo_operacion=cancelada.
      Las reinversiones no existen fuera de notas, así que no hace falta filtrarlas.
    """
    if fecha is None:
        fecha = pd.Timestamp.today().normalize()
    fecha = pd.Timestamp(fecha).normalize()
    trabajo = df_inv.copy()

    # Filtro base: activas en la fecha
    activas = trabajo[
        (trabajo["fecha_inversion"].notna()) &
        (trabajo["fecha_inversion"] <= fecha) &
        (trabajo["fecha_final_inversion"].isna() | (trabajo["fecha_final_inversion"] >= fecha))
    ].copy()

    if activas.empty:
        return activas

    es_cancelada = activas["tipo_operacion"].apply(limpiar_texto) == "cancelada"

    # Excluir canceladas en todos los tipos (notas y empresprivada).
    # Las CANCELADAS no cuentan como capital activo aunque la fecha_final sea futura,
    # porque el inversor ya salió. Esto alinea el Dashboard con la suma de extractos.
    resultado = activas[~es_cancelada].copy()
    return resultado


def tarjeta_kpi(titulo, valor, subtitulo="", estado="normal"):
    colores = {"normal": ("#ffffff", "#0e2338"), "positivo": ("#edf7ed", "#166534"), "riesgo": ("#fff4e5", "#b45309"), "negativo": ("#fee2e2", "#991b1b")}
    fondo, color = colores.get(estado, colores["normal"])
    st.markdown(
        f"""
        <div style="background:{fondo};padding:22px 24px;border-radius:22px;border:1px solid rgba(191,154,95,0.24);box-shadow:0 12px 32px rgba(15,35,55,0.08);min-height:130px;">
            <div style="font-size:13px;color:#667085;font-weight:700;text-transform:uppercase;">{titulo}</div>
            <div style="font-size:32px;color:{color};font-weight:850;margin-top:8px;">{valor}</div>
            <div style="font-size:13px;color:#667085;margin-top:8px;">{subtitulo}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def _fmt_market_cap(valor):
    """Formatea capitalización bursátil en notación compacta (1,69 B$ / 174,9 MM$), con coma
    decimal en español, para que las tarjetas de subyacente se lean como un panel financiero real."""
    if valor is None or pd.isna(valor):
        return "N/D"
    valor = float(valor)
    if valor >= 1e12:
        return f"{valor / 1e12:.2f}".replace(".", ",") + " B$"
    if valor >= 1e9:
        return f"{valor / 1e9:.2f}".replace(".", ",") + " MM$"
    if valor >= 1e6:
        return f"{valor / 1e6:.1f}".replace(".", ",") + " M$"
    return f"${valor:,.0f}"


def _tarjeta_subyacente_html(datos: dict, notas_ticker: pd.DataFrame = None) -> str:
    """HTML de una tarjeta de subyacente al estilo 'panel de cartera': precio, variación del día,
    barra de posición en el rango de 52 semanas, y estadísticas clave (cap., PER, distancia a
    la barrera más próxima de tus notas si aplica). Usa la misma paleta marino/dorado que
    tarjeta_kpi para que no se vea como un componente aparte pegado con pegamento."""
    ticker = datos.get("ticker", "")
    nombre = datos.get("nombre") or ticker
    precio = datos.get("precio_actual")
    var_dia = datos.get("variacion_dia_pct")
    minimo = datos.get("precio_min_52sem")
    maximo = datos.get("precio_max_52sem")
    posicion = datos.get("posicion_en_rango_52sem_pct")
    pct_desde_max = datos.get("pct_desde_max_52sem")
    market_cap = datos.get("market_cap")
    pe = datos.get("trailing_pe")

    subida = (var_dia is not None) and var_dia >= 0
    color_chip = "#166534" if subida else "#991b1b"
    fondo_chip = "#edf7ed" if subida else "#fee2e2"
    signo = "+" if subida else ""
    chip_html = (
        f"<span style='display:inline-block;padding:3px 10px;border-radius:8px;font-size:13px;"
        f"font-weight:700;background:{fondo_chip};color:{color_chip};font-family:ui-monospace,monospace;'>"
        f"{signo}{var_dia:.2f}%</span>" if var_dia is not None else ""
    )

    if posicion is not None and minimo is not None and maximo is not None:
        pos_clamp = max(0, min(100, posicion))
        barra_html = (
            f'<div style="margin-top:14px;">'
            f'<div style="display:flex;justify-content:space-between;font-family:ui-monospace,monospace;font-size:11px;color:#667085;margin-bottom:5px;">'
            f'<span>${minimo:,.2f}</span><span style="font-size:10px;letter-spacing:.08em;text-transform:uppercase;color:#9aa4b2;">rango 52 semanas · {pos_clamp:.0f}% del recorrido</span><span>${maximo:,.2f}</span>'
            f'</div>'
            f'<div style="position:relative;height:7px;border-radius:4px;background:linear-gradient(90deg,#eef1f6,#f0e4ce);border:1px solid #e3e7ee;">'
            f'<span style="position:absolute;top:50%;left:{pos_clamp:.1f}%;width:3px;height:17px;border-radius:2px;background:#0e2338;transform:translate(-50%,-50%);"></span>'
            f'</div>'
            f'</div>'
        )
    else:
        barra_html = "<div style='margin-top:10px;font-size:12px;color:#9aa4b2;'>Sin histórico suficiente para el rango de 52 semanas.</div>"

    stats = []
    stats.append(("Capitalización", _fmt_market_cap(market_cap)))
    stats.append(("PER", f"{pe:.1f}" if pe else "N/D"))
    stats.append(("Desde máx. 52 sem.", f"{pct_desde_max:+.1f}%" if pct_desde_max is not None else "N/D"))

    if notas_ticker is not None and not notas_ticker.empty and precio:
        peor_margen = None
        for _, r in notas_ticker.iterrows():
            barrera_cupon = r.get("barrera_cupon")
            precio_compra = r.get("precio_compra")
            if pd.notna(barrera_cupon) and pd.notna(precio_compra) and precio_compra:
                margen = (precio / (precio_compra * barrera_cupon) - 1) * 100
                if peor_margen is None or margen < peor_margen:
                    peor_margen = margen
        if peor_margen is not None:
            stats.append(("Margen a barrera (peor nota)", f"{peor_margen:+.1f}%"))

    stats_html = "".join(
        f"<div style='display:flex;flex-direction:column;gap:1px;'>"
        f"<span style='font-size:10px;letter-spacing:.08em;text-transform:uppercase;color:#9aa4b2;font-weight:600;'>{k}</span>"
        f"<span style='font-family:ui-monospace,monospace;font-size:13.5px;color:#0e2338;'>{v}</span></div>"
        for k, v in stats
    )

    return f"""
    <div style="background:#fff;border:1px solid #e3e7ee;border-radius:14px;padding:20px 22px;
                box-shadow:0 12px 32px rgba(15,35,55,0.06);display:flex;flex-direction:column;gap:14px;">
      <div style="display:flex;justify-content:space-between;align-items:flex-start;gap:12px;">
        <div>
          <p style="margin:0;font-family:ui-monospace,monospace;font-size:16px;font-weight:700;letter-spacing:.03em;color:#0e2338;">{ticker}</p>
          <p style="margin:2px 0 0;font-size:13px;color:#667085;">{nombre}</p>
        </div>
        <div style="text-align:right;">
          <p style="margin:0;font-family:ui-monospace,monospace;font-size:26px;font-weight:600;color:#0e2338;">
            {f"${precio:,.2f}" if precio else "N/D"}</p>
          <div style="margin-top:4px;">{chip_html}</div>
        </div>
      </div>
      {barra_html}
      <div style="display:flex;gap:22px;flex-wrap:wrap;border-top:1px solid #e3e7ee;padding-top:12px;">
        {stats_html}
      </div>
    </div>
    """


def grafico_precio_subyacente(datos: dict, precio_contingencia: float = None):
    """Gráfico de evolución de precio (~1 año) del subyacente, con línea de la barrera de
    contingencia marcada si la ficha se está viendo desde el contexto de una nota concreta."""
    if px is None:
        st.warning("Falta plotly. Añade plotly a requirements.txt.")
        return None
    historico = datos.get("historico_precios")
    if not historico:
        st.info("Sin histórico de precio disponible para graficar.")
        return None
    df_hist = pd.DataFrame(historico)
    df_hist["fecha"] = pd.to_datetime(df_hist["fecha"])
    fig = px.line(df_hist, x="fecha", y="close", title=None)
    fig.update_traces(line=dict(color="#9A6B24", width=2), hovertemplate="%{x|%d %b %Y}<br>$%{y:,.2f}<extra></extra>")
    if precio_contingencia:
        fig.add_hline(
            y=precio_contingencia, line_dash="dash", line_color="#B03A2E", line_width=1.5,
            annotation_text=f"Barrera contingencia ${precio_contingencia:,.2f}",
            annotation_position="top left", annotation_font_color="#B03A2E", annotation_font_size=11,
        )
    fig.update_layout(
        height=280, margin=dict(l=10, r=10, t=10, b=10),
        paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Archivo, sans-serif", size=12, color="#58616E"),
        xaxis=dict(title=None, showgrid=False), yaxis=dict(title=None, showgrid=True, gridcolor="#EAEEF3"),
    )
    st.plotly_chart(fig, use_container_width=True)
    return fig


def validar_base_datos(df_inv, df_cal, df_control):
    resultados = []
    def add(nombre, cantidad, gravedad):
        resultados.append({"Validación": nombre, "Incidencias": int(cantidad), "Estado": gravedad if cantidad > 0 else "OK"})
    add("Inversiones sin fecha de inversión", df_inv["fecha_inversion"].isna().sum() if "fecha_inversion" in df_inv.columns else len(df_inv), "ALTA")
    add("Inversiones sin capital invertido", (df_inv["capital_invertido"].fillna(0) <= 0).sum() if "capital_invertido" in df_inv.columns else len(df_inv), "ALTA")
    add("Inversiones sin inversor", (df_inv["inversor"].fillna("").astype(str).str.strip() == "").sum() if "inversor" in df_inv.columns else len(df_inv), "MEDIA")
    notas_filtradas = filtrar_notas(df_inv) if not df_inv.empty else pd.DataFrame()
    add("Notas sin número detectado", notas_filtradas["nota_num"].isna().sum() if not notas_filtradas.empty and "nota_num" in notas_filtradas.columns else 0, "ALTA")
    add("Eventos de calendario sin fecha", df_cal["fecha"].isna().sum() if "fecha" in df_cal.columns else 0, "MEDIA")
    add("Control de notas sin ticker", (df_control["ticker"].fillna("").astype(str).str.strip() == "").sum() if "ticker" in df_control.columns else 0, "ALTA")
    add("Control de notas sin precio de compra", df_control["precio_compra"].isna().sum() if "precio_compra" in df_control.columns else 0, "ALTA")
    return pd.DataFrame(resultados)


def detectar_alertas_financieras(df_inv, df_cal, df_control):
    """
    Alertas mejoradas para notas:
    - Evento próximo: observaciones y pagos cercanos.
    - Observación negativa real: solo si ya pasó y el precio está bajo barrera.
    - Observación sin dato: se avisa, pero no bloquea el cobro previsto.
    - Pago bloqueado: pago futuro o histórico cuya observación previa fue negativa real.
    - Pago previsto: pago futuro con observación pendiente o sin dato.
    """
    hoy = pd.Timestamp.today().normalize()
    alertas = []

    def add(tipo, detalle, fecha, prioridad, nota=""):
        alertas.append({
            "Tipo": tipo,
            "Nota": nota,
            "Detalle": detalle,
            "Fecha": pd.Timestamp(fecha).strftime("%d/%m/%Y") if pd.notna(fecha) else "",
            "Prioridad": prioridad,
        })

    if df_cal is not None and not df_cal.empty and "fecha" in df_cal.columns:
        eventos_7 = df_cal[
            (df_cal["fecha"].notna())
            & (df_cal["fecha"] >= hoy)
            & (df_cal["fecha"] <= hoy + pd.Timedelta(days=7))
        ].copy().sort_values(["fecha", "tipo_evento", "nota"])

        for _, row in eventos_7.iterrows():
            add(
                "Evento próximo",
                f"{row.get('tipo_evento', '')} de NOTA {row.get('nota', '')}",
                row.get("fecha"),
                "MEDIA",
                row.get("nota", ""),
            )

        # Observaciones ya vencidas o de hoy: revisar si fueron negativas reales o sin dato.
        observaciones_vencidas = df_cal[
            (df_cal["tipo_evento"] == "OBSERVACION")
            & (df_cal["fecha"].notna())
            & (df_cal["fecha"] <= hoy)
        ].copy().sort_values("fecha")

        for _, row in observaciones_vencidas.iterrows():
            nota = row.get("nota")
            fecha_obs = row.get("fecha")
            if pd.isna(nota):
                continue
            nota_int = int(nota)
            resultado, detalle = evaluar_nota_en_fecha(df_control, nota_int, fecha_obs, preferida="contingency")
            detalle_txt = resumen_detalle_observacion(detalle)

            if resultado == "NEGATIVA":
                add(
                    "Nota negativa real",
                    f"NOTA {nota_int}: observación negativa. No debe contarse el cobro. {detalle_txt}",
                    fecha_obs,
                    "ALTA",
                    nota_int,
                )
            elif resultado == "SIN DATO":
                add(
                    "Revisar dato faltante",
                    f"NOTA {nota_int}: la observación ya pasó, pero falta precio. Se mantiene como cobro previsto hasta revisar. {detalle_txt}",
                    fecha_obs,
                    "MEDIA",
                    nota_int,
                )
            elif resultado in ["SIN_CONTROL", "SIN_COLUMNAS", "SIN_OBSERVACION"]:
                add(
                    "Revisar configuración",
                    f"NOTA {nota_int}: no se ha podido evaluar correctamente ({resultado}).",
                    fecha_obs,
                    "ALTA",
                    nota_int,
                )

        # Observaciones futuras a 30 días: seguimiento.
        observaciones_futuras = df_cal[
            (df_cal["tipo_evento"] == "OBSERVACION")
            & (df_cal["fecha"].notna())
            & (df_cal["fecha"] > hoy)
            & (df_cal["fecha"] <= hoy + pd.Timedelta(days=30))
        ].copy().sort_values("fecha")

        for _, row in observaciones_futuras.iterrows():
            add(
                "Observación pendiente",
                f"NOTA {row.get('nota', '')}: observación futura. Se cuenta como cobro previsto hasta que llegue la fecha.",
                row.get("fecha"),
                "BAJA",
                row.get("nota", ""),
            )

        # Pagos próximos: indicar si están habilitados, previstos o bloqueados.
        pagos_30 = df_cal[
            (df_cal["tipo_evento"] == "PAGO")
            & (df_cal["fecha"].notna())
            & (df_cal["fecha"] >= hoy)
            & (df_cal["fecha"] <= hoy + pd.Timedelta(days=30))
        ].copy().sort_values("fecha")

        for _, row in pagos_30.iterrows():
            nota = row.get("nota")
            fecha_pago = row.get("fecha")
            if pd.isna(nota):
                continue
            nota_int = int(nota)
            fecha_obs = obtener_observacion_previa_nota(df_cal, nota_int, fecha_pago)
            resultado, detalle = evaluar_nota_en_fecha(df_control, nota_int, fecha_obs, preferida="contingency") if fecha_obs is not None else ("SIN_OBSERVACION", pd.DataFrame())
            detalle_txt = resumen_detalle_observacion(detalle)

            if resultado == "NEGATIVA":
                add(
                    "Pago bloqueado",
                    f"NOTA {nota_int}: pago próximo bloqueado por observación negativa. {detalle_txt}",
                    fecha_pago,
                    "ALTA",
                    nota_int,
                )
            elif resultado in ["PENDIENTE", "SIN DATO", "NO_EVALUADA"]:
                add(
                    "Pago previsto",
                    f"NOTA {nota_int}: pago próximo contado como previsto. Estado observación: {resultado}. {detalle_txt}",
                    fecha_pago,
                    "MEDIA",
                    nota_int,
                )
            elif resultado == "POSITIVA":
                add(
                    "Pago habilitado",
                    f"NOTA {nota_int}: pago próximo habilitado por observación positiva. {detalle_txt}",
                    fecha_pago,
                    "BAJA",
                    nota_int,
                )

    validaciones = validar_base_datos(df_inv, df_cal, df_control)
    errores_altos = validaciones[(validaciones["Incidencias"] > 0) & (validaciones["Estado"] == "ALTA")]
    for _, row in errores_altos.iterrows():
        add(
            "Validación crítica",
            f"{row['Validación']}: {row['Incidencias']} incidencias",
            hoy,
            "ALTA",
            "",
        )

    if not alertas:
        return pd.DataFrame(columns=["Tipo", "Nota", "Detalle", "Fecha", "Prioridad"])

    orden = {"ALTA": 0, "MEDIA": 1, "BAJA": 2}
    out = pd.DataFrame(alertas)
    out["orden_prioridad"] = out["Prioridad"].map(orden).fillna(9)
    out["fecha_orden"] = pd.to_datetime(out["Fecha"], errors="coerce", dayfirst=True)
    out = out.sort_values(["orden_prioridad", "fecha_orden", "Tipo"]).drop(columns=["orden_prioridad", "fecha_orden"])
    return out


def calcular_rentabilidad_inversiones_mes(df_inv, df_cal, df_control, anio: int, mes: int, prorratear_notas: bool = False) -> pd.DataFrame:
    """
    Construye una tabla homogénea de rentabilidad mensual por inversión.
    - rentabilidad_beneficio_mes: beneficio empresa / capital.
    - rentabilidad_beneficio_anualizada: rentabilidad mensual x 12.
    - rentabilidad_pagada_inversor_mes: pago inversor / capital.
    - rentabilidad_pagada_inversor_anualizada: rentabilidad mensual pagada x 12.

    prorratear_notas: si True, el cobro de notas se reparte por devengo mensual en vez de por
    evento de calendario (ver resumen_notas_mes). Solo afecta a las filas 'notas'.
    """
    filas = []

    # Notas estructuradas
    _, _, _, detalle_notas, _ = resumen_notas_mes(df_inv, df_cal, df_control, anio, mes, prorratear=prorratear_notas)
    if detalle_notas is not None and not detalle_notas.empty:
        for _, row in detalle_notas.iterrows():
            capital = float(row.get("capital_invertido", 0) or 0)
            cobro = float(row.get("cobro_compania", 0) or 0)
            pago = float(row.get("pago_inversor", 0) or 0)
            beneficio = float(row.get("beneficio_empresa", 0) or 0)
            filas.append({
                "activo": "notas",
                "nombre_activo": f"NOTA {row.get('nota', '')}",
                "id_inversion": row.get("id_inversion", ""),
                "inversor": row.get("inversor", ""),
                "capital": capital,
                "cobro_compania_mes": cobro,
                "pago_inversor_mes": pago,
                "beneficio_empresa_mes": beneficio,
                "resultado_observacion": row.get("resultado_observacion", ""),
                "rentabilidad_beneficio_mes": beneficio / capital if capital else 0,
                "rentabilidad_beneficio_anualizada": (beneficio / capital * 12) if capital else 0,
                "rentabilidad_pagada_inversor_mes": pago / capital if capital else 0,
                "rentabilidad_pagada_inversor_anualizada": (pago / capital * 12) if capital else 0,
            })

    # Activos con rentabilidad fija / operativa
    for activo, tasa in [("paraguay", TASA_ANUAL_PARAGUAY), ("bolivia", TASA_ANUAL_BOLIVIA), ("motoclick", TASA_ANUAL_MOTOCLICK), ("futbol", TASA_ANUAL_FUTBOL), ("bitcoin", TASA_ANUAL_BITCOIN)]:
        det = detalle_activo_mes(df_inv, activo, tasa, anio, mes)
        if det is None or det.empty:
            continue
        # Para MotoClick aplicar ajuste por devoluciones/reinversiones reales
        if activo == "motoclick":
            det["activo"] = "motoclick"
            det = ajustar_ingreso_motoclick(det, df_inv, anio, mes)
        for _, row in det.iterrows():
            capital = float(row.get("capital_invertido", 0) or 0)
            cobro = float(row.get("ingreso_bruto", 0) or 0)
            pago = float(row.get("pago_inversor_mes", 0) or 0)
            beneficio = float(row.get("beneficio_empresa_mes", 0) or 0)
            filas.append({
                "activo": activo,
                "nombre_activo": activo,
                "id_inversion": row.get("id_inversion", ""),
                "inversor": row.get("inversor", ""),
                "capital": capital,
                "cobro_compania_mes": cobro,
                "pago_inversor_mes": pago,
                "beneficio_empresa_mes": beneficio,
                "resultado_observacion": "NO APLICA",
                "rentabilidad_beneficio_mes": beneficio / capital if capital else 0,
                "rentabilidad_beneficio_anualizada": (beneficio / capital * 12) if capital else 0,
                "rentabilidad_pagada_inversor_mes": pago / capital if capital else 0,
                "rentabilidad_pagada_inversor_anualizada": (pago / capital * 12) if capital else 0,
            })

    # Préstamos internos (tasa propia por fila, no una constante global — ver detalle_prestamos_mes)
    det_prest = detalle_prestamos_mes(df_inv, anio, mes)
    if det_prest is not None and not det_prest.empty:
        for _, row in det_prest.iterrows():
            capital = float(row.get("capital_invertido", 0) or 0)
            cobro = float(row.get("ingreso_bruto", 0) or 0)
            pago = float(row.get("pago_inversor_mes", 0) or 0)
            beneficio = float(row.get("beneficio_empresa_mes", 0) or 0)
            filas.append({
                "activo": "prestamo",
                "nombre_activo": "Préstamo interno",
                "id_inversion": row.get("id_inversion", ""),
                "inversor": row.get("inversor", ""),
                "capital": capital,
                "cobro_compania_mes": cobro,
                "pago_inversor_mes": pago,
                "beneficio_empresa_mes": beneficio,
                "resultado_observacion": "NO APLICA",
                "rentabilidad_beneficio_mes": beneficio / capital if capital else 0,
                "rentabilidad_beneficio_anualizada": (beneficio / capital * 12) if capital else 0,
                "rentabilidad_pagada_inversor_mes": pago / capital if capital else 0,
                "rentabilidad_pagada_inversor_anualizada": (pago / capital * 12) if capital else 0,
            })

    return pd.DataFrame(filas)


def preparar_tabla_rentabilidad(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return df
    out = df.copy()
    for col in ["capital", "cobro_compania_mes", "pago_inversor_mes", "beneficio_empresa_mes"]:
        if col in out.columns:
            out[col] = out[col].map(fmt)
    for col in [
        "rentabilidad_beneficio_mes",
        "rentabilidad_beneficio_anualizada",
        "rentabilidad_pagada_inversor_mes",
        "rentabilidad_pagada_inversor_anualizada",
    ]:
        if col in out.columns:
            out[col] = out[col].map(fmt_pct)
    return out


def _precio_yfinance_history(ticker):
    """Intento 1: historial de 5 días (método actual)."""
    try:
        hist = yf.Ticker(ticker).history(period="5d")
        if hist is not None and not hist.empty:
            cierre = hist["Close"].dropna()
            if not cierre.empty:
                return float(cierre.iloc[-1])
    except Exception:
        pass
    return None


def _precio_yfinance_fastinfo(ticker):
    """Intento 2: endpoint distinto de yfinance (fast_info) — a veces funciona cuando history() falla."""
    try:
        fi = yf.Ticker(ticker).fast_info
        precio = fi.get("lastPrice") or fi.get("last_price") or fi.get("regularMarketPrice")
        if precio:
            return float(precio)
    except Exception:
        pass
    return None


def _precio_stooq(ticker):
    """Intento 3: Stooq, fuente alternativa gratuita sin API key. Fallback final si yfinance falla."""
    try:
        for candidato in [ticker, f"{ticker}.US"]:
            url = f"https://stooq.com/q/l/?s={candidato.lower()}&f=sd2t2ohlcv&h&e=csv"
            r = requests.get(url, timeout=8)
            if r.status_code == 200:
                lineas = r.text.strip().split("\n")
                if len(lineas) >= 2:
                    campos = lineas[1].split(",")
                    if len(campos) >= 7 and campos[6] not in ("N/D", ""):
                        precio = float(campos[6])
                        if precio > 0:
                            return precio
    except Exception:
        pass
    return None


def obtener_precio_actual_con_fallback(ticker):
    """Precio actual con cadena de respaldo: yfinance (history) -> yfinance (fast_info) -> Stooq.
    Solo devuelve None si las tres fuentes fallan — evita el 'SIN PRECIO DISPONIBLE' innecesario."""
    ticker = str(ticker).strip().upper()
    for fuente in (_precio_yfinance_history, _precio_yfinance_fastinfo, _precio_stooq):
        precio = fuente(ticker)
        if precio is not None and precio > 0:
            return precio
    return None


def clasificar_alerta_riesgo(variacion_pct):
    """Clasifica el riesgo de un ticker según su variación % desde el precio de compra.
    Criterio único y simple (a pedido de Yuri, tras sopesar pros/contras): 🔴 EN RIESGO si la
    variación es ≤ -30%. Se eligió variación fija en vez de margen a la barrera porque:
    (a) es mucho más fácil de verificar a mano contra el Excel, y (b) es más robusto — depende
    solo de precio_compra (que se consulta en vivo con yfinance al crear la nota y casi nunca
    falla), no de barrera_cupon/barrera_capital (campos que hemos visto mal extraídos varias veces).
    El margen a la barrera de cada nota se sigue calculando y mostrando como dato adicional en el
    contexto de la IA, pero YA NO decide si una nota entra en la lista de riesgo.
    """
    if pd.isna(variacion_pct):
        return "SIN DATO"
    try:
        variacion_pct = float(variacion_pct)
    except Exception:
        return "SIN DATO"
    if variacion_pct <= -30:
        return "ROJO"
    return "OK"


@st.cache_data(show_spinner=False, ttl=1800)
def construir_resumen_actual_notas_alertas(df_control: pd.DataFrame) -> pd.DataFrame:
    """Calcula precios actuales, variación y alerta por ticker/nota para la sección Notas y el dashboard.
    El criterio de riesgo (alerta_riesgo) es variación ≤ -30% desde precio_compra. El margen a la
    barrera de contingencia se calcula igual y se deja en margen_a_barrera_% como dato adicional,
    pero un ticker sin barrera cargada en CONTROL_NOTAS igual se evalúa por riesgo (no se descarta)."""
    if yf is None or df_control is None or df_control.empty:
        return pd.DataFrame()

    control = df_control.copy()
    barrera_col = next((c for c in ["contingency", "barrera_capital", "barrera_cupon"] if c in control.columns), None)
    faltan = [c for c in ["nota", "ticker", "precio_compra"] if c not in control.columns]
    if faltan:
        return pd.DataFrame()

    control["nota"] = pd.to_numeric(control["nota"], errors="coerce")
    control["ticker"] = control["ticker"].astype(str).str.strip().str.upper()
    control["precio_compra"] = pd.to_numeric(control["precio_compra"], errors="coerce")
    if barrera_col is not None:
        control[barrera_col] = pd.to_numeric(control[barrera_col], errors="coerce").apply(lambda x: x / 100 if pd.notna(x) and x > 1 else x)
    control = control.dropna(subset=["nota", "ticker", "precio_compra"]).copy()

    # Antes se pedía el precio ticker por ticker, en secuencia, probando hasta 3 fuentes cada
    # uno — con ~30-40 tickers eso son decenas de llamadas de red seguidas, y es la causa real
    # de la lentitud en el primer login (cuando esta caché de 30 min está fría), no la fuente de
    # datos (Drive/Postgres). Se piden todos los tickers ÚNICOS en paralelo (además, si el mismo
    # ticker aparece en varias notas, ahora solo se pide una vez en vez de una vez por nota).
    tickers_unicos = [t for t in control["ticker"].unique().tolist() if t]
    precios_por_ticker = {}
    if tickers_unicos:
        import concurrent.futures
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(10, len(tickers_unicos))) as executor:
            futuros = {executor.submit(obtener_precio_actual_con_fallback, t): t for t in tickers_unicos}
            for futuro in concurrent.futures.as_completed(futuros):
                t = futuros[futuro]
                try:
                    precios_por_ticker[t] = futuro.result()
                except Exception:
                    precios_por_ticker[t] = None

    filas = []
    for _, row in control.iterrows():
        ticker = row["ticker"]
        precio_actual = precios_por_ticker.get(ticker)

        precio_compra = float(row["precio_compra"])
        barrera = float(row[barrera_col]) if barrera_col is not None and pd.notna(row.get(barrera_col)) else None
        precio_contingencia = (precio_compra * barrera) if barrera is not None else None
        variacion = None if precio_actual is None else ((precio_actual - precio_compra) / precio_compra) * 100
        margen_a_barrera = None if (precio_actual is None or precio_contingencia is None) else ((precio_actual - precio_contingencia) / precio_contingencia) * 100
        alerta_riesgo = clasificar_alerta_riesgo(variacion)
        estado_barrera = "SIN DATO" if (precio_actual is None or precio_contingencia is None) else ("OK" if precio_actual >= precio_contingencia else "RIESGO")
        # Cuántos puntos porcentuales le quedan de colchón antes de tocar el umbral de riesgo (-30%).
        # Si ya está en ROJO, sale ≤ 0 (ya lo cruzó).
        distancia_a_rojo = None if variacion is None else (variacion - (-30))

        filas.append({
            "nota": int(row["nota"]),
            "ticker": ticker,
            "precio_compra": precio_compra,
            "precio_actual": precio_actual,
            "variacion_%": variacion,
            "distancia_a_rojo_%": distancia_a_rojo,
            "precio_contingencia": precio_contingencia,
            "margen_a_barrera_%": margen_a_barrera,
            "estado_barrera": estado_barrera,
            "alerta_riesgo": alerta_riesgo,
        })

    return pd.DataFrame(filas)


def resumen_alertas_por_nota(resumen_notas_actual: pd.DataFrame) -> pd.DataFrame:
    """Resume alertas amarillas/rojas por nota, tomando el ticker con PEOR margen a su barrera
    de contingencia (el criterio único de riesgo: cuánto colchón le queda a esa nota concreta)."""
    cols_vacio = ["nota", "alerta", "peor_margen_%", "peor_variacion_%", "tickers"]
    if resumen_notas_actual is None or resumen_notas_actual.empty:
        return pd.DataFrame(columns=cols_vacio)
    alertas = resumen_notas_actual[resumen_notas_actual["alerta_riesgo"] == "ROJO"].copy()
    if alertas.empty:
        return pd.DataFrame(columns=cols_vacio)

    filas = []
    for nota, grupo in alertas.groupby("nota"):
        grupo = grupo.copy()
        peor = grupo.sort_values("variacion_%", ascending=True).iloc[0]
        filas.append({
            "nota": int(nota),
            "alerta": peor["alerta_riesgo"],
            "peor_margen_%": peor.get("margen_a_barrera_%", None),
            "peor_variacion_%": peor.get("variacion_%", None),
            "tickers": ", ".join(grupo["ticker"].astype(str).unique()),
        })
    out = pd.DataFrame(filas)
    return out.sort_values("peor_variacion_%")


def evaluar_calls_posibles_notas(df_cal: pd.DataFrame, df_control: pd.DataFrame, df_calls: pd.DataFrame,
                                   resumen_riesgo: pd.DataFrame) -> list[dict]:
    """Devuelve SOLO las notas donde un CALL es objetivamente posible en su próxima fecha de call.

    Criterio (a pedido de Yuri): una nota es 'CALL POSIBLE' cuando, a la vez:
      1) Tiene una fecha de call futura en CALENDARIO_CALLS.
      2) TODOS los cupones pagados hasta hoy fueron POSITIVA (nunca hubo observación NEGATIVA /
         barrera de cupón rota) — se reevalúa cada pago histórico con evaluar_nota_en_fecha,
         igual que hace preparar_calendario_integrado_notas.
      3) TODOS los tickers de la nota están HOY en positivo (variación_% >= 0) frente a su
         precio_compra — usando resumen_riesgo, que viene de construir_resumen_actual_notas_alertas
         (mismos precios que el semáforo, para que esto nunca pueda divergir de lo que ve Yuri).
    Si falta algún dato de precio de algún ticker de la nota, esa nota NO se reporta (no se puede
    confirmar la condición 3 con datos incompletos). Las notas que no cumplen NO se incluyen en el
    resultado — este es el criterio explícito de Yuri: solo enseñar las que sí son posible call."""
    resultado = []
    if df_calls is None or df_calls.empty or df_cal is None or df_cal.empty or df_control is None or df_control.empty:
        return resultado

    hoy = pd.Timestamp.today().normalize()

    calls = df_calls.copy()
    col_fecha_call = "fecha_call" if "fecha_call" in calls.columns else ("fecha" if "fecha" in calls.columns else None)
    if col_fecha_call is None or "nota" not in calls.columns:
        return resultado
    calls[col_fecha_call] = pd.to_datetime(calls[col_fecha_call], errors="coerce", dayfirst=True).dt.normalize()
    calls["nota"] = pd.to_numeric(calls["nota"], errors="coerce")
    calls_fut = calls.dropna(subset=[col_fecha_call, "nota"])
    calls_fut = calls_fut[calls_fut[col_fecha_call] >= hoy].sort_values(col_fecha_call)
    if calls_fut.empty:
        return resultado

    prox_call_por_nota = calls_fut.groupby("nota")[col_fecha_call].first()

    df_cal_local = df_cal.copy()
    df_cal_local["fecha"] = pd.to_datetime(df_cal_local["fecha"], errors="coerce")

    for nota_f, fecha_call in prox_call_por_nota.items():
        nota_id = int(nota_f)
        dias_restantes = (fecha_call - hoy).days

        # 1) Histórico de cupones: TODOS los pagos pasados de esta nota deben ser != NEGATIVA
        pagos_pasados = df_cal_local[
            (df_cal_local["nota"] == nota_id) &
            (df_cal_local["tipo_evento"] == "PAGO") &
            (df_cal_local["fecha"].notna()) &
            (df_cal_local["fecha"] <= hoy)
        ].sort_values("fecha")

        cupon_siempre_pagado = True
        n_cupones_evaluados = 0
        for _, row_pago in pagos_pasados.iterrows():
            fecha_pago = pd.Timestamp(row_pago["fecha"]).normalize()
            fecha_obs = obtener_observacion_previa_nota(df_cal_local, nota_id, fecha_pago)
            if fecha_obs is None:
                continue
            resultado_obs, _ = evaluar_nota_en_fecha(df_control, nota_id, fecha_obs, preferida="contingency")
            n_cupones_evaluados += 1
            if resultado_obs == "NEGATIVA":
                cupon_siempre_pagado = False
                break

        if not cupon_siempre_pagado or n_cupones_evaluados == 0:
            continue

        # 2) Todos los tickers de la nota en positivo HOY vs precio_compra (mismos precios que el semáforo)
        sub_control = df_control[pd.to_numeric(df_control.get("nota"), errors="coerce") == nota_id]
        if sub_control.empty:
            continue
        tickers_nota = sub_control["ticker"].astype(str).str.strip().str.upper().unique().tolist()

        if resumen_riesgo is None or resumen_riesgo.empty:
            continue
        filas_nota = resumen_riesgo[resumen_riesgo["nota"] == nota_id]
        if filas_nota.empty or len(filas_nota) < len(tickers_nota):
            continue  # faltan datos de precio de algún ticker: no se puede confirmar, no se reporta

        variaciones = pd.to_numeric(filas_nota["variacion_%"], errors="coerce")
        if variaciones.isna().any() or (variaciones < 0).any():
            continue  # algún ticker sin dato o en negativo: no es call posible

        detalle_precios = "; ".join(
            f"{r['ticker']}: {r['variacion_%']:+.1f}%" for _, r in filas_nota.iterrows()
        )
        resultado.append({
            "nota": nota_id,
            "fecha_call": fecha_call,
            "dias_restantes": dias_restantes,
            "peor_variacion_%": float(variaciones.min()),
            "detalle_precios": detalle_precios,
            "n_cupones_pagados_sin_fallo": n_cupones_evaluados,
        })

    resultado.sort(key=lambda d: d["dias_restantes"])
    return resultado


def construir_semaforo_consolidado_notas(resumen_notas_actual: pd.DataFrame) -> pd.DataFrame:
    """Semáforo por nota (una fila por nota, no por ticker): para cada nota se toma el ticker con
    PEOR variación (el que manda en un worst-of), y se muestra su precio actual, su precio de
    contingencia (barrera) y el margen % entre ambos. A diferencia de resumen_alertas_por_nota,
    esta función devuelve TODAS las notas (verdes y rojas), para tener de un vistazo el estado
    completo de la cartera de notas sin entrar nota por nota."""
    cols_vacio = ["nota", "alerta", "peor_ticker", "peor_variacion_%", "precio_actual", "precio_contingencia", "margen_a_barrera_%", "distancia_a_rojo_%", "n_tickers"]
    if resumen_notas_actual is None or resumen_notas_actual.empty:
        return pd.DataFrame(columns=cols_vacio)

    filas = []
    for nota, grupo in resumen_notas_actual.groupby("nota"):
        grupo_valido = grupo.dropna(subset=["variacion_%"])
        if grupo_valido.empty:
            filas.append({
                "nota": int(nota), "alerta": "SIN DATO", "peor_ticker": ", ".join(grupo["ticker"].astype(str).unique()),
                "peor_variacion_%": None, "precio_actual": None, "precio_contingencia": None,
                "margen_a_barrera_%": None, "distancia_a_rojo_%": None, "n_tickers": len(grupo),
            })
            continue
        peor = grupo_valido.sort_values("variacion_%", ascending=True).iloc[0]
        filas.append({
            "nota": int(nota),
            "alerta": peor["alerta_riesgo"],
            "peor_ticker": peor["ticker"],
            "peor_variacion_%": peor["variacion_%"],
            "precio_actual": peor.get("precio_actual"),
            "precio_contingencia": peor.get("precio_contingencia"),
            "margen_a_barrera_%": peor.get("margen_a_barrera_%"),
            "distancia_a_rojo_%": peor.get("distancia_a_rojo_%"),
            "n_tickers": len(grupo),
        })
    out = pd.DataFrame(filas)
    return out.sort_values("peor_variacion_%", na_position="last").reset_index(drop=True) if not out.empty else out


def colorear_semaforo_consolidado(row):
    alerta = row.get("alerta", "")
    if alerta == "ROJO":
        return ["background-color: #fee2e2; color: #7f1d1d; font-weight: 700"] * len(row)
    dist = row.get("distancia_a_rojo_%")
    if pd.notna(dist) and dist < 10:
        # Todavía en OK pero a menos de 10 puntos de convertirse en ROJO: aviso temprano.
        return ["background-color: #fef3c7; color: #78350f; font-weight: 700"] * len(row)
    return [""] * len(row)


def colorear_filas_alerta_notas(row):
    alerta = row.get("alerta_riesgo", "")
    if alerta == "ROJO":
        return ["background-color: #fee2e2; color: #7f1d1d; font-weight: 700"] * len(row)
    if alerta == "AMARILLO":
        return ["background-color: #fef3c7; color: #78350f; font-weight: 700"] * len(row)
    return [""] * len(row)


def inicio_semana_lunes(fecha):
    fecha = pd.Timestamp(fecha).normalize()
    return fecha - pd.Timedelta(days=fecha.weekday())


def resumen_cobros_semanales_mes_notas(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, anio: int, mes: int) -> pd.DataFrame:
    """Agrupa los cobros de notas por semanas naturales lunes-domingo dentro de un mes."""
    pagos_mes = pagos_notas_mes(df_cal, anio, mes)
    detalle = preparar_detalle_notas(df_inv, pagos_mes, df_cal=df_cal, df_control=df_control)
    if detalle is None or detalle.empty:
        return pd.DataFrame(columns=["semana", "nota", "fecha_pago", "cobro_compania"])

    trabajo = detalle.copy()
    trabajo["fecha_pago"] = pd.to_datetime(trabajo["fecha_pago"], errors="coerce").dt.normalize()
    trabajo = trabajo[trabajo["fecha_pago"].notna()].copy()
    trabajo["inicio_semana"] = trabajo["fecha_pago"].apply(inicio_semana_lunes)
    trabajo["fin_semana"] = trabajo["inicio_semana"] + pd.Timedelta(days=6)
    trabajo["semana"] = trabajo.apply(
        lambda r: f"Semana del {r['inicio_semana'].day} - {r['fin_semana'].day} de {nombre_mes_es(int(r['inicio_semana'].month))}",
        axis=1,
    )

    resumen = trabajo.groupby(["inicio_semana", "fin_semana", "semana", "nota", "fecha_pago"], as_index=False)["cobro_compania"].sum()
    resumen = resumen.sort_values(["inicio_semana", "fecha_pago", "nota"])
    return resumen[["semana", "nota", "fecha_pago", "cobro_compania"]]


def mostrar_cobros_semanales_dashboard(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, anio: int, mes: int):
    st.markdown("### Cobros semanales del mes")
    st.caption("Cada semana muestra el total previsto y un desplegable con el desglose por nota.")

    tabla_semanal = resumen_cobros_semanales_mes_notas(df_inv, df_cal, df_control, anio, mes)
    if tabla_semanal.empty:
        st.info("No hay cobros de notas previstos para ese mes.")
        return

    resumen_semana = (
        tabla_semanal
        .groupby("semana", as_index=False)["cobro_compania"]
        .sum()
        .rename(columns={"cobro_compania": "total_semana"})
    )

    st.dataframe(preparar_tabla_monetaria(resumen_semana, ["total_semana"]), use_container_width=True)

    for _, fila_semana in resumen_semana.iterrows():
        semana = fila_semana["semana"]
        total = float(fila_semana["total_semana"] or 0)
        detalle = tabla_semanal[tabla_semanal["semana"] == semana].copy()
        detalle = (
            detalle.groupby(["nota", "fecha_pago"], as_index=False)["cobro_compania"]
            .sum()
            .sort_values(["fecha_pago", "nota"])
        )
        detalle["nota"] = detalle["nota"].apply(lambda x: f"NOTA {int(x)}" if pd.notna(x) else "NOTA")
        with st.expander(f"{semana} · Total {fmt(total)}", expanded=False):
            st.dataframe(preparar_tabla_monetaria(detalle, ["cobro_compania"]), use_container_width=True)



def obtener_resumen_dashboard(df_inv, df_cal, df_control, anio: int | None = None, mes: int | None = None, vista_activo: str = "General", incluir_chaparro: bool = True, prorratear_notas: bool = False):
    hoy_real = pd.Timestamp.today().normalize()
    if anio is None:
        anio = hoy_real.year
    if mes is None:
        mes = hoy_real.month
    fecha_analisis = pd.Timestamp(int(anio), int(mes), ultimo_dia_mes(int(anio), int(mes))).normalize()
    # El filtro chaparro solo afecta al capital activo mostrado, NO a los cálculos de cobros.
    # resumen_notas_mes y detalle_activo_mes usan df_inv completo (igual que Consultas).
    # Así se garantiza que Dashboard y Consultas siempre muestren los mismos importes.
    df_inv_filtrado = aplicar_filtro_chaparro_fernandez(df_inv, incluir_chaparro)
    activas = inversiones_activas_global(df_inv_filtrado, fecha_analisis)
    if not activas.empty:
        activas["activo"] = activas.apply(detectar_activo, axis=1)
    capital_total = activas["capital_invertido"].sum() if not activas.empty else 0
    c_notas, p_notas, b_notas, detalle_notas, _ = resumen_notas_mes(df_inv, df_cal, df_control, int(anio), int(mes), prorratear=prorratear_notas)
    detalles_fijos = []
    for activo, tasa in [("paraguay", TASA_ANUAL_PARAGUAY), ("bolivia", TASA_ANUAL_BOLIVIA), ("motoclick", TASA_ANUAL_MOTOCLICK), ("futbol", TASA_ANUAL_FUTBOL), ("bitcoin", TASA_ANUAL_BITCOIN)]:
        det = detalle_activo_mes(df_inv, activo, tasa, int(anio), int(mes))
        if not det.empty:
            det["activo"] = activo
            detalles_fijos.append(det)
    det_prestamos = detalle_prestamos_mes(df_inv, int(anio), int(mes))
    if not det_prestamos.empty:
        det_prestamos["activo"] = "prestamo"
        detalles_fijos.append(det_prestamos)
    d_fijos = pd.concat(detalles_fijos, ignore_index=True) if detalles_fijos else pd.DataFrame()
    d_fijos = ajustar_ingreso_motoclick(d_fijos, df_inv, int(anio), int(mes))
    cobro_fijos = d_fijos["ingreso_bruto"].sum() if not d_fijos.empty else 0
    pago_fijos = d_fijos["pago_inversor_mes"].sum() if not d_fijos.empty else 0
    beneficio_fijos = d_fijos["beneficio_empresa_mes"].sum() if not d_fijos.empty else 0

    cobro_total_mes = c_notas + cobro_fijos
    pago_total_mes = p_notas + pago_fijos
    beneficio_total_mes = b_notas + beneficio_fijos

    rentabilidad_beneficio_mes = beneficio_total_mes / capital_total if capital_total else 0
    rentabilidad_beneficio_anualizada = rentabilidad_beneficio_mes * 12
    rentabilidad_pagada_inversor_mes = pago_total_mes / capital_total if capital_total else 0
    rentabilidad_pagada_inversor_anualizada = rentabilidad_pagada_inversor_mes * 12

    rentabilidad_inversiones = calcular_rentabilidad_inversiones_mes(df_inv, df_cal, df_control, int(anio), int(mes), prorratear_notas=prorratear_notas)

    if not rentabilidad_inversiones.empty:
        rentabilidad_por_activo = rentabilidad_inversiones.groupby("activo", as_index=False).agg(
            capital=("capital", "sum"),
            cobro_compania_mes=("cobro_compania_mes", "sum"),
            pago_inversor_mes=("pago_inversor_mes", "sum"),
            beneficio_empresa_mes=("beneficio_empresa_mes", "sum"),
        )
        rentabilidad_por_activo["rentabilidad_beneficio_mes"] = rentabilidad_por_activo.apply(lambda r: r["beneficio_empresa_mes"] / r["capital"] if r["capital"] else 0, axis=1)
        rentabilidad_por_activo["rentabilidad_beneficio_anualizada"] = rentabilidad_por_activo["rentabilidad_beneficio_mes"] * 12
        rentabilidad_por_activo["rentabilidad_pagada_inversor_mes"] = rentabilidad_por_activo.apply(lambda r: r["pago_inversor_mes"] / r["capital"] if r["capital"] else 0, axis=1)
        rentabilidad_por_activo["rentabilidad_pagada_inversor_anualizada"] = rentabilidad_por_activo["rentabilidad_pagada_inversor_mes"] * 12
    else:
        rentabilidad_por_activo = pd.DataFrame()

    # Si el dashboard se filtra por activo, recalculamos los KPIs sobre ese bloque concreto.
    mapa_vista_activo = {
        "Notas": "notas",
        "Fútbol": "futbol",
        "MotoClick": "motoclick",
        "Paraguay": "paraguay",
        "Bolivia": "bolivia",
        "Bitcoin": "bitcoin",
        "Préstamo": "prestamo",
    }
    activo_filtrado = mapa_vista_activo.get(str(vista_activo), None)
    if activo_filtrado:
        activas = activas[activas["activo"] == activo_filtrado].copy() if not activas.empty and "activo" in activas.columns else pd.DataFrame()
        capital_total = activas["capital_invertido"].sum() if not activas.empty else 0

        rentabilidad_inversiones = rentabilidad_inversiones[rentabilidad_inversiones["activo"] == activo_filtrado].copy() if not rentabilidad_inversiones.empty and "activo" in rentabilidad_inversiones.columns else pd.DataFrame()
        cobro_total_mes = float(rentabilidad_inversiones["cobro_compania_mes"].sum()) if not rentabilidad_inversiones.empty else 0.0
        pago_total_mes = float(rentabilidad_inversiones["pago_inversor_mes"].sum()) if not rentabilidad_inversiones.empty else 0.0
        beneficio_total_mes = float(rentabilidad_inversiones["beneficio_empresa_mes"].sum()) if not rentabilidad_inversiones.empty else 0.0

        rentabilidad_beneficio_mes = beneficio_total_mes / capital_total if capital_total else 0
        rentabilidad_beneficio_anualizada = rentabilidad_beneficio_mes * 12
        rentabilidad_pagada_inversor_mes = pago_total_mes / capital_total if capital_total else 0
        rentabilidad_pagada_inversor_anualizada = rentabilidad_pagada_inversor_mes * 12

        if not rentabilidad_inversiones.empty:
            rentabilidad_por_activo = rentabilidad_inversiones.groupby("activo", as_index=False).agg(
                capital=("capital", "sum"),
                cobro_compania_mes=("cobro_compania_mes", "sum"),
                pago_inversor_mes=("pago_inversor_mes", "sum"),
                beneficio_empresa_mes=("beneficio_empresa_mes", "sum"),
            )
            rentabilidad_por_activo["rentabilidad_beneficio_mes"] = rentabilidad_por_activo.apply(lambda r: r["beneficio_empresa_mes"] / r["capital"] if r["capital"] else 0, axis=1)
            rentabilidad_por_activo["rentabilidad_beneficio_anualizada"] = rentabilidad_por_activo["rentabilidad_beneficio_mes"] * 12
            rentabilidad_por_activo["rentabilidad_pagada_inversor_mes"] = rentabilidad_por_activo.apply(lambda r: r["pago_inversor_mes"] / r["capital"] if r["capital"] else 0, axis=1)
            rentabilidad_por_activo["rentabilidad_pagada_inversor_anualizada"] = rentabilidad_por_activo["rentabilidad_pagada_inversor_mes"] * 12
        else:
            rentabilidad_por_activo = pd.DataFrame()

    eventos_futuros = df_cal[(df_cal["fecha"].notna()) & (df_cal["fecha"] >= fecha_analisis)].copy().sort_values("fecha") if not df_cal.empty else pd.DataFrame()
    return {
        "activas": activas,
        "capital_total": capital_total,
        "cobro_total_mes": cobro_total_mes,
        "pago_total_mes": pago_total_mes,
        "beneficio_total_mes": beneficio_total_mes,
        "rentabilidad_beneficio_mes": rentabilidad_beneficio_mes,
        "rentabilidad_beneficio_anualizada": rentabilidad_beneficio_anualizada,
        "rentabilidad_pagada_inversor_mes": rentabilidad_pagada_inversor_mes,
        "rentabilidad_pagada_inversor_anualizada": rentabilidad_pagada_inversor_anualizada,
        "rentabilidad_inversiones": rentabilidad_inversiones,
        "rentabilidad_por_activo": rentabilidad_por_activo,
        "eventos_futuros": eventos_futuros,
        "detalle_notas": detalle_notas,
        "detalle_fijos": d_fijos,
    }


def grafico_capital_por_activo(activas):
    if px is None:
        st.warning("Falta plotly. Añade plotly a requirements.txt.")
        return
    if activas.empty:
        st.info("No hay inversiones activas para graficar.")
        return
    resumen = activas.groupby("activo", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False)
    fig = px.pie(resumen, names="activo", values="capital", hole=0.45, title="Distribución del capital activo por activo")
    fig.update_layout(height=420, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(family="Inter", size=13), title_font=dict(size=20))
    st.plotly_chart(fig, use_container_width=True)


def grafico_capital_por_inversor(activas):
    if px is None:
        st.warning("Falta plotly. Añade plotly a requirements.txt.")
        return
    if activas.empty:
        st.info("No hay inversiones activas para graficar.")
        return
    resumen = activas.groupby("inversor", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False).head(10)
    fig = px.bar(resumen, x="inversor", y="capital", title="Top inversores por capital activo", text_auto=".2s")
    fig.update_layout(height=420, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(family="Inter", size=13), title_font=dict(size=20), xaxis_title="Inversor", yaxis_title="Capital activo")
    st.plotly_chart(fig, use_container_width=True)


def grafico_beneficio_mensual(df_inv_calculo, df_cal, df_control, prorratear_notas: bool = False):
    if px is None:
        st.warning("Falta plotly. Añade plotly a requirements.txt.")
        return
    hoy = pd.Timestamp.today().normalize()
    filas = []
    for i in range(11, -1, -1):
        fecha = hoy - pd.DateOffset(months=i)
        anio, mes = fecha.year, fecha.month
        _, _, b_notas, _, _ = resumen_notas_mes(df_inv, df_cal, df_control, anio, mes, prorratear=prorratear_notas)
        detalles_fijos = []
        for activo, tasa in [("paraguay", TASA_ANUAL_PARAGUAY), ("bolivia", TASA_ANUAL_BOLIVIA), ("motoclick", TASA_ANUAL_MOTOCLICK), ("futbol", TASA_ANUAL_FUTBOL), ("bitcoin", TASA_ANUAL_BITCOIN)]:
            det = detalle_activo_mes(df_inv, activo, tasa, anio, mes)
            if not det.empty:
                detalles_fijos.append(det)
        d_fijos = pd.concat(detalles_fijos, ignore_index=True) if detalles_fijos else pd.DataFrame()
        d_fijos = ajustar_ingreso_motoclick(d_fijos, df_inv, anio, mes)
        b_fijos = d_fijos["beneficio_empresa_mes"].sum() if not d_fijos.empty else 0
        filas.append({"mes": f"{mes:02d}/{anio}", "beneficio": b_notas + b_fijos})
    data = pd.DataFrame(filas)
    fig = px.line(data, x="mes", y="beneficio", markers=True, title="Evolución del beneficio mensual estimado")
    fig.update_layout(height=420, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font=dict(family="Inter", size=13), title_font=dict(size=20), xaxis_title="Mes", yaxis_title="Beneficio")
    st.plotly_chart(fig, use_container_width=True)


def graficos_comparativos_notas(resultados_notas: list):
    """Gráficos comparativos entre las notas candidatas del comparador: barras de score y barras
    agrupadas de las tres probabilidades clave (cupón, call, pérdida de capital). Reutiliza
    resultados_notas tal cual lo produce _tab_comparador_notas — no recalcula nada."""
    if px is None:
        st.warning("Falta plotly. Añade plotly a requirements.txt.")
        return
    if not resultados_notas:
        return

    col1, col2 = st.columns(2)
    with col1:
        df_score = pd.DataFrame([{"Nota": r["nombre"], "Score": r["score"]} for r in resultados_notas])
        fig_score = px.bar(df_score, x="Nota", y="Score", text_auto=".0f", title="Score comparado (0-100)")
        fig_score.update_traces(marker_color="#9A6B24")
        fig_score.update_layout(
            height=340, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
            font=dict(family="Archivo, sans-serif", size=12), title_font=dict(size=15),
            yaxis=dict(range=[0, 100]), showlegend=False,
        )
        st.plotly_chart(fig_score, use_container_width=True)

    with col2:
        filas_prob = []
        for r in resultados_notas:
            filas_prob.append({"Nota": r["nombre"], "Métrica": "Prob. cupón/periodo", "Valor": r["prob_cupon_media"] * 100})
            filas_prob.append({"Nota": r["nombre"], "Métrica": "Prob. call (total)", "Valor": r["prob_call_total"] * 100})
            filas_prob.append({"Nota": r["nombre"], "Métrica": "Prob. pérdida capital", "Valor": r["prob_perdida_capital"] * 100})
        df_prob = pd.DataFrame(filas_prob)
        fig_prob = px.bar(
            df_prob, x="Nota", y="Valor", color="Métrica", barmode="group", text_auto=".0f",
            title="Probabilidades clave por nota (%)",
            color_discrete_map={"Prob. cupón/periodo": "#0E7C5A", "Prob. call (total)": "#9A6B24", "Prob. pérdida capital": "#B03A2E"},
        )
        fig_prob.update_layout(
            height=340, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
            font=dict(family="Archivo, sans-serif", size=12), title_font=dict(size=15),
            legend=dict(orientation="h", yanchor="bottom", y=-0.35), yaxis_title="%",
        )
        st.plotly_chart(fig_prob, use_container_width=True)

    df_renta = pd.DataFrame([{"Nota": r["nombre"], "Rentabilidad neta esperada (anual %)": r["rentabilidad_esperada_neta"] * 100} for r in resultados_notas])
    fig_renta = px.bar(df_renta, x="Nota", y="Rentabilidad neta esperada (anual %)", text_auto=".2f", title="Rentabilidad neta esperada anualizada")
    fig_renta.update_traces(marker_color=["#0E7C5A" if v >= 0 else "#B03A2E" for v in df_renta["Rentabilidad neta esperada (anual %)"]])
    fig_renta.update_layout(
        height=320, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Archivo, sans-serif", size=12), title_font=dict(size=15), showlegend=False,
    )
    st.plotly_chart(fig_renta, use_container_width=True)





def etiqueta_tipo_interes(valor) -> str:
    """Convierte un interés decimal o porcentual en una etiqueta limpia."""
    try:
        v = float(valor)
        pct = v * 100 if abs(v) <= 1 else v
        if abs(pct - round(pct)) < 1e-9:
            return f"{int(round(pct))}%"
        return f"{pct:.2f}%".replace(".00%", "%")
    except Exception:
        return "SIN TIPO"


def construir_desglose_notas_por_tipo_inversor(df_inv: pd.DataFrame, detalle_notas: pd.DataFrame, fecha_analisis) -> pd.DataFrame:
    """Resume NOTAS por tipo pagado al inversor: capital activo, cobros, pagos y beneficio.

    Capital invertido: capital activo vivo al cierre del periodo seleccionado.
    Cobros/pagos/beneficio: importes del mes seleccionado construidos desde CALENDARIO_NOTAS.
    """
    fecha_analisis = pd.Timestamp(fecha_analisis).normalize()

    notas = filtrar_notas(df_inv)
    if notas is None or notas.empty:
        return pd.DataFrame(columns=[
            "tipo_inversor", "num_inversiones", "capital_invertido", "cobro_compania_mes",
            "pago_inversor_mes", "beneficio_empresa_mes", "rentabilidad_bruta_mes",
            "rentabilidad_bruta_anualizada", "coste_inversor_mes", "coste_inversor_anualizado",
            "margen_beneficio_mes", "margen_beneficio_anualizado",
        ])

    notas_activas = notas[
        (notas["fecha_inversion"].notna())
        & (notas["fecha_inversion"] <= fecha_analisis)
        & (notas["fecha_final_inversion"].isna() | (notas["fecha_final_inversion"] >= fecha_analisis))
    ].copy()

    if not notas_activas.empty:
        notas_activas["interes_inversor_anual"] = pd.to_numeric(notas_activas["interes_inversor_anual"], errors="coerce").fillna(0)
        capital_por_tipo = notas_activas.groupby("interes_inversor_anual", as_index=False).agg(
            num_inversiones=("id_inversion", "count"),
            capital_invertido=("capital_invertido", "sum"),
        )
    else:
        capital_por_tipo = pd.DataFrame(columns=["interes_inversor_anual", "num_inversiones", "capital_invertido"])

    if detalle_notas is not None and not detalle_notas.empty:
        det = detalle_notas.copy()
        det["interes_inversor_anual"] = pd.to_numeric(det["interes_inversor_anual"], errors="coerce").fillna(0)
        flujo_por_tipo = det.groupby("interes_inversor_anual", as_index=False).agg(
            cobro_compania_mes=("cobro_compania", "sum"),
            pago_inversor_mes=("pago_inversor", "sum"),
            beneficio_empresa_mes=("beneficio_empresa", "sum"),
        )
    else:
        flujo_por_tipo = pd.DataFrame(columns=["interes_inversor_anual", "cobro_compania_mes", "pago_inversor_mes", "beneficio_empresa_mes"])

    resumen = capital_por_tipo.merge(flujo_por_tipo, on="interes_inversor_anual", how="outer")
    for col in ["num_inversiones", "capital_invertido", "cobro_compania_mes", "pago_inversor_mes", "beneficio_empresa_mes"]:
        if col in resumen.columns:
            resumen[col] = pd.to_numeric(resumen[col], errors="coerce").fillna(0)

    if resumen.empty:
        return resumen

    resumen["tipo_inversor"] = resumen["interes_inversor_anual"].apply(etiqueta_tipo_interes)
    resumen["rentabilidad_bruta_mes"] = resumen.apply(lambda r: r["cobro_compania_mes"] / r["capital_invertido"] if r["capital_invertido"] else 0, axis=1)
    resumen["rentabilidad_bruta_anualizada"] = resumen["rentabilidad_bruta_mes"] * 12
    resumen["coste_inversor_mes"] = resumen.apply(lambda r: r["pago_inversor_mes"] / r["capital_invertido"] if r["capital_invertido"] else 0, axis=1)
    resumen["coste_inversor_anualizado"] = resumen["coste_inversor_mes"] * 12
    resumen["margen_beneficio_mes"] = resumen.apply(lambda r: r["beneficio_empresa_mes"] / r["capital_invertido"] if r["capital_invertido"] else 0, axis=1)
    resumen["margen_beneficio_anualizado"] = resumen["margen_beneficio_mes"] * 12

    total = {
        "interes_inversor_anual": 999,
        "tipo_inversor": "TOTAL",
        "num_inversiones": resumen["num_inversiones"].sum(),
        "capital_invertido": resumen["capital_invertido"].sum(),
        "cobro_compania_mes": resumen["cobro_compania_mes"].sum(),
        "pago_inversor_mes": resumen["pago_inversor_mes"].sum(),
        "beneficio_empresa_mes": resumen["beneficio_empresa_mes"].sum(),
    }
    total["rentabilidad_bruta_mes"] = total["cobro_compania_mes"] / total["capital_invertido"] if total["capital_invertido"] else 0
    total["rentabilidad_bruta_anualizada"] = total["rentabilidad_bruta_mes"] * 12
    total["coste_inversor_mes"] = total["pago_inversor_mes"] / total["capital_invertido"] if total["capital_invertido"] else 0
    total["coste_inversor_anualizado"] = total["coste_inversor_mes"] * 12
    total["margen_beneficio_mes"] = total["beneficio_empresa_mes"] / total["capital_invertido"] if total["capital_invertido"] else 0
    total["margen_beneficio_anualizado"] = total["margen_beneficio_mes"] * 12

    resumen = pd.concat([resumen, pd.DataFrame([total])], ignore_index=True)
    resumen = resumen.sort_values("interes_inversor_anual").reset_index(drop=True)
    columnas = [
        "tipo_inversor", "num_inversiones", "capital_invertido", "cobro_compania_mes",
        "pago_inversor_mes", "beneficio_empresa_mes", "rentabilidad_bruta_mes",
        "rentabilidad_bruta_anualizada", "coste_inversor_mes", "coste_inversor_anualizado",
        "margen_beneficio_mes", "margen_beneficio_anualizado",
    ]
    return resumen[columnas]


def preparar_tabla_tipo_inversor_notas(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return df
    out = df.copy()
    for col in ["capital_invertido", "cobro_compania_mes", "pago_inversor_mes", "beneficio_empresa_mes"]:
        if col in out.columns:
            out[col] = out[col].map(fmt)
    for col in [
        "rentabilidad_bruta_mes", "rentabilidad_bruta_anualizada",
        "coste_inversor_mes", "coste_inversor_anualizado",
        "margen_beneficio_mes", "margen_beneficio_anualizado",
    ]:
        if col in out.columns:
            out[col] = out[col].map(fmt_pct)
    return out


def mostrar_desglose_notas_por_tipo_inversor(df_inv: pd.DataFrame, detalle_notas: pd.DataFrame, fecha_analisis, anio: int, mes: int):
    """Pinta en el dashboard el desglose de notas por 7,5%, 10%, 15%, etc."""
    st.markdown("### Notas por tipo pagado al inversor")
    st.caption("Capital activo de notas y resultado mensual separado por el interés pactado con cada inversor.")
    tabla = construir_desglose_notas_por_tipo_inversor(df_inv, detalle_notas, fecha_analisis)
    if tabla is None or tabla.empty:
        st.info("No hay notas activas o pagos de notas para mostrar en este periodo.")
        return

    columnas = [
        "tipo_inversor", "num_inversiones", "capital_invertido", "cobro_compania_mes",
        "pago_inversor_mes", "beneficio_empresa_mes", "rentabilidad_bruta_anualizada",
        "coste_inversor_anualizado", "margen_beneficio_anualizado",
    ]
    st.dataframe(preparar_tabla_tipo_inversor_notas(tabla[columnas]), use_container_width=True)

    excel_bytes = BytesIO()
    with pd.ExcelWriter(excel_bytes, engine="openpyxl") as writer:
        tabla.to_excel(writer, index=False, sheet_name="NOTAS_TIPO_INVERSOR")
    st.download_button(
        "Descargar desglose de notas por tipo inversor",
        data=excel_bytes.getvalue(),
        file_name=f"notas_por_tipo_inversor_{anio}_{mes:02d}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

def mostrar_rentabilidad_por_activo_dashboard(tabla_activo: pd.DataFrame):
    """Muestra tarjetas y tabla de rentabilidad por activo directamente en el dashboard."""
    st.markdown("### Rentabilidad por activo")
    st.caption("Beneficio y coste de inversores por cada tipo de activo en el mes actual.")

    if tabla_activo is None or tabla_activo.empty:
        st.info("No hay datos de rentabilidad por activo para este mes.")
        return

    tabla = tabla_activo.copy().sort_values("beneficio_empresa_mes", ascending=False)

    cols = st.columns(min(4, len(tabla)))
    for i, (_, row) in enumerate(tabla.iterrows()):
        activo = str(row.get("activo", "Activo")).upper()
        capital = float(row.get("capital", 0) or 0)
        beneficio = float(row.get("beneficio_empresa_mes", 0) or 0)
        rent_mes = float(row.get("rentabilidad_beneficio_mes", 0) or 0)
        rent_anual = float(row.get("rentabilidad_beneficio_anualizada", 0) or 0)
        pagado_mes = float(row.get("rentabilidad_pagada_inversor_mes", 0) or 0)
        pagado_anual = float(row.get("rentabilidad_pagada_inversor_anualizada", 0) or 0)

        subtitulo = (
            f"Capital {fmt(capital)} · Beneficio {fmt(beneficio)}<br>"
            f"Pago inversores {fmt_pct(pagado_mes)} mes / {fmt_pct(pagado_anual)} anual"
        )
        estado = "positivo" if beneficio >= 0 else "negativo"
        with cols[i % len(cols)]:
            tarjeta_kpi(
                f"{activo} · rent. beneficio",
                f"{fmt_pct(rent_mes)} / {fmt_pct(rent_anual)} anual",
                subtitulo,
                estado,
            )

    with st.expander("Ver tabla completa de rentabilidad por activo", expanded=True):
        columnas = [
            "activo",
            "capital",
            "cobro_compania_mes",
            "pago_inversor_mes",
            "beneficio_empresa_mes",
            "rentabilidad_beneficio_mes",
            "rentabilidad_beneficio_anualizada",
            "rentabilidad_pagada_inversor_mes",
            "rentabilidad_pagada_inversor_anualizada",
        ]
        columnas = [c for c in columnas if c in tabla.columns]
        st.dataframe(preparar_tabla_rentabilidad(tabla[columnas]), use_container_width=True)


# =========================
# HISTÓRICO Y PROYECCIONES
# =========================
def primer_dia_mes_fecha(fecha):
    fecha = pd.Timestamp(fecha).normalize()
    return pd.Timestamp(fecha.year, fecha.month, 1)


def etiqueta_mes(fecha):
    fecha = pd.Timestamp(fecha).normalize()
    return f"{fecha.year}-{fecha.month:02d}"


def rango_meses(fecha_inicio, fecha_fin):
    inicio = primer_dia_mes_fecha(fecha_inicio)
    fin = primer_dia_mes_fecha(fecha_fin)
    meses = []
    actual = inicio
    while actual <= fin:
        meses.append(actual)
        actual = actual + pd.DateOffset(months=1)
    return meses


def fecha_minima_sistema(df_inv: pd.DataFrame, df_cal: pd.DataFrame):
    fechas = []
    if df_inv is not None and not df_inv.empty and "fecha_inversion" in df_inv.columns:
        serie = pd.to_datetime(df_inv["fecha_inversion"], errors="coerce").dropna()
        if not serie.empty:
            fechas.append(serie.min())
    if df_cal is not None and not df_cal.empty and "fecha" in df_cal.columns:
        serie = pd.to_datetime(df_cal["fecha"], errors="coerce").dropna()
        if not serie.empty:
            fechas.append(serie.min())
    if fechas:
        return min(fechas).normalize()
    return pd.Timestamp.today().normalize()


def construir_movimientos_historico_proyeccion(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, fecha_inicio, fecha_fin, incluir_chaparro: bool = True) -> pd.DataFrame:
    """Construye movimientos mensuales de cobros, pagos y beneficio desde inicio y hacia futuro.

    Reglas:
    - NOTAS: usa CALENDARIO_NOTAS. Cada PAGO genera cobro compañía, pago inversor y beneficio por inversión.
    - Paraguay, MotoClick y Fútbol: devenga mes a mes según fecha_inversion / fecha_final_inversion.
    - Histórico/proyección se clasifica según si el mes es anterior o posterior al mes actual.
    """
    df_inv = aplicar_filtro_chaparro_fernandez(df_inv, incluir_chaparro)
    filas = []
    hoy = pd.Timestamp.today().normalize()

    for fecha_mes in rango_meses(fecha_inicio, fecha_fin):
        anio = int(fecha_mes.year)
        mes = int(fecha_mes.month)
        fin_mes = pd.Timestamp(anio, mes, ultimo_dia_mes(anio, mes)).normalize()
        tipo_dato = "HISTÓRICO" if fin_mes <= hoy else "PROYECCIÓN"
        mes_label = etiqueta_mes(fecha_mes)

        # 1) Notas: se calculan únicamente cuando hay evento PAGO en calendario.
        _, _, _, detalle_notas, _ = resumen_notas_mes(df_inv, df_cal, df_control, anio, mes)
        if detalle_notas is not None and not detalle_notas.empty:
            for _, row in detalle_notas.iterrows():
                nota = row.get("nota", "")
                filas.append({
                    "mes_fecha": fecha_mes,
                    "mes": mes_label,
                    "tipo_dato": tipo_dato,
                    "activo": "notas",
                    "nombre_activo": f"NOTA {nota}",
                    "nota": nota,
                    "id_inversion": row.get("id_inversion", ""),
                    "inversor": row.get("inversor", ""),
                    "capital_base": float(row.get("capital_invertido", 0) or 0),
                    "cobrado_compania": float(row.get("cobro_compania", 0) or 0),
                    "pagado_inversores": float(row.get("pago_inversor", 0) or 0),
                    "beneficio_empresa": float(row.get("beneficio_empresa", 0) or 0),
                    "resultado_observacion": row.get("resultado_observacion", ""),
                })

        # 2) Activos con ingreso fijo o operativo.
        for activo, tasa in [("paraguay", TASA_ANUAL_PARAGUAY), ("bolivia", TASA_ANUAL_BOLIVIA), ("motoclick", TASA_ANUAL_MOTOCLICK), ("futbol", TASA_ANUAL_FUTBOL), ("bitcoin", TASA_ANUAL_BITCOIN)]:
            det = detalle_activo_mes(df_inv, activo, tasa, anio, mes)
            if det is None or det.empty:
                continue
            for _, row in det.iterrows():
                cobro = float(row.get("ingreso_bruto", 0) or 0)
                pago = float(row.get("pago_inversor_mes", 0) or 0)
                filas.append({
                    "mes_fecha": fecha_mes,
                    "mes": mes_label,
                    "tipo_dato": tipo_dato,
                    "activo": activo,
                    "nombre_activo": activo,
                    "nota": "",
                    "id_inversion": row.get("id_inversion", ""),
                    "inversor": row.get("inversor", ""),
                    "capital_base": float(row.get("capital_invertido", 0) or 0),
                    "cobrado_compania": cobro,
                    "pagado_inversores": pago,
                    "beneficio_empresa": cobro - pago,
                    "resultado_observacion": "NO APLICA",
                })

    if not filas:
        return pd.DataFrame(columns=[
            "mes_fecha", "mes", "tipo_dato", "activo", "nombre_activo", "nota", "id_inversion", "inversor",
            "capital_base", "cobrado_compania", "pagado_inversores", "beneficio_empresa", "resultado_observacion"
        ])

    out = pd.DataFrame(filas)
    out = out.sort_values(["mes_fecha", "activo", "nombre_activo", "inversor", "id_inversion"]).reset_index(drop=True)
    return out


def resumir_movimientos_por_mes(movimientos: pd.DataFrame) -> pd.DataFrame:
    if movimientos is None or movimientos.empty:
        return pd.DataFrame(columns=["mes", "tipo_dato", "cobrado_compania", "pagado_inversores", "beneficio_empresa"])
    resumen = movimientos.groupby(["mes_fecha", "mes", "tipo_dato"], as_index=False).agg(
        cobrado_compania=("cobrado_compania", "sum"),
        pagado_inversores=("pagado_inversores", "sum"),
        beneficio_empresa=("beneficio_empresa", "sum"),
    )
    return resumen.sort_values("mes_fecha")


def resumir_movimientos_por_mes_activo(movimientos: pd.DataFrame) -> pd.DataFrame:
    if movimientos is None or movimientos.empty:
        return pd.DataFrame(columns=["mes", "tipo_dato", "activo", "cobrado_compania", "pagado_inversores", "beneficio_empresa"])
    resumen = movimientos.groupby(["mes_fecha", "mes", "tipo_dato", "activo"], as_index=False).agg(
        cobrado_compania=("cobrado_compania", "sum"),
        pagado_inversores=("pagado_inversores", "sum"),
        beneficio_empresa=("beneficio_empresa", "sum"),
    )
    return resumen.sort_values(["mes_fecha", "activo"])


def resumir_movimientos_por_inversion(movimientos: pd.DataFrame) -> pd.DataFrame:
    if movimientos is None or movimientos.empty:
        return pd.DataFrame(columns=["activo", "nombre_activo", "id_inversion", "inversor", "cobrado_compania", "pagado_inversores", "beneficio_empresa"])
    resumen = movimientos.groupby(["activo", "nombre_activo", "id_inversion", "inversor"], as_index=False).agg(
        capital_base=("capital_base", "max"),
        cobrado_compania=("cobrado_compania", "sum"),
        pagado_inversores=("pagado_inversores", "sum"),
        beneficio_empresa=("beneficio_empresa", "sum"),
    )
    return resumen.sort_values("beneficio_empresa", ascending=False)


def preparar_tabla_movimientos(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return df
    out = df.copy()
    for col in ["mes_fecha"]:
        if col in out.columns:
            out[col] = pd.to_datetime(out[col], errors="coerce").dt.strftime("%d/%m/%Y")
    for col in ["capital_base", "cobrado_compania", "pagado_inversores", "beneficio_empresa", "cobro_simulado", "pago_simulado", "beneficio_simulado", "cobrado_total_con_simulacion", "pagado_total_con_simulacion", "beneficio_total_con_simulacion"]:
        if col in out.columns:
            out[col] = out[col].map(fmt)
    return out


def construir_simulacion_capital_extra(fecha_inicio, meses_duracion: int, capital_extra: float, tasa_cobro_anual: float, tasa_pago_anual: float, fecha_fin_rango) -> pd.DataFrame:
    if capital_extra <= 0 or meses_duracion <= 0:
        return pd.DataFrame()
    fecha_inicio = primer_dia_mes_fecha(fecha_inicio)
    fecha_fin_rango = primer_dia_mes_fecha(fecha_fin_rango)
    filas = []
    for i in range(int(meses_duracion)):
        fecha_mes = fecha_inicio + pd.DateOffset(months=i)
        if fecha_mes > fecha_fin_rango:
            break
        cobro = float(capital_extra) * float(tasa_cobro_anual) / 12
        pago = float(capital_extra) * float(tasa_pago_anual) / 12
        filas.append({
            "mes_fecha": fecha_mes,
            "mes": etiqueta_mes(fecha_mes),
            "capital_simulado": float(capital_extra),
            "cobro_simulado": cobro,
            "pago_simulado": pago,
            "beneficio_simulado": cobro - pago,
        })
    return pd.DataFrame(filas)


def exportar_historico_proyeccion_excel(resumen_mes, resumen_activo, resumen_inversion, detalle, simulacion=None) -> bytes:
    salida = BytesIO()
    with pd.ExcelWriter(salida, engine="openpyxl") as writer:
        resumen_mes.to_excel(writer, index=False, sheet_name="RESUMEN_MENSUAL")
        resumen_activo.to_excel(writer, index=False, sheet_name="MES_ACTIVO")
        resumen_inversion.to_excel(writer, index=False, sheet_name="POR_INVERSION")
        detalle.to_excel(writer, index=False, sheet_name="DETALLE")
        if simulacion is not None and not simulacion.empty:
            simulacion.to_excel(writer, index=False, sheet_name="SIMULACION")
    return salida.getvalue()


def seccion_historico_y_proyecciones():
    df_inv, df_cal, df_control = cargar_excel_completo()
    st.markdown("## Histórico y proyecciones")
    st.caption("Control mensual de cuánto se cobra, cuánto se paga y qué beneficio queda desde el inicio, con proyección futura y simulador de nuevas inversiones.")

    hoy = pd.Timestamp.today().normalize()
    fecha_inicio_default = fecha_minima_sistema(df_inv, df_cal)
    fecha_fin_default = hoy + pd.DateOffset(months=12)

    c1, c2, c3, c4 = st.columns([1, 1, 1, 1.2])
    fecha_inicio = pd.Timestamp(c1.date_input("Desde", value=fecha_inicio_default.date(), key="hist_proj_desde")).normalize()
    fecha_fin = pd.Timestamp(c2.date_input("Hasta", value=fecha_fin_default.date(), key="hist_proj_hasta")).normalize()
    activo_filtro = c3.selectbox("Activo", ["Todos", "notas", "paraguay", "bolivia", "motoclick", "futbol", "bitcoin"], key="hist_proj_activo")
    incluir_chaparro = c4.checkbox(
        "Incluir Chaparro Fernández",
        value=False,
        key="hist_proj_incluir_chaparro",
        help="Si está desactivado, Chaparro Fernández queda fuera de capital, cobros, pagos y beneficio. Si está activado, se incluye como interno: pago inversor = cobro de la nota y beneficio = 0.",
    )

    if fecha_fin < fecha_inicio:
        st.error("La fecha final no puede ser anterior a la fecha inicial.")
        return

    df_inv_marcado = aplicar_filtro_chaparro_fernandez(df_inv, True)
    inversiones_chaparro = df_inv_marcado[df_inv_marcado.get("es_chaparro_fernandez", False) == True].copy() if not df_inv_marcado.empty else pd.DataFrame()
    if not inversiones_chaparro.empty:
        with st.expander("Ver inversiones detectadas como Chaparro Fernández", expanded=False):
            columnas_auditoria = [c for c in ["id_inversion", "inversor", "tipo_inversion", "subtipo_inversion", "nombre_activo", "capital_invertido", "interes_inversor_anual", "fecha_inversion", "fecha_final_inversion"] if c in inversiones_chaparro.columns]
            st.dataframe(preparar_tabla_monetaria(inversiones_chaparro[columnas_auditoria], ["capital_invertido"]), use_container_width=True)

    with st.spinner("Calculando histórico y proyecciones..."):
        movimientos = construir_movimientos_historico_proyeccion(
            df_inv,
            df_cal,
            df_control,
            fecha_inicio,
            fecha_fin,
            incluir_chaparro=incluir_chaparro,
        )

    if activo_filtro != "Todos" and not movimientos.empty:
        movimientos = movimientos[movimientos["activo"] == activo_filtro].copy()

    resumen_mes = resumir_movimientos_por_mes(movimientos)
    resumen_activo = resumir_movimientos_por_mes_activo(movimientos)
    resumen_inversion = resumir_movimientos_por_inversion(movimientos)

    historico = movimientos[movimientos["tipo_dato"] == "HISTÓRICO"].copy() if not movimientos.empty else pd.DataFrame()
    futuro = movimientos[movimientos["tipo_dato"] == "PROYECCIÓN"].copy() if not movimientos.empty else pd.DataFrame()

    total_cobrado_hist = float(historico["cobrado_compania"].sum()) if not historico.empty else 0.0
    total_pagado_hist = float(historico["pagado_inversores"].sum()) if not historico.empty else 0.0
    total_beneficio_hist = float(historico["beneficio_empresa"].sum()) if not historico.empty else 0.0

    total_cobrado_fut = float(futuro["cobrado_compania"].sum()) if not futuro.empty else 0.0
    total_pagado_fut = float(futuro["pagado_inversores"].sum()) if not futuro.empty else 0.0
    total_beneficio_fut = float(futuro["beneficio_empresa"].sum()) if not futuro.empty else 0.0

    st.markdown("### Totales")
    k1, k2, k3 = st.columns(3)
    k1.metric("Cobrado histórico", fmt(total_cobrado_hist))
    k2.metric("Pagado histórico", fmt(total_pagado_hist))
    k3.metric("Beneficio histórico", fmt(total_beneficio_hist))

    k4, k5, k6 = st.columns(3)
    k4.metric("Cobro proyectado", fmt(total_cobrado_fut))
    k5.metric("Pago proyectado", fmt(total_pagado_fut))
    k6.metric("Beneficio proyectado", fmt(total_beneficio_fut))

    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "Mes a mes",
        "Mes + activo",
        "Por inversión",
        "Simulador",
        "Detalle completo",
    ])

    with tab1:
        st.caption("Aquí ves, mes a mes, cuánto se cobra, cuánto se paga y qué beneficio queda.")
        if resumen_mes.empty:
            st.info("No hay movimientos para el rango seleccionado.")
        else:
            st.dataframe(preparar_tabla_movimientos(resumen_mes), use_container_width=True)
            if px is not None:
                fig = px.bar(
                    resumen_mes,
                    x="mes",
                    y=["cobrado_compania", "pagado_inversores", "beneficio_empresa"],
                    title="Cobrado, pagado y beneficio por mes",
                    barmode="group",
                )
                fig.update_layout(height=420, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)")
                st.plotly_chart(fig, use_container_width=True)

    with tab2:
        st.caption("El mismo cálculo mensual, pero separado por activo.")
        if resumen_activo.empty:
            st.info("No hay movimientos por activo para el rango seleccionado.")
        else:
            st.dataframe(preparar_tabla_movimientos(resumen_activo), use_container_width=True)

    with tab3:
        st.caption("Total cobrado, pagado y beneficio acumulado por cada inversión dentro del rango seleccionado.")
        if resumen_inversion.empty:
            st.info("No hay movimientos por inversión para el rango seleccionado.")
        else:
            st.dataframe(preparar_tabla_movimientos(resumen_inversion), use_container_width=True)

    with tab4:
        st.caption("Simula una nueva entrada de capital. Ejemplo: 200.000 al 25% de cobro y pagando 10% al inversor.")
        s1, s2, s3, s4 = st.columns(4)
        capital_extra = float(s1.number_input("Capital nuevo", min_value=0.0, value=200000.0, step=10000.0, key="sim_capital_nuevo"))
        tasa_cobro_pct = float(s2.number_input("% cobro anual compañía", min_value=0.0, value=25.0, step=0.5, key="sim_cobro_pct"))
        tasa_pago_pct = float(s3.number_input("% pago anual inversor", min_value=0.0, value=10.0, step=0.5, key="sim_pago_pct"))
        meses_duracion = int(s4.number_input("Meses de duración", min_value=1, max_value=120, value=12, step=1, key="sim_meses"))

        s5, s6 = st.columns(2)
        fecha_inicio_sim = pd.Timestamp(s5.date_input("Inicio simulación", value=hoy.date(), key="sim_fecha_inicio")).normalize()
        nombre_sim = s6.text_input("Nombre escenario", value="Nueva inversión simulada", key="sim_nombre")

        simulacion = construir_simulacion_capital_extra(
            fecha_inicio=fecha_inicio_sim,
            meses_duracion=meses_duracion,
            capital_extra=capital_extra,
            tasa_cobro_anual=tasa_cobro_pct / 100,
            tasa_pago_anual=tasa_pago_pct / 100,
            fecha_fin_rango=fecha_fin,
        )

        if simulacion.empty:
            st.info("La simulación no genera movimientos dentro del rango seleccionado.")
        else:
            total_sim_cobro = float(simulacion["cobro_simulado"].sum())
            total_sim_pago = float(simulacion["pago_simulado"].sum())
            total_sim_beneficio = float(simulacion["beneficio_simulado"].sum())
            mensual_cobro = float(capital_extra * (tasa_cobro_pct / 100) / 12)
            mensual_pago = float(capital_extra * (tasa_pago_pct / 100) / 12)
            mensual_beneficio = mensual_cobro - mensual_pago

            p1, p2, p3 = st.columns(3)
            p1.metric("Cobro mensual simulado", fmt(mensual_cobro))
            p2.metric("Pago mensual simulado", fmt(mensual_pago))
            p3.metric("Beneficio mensual simulado", fmt(mensual_beneficio))

            p4, p5, p6 = st.columns(3)
            p4.metric("Cobro total simulado", fmt(total_sim_cobro))
            p5.metric("Pago total simulado", fmt(total_sim_pago))
            p6.metric("Beneficio total simulado", fmt(total_sim_beneficio))

            base_mes = resumen_mes[["mes_fecha", "mes", "tipo_dato", "cobrado_compania", "pagado_inversores", "beneficio_empresa"]].copy() if not resumen_mes.empty else pd.DataFrame()
            if base_mes.empty:
                base_mes = pd.DataFrame(columns=["mes_fecha", "mes", "tipo_dato", "cobrado_compania", "pagado_inversores", "beneficio_empresa"])
            combinado = base_mes.merge(simulacion, on=["mes_fecha", "mes"], how="outer")
            combinado["tipo_dato"] = combinado["tipo_dato"].fillna("PROYECCIÓN")
            for col in ["cobrado_compania", "pagado_inversores", "beneficio_empresa", "cobro_simulado", "pago_simulado", "beneficio_simulado"]:
                if col in combinado.columns:
                    combinado[col] = pd.to_numeric(combinado[col], errors="coerce").fillna(0)
            combinado["cobrado_total_con_simulacion"] = combinado["cobrado_compania"] + combinado["cobro_simulado"]
            combinado["pagado_total_con_simulacion"] = combinado["pagado_inversores"] + combinado["pago_simulado"]
            combinado["beneficio_total_con_simulacion"] = combinado["beneficio_empresa"] + combinado["beneficio_simulado"]
            combinado = combinado.sort_values("mes_fecha")

            st.markdown(f"#### Resultado combinado: {nombre_sim}")
            columnas_sim = [
                "mes", "tipo_dato", "cobrado_compania", "pagado_inversores", "beneficio_empresa",
                "cobro_simulado", "pago_simulado", "beneficio_simulado",
                "cobrado_total_con_simulacion", "pagado_total_con_simulacion", "beneficio_total_con_simulacion",
            ]
            st.dataframe(preparar_tabla_movimientos(combinado[columnas_sim]), use_container_width=True)

    with tab5:
        st.caption("Detalle línea a línea usado para construir todos los cálculos.")
        if movimientos.empty:
            st.info("No hay detalle para el rango seleccionado.")
        else:
            columnas = [
                "mes", "tipo_dato", "activo", "nombre_activo", "nota", "id_inversion", "inversor", "capital_base",
                "cobrado_compania", "pagado_inversores", "beneficio_empresa", "resultado_observacion"
            ]
            columnas = [c for c in columnas if c in movimientos.columns]
            st.dataframe(preparar_tabla_movimientos(movimientos[columnas]), use_container_width=True)

    st.markdown("### Exportar")
    excel_bytes = exportar_historico_proyeccion_excel(resumen_mes, resumen_activo, resumen_inversion, movimientos)
    st.download_button(
        "Descargar histórico y proyecciones en Excel",
        data=excel_bytes,
        file_name=f"historico_proyecciones_{fecha_inicio.strftime('%Y%m%d')}_{fecha_fin.strftime('%Y%m%d')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def boton_descarga_excel(df: pd.DataFrame, nombre_archivo: str, label: str = "⬇️ Descargar Excel"):
    """Muestra un botón para descargar cualquier DataFrame como .xlsx."""
    import io
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Datos")
    buf.seek(0)
    st.download_button(
        label=label,
        data=buf.getvalue(),
        file_name=nombre_archivo,
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        key=f"dl_{nombre_archivo}_{id(df)}",
    )


@st.cache_data(show_spinner=False, ttl=1800)
def obtener_estado_bitcoin_etf():
    """
    Precio actual y variación del ETF de Bitcoin (IBIT) respecto al precio de compra,
    más histórico de precios de los últimos 6 meses para el gráfico desplegable.
    """
    resultado = {
        "precio_actual": None,
        "variacion_%": None,
        "valor_actual_posicion": None,
        "ganancia_perdida": None,
        "historico": pd.DataFrame(),
    }
    if yf is None:
        return resultado
    try:
        ticker_obj = yf.Ticker(BITCOIN_ETF_TICKER)
        hist = ticker_obj.history(period="6mo")
        if hist is not None and not hist.empty:
            cierres = hist["Close"].dropna()
            if not cierres.empty:
                precio_actual = float(cierres.iloc[-1])
                resultado["precio_actual"] = precio_actual
                resultado["variacion_%"] = (precio_actual - BITCOIN_ETF_PRECIO_COMPRA) / BITCOIN_ETF_PRECIO_COMPRA * 100
                valor_actual = precio_actual * BITCOIN_ETF_UNIDADES
                resultado["valor_actual_posicion"] = valor_actual
                resultado["ganancia_perdida"] = valor_actual - BITCOIN_ETF_COSTE_REAL
                cierres.index = pd.to_datetime(cierres.index)
                if cierres.index.tz is not None:
                    cierres.index = cierres.index.tz_localize(None)
                resultado["historico"] = cierres
    except Exception:
        pass
    return resultado


def calcular_concentracion_cartera(df_inv: pd.DataFrame, df_control: pd.DataFrame) -> dict:
    """
    Concentración del capital activo hoy por (1) emisor/banco de las notas y (2) sector/industria
    de las compañías subyacentes. El sector/industria viene de Yahoo Finance (clasificación real,
    no una categoría "IA" inventada — Yahoo no tiene esa granularidad, se usa su campo 'industry',
    que es más fino que 'sector' pero tampoco es una etiqueta de "inteligencia artificial" como tal).
    """
    hoy = pd.Timestamp.today().normalize()
    resultado = {"por_emisor": pd.DataFrame(), "por_sector": pd.DataFrame(), "por_industria": pd.DataFrame(), "notas_sin_emisor": []}
    if df_control is None or df_control.empty or "nota" not in df_control.columns:
        return resultado

    control = df_control.copy()
    control["nota"] = pd.to_numeric(control["nota"], errors="coerce")
    control = control.dropna(subset=["nota"])

    filas_emisor, filas_sector, filas_industria, notas_sin_emisor = [], [], [], []
    for nota_num in sorted(control["nota"].unique()):
        filas_nota = control[control["nota"] == nota_num]
        nombre_activo = f"NOTA_{int(nota_num):02d}"
        capital_nota = capital_activo_en_fecha(
            df_inv[df_inv.get("nombre_activo", pd.Series(dtype=str)).astype(str).str.upper().str.replace(" ", "_") == nombre_activo],
            hoy,
        )
        if capital_nota <= 0:
            continue

        emisor = filas_nota.get("emisor", pd.Series(dtype=str)).dropna()
        emisor = str(emisor.iloc[0]).strip() if not emisor.empty and str(emisor.iloc[0]).strip() else None
        if emisor:
            filas_emisor.append({"Emisor": emisor, "Capital": capital_nota, "Nota": int(nota_num)})
        else:
            notas_sin_emisor.append(int(nota_num))
            filas_emisor.append({"Emisor": "Sin emisor registrado", "Capital": capital_nota, "Nota": int(nota_num)})

        primer_ticker = str(filas_nota.iloc[0].get("ticker", "")).strip().upper() if not filas_nota.empty else ""
        if primer_ticker:
            fd = obtener_datos_fundamentales(primer_ticker)
            sector = fd.get("sector") or "Sin sector (Yahoo Finance no lo devolvió)"
            industria = fd.get("industria") or "Sin industria (Yahoo Finance no lo devolvió)"
            filas_sector.append({"Sector": sector, "Capital": capital_nota, "Nota": int(nota_num), "Ticker representativo": primer_ticker})
            filas_industria.append({"Industria": industria, "Capital": capital_nota, "Nota": int(nota_num), "Ticker representativo": primer_ticker})

    if filas_emisor:
        df_e = pd.DataFrame(filas_emisor).groupby("Emisor", as_index=False)["Capital"].sum().sort_values("Capital", ascending=False)
        total_e = df_e["Capital"].sum()
        df_e["% del total"] = (df_e["Capital"] / total_e * 100).round(1)
        resultado["por_emisor"] = df_e

    if filas_sector:
        df_s = pd.DataFrame(filas_sector).groupby("Sector", as_index=False)["Capital"].sum().sort_values("Capital", ascending=False)
        total_s = df_s["Capital"].sum()
        df_s["% del total"] = (df_s["Capital"] / total_s * 100).round(1)
        resultado["por_sector"] = df_s

    if filas_industria:
        df_i = pd.DataFrame(filas_industria).groupby("Industria", as_index=False)["Capital"].sum().sort_values("Capital", ascending=False)
        total_i = df_i["Capital"].sum()
        df_i["% del total"] = (df_i["Capital"] / total_i * 100).round(1)
        resultado["por_industria"] = df_i

    resultado["notas_sin_emisor"] = notas_sin_emisor
    return resultado


def _tab_concentracion_cartera(df_inv: pd.DataFrame, df_control: pd.DataFrame):
    st.markdown("### 🧭 Concentración de la cartera de notas")
    st.caption(
        "Capital activo hoy repartido por emisor (banco), por sector y por industria de la compañía "
        "subyacente (clasificación real de Yahoo Finance — la industria es más fina que el sector, ej. "
        "sector 'Technology' → industria 'Semiconductors' o 'Software - Infrastructure')."
    )
    with st.spinner("Calculando concentración..."):
        conc = calcular_concentracion_cartera(df_inv, df_control)

    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("#### Por emisor (banco)")
        if conc["por_emisor"].empty:
            st.info("No hay datos suficientes (revisa que las notas tengan emisor guardado en CONTROL_NOTAS).")
        else:
            df_mostrar = conc["por_emisor"][["Emisor", "Capital", "% del total"]].copy()
            df_mostrar["Capital"] = df_mostrar["Capital"].apply(lambda x: fmt(x))
            df_mostrar["% del total"] = df_mostrar["% del total"].apply(lambda x: f"{x:.1f}%")
            st.dataframe(df_mostrar, use_container_width=True, hide_index=True)
            st.bar_chart(conc["por_emisor"].set_index("Emisor")["Capital"])

    with col2:
        st.markdown("#### Por sector (compañía subyacente)")
        if conc["por_sector"].empty:
            st.info("No hay datos suficientes.")
        else:
            df_mostrar = conc["por_sector"][["Sector", "Capital", "% del total"]].copy()
            df_mostrar["Capital"] = df_mostrar["Capital"].apply(lambda x: fmt(x))
            df_mostrar["% del total"] = df_mostrar["% del total"].apply(lambda x: f"{x:.1f}%")
            st.dataframe(df_mostrar, use_container_width=True, hide_index=True)
            st.bar_chart(conc["por_sector"].set_index("Sector")["Capital"])

    with col3:
        st.markdown("#### Por industria (más granular)")
        if conc["por_industria"].empty:
            st.info("No hay datos suficientes.")
        else:
            df_mostrar = conc["por_industria"][["Industria", "Capital", "% del total"]].copy()
            df_mostrar["Capital"] = df_mostrar["Capital"].apply(lambda x: fmt(x))
            df_mostrar["% del total"] = df_mostrar["% del total"].apply(lambda x: f"{x:.1f}%")
            st.dataframe(df_mostrar, use_container_width=True, hide_index=True)
            st.bar_chart(conc["por_industria"].set_index("Industria")["Capital"])

    if conc.get("notas_sin_emisor"):
        st.warning(
            f"⚠️ Estas notas no tienen emisor guardado y aparecen como 'Sin emisor registrado': "
            f"{', '.join(str(n) for n in conc['notas_sin_emisor'])}. Es porque se añadieron antes de "
            f"guardar el emisor — puedes rellenarlo directamente en la hoja CONTROL_NOTAS (columna EMISOR)."
        )


def tarjeta_bitcoin_etf():
    """Tarjeta desplegable del Dashboard con precio actualizado y gráfico del ETF de Bitcoin (IBIT)."""
    estado = obtener_estado_bitcoin_etf()
    precio_actual = estado.get("precio_actual")
    variacion = estado.get("variacion_%")
    valor_actual = estado.get("valor_actual_posicion")
    ganancia = estado.get("ganancia_perdida")

    if precio_actual is None:
        st.info(f"No se pudo obtener el precio actual de {BITCOIN_ETF_TICKER} en este momento.")
        return

    color_estado = "positivo" if variacion is not None and variacion >= 0 else "negativo"
    flecha = "▲" if variacion is not None and variacion >= 0 else "▼"

    with st.expander(f"₿ {BITCOIN_ETF_NOMBRE} ({BITCOIN_ETF_TICKER}) — {precio_actual:,.2f} $ {flecha} {variacion:+.2f}%", expanded=False):
        b1, b2, b3, b4 = st.columns(4)
        with b1:
            tarjeta_kpi("Precio de compra", f"${BITCOIN_ETF_PRECIO_COMPRA:,.4f}", f"{BITCOIN_ETF_UNIDADES:,} unidades · 30/03/2026", "normal")
        with b2:
            tarjeta_kpi("Precio actual", f"${precio_actual:,.4f}", "Última cotización disponible", color_estado)
        with b3:
            tarjeta_kpi("Variación", f"{variacion:+.2f}%", "vs. precio de compra", color_estado)
        with b4:
            tarjeta_kpi("Valor posición hoy", fmt(valor_actual), f"{'Ganancia' if ganancia >= 0 else 'Pérdida'}: {fmt(ganancia)}", color_estado)

        historico = estado.get("historico")
        if historico is not None and not historico.empty:
            df_chart = historico.rename("Precio IBIT").to_frame()
            df_chart["Precio de compra"] = BITCOIN_ETF_PRECIO_COMPRA
            st.line_chart(df_chart, height=280)
        st.caption(f"Compra: {BITCOIN_ETF_UNIDADES:,} unidades a ${BITCOIN_ETF_PRECIO_COMPRA:,.4f} el 30/03/2026 · Coste total ${BITCOIN_ETF_COSTE_REAL:,.2f} · Capital fondo ${BITCOIN_ETF_CAPITAL_INVERTIDO:,.2f}")


def dashboard_financiero():
    df_inv, df_cal, df_control = cargar_excel_completo()


    st.markdown("## Dashboard financiero")
    st.caption("Panel ejecutivo de capital activo, cobros, pagos, beneficio y rentabilidades.")

    hoy = pd.Timestamp.today().normalize()
    col_activo, col_periodo_1, col_periodo_2, col_chaparro, col_prorrateo, col_devengo = st.columns([1.2, 0.8, 0.8, 1.0, 1.0, 1.0])
    vista_dashboard = col_activo.selectbox(
        "Dashboard",
        ["General", "Notas", "Fútbol", "MotoClick", "Paraguay", "Bolivia", "Bitcoin", "Préstamo"],
        key="dashboard_vista_activo",
    )
    incluir_chaparro = col_chaparro.checkbox(
        "Incluir Chaparro Fernández",
        value=False,
        key="dashboard_incluir_chaparro",
        help="Si está desactivado, Chaparro Fernández queda fuera de capital, cobros, pagos y rentabilidad. Si está activado, se incluye como interno: pago inversor = cobro de la nota y beneficio = 0.",
    )
    prorratear_notas = col_prorrateo.checkbox(
        "Prorratear cobro de notas",
        value=False,
        key="dashboard_prorratear_notas",
        help="Si está activado, el cobro de notas trimestrales/semestrales se reparte a partes iguales entre los meses del periodo (devengo mensual) en vez de aparecer todo de golpe en el mes de cobro real. Solo afecta a los INGRESOS de notas mostrados aquí y en el desglose de Notas; el pago a inversores no cambia, y no afecta a Fútbol, MotoClick, Paraguay, Bolivia ni Bitcoin.",
    )
    incluir_devengado = col_devengo.checkbox(
        "Incluir interés devengado",
        value=False,
        key="dashboard_incluir_devengado",
        help="Si está activado, añade una tarjeta con el cobro real del mes + el interés corrido de cada nota desde su último pago real hasta fin del mes seleccionado (ACT/365), y un detalle por nota. Es específico del mes elegido arriba, no se acumula entre meses. No afecta al pago a inversores ni al beneficio.",
    )
    anio_dashboard = int(col_periodo_1.number_input(
        "Año del dashboard",
        min_value=2020,
        max_value=2100,
        value=hoy.year,
        key="dashboard_anio_general",
    ))
    mes_dashboard = int(col_periodo_2.number_input(
        "Mes del dashboard",
        min_value=1,
        max_value=12,
        value=hoy.month,
        key="dashboard_mes_general",
    ))
    st.caption(
        f"Vista seleccionada: {vista_dashboard} · Periodo: {nombre_mes_es(mes_dashboard)} {anio_dashboard} · "
        f"Chaparro Fernández: {'incluido' if incluir_chaparro else 'excluido'} · "
        f"Cobro de notas: {'prorrateado (devengo mensual)' if prorratear_notas else 'real (calendario)'}"
    )

    df_inv_marcado = aplicar_filtro_chaparro_fernandez(df_inv, True)
    inversiones_chaparro = df_inv_marcado[df_inv_marcado.get("es_chaparro_fernandez", False) == True].copy() if not df_inv_marcado.empty else pd.DataFrame()
    if not inversiones_chaparro.empty:
        with st.expander("Ver inversiones detectadas como Chaparro Fernández", expanded=False):
            columnas_auditoria = [c for c in ["id_inversion", "inversor", "tipo_inversion", "subtipo_inversion", "nombre_activo", "capital_invertido", "interes_inversor_anual", "fecha_inversion", "fecha_final_inversion"] if c in inversiones_chaparro.columns]
            st.dataframe(preparar_tabla_monetaria(inversiones_chaparro[columnas_auditoria], ["capital_invertido"]), use_container_width=True)

    df_control_dashboard = obtener_control_notas_activas(df_inv, df_control)
    resumen_notas_actual = construir_resumen_actual_notas_alertas(df_control_dashboard)
    alertas_notas = resumen_alertas_por_nota(resumen_notas_actual)
    if not alertas_notas.empty:
        rojas = int((alertas_notas["alerta"] == "ROJO").sum())
        st.error(f"Alertas de notas: {rojas} nota(s) en riesgo (variación ≤ -30% en al menos un ticker).")
        with st.expander("Ver alertas de notas por variación", expanded=False):
            tabla_alertas = alertas_notas.copy()
            tabla_alertas["peor_variacion_%"] = tabla_alertas["peor_variacion_%"].apply(lambda x: f"{float(x):.2f}%" if pd.notna(x) else "Sin dato")
            st.dataframe(tabla_alertas, use_container_width=True)
            boton_descarga_excel(alertas_notas, "alertas_notas.xlsx")

    resumen = obtener_resumen_dashboard(
        df_inv,
        df_cal,
        df_control,
        anio_dashboard,
        mes_dashboard,
        vista_dashboard,
        incluir_chaparro=incluir_chaparro,
        prorratear_notas=prorratear_notas,
    )
    df_inv_calculo = aplicar_filtro_chaparro_fernandez(df_inv, incluir_chaparro)

    detalle_devengo = pd.DataFrame()
    total_devengado = 0.0
    if incluir_devengado:
        detalle_devengo = interes_devengado_no_cobrado_notas(df_inv, df_cal, df_control, anio_dashboard, mes_dashboard)
        total_devengado = float(detalle_devengo["interes_devengado_no_cobrado"].sum()) if not detalle_devengo.empty else 0.0

    cols_kpi = st.columns(4)
    with cols_kpi[0]:
        tarjeta_kpi("Capital activo total", fmt(resumen["capital_total"]), "Capital actualmente vivo", "normal")
    with cols_kpi[1]:
        if incluir_devengado:
            tarjeta_kpi(
                "Cobro estimado mes",
                fmt(resumen["cobro_total_mes"]),
                f"+ {fmt(total_devengado)} devengado no cobrado",
                "positivo",
            )
        else:
            tarjeta_kpi("Cobro estimado mes", fmt(resumen["cobro_total_mes"]), "Ingresos brutos esperados", "positivo")
    with cols_kpi[2]:
        tarjeta_kpi("Pago inversores mes", fmt(resumen["pago_total_mes"]), "Obligaciones estimadas", "riesgo")
    with cols_kpi[3]:
        if incluir_devengado:
            beneficio_con_devengado = resumen["cobro_total_mes"] + total_devengado - resumen["pago_total_mes"]
            estado = "positivo" if beneficio_con_devengado >= 0 else "negativo"
            tarjeta_kpi(
                "Beneficio estimado mes",
                fmt(beneficio_con_devengado),
                "Margen neto estimado (incluye devengado no cobrado)",
                estado,
            )
        else:
            estado = "positivo" if resumen["beneficio_total_mes"] >= 0 else "negativo"
            tarjeta_kpi("Beneficio estimado mes", fmt(resumen["beneficio_total_mes"]), "Margen neto estimado", estado)

    if incluir_devengado:
        if not detalle_devengo.empty:
            with st.expander("Ver detalle del interés devengado por nota", expanded=False):
                tabla_devengo = detalle_devengo.copy().sort_values("interes_devengado_no_cobrado", ascending=False)
                tabla_devengo["ultimo_pago_cobrado"] = pd.to_datetime(tabla_devengo["ultimo_pago_cobrado"]).dt.strftime("%d/%m/%Y")
                tabla_devengo["interes_nota_anual"] = tabla_devengo["interes_nota_anual"].apply(lambda x: f"{x*100:.2f}%")
                tabla_devengo["capital_invertido"] = tabla_devengo["capital_invertido"].apply(fmt)
                tabla_devengo["interes_devengado_no_cobrado"] = tabla_devengo["interes_devengado_no_cobrado"].apply(fmt)
                st.dataframe(tabla_devengo, use_container_width=True, hide_index=True)
                boton_descarga_excel(detalle_devengo, "interes_devengado_no_cobrado.xlsx")
        else:
            st.caption(f"No hay interés corrido de notas para {nombre_mes_es(mes_dashboard)} {anio_dashboard}.")

    # ── Cobros semanales: justo debajo de las tarjetas de ingresos/gastos/beneficio, visible
    # sin tener que desplegar nada más abajo. Siempre usa el dataframe completo (df_inv), sin
    # aplicar el filtro del checkbox "Incluir Chaparro Fernández" — ver nota más abajo.
    if vista_dashboard in ["General", "Notas"]:
        mostrar_cobros_semanales_dashboard(df_inv, df_cal, df_control, anio_dashboard, mes_dashboard)

    st.markdown("### Rentabilidad del mes")
    r1, r2, r3, r4 = st.columns(4)
    with r1:
        tarjeta_kpi("Rent. beneficio mensual", fmt_pct(resumen["rentabilidad_beneficio_mes"]), "Beneficio / capital activo", "positivo" if resumen["rentabilidad_beneficio_mes"] >= 0 else "negativo")
    with r2:
        tarjeta_kpi("Rent. beneficio anualizada", fmt_pct(resumen["rentabilidad_beneficio_anualizada"]), "Mensual x 12", "positivo" if resumen["rentabilidad_beneficio_anualizada"] >= 0 else "negativo")
    with r3:
        tarjeta_kpi("% pagado inversores mes", fmt_pct(resumen["rentabilidad_pagada_inversor_mes"]), "Pago inversores / capital", "riesgo")
    with r4:
        tarjeta_kpi("% pagado inversores anual", fmt_pct(resumen["rentabilidad_pagada_inversor_anualizada"]), "Coste anualizado del capital", "riesgo")

    if vista_dashboard == "General":
        tarjeta_bitcoin_etf()

        st.markdown("---")
        with st.expander("🧭 Concentración de la cartera (emisor y sector)", expanded=False):
            _tab_concentracion_cartera(df_inv, df_control)

    if vista_dashboard in ["General", "Notas"]:
        fecha_analisis_notas = pd.Timestamp(anio_dashboard, mes_dashboard, ultimo_dia_mes(anio_dashboard, mes_dashboard)).normalize()
        mostrar_desglose_notas_por_tipo_inversor(
            df_inv_calculo,
            resumen.get("detalle_notas", pd.DataFrame()),
            fecha_analisis_notas,
            anio_dashboard,
            mes_dashboard,
        )

    mostrar_rentabilidad_por_activo_dashboard(resumen.get("rentabilidad_por_activo", pd.DataFrame()))

    # ── Desplegable movimientos MotoClick ─────────────────────────────────
    if vista_dashboard == "MotoClick":
        df_mov_display = cargar_movimientos_motoclick()
        with st.expander("📋 Devoluciones y reinversiones MotoClick", expanded=False):
            if df_mov_display.empty:
                st.info("No hay movimientos registrados en la hoja MOVIMIENTOS_MOTOCLICK.")
            else:
                # Filtrar solo el mes seleccionado si hay datos, pero mostrar todos con columna mes
                df_mov_display = df_mov_display.copy()
                df_mov_display["mes"] = df_mov_display["fecha"].dt.strftime("%m/%Y")
                df_mov_mes = df_mov_display[
                    (df_mov_display["fecha"].dt.year == anio_dashboard) &
                    (df_mov_display["fecha"].dt.month == mes_dashboard)
                ]
                total_dev = df_mov_mes.loc[df_mov_mes["tipo"] == "DEVOLUCION", "importe"].sum()
                total_reinv = df_mov_mes.loc[df_mov_mes["tipo"] == "REINVERSION", "importe"].sum()
                ajuste_neto = total_reinv - total_dev

                col_d, col_r, col_n = st.columns(3)
                col_d.metric("Devoluciones del mes", f"${total_dev:,.0f}", delta=None)
                col_r.metric("Reinversiones del mes", f"${total_reinv:,.0f}", delta=None)
                col_n.metric("Ajuste neto", f"${ajuste_neto:,.0f}",
                             delta="capital reducido" if ajuste_neto < 0 else "capital añadido")

                st.caption(f"Todos los movimientos registrados (mes seleccionado: {mes_dashboard:02d}/{anio_dashboard})")
                tabla_movs = df_mov_display[["fecha","tipo","importe","descripcion","mes"]].copy()
                tabla_movs["fecha"] = tabla_movs["fecha"].dt.strftime("%d/%m/%Y")
                tabla_movs["importe"] = tabla_movs["importe"].apply(lambda x: f"${x:,.0f}")
                # Colorear el mes seleccionado
                def highlight_mes(row):
                    if row["mes"] == f"{mes_dashboard:02d}/{anio_dashboard}":
                        return ["background-color: #fff9e6"] * len(row)
                    return [""] * len(row)
                st.dataframe(tabla_movs.style.apply(highlight_mes, axis=1),
                             use_container_width=True, hide_index=True)
                boton_descarga_excel(
                    df_mov_display[["fecha","tipo","importe","descripcion"]],
                    f"movimientos_motoclick_{anio_dashboard}_{mes_dashboard:02d}.xlsx",
                    "⬇️ Descargar movimientos"
                )

    st.markdown("---")

    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "Capital por activo",
        "Capital por inversor",
        "Beneficio mensual",
        "Rentabilidad por activo",
        "Rentabilidad por inversión",
    ])
    with tab1:
        grafico_capital_por_activo(resumen["activas"])
    with tab2:
        grafico_capital_por_inversor(resumen["activas"])
    with tab3:
        grafico_beneficio_mensual(df_inv_calculo, df_cal, df_control, prorratear_notas=prorratear_notas)
    with tab4:
        st.caption("Resumen del periodo seleccionado por tipo de activo. La rentabilidad anualizada es la rentabilidad mensual multiplicada por 12.")
        tabla_activo = resumen.get("rentabilidad_por_activo", pd.DataFrame())
        if tabla_activo is None or tabla_activo.empty:
            st.info("No hay datos de rentabilidad por activo para este mes.")
        else:
            st.dataframe(preparar_tabla_rentabilidad(tabla_activo), use_container_width=True)
            boton_descarga_excel(tabla_activo, f"rentabilidad_activo_{anio_dashboard}_{mes_dashboard:02d}.xlsx")
    with tab5:
        st.caption("Detalle inversión por inversión: cuánto genera, cuánto se paga al inversor y qué rentabilidad de beneficio deja.")
        tabla_inv = resumen.get("rentabilidad_inversiones", pd.DataFrame())
        if tabla_inv is None or tabla_inv.empty:
            st.info("No hay datos de rentabilidad por inversión para este mes.")
        else:
            columnas = [
                "activo", "nombre_activo", "id_inversion", "inversor", "capital",
                "cobro_compania_mes", "pago_inversor_mes", "beneficio_empresa_mes",
                "rentabilidad_beneficio_mes", "rentabilidad_beneficio_anualizada",
                "rentabilidad_pagada_inversor_mes", "rentabilidad_pagada_inversor_anualizada",
                "resultado_observacion",
            ]
            columnas = [c for c in columnas if c in tabla_inv.columns]
            st.dataframe(preparar_tabla_rentabilidad(tabla_inv[columnas]), use_container_width=True)
            boton_descarga_excel(tabla_inv[columnas], f"rentabilidad_inversiones_{anio_dashboard}_{mes_dashboard:02d}.xlsx")


def centro_control_inversiones():
    df_inv, _, _ = cargar_excel_completo()
    st.markdown("## Centro de control de inversiones")
    st.caption("Consulta profesional por filtros.")
    hoy = pd.Timestamp.today().normalize()
    c1, c2, c3, c4, c5, c6 = st.columns(6)
    activo = c1.selectbox("Activo", ["Todos", "notas", "paraguay", "bolivia", "motoclick", "futbol", "bitcoin", "otros"])
    inversores = ["Todos"] + sorted([x for x in df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).unique() if x.strip()])
    inversor = c2.selectbox("Inversor", inversores)
    anio_cc = int(c3.number_input("Año", min_value=2020, max_value=2100, value=hoy.year, key="cc_anio"))
    mes_cc = int(c4.number_input("Mes", min_value=1, max_value=12, value=hoy.month, key="cc_mes"))
    incluir_chaparro = c5.checkbox("Incluir Chaparro Fernández", value=False, key="cc_incluir_chaparro")
    fecha = pd.Timestamp(anio_cc, mes_cc, ultimo_dia_mes(anio_cc, mes_cc)).normalize()
    c6.markdown(f"<br><small>📅 Cierre: **{fecha.strftime('%d/%m/%Y')}**</small>", unsafe_allow_html=True)
    df_inv_filtrado = aplicar_filtro_chaparro_fernandez(df_inv, incluir_chaparro)
    activas = inversiones_activas_global(df_inv_filtrado, fecha)
    if not activas.empty:
        activas["activo"] = activas.apply(detectar_activo, axis=1)
    if activo != "Todos":
        activas = activas[activas["activo"] == activo]
    if inversor != "Todos":
        activas = activas[activas["inversor"].astype(str).str.lower() == inversor.lower()]
    capital = activas["capital_invertido"].sum() if not activas.empty else 0
    num_inversiones = len(activas)
    ticket_medio = capital / num_inversiones if num_inversiones else 0
    c1, c2, c3 = st.columns(3)
    c1.metric("Capital activo", fmt(capital))
    c2.metric("Inversiones activas", num_inversiones)
    c3.metric("Ticket medio", fmt(ticket_medio))
    if activas.empty:
        st.info("No hay inversiones activas con estos filtros.")
        return
    tab1, tab2, tab3 = st.tabs(["Detalle", "Por activo", "Por inversor"])
    with tab1:
        st.dataframe(preparar_tabla_monetaria(activas, ["capital_invertido", "interes_inversor_anual", "interes_nota_anual"]), use_container_width=True)
        boton_descarga_excel(activas, f"centro_control_{anio_cc}_{mes_cc:02d}.xlsx")
    with tab2:
        resumen_activo = activas.groupby("activo", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False)
        st.dataframe(preparar_tabla_monetaria(resumen_activo, ["capital"]), use_container_width=True)
        boton_descarga_excel(resumen_activo, f"capital_por_activo_{anio_cc}_{mes_cc:02d}.xlsx")
    with tab3:
        resumen_inv = activas.groupby("inversor", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False)
        st.dataframe(preparar_tabla_monetaria(resumen_inv, ["capital"]), use_container_width=True)
        boton_descarga_excel(resumen_inv, f"capital_por_inversor_{anio_cc}_{mes_cc:02d}.xlsx")


# =========================
# SECCIONES ORIGINALES MEJORADAS
# =========================
def seccion_activo(nombre_visible: str, activo_key: str, tasa_anual: float, incluir_ingresado_desde_inicio: bool = False):
    df_inv, _, _ = cargar_excel_completo()
    st.header(f"📌 Consultas {nombre_visible}")
    opciones = [
        f"¿Cuánto ingresará {nombre_visible} en un mes?",
        "¿Cuánto cobrará cada inversor ese mes?",
        "¿Cuánto cobrará un inversor concreto ese mes?",
        "¿Cuál será el beneficio de la empresa ese mes?",
        "¿Cuál es el total pagado a inversores desde el inicio?",
        "¿Cuánto ha ingresado la compañía desde el inicio?",
        "¿Cuál es el beneficio total acumulado desde el inicio?",
        f"¿Cuánto capital hay actualmente activo en {nombre_visible} hoy?",
        f"¿Cuánto capital había activo en {nombre_visible} en un mes concreto?",
    ]
    consulta = st.selectbox("Elige una pregunta", opciones)
    necesita_mes = consulta in opciones[:4] or consulta == f"¿Cuánto capital había activo en {nombre_visible} en un mes concreto?"
    anio = mes = None
    if necesita_mes:
        c1, c2 = st.columns(2)
        anio = int(c1.number_input("Año", 2020, 2100, pd.Timestamp.today().year, key=f"{activo_key}_anio"))
        mes = int(c2.number_input("Mes", 1, 12, pd.Timestamp.today().month, key=f"{activo_key}_mes"))
    nombre_inversor = None
    if consulta == "¿Cuánto cobrará un inversor concreto ese mes?":
        inversores = sorted([x for x in df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).unique() if x.strip()])
        nombre_inversor = st.selectbox("Inversor", inversores) if inversores else st.text_input("Inversor")
    def _det_mes(a, m):
        """Devuelve detalle del activo con ingreso MotoClick ajustado si corresponde."""
        det = detalle_activo_mes(df_inv, activo_key, tasa_anual, a, m)
        if activo_key == "motoclick" and not det.empty:
            det["activo"] = "motoclick"
            det = ajustar_ingreso_motoclick(det, df_inv, a, m)
        return det

    if st.button("Calcular", key=f"calc_{activo_key}_{consulta}"):
        if consulta == f"¿Cuánto ingresará {nombre_visible} en un mes?":
            detalle = _det_mes(anio, mes)
            mostrar_metricas(f"Resultado {nombre_mes_es(mes)} {anio}", [("Ingreso bruto", fmt(detalle["ingreso_bruto"].sum() if not detalle.empty else 0))])
            if not detalle.empty:
                st.dataframe(preparar_tabla_monetaria(detalle, ["capital_invertido", "ingreso_bruto", "pago_inversor_mes", "beneficio_empresa_mes"]), use_container_width=True)
                boton_descarga_excel(detalle, f"ingreso_{activo_key}_{anio}_{mes:02d}.xlsx")
        elif consulta == "¿Cuánto cobrará cada inversor ese mes?":
            detalle = _det_mes(anio, mes)
            if detalle.empty:
                st.info("No hay cobros de inversores para ese mes.")
            else:
                resumen = detalle.groupby("inversor", as_index=False)["pago_inversor_mes"].sum().rename(columns={"pago_inversor_mes": "cobro_mes"}).sort_values("cobro_mes", ascending=False)
                st.dataframe(preparar_tabla_monetaria(resumen, ["cobro_mes"]), use_container_width=True)
        elif consulta == "¿Cuánto cobrará un inversor concreto ese mes?":
            detalle = _det_mes(anio, mes)
            filtrado = detalle[detalle["inversor"].astype(str).str.lower() == str(nombre_inversor).strip().lower()] if not detalle.empty else pd.DataFrame()
            mostrar_metricas("Resultado", [(f"Cobro de {nombre_inversor}", fmt(filtrado["pago_inversor_mes"].sum() if not filtrado.empty else 0))])
        elif consulta == "¿Cuál será el beneficio de la empresa ese mes?":
            detalle = _det_mes(anio, mes)
            mostrar_metricas(f"Resultado {nombre_mes_es(mes)} {anio}", [("Beneficio empresa", fmt(detalle["beneficio_empresa_mes"].sum() if not detalle.empty else 0))])
            if not detalle.empty:
                st.dataframe(preparar_tabla_monetaria(detalle, ["capital_invertido", "ingreso_bruto", "pago_inversor_mes", "beneficio_empresa_mes"]), use_container_width=True)
                boton_descarga_excel(detalle, f"beneficio_{activo_key}_{anio}_{mes:02d}.xlsx")
        elif consulta == "¿Cuál es el total pagado a inversores desde el inicio?":
            mostrar_metricas("Resultado", [("Total pagado", fmt(total_pagado_activo_desde_inicio(df_inv, activo_key, tasa_anual)))])
        elif consulta == "¿Cuánto ha ingresado la compañía desde el inicio?":
            mostrar_metricas("Resultado", [("Total ingresado", fmt(total_ingresado_activo_desde_inicio(df_inv, activo_key, tasa_anual)))])
        elif consulta == "¿Cuál es el beneficio total acumulado desde el inicio?":
            ingreso = total_ingresado_activo_desde_inicio(df_inv, activo_key, tasa_anual)
            pagado  = total_pagado_activo_desde_inicio(df_inv, activo_key, tasa_anual)
            mostrar_metricas("Resultado", [("Beneficio acumulado", fmt(ingreso - pagado))])
        elif consulta == f"¿Cuánto capital hay actualmente activo en {nombre_visible} hoy?":
            bruto = capital_activo_en_fecha(df_inv, pd.Timestamp.today(), activo_key, False)
            real = capital_activo_en_fecha(df_inv, pd.Timestamp.today(), activo_key, True)
            mostrar_metricas("Resultado", [("Capital activo", fmt(bruto)), ("Capital activo real", fmt(real))])
        elif consulta == f"¿Cuánto capital había activo en {nombre_visible} en un mes concreto?":
            fecha = pd.Timestamp(anio, mes, ultimo_dia_mes(anio, mes))
            bruto = capital_activo_en_fecha(df_inv, fecha, activo_key, False)
            real = capital_activo_en_fecha(df_inv, fecha, activo_key, True)
            mostrar_metricas(f"Cierre {nombre_mes_es(mes)} {anio}", [("Capital activo", fmt(bruto)), ("Capital activo real", fmt(real))])
        elif consulta == "¿Cuánto ha ingresado la compañía desde el inicio?":
            mostrar_metricas("Resultado", [("Total ingresado", fmt(total_ingresado_activo_desde_inicio(df_inv, activo_key, tasa_anual)))])
        elif consulta == "¿Cuál es el beneficio total acumulado desde el inicio?":
            ingreso = total_ingresado_activo_desde_inicio(df_inv, activo_key, tasa_anual)
            pagado = total_pagado_activo_desde_inicio(df_inv, activo_key, tasa_anual)
            mostrar_metricas("Resultado", [("Beneficio acumulado", fmt(ingreso - pagado))])


def seccion_notas():
    df_inv, df_cal, df_control = cargar_excel_completo()
    st.header("🧾 Consultas Notas")
    consulta = st.selectbox("Elige una pregunta", [
        "¿Cuánto cobrará la compañía en un mes de notas?", "¿Cuánto se pagará a inversores en un mes de notas?", "¿Cuál será el beneficio de la empresa en un mes de notas?",
        "¿Cuánto cobrará cada inversor ese mes?", "¿Cuánto cobrará un inversor concreto ese mes?", "¿Cuánto ha cobrado la compañía desde el inicio?",
        "¿Cuánto se ha pagado a inversores desde el inicio?", "¿Cuál es el beneficio total desde el inicio?", "¿Cuál es el próximo pago de una nota?",
        "¿Cuál es la próxima observación de una nota?", "¿Cuánto capital hay invertido en total?", "¿Cuánto capital hay actualmente activo?",
        "¿Cuánto capital tiene un inversor?", "¿Cuánto capital activo tiene un inversor?", "Ver ranking de capital por inversor", "Ver ranking de capital activo",
    ])
    consultas_mes = ["¿Cuánto cobrará la compañía en un mes de notas?", "¿Cuánto se pagará a inversores en un mes de notas?", "¿Cuál será el beneficio de la empresa en un mes de notas?", "¿Cuánto cobrará cada inversor ese mes?", "¿Cuánto cobrará un inversor concreto ese mes?"]
    anio = mes = None
    if consulta in consultas_mes:
        c1, c2 = st.columns(2)
        anio = int(c1.number_input("Año", 2020, 2100, pd.Timestamp.today().year, key="notas_anio"))
        mes = int(c2.number_input("Mes", 1, 12, pd.Timestamp.today().month, key="notas_mes"))
    inversores = sorted([x for x in df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).unique() if x.strip()])
    nombre_inversor = None
    if consulta in ["¿Cuánto cobrará un inversor concreto ese mes?", "¿Cuánto capital tiene un inversor?", "¿Cuánto capital activo tiene un inversor?"]:
        nombre_inversor = st.selectbox("Inversor", inversores) if inversores else st.text_input("Inversor")
    nota = None
    if consulta in ["¿Cuál es el próximo pago de una nota?", "¿Cuál es la próxima observación de una nota?"]:
        notas_disponibles = sorted([int(x) for x in df_cal.get("nota", pd.Series(dtype="Int64")).dropna().unique()])
        nota = st.selectbox("Número de nota", notas_disponibles) if notas_disponibles else int(st.number_input("Número de nota", min_value=1, value=1))
    if st.button("Calcular", key=f"calc_notas_{consulta}"):
        if consulta in consultas_mes:
            total_cobrado, total_pagado, total_beneficio, detalle, pagos = resumen_notas_mes(df_inv, df_cal, df_control, anio, mes)
            if consulta == "¿Cuánto cobrará la compañía en un mes de notas?":
                mostrar_metricas(f"Resultado {nombre_mes_es(mes)} {anio}", [("Cobra compañía", fmt(total_cobrado))])
                resumen_cuentas = resumen_por_cuenta_cobro(detalle)
                if not resumen_cuentas.empty:
                    st.dataframe(preparar_tabla_monetaria(resumen_cuentas, ["cobro_compania"]), use_container_width=True)
                    boton_descarga_excel(resumen_cuentas, f"cobros_por_cuenta_{anio}_{mes:02d}.xlsx")
            elif consulta == "¿Cuánto se pagará a inversores en un mes de notas?":
                mostrar_metricas(f"Resultado {nombre_mes_es(mes)} {anio}", [("Pago inversores", fmt(total_pagado))])
            elif consulta == "¿Cuál será el beneficio de la empresa en un mes de notas?":
                mostrar_metricas(f"Resultado {nombre_mes_es(mes)} {anio}", [("Beneficio empresa", fmt(total_beneficio))])
            elif consulta == "¿Cuánto cobrará cada inversor ese mes?":
                resumen = detalle.groupby("inversor", as_index=False)["pago_inversor"].sum().rename(columns={"pago_inversor": "cobro_mes"}).sort_values("cobro_mes", ascending=False) if not detalle.empty else pd.DataFrame()
                st.dataframe(preparar_tabla_monetaria(resumen, ["cobro_mes"]), use_container_width=True) if not resumen.empty else st.info("No hay cobros de inversores para ese mes.")
            elif consulta == "¿Cuánto cobrará un inversor concreto ese mes?":
                filtrado = detalle[detalle["inversor"].astype(str).str.lower() == str(nombre_inversor).strip().lower()] if not detalle.empty else pd.DataFrame()
                mostrar_metricas("Resultado", [(f"Cobro de {nombre_inversor}", fmt(filtrado["pago_inversor"].sum() if not filtrado.empty else 0))])
            if not pagos.empty:
                with st.expander("Ver pagos detectados"):
                    st.dataframe(preparar_tabla_monetaria(pagos, []), use_container_width=True)
            if not detalle.empty:
                with st.expander("Ver detalle por nota e inversión"):
                    st.dataframe(preparar_tabla_monetaria(detalle, ["capital_invertido", "cobro_compania", "pago_inversor", "beneficio_empresa"]), use_container_width=True)
                    boton_descarga_excel(detalle, f"detalle_notas_mes_{anio}_{mes:02d}.xlsx")
        elif consulta == "¿Cuánto ha cobrado la compañía desde el inicio?":
            detalle = preparar_detalle_notas(df_inv, pagos_notas_hasta_hoy(df_cal), df_cal=df_cal, df_control=df_control)
            mostrar_metricas("Resultado", [("Total cobrado compañía", fmt(detalle["cobro_compania"].sum() if not detalle.empty else 0))])
        elif consulta == "¿Cuánto se ha pagado a inversores desde el inicio?":
            detalle = preparar_detalle_notas(df_inv, pagos_notas_hasta_hoy(df_cal), df_cal=df_cal, df_control=df_control)
            mostrar_metricas("Resultado", [("Total pagado inversores", fmt(detalle["pago_inversor"].sum() if not detalle.empty else 0))])
        elif consulta == "¿Cuál es el beneficio total desde el inicio?":
            detalle = preparar_detalle_notas(df_inv, pagos_notas_hasta_hoy(df_cal), df_cal=df_cal, df_control=df_control)
            mostrar_metricas("Resultado", [("Beneficio total", fmt(detalle["beneficio_empresa"].sum() if not detalle.empty else 0))])
        elif consulta == "¿Cuál es el próximo pago de una nota?":
            fecha = proximo_evento_nota(df_cal, int(nota), "PAGO")
            st.success(f"El próximo pago de la nota {nota} es el {pd.Timestamp(fecha).strftime('%d/%m/%Y')}") if fecha is not None else st.info("No hay pagos futuros para esa nota.")
        elif consulta == "¿Cuál es la próxima observación de una nota?":
            fecha = proximo_evento_nota(df_cal, int(nota), "OBSERVACION")
            st.success(f"La próxima observación de la nota {nota} es el {pd.Timestamp(fecha).strftime('%d/%m/%Y')}") if fecha is not None else st.info("No hay observaciones futuras para esa nota.")
        elif consulta == "¿Cuánto capital hay invertido en total?":
            mostrar_metricas("Resultado", [("Capital total invertido", fmt(filtrar_notas(df_inv)["capital_invertido"].sum()))])
        elif consulta == "¿Cuánto capital hay actualmente activo?":
            trabajo = filtrar_notas(df_inv); hoy = pd.Timestamp.today().normalize()
            activas = trabajo[(trabajo["fecha_inversion"].notna()) & (trabajo["fecha_inversion"] <= hoy) & (trabajo["fecha_final_inversion"].isna() | (trabajo["fecha_final_inversion"] >= hoy))]
            mostrar_metricas("Resultado", [("Capital activo hoy", fmt(activas["capital_invertido"].sum() if not activas.empty else 0))])
        elif consulta == "¿Cuánto capital tiene un inversor?":
            trabajo = filtrar_notas(df_inv); filtrado = trabajo[trabajo["inversor"].astype(str).str.lower() == str(nombre_inversor).strip().lower()]
            mostrar_metricas("Resultado", [(f"Capital total de {nombre_inversor}", fmt(filtrado["capital_invertido"].sum() if not filtrado.empty else 0))])
        elif consulta == "¿Cuánto capital activo tiene un inversor?":
            trabajo = filtrar_notas(df_inv); hoy = pd.Timestamp.today().normalize()
            filtrado = trabajo[(trabajo["inversor"].astype(str).str.lower() == str(nombre_inversor).strip().lower()) & (trabajo["fecha_inversion"].notna()) & (trabajo["fecha_inversion"] <= hoy) & (trabajo["fecha_final_inversion"].isna() | (trabajo["fecha_final_inversion"] >= hoy))]
            mostrar_metricas("Resultado", [(f"Capital activo de {nombre_inversor}", fmt(filtrado["capital_invertido"].sum() if not filtrado.empty else 0))])
        elif consulta == "Ver ranking de capital por inversor":
            st.dataframe(preparar_tabla_monetaria(resumen_capital_por_inversor_notas(df_inv, False), ["capital"]), use_container_width=True)
        elif consulta == "Ver ranking de capital activo":
            st.dataframe(preparar_tabla_monetaria(resumen_capital_por_inversor_notas(df_inv, True), ["capital"]), use_container_width=True)


def _cargar_pdfs_notas() -> dict:
    """Lee todos los PDFs de la carpeta notas_pdfs/ y los devuelve en base64."""
    import os, base64, tempfile, re as _re2
    pdfs = {}
    # Buscar en carpeta temporal (subidos desde la app) y en carpeta del repo (GitHub)
    carpetas = [
        os.path.join(tempfile.gettempdir(), "notas_pdfs_cf"),
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "notas_pdfs"),
        os.path.dirname(os.path.abspath(__file__)),  # raíz del repo como fallback
    ]
    for carpeta in carpetas:
        if not os.path.exists(carpeta) or not os.path.isdir(carpeta):
            continue
        for fname in sorted(os.listdir(carpeta)):
            if fname.lower().endswith(".pdf"):
                # Normalizar nombre: "NOTA 10.pdf" -> "nota10"
                clave = _re2.sub(r"[^a-z0-9]", "", fname.lower().replace(".pdf",""))
                if clave not in {_re2.sub(r"[^a-z0-9]","",k.lower().replace(".pdf","")) for k in pdfs}:
                    try:
                        with open(os.path.join(carpeta, fname), "rb") as f:
                            pdfs[fname] = base64.standard_b64encode(f.read()).decode()
                    except Exception:
                        pass
    return pdfs


def _tab_asistente_ia_notas(df_inv, df_cal, df_control):
    """Pestaña: chat IA especializado en los PDFs de las notas estructuradas."""
    st.caption("Pregúntame cualquier cosa sobre las notas: cobros, calls, barreras, vencimientos...")

    import os, tempfile
    carpeta_pdfs = os.path.join(tempfile.gettempdir(), "notas_pdfs_cf")
    os.makedirs(carpeta_pdfs, exist_ok=True)
    pdfs_existentes = sorted([f for f in os.listdir(carpeta_pdfs) if f.lower().endswith(".pdf")])

    with st.expander(f"📂 PDFs cargados ({len(pdfs_existentes)})", expanded=not pdfs_existentes):
        if pdfs_existentes:
            col_l, col_d = st.columns([3, 1])
            with col_l:
                for pdf in pdfs_existentes:
                    st.markdown(f"✅ `{pdf}`")
            with col_d:
                pdf_borrar = st.selectbox("Borrar", ["—"] + pdfs_existentes, key="pdf_borrar_sel")
                if pdf_borrar != "—" and st.button("🗑️ Borrar", key="btn_borrar_pdf"):
                    os.remove(os.path.join(carpeta_pdfs, pdf_borrar))
                    st.rerun()
        else:
            st.info("Sube los PDFs de las notas para activar el asistente.")

        subidos = st.file_uploader(
            "Subir PDFs (nota1.pdf, nota2.pdf...)",
            type=["pdf"], accept_multiple_files=True, key="uploader_pdfs_notas"
        )
        if subidos:
            for f in subidos:
                with open(os.path.join(carpeta_pdfs, f.name), "wb") as out:
                    out.write(f.read())
            st.success(f"✅ {len(subidos)} PDF(s) guardado(s).")
            st.rerun()

    if "chat_notas_ia" not in st.session_state:
        st.session_state["chat_notas_ia"] = []

    for msg in st.session_state["chat_notas_ia"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if not st.session_state["chat_notas_ia"]:
        st.markdown("**Preguntas frecuentes:**")
        sugs = ["¿Cuándo es el próximo cobro?", "¿Cuándo es el próximo call?",
                "¿Cuánto cobraremos este mes?", "¿Qué notas están cerca de la barrera?",
                "Dame un resumen de todas las notas activas"]
        cols = st.columns(len(sugs))
        for i, s in enumerate(sugs):
            if cols[i].button(s, key=f"sug_nota_{i}", use_container_width=True):
                st.session_state["chat_notas_ia"].append({"role": "user", "content": s})
                st.rerun()

    pregunta = st.chat_input("Escribe tu pregunta sobre las notas...", key="chat_input_notas")
    if pregunta:
        st.session_state["chat_notas_ia"].append({"role": "user", "content": pregunta})
        st.rerun()

    if st.session_state["chat_notas_ia"] and st.session_state["chat_notas_ia"][-1]["role"] == "user":
        ultima = st.session_state["chat_notas_ia"][-1]["content"]
        with st.chat_message("assistant"):
            with st.spinner("Analizando documentos..."):
                try:
                    import requests as _req
                    pdfs_b64 = _cargar_pdfs_notas()
                    hoy = pd.Timestamp.today().normalize()
                    ctx_lines = [f"Fecha de hoy: {hoy.strftime('%d/%m/%Y')}"]
                    # Solo próximos 15 eventos para minimizar tokens
                    if df_cal is not None and not df_cal.empty:
                        df_c2 = df_cal.copy()
                        df_c2["fecha"] = pd.to_datetime(df_c2["fecha"], errors="coerce")
                        prox = df_c2[df_c2["fecha"] >= hoy].sort_values("fecha").head(15)
                        cols_cal = [c for c in ["fecha","nota","tipo_evento","importe_cobro","importe_pago_inversor"] if c in prox.columns]
                        ctx_lines += ["", "=== PRÓXIMOS EVENTOS ===", prox[cols_cal].to_string(index=False)]
                    # Control notas: incluye emisor y cupón anual — antes faltaban aquí, por lo que
                    # el asistente no tenía forma de saber el cupón real de una nota confirmada y
                    # podía inventarlo o sacarlo de un PDF equivocado.
                    if df_control is not None and not df_control.empty:
                        cols_ctrl = [c for c in ["nota","ticker","emisor","precio_compra","contingency","barrera_capital","fecha_vencimiento","proximo_call"] if c in df_control.columns]
                        if cols_ctrl:
                            ctx_lines += ["", "=== NOTAS (resumen, fuente: CONTROL_NOTAS) ===", df_control[cols_ctrl].to_string(index=False)]
                    # Cupón anual por nota — fuente: INVERSIONES (interes_nota_anual), la misma
                    # columna que usa el resto de la app. Es la ÚNICA fuente válida de cupón.
                    if df_inv is not None and not df_inv.empty and "nombre_activo" in df_inv.columns and "interes_nota_anual" in df_inv.columns:
                        try:
                            df_notas_cupon = df_inv[df_inv["nombre_activo"].astype(str).str.upper().str.startswith("NOTA")].copy()
                            cupon_por_nota = df_notas_cupon.groupby("nombre_activo")["interes_nota_anual"].first().dropna()
                            if not cupon_por_nota.empty:
                                lineas_cupon = [f"  {n}: {v*100:.2f}% anual" for n, v in cupon_por_nota.items()]
                                ctx_lines += ["", "=== CUPÓN ANUAL POR NOTA (fuente: INVERSIONES, USA SIEMPRE ESTE DATO) ==="] + lineas_cupon
                        except Exception:
                            pass
                    contexto = "\n".join(ctx_lines)

                    # Filtrar PDFs: SOLO si la pregunta menciona un número de nota concreto que
                    # coincide con un PDF cargado. Antes, si no había número explícito pero la
                    # pregunta tenía palabras como "próximo"/"cobro"/"call", se mandaban los primeros
                    # 3 PDFs de la carpeta al azar — podían ser de notas totalmente distintas a la
                    # pregunta, y el modelo terminaba describiendo cupones/fechas de la nota
                    # equivocada. Ahora, sin número explícito, no se manda ningún PDF: se responde
                    # solo con los datos ya verificados del Excel (arriba), que son la fuente fiable.
                    import re as _re
                    nums_mencionados = _re.findall(r"nota[_\s]*(\d+)", ultima.lower())
                    pdfs_filtrados = {}
                    if nums_mencionados:
                        # Normalizar nombre del PDF para comparar: "NOTA 10.pdf" -> "nota10"
                        def _norm(s): return _re.sub(r"[^a-z0-9]","",s.lower().replace(".pdf",""))
                        pdfs_filtrados = {k: v for k, v in pdfs_b64.items()
                                          if any(_norm(k) == f"nota{n}" for n in nums_mencionados)}
                        if not pdfs_filtrados:
                            contexto += (
                                f"\n\n[AVISO: se menciona la nota {', '.join(nums_mencionados)} pero no hay "
                                f"un PDF cargado con ese nombre — usa solo los datos del Excel de arriba, y si "
                                f"falta algo que solo está en el documento oficial, dilo explícitamente en vez "
                                f"de adivinarlo con otro PDF.]"
                            )

                    contenido = []
                    for nombre_pdf, pdf_b64 in pdfs_filtrados.items():
                        contenido.append({"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64}, "title": nombre_pdf.replace(".pdf","").upper()})
                    contenido.append({"type": "text", "text": f"DATOS DEL EXCEL:\n\n{contexto}\n\n---\nPREGUNTA: {ultima}"})

                    historial = []
                    for m in st.session_state["chat_notas_ia"][:-1][-6:]:
                        historial.append({"role": m["role"], "content": m["content"]})
                    historial.append({"role": "user", "content": contenido})

                    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
                    resp = _req.post("https://api.anthropic.com/v1/messages",
                        headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
                        json={"model": "claude-sonnet-4-5", "max_tokens": 1000,
                              "system": 'Eres el asistente financiero de Chaparro Fernández Wealth Management, un fondo de inversión privado. Respondes preguntas de socios e inversores con total precisión sobre el estado del fondo.\n\n== ESTRUCTURA DEL NEGOCIO ==\nEl fondo capta capital de inversores, lo invierte en activos y paga a cada inversor un interés fijo anual. El beneficio es la diferencia entre lo que rinden los activos y lo que se paga a inversores.\n\n== ACTIVOS Y TASAS ==\nParaguay: 15% | Bolivia: 15% | MotoClick: 25% | Fútbol: 15% | Bitcoin: 20% | Notas: tasa variable por nota\n\n== INVERSORES Y TASAS ==\nLEO: 10% | JORDI CHAPARRO: 15% | YURI FERNANDEZ: 15%\nROBERTO VISCAFE: 5% hasta 31/01/2026, 7.5% desde 01/02/2026 hasta 30/06/2026, 10% desde 01/07/2026\nCROWE BOLIVIA: 5% hasta 31/01/2026, 7.5% desde 01/02/2026 hasta 30/06/2026, 10% desde 01/07/2026\nJR REAL ESTATE: 5% hasta 31/01/2026, 7.5% desde 01/02/2026 hasta 30/06/2026, 10% desde 01/07/2026\n2012 JACC GROUP: 10% | PEDRO MAGAÑA: 10% | PAM: 10%\nCHAPARRO FERNANDEZ: 0% — sociedad gestora, no recibe pago como inversor\nGOLDEN BRICKS: 10% | TERESA: 10% | JEP: 15%\nJORDI ESPECIAL: 10% | EVA CHAPARRO: 15% | PAOLA CHAPARRO: 15% | JAPAN JORDI: 15%\n\n== REGLAS DE CÁLCULO ==\n1. Pago inversor = capital x tasa_inversor / 12 x pro-rata días del mes\n2. El pago es MENSUAL y FIJO independiente del calendario de cobros de notas\n3. Para notas: el cobro de la empresa sigue CALENDARIO_NOTAS; el pago al inversor es siempre mensual\n4. Reinversiones cuentan para cobro empresa Y pago inversor\n5. CHAPARRO FERNANDEZ: pago=0, todo cobro es beneficio de la empresa\n6. NOTA_10: pago trimestral\n\n== REGLA ANTI-INVENCIÓN (cupones, fechas, emisor) ==\nEl cupón/tasa anual y el emisor de cada nota SOLO son válidos si vienen de los bloques "CUPÓN ANUAL POR NOTA" o "NOTAS (resumen, fuente: CONTROL_NOTAS)" del contexto, o de un PDF que sí esté adjunto en este mensaje. Las fechas de observación/pago/call SOLO son válidas si vienen del bloque "PRÓXIMOS EVENTOS". Si te preguntan por un dato de una nota y no lo ves en ninguna de esas fuentes (por ejemplo, porque no se adjuntó el PDF de esa nota), dilo explícitamente ("no tengo ese dato para la Nota X en este momento — súbeme su PDF o revisa CONTROL_NOTAS") en vez de aproximarlo, deducirlo de otra nota parecida o inventarlo.\n\n== FORMATO ==\nResponde SIEMPRE en español. Sé conciso: da el dato pedido directamente. Si piden detalle, entonces desarrolla. Fechas DD/MM/YYYY, importes con $ y 2 decimales.',
                              "messages": historial}, timeout=60)
                    data = resp.json()
                    respuesta = "".join(b.get("text","") for b in data.get("content",[]) if b.get("type")=="text")
                    if not respuesta:
                        respuesta = f"Error API: {data.get('error',{}).get('message', str(data))}"
                except Exception as e:
                    respuesta = f"Error: {e}"
                st.markdown(_md_seguro(respuesta))
        st.session_state["chat_notas_ia"].append({"role": "assistant", "content": respuesta})

    if st.session_state["chat_notas_ia"]:
        if st.button("🗑️ Limpiar conversación", key="btn_limpiar_notas"):
            st.session_state["chat_notas_ia"] = []
            st.rerun()


def extraer_datos_nota_con_ia(pdf_bytes: bytes) -> dict:
    """
    Envía el PDF oficial de una nota estructurada a la API de Claude y le pide
    que devuelva ÚNICAMENTE JSON con los datos clave (tickers, precios iniciales,
    barreras, cupón, vencimiento y calendario observación/pago).
    Si algún dato no se puede determinar con confianza, la IA debe usar el string
    "REVISAR" en vez de inventarlo, para que el usuario lo revise en la previsualización.
    """
    import base64
    import json as _json

    # Validación previa: si el archivo no es un PDF real, o está corrupto/cifrado, la API de
    # Claude lo rechaza con un mensaje genérico ("The PDF specified was not valid") que no dice
    # POR QUÉ. Detectamos estos casos antes de gastar una llamada a la API, para dar un mensaje
    # que sí sirva para actuar.
    if pdf_bytes and not pdf_bytes.startswith(b"%PDF-"):
        # Caso conocido: el archivo trae pegado un envoltorio HTTP multipart/form-data alrededor
        # del PDF real (típico de algunas herramientas de descarga/proxy) — el PDF de verdad
        # empieza en "%PDF-" y termina en "%%EOF" en algún punto interior. Si lo encontramos,
        # lo recortamos automáticamente en vez de rechazar el archivo.
        _inicio = pdf_bytes.find(b"%PDF-")
        _fin = pdf_bytes.rfind(b"%%EOF")
        if _inicio != -1 and _fin != -1 and _fin > _inicio:
            pdf_bytes = pdf_bytes[_inicio:_fin + len(b"%%EOF")]

    if not pdf_bytes or not pdf_bytes.startswith(b"%PDF-"):
        return {"error": (
            "El archivo no parece ser un PDF válido (no empieza con la firma %PDF-). "
            "Puede que se haya subido corrupto, incompleto, o que en realidad sea otro tipo de "
            "archivo con extensión .pdf. Probá abrirlo en tu computadora para confirmar que se ve "
            "bien, y si hace falta, exportalo/guardalo de nuevo como PDF y volvé a subirlo."
        )}

    # La API de Claude rechaza cualquier PDF de más de 100 páginas (límite duro, no de tokens).
    # Los pricing supplements suelen traer anexos legales larguísimos — si pasa de 100 páginas,
    # extraemos el texto en Python y lo mandamos como texto plano en vez de como documento.
    usar_texto_plano = False
    texto_extraido = None
    try:
        from pypdf import PdfReader
        lector = PdfReader(BytesIO(pdf_bytes))
        if lector.is_encrypted:
            # Muchos pricing supplements vienen con cifrado de SOLO PERMISOS (sin contraseña
            # real para abrirlos). decrypt("") los desbloquea sin problema, pero la propiedad
            # is_encrypted de pypdf se queda en True para siempre aunque ya esté desbloqueado
            # -- por eso NO hay que volver a mirar is_encrypted, sino el resultado de decrypt().
            resultado_decrypt = 0
            try:
                resultado_decrypt = lector.decrypt("")
            except Exception:
                resultado_decrypt = 0
            if not resultado_decrypt:
                return {"error": (
                    "El PDF está protegido con una contraseña real (no solo permisos), y no se "
                    "puede leer así. Quitale la contraseña (por ejemplo abriéndolo y exportándolo "
                    "de nuevo sin seguridad) y volvé a subirlo."
                )}
        n_paginas = len(lector.pages)
        if n_paginas > 100:
            usar_texto_plano = True
            texto_extraido = "\n\n".join((pagina.extract_text() or "") for pagina in lector.pages)
    except Exception as _e_pdf:
        return {"error": (
            f"El PDF no se pudo abrir para procesarlo (archivo probablemente corrupto o mal generado): {_e_pdf}. "
            "Probá volver a exportar/descargar el documento original y subirlo de nuevo."
        )}

    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")

    system = """Extraes datos estructurados de documentos oficiales (pricing supplement) de notas estructuradas (structured notes).

Devuelve ÚNICAMENTE un JSON válido, sin texto antes ni después, sin backticks de markdown, con este esquema exacto:

{
  "emisor": "string, ej. BNP Paribas",
  "cusip": "string, el identificador CUSIP de 9 caracteres de la nota (letras y números, ej. 05619JGD6). Busca la etiqueta 'CUSIP' o 'CUSIP No.' en la portada o en la sección de términos generales del documento. Si el documento solo trae ISIN y no CUSIP, usa el ISIN aquí tal cual. Si no aparece ninguno de los dos, usa 'REVISAR'.",
  "cupon_anual_pct": number entre 0 y 1 (ej. 0.375 para 37.5%),
  "fecha_vencimiento": "YYYY-MM-DD",
  "fecha_inicio_nota": "YYYY-MM-DD, la Initial Valuation Date (o Pricing Date si no hay Initial Valuation Date separada) de la nota",
  "tiene_memoria": true o false,
  "tiene_one_star": true o false,
  "tickers": [
    {"ticker": "string en mayúsculas", "barrera_cupon_pct": number entre 0 y 1, "barrera_capital_pct": number entre 0 y 1, "call_level_pct": number, normalmente 1.0 si el call requiere estar al 100% o más del precio inicial}
  ],
  "calendario": [
    {"observacion": "YYYY-MM-DD", "pago": "YYYY-MM-DD"}
  ],
  "fechas_call": ["YYYY-MM-DD"]
}

REGLAS:
- "barrera_cupon_pct" es el umbral mínimo (como fracción del precio inicial) para que se pague el cupón condicional ese periodo (ej. 0.60 para "60% of Initial Price").
- "barrera_capital_pct" es el umbral para el capital al vencimiento (normalmente igual a barrera_cupon_pct, a veces distinto).
- "call_level_pct" es el umbral (como fracción del precio inicial) que activa la llamada anticipada (call) — normalmente 1.0 (100% del precio inicial).
- "calendario" debe incluir TODAS las fechas de Coupon/Review/Contingent Interest Observation Date y su correspondiente Coupon/Interest/Contingent Interest Payment Date que aparezcan en el documento, en el mismo orden. Hay TRES formatos posibles según el emisor, identifica cuál aplica:
  a) TABLA de dos columnas ya emparejadas (ej. BNP Paribas: "Coupon Valuation Date | Coupon Payment Date") — léela directamente.
  b) DOS LISTAS separadas por comas en un párrafo (ej. JPMorgan: "Review Dates: fecha1, fecha2..." y "Interest Payment Dates: fecha1, fecha2..." por separado) — empareja por posición (1ª con 1ª, 2ª con 2ª, etc.), ambas listas tienen el mismo número de fechas.
  c) REGLA DE REPETICIÓN en texto, sin fechas explícitas (ej. TD Bank: "Monthly, on the 2nd calendar day of each month, commencing on January 2, 2026 and ending on December 4, 2028"). En este caso DEBES GENERAR tú mismo la lista completa de fechas aplicando la regla exactamente (mismo día del mes, cada mes/trimestre según corresponda, desde la fecha de inicio hasta la de fin, ambas incluidas). Para la fecha de pago, si el texto dice algo como "the third Business Day following the relevant Observation Date", suma esos días naturales a la fecha de observación como aproximación razonable (no necesitas calcular festivos bancarios exactos) — EXCEPTO la última fecha de pago, que casi siempre es explícitamente la Maturity Date, usa esa fecha tal cual en vez de la aproximación.
- La terminología de barreras y call varía mucho entre emisores — no busques una frase exacta, busca el CONCEPTO: la "barrera de cupón" puede llamarse "Conditional Coupon Barrier Price", "Interest Barrier", "Trigger Value" o "Contingent Interest Barrier Value" según el emisor; la "barrera de capital" puede ser "Barrier Price" o "Barrier Value"; el "nivel de call" puede ser "Initial Value" (call si el precio vuelve a superarlo, sin nombre propio) o tener nombre explícito como "Call Threshold Value". Identifica el concepto correcto aunque el nombre exacto cambie.
- "fechas_call" es la lista de fechas en las que la nota PODRÍA ser llamada anticipadamente por el emisor. Hay tres formatos posibles, igual que con el calendario de cupón — busca el que aplique:
  a) Una tabla separada llamada "Early Redemption Date", "Call Settlement Date" o similar, con fechas explícitas — usa esas fechas directamente.
  b) Si el texto dice que el call se paga en "the first/next Interest Payment Date immediately following" cada Review/Observation Date (excluyendo la primera y la última) — usa esas mismas fechas de PAGO del calendario, excluyendo la primera y la última.
  c) Si las "Call Observation Dates" se describen como una REGLA DE REPETICIÓN propia y distinta a la del cupón (ej. TD Bank: "Quarterly, on the 2nd calendar day of each March, June, September and December, commencing on..."), genera tú mismo esa lista de fechas aplicando la regla (nota que puede tener una periodicidad distinta a la del cupón — ej. cupón mensual pero call solo trimestral), y usa como fecha de call la correspondiente Payment/Call Payment Date de cada una (aplicando la misma lógica de aproximación de días hábiles que en el calendario).
- "fecha_inicio_nota" es la fecha exacta que se usará para consultar el precio de cierre REAL de mercado de cada ticker (no hace falta que extraigas ningún precio en dólares del texto — solo esta fecha). Aunque el documento sea un borrador "SUBJECT TO COMPLETION" con precios en $[●], la fecha de Initial Valuation Date suele estar indicada igualmente (a veces entre corchetes tipo "September [9], 2025", en cuyo caso usa esa fecha igualmente).
- "tiene_memoria": true SOLO si el documento describe un cupón "memoria" — es decir, si un periodo no se cumple la condición de barrera, ese cupón no se pierde, sino que se acumula y se paga junto con un cupón futuro cuando sí se cumpla. Reconócelo por fórmulas del tipo "N x cupón% x (1 + T)" donde T es el número de periodos sin pago desde el último cupón pagado, o por texto explícito tipo "Memory Coupon" / "Memory Interest". Si el cupón de cada periodo es independiente (lo que no se cobra un mes se pierde para siempre, sin acumular), pon false.
- "tiene_one_star": true SOLO si la condición de pago del cupón (o del call) requiere que BASTE CON UNA SOLA de las acciones subyacentes esté por encima de la barrera para que se pague (en vez de exigir que TODAS lo estén, que es lo habitual en notas "worst-of"). Reconócelo por frases tipo "if the Closing Price of at least one Underlying Asset is greater than or equal to..." en vez de "if the Closing Price of each Underlying Asset is greater than or equal to...". Si la condición exige que TODAS las acciones estén por encima (worst-of, lo más común), pon false.
- Si un dato concreto no aparece en el documento o no estás seguro, usa el string "REVISAR" en ese campo en vez de inventar un número o fecha.
- No añadas ningún campo que no esté en el esquema. No expliques nada, solo el JSON."""

    if usar_texto_plano:
        contenido = [
            {"type": "text", "text": f"Documento (texto extraído, {n_paginas} páginas — demasiado largo para enviarlo como PDF; los términos clave y el calendario siempre están en las primeras páginas):\n\n{texto_extraido[:250000]}"},
            {"type": "text", "text": "Extrae los datos de esta nota estructurada según el esquema JSON indicado."},
        ]
    else:
        pdf_b64 = base64.b64encode(pdf_bytes).decode("utf-8")
        contenido = [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64}, "title": "NOTA_PDF"},
            {"type": "text", "text": "Extrae los datos de esta nota estructurada según el esquema JSON indicado."},
        ]

    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
        json={
            "model": "claude-sonnet-4-5",
            "max_tokens": 4000,
            "system": system,
            "messages": [{"role": "user", "content": contenido}],
        },
        timeout=90,
    )
    data = resp.json()
    texto = "".join(b.get("text", "") for b in data.get("content", []) if b.get("type") == "text")
    if not texto:
        return {"error": data.get("error", {}).get("message", "La IA no devolvió ninguna respuesta.")}

    texto_limpio = texto.strip()
    if texto_limpio.startswith("```"):
        texto_limpio = texto_limpio.strip("`")
        if texto_limpio.lower().startswith("json"):
            texto_limpio = texto_limpio[4:]
    try:
        return _json.loads(texto_limpio)
    except Exception as e:
        return {"error": f"No se pudo interpretar la respuesta de la IA como JSON: {e}", "respuesta_cruda": texto}


def extraer_datos_gasto_con_ia(pdf_bytes: bytes) -> dict:
    """
    Envía una factura/recibo en PDF (Railway, GoDaddy, Anthropic, Twilio, etc.) a Claude
    y le pide que devuelva ÚNICAMENTE JSON con los datos clave: proveedor, concepto,
    importe, moneda y fecha. Si algún dato no se puede determinar, usa "REVISAR" para
    que el usuario lo revise a mano en vez de que la IA invente un número.
    """
    import base64
    import json as _json

    if pdf_bytes and not pdf_bytes.startswith(b"%PDF-"):
        _inicio = pdf_bytes.find(b"%PDF-")
        _fin = pdf_bytes.rfind(b"%%EOF")
        if _inicio != -1 and _fin != -1 and _fin > _inicio:
            pdf_bytes = pdf_bytes[_inicio:_fin + len(b"%%EOF")]

    if not pdf_bytes or not pdf_bytes.startswith(b"%PDF-"):
        return {"error": "El archivo no parece ser un PDF válido (no empieza con la firma %PDF-)."}

    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")

    system = """Extraes datos de una factura o recibo en PDF (de proveedores tecnológicos como Railway,
GoDaddy, Anthropic, Twilio, Google, Microsoft, etc.).

Devuelve ÚNICAMENTE un JSON válido, sin texto antes ni después, sin backticks de markdown, con este esquema exacto:

{
  "proveedor": "string, nombre del proveedor/empresa que emite la factura (ej. Railway, GoDaddy, Anthropic)",
  "concepto": "string corto (máx 8 palabras) describiendo qué se pagó (ej. 'Plan Hobby - agosto 2026', 'Registro de dominio .com 3 años')",
  "importe": number, el importe TOTAL pagado (el que el cliente realmente paga, con impuestos incluidos si los hay), sin símbolo de moneda,
  "moneda": "string, código de 3 letras: EUR, USD, etc.",
  "fecha": "YYYY-MM-DD, la fecha de la factura o del cargo",
  "numero_factura": "string, el número/identificador único de factura o recibo tal como aparece en el documento (ej. 'INV-2026-08-1234', 'Receipt #8-8348'). Si no aparece ninguno, usa 'REVISAR'"
}

REGLAS:
- "importe" es SIEMPRE el total final pagado, no un subtotal ni una tarifa de lista.
- "numero_factura" cópialo EXACTAMENTE tal como aparece en el documento (incluyendo prefijos, guiones, #, etc.) — se usa para detectar facturas duplicadas, así que la exactitud importa más que la limpieza.
- Si algún dato no aparece en el documento o no estás seguro, usa el string "REVISAR" en ese campo (o -1 para importe si no se puede determinar) en vez de inventar.
- No añadas ningún campo que no esté en el esquema. No expliques nada, solo el JSON."""

    pdf_b64 = base64.b64encode(pdf_bytes).decode("utf-8")
    contenido = [
        {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64}, "title": "FACTURA_PDF"},
        {"type": "text", "text": "Extrae los datos de esta factura según el esquema JSON indicado."},
    ]

    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
        json={
            "model": "claude-sonnet-4-5",
            "max_tokens": 1000,
            "system": system,
            "messages": [{"role": "user", "content": contenido}],
        },
        timeout=90,
    )
    data = resp.json()
    texto = "".join(b.get("text", "") for b in data.get("content", []) if b.get("type") == "text")
    if not texto:
        return {"error": data.get("error", {}).get("message", "La IA no devolvió ninguna respuesta.")}

    texto_limpio = texto.strip()
    if texto_limpio.startswith("```"):
        texto_limpio = texto_limpio.strip("`")
        if texto_limpio.lower().startswith("json"):
            texto_limpio = texto_limpio[4:]
    try:
        return _json.loads(texto_limpio)
    except Exception as e:
        return {"error": f"No se pudo interpretar la respuesta de la IA como JSON: {e}", "respuesta_cruda": texto}


CARPETA_DATOS = "/data"  # volumen persistente de Railway (mount path: /data)
CARPETA_FACTURAS = "/data/facturas"
ARCHIVO_GASTOS_JSON = "/data/gastos_plataforma.json"
CARPETA_PDFS_NOTAS = "/data/notas_pdfs"

# Carpeta "PDFs Notas - CF Wealth" en Google Drive, compartida en modo escritor con la cuenta
# de servicio del fondo (app-fondo-writer@fondo-write.iam.gserviceaccount.com). Sirve como
# respaldo de los PDFs fuera del volumen de Railway y como origen alternativo si el volumen
# se pierde o se recrea desde cero.
GDRIVE_CARPETA_PDFS_NOTAS = "1j2Bj1sq_80agm47RRJ_zo953wigcU9qZ"


def _servicio_drive():
    """Crea un cliente autenticado de la API de Drive con la cuenta de servicio del fondo
    (la misma que ya usa la app para leer/escribir inversiones.xlsx). Devuelve None si no
    hay credenciales configuradas, sin lanzar excepción."""
    if "gcp_service_account" not in st.secrets:
        return None
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        credenciales = service_account.Credentials.from_service_account_info(
            dict(st.secrets["gcp_service_account"]),
            scopes=["https://www.googleapis.com/auth/drive"],
        )
        return build("drive", "v3", credentials=credenciales)
    except Exception:
        return None


def _subir_pdf_nota_a_drive(numero_nota: int, pdf_bytes: bytes) -> bool:
    """Sube (o actualiza si ya existe) el PDF de una nota en la carpeta de respaldo de Google
    Drive. Es un respaldo best-effort: si falla, nunca debe bloquear el guardado principal
    en el volumen de Railway, así que nunca lanza excepción hacia afuera."""
    servicio = _servicio_drive()
    if servicio is None:
        return False
    try:
        from googleapiclient.http import MediaIoBaseUpload
        nombre_archivo = f"nota_{int(numero_nota):02d}.pdf"
        existentes = servicio.files().list(
            q=f"name='{nombre_archivo}' and '{GDRIVE_CARPETA_PDFS_NOTAS}' in parents and trashed=false",
            fields="files(id)",
        ).execute().get("files", [])
        media = MediaIoBaseUpload(BytesIO(pdf_bytes), mimetype="application/pdf", resumable=False)
        if existentes:
            servicio.files().update(fileId=existentes[0]["id"], media_body=media, fields="id").execute()
        else:
            servicio.files().create(
                body={"name": nombre_archivo, "parents": [GDRIVE_CARPETA_PDFS_NOTAS]},
                media_body=media, fields="id",
            ).execute()
        return True
    except Exception:
        return False


def _descargar_pdf_nota_de_drive(numero_nota: int) -> bytes | None:
    """Descarga el PDF de respaldo de una nota desde Drive, si existe. Se usa como último
    recurso cuando el archivo no está en el volumen local de Railway."""
    servicio = _servicio_drive()
    if servicio is None:
        return None
    try:
        from googleapiclient.http import MediaIoBaseDownload
        nombre_archivo = f"nota_{int(numero_nota):02d}.pdf"
        encontrados = servicio.files().list(
            q=f"name='{nombre_archivo}' and '{GDRIVE_CARPETA_PDFS_NOTAS}' in parents and trashed=false",
            fields="files(id)",
        ).execute().get("files", [])
        if not encontrados:
            return None
        buffer = BytesIO()
        descargador = MediaIoBaseDownload(buffer, servicio.files().get_media(fileId=encontrados[0]["id"]))
        listo = False
        while not listo:
            _, listo = descargador.next_chunk()
        return buffer.getvalue()
    except Exception:
        return None


def guardar_pdf_nota(numero_nota: int, pdf_bytes: bytes) -> str:
    """Guarda el PDF oficial de una nota en el volumen persistente de Railway (/data), para
    poder auditarla más adelante sin tener que volver a pedírselo a nadie. Devuelve la ruta
    guardada, o cadena vacía si no se pudo guardar (por ejemplo, en local sin volumen montado).
    Además, sube una copia de respaldo a Google Drive (best-effort: si falla, no afecta al
    guardado principal ni al flujo de la nota)."""
    ruta = ""
    try:
        os.makedirs(CARPETA_PDFS_NOTAS, exist_ok=True)
        ruta = os.path.join(CARPETA_PDFS_NOTAS, f"nota_{int(numero_nota):02d}.pdf")
        with open(ruta, "wb") as f:
            f.write(pdf_bytes)
    except Exception:
        ruta = ""
    try:
        _subir_pdf_nota_a_drive(numero_nota, pdf_bytes)
    except Exception:
        pass
    try:
        from postgres_writer import guardar_pdf_nota_postgres
        guardar_pdf_nota_postgres(numero_nota, pdf_bytes)
    except Exception:
        pass
    return ruta


def leer_pdf_nota_guardado(numero_nota: int) -> bytes | None:
    """Lee el PDF guardado de una nota. Primero intenta el volumen local de Railway (rápido);
    si no está ahí, cae al respaldo en Google Drive y lo vuelve a guardar localmente de paso,
    para no tener que volver a descargarlo la próxima vez."""
    ruta = os.path.join(CARPETA_PDFS_NOTAS, f"nota_{int(numero_nota):02d}.pdf")
    try:
        with open(ruta, "rb") as f:
            return f.read()
    except Exception:
        pass
    try:
        from postgres_reader import leer_pdf_nota_postgres
        pdf_pg = leer_pdf_nota_postgres(numero_nota)
        if pdf_pg:
            try:
                os.makedirs(CARPETA_PDFS_NOTAS, exist_ok=True)
                with open(ruta, "wb") as f:
                    f.write(pdf_pg)
            except Exception:
                pass
            return pdf_pg
    except Exception:
        pass
    pdf_drive = _descargar_pdf_nota_de_drive(numero_nota)
    if pdf_drive:
        try:
            os.makedirs(CARPETA_PDFS_NOTAS, exist_ok=True)
            with open(ruta, "wb") as f:
                f.write(pdf_drive)
        except Exception:
            pass
        return pdf_drive
    return None


def _guardar_pdf_factura(pdf_bytes: bytes, nombre_sugerido: str) -> Optional[str]:
    """Guarda el PDF en el volumen persistente de Railway (/data/facturas) y devuelve el
    nombre de archivo guardado, o None si el volumen no está montado (por ejemplo en local
    o si Railway todavía no tiene el volumen configurado) — en ese caso no falla, solo
    guarda los datos sin el PDF."""
    try:
        os.makedirs(CARPETA_FACTURAS, exist_ok=True)
        marca = datetime.now().strftime("%Y%m%d%H%M%S")
        nombre_limpio = re.sub(r"[^a-zA-Z0-9_.-]", "_", nombre_sugerido)[:60]
        nombre_archivo = f"{marca}_{nombre_limpio}.pdf"
        with open(os.path.join(CARPETA_FACTURAS, nombre_archivo), "wb") as f:
            f.write(pdf_bytes)
        return nombre_archivo
    except Exception:
        return None


def _leer_gastos_json() -> list:
    """Lee la lista de gastos guardada en el volumen persistente (/data/gastos_plataforma.json).
    Si el archivo no existe todavía (primer uso), devuelve una lista vacía sin fallar."""
    try:
        if not os.path.isfile(ARCHIVO_GASTOS_JSON):
            return []
        with open(ARCHIVO_GASTOS_JSON, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []


def _escribir_gastos_json(lista: list) -> None:
    os.makedirs(CARPETA_DATOS, exist_ok=True)
    with open(ARCHIVO_GASTOS_JSON, "w", encoding="utf-8") as f:
        json.dump(lista, f, ensure_ascii=False, indent=2)


def _hash_pdf(pdf_bytes: bytes) -> str:
    import hashlib
    return hashlib.sha256(pdf_bytes).hexdigest()


def detectar_gastos_duplicados(gasto: dict, pdf_bytes: Optional[bytes] = None) -> list:
    """Compara un gasto (todavía sin guardar) contra los ya guardados y devuelve una lista
    de avisos de posible duplicado, de más a menos fiable:
      - 'exacto': el PDF es byte a byte idéntico a uno ya guardado (mismo archivo subido 2 veces)
      - 'numero_factura': mismo proveedor + mismo número de factura ya registrado
      - 'probable': mismo proveedor + mismo importe + misma fecha (podría ser coincidencia,
        pero es raro que dos facturas distintas cuadren exacto en los tres campos)
    No bloquea el guardado — solo informa, la decisión final es de quien lo guarda."""
    avisos = []
    lista = _leer_gastos_json()
    if not lista:
        return avisos

    hash_nuevo = _hash_pdf(pdf_bytes) if pdf_bytes else None
    numero_nuevo = str(gasto.get("numero_factura", "")).strip()
    proveedor_nuevo = str(gasto.get("proveedor", "")).strip().lower()
    importe_nuevo = gasto.get("importe")
    fecha_nueva = str(gasto.get("fecha", "")).strip()

    for g in lista:
        proveedor_existente = str(g.get("proveedor", "")).strip().lower()

        if hash_nuevo:
            nombre_pdf = g.get("archivo_pdf", "")
            ruta = os.path.join(CARPETA_FACTURAS, nombre_pdf) if nombre_pdf else None
            if ruta and os.path.isfile(ruta):
                try:
                    with open(ruta, "rb") as f:
                        if _hash_pdf(f.read()) == hash_nuevo:
                            avisos.append({
                                "tipo": "exacto",
                                "mensaje": f"Este PDF ya está guardado tal cual — factura de {g.get('proveedor')} del {g.get('fecha')} por {g.get('importe')} {g.get('moneda', 'EUR')}.",
                            })
                            continue
                except Exception:
                    pass

        if (numero_nuevo and numero_nuevo != "REVISAR"
                and proveedor_existente == proveedor_nuevo
                and str(g.get("numero_factura", "")).strip() == numero_nuevo):
            avisos.append({
                "tipo": "numero_factura",
                "mensaje": f"Ya hay una factura de {g.get('proveedor')} con el mismo número ({numero_nuevo}), del {g.get('fecha')}.",
            })
            continue

        if (proveedor_existente == proveedor_nuevo
                and g.get("importe") == importe_nuevo
                and str(g.get("fecha", "")).strip() == fecha_nueva
                and proveedor_nuevo):
            avisos.append({
                "tipo": "probable",
                "mensaje": f"Ya hay un gasto de {g.get('proveedor')} con el mismo importe ({importe_nuevo} {g.get('moneda', 'EUR')}) y la misma fecha ({fecha_nueva}).",
            })

    return avisos


def guardar_gasto_plataforma(gasto: dict, pdf_bytes: Optional[bytes] = None) -> tuple[bool, str]:
    """Guarda el gasto directamente en el volumen persistente de Railway (/data), como un
    archivo JSON — totalmente independiente del Excel/Google Drive, para no depender de la
    credencial de servicio. Si se pasa el PDF original, también se guarda ahí mismo."""
    try:
        nombre_pdf_guardado = ""
        if pdf_bytes:
            nombre_pdf_guardado = _guardar_pdf_factura(pdf_bytes, gasto.get("proveedor", "factura")) or ""

        lista = _leer_gastos_json()
        nuevo_id = (max((g.get("id", 0) for g in lista), default=0)) + 1
        lista.append({
            "id": nuevo_id,
            "fecha": gasto.get("fecha"),
            "proveedor": gasto.get("proveedor"),
            "concepto": gasto.get("concepto"),
            "categoria": gasto.get("categoria"),
            "importe": gasto.get("importe"),
            "moneda": gasto.get("moneda", "EUR"),
            "numero_factura": gasto.get("numero_factura", ""),
            "archivo_pdf": nombre_pdf_guardado,
            "registrado_por": gasto.get("registrado_por", ""),
            "registrado_en": datetime.now().isoformat(),
        })
        _escribir_gastos_json(lista)

        mensaje = "Guardado en el disco persistente de la app."
        if pdf_bytes and not nombre_pdf_guardado:
            mensaje += " (El PDF no se pudo guardar — ¿está creado el volumen /data en Railway?)"
        return True, mensaje
    except Exception as e:
        return False, f"Error al guardar: {type(e).__name__}: {e}"


def eliminar_gasto_plataforma(id_gasto: int) -> tuple[bool, str]:
    """Elimina un gasto por su id dentro de /data/gastos_plataforma.json."""
    try:
        lista = _leer_gastos_json()
        nueva_lista = [g for g in lista if g.get("id") != id_gasto]
        if len(nueva_lista) == len(lista):
            return False, "Ese gasto ya no existe (puede que la lista se haya actualizado)."
        _escribir_gastos_json(nueva_lista)
        return True, "Eliminado."
    except Exception as e:
        return False, f"Error al eliminar: {type(e).__name__}: {e}"


def _leer_gastos_plataforma_cached() -> pd.DataFrame:
    lista = _leer_gastos_json()
    if not lista:
        return pd.DataFrame(columns=["id", "fecha", "proveedor", "concepto", "categoria", "importe", "moneda", "archivo_pdf"])
    df = pd.DataFrame(lista)
    df["fecha"] = pd.to_datetime(df["fecha"], errors="coerce")
    df["importe"] = pd.to_numeric(df["importe"], errors="coerce").fillna(0)
    return df.reset_index(drop=True)


TASAS_USD_RESPALDO = {"EUR": 1.08, "GBP": 1.27, "USD": 1.0}  # respaldo si la API de cambio falla


@st.cache_data(ttl=6 * 3600, show_spinner=False)
def _obtener_tasas_cambio_a_usd() -> dict:
    """Tipo de cambio de cada moneda a USD, actualizado cada 6h. Si la API externa falla
    (sin conexión, caída, etc.) usa un valor de respaldo fijo para no romper la app."""
    tasas = dict(TASAS_USD_RESPALDO)
    try:
        resp = requests.get("https://api.frankfurter.app/latest?from=EUR&to=USD,GBP", timeout=8)
        datos = resp.json().get("rates", {})
        if "USD" in datos:
            tasas["EUR"] = datos["USD"]
        if "GBP" in datos and datos["GBP"] > 0:
            tasas["GBP"] = datos["USD"] / datos["GBP"]  # EUR->USD entre EUR->GBP = GBP->USD
    except Exception:
        pass
    return tasas


def convertir_a_usd(importe: float, moneda: str) -> float:
    tasas = _obtener_tasas_cambio_a_usd()
    return importe * tasas.get(str(moneda).upper(), 1.0)


def seccion_gastos_plataforma():
    """Panel de administrador (solo Yuri): sube una factura en PDF, la IA extrae los datos
    clave, y se guarda en la hoja GASTOS_PLATAFORMA — con vista mensual del total gastado
    en mantener la plataforma (Railway, dominio, Anthropic, etc.)."""
    st.markdown("## 💰 Gastos de la plataforma")
    st.caption(
        "Sube la factura en PDF (Railway, GoDaddy, Anthropic, Twilio...) y la IA extrae "
        "proveedor, importe y fecha automáticamente. El PDF en sí no se guarda en la app "
        "— solo los datos; conserva tú el archivo original en tu propia carpeta."
    )

    if st.session_state.get("gasto_guardado_ok"):
        st.success(st.session_state.pop("gasto_guardado_ok"))

    pdf_subido = st.file_uploader("Factura en PDF", type=["pdf"], key="uploader_gasto_pdf")

    if pdf_subido is not None:
        if st.button("🔍 Extraer datos con IA", key="btn_extraer_gasto"):
            with st.spinner("Leyendo la factura..."):
                extraido = extraer_datos_gasto_con_ia(pdf_subido.getvalue())
            if "error" in extraido:
                st.error(extraido["error"])
            else:
                st.session_state["gasto_extraido"] = extraido
                st.success("Datos extraídos — revísalos abajo antes de guardar.")

    if "gasto_extraido" in st.session_state:
        ext = st.session_state["gasto_extraido"]
        st.markdown("#### Revisa y confirma")
        c1, c2 = st.columns(2)
        proveedor = c1.text_input("Proveedor", value=str(ext.get("proveedor", "")))
        concepto = c2.text_input("Concepto", value=str(ext.get("concepto", "")))
        c3, c4, c5 = st.columns(3)
        importe_val = ext.get("importe", 0)
        try:
            importe_val = float(importe_val)
        except (TypeError, ValueError):
            importe_val = 0.0
        importe = c3.number_input("Importe", value=importe_val, step=0.01, min_value=0.0)
        moneda = c4.selectbox("Moneda", ["EUR", "USD", "GBP"], index=0 if ext.get("moneda") != "USD" else 1)
        categoria = c5.selectbox("Categoría", ["Hosting", "Dominio", "IA", "Email", "Telefonía", "Otros"])
        fecha_val = ext.get("fecha")
        try:
            fecha_dt = pd.to_datetime(fecha_val).date()
        except Exception:
            fecha_dt = datetime.now().date()
        fecha = st.date_input("Fecha de la factura", value=fecha_dt)
        numero_factura = st.text_input("Nº de factura (opcional, ayuda a detectar duplicados)", value=str(ext.get("numero_factura", "")) if ext.get("numero_factura") != "REVISAR" else "")
        gasto_actual = {
            "proveedor": proveedor, "concepto": concepto, "importe": importe,
            "moneda": moneda, "categoria": categoria, "fecha": str(fecha),
            "numero_factura": numero_factura,
        }
        avisos_duplicado = detectar_gastos_duplicados(gasto_actual, pdf_bytes=pdf_subido.getvalue() if pdf_subido is not None else None)
        if avisos_duplicado:
            for aviso in avisos_duplicado:
                if aviso["tipo"] == "exacto":
                    st.error(f"🔁 Posible duplicado exacto: {aviso['mensaje']}")
                elif aviso["tipo"] == "numero_factura":
                    st.error(f"🔁 Posible duplicado (mismo nº de factura): {aviso['mensaje']}")
                else:
                    st.warning(f"⚠️ Podría ser un duplicado: {aviso['mensaje']}")
            confirmar_duplicado = st.checkbox("Sé que puede ser un duplicado, guárdalo de todas formas", key="check_confirmar_duplicado")
        else:
            confirmar_duplicado = True

        if st.button("✅ Guardar gasto", key="btn_guardar_gasto", disabled=bool(avisos_duplicado) and not confirmar_duplicado):
            try:
                with st.spinner("Guardando..."):
                    exito, mensaje = guardar_gasto_plataforma(
                        gasto_actual | {"registrado_por": str(st.session_state.get("usuario", ""))},
                        pdf_bytes=pdf_subido.getvalue() if pdf_subido is not None else None,
                    )
                if exito:
                    st.session_state["gasto_guardado_ok"] = f"✅ Gasto guardado. {mensaje}"
                    del st.session_state["gasto_extraido"]
                    st.cache_data.clear()
                    st.rerun()
                else:
                    st.error(f"No se pudo guardar: {mensaje}")
            except Exception as e:
                st.error(f"Fallo inesperado al guardar (esto no debería pasar — copia este mensaje si le escribes a soporte): {type(e).__name__}: {e}")
                import traceback
                st.code(traceback.format_exc())

    st.divider()
    st.markdown("### Historial por mes")
    df_g = _leer_gastos_plataforma_cached()
    if df_g.empty:
        st.info("Todavía no hay ningún gasto registrado.")
        return

    df_g["mes"] = df_g["fecha"].dt.strftime("%Y-%m")
    meses_disponibles = sorted(df_g["mes"].dropna().unique(), reverse=True)
    mes_sel = st.selectbox("Mes", ["Todos"] + list(meses_disponibles))
    df_f = df_g if mes_sel == "Todos" else df_g[df_g["mes"] == mes_sel]
    df_f = df_f.copy()
    df_f["importe_usd"] = df_f.apply(lambda r: convertir_a_usd(r["importe"], r.get("moneda", "USD")), axis=1)

    mostrar_metricas("Resumen (en USD)", [
        ("Total", f"${df_f['importe_usd'].sum():,.2f}"),
        ("Facturas", f"{len(df_f)}"),
        ("Media por factura", f"${(df_f['importe_usd'].sum() / len(df_f)) if len(df_f) else 0:,.2f}"),
    ])
    st.caption(
        "Los importes originales en EUR/GBP se convierten a USD con el tipo de cambio del "
        "día (actualizado cada 6h). El importe original se conserva y se muestra en el detalle."
    )

    st.markdown("#### Por categoría")
    resumen_cat = df_f.groupby("categoria")["importe_usd"].sum().reset_index().sort_values("importe_usd", ascending=False)
    resumen_cat = resumen_cat.rename(columns={"importe_usd": "importe"})
    resumen_cat["importe"] = resumen_cat["importe"].map(lambda v: f"${v:,.2f}")
    st.dataframe(resumen_cat, use_container_width=True, hide_index=True)

    st.markdown("#### Detalle")
    detalle = df_f.sort_values("fecha", ascending=False).reset_index(drop=True)
    for _, fila in detalle.iterrows():
        c1, c2, c3, c4, c5, c6 = st.columns([2, 3, 2, 2, 1, 1])
        c1.write(fila["proveedor"])
        c2.write(fila["concepto"])
        c3.write(fila["categoria"])
        moneda_original = str(fila.get("moneda", "USD"))
        if moneda_original == "USD":
            c4.write(f"${fila['importe']:,.2f}")
        else:
            c4.write(f"${fila['importe_usd']:,.2f}")
            c4.caption(f"({fila['importe']:,.2f} {moneda_original} orig.)")
        nombre_pdf = str(fila.get("archivo_pdf", "") or "")
        ruta_pdf = os.path.join(CARPETA_FACTURAS, nombre_pdf) if nombre_pdf else None
        if ruta_pdf and os.path.isfile(ruta_pdf):
            with open(ruta_pdf, "rb") as f_pdf:
                c5.download_button("📄", data=f_pdf.read(), file_name=nombre_pdf, mime="application/pdf", key=f"dl_gasto_{fila['id']}")
        else:
            c5.write("—")
        if c6.button("🗑️", key=f"del_gasto_{fila['id']}"):
            exito, mensaje = eliminar_gasto_plataforma(int(fila["id"]))
            if exito:
                st.rerun()
            else:
                st.warning(mensaje)


def _reconstruir_extraido_desde_tablas(extraido: dict, df_control_editado: pd.DataFrame, df_cal_editado: pd.DataFrame, df_calls_editado: pd.DataFrame) -> dict:
    """
    Reconstruye la estructura 'extraido' (la misma que devuelve la IA) tomando como
    fuente de verdad las correcciones manuales que el usuario ya hizo en las tablas
    editables (st.data_editor), para poder volver a guardarlas como borrador y que
    sobrevivan a un redeploy sin perder el trabajo de revisión.
    """
    import itertools
    nuevo = dict(extraido)  # conserva emisor/cupon/fecha_vencimiento/fecha_inicio_nota si no están en tablas

    tickers = []
    for _, fila in df_control_editado.iterrows():
        tickers.append({
            "ticker": fila.get("TICKER"),
            "barrera_cupon_pct": fila.get("BARRERA_CUPON"),
            "call_level_pct": fila.get("CALL_LEVEL"),
            "barrera_capital_pct": fila.get("BARRERA_CAPITAL"),
        })
    nuevo["tickers"] = tickers

    if not df_control_editado.empty:
        emisor_col = df_control_editado.get("EMISOR")
        if emisor_col is not None and not emisor_col.dropna().empty:
            nuevo["emisor"] = emisor_col.dropna().iloc[0]
        for campo_bool, col in [("tiene_memoria", "TIENE_MEMORIA"), ("tiene_one_star", "TIENE_ONE_STAR")]:
            col_serie = df_control_editado.get(col)
            if col_serie is not None and not col_serie.dropna().empty:
                nuevo[campo_bool] = str(col_serie.dropna().iloc[0]).strip().upper() == "SI"

    obs_rows = df_cal_editado[df_cal_editado.get("TIPO_EVENTO") == "OBSERVACION"]["FECHA"].tolist() if not df_cal_editado.empty else []
    pago_rows = df_cal_editado[df_cal_editado.get("TIPO_EVENTO") == "PAGO"]["FECHA"].tolist() if not df_cal_editado.empty else []
    nuevo["calendario"] = [
        {"observacion": o, "pago": p} for o, p in itertools.zip_longest(obs_rows, pago_rows, fillvalue=None)
    ]

    nuevo["fechas_call"] = df_calls_editado["FECHA_CALL"].tolist() if not df_calls_editado.empty and "FECHA_CALL" in df_calls_editado.columns else []

    return nuevo


def _tab_añadir_nota_nueva(df_control: pd.DataFrame, df_cal: pd.DataFrame, df_calls: pd.DataFrame = None,
                            tipo_operacion: str = "NUEVA", df_inv: pd.DataFrame = None):
    st.caption(
        "Sube el documento oficial (pricing supplement) de una nota nueva. La IA extrae automáticamente "
        "tickers, precios iniciales, barreras, cupón y el calendario de observación/pago. "
        "Revisa siempre la previsualización antes de guardar — cualquier dato marcado con ⚠️ REVISAR "
        "no se pudo determinar con confianza y hay que rellenarlo a mano.\n\n"
        "Cada número de nota mantiene su propia previsualización guardada: si cambias de número, "
        "verás el formulario en blanco para una nota sin empezar, o los datos ya extraídos si vuelves "
        "a un número en el que ya trabajaste."
    )

    st.markdown("---")

    if "notas_wizard_datos" not in st.session_state:
        st.session_state["notas_wizard_datos"] = {}
    almacen = st.session_state["notas_wizard_datos"]

    numero_nota = int(st.number_input("Número de nota (ej. 28)", min_value=1, max_value=999, step=1, key="nueva_nota_numero"))

    # A diferencia de antes, YA NO se carga solo un borrador guardado de un intento anterior —
    # se avisa y se deja elegir, para no arrancar con datos de otra sesión sin darte cuenta.
    # "Ignorarlo" NO toca el Excel (a propósito): solo lo recuerda en esta sesión (session_state),
    # para que funcione bien incluso sin credenciales de Drive configuradas — si borráramos el
    # borrador del Excel aquí, la siguiente recarga automática desde Drive lo volvería a traer.
    if "notas_borradores_ignorados" not in st.session_state:
        st.session_state["notas_borradores_ignorados"] = set()
    borradores_ignorados = st.session_state["notas_borradores_ignorados"]

    if numero_nota not in almacen and numero_nota not in borradores_ignorados:
        borrador_disponible = cargar_borrador_nota("nueva", numero_nota)
        if borrador_disponible:
            st.info(f"📂 Hay un borrador guardado de un intento anterior para la Nota {numero_nota}.")
            col_cargar, col_ignorar = st.columns(2)
            with col_cargar:
                if st.button(f"📂 Cargar ese borrador", key=f"cargar_borrador_nota_{numero_nota}"):
                    almacen[numero_nota] = borrador_disponible
                    st.rerun()
            with col_ignorar:
                if st.button("🗑️ Ignorarlo y empezar en blanco", key=f"ignorar_borrador_nota_{numero_nota}"):
                    borradores_ignorados.add(numero_nota)
                    st.rerun()
            return

    extraido_guardado = almacen.get(numero_nota)

    if extraido_guardado:
        st.info(f"📂 Ya hay datos extraídos guardados para la Nota {numero_nota}. Puedes revisarlos abajo, o subir un PDF nuevo para reemplazarlos.")

    pdf_subido = st.file_uploader(f"Documento oficial de la Nota {numero_nota} (PDF)", type=["pdf"], key=f"nueva_nota_pdf_{numero_nota}")
    verificar_doble = st.checkbox(
        "🔍 Verificar automáticamente con una segunda lectura independiente (recomendado — tarda un poco más, pero detecta si la IA se confunde en una fecha o una tasa)",
        value=True, key=f"nueva_nota_doble_lectura_{numero_nota}",
    )

    if st.button("🔎 Extraer datos con IA", type="primary", disabled=pdf_subido is None):
        pdf_bytes_subido = pdf_subido.read()
        with st.spinner("Leyendo el documento y extrayendo los datos..."):
            resultado = extraer_datos_nota_con_ia(pdf_bytes_subido)
        if "error" in resultado:
            st.error(f"No se pudieron extraer los datos: {resultado['error']}")
        else:
            almacen[numero_nota] = resultado
            almacen_pdfs_notas = st.session_state.setdefault("nota_pdf_bytes", {})
            almacen_pdfs_notas[numero_nota] = pdf_bytes_subido
            st.session_state.pop(f"nota_doble_lectura_resultado_{numero_nota}", None)
            if verificar_doble:
                with st.spinner("Verificando con una segunda lectura independiente..."):
                    segundo_resultado = extraer_datos_nota_con_ia(pdf_bytes_subido)
                if "error" not in segundo_resultado:
                    st.session_state[f"nota_doble_lectura_resultado_{numero_nota}"] = segundo_resultado
            st.success(f"Datos extraídos para la Nota {numero_nota}. Revisa la previsualización abajo.")
            guardar_borrador_nota("nueva", numero_nota, resultado)

    extraido = almacen.get(numero_nota)
    if not extraido:
        return

    # ── Verificación automática: comparar la 1ª lectura contra una 2ª lectura independiente ──
    segundo = st.session_state.get(f"nota_doble_lectura_resultado_{numero_nota}")
    if segundo:
        campos_clave_verificar = [
            ("emisor", "Emisor"), ("cupon_anual_pct", "Cupón anual (%)"),
            ("fecha_vencimiento", "Fecha de vencimiento"), ("fecha_inicio_nota", "Fecha de inicio"),
        ]
        discrepancias_doble = []
        for clave_campo, etiqueta in campos_clave_verificar:
            v1, v2 = extraido.get(clave_campo), segundo.get(clave_campo)
            if v1 is None and v2 is None:
                continue
            if str(v1).strip().upper() != str(v2).strip().upper():
                discrepancias_doble.append({"Campo": etiqueta, "1ª lectura": v1, "2ª lectura": v2})
        tickers1 = {str(t.get("ticker", "")).strip().upper(): t for t in extraido.get("tickers", [])}
        tickers2 = {str(t.get("ticker", "")).strip().upper(): t for t in segundo.get("tickers", [])}
        for tk in set(tickers1) | set(tickers2):
            t1, t2 = tickers1.get(tk, {}), tickers2.get(tk, {})
            for campo_ticker in ["precio_compra", "barrera_cupon", "barrera_capital", "call_level"]:
                v1, v2 = t1.get(campo_ticker), t2.get(campo_ticker)
                if v1 is None and v2 is None:
                    continue
                if str(v1).strip().upper() != str(v2).strip().upper():
                    discrepancias_doble.append({"Campo": f"{tk} — {campo_ticker}", "1ª lectura": v1, "2ª lectura": v2})
        if discrepancias_doble:
            st.warning(
                f"⚠️ **La verificación automática detectó {len(discrepancias_doble)} diferencia(s)** entre "
                "la primera y la segunda lectura del PDF — revísalas con especial cuidado antes de guardar, "
                "puede que la IA se haya confundido en algún dato:"
            )
            st.dataframe(pd.DataFrame(discrepancias_doble), use_container_width=True, hide_index=True)
        else:
            st.success("✅ Verificación automática: la segunda lectura coincide con la primera en emisor, cupón, fechas y datos de tickers.")

    def _marcar(valor):
        return "⚠️ REVISAR" if (valor is None or str(valor).strip().upper() == "REVISAR") else valor

    st.markdown("---")
    st.markdown(f"### Nota {numero_nota}")
    st.markdown(f"**Emisor:** {_marcar(extraido.get('emisor'))}  |  **CUSIP:** {_marcar(extraido.get('cusip'))}  |  **Cupón anual:** {_marcar(extraido.get('cupon_anual_pct'))}  |  **Vencimiento:** {_marcar(extraido.get('fecha_vencimiento'))}")
    etiquetas = []
    if extraido.get("tiene_memoria") is True:
        etiquetas.append("🧠 Memoria")
    if extraido.get("tiene_one_star") is True:
        etiquetas.append("⭐ One-star")
    if etiquetas:
        st.markdown("  |  ".join(f"**{e}**" for e in etiquetas))

    fecha_inicio_nota = extraido.get("fecha_inicio_nota")
    fecha_inicio_valida = None
    if fecha_inicio_nota and str(fecha_inicio_nota).strip().upper() != "REVISAR":
        fecha_inicio_valida = pd.to_datetime(fecha_inicio_nota, errors="coerce")
    st.markdown(f"**Fecha de inicio de la nota (Initial Valuation Date):** {_marcar(fecha_inicio_nota)}")

    st.markdown("#### Tickers y barreras (se guardará en CONTROL_NOTAS)")
    st.caption("El precio de compra se consulta en vivo con el cierre real de mercado de cada ticker en la fecha de inicio — no se extrae del texto del PDF.")
    filas_control = []
    for t in extraido.get("tickers", []):
        ticker = str(t.get("ticker", "")).strip().upper()
        precio_real = None
        if fecha_inicio_valida is not None and pd.notna(fecha_inicio_valida) and ticker:
            precio_real = obtener_cierre_ticker_fecha(ticker, fecha_inicio_valida)
        filas_control.append({
            "NOTA": numero_nota,
            "TICKER": _marcar(t.get("ticker")),
            "PRECIO_COMPRA": _marcar(precio_real),
            "BARRERA_CUPON": _marcar(t.get("barrera_cupon_pct")),
            "PRECIO_ACTUAL": None,
            "CALL_LEVEL": _marcar(t.get("call_level_pct")),
            "BARRERA_CAPITAL": _marcar(t.get("barrera_capital_pct")),
            "TIENE_MEMORIA": "SI" if extraido.get("tiene_memoria") is True else "NO",
            "TIENE_ONE_STAR": "SI" if extraido.get("tiene_one_star") is True else "NO",
            "EMISOR": _marcar(extraido.get("emisor")),
            "CUSIP": _marcar(extraido.get("cusip")),
        })
    df_control_preview = pd.DataFrame(filas_control)
    df_control_editado = st.data_editor(df_control_preview, use_container_width=True, num_rows="dynamic", key=f"editor_control_nota_{numero_nota}")

    st.markdown("#### Calendario de observación/pago (se guardará en CALENDARIO_NOTAS)")
    filas_cal = []
    for evento in extraido.get("calendario", []):
        filas_cal.append({"NOTA": numero_nota, "TIPO_EVENTO": "OBSERVACION", "FECHA": _marcar(evento.get("observacion"))})
        filas_cal.append({"NOTA": numero_nota, "TIPO_EVENTO": "PAGO", "FECHA": _marcar(evento.get("pago"))})
    df_cal_preview = pd.DataFrame(filas_cal)
    df_cal_editado = st.data_editor(df_cal_preview, use_container_width=True, num_rows="dynamic", key=f"editor_cal_nota_{numero_nota}")

    # Fecha de inicio sugerida para INVERSIONES: primer PAGO menos 1 mes
    fechas_pago = pd.to_datetime(
        [e.get("pago") for e in extraido.get("calendario", []) if e.get("pago") and str(e.get("pago")).strip().upper() != "REVISAR"],
        errors="coerce",
    )
    fechas_pago = fechas_pago.dropna() if len(fechas_pago) else fechas_pago
    if len(fechas_pago) > 0:
        primer_pago = min(fechas_pago)
        fecha_inversion_sugerida = primer_pago - pd.DateOffset(months=1)
        st.info(
            f"📌 **Fecha de inicio para INVERSIONES (pago a inversores):** {fecha_inversion_sugerida.strftime('%d/%m/%Y')} "
            f"— calculada como la primera fecha de PAGO ({primer_pago.strftime('%d/%m/%Y')}) menos 1 mes. "
            "Usa esta fecha como `fecha_inversion` al dar de alta al inversor de esta nota en `INVERSIONES` "
            "(este wizard no toca esa hoja, hazlo desde Centro de Control)."
        )

    st.markdown("#### Fechas de posible call (se guardará en CALENDARIO_CALLS)")
    filas_calls = [{"NOTA": numero_nota, "FECHA_CALL": _marcar(f), "ESTADO": None, "OBSERVACIONES": None} for f in extraido.get("fechas_call", [])]
    df_calls_preview = pd.DataFrame(filas_calls)
    df_calls_editado = st.data_editor(df_calls_preview, use_container_width=True, num_rows="dynamic", key=f"editor_calls_nota_{numero_nota}")

    hay_revisar = (
        df_control_editado.astype(str).apply(lambda c: c.str.contains("REVISAR", na=False)).any().any()
        or df_cal_editado.astype(str).apply(lambda c: c.str.contains("REVISAR", na=False)).any().any()
        or df_calls_editado.astype(str).apply(lambda c: c.str.contains("REVISAR", na=False)).any().any()
    )
    if hay_revisar:
        st.warning("Hay campos marcados con ⚠️ REVISAR — corrígelos en las tablas de arriba antes de guardar.")

    # ── Comparación contra lo que YA está guardado en el Excel para esta nota ──────────
    # Guardar ahora REEMPLAZA cualquier fila vieja de esta nota (ver _quitar_nota_vieja más abajo),
    # así que si ya había datos guardados, hay que mostrarlos lado a lado con los nuevos ANTES de
    # dejar guardar — si no, se puede perder sin darse cuenta una corrección manual que se había
    # hecho directo en el Excel y que esta nueva extracción no tiene.
    control_existente_nota = df_control[pd.to_numeric(df_control.get("nota"), errors="coerce") == numero_nota].copy() if df_control is not None and not df_control.empty and "nota" in df_control.columns else pd.DataFrame()
    hay_discrepancia = False
    confirmar_reemplazo = True  # si no hay nada guardado todavía, no hace falta confirmar nada
    if not control_existente_nota.empty:
        campos_comparar = [
            ("precio_compra", "PRECIO_COMPRA"), ("barrera_cupon", "BARRERA_CUPON"),
            ("barrera_capital", "BARRERA_CAPITAL"), ("call_level", "CALL_LEVEL"),
            ("emisor", "EMISOR"), ("tiene_memoria", "TIENE_MEMORIA"), ("tiene_one_star", "TIENE_ONE_STAR"),
            ("cusip", "CUSIP"),
        ]
        filas_comparacion = []
        for _, fila_nueva in df_control_editado.iterrows():
            ticker_nuevo = str(fila_nueva.get("TICKER", "")).strip().upper()
            fila_vieja = control_existente_nota[control_existente_nota["ticker"].astype(str).str.strip().str.upper() == ticker_nuevo]
            for col_vieja, col_nueva in campos_comparar:
                valor_viejo = fila_vieja.iloc[0][col_vieja] if not fila_vieja.empty and col_vieja in fila_vieja.columns else "(no existía)"
                valor_nuevo = fila_nueva.get(col_nueva)
                iguales = str(valor_viejo).strip().upper() == str(valor_nuevo).strip().upper()
                if not iguales:
                    filas_comparacion.append({
                        "ticker": ticker_nuevo, "campo": col_nueva,
                        "valor actual en Excel": valor_viejo, "valor nuevo extraído": valor_nuevo,
                    })
        if filas_comparacion:
            hay_discrepancia = True
            st.warning(
                f"⚠️ La Nota {numero_nota} ya tiene datos guardados en el Excel y hay {len(filas_comparacion)} "
                "diferencia(s) con lo que se acaba de extraer/editar arriba. Si guardás, esto REEMPLAZA lo que "
                "ya había — revisá la comparación antes de confirmar:"
            )
            st.dataframe(pd.DataFrame(filas_comparacion), use_container_width=True, hide_index=True)
            confirmar_reemplazo = st.checkbox(
                f"Revisé las diferencias de arriba y confirmo que quiero reemplazar los datos guardados de la Nota {numero_nota}",
                key=f"confirmar_reemplazo_nota_{numero_nota}",
            )
        else:
            st.caption(f"✅ Sin diferencias: coincide con lo que ya había guardado en el Excel para la Nota {numero_nota}.")

    # ── Origen de la reinversión (solo si tipo_operacion == REINVERSION) ───────────────
    ids_origen_nota = []
    if tipo_operacion == "REINVERSION":
        st.markdown("---")
        st.markdown("#### 🔁 Origen de la reinversión")
        st.caption(
            "Podés elegir una o varias posiciones de origen (por ejemplo, si el capital de esta nota "
            "nueva viene de varias posiciones que se cerraron juntas). Se guarda en id_inversion_origen "
            "de cada fila de inversor (separadas por coma si son varias) — no modifica ni cierra las "
            "posiciones originales. Si además hay que cerrarlas formalmente, hacelo por separado con "
            "una operación CANCELADA."
        )
        activas_origen_nota = _posiciones_activas_para_cerrar(df_inv)
        if not activas_origen_nota.empty:
            activas_origen_nota = activas_origen_nota.copy()
            activas_origen_nota["_etiqueta"] = activas_origen_nota.apply(
                lambda r: f"{r.get('id_inversion','?')} | {r.get('inversor','?')} | {r.get('nombre_activo','?')} | ${float(r.get('capital_invertido',0) or 0):,.2f}",
                axis=1,
            )
            etiquetas_origen_nota = st.multiselect(
                "Posición(es) de origen (opcional, se aplica a todos los inversores de abajo)",
                activas_origen_nota["_etiqueta"].tolist(), key=f"nota_origen_reinv_multisel_{numero_nota}",
            )
            if etiquetas_origen_nota:
                ids_origen_nota = [
                    str(activas_origen_nota[activas_origen_nota["_etiqueta"] == et].iloc[0].get("id_inversion", ""))
                    for et in etiquetas_origen_nota
                ]
        else:
            st.caption("No se encontraron posiciones activas (NUEVA/REINVERSION sin fecha_final_inversion) para vincular como origen.")

    # ── Inversor(es) de esta nota: mismo paso, sin cambiar de pestaña ──────────────────
    st.markdown("---")
    st.markdown("#### 💰 Inversor(es) de esta nota")
    st.caption(
        "Regístralos aquí mismo — al guardar, se escriben a la vez en CONTROL_NOTAS/CALENDARIO_NOTAS "
        "**y** en INVERSIONES, sin tener que ir luego a '➕ Nueva inversión'."
    )
    clave_num_inv_nota = f"nota_num_inversores_{numero_nota}"
    if clave_num_inv_nota not in st.session_state:
        st.session_state[clave_num_inv_nota] = 1

    try:
        inversores_conocidos_nota = sorted(set(list(USUARIOS_INVERSORES.keys()) + ["CHAPARRO FERNANDEZ"]))
    except NameError:
        inversores_conocidos_nota = []

    fecha_sugerida_inv, _periodicidad_inv, _fuente_inv = _sugerir_fecha_inicio_nota(int(numero_nota), df_cal)
    valor_default_fecha_inv = fecha_sugerida_inv.date() if fecha_sugerida_inv is not None else pd.Timestamp.today().date()
    if fecha_sugerida_inv is not None:
        periodo_txt_inv = {1: "mensual", 3: "trimestral", 6: "semestral"}.get(_periodicidad_inv, f"cada {_periodicidad_inv} meses")
        st.caption(f"📌 Fecha de inicio sugerida para la inversión: **{fecha_sugerida_inv.strftime('%d/%m/%Y')}** (primer pago menos 1 periodo — {periodo_txt_inv} —, de {_fuente_inv}).")
    else:
        st.caption("No se encontró un calendario de PAGO para sugerir la fecha — revisa/ajusta la fecha de inicio a mano abajo (primer pago menos 1 periodo).")
    valor_fecha_inicio_inv = st.date_input(
        "Fecha de inicio de la inversión (fecha_inversion)", value=valor_default_fecha_inv,
        key=f"nota_fecha_inicio_inv_{numero_nota}",
        help="Prellenada con la sugerencia (primer pago menos 1 periodo) — puedes escribir otra fecha a mano si corresponde.",
    )
    filas_inversores_nota = []
    for i in range(st.session_state[clave_num_inv_nota]):
        with st.container(border=True):
            st.markdown(f"**Inversor #{i + 1}**")
            c1, c2, c3 = st.columns(3)
            with c1:
                inv_sel = st.selectbox("Inversor", inversores_conocidos_nota + ["Otro (escribir)"], key=f"nota_inv_sel_{numero_nota}_{i}")
                inv_final = st.text_input("Escribe el nombre (exacto, mayúsculas)", key=f"nota_inv_libre_{numero_nota}_{i}") if inv_sel == "Otro (escribir)" else inv_sel
            with c2:
                capital_inv_i = st.number_input("Capital ($)", min_value=0.0, step=1000.0, format="%.2f", key=f"nota_inv_capital_{numero_nota}_{i}")
            with c3:
                tasa_inv_nota_i = st.number_input("Interés anual AL INVERSOR (%)", min_value=0.0, max_value=100.0, step=0.5, format="%.2f", key=f"nota_inv_tasa_{numero_nota}_{i}")
            c4, c5 = st.columns(2)
            with c4:
                id_sugerido_nota_i = _siguiente_id_inversion(df_inv, offset=i)
                if id_sugerido_nota_i:
                    st.caption(f"📌 Sugerido: **{id_sugerido_nota_i}**")
                id_inv_nota_i = st.text_input(
                    "id_inversion (ej. OP005)",
                    value=id_sugerido_nota_i, key=f"nota_inv_id_{numero_nota}_{i}",
                    help="Se sugiere automáticamente el siguiente número tras el último usado — puedes borrarlo y escribir otro a mano.",
                )
            with c5:
                email_inv_i = st.text_input("email (opcional)", key=f"nota_inv_email_{numero_nota}_{i}")
            cuenta_cobro_i = st.selectbox(
                "¿En qué cuenta se cobra esta nota?", ["COMPAÑÍA", "JORDI"],
                key=f"nota_inv_cuenta_cobro_{numero_nota}_{i}",
                help="COMPAÑÍA: el cobro de esta nota entra en las cuentas del fondo. JORDI: entra en la cuenta personal de Jordi (afecta a la deuda con Jordi).",
            )

            if st.session_state[clave_num_inv_nota] > 1:
                if st.button(f"🗑️ Quitar inversor #{i + 1}", key=f"nota_inv_quitar_{numero_nota}_{i}"):
                    st.session_state[clave_num_inv_nota] -= 1
                    st.rerun()

            filas_inversores_nota.append({
                "id_inversion": id_inv_nota_i, "inversor": inv_final, "tipo_inversion": "nota",
                "subtipo_inversion": "ESTRUCTURADA", "nombre_activo": f"NOTA_{int(numero_nota):02d}",
                "metodo_calculo": "NOTA", "cuenta_cobro": cuenta_cobro_i, "activo_generador_interes": "SI",
                "fecha_inversion": str(valor_fecha_inicio_inv), "fecha_final_inversion": "", "motivo": "",
                "capital_invertido": capital_inv_i,
                "interes_nota_anual": round(float(extraido.get("cupon_anual_pct") or 0), 6) if str(extraido.get("cupon_anual_pct", "")).strip().upper() != "REVISAR" else 0,
                "interes_inversor_anual": round(tasa_inv_nota_i / 100.0, 6),
                "tipo_operacion": tipo_operacion,
                "id_inversion_origen": ", ".join(ids_origen_nota) if ids_origen_nota else "",
                "capital_nuevo_real": "si" if tipo_operacion == "NUEVA" else "no",
                "email": email_inv_i, "pago_intereses": "reinvierte",
            })

    if st.button("➕ Añadir otro inversor a esta nota", key=f"nota_inv_add_{numero_nota}"):
        st.session_state[clave_num_inv_nota] += 1
        st.rerun()

    col_guardar, col_avance, col_descartar = st.columns([1, 1, 1])
    with col_avance:
        if st.button(f"📌 Guardar avance de mis correcciones", help="Guarda en Google Drive lo que ya has corregido en las tablas de arriba, aunque todavía queden ⚠️ REVISAR. Útil antes de tocar el código, para no perder tus correcciones al redesplegar."):
            extraido_actualizado = _reconstruir_extraido_desde_tablas(extraido, df_control_editado, df_cal_editado, df_calls_editado)
            almacen[numero_nota] = extraido_actualizado
            with st.spinner("Guardando avance..."):
                guardar_borrador_nota("nueva", numero_nota, extraido_actualizado)
    with col_guardar:
        if hay_discrepancia and not confirmar_reemplazo:
            st.button(f"💾 Guardar Nota {numero_nota} en el Excel", type="primary", disabled=True, help="Marcá el checkbox de arriba para confirmar que revisaste las diferencias.")
        elif st.button(f"💾 Guardar Nota {numero_nota} en el Excel", type="primary", disabled=hay_revisar):
            hojas = leer_todas_las_hojas_excel()
            if not hojas or "CONTROL_NOTAS" not in hojas or "CALENDARIO_NOTAS" not in hojas:
                st.error("No se pudo leer el Excel actual para guardar los cambios.")
            else:
                # Si esta nota ya se había guardado antes (ej. re-guardado tras corregir datos),
                # borramos primero sus filas viejas en cada hoja para no duplicar — quedarse con
                # dos tandas mezcladas (una vieja incompleta + una nueva corregida) es peor que
                # sobreescribir, porque un análisis de riesgo/cobros usaría ambas a la vez.
                def _quitar_nota_vieja(df, col_nota="NOTA"):
                    if df is None or df.empty or col_nota not in df.columns:
                        return df
                    return df[pd.to_numeric(df[col_nota], errors="coerce") != numero_nota].copy()

                filas_viejas_control = int((pd.to_numeric(hojas["CONTROL_NOTAS"].get("NOTA"), errors="coerce") == numero_nota).sum()) if "NOTA" in hojas["CONTROL_NOTAS"].columns else 0
                if filas_viejas_control > 0:
                    st.info(f"Ya había {filas_viejas_control} fila(s) guardada(s) antes para la Nota {numero_nota} — se reemplazan por esta versión, no se duplican.")

                hojas["CONTROL_NOTAS"] = pd.concat([_quitar_nota_vieja(hojas["CONTROL_NOTAS"]), df_control_editado], ignore_index=True)
                hojas["CALENDARIO_NOTAS"] = pd.concat([_quitar_nota_vieja(hojas["CALENDARIO_NOTAS"]), df_cal_editado], ignore_index=True)
                if "CALENDARIO_CALLS" in hojas and not df_calls_editado.empty:
                    hojas["CALENDARIO_CALLS"] = pd.concat([_quitar_nota_vieja(hojas["CALENDARIO_CALLS"]), df_calls_editado], ignore_index=True)
                # Quitamos el borrador de esta misma tanda de hojas para no escribir el Excel dos veces
                if "BORRADORES_NOTAS" in hojas and not hojas["BORRADORES_NOTAS"].empty:
                    df_b = hojas["BORRADORES_NOTAS"]
                    hojas["BORRADORES_NOTAS"] = df_b[~((df_b["TIPO"] == "nueva") & (pd.to_numeric(df_b["NOTA"], errors="coerce") == numero_nota))]

                # ── Inversor(es) de esta nota, si se rellenaron campos ──────────────────
                filas_inv_validas = [f for f in filas_inversores_nota if str(f.get("inversor", "")).strip() and str(f.get("id_inversion", "")).strip() and float(f.get("capital_invertido") or 0) > 0]
                mensaje_inversores = ""
                if filas_inv_validas:
                    if "INVERSIONES" not in hojas:
                        st.warning("⚠️ No se encontró la hoja INVERSIONES — la nota se guardó, pero los inversores NO se pudieron registrar. Usa '➕ Nueva inversión' para añadirlos a mano.")
                    else:
                        mapa_inv = _mapa_columnas_reales(hojas, "INVERSIONES")
                        if not mapa_inv:
                            # Plan B: las columnas reales de INVERSIONES ya son minúsculas/snake_case
                            # (confirmado en el propio Excel) — usamos esos nombres directamente en
                            # vez de depender de la detección automática, que aquí está fallando.
                            _df_inv_diag = hojas.get("INVERSIONES")
                            with st.expander("🔍 Detalle técnico (para diagnosticar más adelante)"):
                                if _df_inv_diag is None:
                                    st.write("hojas['INVERSIONES'] es None.")
                                else:
                                    st.write(f"Tipo: {type(_df_inv_diag)}")
                                    st.write(f"Forma (filas, columnas): {getattr(_df_inv_diag, 'shape', 'sin atributo shape')}")
                                    st.write(f"Columnas encontradas: {list(getattr(_df_inv_diag, 'columns', []))}")
                            columnas_esperadas = ["id_inversion", "inversor", "tipo_inversion", "subtipo_inversion",
                                                   "nombre_activo", "metodo_calculo", "cuenta_cobro", "activo_generador_interes",
                                                   "fecha_inversion", "fecha_final_inversion", "motivo", "capital_invertido",
                                                   "interes_nota_anual", "interes_inversor_anual", "tipo_operacion",
                                                   "id_inversion_origen", "capital_nuevo_real", "email", "pago_intereses"]
                            mapa_inv = {c: c for c in columnas_esperadas}
                        filas_reales_inv = []
                        columnas_faltantes_inv = set()
                        for fila_logica in filas_inv_validas:
                            fila_real = {}
                            for clave_logica, valor in fila_logica.items():
                                col_real = mapa_inv.get(clave_logica)
                                if col_real:
                                    fila_real[col_real] = valor
                                else:
                                    columnas_faltantes_inv.add(clave_logica)
                            filas_reales_inv.append(fila_real)
                        hojas["INVERSIONES"] = pd.concat([hojas["INVERSIONES"], pd.DataFrame(filas_reales_inv)], ignore_index=True)
                        mensaje_inversores = f" y {len(filas_reales_inv)} inversor(es) en INVERSIONES"
                        if columnas_faltantes_inv:
                            st.warning(f"⚠️ Estas columnas no existen tal cual en INVERSIONES y no se escribieron: {', '.join(sorted(columnas_faltantes_inv))}.")

                        # Vínculo adicional en la hoja histórica REINVERSIONES (trazabilidad, igual
                        # que en '➕ Nueva inversión') — soporta varios orígenes por fila (separados
                        # por coma en id_inversion_origen), escribiendo un enlace por cada uno.
                        if tipo_operacion == "REINVERSION" and "REINVERSIONES" in hojas:
                            mapa_reinv = _mapa_columnas_reales(hojas, "REINVERSIONES")
                            col_origen = mapa_reinv.get("id_inversion_origen")
                            col_destino = mapa_reinv.get("id_inversion_destino")
                            if col_origen and col_destino:
                                filas_reinv = []
                                for fila_logica in filas_inv_validas:
                                    ids_origen_f = str(fila_logica.get("id_inversion_origen") or "").strip()
                                    if not ids_origen_f:
                                        continue
                                    for id_origen_uno in [x.strip() for x in ids_origen_f.split(",") if x.strip()]:
                                        fila_reinv = {col_origen: id_origen_uno, col_destino: fila_logica.get("id_inversion", "")}
                                        if mapa_reinv.get("fecha"):
                                            fila_reinv[mapa_reinv["fecha"]] = pd.Timestamp(fila_logica.get("fecha_inversion"))
                                        if mapa_reinv.get("importe"):
                                            fila_reinv[mapa_reinv["importe"]] = fila_logica.get("capital_invertido")
                                        filas_reinv.append(fila_reinv)
                                if filas_reinv:
                                    hojas["REINVERSIONES"] = pd.concat([hojas["REINVERSIONES"], pd.DataFrame(filas_reinv)], ignore_index=True)
                            else:
                                st.info("No se encontraron las columnas id_inversion_origen/id_inversion_destino en REINVERSIONES — añade el vínculo a mano si hace falta.")

                guardar_excel_completo_desde_hojas(hojas)

                # ── Guardar el PDF en el volumen persistente ─────────────────────────────
                pdf_bytes_guardar = st.session_state.get("nota_pdf_bytes", {}).get(numero_nota)
                if pdf_bytes_guardar:
                    ruta_pdf_guardada = guardar_pdf_nota(numero_nota, pdf_bytes_guardar)
                    if ruta_pdf_guardada:
                        st.caption(f"📎 PDF original guardado en el servidor ({ruta_pdf_guardada}).")

                del almacen[numero_nota]
                st.session_state[clave_num_inv_nota] = 1
                st.success(f"Nota {numero_nota} guardada en CONTROL_NOTAS, CALENDARIO_NOTAS y CALENDARIO_CALLS{mensaje_inversores}.")
    with col_descartar:
        if st.button(f"🗑️ Descartar previsualización de la Nota {numero_nota}"):
            del almacen[numero_nota]
            borrar_borrador_nota("nueva", numero_nota)
            st.rerun()


def _tab_auditar_nota(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, df_calls: pd.DataFrame):
    st.caption(
        "Sube el documento oficial de una nota que ya está en el Excel. La IA lo lee de nuevo y Python "
        "compara, campo a campo, lo que dice el PDF contra lo que hay guardado en CONTROL_NOTAS, "
        "CALENDARIO_NOTAS y CALENDARIO_CALLS — así detectas tipos de interés, fechas de cobro, "
        "observación o call mal cargados, sin que la IA tenga que comparar números de memoria."
    )

    notas_existentes = sorted(
        int(n) for n in pd.to_numeric(df_control.get("nota", pd.Series(dtype=float)), errors="coerce").dropna().unique()
    ) if df_control is not None and not df_control.empty else []

    if not notas_existentes:
        st.info("No hay notas en CONTROL_NOTAS todavía para auditar.")
        return

    # --- Leer auditoría sin PDF: lo que ya quedó guardado con "Guardar esta auditoría" ---
    # Al guardar la auditoría se borra el borrador (para no dejar dos copias del mismo dato), así
    # que la única forma de volver a verla es leerla de vuelta de AUDITORIA_NOTAS, sin PDF de nuevo.
    with st.expander("📂 Leer auditoría sin PDF (ya guardada en el Excel)", expanded=False):
        nota_consulta = st.selectbox("Nota a consultar", notas_existentes, key="consultar_auditoria_numero")
        df_audit = _leer_auditoria_notas_cached()
        if df_audit.empty or "nota" not in df_audit.columns:
            st.caption("Todavía no se ha guardado ninguna auditoría para ninguna nota.")
        else:
            audit_n = df_audit[pd.to_numeric(df_audit["nota"], errors="coerce") == nota_consulta].copy()
            if audit_n.empty:
                st.caption(f"Todavía no se ha guardado ninguna auditoría para la Nota {nota_consulta}.")
            else:
                audit_n["fecha_auditoria"] = pd.to_datetime(audit_n.get("fecha_auditoria"), errors="coerce")
                ultima_fecha = audit_n["fecha_auditoria"].max()
                audit_ultima = audit_n[audit_n["fecha_auditoria"] == ultima_fecha]
                st.caption(f"Auditada el {ultima_fecha.strftime('%d/%m/%Y %H:%M') if pd.notna(ultima_fecha) else '(fecha desconocida)'}")
                cols_audit = [c for c in ["campo", "en_excel", "en_pdf", "estado"] if c in audit_ultima.columns]
                st.dataframe(audit_ultima[cols_audit], use_container_width=True, hide_index=True)
        st.caption("👁️ Solo lectura. Para volver a comparar contra el PDF oficial y detectar nuevas discrepancias, usa el auditor de abajo (necesita subir el PDF otra vez).")

    st.markdown("---")
    numero_nota = st.selectbox("Nota a auditar", notas_existentes, key="auditar_nota_numero")

    pdf_guardado_bytes = leer_pdf_nota_guardado(numero_nota)
    usar_pdf_guardado = False
    if pdf_guardado_bytes:
        usar_pdf_guardado = st.checkbox(
            f"📎 Usar el PDF ya guardado en el servidor para la Nota {numero_nota} (no hace falta volver a subirlo)",
            value=True, key=f"auditar_usar_guardado_{numero_nota}",
        )
    if usar_pdf_guardado:
        pdf_subido = None
        pdf_bytes_para_auditar = pdf_guardado_bytes
        st.caption("Usando el PDF guardado en el servidor.")
    else:
        pdf_subido = st.file_uploader(f"Documento oficial de la Nota {numero_nota} (PDF)", type=["pdf"], key=f"auditar_pdf_{numero_nota}")
        pdf_bytes_para_auditar = pdf_subido.read() if pdf_subido is not None else None

    if "auditoria_extraidos" not in st.session_state:
        st.session_state["auditoria_extraidos"] = {}
    almacen_auditoria = st.session_state["auditoria_extraidos"]

    if numero_nota not in almacen_auditoria:
        borrador = cargar_borrador_nota("auditar", numero_nota)
        if borrador:
            almacen_auditoria[numero_nota] = borrador
            st.info(f"📂 Se recuperó una auditoría en curso para la Nota {numero_nota} (guardada automáticamente).")

    if st.button("🔍 Auditar nota", type="primary", disabled=pdf_bytes_para_auditar is None):
        with st.spinner("Leyendo el documento y comparando contra el Excel..."):
            resultado = _ejecutar_auditoria_ia(numero_nota, pdf_bytes_para_auditar)
        if "error" in resultado:
            st.error(f"No se pudieron extraer los datos del PDF: {resultado['error']}")
            return

    _render_comparacion_auditoria(numero_nota, df_inv, df_cal, df_control, df_calls)


def _ejecutar_auditoria_ia(numero_nota: int, pdf_bytes: bytes) -> dict:
    """Llama a la IA para extraer los datos del PDF de una nota, guarda el resultado en el
    almacén de sesión 'auditoria_extraidos' (para que _render_comparacion_auditoria lo pinte) y
    en el borrador automático de Excel. Reutilizada tanto por el botón manual de la pestaña
    'Auditar nota' como por el guardado automático desde 'PDFs de notas'. Devuelve el dict de
    la IA tal cual (con clave 'error' si algo falló)."""
    if "auditoria_extraidos" not in st.session_state:
        st.session_state["auditoria_extraidos"] = {}
    resultado = extraer_datos_nota_con_ia(pdf_bytes)
    if "error" not in resultado:
        st.session_state["auditoria_extraidos"][numero_nota] = resultado
        guardar_borrador_nota("auditar", numero_nota, resultado)
    return resultado


def _render_comparacion_auditoria(numero_nota: int, df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, df_calls: pd.DataFrame):
    """Pinta la comparación campo a campo (PDF vs Excel) y la validación económica para una
    nota, a partir de lo que ya haya en st.session_state['auditoria_extraidos'][numero_nota].
    No hace nada si todavía no hay ninguna extracción guardada para esa nota."""
    almacen_auditoria = st.session_state.get("auditoria_extraidos", {})
    extraido = almacen_auditoria.get(numero_nota)
    if not extraido:
        return

    st.markdown("---")
    etiquetas = []
    if extraido.get("tiene_memoria") is True:
        etiquetas.append("🧠 Memoria")
    if extraido.get("tiene_one_star") is True:
        etiquetas.append("⭐ One-star")
    if etiquetas:
        st.markdown("  |  ".join(f"**{e}**" for e in etiquetas))

    def _es_revisar(v):
        return v is None or str(v).strip().upper() == "REVISAR"

    filas_diff = []

    def _comparar(campo, valor_excel, valor_pdf, tolerancia=None):
        if _es_revisar(valor_pdf):
            estado = "⚠️ IA no pudo leerlo del PDF"
        elif valor_excel is None or (isinstance(valor_excel, float) and pd.isna(valor_excel)):
            estado = "ℹ️ No hay dato en Excel"
        else:
            try:
                ve, vp = float(valor_excel), float(valor_pdf)
                coincide = abs(ve - vp) <= (tolerancia or 0.001)
            except (TypeError, ValueError):
                coincide = str(valor_excel).strip().upper() == str(valor_pdf).strip().upper()
            estado = "✅ Coincide" if coincide else "❌ DISCREPANCIA"
        filas_diff.append({"Campo": campo, "En Excel": valor_excel, "En el PDF": valor_pdf, "Estado": estado})

    # --- Tickers y barreras (CONTROL_NOTAS) ---
    control_nota = df_control[pd.to_numeric(df_control["nota"], errors="coerce") == numero_nota].copy() if "nota" in df_control.columns else pd.DataFrame()
    for t in extraido.get("tickers", []):
        ticker = str(t.get("ticker", "")).strip().upper()
        fila_excel = control_nota[control_nota.get("ticker", pd.Series(dtype=str)).astype(str).str.upper() == ticker]
        if fila_excel.empty:
            filas_diff.append({"Campo": f"{ticker}: presencia", "En Excel": "No existe", "En el PDF": ticker, "Estado": "❌ DISCREPANCIA — ticker no está en CONTROL_NOTAS"})
            continue
        fila_excel = fila_excel.iloc[0]
        _comparar(f"{ticker}: barrera cupón", fila_excel.get("barrera_cupon"), t.get("barrera_cupon_pct"))
        col_bc = "barrera_capital" if "barrera_capital" in fila_excel.index else "contingency"
        _comparar(f"{ticker}: barrera capital", fila_excel.get(col_bc), t.get("barrera_capital_pct"))
        _comparar(f"{ticker}: call level", fila_excel.get("call_level"), t.get("call_level_pct"))

    # --- CUSIP (comparado contra CONTROL_NOTAS de esa nota) ---
    if not control_nota.empty:
        cusip_excel = control_nota.iloc[0].get("cusip") if "cusip" in control_nota.columns else None
        cusip_excel = cusip_excel if cusip_excel and str(cusip_excel).strip() else None
        _comparar("CUSIP", cusip_excel, extraido.get("cusip"))

    # --- Cupón anual (comparado contra INVERSIONES de esa nota) ---
    nombre_nota = f"NOTA_{numero_nota:02d}"
    inv_nota = df_inv[df_inv.get("nombre_activo", pd.Series(dtype=str)).astype(str).str.upper().str.replace(" ", "_") == nombre_nota] if "nombre_activo" in df_inv.columns else pd.DataFrame()
    if not inv_nota.empty and "interes_nota_anual" in inv_nota.columns:
        tasas_excel = inv_nota["interes_nota_anual"].dropna().unique()
        tasa_excel = float(tasas_excel[0]) if len(tasas_excel) else None
        _comparar("Cupón anual", tasa_excel, extraido.get("cupon_anual_pct"))

    # --- Calendario de observación/pago (CALENDARIO_NOTAS) ---
    cal_nota = df_cal[df_cal.get("nota") == numero_nota].copy() if "nota" in df_cal.columns else pd.DataFrame()
    obs_excel = set(pd.to_datetime(cal_nota[cal_nota["tipo_evento"] == "OBSERVACION"]["fecha"], errors="coerce").dt.strftime("%Y-%m-%d").dropna()) if not cal_nota.empty else set()
    pago_excel = set(pd.to_datetime(cal_nota[cal_nota["tipo_evento"] == "PAGO"]["fecha"], errors="coerce").dt.strftime("%Y-%m-%d").dropna()) if not cal_nota.empty else set()
    obs_pdf = set(e.get("observacion") for e in extraido.get("calendario", []) if e.get("observacion") and not _es_revisar(e.get("observacion")))
    pago_pdf = set(e.get("pago") for e in extraido.get("calendario", []) if e.get("pago") and not _es_revisar(e.get("pago")))

    faltan_obs_excel = sorted(obs_pdf - obs_excel)
    sobran_obs_excel = sorted(obs_excel - obs_pdf)
    faltan_pago_excel = sorted(pago_pdf - pago_excel)
    sobran_pago_excel = sorted(pago_excel - pago_pdf)

    if faltan_obs_excel or sobran_obs_excel or faltan_pago_excel or sobran_pago_excel:
        if faltan_obs_excel:
            filas_diff.append({"Campo": "Fechas OBSERVACION en el PDF pero no en Excel", "En Excel": "—", "En el PDF": ", ".join(faltan_obs_excel), "Estado": "❌ DISCREPANCIA"})
        if sobran_obs_excel:
            filas_diff.append({"Campo": "Fechas OBSERVACION en Excel pero no en el PDF", "En Excel": ", ".join(sobran_obs_excel), "En el PDF": "—", "Estado": "❌ DISCREPANCIA"})
        if faltan_pago_excel:
            filas_diff.append({"Campo": "Fechas PAGO en el PDF pero no en Excel", "En Excel": "—", "En el PDF": ", ".join(faltan_pago_excel), "Estado": "❌ DISCREPANCIA"})
        if sobran_pago_excel:
            filas_diff.append({"Campo": "Fechas PAGO en Excel pero no en el PDF", "En Excel": ", ".join(sobran_pago_excel), "En el PDF": "—", "Estado": "❌ DISCREPANCIA"})
    else:
        filas_diff.append({"Campo": "Calendario observación/pago", "En Excel": f"{len(obs_excel)} obs / {len(pago_excel)} pagos", "En el PDF": f"{len(obs_pdf)} obs / {len(pago_pdf)} pagos", "Estado": "✅ Coincide"})

    # --- Fechas de call (CALENDARIO_CALLS) ---
    calls_nota = df_calls[df_calls.get("nota") == numero_nota].copy() if df_calls is not None and not df_calls.empty and "nota" in df_calls.columns else pd.DataFrame()
    calls_excel = set(pd.to_datetime(calls_nota["fecha_call"], errors="coerce").dt.strftime("%Y-%m-%d").dropna()) if not calls_nota.empty else set()
    calls_pdf = set(f for f in extraido.get("fechas_call", []) if f and not _es_revisar(f))
    faltan_calls_excel = sorted(calls_pdf - calls_excel)
    sobran_calls_excel = sorted(calls_excel - calls_pdf)
    if faltan_calls_excel or sobran_calls_excel:
        if faltan_calls_excel:
            filas_diff.append({"Campo": "Fechas de CALL en el PDF pero no en Excel", "En Excel": "—", "En el PDF": ", ".join(faltan_calls_excel), "Estado": "❌ DISCREPANCIA"})
        if sobran_calls_excel:
            filas_diff.append({"Campo": "Fechas de CALL en Excel pero no en el PDF", "En Excel": ", ".join(sobran_calls_excel), "En el PDF": "—", "Estado": "❌ DISCREPANCIA"})
    else:
        filas_diff.append({"Campo": "Fechas de posible call", "En Excel": f"{len(calls_excel)} fechas", "En el PDF": f"{len(calls_pdf)} fechas", "Estado": "✅ Coincide"})

    df_diff = pd.DataFrame(filas_diff)
    n_discrepancias = int(df_diff["Estado"].astype(str).str.startswith("❌").sum())
    if n_discrepancias > 0:
        st.error(f"Se encontraron {n_discrepancias} discrepancia(s) entre el Excel y el documento oficial.")
    else:
        st.success("No se encontraron discrepancias entre el Excel y el documento oficial.")

    st.dataframe(df_diff, use_container_width=True, hide_index=True)

    # --- Validación económica: capital, tasa, cobro, pago a inversores y beneficio ---
    st.markdown("---")
    st.markdown("#### Validación económica (capital, cobro, pago a inversores y beneficio)")

    fecha_ref = proximo_evento_nota(df_cal, numero_nota, "PAGO")
    if fecha_ref is None:
        pagos_pasados = df_cal[(df_cal.get("nota") == numero_nota) & (df_cal.get("tipo_evento") == "PAGO") & (df_cal["fecha"].notna())].sort_values("fecha") if "nota" in df_cal.columns else pd.DataFrame()
        fecha_ref = pagos_pasados.iloc[-1]["fecha"] if not pagos_pasados.empty else None

    capital_total_guardar = cobro_total_guardar = pago_total_guardar = beneficio_total_guardar = None

    if fecha_ref is None:
        st.info("No hay fechas de pago en el calendario de esta nota para poder validar el cálculo económico.")
    else:
        activas = inversiones_activas_para_nota(df_inv, numero_nota, fecha_ref)
        if activas.empty:
            st.info(f"No hay inversores activos en la Nota {numero_nota} en la fecha {pd.Timestamp(fecha_ref).strftime('%d/%m/%Y')}.")
        else:
            capital_total = float(activas["capital_invertido"].sum())
            capital_total_guardar = capital_total
            tasas_nota_distintas = sorted(activas["interes_nota_anual"].dropna().unique())
            st.markdown(f"**Capital total invertido:** {fmt(capital_total)} ({len(activas)} posiciones) — fecha de referencia: {pd.Timestamp(fecha_ref).strftime('%d/%m/%Y')}")

            if len(tasas_nota_distintas) > 1:
                st.warning(f"⚠️ Hay más de una tasa de nota distinta entre los inversores de esta nota en INVERSIONES: {[f'{t*100:.3f}%' for t in tasas_nota_distintas]}. Debería ser una sola.")

            cupon_pdf = extraido.get("cupon_anual_pct")
            if tasas_nota_distintas and not _es_revisar(cupon_pdf):
                for tasa in tasas_nota_distintas:
                    if abs(float(tasa) - float(cupon_pdf)) > 0.001:
                        st.error(f"❌ La tasa de nota en INVERSIONES ({float(tasa)*100:.3f}%) no coincide con el cupón anual del PDF ({float(cupon_pdf)*100:.3f}%).")

            # Recalcula cobro/pago/beneficio con la MISMA lógica que usa el Dashboard,
            # para validar de extremo a extremo (incluye la periodicidad detectada y la evaluación de barrera).
            df_pagos_evento = pd.DataFrame([{"nota": numero_nota, "tipo_evento": "PAGO", "fecha": pd.Timestamp(fecha_ref)}])
            detalle = preparar_detalle_notas(df_inv, df_pagos_evento, df_cal=df_cal, df_control=df_control)
            detalle_nota = detalle[detalle["nota"] == numero_nota] if not detalle.empty else pd.DataFrame()
            if detalle_nota.empty:
                st.info("No se pudo calcular el desglose económico para esta fecha.")
            else:
                cols_mostrar = ["inversor", "capital_invertido", "interes_nota_anual", "interes_inversor_anual", "cobro_compania", "pago_inversor", "beneficio_empresa", "resultado_observacion"]
                cols_mostrar = [c for c in cols_mostrar if c in detalle_nota.columns]
                st.dataframe(
                    preparar_tabla_monetaria(detalle_nota[cols_mostrar], ["capital_invertido", "cobro_compania", "pago_inversor", "beneficio_empresa"]),
                    use_container_width=True, hide_index=True,
                )
                st.markdown(
                    f"**Totales calculados para el {pd.Timestamp(fecha_ref).strftime('%d/%m/%Y')}:** "
                    f"Cobro compañía {fmt(detalle_nota['cobro_compania'].sum())} · "
                    f"Pago a inversores {fmt(detalle_nota['pago_inversor'].sum())} · "
                    f"Beneficio {fmt(detalle_nota['beneficio_empresa'].sum())}"
                )
                capital_total_guardar = capital_total
                cobro_total_guardar = float(detalle_nota["cobro_compania"].sum())
                pago_total_guardar = float(detalle_nota["pago_inversor"].sum())
                beneficio_total_guardar = float(detalle_nota["beneficio_empresa"].sum())
    st.markdown("---")
    if st.button(f"💾 Guardar esta auditoría de la Nota {numero_nota}", type="primary"):
        ahora = pd.Timestamp.now()
        filas_guardar = []
        for fila in filas_diff:
            filas_guardar.append({
                "FECHA_AUDITORIA": ahora,
                "NOTA": numero_nota,
                "CAMPO": fila["Campo"],
                "EN_EXCEL": fila["En Excel"],
                "EN_PDF": fila["En el PDF"],
                "ESTADO": fila["Estado"],
            })
        filas_guardar.append({
            "FECHA_AUDITORIA": ahora,
            "NOTA": numero_nota,
            "CAMPO": "RESUMEN — Memoria / One-star",
            "EN_EXCEL": "",
            "EN_PDF": ", ".join(etiquetas) if etiquetas else "Ninguna",
            "ESTADO": "ℹ️ Informativo",
        })
        filas_guardar.append({
            "FECHA_AUDITORIA": ahora,
            "NOTA": numero_nota,
            "CAMPO": "RESUMEN — Capital / Cobro / Pago / Beneficio",
            "EN_EXCEL": "",
            "EN_PDF": (
                f"Capital {fmt(capital_total_guardar)} · Cobro {fmt(cobro_total_guardar)} · "
                f"Pago inv. {fmt(pago_total_guardar)} · Beneficio {fmt(beneficio_total_guardar)}"
                if capital_total_guardar is not None else "No calculado"
            ),
            "ESTADO": "ℹ️ Informativo",
        })

        hojas = leer_todas_las_hojas_excel()
        df_nuevo = pd.DataFrame(filas_guardar)
        if "AUDITORIA_NOTAS" in hojas and not hojas["AUDITORIA_NOTAS"].empty:
            hojas["AUDITORIA_NOTAS"] = pd.concat([hojas["AUDITORIA_NOTAS"], df_nuevo], ignore_index=True)
        else:
            hojas["AUDITORIA_NOTAS"] = df_nuevo
        if "BORRADORES_NOTAS" in hojas and not hojas["BORRADORES_NOTAS"].empty:
            df_b = hojas["BORRADORES_NOTAS"]
            hojas["BORRADORES_NOTAS"] = df_b[~((df_b["TIPO"] == "auditar") & (pd.to_numeric(df_b["NOTA"], errors="coerce") == numero_nota))]
        guardar_excel_completo_desde_hojas(hojas)
        del almacen_auditoria[numero_nota]
        st.success(f"Auditoría de la Nota {numero_nota} guardada en la hoja AUDITORIA_NOTAS.")


def _tab_ficha_compania(df_control: pd.DataFrame):
    st.caption(
        "Análisis fundamental de una compañía subyacente: precio, rango de 52 semanas, PER, "
        "consenso de analistas, volatilidad, próximo earnings, gráfico de evolución y noticias "
        "recientes que avalan el análisis. **Todo dato de mercado es real (Yahoo Finance)** — la "
        "única parte generada por IA es el resumen de noticias, síntesis de fuentes públicas, "
        "nunca una predicción de precio."
    )

    tickers_disponibles = sorted(df_control["ticker"].dropna().unique()) if df_control is not None and not df_control.empty and "ticker" in df_control.columns else []
    col1, col2 = st.columns([2, 1])
    with col1:
        ticker_elegido = st.selectbox("Ticker a analizar (de tus notas actuales)", tickers_disponibles) if tickers_disponibles else None
    with col2:
        ticker_manual = st.text_input("...o escribe cualquier otro ticker")
    ticker = (ticker_manual.strip().upper() if ticker_manual.strip() else ticker_elegido)

    if not ticker:
        st.info("Elige o escribe un ticker para analizar.")
        return

    if st.button(f"🔎 Analizar {ticker}", type="primary"):
        with st.spinner(f"Consultando datos de mercado de {ticker}..."):
            datos = obtener_datos_fundamentales(ticker)
        st.session_state[f"ficha_{ticker}"] = datos
        st.session_state.pop(f"ficha_{ticker}_noticias", None)  # fuerza rebuscar noticias tras un reanálisis

    datos = st.session_state.get(f"ficha_{ticker}")
    if not datos:
        return

    if datos.get("error") and not datos.get("precio_actual"):
        st.error(f"No se pudieron obtener datos de {ticker}: {datos['error']}")
        return

    # --- Notas actuales que incluyen este ticker (se calcula antes para poder marcar la barrera en el gráfico) ---
    notas_con_ticker = pd.DataFrame()
    if df_control is not None and not df_control.empty and "ticker" in df_control.columns:
        notas_con_ticker = df_control[df_control["ticker"] == ticker].copy()

    precio_contingencia_peor = None
    if not notas_con_ticker.empty and datos.get("precio_actual"):
        margenes = []
        for _, r in notas_con_ticker.iterrows():
            barrera_cupon, precio_compra = r.get("barrera_cupon"), r.get("precio_compra")
            if pd.notna(barrera_cupon) and pd.notna(precio_compra) and precio_compra:
                pc = precio_compra * barrera_cupon
                margenes.append((datos["precio_actual"] / pc - 1, pc))
        if margenes:
            precio_contingencia_peor = min(margenes, key=lambda x: x[0])[1]

    st.markdown("---")
    st.markdown(_tarjeta_subyacente_html(datos, notas_con_ticker), unsafe_allow_html=True)

    if datos.get("aviso_analistas"):
        st.info(f"ℹ️ {datos['aviso_analistas']}")
    if datos.get("error"):
        st.error(f"⚠️ {datos['error']}")
        if st.button("🔄 Reintentar", key=f"retry_{ticker}"):
            obtener_datos_fundamentales.clear()
            st.rerun()

    st.markdown("#### 📈 Evolución del precio (12 meses)")
    grafico_precio_subyacente(datos, precio_contingencia=precio_contingencia_peor)
    if precio_contingencia_peor:
        st.caption(f"Línea roja: barrera de contingencia de la nota con menos margen (${precio_contingencia_peor:,.2f}).")

    c1, c2 = st.columns(2)
    with c1:
        tarjeta_kpi(
            "Precio objetivo (consenso)",
            f"${datos['target_medio']:,.2f}" if datos.get("target_medio") else "N/D",
            (f"{(datos['target_medio'] / datos['precio_actual'] - 1) * 100:+.1f}% vs precio actual"
             if datos.get("target_medio") and datos.get("precio_actual") else "Ver aviso arriba"),
            "positivo" if datos.get("target_medio") and datos.get("precio_actual") and datos["target_medio"] >= datos["precio_actual"] else "normal",
        )
    with c2:
        tarjeta_kpi(
            "Volatilidad anual histórica", f"{datos['volatilidad_anual_pct']:.1f}%" if datos.get("volatilidad_anual_pct") else "N/D",
            "Últimos 12 meses", "riesgo" if (datos.get("volatilidad_anual_pct") or 0) > 40 else "normal",
        )

    if datos.get("target_alto") or datos.get("target_bajo"):
        st.caption(f"Rango de analistas: ${datos.get('target_bajo', 0):,.2f} — ${datos.get('target_alto', 0):,.2f}  |  {datos.get('n_analistas', '?')} analistas  |  Recomendación consenso: {datos.get('recomendacion', 'N/D')}")
    if datos.get("proxima_fecha_resultados"):
        st.warning(f"📅 Próxima fecha de resultados (earnings): {datos['proxima_fecha_resultados']} — la volatilidad suele dispararse alrededor de esta fecha.")

    # Notas actuales que incluyen este ticker, y su posición respecto a las barreras
    if not notas_con_ticker.empty and datos.get("precio_actual"):
        st.markdown("#### Tus notas con este ticker")
        filas = []
        for _, r in notas_con_ticker.iterrows():
            precio_compra = r.get("precio_compra")
            if pd.isna(precio_compra) or not precio_compra:
                continue
            pct_vs_inicial = (datos["precio_actual"] / precio_compra - 1) * 100
            barrera_cupon = r.get("barrera_cupon")
            margen_barrera = (datos["precio_actual"] / (precio_compra * barrera_cupon) - 1) * 100 if pd.notna(barrera_cupon) else None
            filas.append({
                "Nota": int(r.get("nota")) if pd.notna(r.get("nota")) else "?",
                "Precio compra": f"${precio_compra:,.2f}",
                "Precio actual": f"${datos['precio_actual']:,.2f}",
                "vs Inicial": f"{pct_vs_inicial:+.1f}%",
                "Margen a barrera cupón": f"{margen_barrera:+.1f}%" if margen_barrera is not None else "N/D",
            })
        if filas:
            st.dataframe(pd.DataFrame(filas), use_container_width=True, hide_index=True)

    # Noticias recientes vía Claude con búsqueda web — feed estructurado, no bloque de texto,
    # cada noticia con su fuente, fecha y enlace directo, igual que el resto de la app.
    st.markdown("#### 📰 Noticias y contexto reciente")
    clave_noticias = f"ficha_{ticker}_noticias"
    if clave_noticias not in st.session_state:
        with st.spinner("Buscando noticias recientes que avalen el análisis..."):
            st.session_state[clave_noticias] = buscar_noticias_libre(f"{datos.get('nombre') or ticker} ({ticker})")
    noticias_estructuradas = st.session_state[clave_noticias]
    _renderizar_tarjetas_noticias(noticias_estructuradas)

    # --- Informe PDF profesional para socios, con tarjeta, gráfico y noticias incluidas ---
    st.markdown("---")
    with st.spinner("Preparando informe PDF..."):
        try:
            pdf_bytes = _generar_informe_subyacente_pdf(datos, noticias_estructuradas, notas_con_ticker)
            st.download_button(
                f"⬇️ Descargar informe PDF de {ticker}", data=pdf_bytes,
                file_name=f"informe_{ticker}_{pd.Timestamp.now().strftime('%Y%m%d')}.pdf",
                mime="application/pdf", type="primary",
            )
        except Exception as e:
            st.error(f"No se pudo generar el PDF: {e}")


@st.cache_data(show_spinner=False, ttl=21600)
def obtener_resumen_noticias_ia(ticker: str, nombre_compania: str) -> str:
    """Usa la API de Claude con la herramienta de búsqueda web para resumir noticias
    recientes relevantes de una compañía. Devuelve texto, nunca cifras de precio inventadas.
    Cacheado 6h: las noticias no cambian minuto a minuto y cada llamada cuesta tokens de API."""
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
            json={
                "model": "claude-sonnet-4-5",
                "max_tokens": 1400,
                "tools": [{"type": "web_search_20250305", "name": "web_search", "max_uses": 8}],
                "system": (
                    "Buscas y resumes noticias RECIENTES (últimos 1-2 meses) que puedan afectar al precio de una acción, "
                    "para un inversor en notas estructuradas ligadas a esa acción. Haz una búsqueda EXHAUSTIVA: "
                    "no te quedes con el primer resultado — cruza varias fuentes (prensa financiera, comunicados de la "
                    "propia compañía, notas de analistas) antes de concluir que no hay nada relevante. Cubre en concreto: "
                    "resultados financieros recientes (con cifras si las encuentras), cambios de rating o precio objetivo "
                    "de analistas, noticias regulatorias o legales, eventos corporativos relevantes (fusiones, lanzamientos "
                    "de producto, cambios de dirección), y cualquier profit warning o guidance revisado. "
                    "Sé conciso en la redacción final (máximo 8-10 líneas) pero exhaustivo en la búsqueda. "
                    "NUNCA inventes un precio objetivo ni una predicción — solo resume hechos y opiniones ya publicadas, citando la fuente. "
                    "Si de verdad no encuentras NADA relevante tras buscar en varias fuentes, dilo claramente en vez de inventar contenido."
                ),
                "messages": [{"role": "user", "content": f"Busca a fondo, cruzando varias fuentes, noticias recientes relevantes sobre {nombre_compania} ({ticker}) que puedan afectar su cotización."}],
            },
            timeout=60,
        )
        data = resp.json()

        # 1) Error a nivel de petición (API key sin permisos, rate limit, etc.)
        if data.get("type") == "error" or data.get("error"):
            msg = data.get("error", {}).get("message", str(data))
            return (f"⚠️ No se pudo buscar noticias: error de la API ({msg}). "
                    f"Si el error menciona permisos o la herramienta 'web_search', puede que la búsqueda web "
                    f"no esté activada para esta cuenta en la Consola de Anthropic (Settings → hay que habilitarla).")

        contenido = data.get("content", [])

        # 2) Error a nivel de herramienta (la petición general fue 200 OK, pero la búsqueda en sí falló)
        for bloque in contenido:
            if bloque.get("type") == "web_search_tool_result":
                resultado_bloque = bloque.get("content", {})
                if isinstance(resultado_bloque, dict) and resultado_bloque.get("type") == "web_search_tool_result_error":
                    codigo_error = resultado_bloque.get("error_code", "desconocido")
                    return (f"⚠️ La búsqueda web falló (código: {codigo_error}). Si es 'unavailable' o similar, "
                            f"probablemente la búsqueda web no está activada para esta cuenta en la Consola de Anthropic.")

        # 3) ¿Llegó a intentar buscar? Si no hay ni un solo intento de búsqueda, algo no está bien configurado.
        hubo_intento_busqueda = any(b.get("type") == "server_tool_use" and b.get("name") == "web_search" for b in contenido)

        texto = "".join(b.get("text", "") for b in contenido if b.get("type") == "text")
        if not texto.strip():
            if not hubo_intento_busqueda:
                return ("⚠️ La IA no llegó a intentar ninguna búsqueda web. Puede que la herramienta de búsqueda "
                        "web no esté activada para esta cuenta — revísalo en la Consola de Anthropic (Settings → Web search).")
            return "No se encontraron noticias recientes relevantes tras buscar."
        return texto.strip()
    except Exception as e:
        return f"⚠️ No se pudieron obtener noticias (error de conexión): {e}"


def _proyectar_dias_earnings(proxima_fecha_dt, meses_horizonte: int) -> list:
    """A partir de la próxima fecha de resultados CONOCIDA (dato real de Yahoo Finance), proyecta
    las siguientes citas trimestrales dentro del horizonte de la nota (~91 días entre informes).
    Solo la primera fecha es un dato real confirmado — las siguientes son una proyección razonable
    asumiendo periodicidad trimestral, que es como reportan la inmensa mayoría de las acciones US."""
    if proxima_fecha_dt is None:
        return []
    try:
        hoy = pd.Timestamp.now().normalize()
        primera = pd.Timestamp(proxima_fecha_dt)
        if primera.tz is not None:
            primera = primera.tz_localize(None)
        primera = primera.normalize()
        dias_primera = (primera - hoy).days
        if dias_primera < 0:
            dias_primera += 91  # si la fecha ya venció (dato desactualizado), salta a la siguiente proyectada
        dias = []
        d = dias_primera
        while d <= meses_horizonte * 30.44:
            if d > 0:
                dias.append(int(d))
            d += 91
        return dias
    except Exception:
        return []


def _generar_horario_eventos(meses_vencimiento: int, periodicidad: int, frecuencia_call: int | None = None) -> list:
    """Genera un calendario sintético de eventos cupón + call + vencimiento, en días desde hoy.

    El cupón se paga cada `periodicidad` meses. El call es una ventana DISTINTA del emisor
    (`frecuencia_call` meses, por defecto igual a la periodicidad si no se especifica) —
    en la práctica muchas notas pagan cupón mensual pero solo son callables cada 3 meses.
    Se excluye la primera y la última ventana de call, como es habitual en el mercado.
    """
    if frecuencia_call is None:
        frecuencia_call = periodicidad
    eventos = []
    mes = periodicidad
    while mes <= meses_vencimiento:
        dias = int(mes * 30.44)
        eventos.append({"dias": max(dias, 1), "tipo": "cupon"})
        mes += periodicidad

    fechas_call = []
    mes = frecuencia_call
    while mes <= meses_vencimiento:
        fechas_call.append(mes)
        mes += frecuencia_call
    for idx, mes_call in enumerate(fechas_call):
        es_primera = (idx == 0)
        es_ultima = (idx == len(fechas_call) - 1)
        if not es_primera and not es_ultima:
            eventos.append({"dias": max(int(mes_call * 30.44), 1), "tipo": "call"})

    eventos.append({"dias": int(meses_vencimiento * 30.44), "tipo": "vencimiento"})
    return eventos


def _informe_ia_a_pdf(titulo: str, texto_markdown: str) -> bytes:
    """Convierte el informe de texto/markdown generado por la IA (fichas de compañía +
    recomendación) a un PDF sencillo con reportlab, sin dependencias externas."""
    import re as _re
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import ParagraphStyle

    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm)
    story = []
    style_titulo = ParagraphStyle("titulo", fontName="Helvetica-Bold", fontSize=16, spaceAfter=10)
    style_h3 = ParagraphStyle("h3", fontName="Helvetica-Bold", fontSize=12, spaceBefore=10, spaceAfter=4, textColor=rl_colors.Color(13/255, 33/255, 55/255))
    style_h4 = ParagraphStyle("h4", fontName="Helvetica-Bold", fontSize=10.5, spaceBefore=8, spaceAfter=3, textColor=rl_colors.Color(26/255, 63/255, 92/255))
    style_normal = ParagraphStyle("normal", fontName="Helvetica", fontSize=9.5, leading=13, spaceAfter=6)
    style_caption = ParagraphStyle("caption", fontName="Helvetica-Oblique", fontSize=8, textColor=rl_colors.grey, spaceAfter=8)

    story.append(Paragraph(titulo, style_titulo))
    story.append(Paragraph(f"Generado el {pd.Timestamp.now().strftime('%d/%m/%Y %H:%M')}", style_caption))

    def _linea_a_html(linea):
        linea = _re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", linea)
        return linea

    for linea in texto_markdown.split("\n"):
        linea_limpia = linea.strip()
        if not linea_limpia:
            story.append(Spacer(1, 3 * mm))
            continue
        if linea_limpia.startswith("#### "):
            story.append(Paragraph(_linea_a_html(linea_limpia[5:]), style_h4))
        elif linea_limpia.startswith("### "):
            story.append(Paragraph(_linea_a_html(linea_limpia[4:]), style_h3))
        elif linea_limpia.startswith("- ") or linea_limpia.startswith("* "):
            story.append(Paragraph("• " + _linea_a_html(linea_limpia[2:]), style_normal))
        else:
            story.append(Paragraph(_linea_a_html(linea_limpia), style_normal))

    doc.build(story)
    return output.getvalue()


def _exportar_grafico_barras_png(categorias: list, valores: list, titulo: str, color: str = "#9A6B24",
                                  ylabel: str = "", pct: bool = True, colores_por_barra: list = None) -> bytes:
    """Barra simple (matplotlib, Agg) para embeber en PDF — usada tanto para el score comparado
    como para la rentabilidad neta esperada. Sin dependencias de sistema (ver nota sobre kaleido
    en _exportar_grafico_precio_png)."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(8.6, 3.0), dpi=150)
        colores = colores_por_barra or [color] * len(categorias)
        ax.bar(categorias, valores, color=colores)
        ax.set_title(titulo, fontsize=11, color="#0E2338", fontweight="bold", loc="left")
        if ylabel:
            ax.set_ylabel(ylabel, fontsize=9)
        ax.spines[["top", "right"]].set_visible(False)
        ax.spines[["left", "bottom"]].set_color("#cccccc")
        ax.grid(axis="y", color="#e3e7ee", linewidth=0.8)
        ax.tick_params(colors="#555555", labelsize=9)
        for i, v in enumerate(valores):
            etiqueta = f"{v:.1f}%" if pct else f"{v:.0f}"
            ax.annotate(etiqueta, (i, v), textcoords="offset points", xytext=(0, 4), ha="center", fontsize=8.5, color="#333")
        fig.tight_layout()
        buf = BytesIO()
        fig.savefig(buf, format="png", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()
    except Exception:
        return None


def _exportar_grafico_barras_agrupadas_png(resultados_notas: list, titulo: str = "Probabilidades clave por nota") -> bytes:
    """Barras agrupadas (prob. cupón / prob. call / prob. pérdida capital) por nota, para el PDF."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as _np

        nombres = [r["nombre"] for r in resultados_notas]
        cupon = [r["prob_cupon_media"] * 100 for r in resultados_notas]
        call = [r["prob_call_total"] * 100 for r in resultados_notas]
        perdida = [r["prob_perdida_capital"] * 100 for r in resultados_notas]

        x = _np.arange(len(nombres))
        w = 0.26
        fig, ax = plt.subplots(figsize=(8.6, 3.2), dpi=150)
        ax.bar(x - w, cupon, width=w, label="Prob. cupón/periodo", color="#0E7C5A")
        ax.bar(x, call, width=w, label="Prob. call (total)", color="#9A6B24")
        ax.bar(x + w, perdida, width=w, label="Prob. pérdida capital", color="#B03A2E")
        ax.set_xticks(x)
        ax.set_xticklabels(nombres, fontsize=9)
        ax.set_title(titulo, fontsize=11, color="#0E2338", fontweight="bold", loc="left")
        ax.spines[["top", "right"]].set_visible(False)
        ax.spines[["left", "bottom"]].set_color("#cccccc")
        ax.grid(axis="y", color="#e3e7ee", linewidth=0.8)
        ax.tick_params(colors="#555555", labelsize=9)
        ax.legend(fontsize=8, loc="upper center", bbox_to_anchor=(0.5, -0.12), ncol=3, frameon=False)
        fig.tight_layout()
        buf = BytesIO()
        fig.savefig(buf, format="png", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()
    except Exception:
        return None


def _generar_informe_comparador_pdf(resultados_notas: list, texto_recomendacion: str,
                                     capital_disponible: float, tasa_inversor_pct: float,
                                     bloques_informe: list) -> bytes:
    """Informe PDF del comparador de notas, con identidad CF Wealth: tabla de resultados, gráfico
    de score y de probabilidades comparadas, ficha + noticias reales (con fuente y fecha) de cada
    subyacente, y la recomendación de reparto razonada. Reemplaza al PDF anterior (texto plano sin
    gráficos) — la parte narrativa (bloques_informe) se sigue generando igual que antes."""
    import re as _re
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, HRFlowable
    from reportlab.lib.styles import ParagraphStyle

    NAVY = rl_colors.Color(14/255, 35/255, 56/255)
    GOLD = rl_colors.Color(154/255, 107/255, 36/255)
    GREY = rl_colors.Color(102/255, 112/255, 133/255)
    LINE = rl_colors.Color(227/255, 231/255, 238/255)

    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=A4, leftMargin=16 * mm, rightMargin=16 * mm, topMargin=16 * mm, bottomMargin=16 * mm)
    story = []

    style_eyebrow = ParagraphStyle("eyebrow", fontName="Helvetica-Bold", fontSize=8.5, textColor=GOLD, spaceAfter=2)
    style_titulo = ParagraphStyle("titulo", fontName="Helvetica-Bold", fontSize=18, textColor=NAVY, spaceAfter=2)
    style_sub = ParagraphStyle("sub", fontName="Helvetica", fontSize=9.5, textColor=GREY, spaceAfter=10)
    style_h3 = ParagraphStyle("h3", fontName="Helvetica-Bold", fontSize=12, spaceBefore=12, spaceAfter=5, textColor=NAVY)
    style_h4 = ParagraphStyle("h4", fontName="Helvetica-Bold", fontSize=10, spaceBefore=7, spaceAfter=3, textColor=rl_colors.Color(26/255, 63/255, 92/255))
    style_normal = ParagraphStyle("normal", fontName="Helvetica", fontSize=9, leading=12.5, spaceAfter=5)
    style_caption = ParagraphStyle("caption", fontName="Helvetica-Oblique", fontSize=7.5, textColor=GREY, spaceAfter=6)

    story.append(Paragraph("CF WEALTH · COMPARADOR DE NOTAS", style_eyebrow))
    story.append(Paragraph("Informe de comparación", style_titulo))
    story.append(Paragraph(
        f"Generado el {pd.Timestamp.now().strftime('%d/%m/%Y %H:%M')} · Capital a repartir: ${capital_disponible:,.2f} · "
        f"Tasa pagada al inversor: {tasa_inversor_pct*100:.2f}% anual · Documento interno de socios.", style_sub))
    story.append(HRFlowable(width="100%", thickness=0.6, color=LINE, spaceAfter=8))

    # --- Tabla resumen ---
    filas_tabla = [["Nota", "Score", "Cupón anual", "Prob. cupón", "Prob. pérdida cap.", "Rent. neta esp."]]
    for r in resultados_notas:
        filas_tabla.append([
            r["nombre"], f"{r['score']:.0f}/100", f"{r['cupon_anual']*100:.1f}%",
            f"{r['prob_cupon_media']*100:.1f}%", f"{r['prob_perdida_capital']*100:.1f}%",
            f"{r['rentabilidad_esperada_neta']*100:.2f}%",
        ])
    tabla = Table(filas_tabla, colWidths=[38*mm, 20*mm, 24*mm, 24*mm, 30*mm, 27*mm])
    tabla.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"), ("FONTSIZE", (0, 0), (-1, -1), 8.5),
        ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white), ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [rl_colors.white, rl_colors.Color(0.97, 0.97, 0.98)]),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6), ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("LINEBELOW", (0, 0), (-1, -1), 0.3, LINE), ("ALIGN", (1, 0), (-1, -1), "CENTER"),
    ]))
    story.append(tabla)

    # --- Gráficos comparativos ---
    png_score = _exportar_grafico_barras_png(
        [r["nombre"] for r in resultados_notas], [r["score"] for r in resultados_notas],
        "Score comparado (0-100)", color="#9A6B24", pct=False,
    )
    png_prob = _exportar_grafico_barras_agrupadas_png(resultados_notas)
    png_renta = _exportar_grafico_barras_png(
        [r["nombre"] for r in resultados_notas], [r["rentabilidad_esperada_neta"] * 100 for r in resultados_notas],
        "Rentabilidad neta esperada anualizada (%)",
        colores_por_barra=["#0E7C5A" if r["rentabilidad_esperada_neta"] >= 0 else "#B03A2E" for r in resultados_notas],
    )
    story.append(Paragraph("Comparativa visual", style_h3))
    for png in (png_score, png_prob, png_renta):
        if png:
            story.append(RLImage(BytesIO(png), width=170*mm, height=170*mm*3.0/8.6))
            story.append(Spacer(1, 3*mm))
    if not (png_score or png_prob or png_renta):
        story.append(Paragraph("Gráficos no disponibles en este PDF (falta matplotlib en el servidor).", style_caption))

    # --- Contenido narrativo: ficha de cada subyacente + noticias con fuente/fecha ---
    def _linea_a_html(linea):
        return _re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", linea)

    style_conclusion = ParagraphStyle("conclusion", fontName="Helvetica-Bold", fontSize=13, textColor=GOLD, spaceBefore=6, spaceAfter=10)

    story.append(Paragraph("Detalle por nota y compañía subyacente", style_h3))
    for linea in "\n\n".join(bloques_informe).split("\n"):
        linea_limpia = linea.strip()
        if not linea_limpia:
            story.append(Spacer(1, 2.5 * mm))
            continue
        if _re.match(r"^\**Conclusi[óo]n:", linea_limpia, _re.IGNORECASE):
            # La línea de conclusión (primera de la recomendación, ver prompt en el comparador)
            # se destaca en dorado y más grande para que ancle visualmente todo el análisis que
            # sigue — el mismo tratamiento que recibe en el PowerPoint, para que ambos documentos
            # transmitan el mismo mensaje de principio a fin, no solo en el párrafo final.
            story.append(Paragraph(_linea_a_html(linea_limpia.strip("*")), style_conclusion))
        elif linea_limpia.startswith("#### "):
            story.append(Paragraph(_linea_a_html(linea_limpia[5:]), style_h4))
        elif linea_limpia.startswith("### "):
            story.append(Paragraph(_linea_a_html(linea_limpia[4:]), style_h3))
        elif linea_limpia.startswith("- ") or linea_limpia.startswith("* "):
            story.append(Paragraph("• " + _linea_a_html(linea_limpia[2:]), style_normal))
        else:
            story.append(Paragraph(_linea_a_html(linea_limpia), style_normal))

    # --- Noticias que avalan el análisis, por nota y ticker ---
    story.append(Paragraph("Noticias que avalan el análisis", style_h3))
    for r in resultados_notas:
        noticias_ticker_dict = r.get("noticias_por_ticker") or {}
        if not noticias_ticker_dict:
            continue
        story.append(Paragraph(r["nombre"], style_h4))
        for ticker, noticias in noticias_ticker_dict.items():
            if not noticias or (isinstance(noticias[0], dict) and "_error" in noticias[0]):
                continue
            for n in noticias[:5]:
                prioridad = str(n.get("prioridad", "media")).upper()
                story.append(Paragraph(f"<b>[{ticker}] {n.get('titular', 'Sin título')}</b>  <font size=7 color='#9A6B24'>[{prioridad}]</font>", style_normal))
                story.append(Paragraph(f"{n.get('fuente', 'Fuente desconocida')} · {n.get('fecha', 'sin fecha')}", style_caption))
                if n.get("url"):
                    story.append(Paragraph(f"<link href='{n['url']}'><font color='#9A6B24'>{n['url']}</font></link>", style_caption))
    story.append(Paragraph(
        "Noticias generadas por IA a partir de búsqueda web pública. Documento interno — no constituye "
        "asesoramiento de inversión ni una recomendación de compra/venta.", style_caption))

    doc.build(story)
    return output.getvalue()


def _exportar_grafico_precio_png(datos: dict, precio_contingencia: float = None) -> bytes:
    """Genera el gráfico de evolución de precio como PNG para incrustarlo en el PDF, usando
    matplotlib (backend Agg, sin GUI y sin dependencias de sistema) en vez de plotly+kaleido:
    kaleido moderno necesita Chrome instalado en el servidor, algo frágil en un contenedor
    Railway — matplotlib genera la imagen de forma puramente local y siempre funciona igual."""
    historico = datos.get("historico_precios")
    if not historico:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.dates as mdates

        df_hist = pd.DataFrame(historico)
        df_hist["fecha"] = pd.to_datetime(df_hist["fecha"])

        fig, ax = plt.subplots(figsize=(9.8, 3.6), dpi=150)
        ax.plot(df_hist["fecha"], df_hist["close"], color="#9A6B24", linewidth=2)
        if precio_contingencia:
            ax.axhline(y=precio_contingencia, color="#B03A2E", linestyle="--", linewidth=1.3)
            ax.annotate(f"Barrera ${precio_contingencia:,.2f}", xy=(df_hist["fecha"].iloc[0], precio_contingencia),
                        xytext=(4, 4), textcoords="offset points", color="#B03A2E", fontsize=9)
        ax.spines[["top", "right"]].set_visible(False)
        ax.spines[["left", "bottom"]].set_color("#cccccc")
        ax.grid(axis="y", color="#e3e7ee", linewidth=0.8)
        ax.tick_params(colors="#555555", labelsize=9)
        ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
        ax.xaxis.set_major_formatter(mdates.DateFormatter("%b %y"))
        fig.tight_layout()

        buf = BytesIO()
        fig.savefig(buf, format="png", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()
    except Exception:
        return None  # el PDF se genera igual, solo sin gráfico, si algo falla


def _generar_informe_subyacente_pdf(datos: dict, noticias: list, notas_ticker: pd.DataFrame = None) -> bytes:
    """Informe PDF profesional de un subyacente para socios (Yuri/Alan/Jordi): ficha con precio,
    rango de 52 semanas y PER, consenso de analistas, gráfico de evolución con la barrera marcada,
    tus notas expuestas a este ticker, y las noticias reales que avalan el análisis (con fuente y
    fecha). Estilo marino/dorado consistente con el resto de informes de la app (comparador,
    memo de directorio)."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, HRFlowable
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_RIGHT

    NAVY = rl_colors.Color(14/255, 35/255, 56/255)
    GOLD = rl_colors.Color(154/255, 107/255, 36/255)
    GREY = rl_colors.Color(102/255, 112/255, 133/255)
    GREEN = rl_colors.Color(22/255, 101/255, 52/255)
    RED = rl_colors.Color(153/255, 27/255, 27/255)
    LINE = rl_colors.Color(227/255, 231/255, 238/255)

    ticker = datos.get("ticker", "")
    nombre = datos.get("nombre") or ticker
    precio = datos.get("precio_actual")

    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm, topMargin=16 * mm, bottomMargin=16 * mm)
    story = []

    style_eyebrow = ParagraphStyle("eyebrow", fontName="Helvetica-Bold", fontSize=8.5, textColor=GOLD, spaceAfter=2)
    style_titulo = ParagraphStyle("titulo", fontName="Helvetica-Bold", fontSize=20, textColor=NAVY, spaceAfter=2)
    style_sub = ParagraphStyle("sub", fontName="Helvetica", fontSize=10, textColor=GREY, spaceAfter=10)
    style_h3 = ParagraphStyle("h3", fontName="Helvetica-Bold", fontSize=12, spaceBefore=14, spaceAfter=6, textColor=NAVY)
    style_normal = ParagraphStyle("normal", fontName="Helvetica", fontSize=9.5, leading=13.5, spaceAfter=5, textColor=rl_colors.Color(0.15, 0.17, 0.2))
    style_caption = ParagraphStyle("caption", fontName="Helvetica-Oblique", fontSize=7.5, textColor=GREY, spaceAfter=8)
    style_precio = ParagraphStyle("precio", fontName="Helvetica-Bold", fontSize=22, textColor=NAVY, alignment=TA_RIGHT)
    style_chip = ParagraphStyle("chip", fontName="Helvetica-Bold", fontSize=10, alignment=TA_RIGHT)

    story.append(Paragraph("CF WEALTH · INFORME DE SUBYACENTE", style_eyebrow))
    story.append(Paragraph(f"{nombre} ({ticker})", style_titulo))
    story.append(Paragraph(
        f"Generado el {pd.Timestamp.now().strftime('%d/%m/%Y %H:%M')} · Documento interno de socios — no es "
        f"asesoramiento de inversión ni una recomendación de compra/venta.", style_sub))
    story.append(HRFlowable(width="100%", thickness=0.6, color=LINE, spaceAfter=10))

    # --- Cabecera de precio ---
    var_dia = datos.get("variacion_dia_pct")
    chip_texto = f"{'+' if (var_dia or 0) >= 0 else ''}{var_dia:.2f}%" if var_dia is not None else "N/D"
    style_chip.textColor = GREEN if (var_dia or 0) >= 0 else RED
    cabecera = Table([[
        Paragraph(f"${precio:,.2f}" if precio else "N/D", style_precio),
        Paragraph(chip_texto + " hoy", style_chip),
    ]], colWidths=[110 * mm, 55 * mm])
    cabecera.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "MIDDLE")]))
    story.append(cabecera)

    # --- Tabla de métricas clave ---
    minimo, maximo = datos.get("precio_min_52sem"), datos.get("precio_max_52sem")
    rango_txt = f"${minimo:,.2f} — ${maximo:,.2f}" if minimo and maximo else "N/D"
    pe_txt = f"{datos['trailing_pe']:.1f}" if datos.get("trailing_pe") else "N/D"
    cap_txt = _fmt_market_cap(datos.get("market_cap"))
    target_txt = f"${datos['target_medio']:,.2f}" if datos.get("target_medio") else "N/D"
    vol_txt = f"{datos['volatilidad_anual_pct']:.1f}%" if datos.get("volatilidad_anual_pct") else "N/D"
    earn_txt = str(datos.get("proxima_fecha_resultados") or "N/D")[:10]

    filas_tabla = [
        ["Rango 52 semanas", rango_txt, "PER", pe_txt],
        ["Capitalización", cap_txt, "Precio objetivo consenso", target_txt],
        ["Volatilidad anual", vol_txt, "Próximo earnings", earn_txt],
    ]
    tabla_metricas = Table(filas_tabla, colWidths=[42 * mm, 45 * mm, 45 * mm, 45 * mm])
    tabla_metricas.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("TEXTCOLOR", (0, 0), (0, -1), GREY), ("TEXTCOLOR", (2, 0), (2, -1), GREY),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica-Bold"), ("FONTNAME", (3, 0), (3, -1), "Helvetica-Bold"),
        ("TEXTCOLOR", (1, 0), (1, -1), NAVY), ("TEXTCOLOR", (3, 0), (3, -1), NAVY),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7), ("TOPPADDING", (0, 0), (-1, -1), 7),
        ("LINEBELOW", (0, 0), (-1, -2), 0.4, LINE),
    ]))
    story.append(Spacer(1, 4 * mm))
    story.append(tabla_metricas)

    # --- Gráfico de evolución ---
    precio_contingencia_peor = None
    if notas_ticker is not None and not notas_ticker.empty and precio:
        margenes = []
        for _, r in notas_ticker.iterrows():
            barrera_cupon, precio_compra = r.get("barrera_cupon"), r.get("precio_compra")
            if pd.notna(barrera_cupon) and pd.notna(precio_compra) and precio_compra:
                pc = precio_compra * barrera_cupon
                margenes.append((precio / pc - 1, pc))
        if margenes:
            precio_contingencia_peor = min(margenes, key=lambda x: x[0])[1]

    png_grafico = _exportar_grafico_precio_png(datos, precio_contingencia_peor)
    story.append(Paragraph("Evolución del precio (12 meses)", style_h3))
    if png_grafico:
        img_buffer = BytesIO(png_grafico)
        story.append(RLImage(img_buffer, width=170 * mm, height=66 * mm))
    else:
        story.append(Paragraph("Gráfico no disponible en este PDF (falta la librería kaleido en el servidor).", style_caption))

    # --- Tus notas con este ticker ---
    if notas_ticker is not None and not notas_ticker.empty and precio:
        story.append(Paragraph("Notas expuestas a este ticker", style_h3))
        filas_notas = [["Nota", "Precio compra", "vs inicial", "Margen a barrera"]]
        for _, r in notas_ticker.iterrows():
            precio_compra = r.get("precio_compra")
            if pd.isna(precio_compra) or not precio_compra:
                continue
            pct_vs_inicial = (precio / precio_compra - 1) * 100
            barrera_cupon = r.get("barrera_cupon")
            margen = (precio / (precio_compra * barrera_cupon) - 1) * 100 if pd.notna(barrera_cupon) else None
            filas_notas.append([
                str(int(r.get("nota"))) if pd.notna(r.get("nota")) else "?",
                f"${precio_compra:,.2f}", f"{pct_vs_inicial:+.1f}%",
                f"{margen:+.1f}%" if margen is not None else "N/D",
            ])
        if len(filas_notas) > 1:
            tabla_notas = Table(filas_notas, colWidths=[25 * mm, 40 * mm, 35 * mm, 40 * mm])
            tabla_notas.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"), ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white), ("BACKGROUND", (0, 0), (-1, 0), NAVY),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [rl_colors.white, rl_colors.Color(0.97, 0.97, 0.98)]),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6), ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("LINEBELOW", (0, 0), (-1, -1), 0.3, LINE),
            ]))
            story.append(tabla_notas)

    # --- Noticias que avalan el informe ---
    story.append(Paragraph("Noticias recientes que avalan el análisis", style_h3))
    if noticias and isinstance(noticias[0], dict) and "_error" in noticias[0]:
        story.append(Paragraph(f"No se pudieron obtener noticias: {noticias[0]['_error']}", style_normal))
    elif not noticias:
        story.append(Paragraph("No se encontraron noticias relevantes recientes tras buscar.", style_normal))
    else:
        for n in noticias[:8]:
            prioridad = str(n.get("prioridad", "media")).upper()
            story.append(Paragraph(f"<b>{n.get('titular', 'Sin título')}</b>  <font size=7 color='#9A6B24'>[{prioridad}]</font>", style_normal))
            story.append(Paragraph(f"{n.get('fuente', 'Fuente desconocida')} · {n.get('fecha', 'sin fecha')}", style_caption))
            story.append(Paragraph(n.get("resumen", ""), style_normal))
            if n.get("url"):
                story.append(Paragraph(f"<link href='{n['url']}'><font color='#9A6B24'>{n['url']}</font></link>", style_caption))
            story.append(Spacer(1, 2 * mm))
    story.append(Paragraph(
        "Noticias generadas por IA a partir de búsqueda web pública, en español y con enlace a la fuente original. "
        "No constituyen recomendación de inversión ni predicción de precio.", style_caption))

    doc.build(story)
    return output.getvalue()


def generar_memo_directorio_notas_pptx(
    resultados_notas: list, texto_recomendacion: str,
    capital_disponible: float, tasa_inversor_pct: float,
    veredicto_nota_recomendada: str = None,
) -> bytes:
    """
    Arma un memo de directorio (.pptx) a partir de lo que YA calculó el comparador de notas
    (resultados_notas, con Monte Carlo + score) y la recomendación de reparto de Claude.
    No recalcula nada — solo formatea lo que ya existe en resultados_notas.

    OJO DE DISEÑO: python-pptx crea tablas y gráficos con el tema azul por defecto de Office
    (bandas de color, bordes de rejilla) que no coincide con nuestra paleta navy — hay que
    despojarlas de ese estilo explícitamente (_strip_table_style) y poner bordes/colores a mano,
    o el resultado se ve "genérico" en vez de a medida.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    import re as _re

    NAVY = RGBColor(0x1E, 0x27, 0x61)
    NAVY_2 = RGBColor(0x2A, 0x34, 0x70)
    ICE = RGBColor(0xCA, 0xDC, 0xFC)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    RED = RGBColor(0xC0, 0x39, 0x2B)
    GOLD = RGBColor(0x9A, 0x6B, 0x24)
    GREEN = RGBColor(0x2E, 0x7D, 0x32)
    AMBER = RGBColor(0xD6, 0x89, 0x10)
    GREY = RGBColor(0x5A, 0x5A, 0x5A)
    GREY_LIGHT = RGBColor(0x99, 0x99, 0x99)
    TEXT = RGBColor(0x22, 0x22, 0x22)
    LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFC)
    BORDER = RGBColor(0xE3, 0xE5, 0xEF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    W, H = 13.333, 7.5

    # ---------- Helpers de bajo nivel: sin estos, python-pptx se ve "de fábrica" ----------
    def _strip_table_style(table):
        """Quita el estilo de tabla azul por defecto de Office (bandas + bordes de rejilla).
        Sin esto, cualquier fill que pongamos por celda queda tapado o mezclado con el tema."""
        tbl = table._tbl
        tblPr = tbl.find(qn("a:tblPr"))
        if tblPr is not None:
            tblPr.set("firstRow", "0")
            tblPr.set("bandRow", "0")
            style_id = tblPr.find(qn("a:tableStyleId"))
            if style_id is not None:
                tblPr.remove(style_id)

    def _cell_border(cell, edges=("bottom",), color=BORDER, width_pt=0.75):
        """Pone SOLO los bordes que le pidamos (por defecto ninguno) — así las tablas quedan
        con líneas horizontales finas tipo editorial, sin rejilla vertical."""
        tag_map = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}
        order = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        hexcolor = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}" if isinstance(color, tuple) else str(color)
        for edge in ("left", "right", "top", "bottom"):
            tag = tag_map[edge]
            existing = tcPr.find(qn(tag))
            if existing is not None:
                tcPr.remove(existing)
            ln = tcPr.makeelement(qn(tag), {"w": str(int(width_pt * 12700)), "cap": "flat", "cmpd": "sng"})
            if edge in edges:
                fill = ln.makeelement(qn("a:solidFill"), {})
                clr = fill.makeelement(qn("a:srgbClr"), {"val": hexcolor})
                fill.append(clr)
                ln.append(fill)
            else:
                ln.append(ln.makeelement(qn("a:noFill"), {}))
            insert_idx = 0
            for i, child in enumerate(tcPr):
                if child.tag.split("}")[-1] in [t.split(":")[-1] for t in order]:
                    insert_idx = i + 1
                else:
                    break
            tcPr.insert(insert_idx, ln)

    def _style_chart_axes(chart, label_color=GREY, grid_color=RGBColor(0xEE, 0xEE, 0xF2)):
        try:
            cat = chart.category_axis
            cat.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xD6)
            cat.format.line.width = Pt(0.75)
            cat.tick_labels.font.size = Pt(10)
            cat.tick_labels.font.color.rgb = label_color
            cat.tick_labels.font.name = "Calibri"
            cat.has_major_gridlines = False
            cat.major_tick_mark = 0
        except Exception:
            pass
        try:
            val = chart.value_axis
            val.format.line.fill.background()
            val.has_major_gridlines = True
            val.major_gridlines.format.line.color.rgb = grid_color
            val.major_gridlines.format.line.width = Pt(0.75)
            val.tick_labels.font.size = Pt(9)
            val.tick_labels.font.color.rgb = label_color
            val.tick_labels.font.name = "Calibri"
            val.major_tick_mark = 0
        except Exception:
            pass

    def add_slide(bg=WHITE):
        s = prs.slides.add_slide(BLANK)
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = bg
        rect.line.fill.background()
        rect.shadow.inherit = False
        s.shapes._spTree.remove(rect._element)
        s.shapes._spTree.insert(2, rect._element)
        return s

    def add_text(s, x, y, w, h, text, size=12, bold=False, italic=False, color=TEXT,
                 align=PP_ALIGN.LEFT, font="Calibri", anchor=None, line_spacing=1.0, letter_spacing=None):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if anchor:
            tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))
        return box

    def add_eyebrow(s, x, y, text, color=NAVY, w=8):
        """Etiqueta pequeña en mayúsculas + barrita de acento debajo — ancla visual que usa
        el memo de referencia para separar secciones, en vez de un simple texto suelto."""
        add_text(s, x, y, w, 0.28, text, size=11.5, bold=True, color=color, letter_spacing=1.5)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.32), Inches(0.45), Pt(2.6))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background(); bar.shadow.inherit = False

    def add_card(s, x, y, w, h, value, label, value_color=NAVY, bg=LIGHT_BG, accent=None):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.adjustments[0] = 0.06
        box.fill.solid(); box.fill.fore_color.rgb = bg
        box.line.color.rgb = BORDER; box.line.width = Pt(0.75)
        box.shadow.inherit = False
        if accent:
            top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.15), Inches(y), Inches(w - 0.3), Pt(3))
            top.fill.solid(); top.fill.fore_color.rgb = accent
            top.line.fill.background(); top.shadow.inherit = False
        add_text(s, x, y + 0.18, w, 0.65, value, size=25, bold=True, color=value_color,
                  align=PP_ALIGN.CENTER, font="Cambria")
        add_text(s, x + 0.15, y + h - 0.48, w - 0.3, 0.4, label, size=10, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.05)

    def add_table(s, x, y, w, headers, rows, col_widths=None, red_col=None, bold_col=None, row_h=0.4, red_rows=None):
        """red_rows: set/list de índices de fila (0-based, sin contar cabecera) cuya celda en
        red_col se pinta de rojo — evita meter marcadores dentro del texto de la celda."""
        red_rows = set(red_rows or [])
        n_rows, n_cols = len(rows) + 1, len(headers)
        tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows))
        tbl = tbl_shape.table
        _strip_table_style(tbl)
        tbl.first_row = False
        tbl.horz_banding = False
        if col_widths:
            total = sum(col_widths)
            for c, cw in zip(tbl.columns, col_widths):
                c.width = Inches(w * cw / total)
        for r_ in tbl.rows:
            r_.height = Inches(row_h)
        for j, htxt in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = htxt
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            _cell_border(cell, edges=(), color=NAVY)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10.5); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
        for i, fila in enumerate(rows):
            es_ultima = (i == len(rows) - 1)
            for j, val in enumerate(fila):
                cell = tbl.cell(i + 1, j)
                cell.text = str(val)
                cell.margin_left = cell.margin_right = Inches(0.08)
                cell.margin_top = cell.margin_bottom = Inches(0.04)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
                _cell_border(cell, edges=() if es_ultima else ("bottom",), color=BORDER, width_pt=0.75)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10.5); p.font.name = "Calibri"
                    p.font.bold = (bold_col is not None and j == bold_col)
                    p.font.color.rgb = RED if (red_col is not None and j == red_col and i in red_rows) else TEXT
        return tbl

    def footer(s, n):
        add_text(s, 0.4, H - 0.38, 6, 0.3, "Chaparro Fernández Wealth · Confidencial", size=8.5, color=GREY_LIGHT)
        add_text(s, W - 1.2, H - 0.38, 0.8, 0.3, f"{n:02d}", size=8.5, color=GREY_LIGHT, align=PP_ALIGN.RIGHT)

    def veredicto_color(p_perdida):
        return RED if p_perdida > 0.35 else (AMBER if p_perdida > 0.2 else GREEN)

    ganador = resultados_notas[0]  # ya viene ordenado por score desc
    pago_map = {1: "Mensual", 3: "Trimestral", 6: "Semestral"}

    # ---------- Slide 1: Portada ----------
    s = add_slide(NAVY)
    linea = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(H - 0.06), prs.slide_width, Pt(4.5))
    linea.fill.solid(); linea.fill.fore_color.rgb = RGBColor(0x4A, 0x56, 0xB8)
    linea.line.fill.background(); linea.shadow.inherit = False
    add_text(s, 0.7, 0.55, 10, 0.32, "MESA DE PRODUCTOS ESTRUCTURADOS", size=12.5, bold=True, color=ICE, letter_spacing=2)
    add_text(s, 0.7, 0.95, 10, 0.4, "Comparador de Notas — Memo de Directorio", size=16, color=WHITE)
    add_text(s, 0.7, 1.5, 11.5, 1.0, f"{len(resultados_notas)} nota(s) candidata(s)", size=40, bold=True, color=WHITE, font="Cambria")
    add_text(s, 0.7, 2.45, 11.5, 0.4,
             f"Capital a repartir: ${capital_disponible:,.0f}   ·   Tasa pagada al inversor: {tasa_inversor_pct*100:.1f}% anual",
             size=13.5, color=ICE)

    cw, gap, x0 = 3.77, 0.28, 0.7
    for i, r in enumerate(resultados_notas[:3]):
        x = x0 + i * (cw + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.15), Inches(cw), Inches(1.5))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = NAVY_2
        card.line.color.rgb = RGBColor(0x3D, 0x48, 0x8F); card.line.width = Pt(0.75)
        card.shadow.inherit = False
        add_text(s, x + 0.18, 3.28, cw - 0.36, 0.35, f"#{i+1}   {r['nombre']}", size=12.5, bold=True, color=ICE)
        add_text(s, x, 3.62, cw, 0.65, f"{r['score']:.0f}", size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Cambria")
        add_text(s, x, 4.28, cw, 0.3, f"/100  ·  {', '.join(r['tickers'])}", size=10.5, color=ICE, align=PP_ALIGN.CENTER)

    banda = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.15), Inches(11.93), Inches(1.15))
    banda.adjustments[0] = 0.06
    banda.line.fill.background(); banda.shadow.inherit = False
    if veredicto_nota_recomendada is None:
        # La recomendación (texto de Claude) concluye NO asignar capital a ninguna nota — el
        # banner tiene que decir justo eso, no puede mostrar "mejor candidata por score": ese
        # ranking es solo un componente más, y mostrarlo aquí como si fuera la conclusión final
        # es lo que generaba la contradicción entre la portada y la recomendación del memo.
        banda.fill.solid(); banda.fill.fore_color.rgb = RED
        add_text(s, 1.0, 5.28, 10, 0.28, "RECOMENDACIÓN FINAL", size=10.5, bold=True, color=WHITE, letter_spacing=1.5)
        add_text(s, 1.0, 5.58, 11, 0.6, "Rechazar todas las notas — ver justificación en la recomendación de reparto", size=18, bold=True, color=WHITE, font="Cambria")
    else:
        banda.fill.solid(); banda.fill.fore_color.rgb = veredicto_color(ganador["prob_perdida_capital"])
        add_text(s, 1.0, 5.28, 10, 0.28, "NOTA RECOMENDADA", size=10.5, bold=True, color=WHITE, letter_spacing=1.5)
        add_text(s, 1.0, 5.58, 11, 0.6, f"{veredicto_nota_recomendada}", size=21, bold=True, color=WHITE, font="Cambria")
    add_text(s, 0.7, H - 0.55, 10, 0.3, f"Chaparro Fernández Wealth · Confidencial · {pd.Timestamp.now().strftime('%d/%m/%Y %H:%M')}",
             size=9, color=RGBColor(0x93, 0x9C, 0xD6))

    # ---------- Slide 2: Tabla comparativa + gráfico de score ----------
    s = add_slide()
    add_eyebrow(s, 0.5, 0.35, "COMPARATIVA")
    add_text(s, 0.5, 0.78, 11.5, 0.5, "Las candidatas, lado a lado", size=22, bold=True, font="Cambria")

    headers = ["Nota", "Tickers", "Cupón anual", "Pago", "Protección", "Prob. cupón/periodo", "Prob. pérdida capital", "Rent. neta esp.", "Score"]
    rows = [[
        r["nombre"], ", ".join(r["tickers"]), f"{r['cupon_anual']*100:.1f}%",
        pago_map.get(r.get("periodicidad"), "—"),
        "Buffer" if r.get("tipo_proteccion") == "buffer" else "Barrera",
        f"{r['prob_cupon_media']*100:.1f}%",
        f"{r['prob_perdida_capital']*100:.1f}%",
        f"{r['rentabilidad_esperada_neta']*100:.2f}%", f"{r['score']:.0f}/100",
    ] for r in resultados_notas]
    filas_riesgo_alto = {i for i, r in enumerate(resultados_notas) if r["prob_perdida_capital"] > 0.3}
    add_table(s, 0.5, 1.55, 12.3, headers, rows, col_widths=[1.1, 1.2, 0.9, 0.9, 0.9, 1.3, 1.3, 1.1, 0.8],
              red_col=6, bold_col=8, row_h=0.42, red_rows=filas_riesgo_alto)

    y_chart = 1.65 + 0.42 * (len(resultados_notas) + 1) + 0.35
    chart_data = CategoryChartData()
    chart_data.categories = [r["nombre"] for r in resultados_notas]
    chart_data.add_series("Score", [round(r["score"], 1) for r in resultados_notas])
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(y_chart), Inches(12.3), Inches(H - y_chart - 0.5), chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 55
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = NAVY
    plot.series[0].format.line.fill.background()
    plot.has_data_labels = True
    plot.data_labels.number_format = "0"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = NAVY
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    _style_chart_axes(chart)
    footer(s, 2)

    # ---------- Slides 3..N: una por nota ----------
    for idx, r in enumerate(resultados_notas):
        s = add_slide()
        add_eyebrow(s, 0.5, 0.35, f"NOTA {idx+1} DE {len(resultados_notas)}")
        add_text(s, 0.5, 0.78, 11.5, 0.5, r["nombre"], size=24, bold=True, font="Cambria")
        subtitulo = (
            f"{', '.join(r['tickers'])}  ·  {r['meses_venc']} meses hasta vencimiento  ·  "
            f"pago {pago_map.get(r.get('periodicidad'), '—').lower()}  ·  call cada {r.get('frecuencia_call', '—')} meses  ·  "
            f"protección {'buffer' if r.get('tipo_proteccion') == 'buffer' else 'barrera'}"
        )
        add_text(s, 0.5, 1.28, 11.5, 0.4, subtitulo, size=12, color=GREY)

        cw2, gap2 = 2.87, 0.24
        add_card(s, 0.5, 1.85, cw2, 1.3, f"{r['score']:.0f}", "score /100", NAVY, accent=NAVY)
        add_card(s, 0.5 + (cw2 + gap2), 1.85, cw2, 1.3, f"{r['prob_cupon_media']*100:.0f}%", "prob. media cupón", NAVY, accent=NAVY)
        color_perdida = RED if r["prob_perdida_capital"] > 0.3 else NAVY
        add_card(s, 0.5 + 2 * (cw2 + gap2), 1.85, cw2, 1.3, f"{r['prob_perdida_capital']*100:.0f}%", "prob. pérdida capital", color_perdida, accent=color_perdida)
        add_card(s, 0.5 + 3 * (cw2 + gap2), 1.85, cw2, 1.3, f"{r['rentabilidad_esperada_neta']*100:.1f}%", "rentabilidad neta esp.", NAVY, accent=NAVY)

        add_eyebrow(s, 0.5, 3.4, "DETALLE POR TICKER (WORST-OF)", color=GREY)
        headers_t = ["Ticker", "Precio actual", "Barrera cupón", "Barrera capital", "Nivel call", "Colchón capital"]
        rows_t = []
        for tk, tf_ in zip(r["tickers"], r["tickers_full"]):
            fund = next((f for f in r["fundamentales"] if f.get("ticker") == tk), {})
            precio = fund.get("precio_actual")
            colchon = (1 - tf_["barrera_capital_pct"]) * 100
            rows_t.append([
                tk, f"${precio:,.2f}" if precio else "—",
                f"{tf_['barrera_cupon_pct']*100:.0f}%", f"{tf_['barrera_capital_pct']*100:.0f}%",
                f"{tf_['call_level_pct']*100:.0f}%", f"{colchon:.1f}%",
            ])
        add_table(s, 0.5, 3.85, 12.3, headers_t, rows_t, row_h=0.4)

        margen = r.get("margen_sobre_inversor", r["cupon_anual"] - tasa_inversor_pct) * 100
        y_nota = 3.85 + 0.4 * (len(rows_t) + 1) + 0.3
        add_text(
            s, 0.5, y_nota, 12.3, 0.6,
            f"Margen bruto teórico sobre la tasa pagada al inversor: {margen:.2f} pts.   ·   "
            f"Pérdida media si incumple: {r['perdida_pct_promedio']:.1f}%.",
            size=11.5, italic=True, color=GREY,
        )
        footer(s, 3 + idx)

    # ---------- Slide: Recomendación ----------
    s = add_slide(NAVY)
    add_text(s, 0.7, 0.6, 10, 0.32, "RECOMENDACIÓN DE REPARTO", size=12.5, bold=True, color=ICE, letter_spacing=2)
    add_text(s, 0.7, 1.05, 11.5, 0.6, f"Cómo repartir ${capital_disponible:,.0f}", size=27, bold=True, color=WHITE, font="Cambria")

    texto_limpio = _re.sub(r"\*\*(.+?)\*\*", r"\1", texto_recomendacion or "No se pudo generar una recomendación.")
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.9))
    tf = box.text_frame
    tf.word_wrap = True
    primero = True
    for linea_txt in texto_limpio.split("\n"):
        linea_txt = linea_txt.strip()
        if not linea_txt:
            continue
        p = tf.paragraphs[0] if primero else tf.add_paragraph()
        primero = False
        es_conclusion = bool(_re.match(r"^Conclusi[óo]n:", linea_txt, _re.IGNORECASE))
        p.line_spacing = 1.22
        p.space_after = Pt(12 if es_conclusion else 7)
        r_ = p.add_run()
        r_.text = ("•  " + linea_txt[2:]) if linea_txt.startswith(("- ", "* ")) else linea_txt
        if es_conclusion:
            # La línea de conclusión (primera del texto, ver prompt) se destaca en dorado y más
            # grande para que ancle visualmente todo el análisis que viene después — así el
            # lector no puede perderla ni confundirla con un punto intermedio del razonamiento.
            r_.font.size = Pt(15)
            r_.font.color.rgb = GOLD
            r_.font.bold = True
        else:
            r_.font.size = Pt(12.5)
            r_.font.color.rgb = ICE if linea_txt.startswith("#") else WHITE
            r_.font.bold = linea_txt.startswith("#")
        r_.font.name = "Calibri"
    add_text(s, 0.7, H - 0.5, 11, 0.3,
             "Monte Carlo con 5.000 escenarios por nota, volatilidad histórica de 12 meses, con correlación real entre tickers cuando está activada.",
             size=8.5, color=RGBColor(0x93, 0x9C, 0xD6))

    # ---------- Slide final: supuestos ----------
    s = add_slide()
    add_eyebrow(s, 0.5, 0.35, "APÉNDICE")
    add_text(s, 0.5, 0.78, 11.5, 0.5, "Metodología y supuestos", size=22, bold=True, font="Cambria")
    bloques = [
        ("Modelo", "Monte Carlo real-world worst-of (simular_montecarlo_nota), 5.000 trayectorias por nota, GBM sin drift."),
        ("Volatilidad", "Histórica realizada de 12 meses, vía yfinance. No incluye skew de opciones."),
        ("Correlación", "No modelada — se asume independencia entre tickers de una misma nota (supuesto simplificado)."),
        ("Barrera / buffer", "Barrera = pérdida total del peor activo si rompe. Buffer = solo se pierde el exceso sobre el colchón. El Monte Carlo ya calcula la pérdida distinto según el tipo elegido por nota."),
        ("Score", "35% prob. cupón + 20% margen sobre tasa inversor + 20% (1 − prob. pérdida capital) + 10% periodicidad + 10% cercanía de call a 3m + 5% tipo de protección."),
        ("Recomendación", "Generada por Claude a partir de las métricas ya calculadas — no recalcula ni inventa números."),
    ]
    y = 1.65
    for k, v in bloques:
        add_text(s, 0.5, y, 2.3, 0.6, k, size=12, bold=True, color=NAVY)
        add_text(s, 2.9, y, 9.9, 0.65, v, size=11, color=TEXT, line_spacing=1.15)
        y += 0.75
    add_text(s, 0.5, 6.85, 12.3, 0.4,
             "Memo generado automáticamente desde el Comparador de Notas. No sustituye la verificación manual contra el Excel.",
             size=9, italic=True, color=GREY_LIGHT)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def generar_ficha_empresa_ia(
    ticker: str, nombre_compania: str,
    barrera_capital_pct: float | None = None, colchon_capital_pct: float | None = None,
    barrera_cupon_pct: float | None = None, colchon_cupon_pct: float | None = None,
) -> str:
    """
    Usa la API de Claude con búsqueda web para explicar, en dos secciones breves,
    a qué se dedica una compañía y qué noticias recientes son relevantes para su cotización.
    Nunca inventa precios objetivo (esos se muestran aparte con datos reales de Yahoo Finance)
    y siempre redacta con sus propias palabras, sin copiar frases textuales de las fuentes.
    """
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
    contexto_barrera = ""
    tercera_seccion = ""
    if barrera_capital_pct is not None and colchon_capital_pct is not None:
        contexto_barrera += (
            f" Dato clave de la nota: la barrera de CAPITAL (contingencia) de esta acción está en el "
            f"{barrera_capital_pct*100:.0f}% de su precio de entrada — la acción tiene un colchón de "
            f"{colchon_capital_pct:.1f}% de caída antes de poner en riesgo el capital (evento único, al vencimiento)."
        )
    if barrera_cupon_pct is not None and colchon_cupon_pct is not None:
        contexto_barrera += (
            f" Además, la barrera de CUPÓN (la que decide si se cobra el interés mensual/periódico) está en el "
            f"{barrera_cupon_pct*100:.0f}% de su precio de entrada — colchón de {colchon_cupon_pct:.1f}% de caída "
            f"antes de perder el cupón de ESE periodo concreto (evento recurrente, cada mes/trimestre, no afecta al capital)."
        )
    if contexto_barrera:
        tercera_seccion = (
            "\n\n**Opinión sobre el riesgo (cupón y capital):** valora, en 3-4 frases y SOLO con lo que encuentres "
            "(precio objetivo de analistas, rango de estimaciones, noticias), la probabilidad de que la acción caiga "
            "lo suficiente como para (a) perder el cupón de un periodo concreto (barrera de cupón, más exigente y "
            "recurrente cada mes) y (b) poner en riesgo el capital al vencimiento (barrera de capital, más laxa pero "
            "solo se juzga una vez, al final). Son dos riesgos distintos con colchones distintos — coméntalos por "
            "separado. El criterio correcto NO es si la acción tiene volatilidad alta o noticias negativas en "
            "abstracto — es si una caída realista se queda por encima de cada barrera. Si el colchón de cupón es "
            "ajustado pero el de capital es amplio, dilo así explícitamente (riesgo real de perder algún cupón "
            "puntual, pero capital protegido); si ambos colchones son amplios frente a las caídas que manejan los "
            "analistas, dilo claramente como algo positivo aunque haya ruido de corto plazo en las noticias."
        )
    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
            json={
                "model": "claude-sonnet-4-5",
                "max_tokens": 1400,
                "tools": [{"type": "web_search_20250305", "name": "web_search", "max_uses": 6}],
                "system": (
                    "Eres un analista senior que prepara fichas rápidas de compañías subyacentes de notas "
                    "estructuradas para un gestor de fondo. Con búsqueda web, redacta en español EXACTAMENTE "
                    "estas secciones, breves y en tus propias palabras (nunca copies frases textuales de "
                    "ninguna fuente, parafrasea siempre):\n\n"
                    "**A qué se dedica:** 2-3 frases explicando el negocio principal y de dónde saca ingresos.\n\n"
                    "**Noticias relevantes recientes:** 4-6 líneas con hechos de los últimos 1-2 meses que puedan "
                    "afectar la cotización (resultados financieros, cambios de rating de analistas, noticias "
                    "regulatorias/legales, eventos corporativos), citando la fuente entre paréntesis."
                    f"{tercera_seccion}\n\n"
                    "Tono profesional, nunca alarmista, sin mayúsculas dramáticas. NUNCA inventes ni menciones un "
                    "precio objetivo o una predicción de precio — eso se muestra aparte con datos reales de mercado. "
                    "Si de verdad no encuentras nada relevante tras buscar, dilo claramente en vez de inventar contenido."
                    f"{contexto_barrera}"
                ),
                "messages": [{"role": "user", "content": f"Prepara la ficha de {nombre_compania} ({ticker})."}],
            },
            timeout=60,
        )
        data = resp.json()
        if data.get("type") == "error" or data.get("error"):
            msg = data.get("error", {}).get("message", str(data))
            return f"⚠️ No se pudo generar la ficha de {ticker}: error de la API ({msg})."
        contenido = data.get("content", [])
        hubo_error_busqueda_critico = False
        for bloque in contenido:
            if bloque.get("type") == "web_search_tool_result":
                resultado_bloque = bloque.get("content", {})
                if isinstance(resultado_bloque, dict) and resultado_bloque.get("type") == "web_search_tool_result_error":
                    codigo_error = resultado_bloque.get("error_code", "desconocido")
                    # "max_uses_exceeded" es benigno: la IA ya hizo varias búsquedas útiles antes de
                    # tocar el límite y normalmente sigue redactando con lo que ya encontró — no hay
                    # que descartar ese texto. Otros códigos (permisos, servicio no disponible, etc.)
                    # sí son un fallo real de la búsqueda.
                    if codigo_error != "max_uses_exceeded":
                        hubo_error_busqueda_critico = True
                        codigo_error_critico = codigo_error
        texto = "".join(b.get("text", "") for b in contenido if b.get("type") == "text").strip()
        if texto:
            return texto
        if hubo_error_busqueda_critico:
            return f"⚠️ La búsqueda web falló para {ticker} (código: {codigo_error_critico})."
        return f"No se encontró información relevante para {ticker} tras buscar."
    except Exception as e:
        return f"⚠️ No se pudo generar la ficha de {ticker} (error de conexión): {e}"


def _tab_comparador_notas(df_control: pd.DataFrame):
    st.caption(
        "Introduce entre 1 y 3 notas candidatas con sus términos, y tu capital disponible. "
        "Se simula cada una con Monte Carlo (usando la volatilidad histórica real de cada ticker) "
        "para estimar probabilidades de cupón, call y pérdida de capital, y se calcula una "
        "propuesta de reparto de capital razonada.\n\n"
        "⚠️ **Esto es una estimación probabilística con supuestos simplificados** — no una predicción "
        "ni garantía de resultado. Por defecto asume que la acción no tiene tendencia (ni sube ni baja); "
        "si tenés convicción propia (o el consenso de analistas la tiene), podés metérsela abajo."
    )

    n_notas = st.radio("¿Cuántas notas quieres comparar?", [1, 2, 3], horizontal=True, index=1)
    notas_input = []

    cols = st.columns(n_notas)
    for i in range(n_notas):
        with cols[i]:
            st.markdown(f"**Nota candidata {i+1}**")
            nombre = st.text_input(f"Nombre/etiqueta", value=f"Nota {i+1}", key=f"comp_nombre_{i}")
            n_tickers = st.number_input(f"Nº de tickers (worst-of)", min_value=1, max_value=3, value=1, key=f"comp_ntick_{i}")
            tickers_nota = []
            for j in range(n_tickers):
                st.markdown(f"_Ticker {j+1}_")
                tk = st.text_input(f"Símbolo", key=f"comp_tk_{i}_{j}").strip().upper()
                bc = st.number_input(f"Barrera cupón (%)", min_value=1, max_value=100, value=60, key=f"comp_bc_{i}_{j}") / 100
                bcap = st.number_input(f"Barrera/buffer capital (%)", min_value=1, max_value=100, value=60, key=f"comp_bcap_{i}_{j}") / 100
                cl = st.number_input(f"Nivel de call (%)", min_value=50, max_value=200, value=100, key=f"comp_cl_{i}_{j}") / 100
                if tk:
                    tickers_nota.append({"ticker": tk, "barrera_cupon_pct": bc, "barrera_capital_pct": bcap, "call_level_pct": cl})
            cupon_anual = st.number_input(f"Cupón anual (%)", min_value=0.0, max_value=100.0, value=25.0, key=f"comp_cupon_{i}") / 100
            meses_venc = st.number_input(f"Meses hasta vencimiento", min_value=1, max_value=60, value=36, key=f"comp_meses_{i}")
            periodicidad = st.selectbox(
                f"Periodicidad de pago del cupón", [1, 3, 6],
                format_func=lambda x: {1: "Mensual", 3: "Trimestral", 6: "Semestral"}[x],
                key=f"comp_period_{i}",
                help="Mensual es preferible: capital rota más rápido y el riesgo de cada evento individual es menor.",
            )
            frecuencia_call = st.selectbox(
                f"Frecuencia de call (meses)", [1, 3, 6, 12], index=1, key=f"comp_callfreq_{i}",
                help="Cada cuánto puede el emisor cancelar la nota. 3 meses suele ser lo preferible — no es el factor más importante.",
            )
            tipo_proteccion = st.selectbox(
                f"Tipo de protección de capital", ["barrera", "buffer"],
                format_func=lambda x: "Barrera (pérdida total si rompe)" if x == "barrera" else "Buffer (protege el % indicado, solo se pierde el exceso)",
                key=f"comp_tipoprot_{i}",
                help=(
                    "Barrera europea: si el peor activo rompe el nivel al vencimiento, el fondo asume la caída COMPLETA "
                    "del activo desde el nivel inicial. Buffer: el fondo está protegido hasta ese %; si la caída es mayor, "
                    "solo se pierde el exceso. A igualdad de nivel nominal, el buffer siempre es menos arriesgado."
                ),
            )
            st.markdown("_Supuestos de la simulación_")
            modo_deriva = st.selectbox(
                "Deriva de precio (tendencia asumida)",
                ["neutral", "analistas", "manual"],
                format_func=lambda x: {
                    "neutral": "Neutral — sin deriva (más conservador)",
                    "analistas": "Consenso de analistas (precio objetivo, ~12 meses)",
                    "manual": "Manual — mi propia convicción",
                }[x],
                key=f"comp_mododeriva_{i}",
                help=(
                    "Por defecto (neutral) el modelo NO asume que la acción vaya a subir ni a bajar — solo "
                    "mide cuánto se mueve. Eso hace que acciones muy volátiles salgan mal aunque el consenso "
                    "las vea subiendo, porque con 30+ observaciones mensuales hasta un camino alcista puede "
                    "tocar la barrera por el camino. Si tenés una visión fundamentada (propia o de analistas), "
                    "metela aquí para que la simulación la refleje."
                ),
            )
            deriva_manual_pct = None
            if modo_deriva == "manual":
                deriva_manual_pct = st.number_input(
                    "Retorno anual esperado (%)", min_value=-50.0, max_value=100.0, value=15.0, step=1.0,
                    key=f"comp_derivamanual_{i}",
                    help="Tu convicción propia, informada por lo que leas en la ficha de noticias/analistas de cada compañía.",
                )
            modelar_earnings = st.checkbox(
                "Modelar shock de earnings", value=True, key=f"comp_earnings_{i}",
                help=(
                    "En vez de repartir el riesgo de resultados trimestrales uniformemente en el año (lo que "
                    "hace que se diluya), concentra una sacudida extra justo en las fechas de earnings, con la "
                    "magnitud REAL que esta acción se ha movido en sus últimos informes — así el modelo distingue "
                    "un mes cualquiera de un mes con resultados."
                ),
            )
            tiene_memoria = st.checkbox(
                "Tiene memoria (cupones perdidos se recuperan)", value=True, key=f"comp_memoria_{i}",
                help=(
                    "Si un periodo NO cumple la barrera de cupón, no se pierde sin más: se acumula, y el primer "
                    "periodo futuro en el que SÍ se cumpla la barrera, se cobran de golpe TODOS los cupones "
                    "acumulados desde el último pago (efecto retroactivo). Solo se pierden definitivamente los que "
                    "queden acumulados sin recuperar hasta el vencimiento. Revisá el term sheet: si dice 'memory "
                    "coupon' o similar, va marcado; si no lo dice, probablemente sea sin memoria."
                ),
            )
            modelar_correlacion = st.checkbox(
                "Modelar correlación real entre tickers", value=True, key=f"comp_correl_{i}",
                help=(
                    "Solo aplica si la nota tiene más de un ticker (worst-of). Calcula la correlación histórica "
                    "real (12 meses de precios diarios) entre los subyacentes y la usa en el Monte Carlo — si "
                    "suelen moverse juntos (correlación positiva, lo habitual entre acciones del mismo sector), "
                    "la probabilidad de que TODOS caigan a la vez es mayor que tratándolos como independientes. "
                    "Desmarcalo para volver al supuesto anterior (independencia total entre tickers)."
                ),
            )
            notas_input.append({
                "nombre": nombre, "tickers": tickers_nota, "cupon_anual": cupon_anual, "meses_venc": meses_venc,
                "periodicidad": periodicidad, "frecuencia_call": frecuencia_call, "tipo_proteccion": tipo_proteccion,
                "modo_deriva": modo_deriva, "deriva_manual_pct": deriva_manual_pct, "modelar_earnings": modelar_earnings,
                "tiene_memoria": tiene_memoria, "modelar_correlacion": modelar_correlacion,
            })

    col_cap, col_tasa = st.columns(2)
    with col_cap:
        capital_disponible = st.number_input("💰 Capital disponible a repartir ($)", min_value=0.0, value=100000.0, step=1000.0)
    with col_tasa:
        tasa_inversor_pct = st.number_input(
            "📤 Tasa que pagamos al inversor (% anual)", min_value=0.0, max_value=100.0, value=15.0, step=0.5,
            help="El margen real de la nota es cupón de la nota MENOS esta tasa. Se usa para juzgar si el interés de la nota compensa lo que le pagamos al inversor.",
        ) / 100

    if not st.button("⚖️ Comparar y recomendar", type="primary"):
        return

    resultados_notas = []
    for nota in notas_input:
        if not nota["tickers"]:
            continue
        with st.spinner(f"Simulando {nota['nombre']}..."):
            tickers_datos = []
            fundamentales_nota = []
            for t in nota["tickers"]:
                fd = obtener_datos_fundamentales(t["ticker"])
                if fd.get("precio_actual") is None or fd.get("volatilidad_anual_pct") is None:
                    st.error(f"No se pudo obtener precio/volatilidad de {t['ticker']} — se omite {nota['nombre']}.")
                    tickers_datos = None
                    break

                # Deriva de precio: 0 (neutral), consenso de analistas (dato real), o convicción manual.
                if nota["modo_deriva"] == "analistas":
                    drift_pct = fd.get("retorno_implicito_analistas_pct") or 0.0
                elif nota["modo_deriva"] == "manual":
                    drift_pct = nota["deriva_manual_pct"] or 0.0
                else:
                    drift_pct = 0.0

                # Shock de earnings: fechas proyectadas + magnitud histórica real de esta acción.
                dias_earnings, salto_pct = [], None
                if nota["modelar_earnings"]:
                    dias_earnings = _proyectar_dias_earnings(fd.get("proxima_fecha_resultados_dt"), nota["meses_venc"])
                    salto_pct = fd.get("salto_medio_earnings_pct")

                tickers_datos.append({
                    "ticker": t["ticker"], "precio_actual": fd["precio_actual"], "precio_inicial": fd["precio_actual"],
                    "volatilidad_anual_pct": fd["volatilidad_anual_pct"],
                    "barrera_cupon_pct": t["barrera_cupon_pct"], "barrera_capital_pct": t["barrera_capital_pct"], "call_level_pct": t["call_level_pct"],
                    "drift_anual_pct": drift_pct, "dias_earnings": dias_earnings, "salto_earnings_pct": salto_pct,
                    "historico_precios": fd.get("historico_precios"),
                })
                fundamentales_nota.append(fd)
            if not tickers_datos:
                continue

            matriz_corr = None
            if nota.get("modelar_correlacion") and len(tickers_datos) > 1:
                matriz_corr = calcular_matriz_correlacion_tickers(tickers_datos)

            eventos = _generar_horario_eventos(nota["meses_venc"], nota["periodicidad"], nota["frecuencia_call"])
            sim = simular_montecarlo_nota(
                tickers_datos, eventos, tipo_proteccion=nota["tipo_proteccion"], tiene_memoria=nota["tiene_memoria"],
                matriz_correlacion=matriz_corr,
            )

            eventos_cupon = [e for e in sim["eventos"] if e["tipo"] == "cupon"]
            prob_cupon_media = float(np.mean([e["probabilidad"] for e in eventos_cupon])) if eventos_cupon else 0.0
            cupon_periodo = nota["cupon_anual"] / (12 / nota["periodicidad"])
            # Con memoria, "cupones_pagados_esperado" ya incluye el cobro retroactivo de periodos
            # perdidos anteriores — por eso se usa esto (y no solo "probabilidad") para la rentabilidad.
            rentabilidad_esperada_cupones = sum(e["cupones_pagados_esperado"] * cupon_periodo for e in eventos_cupon) * (12 / nota["periodicidad"]) / (nota["meses_venc"] / 12)
            perdida_esperada_anualizada = sim["probabilidad_perdida_capital"] * sim["perdida_pct_promedio_si_incumple"] / 100 / (nota["meses_venc"] / 12)
            rentabilidad_neta_esperada = rentabilidad_esperada_cupones - perdida_esperada_anualizada
            margen_sobre_inversor = nota["cupon_anual"] - tasa_inversor_pct

            resultados_notas.append({
                "nombre": nota["nombre"], "tickers": [t["ticker"] for t in nota["tickers"]],
                "tickers_full": nota["tickers"],
                "cupon_anual": nota["cupon_anual"], "meses_venc": nota["meses_venc"],
                "periodicidad": nota["periodicidad"], "frecuencia_call": nota["frecuencia_call"],
                "tipo_proteccion": nota["tipo_proteccion"],
                "modo_deriva": nota["modo_deriva"], "modelar_earnings": nota["modelar_earnings"],
                "tiene_memoria": nota["tiene_memoria"],
                "matriz_correlacion": sim.get("matriz_correlacion"),
                "cupones_totales_esperados": sim["cupones_totales_esperados"],
                "cupones_totales_periodos": len(eventos_cupon),
                "cupones_perdidos_definitivo_esperado": sim.get("cupones_perdidos_definitivo_esperado"),
                "tickers_datos_sim": tickers_datos,  # incluye drift_anual_pct/salto_earnings_pct usados, para el informe
                "prob_cupon_media": prob_cupon_media, "prob_call_total": sim["probabilidad_call_total"],
                "prob_perdida_capital": sim["probabilidad_perdida_capital"], "perdida_pct_promedio": sim["perdida_pct_promedio_si_incumple"],
                "rentabilidad_esperada_neta": rentabilidad_neta_esperada, "margen_sobre_inversor": margen_sobre_inversor,
                "fundamentales": fundamentales_nota,
            })

    if not resultados_notas:
        st.warning("No se pudo simular ninguna nota — revisa los tickers.")
        return

    # --- Score compuesto 0-100 para rankear las candidatas de un vistazo ---
    #
    # Pesos y su porqué (así se decidió con el fondo, no son arbitrarios):
    #  35% Probabilidad de cobrar el cupón (worst-of, por periodo)  -> lo MÁS importante: es el evento
    #      recurrente que paga el spread mes a mes. Esta probabilidad YA viene de Monte Carlo evaluando
    #      la barrera de cupón concreta de cada ticker — si el estudio dice que casi nunca la rompe,
    #      esta nota puntúa alto AUNQUE la acción sea muy volátil o pueda caer mucho en términos
    #      absolutos. No se juzga por volatilidad genérica, se juzga por si la barrera se rompe o no.
    #  20% Margen sobre la tasa pagada al inversor, normalizado entre candidatas. Un cupón que apenas
    #      supera lo que pagamos al inversor no compensa el riesgo — el fondo absorbe toda la pérdida
    #      si la nota falla, así que un 10% de margen no justifica arriesgar mucho capital.
    #  20% (1 - probabilidad de pérdida de capital), YA AJUSTADA por si la protección es barrera o
    #      buffer (el Monte Carlo calcula la pérdida esperada distinto en cada caso: con buffer solo
    #      se pierde el exceso sobre el colchón, con barrera se pierde la caída completa).
    #  10% Periodicidad de pago: mensual puntúa más que trimestral, que puntúa más que semestral
    #      (mismo capital rotando y cobrando más veces reduce el riesgo de cada evento individual).
    #  10% Frecuencia de call cercana a 3 meses (preferible, pero es el factor menos importante).
    #   5% Tipo de protección de capital: buffer puntúa algo mejor que barrera a igualdad de nivel
    #      nominal, porque amortigua la cola de pérdida en vez de transferirla completa.
    rentas = [r["rentabilidad_esperada_neta"] for r in resultados_notas]
    renta_min, renta_max = min(rentas), max(rentas)
    margenes = [r["margen_sobre_inversor"] for r in resultados_notas]
    margen_min, margen_max = min(margenes), max(margenes)

    def _score_periodicidad(periodicidad):
        return {1: 1.0, 3: 0.6, 6: 0.3}.get(periodicidad, 0.5)

    def _score_call(frecuencia_call):
        # 3 meses = óptimo; cuanto más se aleja (en cualquier dirección), menos puntúa.
        return max(0.0, 1 - abs(frecuencia_call - 3) / 12)

    for r in resultados_notas:
        margen_norm = 0.5 if margen_max == margen_min else (r["margen_sobre_inversor"] - margen_min) / (margen_max - margen_min)
        # Si el margen es negativo o casi nulo, penaliza fuerte independientemente de la normalización relativa
        penalizacion_margen_bajo = 1.0 if r["margen_sobre_inversor"] > 0.05 else max(0.0, r["margen_sobre_inversor"] / 0.05)
        score = 100 * (
            0.35 * r["prob_cupon_media"]
            + 0.20 * margen_norm * penalizacion_margen_bajo
            + 0.20 * (1 - r["prob_perdida_capital"])
            + 0.10 * _score_periodicidad(r["periodicidad"])
            + 0.10 * _score_call(r["frecuencia_call"])
            + 0.05 * (1.0 if r["tipo_proteccion"] == "buffer" else 0.6)
        )
        r["score"] = max(0.0, min(100.0, score))
    resultados_notas.sort(key=lambda r: r["score"], reverse=True)

    st.markdown("---")
    st.markdown("### 📊 Resultados de la simulación")
    st.caption(
        "**Score** (0-100): 35% probabilidad de cobrar el cupón (lo más importante — se juzga por si "
        "rompe SU barrera, no por volatilidad general), 20% margen sobre la tasa que pagamos al "
        "inversor, 20% probabilidad de NO perder capital (ajustada por barrera/buffer), 10% "
        "periodicidad de pago (mensual > trimestral > semestral), 10% cercanía del call a 3 meses, "
        "5% tipo de protección (buffer > barrera). Ordenado de mejor a peor — ayuda a rankear rápido, "
        "no sustituye leer el detalle de cada nota."
    )
    df_resumen = pd.DataFrame([{
        "🏆 Score": f"{r['score']:.0f}/100", "Nota": r["nombre"], "Tickers": ", ".join(r["tickers"]), "Cupón anual": f"{r['cupon_anual']*100:.1f}%",
        "Margen vs. inversor": f"{r['margen_sobre_inversor']*100:.2f}%",
        "Pago": {1: "Mensual", 3: "Trimestral", 6: "Semestral"}[r["periodicidad"]],
        "Memoria": "Sí" if r.get("tiene_memoria") else "No",
        "Call cada": f"{r['frecuencia_call']} m",
        "Protección": "Buffer" if r["tipo_proteccion"] == "buffer" else "Barrera",
        "Prob. media cupón/periodo": f"{r['prob_cupon_media']*100:.1f}%", "Prob. call (total)": f"{r['prob_call_total']*100:.1f}%",
        "Cupones esperados": f"{r['cupones_totales_esperados']:.1f}/{r['cupones_totales_periodos']}",
        "Prob. pérdida capital": f"{r['prob_perdida_capital']*100:.1f}%", "Pérdida media si incumple": f"{r['perdida_pct_promedio']:.1f}%",
        "Rentabilidad neta esperada (anual)": f"{r['rentabilidad_esperada_neta']*100:.2f}%",
    } for r in resultados_notas])
    st.dataframe(df_resumen, use_container_width=True, hide_index=True)

    st.markdown("#### 📈 Comparativa visual")
    graficos_comparativos_notas(resultados_notas)

    # --- Ficha por compañía subyacente: a qué se dedica, noticias, precio objetivo (dato real) ---
    st.markdown("---")
    st.markdown("### 🏢 Compañías subyacentes de cada nota")
    bloques_informe = [f"# Informe de comparación de notas — {pd.Timestamp.now().strftime('%d/%m/%Y')}"]
    deriva_txt_map = {
        "neutral": "neutral (sin deriva)",
        "analistas": "consenso de analistas (por ticker)",
        "manual": "convicción manual",
    }
    for r in resultados_notas:
        st.markdown(f"#### {r['nombre']}  ·  Score {r['score']:.0f}/100")
        margen_vs_inversor = r["margen_sobre_inversor"]
        pago_txt = {1: "mensual", 3: "trimestral", 6: "semestral"}[r["periodicidad"]]
        proteccion_txt = "buffer (protege el %, pérdida solo del exceso)" if r["tipo_proteccion"] == "buffer" else "barrera (pérdida total si rompe)"
        deriva_txt = deriva_txt_map.get(r.get("modo_deriva"), "neutral (sin deriva)")
        earnings_txt = "sí" if r.get("modelar_earnings") else "no"
        memoria_txt = "con memoria (recupera cupones perdidos)" if r.get("tiene_memoria") else "sin memoria (cupón perdido, perdido)"
        st.caption(
            f"Cupón nota: {r['cupon_anual']*100:.2f}% anual, pago {pago_txt} · Pagamos al inversor: {tasa_inversor_pct*100:.2f}% anual · "
            f"Margen bruto teórico: {margen_vs_inversor*100:.2f}% · Call cada {r['frecuencia_call']} meses · Protección de capital: {proteccion_txt} · "
            f"Deriva simulación: {deriva_txt} · Shock de earnings: {earnings_txt} · Cupón: {memoria_txt} · "
            f"Cupones esperados: {r['cupones_totales_esperados']:.1f} de {r['cupones_totales_periodos']} periodos"
        )

        # Matriz de correlación real entre los tickers de esta nota (si aplica: >1 ticker y activada)
        matriz_corr = r.get("matriz_correlacion")
        if matriz_corr is not None and len(r["tickers"]) > 1 and not np.allclose(matriz_corr, np.eye(len(r["tickers"]))):
            st.caption("🔗 Correlación histórica real entre los subyacentes (12 meses), usada en el Monte Carlo:")
            df_corr = pd.DataFrame(matriz_corr, index=r["tickers"], columns=r["tickers"]).round(2)
            st.dataframe(df_corr.style.background_gradient(cmap="RdYlGn_r", vmin=-1, vmax=1).format("{:.2f}"), use_container_width=True)

        bloques_informe.append(f"### {r['nombre']} (score {r['score']:.0f}/100)")
        bloques_informe.append(
            f"Cupón nota: {r['cupon_anual']*100:.2f}% anual, pago {pago_txt}. Pagamos al inversor: {tasa_inversor_pct*100:.2f}% anual. "
            f"Margen bruto teórico: {margen_vs_inversor*100:.2f}%. Call cada {r['frecuencia_call']} meses. Protección de capital: {proteccion_txt}. "
            f"Deriva usada en la simulación: {deriva_txt}. Shock de earnings modelado: {earnings_txt}. Cupón {memoria_txt}. "
            f"Cupones totales esperados a lo largo de la vida de la nota: {r['cupones_totales_esperados']:.1f} de {r['cupones_totales_periodos']} periodos posibles."
        )
        colchones_nota = []
        colchones_cupon_nota = []
        noticias_por_ticker = {}
        for t_full, fd in zip(r["tickers_full"], r["fundamentales"]):
            ticker = t_full["ticker"]
            barrera_capital_pct = t_full["barrera_capital_pct"]
            barrera_cupon_pct = t_full["barrera_cupon_pct"]
            colchon_pct = (1 - barrera_capital_pct) * 100
            colchon_cupon_pct = (1 - barrera_cupon_pct) * 100
            colchones_nota.append(colchon_pct)
            colchones_cupon_nota.append(colchon_cupon_pct)
            with st.expander(f"📊 {fd.get('nombre') or ticker} ({ticker})", expanded=False):
                c1, c2, c3, c4 = st.columns(4)
                with c1:
                    tarjeta_kpi("Precio actual", f"${fd['precio_actual']:,.2f}" if fd.get("precio_actual") else "N/D", fd.get("sector") or "", "normal")
                with c2:
                    if fd.get("target_medio"):
                        var_target = (fd["target_medio"] / fd["precio_actual"] - 1) * 100 if fd.get("precio_actual") else None
                        sub = f"{var_target:+.1f}% vs precio actual" if var_target is not None else ""
                        tarjeta_kpi("Precio objetivo (consenso analistas)", f"${fd['target_medio']:,.2f}", sub, "positivo" if (var_target or 0) >= 0 else "negativo")
                    else:
                        tarjeta_kpi("Precio objetivo (consenso analistas)", "N/D", "Sin cobertura o dato no disponible", "normal")
                with c3:
                    tarjeta_kpi("Colchón hasta barrera de CUPÓN", f"{colchon_cupon_pct:.1f}%", f"Barrera al {barrera_cupon_pct*100:.0f}% · afecta al interés mensual/periódico", "normal")
                with c4:
                    tarjeta_kpi("Colchón hasta barrera de CAPITAL", f"{colchon_pct:.1f}%", f"Barrera al {barrera_capital_pct*100:.0f}% · solo afecta al vencimiento", "normal")
                if fd.get("target_alto") or fd.get("target_bajo"):
                    st.caption(f"Rango de analistas: ${fd.get('target_bajo', 0):,.2f} — ${fd.get('target_alto', 0):,.2f}  |  {fd.get('n_analistas', '?')} analistas  |  Recomendación consenso: {fd.get('recomendacion', 'N/D')}")
                if fd.get("aviso_analistas"):
                    st.info(f"ℹ️ {fd['aviso_analistas']}")
                c5, c6 = st.columns(2)
                with c5:
                    if fd.get("salto_medio_earnings_pct"):
                        tarjeta_kpi(
                            "Salto medio el día de earnings",
                            f"±{fd['salto_medio_earnings_pct']:.1f}%",
                            f"Medido sobre sus últimos {fd.get('n_earnings_medidos', '?')} informes trimestrales",
                            "normal",
                        )
                    else:
                        tarjeta_kpi("Salto medio el día de earnings", "N/D", "Sin datos históricos de earnings suficientes", "normal")
                with c6:
                    if fd.get("proxima_fecha_resultados"):
                        tarjeta_kpi("Próximos resultados", str(fd["proxima_fecha_resultados"])[:10], "Fecha estimada por Yahoo Finance", "normal")
                    else:
                        tarjeta_kpi("Próximos resultados", "N/D", "Sin fecha confirmada", "normal")
                with st.spinner(f"Buscando información y noticias de {ticker}..."):
                    ficha_texto = generar_ficha_empresa_ia(ticker, fd.get("nombre") or ticker, barrera_capital_pct, colchon_pct, barrera_cupon_pct, colchon_cupon_pct)
                st.markdown(_md_seguro(ficha_texto))
                st.caption("⚠️ La explicación, las noticias y la opinión de riesgo son una síntesis de IA en base a fuentes públicas — el precio, el precio objetivo y los colchones de arriba sí son datos reales/calculados.")

                # Noticias estructuradas (fuente, fecha y enlace por artículo) que avalan el análisis —
                # complementan la síntesis narrativa de arriba con las noticias concretas, citables una a una.
                st.markdown("##### 📰 Noticias que avalan este análisis")
                noticias_ticker = buscar_noticias_libre(f"{fd.get('nombre') or ticker} ({ticker})")
                noticias_por_ticker[ticker] = noticias_ticker
                _renderizar_tarjetas_noticias(noticias_ticker)

                bloques_informe.append(f"#### {fd.get('nombre') or ticker} ({ticker})")
                bloques_informe.append(
                    f"Precio actual: ${fd.get('precio_actual', 0):,.2f} · Precio objetivo consenso analistas: "
                    f"${fd.get('target_medio', 0):,.2f}" if fd.get("target_medio") else f"Precio actual: ${fd.get('precio_actual', 0):,.2f} · Sin precio objetivo de consenso disponible"
                )
                bloques_informe.append(f"Colchón hasta barrera de cupón: {colchon_cupon_pct:.1f}% (barrera al {barrera_cupon_pct*100:.0f}%). Colchón hasta barrera de capital: {colchon_pct:.1f}% (barrera al {barrera_capital_pct*100:.0f}%).")
                bloques_informe.append(ficha_texto)
        r["noticias_por_ticker"] = noticias_por_ticker

    # Recomendación de reparto: razonada por Claude, usando SOLO los números ya calculados
    with st.spinner("Generando recomendación..."):
        resumen_texto = "\n".join(
            f"- {r['nombre']} ({', '.join(r['tickers'])}): cupón {r['cupon_anual']*100:.1f}% anual, "
            f"pago {'mensual' if r['periodicidad']==1 else ('trimestral' if r['periodicidad']==3 else 'semestral')}, "
            f"call cada {r['frecuencia_call']} meses, protección de capital tipo {r['tipo_proteccion']}, "
            f"deriva de precio usada en la simulación: {deriva_txt_map.get(r.get('modo_deriva'), 'neutral')}, "
            f"shock de earnings modelado: {'sí' if r.get('modelar_earnings') else 'no'}, "
            f"cupón con memoria: {'sí, recupera cupones perdidos retroactivamente' if r.get('tiene_memoria') else 'no, cupón perdido es definitivo'}, "
            f"cupones totales esperados a lo largo de la vida: {r['cupones_totales_esperados']:.1f} de {r['cupones_totales_periodos']} periodos posibles, "
            f"score compuesto {r['score']:.0f}/100, "
            f"prob. media de cobro de cupón por periodo {r['prob_cupon_media']*100:.1f}%, prob. call total {r['prob_call_total']*100:.1f}%, "
            f"prob. pérdida de capital al vencimiento {r['prob_perdida_capital']*100:.1f}% (pérdida media si ocurre: {r['perdida_pct_promedio']:.1f}%, "
            f"ya calculada según si la protección es barrera o buffer), "
            f"rentabilidad neta esperada anualizada {r['rentabilidad_esperada_neta']*100:.2f}%, "
            f"margen bruto teórico sobre lo que pagamos al inversor: {r['margen_sobre_inversor']*100:.2f}%, "
            f"colchón hasta la barrera de capital por ticker: {', '.join(f'{tk}={(1 - t['barrera_capital_pct'])*100:.1f}%' for tk, t in zip(r['tickers'], r['tickers_full']))}, "
            f"colchón hasta la barrera de cupón por ticker: {', '.join(f'{tk}={(1 - t['barrera_cupon_pct'])*100:.1f}%' for tk, t in zip(r['tickers'], r['tickers_full']))}"
            for r in resultados_notas
        )
        api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
        try:
            resp = requests.post(
                "https://api.anthropic.com/v1/messages",
                headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
                json={
                    "model": "claude-sonnet-4-5", "max_tokens": 1400,
                    "system": (
                        "Eres un analista de notas estructuradas para un fondo que se financia captando capital de "
                        "inversores a una tasa fija y desplegándolo en estas notas — el beneficio del fondo es el "
                        "spread entre lo que cobra la nota y lo que paga al inversor. Se te dan métricas YA CALCULADAS "
                        "(probabilidades de Monte Carlo, score compuesto, rentabilidad esperada, margen bruto teórico "
                        f"sobre la tasa del inversor —que es {tasa_inversor_pct*100:.2f}% anual—, periodicidad de pago, "
                        "frecuencia de call, tipo de protección de capital, la deriva de precio usada en la simulación "
                        "(neutral / consenso de analistas / convicción manual) y si se modeló un shock de earnings, "
                        "y el colchón hasta la barrera de capital y "
                        "hasta la barrera de cupón de cada ticker) de varias notas candidatas — NO recalcules ni "
                        "inventes ningún número, solo razona sobre los que se te dan. IMPORTANTE: las probabilidades "
                        "de cobro y de pérdida YA incorporan la deriva y el shock de earnings si se indicaron — no las "
                        "corrijas mentalmente asumiendo que son 'sin dirección', solo dilo si la nota usó deriva "
                        "neutral. Aplica ESTAS reglas de decisión, "
                        "en este orden de importancia:\n\n"
                        "1) LO MÁS IMPORTANTE es la probabilidad de cobrar el cupón, porque es el evento recurrente "
                        "que genera el spread mes a mes. Esa probabilidad ya viene calculada evaluando la barrera de "
                        "CUPÓN concreta de cada ticker. El criterio correcto NUNCA es 'esta acción es muy volátil' o "
                        "'esta acción puede caer mucho' en abstracto — es si, según el estudio, una caída realista y "
                        "plausible se queda por encima de SU barrera de cupón. Si el colchón es amplio frente a lo que "
                        "manejan los analistas y la probabilidad de cobro es alta, esa nota es atractiva en ese frente "
                        "aunque otro ticker de otra nota parezca 'menos volátil' en términos generales — NO compares "
                        "acciones por cuánto pueden bajar en abstracto, compara por si rompen SU barrera concreta.\n\n"
                        "2) El margen sobre la tasa del inversor tiene que ser SUSTANCIAL para justificar el riesgo — "
                        "un cupón que apenas supera lo que pagamos al inversor (ej. margen de 10 puntos) NO es "
                        "atractivo aunque la probabilidad de cobro sea alta, porque el fondo absorbe TODA la pérdida "
                        "de capital si la nota falla, mientras que al inversor se le sigue pagando igual. No arriesgues "
                        "capital del fondo por un margen pequeño.\n\n"
                        "3) Hay DOS barreras distintas y hay que juzgarlas por separado — la de CUPÓN (recurrente, "
                        "normalmente más exigente) y la de CAPITAL (una sola vez, al vencimiento, normalmente más "
                        "laxa). Dentro de la barrera de capital, distingue explícitamente BARRERA (si rompe, el fondo "
                        "asume la caída completa del peor activo, sin suelo) de BUFFER (el fondo está protegido hasta "
                        "el % indicado y solo pierde el exceso) — a igualdad de nivel nominal, el buffer es "
                        "estructuralmente menos arriesgado y debe valorarse mejor, y la pérdida media si incumple que "
                        "se te da YA refleja esa diferencia.\n\n"
                        "4) La periodicidad de pago importa: mensual es preferible a trimestral, que es preferible a "
                        "semestral (capital rotando más rápido, menor riesgo por evento individual) — factor "
                        "secundario, no decisivo por sí solo.\n\n"
                        "5) La frecuencia de call cercana a 3 meses es preferible — factor MENOR, casi un desempate, "
                        "no debe pesar más que el cupón, el margen o el riesgo de capital.\n\n"
                        "ESTRUCTURA OBLIGATORIA DE LA RESPUESTA — esto es tan importante como el análisis en sí: "
                        "la PRIMERA línea de toda tu respuesta, antes de cualquier título, tabla o punto de análisis, "
                        "tiene que ser el veredicto final resumido en una frase, así (formato exacto): "
                        "'**Conclusión: <frase breve, ej. \"Rechazar ambas notas\" o \"Invertir $X en Nota Y\">**'. "
                        "A partir de ahí, absolutamente TODO el análisis punto por punto que escribas tiene que sonar "
                        "coherente con esa conclusión desde la primera frase de cada punto — nunca escribas un punto "
                        "que, leído solo y sin el resto del informe, parezca apuntar a una conclusión distinta de la "
                        "que ya diste en la primera línea. Si un criterio concreto favorece a la nota que NO vas a "
                        "recomendar, preséntalo siempre como una observación parcial ligada de inmediato a por qué NO "
                        "cambia la conclusión general (ej. 'Nota 2 cobra el cupón con más probabilidad que Nota 1, "
                        "pero esto no basta para recomendarla porque su margen es insuficiente frente al riesgo que "
                        "asume el fondo') — nunca lo presentes como si fuera el argumento que va ganando, aunque el "
                        "resto del informe lo vaya a matizar después. Alguien que lea solo la primera línea, o solo "
                        "un punto intermedio, o solo el cierre, tiene que llevarse siempre el mismo mensaje — este "
                        "informe se reparte tal cual en PDF y en PowerPoint para el directorio, y ambos usan este "
                        "mismo texto, así que una contradicción interna aquí se ve en los dos documentos a la vez.\n\n"
                        "Puede recomendar repartir entre varias notas o poner todo el capital en una sola si está "
                        f"claramente justificado. Capital disponible total: ${capital_disponible:,.2f}. Da una "
                        "recomendación de reparto de capital en dólares para cada nota, con el razonamiento concreto "
                        "siguiendo el orden de prioridad anterior, en español, conciso pero completo.\n\n"
                        "Termina tu respuesta con una última línea, sola, en este formato EXACTO (sin negrita, sin "
                        "texto adicional en esa línea) para que un programa la pueda leer automáticamente:\n"
                        "VEREDICTO: <nombre exacto de la nota recomendada, o 'ninguna' si tu recomendación final es "
                        "no asignar capital a ninguna nota>\n"
                        "Esta línea VEREDICTO tiene que decir exactamente lo mismo que tu 'Conclusión' de la primera "
                        "línea — si recomiendas rechazar todas, pon 'ninguna'; si repartes entre varias sin que una "
                        "domine claramente, pon también 'ninguna'."
                    ),
                    "messages": [{"role": "user", "content": f"Notas candidatas:\n{resumen_texto}\n\nRecomienda cómo repartir el capital disponible."}],
                },
                timeout=60,
            )
            data = resp.json()
            texto_recomendacion = "".join(b.get("text", "") for b in data.get("content", []) if b.get("type") == "text")
        except Exception as e:
            texto_recomendacion = f"[No se pudo generar la recomendación: {e}]"

    # --- Veredicto estructurado, extraído de la última línea de la recomendación (ver prompt
    # arriba) — esto es lo que decide el banner del memo, NO el ganador del Score por separado,
    # para que la portada del PowerPoint nunca contradiga lo que dice el texto de la recomendación. ---
    veredicto_nota_recomendada = None
    _match_veredicto = re.search(r"VEREDICTO:\s*(.+?)\s*$", (texto_recomendacion or "").strip(), re.IGNORECASE | re.MULTILINE)
    if _match_veredicto:
        _veredicto_txt = _match_veredicto.group(1).strip()
        texto_recomendacion = texto_recomendacion[:_match_veredicto.start()].rstrip()  # se quita del texto visible, es solo para el banner
        if _veredicto_txt.lower() not in ("ninguna", "ningún", "ninguno", "ninguna nota", "n/a"):
            veredicto_nota_recomendada = next(
                (r["nombre"] for r in resultados_notas if r["nombre"].strip().lower() == _veredicto_txt.strip().lower()),
                _veredicto_txt,
            )

    st.markdown("### 🎯 Recomendación de reparto")
    st.markdown(_md_seguro(texto_recomendacion) or "No se pudo generar una recomendación.")
    st.caption("Simulación Monte Carlo con 5.000 escenarios por nota, volatilidad histórica de 12 meses, con correlación real entre tickers cuando está activada (ver matriz por nota arriba).")

    bloques_informe.append("### 🎯 Recomendación de reparto")
    bloques_informe.append(texto_recomendacion or "No se pudo generar una recomendación.")

    # --- Descarga del informe completo: PDF + PowerPoint juntos en un único clic ---
    # Un solo st.download_button solo puede entregar un archivo, así que se empaquetan los
    # dos (PDF + PPTX) en un .zip.
    st.markdown("---")
    with st.spinner("Preparando informe (PDF + PowerPoint)..."):
        pdf_bytes = _generar_informe_comparador_pdf(resultados_notas, texto_recomendacion, capital_disponible, tasa_inversor_pct, bloques_informe)
        zip_buffer = BytesIO()
        try:
            pptx_bytes = generar_memo_directorio_notas_pptx(
                resultados_notas, texto_recomendacion, capital_disponible, tasa_inversor_pct,
                veredicto_nota_recomendada=veredicto_nota_recomendada,
            )
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                zf.writestr(f"comparador_notas_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.pdf", pdf_bytes)
                zf.writestr(f"memo_directorio_notas_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.pptx", pptx_bytes)
            st.download_button(
                "⬇️ Descargar informe (PDF + PowerPoint)",
                data=zip_buffer.getvalue(),
                file_name=f"informe_comparador_notas_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.zip",
                mime="application/zip",
                type="primary",
            )
        except Exception as e:
            st.error(f"No se pudo generar el PowerPoint: {e}. Descargando solo el PDF.")
            st.download_button(
                "⬇️ Descargar este informe en PDF",
                data=pdf_bytes,
                file_name=f"comparador_notas_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.pdf",
                mime="application/pdf",
            )



def _tab_analisis_nota_existente(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, df_calls: pd.DataFrame):
    st.caption(
        "Elige una nota que ya tengas en el Excel — se cogen automáticamente sus tickers, barreras, "
        "cupón y **las fechas reales de tu calendario** (no fechas sintéticas), y se hace el análisis "
        "completo: fundamentales, probabilidades por Monte Carlo y noticias recientes de cada ticker."
    )

    notas_existentes = sorted(
        int(n) for n in pd.to_numeric(df_control.get("nota", pd.Series(dtype=float)), errors="coerce").dropna().unique()
    ) if df_control is not None and not df_control.empty else []
    if not notas_existentes:
        st.info("No hay notas en CONTROL_NOTAS todavía.")
        return

    numero_nota = st.selectbox("Nota a analizar", notas_existentes, key="analisis_nota_numero")

    if not st.button(f"🔬 Analizar Nota {numero_nota} en profundidad", type="primary"):
        return

    control_nota = df_control[pd.to_numeric(df_control["nota"], errors="coerce") == numero_nota]
    if control_nota.empty:
        st.error("No hay datos de tickers/barreras para esta nota en CONTROL_NOTAS.")
        return

    nombre_nota = f"NOTA_{numero_nota:02d}"
    inv_nota = df_inv[df_inv.get("nombre_activo", pd.Series(dtype=str)).astype(str).str.upper().str.replace(" ", "_") == nombre_nota]
    cupon_anual = float(inv_nota["interes_nota_anual"].dropna().iloc[0]) if not inv_nota.empty and inv_nota["interes_nota_anual"].notna().any() else None

    # Construir eventos REALES desde el calendario (no sintéticos)
    hoy = pd.Timestamp.today().normalize()
    cal_nota = df_cal[pd.to_numeric(df_cal.get("nota"), errors="coerce") == numero_nota] if df_cal is not None and not df_cal.empty else pd.DataFrame()
    eventos = []
    obs_futuras = cal_nota[(cal_nota["tipo_evento"] == "OBSERVACION") & (cal_nota["fecha"] >= hoy)] if not cal_nota.empty else pd.DataFrame()
    for _, r in obs_futuras.iterrows():
        dias = (r["fecha"] - hoy).days
        if dias > 0:
            eventos.append({"dias": dias, "tipo": "cupon"})

    calls_nota = df_calls[pd.to_numeric(df_calls.get("nota"), errors="coerce") == numero_nota] if df_calls is not None and not df_calls.empty else pd.DataFrame()
    calls_futuras = pd.DataFrame()
    if not calls_nota.empty and "fecha_call" in calls_nota.columns:
        calls_nota = calls_nota.copy()
        calls_nota["fecha_call"] = pd.to_datetime(calls_nota["fecha_call"], errors="coerce", dayfirst=True)
        calls_futuras = calls_nota[calls_nota["fecha_call"] >= hoy]
        for _, r in calls_futuras.iterrows():
            dias = (r["fecha_call"] - hoy).days
            if dias > 0:
                eventos.append({"dias": dias, "tipo": "call"})

    if not eventos:
        st.warning("No hay eventos futuros en el calendario de esta nota (¿ya venció o está descalendarizada?).")
        return

    dias_vencimiento = max(e["dias"] for e in eventos)
    eventos.append({"dias": dias_vencimiento, "tipo": "vencimiento"})

    st.markdown("---")
    st.markdown(f"## Nota {numero_nota}" + (f" — Cupón {cupon_anual*100:.2f}% anual" if cupon_anual else ""))
    st.caption(f"{len(obs_futuras)} observaciones futuras, {len(calls_futuras)} fechas de posible call, vencimiento en {dias_vencimiento} días.")

    tickers_datos = []
    datos_por_ticker = {}
    for _, fila in control_nota.iterrows():
        ticker = str(fila.get("ticker", "")).strip().upper()
        if not ticker:
            continue
        with st.spinner(f"Analizando {ticker}..."):
            fd = obtener_datos_fundamentales(ticker)
        datos_por_ticker[ticker] = fd
        if fd.get("precio_actual") is None or fd.get("volatilidad_anual_pct") is None:
            st.error(f"No se pudo obtener precio/volatilidad de {ticker} — se omite del cálculo de probabilidades.")
            continue
        precio_inicial = float(fila.get("precio_compra")) if pd.notna(fila.get("precio_compra")) else fd["precio_actual"]
        tickers_datos.append({
            "ticker": ticker, "precio_actual": fd["precio_actual"], "precio_inicial": precio_inicial,
            "volatilidad_anual_pct": fd["volatilidad_anual_pct"],
            "barrera_cupon_pct": float(fila.get("barrera_cupon", 0.6)),
            "barrera_capital_pct": float(fila.get("barrera_capital", fila.get("barrera_cupon", 0.6))),
            "call_level_pct": float(fila.get("call_level", 1.0)),
        })

    # --- Resumen por ticker: fundamentales + posición vs barreras + noticias ---
    for _, fila in control_nota.iterrows():
        ticker = str(fila.get("ticker", "")).strip().upper()
        fd = datos_por_ticker.get(ticker, {})
        if not ticker:
            continue
        st.markdown(f"### {fd.get('nombre') or ticker} ({ticker})")

        if fd.get("error"):
            st.error(f"⚠️ {fd['error']}")
            if st.button("🔄 Reintentar este ticker", key=f"retry_nota_{ticker}"):
                obtener_datos_fundamentales.clear()
                st.rerun()

        c1, c2, c3, c4 = st.columns(4)
        precio_compra = fila.get("precio_compra")
        with c1:
            tarjeta_kpi("Precio compra (nota)", f"${precio_compra:,.2f}" if pd.notna(precio_compra) else "N/D", "", "normal")
        with c2:
            var_actual = (fd["precio_actual"] / precio_compra - 1) * 100 if fd.get("precio_actual") and pd.notna(precio_compra) else None
            tarjeta_kpi("Precio actual", f"${fd.get('precio_actual', 0):,.2f}" if fd.get("precio_actual") else "N/D",
                        f"{var_actual:+.1f}% vs compra" if var_actual is not None else "", "positivo" if (var_actual or 0) >= 0 else "negativo")
        with c3:
            tarjeta_kpi("Precio objetivo analistas", f"${fd.get('target_medio', 0):,.2f}" if fd.get("target_medio") else "N/D",
                        f"{fd.get('n_analistas', '?')} analistas" if fd.get("target_medio") else "Ver aviso", "normal")
        with c4:
            tarjeta_kpi("Volatilidad anual", f"{fd.get('volatilidad_anual_pct', 0):.1f}%" if fd.get("volatilidad_anual_pct") else "N/D", "", "riesgo" if (fd.get("volatilidad_anual_pct") or 0) > 40 else "normal")
        if fd.get("aviso_analistas"):
            st.info(f"ℹ️ {fd['aviso_analistas']}")
        if fd.get("proxima_fecha_resultados"):
            st.caption(f"📅 Próxima fecha de resultados: {fd['proxima_fecha_resultados']}")

        with st.spinner(f"Buscando noticias de {ticker}..."):
            noticias = obtener_resumen_noticias_ia(ticker, fd.get("nombre") or ticker)
        st.markdown(_md_seguro(noticias))
        st.markdown("---")

    # --- Probabilidades Monte Carlo para la nota completa (worst-of) ---
    if tickers_datos:
        tiene_memoria_nota = False
        if not control_nota.empty and "tiene_memoria" in control_nota.columns:
            val_memoria = control_nota["tiene_memoria"].dropna()
            if not val_memoria.empty:
                tiene_memoria_nota = str(val_memoria.iloc[0]).strip().upper() == "SI"
        st.caption(f"Cupón con memoria (dato de CONTROL_NOTAS): {'Sí' if tiene_memoria_nota else 'No'}")
        with st.spinner("Ejecutando simulación Monte Carlo..."):
            sim = simular_montecarlo_nota(tickers_datos, eventos, tiene_memoria=tiene_memoria_nota)
        st.markdown("### 🎲 Probabilidades (simulación Monte Carlo, worst-of)")
        eventos_cupon = [e for e in sim["eventos"] if e["tipo"] == "cupon"]
        eventos_call = [e for e in sim["eventos"] if e["tipo"] == "call"]
        c1, c2, c3 = st.columns(3)
        with c1:
            tarjeta_kpi("Prob. próxima observación cumpla barrera", f"{eventos_cupon[0]['probabilidad']*100:.1f}%" if eventos_cupon else "N/D", "Cobro del próximo cupón", "positivo")
        with c2:
            tarjeta_kpi("Prob. call total (todas las fechas)", f"{sim['probabilidad_call_total']*100:.1f}%", f"{len(eventos_call)} fechas de call evaluadas", "normal")
        with c3:
            tarjeta_kpi("Prob. pérdida de capital al vencimiento", f"{sim['probabilidad_perdida_capital']*100:.1f}%",
                        f"Pérdida media si ocurre: {sim['perdida_pct_promedio_si_incumple']:.1f}%", "riesgo" if sim["probabilidad_perdida_capital"] > 0.15 else "normal")

        if len(eventos_cupon) > 1:
            df_prob_cupon = pd.DataFrame([{"Días": e["dias"], "Prob. cupón (%)": round(e["probabilidad"] * 100, 1)} for e in eventos_cupon])
            df_prob_cupon = df_prob_cupon.sort_values("Días").reset_index(drop=True)
            st.line_chart(df_prob_cupon, x="Días", y="Prob. cupón (%)", height=300)
            with st.expander("Ver datos exactos de la gráfica"):
                st.dataframe(df_prob_cupon, use_container_width=True, hide_index=True)
        if tiene_memoria_nota:
            st.metric("Cupones totales esperados (con efecto retroactivo de memoria)",
                      f"{sim['cupones_totales_esperados']:.1f} de {len(eventos_cupon)} periodos posibles")
        st.caption(
            "⚠️ Estimación con volatilidad histórica, sin correlación entre tickers — no es una predicción, es una "
            "probabilidad bajo supuestos. Esta pantalla usa deriva neutral (sin dirección) y no modela shocks de "
            "earnings; para esas opciones usá el Comparador de Notas."
        )
    else:
        st.error(
            "⚠️ No se pudo ejecutar la simulación Monte Carlo porque no se consiguió precio/volatilidad "
            "de ningún ticker de esta nota (fallo de conexión con Yahoo Finance). Revisa los avisos de arriba "
            "y prueba a reintentar en unos minutos."
        )


def obtener_control_notas_activas(df_inv: pd.DataFrame, df_control: pd.DataFrame) -> pd.DataFrame:
    """Devuelve CONTROL_NOTAS excluyendo las notas ya cerradas por call (o call final/vencimiento).
    Extraído como helper reutilizable para que el semáforo, el dashboard de noticias y el
    calendario de earnings usen siempre el mismo criterio de 'nota activa'."""
    if df_control is None or df_control.empty or "nota" not in df_control.columns:
        return df_control
    hoy = pd.Timestamp.today().normalize()
    notas_con_call = set()
    if df_inv is not None and "nombre_activo" in df_inv.columns and "motivo" in df_inv.columns:
        motivo_normalizado = df_inv["motivo"].astype(str).str.lower().str.strip()
        # "call" = llamada anticipada por el emisor. "call final" = vencimiento natural sin
        # call previo, tratado igual que un call (capital reinvertido, nota cerrada).
        notas_df = df_inv[motivo_normalizado.isin(["call", "call final"])].copy()
        notas_df["fecha_final_inversion"] = pd.to_datetime(notas_df["fecha_final_inversion"], errors="coerce")
        notas_df["nota_num"] = notas_df["nombre_activo"].apply(extraer_numero_nota)
        for nota_num, grupo in notas_df.groupby("nota_num"):
            if pd.notna(nota_num) and (grupo["fecha_final_inversion"].notna() & (grupo["fecha_final_inversion"] <= hoy)).all():
                notas_con_call.add(int(nota_num))
    return df_control[~df_control["nota"].isin(notas_con_call)].copy() if notas_con_call else df_control


@st.cache_data(show_spinner=False, ttl=21600)
def obtener_screening_noticias_ia(tickers_nombres: tuple) -> dict:
    """Escaneo RÁPIDO y barato de varias compañías a la vez: para cada una, decide si hay o no
    hay noticias relevantes recientes (sin desarrollarlas todavía). Sirve para elegir luego solo
    las compañías que merece la pena leer en profundidad con obtener_resumen_noticias_ia, y así
    gastar muchos menos tokens que haciendo un resumen completo de todas de entrada.
    tickers_nombres: tupla de tuplas (ticker, nombre_compania), para que sea cacheable."""
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
    lista_texto = "\n".join(f"- {t} ({n})" for t, n in tickers_nombres)
    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
            json={
                "model": "claude-sonnet-4-5",
                "max_tokens": 1200,
                "tools": [{"type": "web_search_20250305", "name": "web_search", "max_uses": 15}],
                "system": (
                    "Haces un ESCANEO RÁPIDO de varias compañías para un inversor en notas estructuradas. "
                    "Para cada compañía de la lista, busca brevemente si ha habido noticias RELEVANTES en las "
                    "últimas 2-4 semanas: sorpresas de resultados, cambios de rating de analistas, problemas "
                    "legales/regulatorios, fusiones/adquisiciones, cambios de dirección, profit warnings. "
                    "No profundices ni escribas resúmenes largos — solo detecta si hay señal o no. "
                    "Responde ÚNICAMENTE con una línea por compañía, en este formato exacto y nada más "
                    "(sin markdown, sin explicación adicional):\n"
                    "TICKER|SI|razón en máximo 12 palabras\n"
                    "TICKER|NO|\n"
                    "Usa SI solo si de verdad encontraste algo concreto y reciente; si no encuentras nada claro, usa NO."
                ),
                "messages": [{"role": "user", "content": f"Escanea estas compañías:\n{lista_texto}"}],
            },
            timeout=90,
        )
        data = resp.json()
        if data.get("type") == "error" or data.get("error"):
            return {"_error": data.get("error", {}).get("message", str(data))}
        contenido = data.get("content", [])
        texto = "".join(b.get("text", "") for b in contenido if b.get("type") == "text")
        resultado = {}
        for linea in texto.strip().splitlines():
            partes = [p.strip() for p in linea.split("|")]
            if len(partes) >= 2:
                tk = partes[0].strip().upper().lstrip("-").strip()
                relevante = partes[1].strip().upper() == "SI"
                razon = partes[2].strip() if len(partes) >= 3 else ""
                resultado[tk] = {"relevante": relevante, "razon": razon}
        return resultado
    except Exception as e:
        return {"_error": str(e)}


@st.cache_data(show_spinner=False, ttl=1800)
def buscar_noticias_libre(query: str) -> list:
    """Búsqueda libre de noticias/investigación vía Claude + web_search — 'preguntá lo que
    quieras', a diferencia de obtener_resumen_noticias_ia (que resume en texto libre y solo
    funciona atado a un ticker de CONTROL_NOTAS). Esto devuelve una lista estructurada con
    título, resumen, fuente, URL y fecha, pensada para enlazar directo a cada artículo/paper.
    Cacheado 30 min (más corto que las 6h de las noticias por ticker): una búsqueda libre es
    más probable que se repita/reformule en la misma sesión."""
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "") or st.secrets.get("anthropic", {}).get("api_key", "")
    prompt = f"""Busca en la web las noticias, informes o documentos más relevantes y recientes
(últimos 30 días salvo que la pregunta pida explícitamente algo histórico) sobre: {query}

Si "{query}" parece un ticker bursátil (letras mayúsculas cortas, ej. HOOD, AAPL, TSLA) trata la
búsqueda como noticias sobre ESA COMPAÑÍA cotizada concreta, no sobre el significado literal de
la palabra. Busca varias veces con términos distintos si la primera búsqueda no da resultados
claros antes de rendirte.

Devuelve ÚNICAMENTE un array JSON (nada de texto antes o después, sin backticks de markdown),
con como máximo 8 elementos, cada uno con este formato exacto:
[{{
  "titular": "...",
  "resumen": "2-3 frases en español, en tus propias palabras, nunca cites texto literal de la fuente",
  "fuente": "nombre del medio o autor",
  "url": "https://...",
  "fecha": "YYYY-MM-DD (fecha real de publicación; si no la sabes con certeza, tu mejor estimación)",
  "prioridad": "alta | media | baja (alta = mueve materialmente la cotización o la tesis de inversión; media = relevante pero no decisivo; baja = contexto o ruido menor)"
}}]

Si tras buscar de verdad no encuentras nada relevante, devuelve un array vacío [] — pero solo
después de haber intentado varias búsquedas, no a la primera."""

    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"},
            json={
                "model": "claude-sonnet-4-5",
                "max_tokens": 2500,
                "tools": [{"type": "web_search_20250305", "name": "web_search", "max_uses": 8}],
                "messages": [{"role": "user", "content": prompt}],
            },
            timeout=60,
        )
        data = resp.json()
    except Exception as e:
        return [{"_error": f"Error de conexión: {e}"}]

    if data.get("type") == "error" or data.get("error"):
        msg = data.get("error", {}).get("message", str(data))
        return [{"_error": f"Error de la API: {msg}. Si menciona permisos o 'web_search', puede que "
                            f"la búsqueda web no esté activada para esta cuenta en la Consola de Anthropic."}]

    contenido = data.get("content", [])
    texto = "".join(b.get("text", "") for b in contenido if b.get("type") == "text").strip()
    if not texto:
        return []

    texto_limpio = re.sub(r"^```(?:json)?", "", texto).strip()
    texto_limpio = re.sub(r"```$", "", texto_limpio).strip()
    if not texto_limpio.startswith("["):
        inicio, fin = texto_limpio.find("["), texto_limpio.rfind("]")
        if inicio != -1 and fin != -1 and fin > inicio:
            texto_limpio = texto_limpio[inicio:fin + 1]

    try:
        resultados = json.loads(texto_limpio)
        if not isinstance(resultados, list):
            raise ValueError("la respuesta no es una lista")
    except (json.JSONDecodeError, ValueError) as e:
        return [{"_error": f"No se pudo interpretar la respuesta de la IA ({e}). Probá reformular la búsqueda."}]

    orden_prioridad = {"alta": 0, "media": 1, "baja": 2}

    def _clave_orden(n):
        fecha = n.get("fecha") or ""
        prioridad = orden_prioridad.get(str(n.get("prioridad", "")).lower(), 1)
        return (fecha, -prioridad)

    return sorted(resultados, key=_clave_orden, reverse=True)


_COLOR_FONDO_PRIORIDAD = {"alta": "#fee2e2", "media": "#fef3c7", "baja": "#e5e7eb"}
_COLOR_TEXTO_PRIORIDAD = {"alta": "#7f1d1d", "media": "#78350f", "baja": "#374151"}


def _renderizar_tarjetas_noticias(resultados: list):
    if resultados and isinstance(resultados[0], dict) and "_error" in resultados[0]:
        st.error(f"⚠️ {resultados[0]['_error']}")
        return
    if not resultados:
        st.info("No se encontraron resultados relevantes tras buscar. Probá reformular la pregunta.")
        return
    for n in resultados:
        prioridad = str(n.get("prioridad", "media")).lower()
        color_fondo = _COLOR_FONDO_PRIORIDAD.get(prioridad, "#e5e7eb")
        color_texto = _COLOR_TEXTO_PRIORIDAD.get(prioridad, "#374151")
        with st.container(border=True):
            c1, c2 = st.columns([5, 1])
            with c1:
                st.markdown(f"**{_md_seguro(n.get('titular', 'Sin título'))}**")
            with c2:
                st.markdown(
                    f"<div style='text-align:right'><span style='background:{color_fondo};color:{color_texto};"
                    f"padding:2px 10px;border-radius:10px;font-size:0.75em;font-weight:700'>"
                    f"{prioridad.upper()}</span></div>",
                    unsafe_allow_html=True,
                )
            st.caption(f"{n.get('fuente', 'Fuente desconocida')} · {n.get('fecha', 'Sin fecha')}")
            st.write(_md_seguro(n.get("resumen", "")))
            if n.get("url"):
                st.link_button("🔗 Leer la fuente completa", n["url"])
    st.caption("⚠️ Resultados generados por IA a partir de búsqueda web pública — no son recomendaciones de inversión.")


def _widget_busqueda_libre_noticias(key_prefix: str):
    """Widget reutilizable de 'preguntá lo que quieras' — se usa tanto en la sección de
    Noticias del menú de administración como en el portal de cada inversor, por eso todas
    las keys de Streamlit llevan key_prefix (para no chocar si ambas viven en la misma sesión,
    ej. Yuri viendo el portal de un inversor)."""
    st.caption(
        "Preguntá lo que quieras — un ticker, una compañía, un sector, un tema de mercado. "
        "Cada resultado enlaza directo a la fuente (noticia, informe, comunicado)."
    )

    with st.form(f"form_busqueda_libre_noticias_{key_prefix}"):
        query = st.text_input(
            "Buscar",
            placeholder="Ej: 'resultados trimestrales de Nvidia', 'HOOD noticias regulatorias', "
                        "'perspectivas del sector bancario europeo'...",
            label_visibility="collapsed",
        )
        enviar = st.form_submit_button("🔍 Buscar", type="primary")

    if enviar:
        if not query or not query.strip():
            st.warning("Escribí algo para buscar.")
        else:
            with st.spinner(f"Buscando '{query.strip()}'..."):
                resultados = buscar_noticias_libre(query.strip())
            _renderizar_tarjetas_noticias(resultados)


def seccion_noticias():
    """Sección de nivel superior en el menú principal — búsqueda libre de noticias/research."""
    st.header("📰 Noticias")
    _widget_busqueda_libre_noticias(key_prefix="admin")


def _tab_dashboard_noticias(df_inv: pd.DataFrame, df_control: pd.DataFrame):
    st.caption(
        "Noticias recientes de las compañías subyacentes de tus notas activas. Para gastar menos tokens de "
        "API, primero puedes hacer un **escaneo rápido** que solo detecta qué compañías tienen novedades "
        "relevantes, y luego pides el resumen completo solo de esas — en vez de buscar a fondo en todas de entrada."
    )
    control_activo = obtener_control_notas_activas(df_inv, df_control)
    if control_activo is None or control_activo.empty or "ticker" not in control_activo.columns:
        st.info("No hay tickers en CONTROL_NOTAS para consultar noticias.")
        return

    tickers_unicos = sorted(control_activo["ticker"].dropna().astype(str).str.strip().str.upper().unique())
    tickers_unicos = [t for t in tickers_unicos if t]
    if not tickers_unicos:
        st.info("No hay tickers válidos en CONTROL_NOTAS.")
        return

    def _notas_de(ticker):
        return sorted(int(n) for n in pd.to_numeric(control_activo[control_activo["ticker"].astype(str).str.upper() == ticker]["nota"], errors="coerce").dropna().unique())

    etiquetas = {t: f"{t} (Nota {', '.join(str(n) for n in _notas_de(t))})" for t in tickers_unicos}
    seleccion = st.multiselect(
        "Compañías a incluir (añade las que quieras)",
        options=tickers_unicos,
        default=[],
        format_func=lambda t: etiquetas.get(t, t),
        help="Vacío por defecto — añade solo las compañías que te interesen, o usa el botón de 'todas' de abajo.",
    )

    col1, col2 = st.columns(2)
    escanear = col1.button("🔍 Escanear novedades (barato)", type="primary", disabled=not seleccion)
    todas_completas = col2.button("📰 Noticias completas de todas las compañías")

    if not seleccion and not todas_completas:
        st.info("Añade compañías arriba para escanear/leer solo esas, o pulsa el botón de todas.")

    if escanear:
        with st.spinner(f"Escaneando {len(seleccion)} compañía(s)..."):
            tickers_nombres = tuple((t, obtener_datos_fundamentales(t).get("nombre") or t) for t in seleccion)
            screening = obtener_screening_noticias_ia(tickers_nombres)
        st.session_state["screening_noticias"] = screening

    screening = st.session_state.get("screening_noticias")

    if screening and "_error" in screening:
        st.error(f"⚠️ Falló el escaneo: {screening['_error']}")
    elif screening:
        con_novedades = [t for t in seleccion if screening.get(t, {}).get("relevante")]
        sin_novedades = [t for t in seleccion if t in screening and not screening.get(t, {}).get("relevante")]
        no_encontrados = [t for t in seleccion if t not in screening]

        if con_novedades:
            st.warning(f"⚠️ {len(con_novedades)} compañía(s) con novedades relevantes:")
            for t in con_novedades:
                st.markdown(f"- **{t}** — {screening[t].get('razon') or 'ver detalle'}")
        else:
            st.success("Sin novedades relevantes detectadas en las compañías escaneadas.")
        if sin_novedades:
            st.caption(f"Sin novedades destacables: {', '.join(sin_novedades)}")
        if no_encontrados:
            st.caption(f"No se pudo escanear: {', '.join(no_encontrados)}")

        if con_novedades:
            elegidas_profundizar = st.multiselect(
                "¿De cuáles quieres leer el resumen completo?",
                options=seleccion,
                default=con_novedades,
                format_func=lambda t: etiquetas.get(t, t),
                key="elegidas_profundizar_noticias",
            )
            if st.button("📖 Leer resumen completo de las elegidas", type="primary") and elegidas_profundizar:
                _mostrar_noticias_completas(elegidas_profundizar, control_activo, _notas_de)

    if todas_completas:
        _mostrar_noticias_completas(tickers_unicos, control_activo, _notas_de)


def _mostrar_noticias_completas(tickers: list, control_activo: pd.DataFrame, _notas_de):
    barra = st.progress(0.0)
    for i, ticker in enumerate(tickers):
        with st.spinner(f"Buscando noticias en profundidad de {ticker}..."):
            fd = obtener_datos_fundamentales(ticker)
            nombre = fd.get("nombre") or ticker
            noticias = obtener_resumen_noticias_ia(ticker, nombre)
        notas_de_este_ticker = _notas_de(ticker)
        titulo = f"{nombre} ({ticker}) — Nota{'s' if len(notas_de_este_ticker) != 1 else ''} {', '.join(str(n) for n in notas_de_este_ticker)}"
        with st.expander(titulo, expanded=True):
            if fd.get("proxima_fecha_resultados"):
                st.caption(f"📅 Próxima fecha de resultados: {fd['proxima_fecha_resultados']}")
            st.markdown(_md_seguro(noticias))
        barra.progress((i + 1) / len(tickers))
    st.caption("⚠️ Resúmenes generados por IA a partir de noticias públicas — no son recomendaciones de inversión ni predicciones de precio.")


def _tab_calendario_earnings(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame):
    st.caption(
        "Próximas fechas de resultados (earnings) de las compañías subyacentes de tus notas activas, "
        "cruzadas con la próxima fecha de observación de cada nota. Un earnings justo antes de una "
        "observación es una señal de alerta real: la volatilidad suele dispararse alrededor de esos días."
    )
    control_activo = obtener_control_notas_activas(df_inv, df_control)
    if control_activo is None or control_activo.empty or "ticker" not in control_activo.columns:
        st.info("No hay tickers en CONTROL_NOTAS para consultar earnings.")
        return

    if not st.button("📅 Consultar calendario de earnings", type="primary"):
        st.info("Pulsa el botón para consultar las próximas fechas de resultados vía Yahoo Finance.")
        return

    hoy = pd.Timestamp.today().normalize()
    filas = []
    tickers_unicos = sorted(control_activo["ticker"].dropna().astype(str).str.strip().str.upper().unique())
    tickers_unicos = [t for t in tickers_unicos if t]
    barra = st.progress(0.0)
    for i, ticker in enumerate(tickers_unicos):
        with st.spinner(f"Consultando {ticker}..."):
            fd = obtener_datos_fundamentales(ticker)
        fecha_earnings = fd.get("proxima_fecha_resultados")
        fecha_earnings_ts = pd.to_datetime(fecha_earnings, errors="coerce") if fecha_earnings else pd.NaT

        notas_ticker = pd.to_numeric(control_activo[control_activo["ticker"].astype(str).str.upper() == ticker]["nota"], errors="coerce").dropna().unique()
        for nota_num in notas_ticker:
            nota_num = int(nota_num)
            fecha_obs = proximo_evento_nota(df_cal, nota_num, "OBSERVACION") if df_cal is not None and not df_cal.empty else None
            fecha_obs_ts = pd.to_datetime(fecha_obs, errors="coerce") if fecha_obs is not None else pd.NaT

            dias_hasta_earnings = (fecha_earnings_ts - hoy).days if pd.notna(fecha_earnings_ts) else None
            alerta_cruce = ""
            if pd.notna(fecha_earnings_ts) and pd.notna(fecha_obs_ts):
                # Earnings cae DENTRO de los 10 días previos a la observación → alerta.
                dias_earnings_a_obs = (fecha_obs_ts - fecha_earnings_ts).days
                if 0 <= dias_earnings_a_obs <= 10:
                    alerta_cruce = "⚠️ Earnings justo antes de observación"

            filas.append({
                "ticker": ticker,
                "compañía": fd.get("nombre") or ticker,
                "nota": nota_num,
                "próximo earnings": fecha_earnings_ts.strftime("%d/%m/%Y") if pd.notna(fecha_earnings_ts) else "Sin dato",
                "días hasta earnings": dias_hasta_earnings if dias_hasta_earnings is not None else "N/D",
                "próxima observación": fecha_obs_ts.strftime("%d/%m/%Y") if pd.notna(fecha_obs_ts) else "Sin dato",
                "alerta": alerta_cruce,
            })
        barra.progress((i + 1) / len(tickers_unicos))

    if not filas:
        st.info("No se encontraron datos de earnings para los tickers de tus notas activas.")
        return

    df_earnings = pd.DataFrame(filas)
    df_earnings["_orden"] = pd.to_datetime(df_earnings["próximo earnings"], errors="coerce", dayfirst=True)
    df_earnings = df_earnings.sort_values("_orden", na_position="last").drop(columns=["_orden"]).reset_index(drop=True)

    n_alertas = int((df_earnings["alerta"] != "").sum())
    if n_alertas:
        st.error(f"⚠️ {n_alertas} caso(s) donde el earnings cae en los 10 días previos a una observación.")
    else:
        st.success("Ningún earnings próximo cae justo antes de una observación (ventana de 10 días).")

    def _colorear_earnings(row):
        if row.get("alerta"):
            return ["background-color: #fee2e2; color: #7f1d1d; font-weight: 700"] * len(row)
        return [""] * len(row)

    st.dataframe(df_earnings.style.apply(_colorear_earnings, axis=1), use_container_width=True, hide_index=True)
    boton_descarga_excel(df_earnings, "calendario_earnings.xlsx")


def _tab_pdfs_notas(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, df_calls: pd.DataFrame):
    st.caption(
        "Documento oficial de cada nota: se guarda en tres sitios al crearla — el volumen del "
        "servidor, Postgres y una carpeta de Google Drive — así el acceso nunca depende de un "
        "único punto de fallo. Puedes verlo o descargarlo aquí en cualquier momento. Para notas "
        "antiguas que no tienen PDF guardado todavía, súbelo desde aquí abajo: al guardarlo, la "
        "IA lo audita automáticamente contra lo que ya hay en el Excel."
    )

    notas_existentes = sorted(
        int(n) for n in pd.to_numeric(df_control.get("nota", pd.Series(dtype=float)), errors="coerce").dropna().unique()
    ) if df_control is not None and not df_control.empty else []

    if not notas_existentes:
        st.info("No hay notas en CONTROL_NOTAS todavía.")
        return

    if st.button("🔄 Sincronizar backup (Drive + Postgres) con los PDFs locales que falten"):
        from postgres_writer import guardar_pdf_nota_postgres
        subidos_drive, subidos_pg, sin_local = 0, 0, 0
        with st.spinner("Revisando y subiendo PDFs pendientes..."):
            for n in notas_existentes:
                ruta_local = os.path.join(CARPETA_PDFS_NOTAS, f"nota_{n:02d}.pdf")
                if not os.path.exists(ruta_local):
                    sin_local += 1
                    continue
                with open(ruta_local, "rb") as f:
                    contenido = f.read()
                subidos_drive += 1 if _subir_pdf_nota_a_drive(n, contenido) else 0
                subidos_pg += 1 if guardar_pdf_nota_postgres(n, contenido) else 0
        mensaje = f"Sincronización completada: {subidos_drive} PDF(s) al día en Drive, {subidos_pg} al día en Postgres."
        if sin_local:
            mensaje += f" {sin_local} nota(s) sin PDF en el servidor (no había nada que subir)."
        st.success(mensaje)

    st.markdown("---")
    st.markdown("#### 🔎 Recuperar CUSIP de notas ya guardadas")
    st.caption(
        "Vuelve a leer el PDF ya guardado de cada nota que todavía no tiene CUSIP en CONTROL_NOTAS "
        "y lo rellena. No toca ningún otro dato de la nota — solo añade el CUSIP."
    )

    control_cusip = df_control.copy() if df_control is not None else pd.DataFrame()
    if "cusip" not in control_cusip.columns:
        control_cusip["cusip"] = None
    notas_sin_cusip = [
        n for n in notas_existentes
        if control_cusip.loc[pd.to_numeric(control_cusip.get("nota"), errors="coerce") == n, "cusip"]
        .apply(lambda v: v is None or pd.isna(v) or str(v).strip() == "" or str(v).strip().upper() == "REVISAR")
        .all()
    ]

    if not notas_sin_cusip:
        st.success("Todas las notas con datos en CONTROL_NOTAS ya tienen CUSIP guardado.")
    else:
        st.caption(f"{len(notas_sin_cusip)} nota(s) sin CUSIP: {', '.join(str(n) for n in notas_sin_cusip)}")
        if st.button("🔍 Extraer CUSIP de esas notas"):
            resultados_backfill = {}
            with st.spinner("Leyendo PDFs y extrayendo CUSIP..."):
                for n in notas_sin_cusip:
                    pdf_n = leer_pdf_nota_guardado(n)
                    if not pdf_n:
                        resultados_backfill[n] = {"cusip": None, "estado": "⚠️ Sin PDF guardado"}
                        continue
                    extraido_n = extraer_datos_nota_con_ia(pdf_n)
                    if "error" in extraido_n:
                        resultados_backfill[n] = {"cusip": None, "estado": f"⚠️ Error: {extraido_n['error']}"}
                        continue
                    cusip_n = extraido_n.get("cusip")
                    if not cusip_n or str(cusip_n).strip().upper() == "REVISAR":
                        resultados_backfill[n] = {"cusip": None, "estado": "⚠️ IA no lo encontró en el PDF"}
                    else:
                        resultados_backfill[n] = {"cusip": str(cusip_n).strip(), "estado": "✅ Extraído"}
            st.session_state["backfill_cusip_resultados"] = resultados_backfill

    resultados_backfill = st.session_state.get("backfill_cusip_resultados")
    if resultados_backfill:
        df_preview_cusip = pd.DataFrame([
            {"Nota": n, "CUSIP extraído": r["cusip"] or "—", "Estado": r["estado"]}
            for n, r in resultados_backfill.items()
        ]).sort_values("Nota")
        st.dataframe(df_preview_cusip, use_container_width=True, hide_index=True)

        n_ok_cusip = sum(1 for r in resultados_backfill.values() if r["cusip"])
        if n_ok_cusip > 0 and st.button(f"💾 Guardar {n_ok_cusip} CUSIP en CONTROL_NOTAS", type="primary"):
            hojas = leer_todas_las_hojas_excel()
            if not hojas or "CONTROL_NOTAS" not in hojas:
                st.error("No se pudo leer el Excel actual para guardar los CUSIP.")
            else:
                df_ctrl = hojas["CONTROL_NOTAS"]
                if "CUSIP" not in df_ctrl.columns:
                    df_ctrl["CUSIP"] = None
                for n, r in resultados_backfill.items():
                    if r["cusip"]:
                        df_ctrl.loc[pd.to_numeric(df_ctrl["NOTA"], errors="coerce") == n, "CUSIP"] = r["cusip"]
                hojas["CONTROL_NOTAS"] = df_ctrl
                guardar_excel_completo_desde_hojas(hojas)
                st.success(f"CUSIP guardado para {n_ok_cusip} nota(s).")
                del st.session_state["backfill_cusip_resultados"]
                st.rerun()

    st.markdown("---")
    for n in notas_existentes:
        pdf_bytes = leer_pdf_nota_guardado(n)
        ruta_local = os.path.join(CARPETA_PDFS_NOTAS, f"nota_{n:02d}.pdf")
        en_local = os.path.exists(ruta_local)
        col_a, col_b, col_c = st.columns([1.5, 2, 2.5])
        col_a.markdown(f"**Nota {n}**")
        if en_local:
            col_b.caption("📁 En servidor")
        elif pdf_bytes:
            col_b.caption("☁️ Solo en Postgres/Drive")
        else:
            col_b.caption("⚠️ Sin PDF guardado")
        with col_c:
            if pdf_bytes:
                st.download_button(
                    "⬇️ Descargar PDF", data=pdf_bytes, file_name=f"nota_{n:02d}.pdf",
                    mime="application/pdf", key=f"descargar_pdf_nota_{n}",
                    use_container_width=True,
                )
        with st.expander(f"➕ Subir o reemplazar el PDF de la Nota {n}"):
            nuevo_pdf = st.file_uploader(
                f"Documento oficial de la Nota {n} (PDF)", type=["pdf"], key=f"subir_pdf_nota_{n}",
            )
            if nuevo_pdf is not None:
                if st.button(f"💾 Guardar y auditar este PDF para la Nota {n}", key=f"guardar_pdf_nota_{n}", type="primary"):
                    with st.spinner("Guardando en servidor, Postgres y Drive, y leyendo el PDF con la IA..."):
                        contenido = nuevo_pdf.read()
                        ruta_guardada = guardar_pdf_nota(n, contenido)
                        resultado_auditoria = _ejecutar_auditoria_ia(n, contenido) if ruta_guardada else None
                    if ruta_guardada:
                        st.success(f"PDF de la Nota {n} guardado en servidor, Postgres y Drive.")
                        if resultado_auditoria and "error" in resultado_auditoria:
                            st.warning(f"El PDF se guardó, pero la IA no pudo auditarlo todavía: {resultado_auditoria['error']}")
                    else:
                        st.error("No se pudo guardar en el volumen del servidor. Revisa que esté montado.")

            _render_comparacion_auditoria(n, df_inv, df_cal, df_control, df_calls)


def seccion_notas_archivo():
    df_inv, df_cal, df_control = cargar_excel_completo()
    df_calls = _leer_calendario_calls_cached()
    if not df_calls.empty and "nota" in df_calls.columns:
        df_calls["nota"] = pd.to_numeric(df_calls["nota"], errors="coerce")
    st.header("🧾 Notas")

    tab_resumen, tab_ficha, tab_analisis, tab_comparador, tab_noticias, tab_earnings, tab_pdfs, tab_auditar = st.tabs([
        "📊 Resumen y alertas",
        "🏢 Ficha de compañía", "🔬 Análisis completo de nota", "⚖️ Comparador de notas",
        "📰 Noticias", "📅 Calendario earnings", "📎 PDFs de notas", "🔍 Auditar nota",
    ])

    with tab_ficha:
        _tab_ficha_compania(df_control)

    with tab_analisis:
        _tab_analisis_nota_existente(df_inv, df_cal, df_control, df_calls)

    with tab_comparador:
        _tab_comparador_notas(df_control)

    with tab_noticias:
        _tab_dashboard_noticias(df_inv, df_control)

    with tab_earnings:
        _tab_calendario_earnings(df_inv, df_cal, df_control)

    with tab_pdfs:
        _tab_pdfs_notas(df_inv, df_cal, df_control, df_calls)

    with tab_auditar:
        _tab_auditar_nota(df_inv, df_cal, df_control, df_calls)

    with tab_resumen:
        st.caption("Resumen de precios actuales, variación, barrera de contingencia y alertas por nota.")

        if yf is None:
            st.error("Falta yfinance. Añade yfinance a requirements.txt.")
            return
        if df_control is None or df_control.empty:
            st.warning("La hoja CONTROL_NOTAS está vacía o no existe.")
            return

        faltan = [c for c in ["nota", "ticker", "precio_compra"] if c not in df_control.columns]
        barrera_col = next((c for c in ["contingency", "barrera_capital", "barrera_cupon"] if c in df_control.columns), None)
        if faltan:
            st.error(f"En CONTROL_NOTAS faltan columnas: {', '.join(faltan)}")
            return
        if barrera_col is None:
            st.error("En CONTROL_NOTAS falta una columna de barrera: CONTINGENCY, BARRERA_CAPITAL o BARRERA_CUPON.")
            return

        df_control_filtrado = obtener_control_notas_activas(df_inv, df_control)

        if st.button("Actualizar precios actuales"):
            st.cache_data.clear()
            st.rerun()

        with st.spinner("Descargando precios actuales..."):
            resumen = construir_resumen_actual_notas_alertas(df_control_filtrado)

        if resumen.empty:
            st.warning("No se pudo generar el resumen.")
            return

        alertas_resumen = resumen_alertas_por_nota(resumen)
        c1, c2, c3 = st.columns(3)
        c1.metric("Notas analizadas", resumen["nota"].nunique())
        c2.metric("Tickers", len(resumen))
        c3.metric("Notas en riesgo (variación ≤ -30%)", int((alertas_resumen["alerta"] == "ROJO").sum()) if not alertas_resumen.empty else 0)

        st.markdown("### 🚦 Semáforo consolidado por nota")
        st.caption(
            "Una fila por nota (no por ticker): se toma el peor ticker de cada nota (el que manda en un worst-of). "
            "🔴 = ya en riesgo (variación ≤ -30% vs precio de compra). 🟡 = todavía OK pero a menos de 10 puntos de entrar en riesgo. "
            "Sin color = margen cómodo. 'margen_a_barrera_%' es la variación entre el precio actual y el precio de contingencia (barrera)."
        )
        semaforo = construir_semaforo_consolidado_notas(resumen)
        if not semaforo.empty:
            st.dataframe(
                semaforo.style.apply(colorear_semaforo_consolidado, axis=1).format({
                    "peor_variacion_%": lambda x: f"{float(x):.2f}%" if pd.notna(x) else "Sin dato",
                    "precio_actual": lambda x: f"${float(x):,.2f}" if pd.notna(x) else "Sin dato",
                    "precio_contingencia": lambda x: f"${float(x):,.2f}" if pd.notna(x) else "Sin dato",
                    "margen_a_barrera_%": lambda x: f"{float(x):+.2f}%" if pd.notna(x) else "Sin dato",
                }),
                use_container_width=True, hide_index=True,
                column_order=["nota", "alerta", "peor_ticker", "peor_variacion_%", "precio_actual", "precio_contingencia", "margen_a_barrera_%", "n_tickers"],
            )
            boton_descarga_excel(semaforo, "semaforo_consolidado_notas.xlsx")
        else:
            st.info("No hay notas con datos suficientes para el semáforo.")

        st.markdown("---")
        st.markdown("### Detalle por ticker")
        tabla = resumen.copy()
        tabla["variacion_%"] = pd.to_numeric(tabla["variacion_%"], errors="coerce")
        tabla["margen_a_barrera_%"] = pd.to_numeric(tabla["margen_a_barrera_%"], errors="coerce")
        columnas_dinero = ["precio_compra", "precio_actual", "precio_contingencia"]
        tabla_mostrar = preparar_tabla_monetaria(tabla, columnas_dinero)
        if "variacion_%" in tabla_mostrar.columns:
            tabla_mostrar["variacion_%"] = tabla["variacion_%"].apply(lambda x: f"{float(x):.2f}%" if pd.notna(x) else "Sin dato")
        if "margen_a_barrera_%" in tabla_mostrar.columns:
            tabla_mostrar["margen_a_barrera_%"] = tabla["margen_a_barrera_%"].apply(lambda x: f"{float(x):.2f}%" if pd.notna(x) else "Sin dato")

        st.dataframe(tabla_mostrar.style.apply(colorear_filas_alerta_notas, axis=1), use_container_width=True)
        boton_descarga_excel(tabla_mostrar, "alertas_observaciones.xlsx")

        st.markdown("### Alertas por variación desde precio de compra")
        st.caption(
            "El criterio de riesgo es la variación % desde el precio de compra de cada ticker: "
            "🔴 En riesgo si es ≤ -30%. La columna 'margen a la barrera' se muestra como dato adicional "
            "de contexto (puede venir 'Sin dato' si la nota no tiene barrera cargada en CONTROL_NOTAS), "
            "pero no decide si una nota entra en esta lista."
        )
        if alertas_resumen.empty:
            st.success("No hay notas en riesgo (variación ≤ -30%) ahora mismo.")
        else:
            rojas = int((alertas_resumen["alerta"] == "ROJO").sum())
            st.error(f"Hay {rojas} nota(s) en riesgo (variación ≤ -30%).")
            alertas_mostrar = alertas_resumen.copy()
            alertas_mostrar["peor_margen_%"] = alertas_mostrar["peor_margen_%"].apply(lambda x: f"{float(x):.2f}%" if pd.notna(x) else "Sin dato")
            alertas_mostrar["peor_variacion_%"] = alertas_mostrar["peor_variacion_%"].apply(lambda x: f"{float(x):.2f}%" if pd.notna(x) else "Sin dato")
            st.dataframe(alertas_mostrar, use_container_width=True)



def seccion_alertas_notas():
    df_inv, df_cal, df_control = cargar_excel_completo()
    st.header("🚨 Alertas Notas")
    fecha = pd.Timestamp(st.date_input("Fecha de consulta", value=pd.Timestamp.today().date())).normalize()
    if df_cal.empty:
        st.warning("No existe la hoja CALENDARIO_NOTAS o está vacía.")
        return
    eventos = df_cal[df_cal["fecha"] == fecha].copy()
    st.subheader(f"Eventos del {fecha.strftime('%d/%m/%Y')}")
    st.dataframe(preparar_tabla_monetaria(eventos, []), use_container_width=True) if not eventos.empty else st.info("No hay observaciones ni pagos para esta fecha.")
    observaciones = eventos[eventos["tipo_evento"] == "OBSERVACION"].copy() if not eventos.empty else pd.DataFrame()
    pagos = eventos[eventos["tipo_evento"] == "PAGO"].copy() if not eventos.empty else pd.DataFrame()
    if not observaciones.empty:
        st.subheader("Evaluación de observaciones")
        for _, row in observaciones.iterrows():
            nota = int(row["nota"])
            resultado, detalle = evaluar_nota_en_fecha(df_control, nota, fecha, preferida="contingency")
            (st.success if resultado == "POSITIVA" else st.error if resultado == "NEGATIVA" else st.warning)(f"NOTA {nota}: {resultado}")
            if not detalle.empty:
                st.dataframe(preparar_tabla_monetaria(detalle, ["precio_compra", "precio_barrera", "cierre_usado"]), use_container_width=True)
    if not pagos.empty:
        st.subheader("Pagos del día")
        for _, row in pagos.iterrows():
            nota = int(row["nota"])
            previas = df_cal[(df_cal["nota"] == nota) & (df_cal["tipo_evento"] == "OBSERVACION") & (df_cal["fecha"] < fecha)].sort_values("fecha")
            if previas.empty:
                st.warning(f"NOTA {nota}: pago hoy, pero no he encontrado observación previa.")
                continue
            fecha_obs = previas.iloc[-1]["fecha"]
            resultado, detalle = evaluar_nota_en_fecha(df_control, nota, fecha_obs, preferida="contingency")
            (st.success if resultado == "POSITIVA" else st.error if resultado == "NEGATIVA" else st.warning)(f"NOTA {nota}: pago hoy. Observación previa {pd.Timestamp(fecha_obs).strftime('%d/%m/%Y')}: {resultado}")
            if not detalle.empty:
                with st.expander(f"Detalle NOTA {nota}"):
                    st.dataframe(preparar_tabla_monetaria(detalle, ["precio_compra", "precio_barrera", "cierre_usado"]), use_container_width=True)


def seccion_alertas_semana():
    _, df_cal, _ = cargar_excel_completo()
    st.header("📆 Alertas Semana")
    fecha_inicio = pd.Timestamp(st.date_input("Fecha de inicio", value=pd.Timestamp.today().date())).normalize()
    fecha_fin = fecha_inicio + pd.Timedelta(days=6)
    st.caption(f"Del {fecha_inicio.strftime('%d/%m/%Y')} al {fecha_fin.strftime('%d/%m/%Y')}")
    eventos = df_cal[(df_cal["fecha"].notna()) & (df_cal["fecha"] >= fecha_inicio) & (df_cal["fecha"] <= fecha_fin)].copy().sort_values(["fecha", "tipo_evento", "nota"])
    if eventos.empty:
        st.info("No hay observaciones ni pagos esta semana.")
        return
    c1, c2 = st.columns(2)
    c1.metric("Observaciones", len(eventos[eventos["tipo_evento"] == "OBSERVACION"]))
    c2.metric("Pagos", len(eventos[eventos["tipo_evento"] == "PAGO"]))
    st.dataframe(preparar_tabla_monetaria(eventos, []), use_container_width=True)


def eventos_calendario_mes(df_cal: pd.DataFrame, anio: int, mes: int) -> pd.DataFrame:
    inicio = pd.Timestamp(anio, mes, 1)
    fin = inicio + pd.offsets.MonthEnd(0)
    eventos = df_cal[(df_cal["fecha"].notna()) & (df_cal["fecha"] >= inicio) & (df_cal["fecha"] <= fin)].copy()
    if eventos.empty:
        return eventos
    eventos["semana_mes"] = ((eventos["fecha"].dt.day - 1) // 7) + 1
    return eventos.sort_values(["fecha", "nota", "tipo_evento"])


def seccion_calendario_notas():
    _, df_cal, _ = cargar_excel_completo()
    st.header("🗓️ Calendario Notas")
    consulta = st.selectbox("Consulta", ["Esta semana", "Mes completo", "Semana concreta de un mes", "Exportar calendario de un mes"])
    if consulta == "Esta semana":
        hoy = pd.Timestamp.today().normalize()
        inicio = hoy - pd.Timedelta(days=hoy.weekday())
        fin = inicio + pd.Timedelta(days=6)
        eventos = df_cal[(df_cal["fecha"].notna()) & (df_cal["fecha"] >= inicio) & (df_cal["fecha"] <= fin)].copy().sort_values(["fecha", "nota", "tipo_evento"])
        st.caption(f"Del {inicio.strftime('%d/%m/%Y')} al {fin.strftime('%d/%m/%Y')}")
        st.dataframe(preparar_tabla_monetaria(eventos, []), use_container_width=True) if not eventos.empty else st.info("No hay eventos esta semana.")
        return
    c1, c2 = st.columns(2)
    anio = int(c1.number_input("Año", 2020, 2100, pd.Timestamp.today().year, key=f"cal_{consulta}_anio"))
    mes = int(c2.number_input("Mes", 1, 12, pd.Timestamp.today().month, key=f"cal_{consulta}_mes"))
    eventos = eventos_calendario_mes(df_cal, anio, mes)
    if consulta == "Mes completo":
        st.subheader(f"Calendario de {nombre_mes_es(mes)} {anio}")
        st.dataframe(preparar_tabla_monetaria(eventos, []), use_container_width=True) if not eventos.empty else st.info("No hay eventos ese mes.")
    elif consulta == "Semana concreta de un mes":
        semana = int(st.number_input("Semana del mes", min_value=1, max_value=5, value=1))
        filtrado = eventos[eventos["semana_mes"] == semana].copy() if not eventos.empty else pd.DataFrame()
        st.dataframe(preparar_tabla_monetaria(filtrado, []), use_container_width=True) if not filtrado.empty else st.info("No hay eventos en esa semana.")
    else:
        if eventos.empty:
            st.info("No hay eventos para exportar en ese mes.")
        else:
            salida = BytesIO()
            exportar = eventos.copy(); exportar["fecha"] = exportar["fecha"].dt.strftime("%d/%m/%Y")
            with pd.ExcelWriter(salida, engine="openpyxl") as writer:
                exportar.to_excel(writer, index=False, sheet_name="CALENDARIO")
            st.download_button("Descargar Excel", data=salida.getvalue(), file_name=f"calendario_notas_{mes}_{anio}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")




def preparar_calendario_integrado_notas(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame, df_calls: pd.DataFrame | None = None, fecha_inicio=None, fecha_fin=None) -> pd.DataFrame:
    """Calendario único de notas: observaciones, pagos cobrables y calls.

    Reglas de pago:
    - Si la observación asociada es NEGATIVA real, el pago NO se muestra.
    - Si está PENDIENTE, POSITIVA, SIN DATO, SIN CONTROL o NO EVALUADA, se muestra como pago previsto/habilitado.
    """
    filas = []

    if fecha_inicio is not None:
        fecha_inicio = pd.Timestamp(fecha_inicio).normalize()
    if fecha_fin is not None:
        fecha_fin = pd.Timestamp(fecha_fin).normalize()

    def dentro_rango(fecha):
        if pd.isna(fecha):
            return False
        fecha = pd.Timestamp(fecha).normalize()
        if fecha_inicio is not None and fecha < fecha_inicio:
            return False
        if fecha_fin is not None and fecha > fecha_fin:
            return False
        return True

    # 1) Observaciones
    if df_cal is not None and not df_cal.empty:
        observaciones = df_cal[(df_cal["tipo_evento"] == "OBSERVACION") & (df_cal["fecha"].notna())].copy()
        for _, row in observaciones.iterrows():
            fecha = row.get("fecha")
            if not dentro_rango(fecha):
                continue
            nota = row.get("nota")
            if pd.isna(nota):
                continue
            nota_int = int(nota)
            resultado, detalle = evaluar_nota_en_fecha(df_control, nota_int, fecha, preferida="contingency") if df_control is not None else ("NO_EVALUADA", pd.DataFrame())
            filas.append({
                "fecha": pd.Timestamp(fecha).normalize(),
                "tipo_evento": "OBSERVACION",
                "nota": nota_int,
                "estado": resultado,
                "monto_cobro": 0.0,
                "detalle": resumen_detalle_observacion(detalle),
            })

        # 2) Pagos: solo se muestran si el cobro_compania total de la nota es > 0
        pagos = df_cal[(df_cal["tipo_evento"] == "PAGO") & (df_cal["fecha"].notna())].copy()
        for _, row in pagos.iterrows():
            fecha = row.get("fecha")
            if not dentro_rango(fecha):
                continue
            nota = row.get("nota")
            if pd.isna(nota):
                continue
            nota_int = int(nota)
            pago_df = pd.DataFrame([row])
            detalle_pago = preparar_detalle_notas(df_inv, pago_df, df_cal=df_cal, df_control=df_control)
            monto = float(detalle_pago["cobro_compania"].sum()) if not detalle_pago.empty and "cobro_compania" in detalle_pago.columns else 0.0
            fecha_obs = obtener_observacion_previa_nota(df_cal, nota_int, fecha)
            resultado_obs = "NO_EVALUADA"
            detalle_obs = pd.DataFrame()
            if fecha_obs is not None and df_control is not None and not df_control.empty:
                resultado_obs, detalle_obs = evaluar_nota_en_fecha(df_control, nota_int, fecha_obs, preferida="contingency")

            # Si la observación fue negativa real o el monto queda a 0, no ponemos el pago en calendario.
            if resultado_obs == "NEGATIVA" or monto <= 0:
                continue

            filas.append({
                "fecha": pd.Timestamp(fecha).normalize(),
                "tipo_evento": "PAGO",
                "nota": nota_int,
                "estado": resultado_obs,
                "monto_cobro": monto,
                "detalle": f"Cobro previsto/habilitado. Observación usada: {pd.Timestamp(fecha_obs).strftime('%d/%m/%Y') if fecha_obs is not None else 'sin observación'}",
            })

    # 3) Calls
    if df_calls is not None and not df_calls.empty:
        calls = df_calls.copy()
        # Normalización flexible por si la hoja tiene fecha_call o fecha.
        if "fecha_call" in calls.columns:
            calls["fecha_call"] = pd.to_datetime(calls["fecha_call"], errors="coerce", dayfirst=True).dt.normalize()
            col_fecha_call = "fecha_call"
        elif "fecha" in calls.columns:
            calls["fecha"] = pd.to_datetime(calls["fecha"], errors="coerce", dayfirst=True).dt.normalize()
            col_fecha_call = "fecha"
        else:
            col_fecha_call = None

        if col_fecha_call is not None:
            if "nota" in calls.columns:
                calls["nota"] = pd.to_numeric(calls["nota"], errors="coerce").astype("Int64")
            for _, row in calls.iterrows():
                fecha = row.get(col_fecha_call)
                if not dentro_rango(fecha):
                    continue
                nota = row.get("nota", pd.NA)
                filas.append({
                    "fecha": pd.Timestamp(fecha).normalize(),
                    "tipo_evento": "CALL",
                    "nota": int(nota) if pd.notna(nota) else "",
                    "estado": str(row.get("estado", "CALL POSIBLE")).upper() if "estado" in calls.columns else "CALL POSIBLE",
                    "monto_cobro": 0.0,
                    "detalle": "Fecha de posible call / cancelación anticipada",
                })

    calendario = pd.DataFrame(filas)
    if calendario.empty:
        return calendario
    orden_tipo = {"OBSERVACION": 1, "CALL": 2, "PAGO": 3}
    calendario["orden_tipo"] = calendario["tipo_evento"].map(orden_tipo).fillna(9)
    calendario = calendario.sort_values(["fecha", "orden_tipo", "nota"]).drop(columns=["orden_tipo"])
    return calendario


def preparar_tabla_calendario_integrado(calendario: pd.DataFrame) -> pd.DataFrame:
    if calendario is None or calendario.empty:
        return calendario
    out = calendario.copy()
    out["fecha"] = pd.to_datetime(out["fecha"], errors="coerce").dt.strftime("%d/%m/%Y")
    if "monto_cobro" in out.columns:
        out["monto_cobro"] = out["monto_cobro"].map(fmt)
    return out


def panel_alertas_y_calendario():
    df_inv, df_cal, df_control = cargar_excel_completo()
    df_calls = _leer_calendario_calls_cached()

    st.markdown("## Alertas y calendario")
    st.caption("Calendario único con observaciones, pagos cobrables y calls. Los pagos con observación negativa real no aparecen como cobro.")

    hoy = pd.Timestamp.today().normalize()

    c1, c2, c3 = st.columns(3)
    vista = c1.selectbox("Vista", ["Próximos 30 días", "Este mes", "Mes concreto", "Rango personalizado"])

    if vista == "Próximos 30 días":
        fecha_inicio = hoy
        fecha_fin = hoy + pd.Timedelta(days=30)
    elif vista == "Este mes":
        fecha_inicio = pd.Timestamp(hoy.year, hoy.month, 1)
        fecha_fin = fecha_inicio + pd.offsets.MonthEnd(0)
    elif vista == "Mes concreto":
        anio = int(c2.number_input("Año", 2020, 2100, hoy.year, key="cal_unico_anio"))
        mes = int(c3.number_input("Mes", 1, 12, hoy.month, key="cal_unico_mes"))
        fecha_inicio = pd.Timestamp(anio, mes, 1)
        fecha_fin = fecha_inicio + pd.offsets.MonthEnd(0)
    else:
        fecha_inicio = pd.Timestamp(c2.date_input("Desde", value=hoy.date(), key="cal_unico_desde")).normalize()
        fecha_fin = pd.Timestamp(c3.date_input("Hasta", value=(hoy + pd.Timedelta(days=30)).date(), key="cal_unico_hasta")).normalize()

    calendario = preparar_calendario_integrado_notas(
        df_inv=df_inv,
        df_cal=df_cal,
        df_control=df_control,
        df_calls=df_calls,
        fecha_inicio=fecha_inicio,
        fecha_fin=fecha_fin,
    )

    alertas = detectar_alertas_financieras(df_inv, df_cal, df_control)

    pagos = calendario[calendario["tipo_evento"] == "PAGO"] if not calendario.empty else pd.DataFrame()
    observaciones = calendario[calendario["tipo_evento"] == "OBSERVACION"] if not calendario.empty else pd.DataFrame()
    calls = calendario[calendario["tipo_evento"] == "CALL"] if not calendario.empty else pd.DataFrame()
    monto_total = float(pagos["monto_cobro"].sum()) if not pagos.empty and "monto_cobro" in pagos.columns else 0.0

    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Observaciones", len(observaciones))
    m2.metric("Pagos cobrables", len(pagos))
    m3.metric("Monto a cobrar", fmt(monto_total))
    m4.metric("Calls", len(calls))

    if not alertas.empty:
        criticas = alertas[alertas["Prioridad"] == "ALTA"] if "Prioridad" in alertas.columns else pd.DataFrame()
        if not criticas.empty:
            st.error(f"Hay {len(criticas)} alertas críticas que requieren revisión.")
        else:
            st.warning(f"Hay {len(alertas)} alertas de seguimiento.")
        with st.expander("Ver alertas del sistema", expanded=False):
            st.dataframe(alertas, use_container_width=True)
    else:
        st.success("No hay alertas activas.")

    st.markdown("### Calendario único")
    if calendario.empty:
        st.info("No hay observaciones, pagos cobrables ni calls en el periodo seleccionado.")
    else:
        filtro_tipo = st.multiselect(
            "Filtrar eventos",
            ["OBSERVACION", "PAGO", "CALL"],
            default=["OBSERVACION", "PAGO", "CALL"],
        )
        tabla = calendario[calendario["tipo_evento"].isin(filtro_tipo)].copy() if filtro_tipo else calendario.copy()
        st.dataframe(preparar_tabla_calendario_integrado(tabla), use_container_width=True)
        boton_descarga_excel(tabla, "calendario_notas.xlsx")

        salida = BytesIO()
        exportar = tabla.copy()
        exportar["fecha"] = pd.to_datetime(exportar["fecha"], errors="coerce").dt.strftime("%d/%m/%Y")
        with pd.ExcelWriter(salida, engine="openpyxl") as writer:
            exportar.to_excel(writer, index=False, sheet_name="CALENDARIO_UNICO")
        st.download_button(
            "Descargar calendario único en Excel",
            data=salida.getvalue(),
            file_name=f"calendario_unico_notas_{fecha_inicio.strftime('%Y%m%d')}_{fecha_fin.strftime('%Y%m%d')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    if px is not None and not calendario.empty:
        st.markdown("### Resumen visual")
        resumen_eventos = calendario.groupby("tipo_evento", as_index=False).size()
        fig = px.bar(resumen_eventos, x="tipo_evento", y="size", title="Eventos por tipo")
        fig.update_layout(height=330, paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", xaxis_title="Tipo", yaxis_title="Cantidad")
        st.plotly_chart(fig, use_container_width=True)

def panel_calidad_datos():
    df_inv, df_cal, df_control = cargar_excel_completo()
    st.markdown("## Calidad de datos")
    validaciones = validar_base_datos(df_inv, df_cal, df_control)
    total_inc = validaciones["Incidencias"].sum()
    criticas = validaciones[(validaciones["Incidencias"] > 0) & (validaciones["Estado"] == "ALTA")]
    c1, c2, c3 = st.columns(3)
    c1.metric("Incidencias totales", int(total_inc))
    c2.metric("Incidencias críticas", len(criticas))
    c3.metric("Validaciones OK", len(validaciones[validaciones["Incidencias"] == 0]))
    if len(criticas) > 0:
        st.error("Hay incidencias críticas que pueden afectar a los cálculos.")
    elif total_inc > 0:
        st.warning("Hay incidencias menores a revisar.")
    else:
        st.success("Base de datos validada correctamente.")
    st.dataframe(validaciones, use_container_width=True)
    boton_descarga_excel(validaciones, "validaciones_notas.xlsx")


def seccion_sistema_fondo():
    df_inv, df_cal, df_control = cargar_excel_completo()
    df_calls = _leer_calendario_calls_cached()
    if not df_calls.empty:
        if "fecha_call" in df_calls.columns:
            df_calls["fecha_call"] = pd.to_datetime(df_calls["fecha_call"], errors="coerce", dayfirst=True).dt.normalize()
        if "nota" in df_calls.columns:
            df_calls["nota"] = pd.to_numeric(df_calls["nota"], errors="coerce").astype("Int64")
    st.header("🏦 Sistema Fondo")
    consulta = st.selectbox("Consulta", ["Panel global", "Capital activo total", "Capital activo por activo", "Capital activo por inversor", "Capital activo de un inversor concreto", "Resumen mensual global", "Validaciones", "Calls de esta semana", "Calls de este mes", "Próximos calls", "Calls vencidos", "Capital desglosado por inversor"])
    if consulta == "Panel global":
        dashboard_financiero()
    elif consulta == "Capital activo total":
        activas = inversiones_activas_global(df_inv)
        mostrar_metricas("Resultado", [("Capital activo total", fmt(activas["capital_invertido"].sum() if not activas.empty else 0))])
    elif consulta == "Capital activo por activo":
        activas = inversiones_activas_global(df_inv)
        if activas.empty:
            st.info("No hay inversiones activas.")
        else:
            activas["activo"] = activas.apply(detectar_activo, axis=1)
            resumen = activas.groupby("activo", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False)
            st.dataframe(preparar_tabla_monetaria(resumen, ["capital"]), use_container_width=True)
    elif consulta == "Capital activo por inversor":
        activas = inversiones_activas_global(df_inv)
        resumen = activas.groupby("inversor", as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"}).sort_values("capital", ascending=False) if not activas.empty else pd.DataFrame()
        st.dataframe(preparar_tabla_monetaria(resumen, ["capital"]), use_container_width=True) if not resumen.empty else st.info("No hay inversiones activas.")
    elif consulta in ["Capital activo de un inversor concreto", "Capital desglosado por inversor"]:
        inversores = sorted([x for x in df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).unique() if x.strip()])
        nombre = st.selectbox("Inversor", inversores) if inversores else st.text_input("Inversor")
        if consulta == "Capital desglosado por inversor":
            c1, c2 = st.columns(2)
            anio = int(c1.number_input("Año", 2020, 2100, pd.Timestamp.today().year))
            mes = int(c2.number_input("Mes", 1, 12, pd.Timestamp.today().month))
            fecha = pd.Timestamp(anio, mes, ultimo_dia_mes(anio, mes))
        else:
            fecha = pd.Timestamp.today().normalize()
        activas = inversiones_activas_global(df_inv, fecha=fecha)
        filtrado = activas[activas["inversor"].astype(str).str.lower() == str(nombre).lower()].copy()
        mostrar_metricas("Resultado", [("Capital activo", fmt(filtrado["capital_invertido"].sum() if not filtrado.empty else 0))])
        if not filtrado.empty:
            filtrado["activo"] = filtrado.apply(detectar_activo, axis=1)
            resumen = filtrado.groupby(["activo", "nombre_activo"], as_index=False)["capital_invertido"].sum().rename(columns={"capital_invertido": "capital"})
            st.dataframe(preparar_tabla_monetaria(resumen, ["capital"]), use_container_width=True)
    elif consulta == "Resumen mensual global":
        c1, c2 = st.columns(2)
        anio = int(c1.number_input("Año", 2020, 2100, pd.Timestamp.today().year))
        mes = int(c2.number_input("Mes", 1, 12, pd.Timestamp.today().month))
        c_notas, p_notas, b_notas, d_notas, _ = resumen_notas_mes(df_inv, df_cal, df_control, anio, mes)
        detalles = []
        for activo, tasa in [("paraguay", TASA_ANUAL_PARAGUAY), ("bolivia", TASA_ANUAL_BOLIVIA), ("motoclick", TASA_ANUAL_MOTOCLICK), ("futbol", TASA_ANUAL_FUTBOL), ("bitcoin", TASA_ANUAL_BITCOIN)]:
            det = detalle_activo_mes(df_inv, activo, tasa, anio, mes)
            if not det.empty:
                det["activo"] = activo; detalles.append(det)
        d_fijos = pd.concat(detalles, ignore_index=True) if detalles else pd.DataFrame()
        d_fijos = ajustar_ingreso_motoclick(d_fijos, df_inv, anio, mes)
        c_fijos = d_fijos["ingreso_bruto"].sum() if not d_fijos.empty else 0
        p_fijos = d_fijos["pago_inversor_mes"].sum() if not d_fijos.empty else 0
        b_fijos = d_fijos["beneficio_empresa_mes"].sum() if not d_fijos.empty else 0
        mostrar_metricas(f"Resumen global {nombre_mes_es(mes)} {anio}", [("Cobro compañía", fmt(c_notas + c_fijos)), ("Pago inversores", fmt(p_notas + p_fijos)), ("Beneficio", fmt(b_notas + b_fijos))])
        if not d_notas.empty:
            with st.expander("Detalle notas"):
                st.dataframe(preparar_tabla_monetaria(d_notas, ["capital_invertido", "cobro_compania", "pago_inversor", "beneficio_empresa"]), use_container_width=True)
                boton_descarga_excel(d_notas, f"detalle_notas_{anio}_{mes:02d}.xlsx")
        if not d_fijos.empty:
            with st.expander("Detalle activos fijos"):
                st.dataframe(preparar_tabla_monetaria(d_fijos, ["capital_invertido", "ingreso_bruto", "pago_inversor_mes", "beneficio_empresa_mes"]), use_container_width=True)
                boton_descarga_excel(d_fijos, f"detalle_fijos_{anio}_{mes:02d}.xlsx")
    elif consulta == "Validaciones":
        df_val = validar_base_datos(df_inv, df_cal, df_control)
        st.dataframe(df_val, use_container_width=True)
        boton_descarga_excel(df_val, "validaciones.xlsx")
    else:
        if df_calls.empty or "fecha_call" not in df_calls.columns:
            st.warning("No existe la hoja CALENDARIO_CALLS o no tiene la columna fecha_call.")
            return
        hoy = pd.Timestamp.today().normalize()
        if consulta == "Calls de esta semana":
            inicio = hoy - pd.Timedelta(days=hoy.weekday()); fin = inicio + pd.Timedelta(days=6)
            res = df_calls[(df_calls["fecha_call"] >= inicio) & (df_calls["fecha_call"] <= fin)].copy()
        elif consulta == "Calls de este mes":
            res = df_calls[(df_calls["fecha_call"].dt.year == hoy.year) & (df_calls["fecha_call"].dt.month == hoy.month)].copy()
        elif consulta == "Próximos calls":
            res = df_calls[df_calls["fecha_call"] >= hoy].copy().sort_values("fecha_call").head(20)
        else:
            res = df_calls[df_calls["fecha_call"] < hoy].copy()
            if "estado" in res.columns:
                res = res[~res["estado"].apply(limpiar_texto).isin(["hecho", "realizado", "ejecutado", "call ejecutado"])]
        st.dataframe(preparar_tabla_monetaria(res, []), use_container_width=True) if not res.empty else st.info("No hay calls para esta consulta.")


# =========================
# EXTRACTOS
# =========================
def formatear_extracto_excel_bytes(contenido_raw: bytes, inversor: str, fecha_corte: datetime, detalle_df: "pd.DataFrame | None" = None) -> bytes:
    """
    Genera el extracto Excel profesional con:
    - Hoja PORTADA con resumen ejecutivo
    - Hoja DETALLE con filas de operación, cierres mensuales, anuales y cierre final
    - (La antigua hoja RESUMEN_MENSUAL se eliminó por duplicar la tabla de la PORTADA)
    - Todo en formato dólar, diseño premium
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
    from openpyxl.utils import get_column_letter
    from openpyxl.drawing.image import Image as XLImage
    import calendar as cal_mod

    # ── Paleta de colores ──────────────────────────────────────────────
    C_AZUL_OSC   = "0D2137"   # cabeceras principales
    C_AZUL_MED   = "1A3F5C"   # cabeceras secundarias
    C_AZUL_CLARO = "D6E9F8"   # fondo filas par detalle
    C_VERDE_OSC  = "1E4620"   # cierre anual texto
    C_VERDE      = "D9EAD3"   # cierre anual fondo
    C_NARANJA_OSC= "7F3F00"   # cierre mensual texto
    C_NARANJA    = "FCE5CD"   # cierre mensual fondo
    C_DORADO_OSC = "4A3000"   # cierre final texto
    C_DORADO     = "FFF2CC"   # cierre final fondo
    C_BLANCO     = "FFFFFF"
    C_GRIS_CLARO = "F7F9FC"
    C_GRIS_MED   = "D9D9D9"

    fmt_usd = '"$"#,##0.00'

    borde_fino  = Side(style="thin",   color=C_GRIS_MED)
    borde_medio = Side(style="medium", color=C_AZUL_MED)
    borde_std   = Border(left=borde_fino, right=borde_fino, top=borde_fino, bottom=borde_fino)
    borde_top   = Border(left=borde_fino, right=borde_fino, top=borde_medio, bottom=borde_fino)

    def fill(hex_color):
        return PatternFill("solid", fgColor=hex_color)

    def font(bold=False, size=11, color="000000", italic=False):
        return Font(name="Calibri", bold=bold, size=size, color=color, italic=italic)

    def aln(h="left", v="center", wrap=False):
        return Alignment(horizontal=h, vertical=v, wrap_text=wrap)

    def set_cell(ws, row, col, value, bold=False, size=11, fcolor=None, tcolor="000000",
                 align_h="left", fmt=None, italic=False, border=None):
        c = ws.cell(row=row, column=col, value=value)
        c.font = font(bold=bold, size=size, color=tcolor, italic=italic)
        c.alignment = aln(h=align_h)
        if fcolor:
            c.fill = fill(fcolor)
        if fmt:
            c.number_format = fmt
        c.border = border if border else borde_std
        return c

    # ── Leer datos originales ─────────────────────────────────────────
    bio = BytesIO(contenido_raw)
    from openpyxl import load_workbook as lw
    wb_orig = lw(bio)

    capital_total = 0.0
    total_intereses = 0.0

    # Leer DETALLE
    ws_det_orig = wb_orig["DETALLE"] if "DETALLE" in wb_orig.sheetnames else None
    det_rows = []
    det_cols = []
    if ws_det_orig:
        headers = [c.value for c in next(ws_det_orig.iter_rows(min_row=1, max_row=1))]
        det_cols = headers
        for row in ws_det_orig.iter_rows(min_row=2, values_only=True):
            det_rows.append(dict(zip(headers, row)))

    # Leer TOTALES_MES
    ws_tot_orig = wb_orig["TOTALES_MES"] if "TOTALES_MES" in wb_orig.sheetnames else None
    tot_rows = []
    if ws_tot_orig:
        for row in ws_tot_orig.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                tot_rows.append({"mes": row[0], "total_mes": row[1] if len(row) > 1 else 0})

    # Ordenar totales_mes cronológicamente
    def mes_sort_key(m):
        try:
            partes = str(m["mes"]).split("/")
            return (int(partes[1]), int(partes[0]))
        except Exception:
            return (9999, 99)
    tot_rows.sort(key=mes_sort_key)

    # Calcular capital activo = operaciones vivas el último día del mes de corte (fecha_corte)
    # Solo NUEVA sin fecha_final o con fecha_final >= fecha_corte
    # Es el mismo criterio que el cierre del último mes
    import calendar as _cal2
    ultimo_dia_corte = _cal2.monthrange(fecha_corte.year, fecha_corte.month)[1]
    fin_corte = datetime(fecha_corte.year, fecha_corte.month, ultimo_dia_corte)
    capital_total = 0.0
    total_intereses = 0.0
    for row in det_rows:
        ffop_str = str(row.get("fecha_fin_op", "") or "")
        cap = float(row.get("capital_invertido", 0) or 0)
        mes_str = str(row.get("mes", "") or "")
        # Solo contar filas del mes de corte (mayo 2026)
        mes_corte_str = f"{fecha_corte.month:02d}/{fecha_corte.year}"
        if mes_str != mes_corte_str:
            continue
        if not ffop_str or ffop_str in ("", "None", "nan"):
            capital_total += cap
        else:
            try:
                ffo = datetime.strptime(ffop_str, "%d/%m/%Y")
                if ffo >= fin_corte:
                    capital_total += cap
            except Exception:
                capital_total += cap
    total_intereses = sum(float(r.get("interes_mes", 0) or 0) for r in det_rows)

    # ── Detectar si este inversor tiene alguna fila con PAGA ──────────
    tiene_pago = any(
        str(r.get("pago_intereses", "REINVIERTE")).strip().upper() == "PAGA"
        for r in det_rows
    )

    # ── Crear workbook nuevo ──────────────────────────────────────────
    wb = Workbook()
    wb.remove(wb.active)

    # ═══════════════════════════════════════════════════════════════
    # HOJA 1: PORTADA
    # ═══════════════════════════════════════════════════════════════
    ws_p = wb.create_sheet("PORTADA")
    ws_p.sheet_view.showGridLines = False
    ws_p.sheet_view.showRowColHeaders = False

    for col_i, w in enumerate([4, 30, 30, 30, 4], 1):
        ws_p.column_dimensions[get_column_letter(col_i)].width = w
    for r_i in range(1, 60):
        ws_p.row_dimensions[r_i].height = 18

    # Banner superior
    ws_p.merge_cells("A1:E3")
    banner = ws_p["A1"]
    banner.value = "CF"
    banner.font = Font(name="Calibri", size=36, bold=True, color=C_BLANCO)
    banner.fill = fill(C_AZUL_OSC)
    banner.alignment = aln(h="center", v="center")
    ws_p.row_dimensions[1].height = 28
    ws_p.row_dimensions[2].height = 28
    ws_p.row_dimensions[3].height = 28

    ws_p.merge_cells("A4:E4")
    sub_b = ws_p["A4"]
    sub_b.value = "Chaparro Fernández Wealth  ·  Extracto de inversiones"
    sub_b.font = Font(name="Calibri", size=12, italic=True, color=C_BLANCO)
    sub_b.fill = fill(C_AZUL_MED)
    sub_b.alignment = aln(h="center", v="center")
    ws_p.row_dimensions[4].height = 22

    # Franja dorada
    ws_p.merge_cells("A5:E5")
    ws_p["A5"].fill = fill("BF9A5F")
    ws_p.row_dimensions[5].height = 5

    # Espacio
    ws_p.row_dimensions[6].height = 12

    # Nombre inversor
    ws_p.merge_cells("B7:D7")
    c_lbl_i = ws_p["B7"]
    c_lbl_i.value = "INVERSOR"
    c_lbl_i.font = Font(name="Calibri", size=9, bold=True, color="888888")
    c_lbl_i.alignment = aln(h="center")

    ws_p.merge_cells("B8:D9")
    c_inv = ws_p["B8"]
    c_inv.value = inversor.upper()
    c_inv.font = Font(name="Calibri", size=22, bold=True, color=C_AZUL_OSC)
    c_inv.alignment = aln(h="center", v="center")
    ws_p.row_dimensions[8].height = 30
    ws_p.row_dimensions[9].height = 30

    ws_p.merge_cells("B10:D10")
    c_fec = ws_p["B10"]
    c_fec.value = f"Fecha de corte: {fecha_corte.strftime('%d/%m/%Y')}"
    c_fec.font = Font(name="Calibri", size=10, italic=True, color="555555")
    c_fec.alignment = aln(h="center")

    ws_p.row_dimensions[11].height = 12

    # Separador
    ws_p.merge_cells("B12:D12")
    ws_p["B12"].fill = fill(C_GRIS_MED)
    ws_p.row_dimensions[12].height = 2

    ws_p.row_dimensions[13].height = 12

    # ── Calcular totales pagado/saldo para lógica de portada y cierres ──
    total_pagado_global = sum(
        float(r.get("interes_mes", 0) or 0)
        for r in det_rows if str(r.get("pago_intereses", "")).strip().upper() == "PAGA"
    )
    saldo_pendiente_global = total_intereses - total_pagado_global

    # Tres KPIs — distintos según si el inversor paga o reinvierte
    if tiene_pago:
        kpi_data = [
            ("CAPITAL ACTIVO",            capital_total,                         C_AZUL_OSC, C_AZUL_CLARO),
            ("INTERESES PENDIENTES",       saldo_pendiente_global,               "1E4620",   C_VERDE),
            ("TOTAL (Cap. + Pendiente)",   capital_total + saldo_pendiente_global, "4A3000", C_DORADO),
        ]
    else:
        kpi_data = [
            ("CAPITAL ACTIVO",        capital_total,                    C_AZUL_OSC, C_AZUL_CLARO),
            ("INTERESES GENERADOS",   total_intereses,                  "1E4620",   C_VERDE),
            ("TOTAL ACUMULADO",       capital_total + total_intereses,  "4A3000",   C_DORADO),
        ]
    for ki, (lbl_k, val_k, txt_k, bg_k) in enumerate(kpi_data):
        col_k = ki + 2  # columnas B, C, D
        c_lk = ws_p.cell(row=14, column=col_k, value=lbl_k)
        c_lk.font = Font(name="Calibri", size=8, bold=True, color=txt_k)
        c_lk.fill = fill(bg_k)
        c_lk.alignment = aln(h="center", v="center")
        c_lk.border = borde_std
        ws_p.row_dimensions[14].height = 20
        c_vk = ws_p.cell(row=15, column=col_k, value=val_k)
        c_vk.font = Font(name="Calibri", size=16, bold=True, color=txt_k)
        c_vk.fill = fill(bg_k)
        c_vk.number_format = fmt_usd
        c_vk.alignment = aln(h="center", v="center")
        c_vk.border = borde_std
        ws_p.row_dimensions[15].height = 34

    ws_p.row_dimensions[16].height = 12

    # Nota descriptiva
    ws_p.merge_cells("B17:D17")
    c_nota = ws_p["B17"]
    c_nota.value = ("Este extracto incluye el detalle de todas sus posiciones activas, "
                    "los intereses devengados mes a mes y el acumulado histórico desde el inicio de su inversión.")
    c_nota.font = Font(name="Calibri", size=9, italic=True, color="555555")
    c_nota.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws_p.row_dimensions[17].height = 40

    ws_p.row_dimensions[18].height = 12

    # Mini tabla resumen mensual en portada
    if tot_rows:
        # Precomputar pagado por mes para inversores con PAGA
        pagado_por_mes: dict = {}
        if tiene_pago:
            for r in det_rows:
                if str(r.get("pago_intereses", "")).strip().upper() == "PAGA":
                    mk_str = str(r.get("mes", ""))
                    pagado_por_mes[mk_str] = pagado_por_mes.get(mk_str, 0.0) + float(r.get("interes_mes", 0) or 0)

        if tiene_pago:
            hdrs_portada = ["MES", "GENERADO ($)", "PAGADO ($)", "SALDO ($)"]
            n_cols_portada = 4
            ws_p.merge_cells("B19:E19")
        else:
            hdrs_portada = ["MES", "INTERESES ($)", "ACUMULADO ($)"]
            n_cols_portada = 3
            ws_p.merge_cells("B19:D19")

        c_hres = ws_p["B19"]
        c_hres.value = "RESUMEN MENSUAL DE INTERESES"
        c_hres.font = Font(name="Calibri", size=10, bold=True, color=C_BLANCO)
        c_hres.fill = fill(C_AZUL_MED)
        c_hres.alignment = aln(h="center")
        ws_p.row_dimensions[19].height = 22

        for ci_h, hdr_h in enumerate(hdrs_portada, 2):
            c_hh = ws_p.cell(row=20, column=ci_h, value=hdr_h)
            c_hh.font = Font(name="Calibri", size=9, bold=True, color=C_AZUL_OSC)
            c_hh.fill = fill(C_AZUL_CLARO)
            c_hh.alignment = aln(h="center")
            c_hh.border = borde_std
        ws_p.row_dimensions[20].height = 20

        acum_p = 0.0
        acum_pagado_p = 0.0
        for ri_p, tr_p in enumerate(tot_rows, 21):
            fondo_p = C_GRIS_CLARO if ri_p % 2 == 0 else C_BLANCO
            val_p = float(tr_p["total_mes"] or 0)
            mes_str_p = str(tr_p["mes"])
            pagado_mes_p = pagado_por_mes.get(mes_str_p, 0.0) if tiene_pago else 0.0
            saldo_mes_p = val_p - pagado_mes_p
            acum_p += val_p
            acum_pagado_p += pagado_mes_p
            ws_p.row_dimensions[ri_p].height = 18

            c_m = ws_p.cell(row=ri_p, column=2, value=mes_str_p)
            c_m.font = font(size=9); c_m.fill = fill(fondo_p)
            c_m.alignment = aln(h="center"); c_m.border = borde_std

            if tiene_pago:
                c_gen = ws_p.cell(row=ri_p, column=3, value=val_p)
                c_gen.font = font(size=9); c_gen.fill = fill(fondo_p)
                c_gen.number_format = fmt_usd; c_gen.alignment = aln(h="right"); c_gen.border = borde_std
                c_pag = ws_p.cell(row=ri_p, column=4, value=-pagado_mes_p)
                c_pag.font = Font(name="Calibri", size=9, color="C00000")
                c_pag.fill = fill(fondo_p)
                c_pag.number_format = fmt_usd; c_pag.alignment = aln(h="right"); c_pag.border = borde_std
                c_sal = ws_p.cell(row=ri_p, column=5, value=saldo_mes_p)
                c_sal.font = Font(name="Calibri", size=9, bold=True, color="1E4620")
                c_sal.fill = fill(C_VERDE if saldo_mes_p == 0 else fondo_p)
                c_sal.number_format = fmt_usd; c_sal.alignment = aln(h="right"); c_sal.border = borde_std
            else:
                c_v2 = ws_p.cell(row=ri_p, column=3, value=val_p)
                c_v2.font = font(size=9); c_v2.fill = fill(fondo_p)
                c_v2.number_format = fmt_usd; c_v2.alignment = aln(h="right"); c_v2.border = borde_std
                c_a2 = ws_p.cell(row=ri_p, column=4, value=acum_p)
                c_a2.font = font(size=9); c_a2.fill = fill(fondo_p)
                c_a2.number_format = fmt_usd; c_a2.alignment = aln(h="right"); c_a2.border = borde_std

        fila_tot_p = 21 + len(tot_rows)
        ws_p.row_dimensions[fila_tot_p].height = 22
        acum_saldo_p = acum_p - acum_pagado_p
        if tiene_pago:
            tot_vals = [("TOTAL", None), (acum_p, fmt_usd), (-acum_pagado_p, fmt_usd), (acum_saldo_p, fmt_usd)]
            for ci_t, (val_t, fmt_t) in enumerate(tot_vals, 2):
                c_t = ws_p.cell(row=fila_tot_p, column=ci_t, value=val_t)
                c_t.font = Font(name="Calibri", size=10, bold=True, color=C_DORADO_OSC)
                c_t.fill = fill(C_DORADO)
                c_t.alignment = aln(h="center" if ci_t == 2 else "right")
                if fmt_t:
                    c_t.number_format = fmt_t
                c_t.border = borde_top
        else:
            for ci_t, val_t in [(2, "TOTAL"), (3, acum_p), (4, acum_p)]:
                c_t = ws_p.cell(row=fila_tot_p, column=ci_t,
                                value=val_t if ci_t > 2 else "TOTAL")
                c_t.font = Font(name="Calibri", size=10, bold=True, color=C_DORADO_OSC)
                c_t.fill = fill(C_DORADO)
                c_t.alignment = aln(h="center" if ci_t == 2 else "right")
                if ci_t > 2:
                    c_t.number_format = fmt_usd
                c_t.border = borde_top

        fila_pie_p = fila_tot_p + 2
    else:
        fila_pie_p = 20

    # Pie de portada
    ws_p.merge_cells(f"A{fila_pie_p}:E{fila_pie_p}")
    c_pie = ws_p[f"A{fila_pie_p}"]
    c_pie.value = "Documento confidencial · Chaparro Fernández Wealth · Generado automáticamente"
    c_pie.font = Font(name="Calibri", size=8, italic=True, color="AAAAAA")
    c_pie.fill = fill(C_AZUL_OSC)
    c_pie.alignment = aln(h="center", v="center")
    ws_p.row_dimensions[fila_pie_p].height = 18

    # ═══════════════════════════════════════════════════════════════
    # HOJA 2: DETALLE con cierres mensuales, anuales y final
    # ═══════════════════════════════════════════════════════════════
    ws_d = wb.create_sheet("DETALLE")
    ws_d.sheet_view.showGridLines = False
    ws_d.freeze_panes = "A5"

    # Columnas visibles: solo las que necesita ver el inversor (sin ocultar nada)
    # COL_CAP = columna del capital, COL_INT = columna del interés
    if tiene_pago:
        col_names  = ["Activo", "Mes", "Fecha inversión", "Capital ($)", "Interés mes ($)", "Estado intereses"]
        col_widths = [22, 14, 16, 16, 16, 28]
        col_map = {
            "nombre_activo": 1, "mes": 2, "fecha_inversion": 3,
            "capital_invertido": 4, "interes_mes": 5, "pago_intereses": 6,
        }
    else:
        col_names  = ["Activo", "Mes", "Fecha inversión", "Capital ($)", "Interés mes ($)"]
        col_widths = [22, 14, 16, 16, 16]
        col_map = {
            "nombre_activo": 1, "mes": 2, "fecha_inversion": 3,
            "capital_invertido": 4, "interes_mes": 5,
        }

    COL_CAP = 4   # columna Capital ($)
    COL_INT = 5   # columna Interés mes ($)
    COL_PAGO = 6  # columna Estado intereses (solo si tiene_pago)

    num_cols = len(col_names)
    for i, w in enumerate(col_widths[:num_cols], 1):
        ws_d.column_dimensions[get_column_letter(i)].width = w

    # Título
    ws_d.merge_cells(f"A1:{get_column_letter(num_cols)}1")
    t = ws_d["A1"]
    t.value = f"DETALLE DEL EXTRACTO — {inversor.upper()}"
    t.font = Font(name="Calibri", size=16, bold=True, color=C_BLANCO)
    t.fill = fill(C_AZUL_OSC)
    t.alignment = aln(h="center", v="center")
    ws_d.row_dimensions[1].height = 42

    ws_d.merge_cells(f"A2:{get_column_letter(num_cols)}2")
    sub = ws_d["A2"]
    sub.value = f"Fecha de corte: {fecha_corte.strftime('%d/%m/%Y')}   |   Inversor: {inversor}"
    sub.font = Font(name="Calibri", size=10, italic=True, color="444444")
    sub.alignment = aln(h="center")
    ws_d.row_dimensions[2].height = 24

    ws_d.row_dimensions[3].height = 8
    for c in range(1, num_cols+1):
        ws_d.cell(row=3, column=c).fill = fill("2E86C1")

    # Cabecera
    for ci, nombre in enumerate(col_names, 1):
        c = ws_d.cell(row=4, column=ci, value=nombre)
        c.font = Font(name="Calibri", size=10, bold=True, color=C_BLANCO)
        c.fill = fill(C_AZUL_MED)
        c.alignment = aln(h="center")
        c.border = borde_std
    ws_d.row_dimensions[4].height = 28

    def mes_key_from_str(mes_str):
        try:
            p = str(mes_str).split("/")
            return (int(p[1]), int(p[0]))
        except Exception:
            return (9999, 99)

    # Agrupar filas por mes
    from collections import defaultdict
    meses_orden = []
    filas_por_mes = defaultdict(list)
    for row in det_rows:
        mk = mes_key_from_str(row.get("mes", ""))
        if mk not in filas_por_mes:
            meses_orden.append(mk)
        filas_por_mes[mk].append(row)
    meses_orden = sorted(set(meses_orden))

    fila_excel = 5
    intereses_acum_anio = 0.0
    capital_anio = 0.0
    anio_actual = None
    intereses_acum_total = 0.0

    for mk in meses_orden:
        anio_mk, mes_mk = mk
        rows_mes = filas_por_mes[mk]

        # Cambio de año: insertar cierre anual del año anterior
        if anio_actual is not None and anio_mk != anio_actual:
            # CIERRE ANUAL
            ws_d.row_dimensions[fila_excel].height = 24
            ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_CAP - 1)
            c = ws_d.cell(row=fila_excel, column=1, value=f"CIERRE {anio_actual}")
            c.font = Font(name="Calibri", size=11, bold=True, color=C_VERDE_OSC)
            c.fill = fill(C_VERDE)
            c.alignment = aln(h="center")
            c.border = borde_top

            cap_c = ws_d.cell(row=fila_excel, column=COL_CAP, value=capital_anio)
            cap_c.font = Font(name="Calibri", size=11, bold=True, color=C_VERDE_OSC)
            cap_c.fill = fill(C_VERDE)
            cap_c.number_format = fmt_usd
            cap_c.alignment = aln(h="right")
            cap_c.border = borde_top

            int_c = ws_d.cell(row=fila_excel, column=COL_INT, value=intereses_acum_anio)
            int_c.font = Font(name="Calibri", size=11, bold=True, color=C_VERDE_OSC)
            int_c.fill = fill(C_VERDE)
            int_c.number_format = fmt_usd
            int_c.alignment = aln(h="right")
            int_c.border = borde_top

            if tiene_pago:
                cx11 = ws_d.cell(row=fila_excel, column=COL_PAGO, value="")
                cx11.fill = fill(C_VERDE); cx11.border = borde_top

            fila_excel += 1
            # Fila acumulado anual
            ws_d.row_dimensions[fila_excel].height = 20
            ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_INT - 1)
            ca = ws_d.cell(row=fila_excel, column=1, value=f"   Capital + Intereses acumulados {anio_actual}")
            ca.font = Font(name="Calibri", size=10, italic=True, color=C_VERDE_OSC)
            ca.fill = fill(C_VERDE)
            ca.alignment = aln(h="right")
            ca.border = borde_std
            total_anio_c = ws_d.cell(row=fila_excel, column=COL_INT, value=capital_anio + intereses_acum_anio)
            total_anio_c.font = Font(name="Calibri", size=10, bold=True, italic=True, color=C_VERDE_OSC)
            total_anio_c.fill = fill(C_VERDE)
            total_anio_c.number_format = fmt_usd
            total_anio_c.alignment = aln(h="right")
            total_anio_c.border = borde_std
            if tiene_pago:
                cx11 = ws_d.cell(row=fila_excel, column=COL_PAGO, value="")
                cx11.fill = fill(C_VERDE); cx11.border = borde_std
            fila_excel += 1
            intereses_acum_anio = 0.0

        anio_actual = anio_mk

        # Filas de detalle del mes
        intereses_mes = 0.0
        capital_mes = 0.0
        par = (fila_excel % 2 == 0)
        for ri, row in enumerate(rows_mes):
            fondo = C_AZUL_CLARO if (fila_excel % 2 == 0) else C_GRIS_CLARO
            ws_d.row_dimensions[fila_excel].height = 26
            pago_int_fila = str(row.get("pago_intereses", "REINVIERTE")).strip().upper()
            for col_key, col_idx in col_map.items():
                val = row.get(col_key, "")
                # Para la columna pago_intereses, transformar el valor a etiqueta visual
                if col_key == "pago_intereses":
                    val = "✅ PAGADO" if pago_int_fila == "PAGA" else "🔄 REINVERTIDO"
                c = ws_d.cell(row=fila_excel, column=col_idx, value=val)
                # Color especial para la celda de estado de intereses
                if col_key == "pago_intereses":
                    if pago_int_fila == "PAGA":
                        c.font = Font(name="Calibri", size=10, bold=True, color="1E4620")
                        c.fill = fill("C6EFCE")
                    else:
                        c.font = Font(name="Calibri", size=10, color="1A3F5C")
                        c.fill = fill(fondo)
                else:
                    c.font = Font(name="Calibri", size=10, color="222222")
                    c.fill = fill(fondo)
                c.border = borde_std
                c.alignment = aln(h="right" if col_idx in [COL_CAP, COL_INT] else "center" if (tiene_pago and col_idx == COL_PAGO) else "left")
                if col_idx == COL_CAP:
                    c.number_format = fmt_usd
                if col_idx == COL_INT:
                    c.number_format = fmt_usd
            v_cap = row.get("capital_invertido", 0) or 0
            v_int = row.get("interes_mes", 0) or 0
            try:
                intereses_mes += float(v_int)
                if float(v_cap) > capital_mes:
                    capital_mes = float(v_cap)
            except Exception:
                pass
            fila_excel += 1

        # Capital activo al cierre del mes = solo operaciones vivas el último día del mes
        # Una operación está viva si su fecha_fin_op >= último día del mes
        import calendar as _cal
        ultimo_dia = _cal.monthrange(anio_mk, mes_mk)[1]
        fin_mes_dt = datetime(anio_mk, mes_mk, ultimo_dia)
        capital_mes_real = 0.0
        for r in rows_mes:
            fecha_fin_op_str = str(r.get("fecha_fin_op", "") or "")
            capital_row = float(r.get("capital_invertido", 0) or 0)
            if not fecha_fin_op_str or fecha_fin_op_str in ("", "None", "nan"):
                # Sin fecha fin → siempre activa
                capital_mes_real += capital_row
            else:
                try:
                    ffo = datetime.strptime(fecha_fin_op_str, "%d/%m/%Y")
                    if ffo >= fin_mes_dt:
                        capital_mes_real += capital_row
                    # Si ffo < fin_mes_dt → cancelada dentro del mes → no suma al capital
                except Exception:
                    capital_mes_real += capital_row
        intereses_mes_real = sum(float(r.get("interes_mes", 0) or 0) for r in rows_mes)

        # CIERRE MENSUAL
        ws_d.row_dimensions[fila_excel].height = 28
        mes_label = f"{mes_mk:02d}/{anio_mk}"
        ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_CAP - 1)
        cm = ws_d.cell(row=fila_excel, column=1, value=f"CIERRE {mes_label}")
        cm.font = Font(name="Calibri", size=10, bold=True, color=C_NARANJA_OSC)
        cm.fill = fill(C_NARANJA)
        cm.alignment = aln(h="center")
        cm.border = borde_top

        cap_cm = ws_d.cell(row=fila_excel, column=COL_CAP, value=capital_mes_real)
        cap_cm.font = Font(name="Calibri", size=10, bold=True, color=C_NARANJA_OSC)
        cap_cm.fill = fill(C_NARANJA)
        cap_cm.number_format = fmt_usd
        cap_cm.alignment = aln(h="right")
        cap_cm.border = borde_top

        int_cm = ws_d.cell(row=fila_excel, column=COL_INT, value=intereses_mes_real)
        int_cm.font = Font(name="Calibri", size=10, bold=True, color=C_NARANJA_OSC)
        int_cm.fill = fill(C_NARANJA)
        int_cm.number_format = fmt_usd
        int_cm.alignment = aln(h="right")
        int_cm.border = borde_top

        # Rellenar columnas extra si tiene_pago
        if tiene_pago:
            intereses_paga_mes = sum(
                float(r.get("interes_mes", 0) or 0)
                for r in rows_mes if str(r.get("pago_intereses", "")).strip().upper() == "PAGA"
            )
            intereses_rein_mes = intereses_mes_real - intereses_paga_mes
            lbl_pago_cm = ws_d.cell(row=fila_excel, column=COL_PAGO,
                value=f"✅ Pagado: ${intereses_paga_mes:,.2f}  |  🔄 Reinvertido: ${intereses_rein_mes:,.2f}")
            lbl_pago_cm.font = Font(name="Calibri", size=9, bold=True, color=C_NARANJA_OSC)
            lbl_pago_cm.fill = fill(C_NARANJA)
            lbl_pago_cm.alignment = aln(h="center")
            lbl_pago_cm.border = borde_top

        fila_excel += 1
        intereses_acum_anio += intereses_mes_real
        intereses_acum_total += intereses_mes_real
        capital_anio = capital_mes_real

    # CIERRE ANUAL del último año
    if anio_actual is not None:
        ws_d.row_dimensions[fila_excel].height = 30
        ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_CAP - 1)
        c = ws_d.cell(row=fila_excel, column=1, value=f"CIERRE {anio_actual}")
        c.font = Font(name="Calibri", size=11, bold=True, color=C_VERDE_OSC)
        c.fill = fill(C_VERDE)
        c.alignment = aln(h="center")
        c.border = borde_top
        cap_c = ws_d.cell(row=fila_excel, column=COL_CAP, value=capital_anio)
        cap_c.font = Font(name="Calibri", size=11, bold=True, color=C_VERDE_OSC)
        cap_c.fill = fill(C_VERDE)
        cap_c.number_format = fmt_usd
        cap_c.alignment = aln(h="right")
        cap_c.border = borde_top
        int_c = ws_d.cell(row=fila_excel, column=COL_INT, value=intereses_acum_anio)
        int_c.font = Font(name="Calibri", size=11, bold=True, color=C_VERDE_OSC)
        int_c.fill = fill(C_VERDE)
        int_c.number_format = fmt_usd
        int_c.alignment = aln(h="right")
        int_c.border = borde_top
        if tiene_pago:
            cx11 = ws_d.cell(row=fila_excel, column=COL_PAGO, value="")
            cx11.fill = fill(C_VERDE); cx11.border = borde_top
        fila_excel += 1

        ws_d.row_dimensions[fila_excel].height = 20
        ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_INT - 1)
        ca = ws_d.cell(row=fila_excel, column=1, value=f"   Capital + Intereses acumulados {anio_actual}")
        ca.font = Font(name="Calibri", size=10, italic=True, color=C_VERDE_OSC)
        ca.fill = fill(C_VERDE)
        ca.alignment = aln(h="right")
        ca.border = borde_std
        total_anio_c = ws_d.cell(row=fila_excel, column=COL_INT, value=capital_anio + intereses_acum_anio)
        total_anio_c.font = Font(name="Calibri", size=10, bold=True, italic=True, color=C_VERDE_OSC)
        total_anio_c.fill = fill(C_VERDE)
        total_anio_c.number_format = fmt_usd
        total_anio_c.alignment = aln(h="right")
        total_anio_c.border = borde_std
        if tiene_pago:
            cx11 = ws_d.cell(row=fila_excel, column=COL_PAGO, value="")
            cx11.fill = fill(C_VERDE); cx11.border = borde_std
        fila_excel += 1
        intereses_acum_anio = 0.0

    # CIERRE FINAL
    ws_d.row_dimensions[fila_excel].height = 36
    ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_CAP - 1)
    cf = ws_d.cell(row=fila_excel, column=1, value=f"CIERRE FINAL  {fecha_corte.strftime('%d/%m/%Y')}")
    cf.font = Font(name="Calibri", size=13, bold=True, color=C_DORADO_OSC)
    cf.fill = fill(C_DORADO)
    cf.alignment = aln(h="center")
    cf.border = borde_top
    cap_cf = ws_d.cell(row=fila_excel, column=COL_CAP, value=capital_total)
    cap_cf.font = Font(name="Calibri", size=13, bold=True, color=C_DORADO_OSC)
    cap_cf.fill = fill(C_DORADO)
    cap_cf.number_format = fmt_usd
    cap_cf.alignment = aln(h="right")
    cap_cf.border = borde_top
    int_cf = ws_d.cell(row=fila_excel, column=COL_INT, value=total_intereses)
    int_cf.font = Font(name="Calibri", size=13, bold=True, color=C_DORADO_OSC)
    int_cf.fill = fill(C_DORADO)
    int_cf.number_format = fmt_usd
    int_cf.alignment = aln(h="right")
    int_cf.border = borde_top
    if tiene_pago:
        total_paga_global = sum(
            float(r.get("interes_mes", 0) or 0)
            for r in det_rows if str(r.get("pago_intereses", "")).strip().upper() == "PAGA"
        )
        total_rein_global = total_intereses - total_paga_global
        cf11 = ws_d.cell(row=fila_excel, column=COL_PAGO,
            value=f"✅ Total pagado: ${total_paga_global:,.2f}  |  🔄 Total reinvertido: ${total_rein_global:,.2f}")
        cf11.font = Font(name="Calibri", size=10, bold=True, color=C_DORADO_OSC)
        cf11.fill = fill(C_DORADO)
        cf11.alignment = aln(h="center")
        cf11.border = borde_top
    fila_excel += 1

    ws_d.row_dimensions[fila_excel].height = 30
    ws_d.merge_cells(start_row=fila_excel, start_column=1, end_row=fila_excel, end_column=COL_INT - 1)
    cfa = ws_d.cell(row=fila_excel, column=1,
                    value="   TOTAL ACUMULADO  (Capital + Intereses)" if not tiene_pago
                    else "   TOTAL (Capital + Intereses pendientes)")
    cfa.font = Font(name="Calibri", size=12, bold=True, color=C_DORADO_OSC)
    cfa.fill = fill(C_DORADO)
    cfa.alignment = aln(h="right")
    cfa.border = borde_std
    total_final_valor = (capital_total + saldo_pendiente_global) if tiene_pago else (capital_total + total_intereses)
    total_cf = ws_d.cell(row=fila_excel, column=COL_INT, value=total_final_valor)
    total_cf.font = Font(name="Calibri", size=12, bold=True, color=C_DORADO_OSC)
    total_cf.fill = fill(C_DORADO)
    total_cf.number_format = fmt_usd
    total_cf.alignment = aln(h="right")
    total_cf.border = borde_std
    if tiene_pago:
        cx11 = ws_d.cell(row=fila_excel, column=COL_PAGO, value="")
        cx11.fill = fill(C_DORADO); cx11.border = borde_std

    # NOTA: se ha quitado la HOJA 3 "RESUMEN MENSUAL" porque duplicaba exactamente
    # la misma tabla "RESUMEN MENSUAL DE INTERESES" que ya aparece en la PORTADA.
    # El Excel ahora tiene solo 2 hojas: PORTADA y DETALLE.
    out = BytesIO()
    wb.save(out)
    return out.getvalue()


def generar_extractos(df_inv: pd.DataFrame, modo: str, inversor_elegido: str | None, anio: int, mes: int, solo_notas: bool = False):
    """Genera extractos para inversores.

    REGLA DEFINITIVA PARA EXTRACTOS:
    - SOLO se tienen en cuenta las filas cuya columna tipo_operacion sea exactamente NUEVA.
    - NO se tienen en cuenta reinversiones, canceladas, vacías ni cualquier otro valor.
    - Las reinversiones no modifican el extracto del inversor: el inversor cobra según su operación matriz NUEVA.
    """
    df = df_inv.copy()

    # Normalizamos columnas de texto necesarias.
    for col in [
        "inversor",
        "tipo_inversion",
        "subtipo_inversion",
        "nombre_activo",
        "tipo_operacion",
        "capital_nuevo_real",
        "motivo",
        "id_inversion",
    ]:
        if col in df.columns:
            df[col] = df[col].fillna("").astype(str).str.strip()

    # ==========================================
    # FILTRO PRINCIPAL DE EXTRACTOS
    # ==========================================
    # Según la regla definida: para extractos SOLO cuenta columna O / tipo_operacion = NUEVA.
    # Todo lo demás queda fuera: reinversion, cancelada, call, vacío, etc.
    if "tipo_operacion" not in df.columns:
        st.error("Falta la columna tipo_operacion en la hoja INVERSIONES. Para generar extractos debe existir y contener 'NUEVA'.")
        return []

    df["tipo_operacion_normalizada"] = df["tipo_operacion"].astype(str).str.strip().str.upper()

    # NUEVA: incluir, calcular hasta fecha de corte (ignorar fecha_final)
    # CANCELADA: incluir, calcular hasta fecha_final_inversion
    # REINVERSION y cualquier otro: excluir
    df = df[df["tipo_operacion_normalizada"].isin(["NUEVA", "CANCELADA"])].copy()

    # Filtro solo notas (subtipo estructurada)
    if solo_notas:
        df = df[df["subtipo_inversion"].str.upper() == "ESTRUCTURADA"].copy()

    if df.empty:
        return []

    if modo == "Un inversor" and inversor_elegido:
        df = df[df["inversor"].str.upper() == inversor_elegido.upper()].copy()

    if df.empty:
        return []

    fecha_corte = datetime(anio, mes, ultimo_dia_mes(anio, mes))

    filas = []
    for _, row in df.iterrows():
        fecha_inicio = row.get("fecha_inversion")
        if pd.isna(fecha_inicio):
            continue

        fecha_inicio_dt = pd.Timestamp(fecha_inicio).to_pydatetime()
        tipo_op = str(row.get("tipo_operacion_normalizada", "")).strip().upper()
        fecha_final_excel = row.get("fecha_final_inversion")

        if tipo_op == "CANCELADA":
            if pd.isna(fecha_final_excel):
                continue
            fecha_fin = min(pd.Timestamp(fecha_final_excel).to_pydatetime(), fecha_corte)
        else:
            # NUEVA: siempre hasta fecha de corte
            fecha_fin = fecha_corte

        if fecha_inicio_dt > fecha_fin:
            continue

        actual = datetime(fecha_inicio_dt.year, fecha_inicio_dt.month, 1)
        fin_mes = datetime(fecha_fin.year, fecha_fin.month, 1)

        while actual <= fin_mes:
            dias_mes = calendar.monthrange(actual.year, actual.month)[1]
            inicio_mes = datetime(actual.year, actual.month, 1)
            fin_mes_real = datetime(actual.year, actual.month, dias_mes)
            inicio_calc = max(fecha_inicio_dt, inicio_mes)
            fin_calc = min(fecha_fin, fin_mes_real)

            if inicio_calc <= fin_calc:
                dias = (fin_calc - inicio_calc).days + 1
                capital = float(row.get("capital_invertido", 0))
                interes_base = float(row.get("interes_inversor_anual", 0))

                # ── Tramo especial Biscafe / Crowe Bolivia / JR Real Estate ──
                # Hasta el 31/01/2026 cobraban al 5%; 01/02/2026-30/06/2026 al 7.5%;
                # desde el 01/07/2026 al 10%. Se aplica siempre, independientemente
                # de cuándo empezó la inversión ni de lo que diga interes_inversor_anual.
                INVERSORES_TRAMO = {"ROBERTO VISCAFE", "CROWE BOLIVIA", "JR REAL ESTATE"}
                inversor_upper = str(row.get("inversor", "")).strip().upper()

                if inversor_upper in INVERSORES_TRAMO:
                    fin_tramo1    = datetime(2026, 1, 31)
                    inicio_tramo2 = datetime(2026, 2, 1)
                    fin_tramo2    = datetime(2026, 6, 30)
                    inicio_tramo3 = datetime(2026, 7, 1)
                    interes_mes = 0.0
                    # Tramo 1: días del mes que caen en o antes del 31/01/2026 → 5%
                    if inicio_calc <= fin_tramo1:
                        fin_t1  = min(fin_calc, fin_tramo1)
                        dias_t1 = (fin_t1 - inicio_calc).days + 1
                        interes_mes += round((capital * 0.05 / 12) * dias_t1 / dias_mes, 2)
                    # Tramo 2: días entre 01/02/2026 y 30/06/2026 → 7.5%
                    ini_t2 = max(inicio_calc, inicio_tramo2)
                    fin_t2 = min(fin_calc, fin_tramo2)
                    if ini_t2 <= fin_t2:
                        dias_t2 = (fin_t2 - ini_t2).days + 1
                        interes_mes += round((capital * 0.075 / 12) * dias_t2 / dias_mes, 2)
                    # Tramo 3: días desde el 01/07/2026 → 10%
                    if fin_calc >= inicio_tramo3:
                        ini_t3  = max(inicio_calc, inicio_tramo3)
                        dias_t3 = (fin_calc - ini_t3).days + 1
                        interes_mes += round((capital * 0.10 / 12) * dias_t3 / dias_mes, 2)
                else:
                    interes_mes = round((capital * interes_base / 12) * dias / dias_mes, 2)
                mes_fecha = datetime(actual.year, actual.month, 1)
                pago_int_raw = str(row.get("pago_intereses", "REINVIERTE")).strip().upper()
                pago_intereses_val = "PAGA" if pago_int_raw == "PAGA" else "REINVIERTE"
                filas.append({
                    "mes_fecha": mes_fecha,
                    "fecha_inversion_orden": fecha_inicio_dt,
                    "inversor": row.get("inversor", ""),
                    "id_inversion": row.get("id_inversion", ""),
                    "tipo_inversion": row.get("tipo_inversion", ""),
                    "subtipo_inversion": row.get("subtipo_inversion", ""),
                    "nombre_activo": row.get("nombre_activo", ""),
                    "mes": f"{actual.month:02d}/{actual.year}",
                    "fecha_inversion": pd.Timestamp(fecha_inicio).strftime("%d/%m/%Y"),
                    "capital_invertido": capital,
                    "dias_devengados": dias,
                    "dias_mes": dias_mes,
                    "interes_mes": interes_mes,
                    "pago_intereses": pago_intereses_val,
                    "fecha_fin_op": fecha_fin.strftime("%d/%m/%Y"),
                })

            actual = datetime(actual.year + 1, 1, 1) if actual.month == 12 else datetime(actual.year, actual.month + 1, 1)

    resultado = pd.DataFrame(filas)
    if resultado.empty:
        return []

    resultado = resultado.sort_values(
        ["inversor", "mes_fecha", "fecha_inversion_orden", "id_inversion", "nombre_activo"],
        ascending=[True, True, True, True, True],
        kind="mergesort",
    ).reset_index(drop=True)

    archivos = []
    for inversor, grupo in resultado.groupby("inversor", sort=True):
        detalle = grupo.copy().sort_values(
            ["mes_fecha", "fecha_inversion_orden", "id_inversion", "nombre_activo"],
            ascending=[True, True, True, True],
            kind="mergesort",
        )

        totales_mes = (
            detalle.groupby(["mes_fecha", "mes"], as_index=False)["interes_mes"]
            .sum()
            .sort_values("mes_fecha")
            .rename(columns={"interes_mes": "total_mes"})
        )
        totales_mes = totales_mes[["mes", "total_mes"]]

        # Capital total activo: NUEVA + REINVERSION activas a fecha de corte
        # Las reinversiones son capital real del inversor aunque no generen intereses propios en el extracto
        base_inversor = df_inv[df_inv["inversor"].astype(str).str.upper() == str(inversor).upper()].copy()
        base_inversor["tipo_op_norm"] = base_inversor["tipo_operacion"].astype(str).str.strip().str.upper()
        base_inversor["fecha_inversion"] = parsear_fecha_robusta(base_inversor["fecha_inversion"])
        base_inversor["fecha_final_inversion"] = parsear_fecha_robusta(base_inversor["fecha_final_inversion"])
        base_inversor["capital_invertido"] = pd.to_numeric(base_inversor["capital_invertido"], errors="coerce").fillna(0)
        fecha_corte_ts = pd.Timestamp(fecha_corte).normalize()
        activas_corte = base_inversor[
            (base_inversor["tipo_op_norm"].isin(["NUEVA", "REINVERSION"]))
            & (base_inversor["fecha_inversion"].notna())
            & (base_inversor["fecha_inversion"] <= fecha_corte_ts)
            & (
                base_inversor["fecha_final_inversion"].isna()
                | (base_inversor["fecha_final_inversion"] >= fecha_corte_ts)
            )
        ].copy()
        capital_total = float(activas_corte["capital_invertido"].sum()) if not activas_corte.empty else 0.0

        detalle_exportar = detalle.drop(columns=["mes_fecha", "fecha_inversion_orden"], errors="ignore")
        # fecha_fin_op se mantiene en el Excel (col 11) para calcular capital activo al cierre mensual

        salida = BytesIO()
        with pd.ExcelWriter(salida, engine="openpyxl") as writer:
            totales_mes.to_excel(writer, sheet_name="TOTALES_MES", index=False)
            detalle_exportar.to_excel(writer, sheet_name="DETALLE", index=False)
        excel_crudo = salida.getvalue()
        nombre_archivo = f"extracto_{str(inversor).upper().replace(' ', '_')}_{fecha_corte.strftime('%d%m%Y')}.xlsx"
        archivos.append((nombre_archivo, formatear_extracto_excel_bytes(excel_crudo, str(inversor), fecha_corte), excel_crudo))
    return archivos

def _tab_importar_extracto_banco(df_control: pd.DataFrame):
    st.caption(
        "Sube el Excel de movimientos que descargas del broker (custodio de las notas "
        "estructuradas). La app clasifica cada fila, la casa con su NOTA por CUSIP cuando "
        "corresponde, y te enseña un resumen antes de guardar nada. Si subes un extracto que se "
        "solapa en fechas con uno ya importado, los movimientos repetidos se detectan solos y no "
        "se duplican."
    )

    archivo = st.file_uploader("Extracto del broker (.xlsx)", type=["xlsx"], key="importar_extracto_banco_file")
    if archivo is None:
        return

    try:
        df_nuevo = parsear_extracto_banco_bytes(archivo.read())
    except Exception as e:
        st.error(f"No se pudo leer el archivo: {e}")
        return

    df_nuevo = matchear_movimientos_con_notas(df_nuevo, df_control)

    ya_importados = _leer_movimientos_banco_cached()
    ids_ya_importados = (
        set(ya_importados["id_movimiento"].astype(str))
        if not ya_importados.empty and "id_movimiento" in ya_importados.columns else set()
    )
    df_nuevo["ya_importado"] = df_nuevo["id_movimiento"].isin(ids_ya_importados)

    n_nuevos = int((~df_nuevo["ya_importado"]).sum())
    n_repetidos = int(df_nuevo["ya_importado"].sum())
    st.info(
        f"{len(df_nuevo)} movimiento(s) leído(s) del archivo: {n_nuevos} nuevo(s), "
        f"{n_repetidos} ya estaban importados (se ignoran)."
    )

    df_mostrar = df_nuevo[~df_nuevo["ya_importado"]].copy()
    ocultar_sweep = st.checkbox(
        "Ocultar barrido de caja (FDIC Sweep) de esta vista previa", value=True, key="ocultar_sweep_importar",
        help="El sweep se guarda igualmente (forma parte del saldo real de la cuenta) — esto solo oculta las filas en esta tabla para que se vea más claro.",
    )
    if ocultar_sweep:
        df_mostrar = df_mostrar[df_mostrar["categoria"] != "RUIDO_SWEEP"]

    if df_mostrar.empty:
        st.caption("Nada que mostrar con el filtro actual.")
    else:
        df_preview = df_mostrar.copy()
        df_preview["Categoría"] = df_preview["categoria"].map(ETIQUETAS_CATEGORIA_BANCO).fillna(df_preview["categoria"])
        df_preview["Nota"] = pd.to_numeric(df_preview["nota_asociada"], errors="coerce").apply(
            lambda n: f"NOTA_{int(n):02d}" if pd.notna(n) else "—"
        )
        cols_preview = ["Fecha", "Categoría", "Monto", "cusip", "Nota", "Descripción de Activo", "Descripción de Transacción"]
        st.dataframe(
            df_preview[cols_preview].rename(columns={"cusip": "CUSIP"}).sort_values("Fecha"),
            use_container_width=True, hide_index=True,
        )

    sin_clasificar = int((df_nuevo.loc[~df_nuevo["ya_importado"], "categoria"] == "SIN_CLASIFICAR").sum())
    if sin_clasificar:
        st.warning(
            f"⚠️ {sin_clasificar} movimiento(s) nuevo(s) no se pudieron clasificar automáticamente "
            "— revísalos en la tabla (desmarca 'Ocultar barrido de caja' si hace falta) antes de guardar."
        )

    if n_nuevos == 0:
        st.caption("No hay movimientos nuevos que guardar.")
        return

    if st.button(f"💾 Guardar {n_nuevos} movimiento(s) nuevo(s) en contabilidad", type="primary"):
        hojas = leer_todas_las_hojas_excel()
        if not hojas:
            st.error("No se pudo leer el Excel actual para guardar.")
            return

        df_guardar = df_nuevo[~df_nuevo["ya_importado"]].copy()
        df_guardar["fecha_importacion"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        df_guardar = df_guardar.rename(columns={
            "Fecha": "fecha", "Monto": "monto",
            "Descripción de Activo": "emisor", "Descripción de Transacción": "descripcion",
        })
        df_guardar["fecha"] = df_guardar["fecha"].dt.strftime("%Y-%m-%d")
        columnas_guardar = [
            "id_movimiento", "fecha", "categoria", "monto", "cusip", "nota_asociada",
            "emisor", "descripcion", "fecha_importacion",
        ]
        df_guardar = df_guardar[columnas_guardar]

        existente = hojas.get(HOJA_MOVIMIENTOS_BANCO)
        if existente is None or existente.empty:
            hojas[HOJA_MOVIMIENTOS_BANCO] = df_guardar
        else:
            existente.columns = [str(c).strip().lower() for c in existente.columns]
            hojas[HOJA_MOVIMIENTOS_BANCO] = pd.concat([existente, df_guardar], ignore_index=True)

        guardar_excel_completo_desde_hojas(hojas)
        st.cache_data.clear()
        st.success(f"{n_nuevos} movimiento(s) guardado(s) en contabilidad.")
        st.rerun()


def _tab_dashboard_contabilidad():
    df_mov = _leer_movimientos_banco_cached()
    if df_mov.empty:
        st.info("Todavía no se ha importado ningún extracto — ve a la pestaña '📥 Importar extracto' para empezar.")
        return

    df_mov = df_mov.copy()
    df_mov["fecha"] = pd.to_datetime(df_mov.get("fecha"), errors="coerce")
    df_mov["monto"] = pd.to_numeric(df_mov.get("monto"), errors="coerce").fillna(0)
    df_mov = df_mov.dropna(subset=["fecha"])
    if df_mov.empty:
        st.info("No hay movimientos con fecha válida guardados todavía.")
        return

    st.caption(
        "Esta vista cubre solo la cuenta de custodia del broker (donde están las notas "
        "estructuradas) — no incluye Paraguay, MotoClick, Fútbol, Bolivia ni Bitcoin, que se "
        "siguen viendo en el Dashboard financiero general."
    )

    c1, c2 = st.columns(2)
    with c1:
        fecha_desde = st.date_input("Desde", value=df_mov["fecha"].min().date(), key="contab_fecha_desde")
    with c2:
        fecha_hasta = st.date_input("Hasta", value=df_mov["fecha"].max().date(), key="contab_fecha_hasta")

    df_rango = df_mov[(df_mov["fecha"] >= pd.Timestamp(fecha_desde)) & (df_mov["fecha"] <= pd.Timestamp(fecha_hasta))]

    # El saldo bancario SÍ incluye el sweep (es efectivo real de la cuenta); los totales por
    # categoría de abajo lo excluyen porque no aporta información económica, solo barrido interno.
    saldo_actual = float(df_mov["monto"].sum())
    ingresos_interes = float(df_rango.loc[df_rango["categoria"] == "INTERES_COBRADO", "monto"].sum())
    comisiones_gastos = float(df_rango.loc[df_rango["categoria"] == "COMISION_GASTO", "monto"].sum())
    aportaciones = float(df_rango.loc[df_rango["categoria"] == "APORTACION_CAPITAL", "monto"].sum())
    altas_inversion = float(df_rango.loc[df_rango["categoria"] == "ALTA_INVERSION", "monto"].sum())
    cancelaciones = float(df_rango.loc[df_rango["categoria"] == "CANCELACION_INVERSION", "monto"].sum())

    st.markdown("#### Resumen")
    c1, c2, c3 = st.columns(3)
    with c1:
        tarjeta_kpi("Saldo de la cuenta (histórico completo)", f"${saldo_actual:,.2f}", "Todos los movimientos importados, incluye sweep")
    with c2:
        tarjeta_kpi("Interés cobrado", f"${ingresos_interes:,.2f}", "En el rango seleccionado", "positivo")
    with c3:
        tarjeta_kpi("Comisiones y gastos", f"${comisiones_gastos:,.2f}", "En el rango seleccionado", "negativo" if comisiones_gastos < 0 else "normal")

    c4, c5, c6 = st.columns(3)
    with c4:
        tarjeta_kpi("Aportaciones de capital", f"${aportaciones:,.2f}", "En el rango seleccionado")
    with c5:
        tarjeta_kpi("Altas de inversión", f"${altas_inversion:,.2f}", "En el rango seleccionado")
    with c6:
        tarjeta_kpi("Cancelaciones / vencimientos", f"${cancelaciones:,.2f}", "En el rango seleccionado")

    st.markdown("---")
    st.markdown("#### Evolución del saldo")
    df_saldo_orden = df_mov.sort_values("fecha").copy()
    df_saldo_orden["saldo_acumulado"] = df_saldo_orden["monto"].cumsum()
    import plotly.graph_objects as go
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=df_saldo_orden["fecha"], y=df_saldo_orden["saldo_acumulado"],
        mode="lines", line=dict(color="#0e2338", width=2), fill="tozeroy",
    ))
    fig.update_layout(height=320, margin=dict(l=10, r=10, t=10, b=10), yaxis_title="Saldo ($)")
    st.plotly_chart(fig, use_container_width=True)

    st.markdown("---")
    st.markdown("#### Movimientos")
    categoria_filtro = st.selectbox(
        "Filtrar por categoría", ["Todas"] + list(ETIQUETAS_CATEGORIA_BANCO.values()), key="contab_filtro_categoria",
    )
    df_tabla = df_rango.copy()
    df_tabla["Categoría"] = df_tabla["categoria"].map(ETIQUETAS_CATEGORIA_BANCO).fillna(df_tabla["categoria"])
    if categoria_filtro != "Todas":
        df_tabla = df_tabla[df_tabla["Categoría"] == categoria_filtro]
    df_tabla["Nota"] = pd.to_numeric(df_tabla.get("nota_asociada"), errors="coerce").apply(
        lambda n: f"NOTA_{int(n):02d}" if pd.notna(n) else "—"
    )
    cols_tabla = [c for c in ["fecha", "Categoría", "monto", "cusip", "Nota", "emisor", "descripcion"] if c in df_tabla.columns]
    st.dataframe(
        df_tabla[cols_tabla].rename(columns={
            "fecha": "Fecha", "monto": "Monto", "cusip": "CUSIP", "emisor": "Emisor", "descripcion": "Descripción",
        }).sort_values("Fecha", ascending=False),
        use_container_width=True, hide_index=True,
    )


def seccion_contabilidad_banco():
    df_inv, df_cal, df_control = cargar_excel_completo()
    st.header("🏦 Contabilidad")

    tab_dashboard, tab_importar = st.tabs(["📊 Contabilidad", "📥 Importar extracto"])
    with tab_dashboard:
        _tab_dashboard_contabilidad()
    with tab_importar:
        _tab_importar_extracto_banco(df_control)


def seccion_extractos():
    df_inv, _, _ = cargar_excel_completo()
    st.header("📤 Extractos")

    # ── Parámetros comunes ────────────────────────────────────────────────────
    modo = st.radio("¿Para quién?", ["Todos los inversores", "Un inversor"], horizontal=True, key="ext_modo")
    inversores = sorted([x for x in df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).unique() if x.strip()])
    inversor = st.selectbox("Inversor", inversores, key="ext_inversor") if modo == "Un inversor" and inversores else None
    c1, c2 = st.columns(2)
    anio = int(c1.number_input("Año", 2020, 2100, pd.Timestamp.today().year, key="ext_anio"))
    mes  = int(c2.number_input("Mes", 1, 12, pd.Timestamp.today().month, key="ext_mes"))

    st.divider()

    # ── Dos pestañas: Descargar / Enviar por email ────────────────────────────
    tab_descargar, tab_email = st.tabs(["⬇️  Descargar", "📧  Enviar por email"])

    with tab_descargar:
        st.caption("Genera el extracto en Excel y descárgalo directamente.")
        ocultar_activo = st.checkbox(
            "🔒 Ocultar en qué está invertido (igual que en el Portal de inversor) — recomendado si se lo vas a dar al inversor",
            value=True, key="ext_ocultar_activo",
        )

        col_btn1, col_btn2 = st.columns(2)
        btn_completo = col_btn1.button("📄 Generar extracto completo", type="primary", key="ext_btn_descargar")
        btn_notas    = col_btn2.button("📑 Generar extracto solo notas", type="secondary", key="ext_btn_notas")

        def _descargar_extractos(solo_notas: bool):
            modo_gen = "Un inversor" if inversor else "Todos"
            archivos = generar_extractos(df_inv, modo_gen, inversor, anio, mes, solo_notas=solo_notas)
            sufijo = "_NOTAS" if solo_notas else ""
            if not archivos:
                st.warning("No se han generado extractos. Revisa el Excel o la fecha seleccionada.")
            elif len(archivos) == 1:
                t = archivos[0]
                nombre = t[0].replace(".xlsx", f"{sufijo}.xlsx")
                contenido = preparar_extracto_privado_inversor(t[1]) if ocultar_activo else t[1]
                st.success(f"Extracto generado: {nombre}")
                key_dl = f"ext_dl_uno{'_notas' if solo_notas else ''}"
                st.download_button("⬇️ Descargar", contenido, file_name=nombre,
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                   key=key_dl)
            else:
                zip_buffer = BytesIO()
                with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                    for t in archivos:
                        nombre_zip = t[0].replace(".xlsx", f"{sufijo}.xlsx")
                        contenido_zip = preparar_extracto_privado_inversor(t[1]) if ocultar_activo else t[1]
                        zf.writestr(nombre_zip, contenido_zip)
                st.success(f"Se han generado {len(archivos)} extractos.")
                key_zip = f"ext_dl_zip{'_notas' if solo_notas else ''}"
                zip_name = f"extractos{'_notas' if solo_notas else ''}_{mes}_{anio}.zip"
                st.download_button("⬇️ Descargar todos en ZIP", zip_buffer.getvalue(),
                                   file_name=zip_name,
                                   mime="application/zip", key=key_zip)

        if btn_completo:
            _descargar_extractos(solo_notas=False)
        if btn_notas:
            _descargar_extractos(solo_notas=True)

    with tab_email:
        st.caption("Genera el extracto en PDF y envíalo directamente al email del inversor.")
        seccion_envio_extractos_email(df_inv, generar_extractos, anio_override=anio, mes_override=mes, inversor_override=inversor, modo_override=modo)



# =========================
# GESTIÓN DE EXCEL DESDE LA APP
# =========================
def leer_todas_las_hojas_excel() -> dict:
    """Lee todas las hojas del archivo Excel para poder conservarlas al guardar.

    BUG REAL encontrado (27/08/2026): antes solo se descargaba de Drive si el archivo local
    todavía no existía. Eso significa que si alguien edita el Google Sheet a mano (por ejemplo,
    borrando filas) y justo después la app escribe algo (crear un préstamo, una nueva inversión,
    etc.), esta función leía la copia local ANTIGUA — sin la edición manual — y esa base
    desactualizada era la que se subía de vuelta a Drive, revirtiendo en silencio el cambio manual.
    Ahora se refresca siempre desde Drive antes de leer, igual que ya hace
    _cargar_excel_completo_desde_drive.
    """
    try:
        descargar_excel_desde_drive()
        hojas = pd.read_excel(ARCHIVO, sheet_name=None)
        return {str(nombre): df for nombre, df in hojas.items()}
    except Exception:
        return {}


def excel_hojas_a_bytes(hojas: dict) -> bytes:
    """Convierte un diccionario de hojas en un Excel descargable."""
    salida = BytesIO()
    with pd.ExcelWriter(salida, engine="openpyxl") as writer:
        for nombre_hoja, df in hojas.items():
            nombre_limpio = str(nombre_hoja)[:31] if str(nombre_hoja).strip() else "Hoja"
            if df is None:
                df = pd.DataFrame()
            df.to_excel(writer, sheet_name=nombre_limpio, index=False)
    return salida.getvalue()


def subir_excel_a_drive(hojas: dict) -> tuple[bool, str]:
    """
    Sube el Excel actualizado (todas las hojas) al Google Sheet del fondo,
    usando la cuenta de servicio configurada en st.secrets['gcp_service_account'].
    Devuelve (exito, mensaje).
    """
    if "gcp_service_account" not in st.secrets:
        return False, "No hay credenciales de Google configuradas en Secrets (falta [gcp_service_account])."
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseUpload
    except ImportError:
        return False, "Faltan las librerías google-api-python-client y google-auth. Añádelas a requirements.txt."

    try:
        credenciales = service_account.Credentials.from_service_account_info(
            dict(st.secrets["gcp_service_account"]),
            scopes=["https://www.googleapis.com/auth/drive"],
        )
        servicio = build("drive", "v3", credentials=credenciales)
        contenido = excel_hojas_a_bytes(hojas)
        media = MediaIoBaseUpload(
            BytesIO(contenido),
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            resumable=False,
        )
        # NOTA: NO se especifica 'body' con mimeType aquí — Drive ya convierte automáticamente
        # el xlsx subido al formato nativo de Google Sheets porque el archivo destino YA es un
        # Google Sheet. Forzar el mimeType de destino explícitamente causaba el error
        # "Invalid MIME type provided for the uploaded content".
        archivo_actualizado = servicio.files().update(
            fileId=GDRIVE_FILE_ID,
            media_body=media,
            fields="id, modifiedTime, mimeType",
        ).execute()
        return True, f"Excel actualizado correctamente en Google Drive (modificado: {archivo_actualizado.get('modifiedTime', '?')})."
    except Exception as e:
        return False, f"No se pudo subir el Excel a Google Drive: {e}"


def guardar_excel_completo_desde_hojas(hojas: dict):
    """Guarda todas las hojas en inversiones.xlsx localmente y, si hay credenciales
    de Google configuradas, también las sube automáticamente al Google Sheet del fondo."""
    contenido = excel_hojas_a_bytes(hojas)
    with open(ARCHIVO, "wb") as f:
        f.write(contenido)
    st.cache_data.clear()

    try:
        from postgres_writer import escribir_hojas_postgres
        escribir_hojas_postgres(hojas)
    except Exception as e:
        print(f"[postgres_writer] no disponible, se omite escritura en paralelo: {e}", file=sys.stderr)

    if "gcp_service_account" in st.secrets:
        exito, mensaje = subir_excel_a_drive(hojas)
        if exito:
            st.success(f"☁️ {mensaje}")
        else:
            st.warning(
                f"⚠️ Se guardó localmente pero no se pudo subir a Google Drive automáticamente: {mensaje}\n\n"
                "Puedes subirlo manualmente desde 'Gestión de Excel' → 'Descargar copia'."
            )


def log_uso_ia(usuario: str, tipo: str, tokens_input: int, tokens_output: int,
                modelo: str = "claude-sonnet-4-5") -> None:
    """
    Registra una llamada al asistente de IA en la hoja LOG_IA_USO, para poder ver después
    qué usuarios usan el asistente y estimar el gasto en dólares de cada uno.
    'tipo' identifica el canal/función: 'fondo' (asistente interno), 'inversor' (portal del
    inversor), etc.

    Igual que guardar_borrador_nota, reescribe y sube todo el Excel — es el mismo patrón que
    ya usa el resto de la app para persistir en Drive. Por eso NUNCA debe romper la conversación
    si falla: se traga cualquier error en silencio (sin st.error), porque perder una fila de
    log es mucho menos grave que cortar una respuesta del asistente al usuario.
    """
    try:
        precio_in, precio_out = PRECIOS_MODELOS_IA.get(modelo, PRECIO_MODELO_IA_POR_DEFECTO)
        coste = (tokens_input / 1_000_000) * precio_in + (tokens_output / 1_000_000) * precio_out

        hojas = leer_todas_las_hojas_excel()
        if not hojas:
            return
        fila_nueva = pd.DataFrame([{
            "USUARIO": usuario,
            "TIPO": tipo,
            "MODELO": modelo,
            "FECHA_HORA": pd.Timestamp.now(),
            "TOKENS_INPUT": int(tokens_input or 0),
            "TOKENS_OUTPUT": int(tokens_output or 0),
            "COSTE_ESTIMADO_USD": round(coste, 6),
        }])
        if HOJA_LOG_IA in hojas and not hojas[HOJA_LOG_IA].empty:
            hojas[HOJA_LOG_IA] = pd.concat([hojas[HOJA_LOG_IA], fila_nueva], ignore_index=True)
        else:
            hojas[HOJA_LOG_IA] = fila_nueva

        contenido = excel_hojas_a_bytes(hojas)
        with open(ARCHIVO, "wb") as f:
            f.write(contenido)
        st.cache_data.clear()
        try:
            from postgres_writer import escribir_hojas_postgres
            escribir_hojas_postgres(hojas)
        except Exception:
            pass
        if "gcp_service_account" in st.secrets:
            subir_excel_a_drive(hojas)
    except Exception:
        pass


@st.cache_data(ttl=120, show_spinner=False)
def _leer_log_ia_cached() -> pd.DataFrame:
    """Versión cacheada (2 min) de la lectura de LOG_IA_USO para el panel de 'Uso IA' —
    evita releer/descargar el Excel en cada rerun de Streamlit."""
    df = leer_hoja_excel(HOJA_LOG_IA)
    if df.empty:
        return df
    df.columns = [str(c).strip().lower() for c in df.columns]
    if "fecha_hora" in df.columns:
        df["fecha_hora"] = pd.to_datetime(df["fecha_hora"], errors="coerce")
    for col in ["tokens_input", "tokens_output", "coste_estimado_usd"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    return df


def seccion_uso_ia():
    """Panel de administrador: qué usuarios usan el asistente de IA y cuánto gasto estimado
    (en dólares) genera cada uno, a partir del log guardado en LOG_IA_USO."""
    st.markdown("## 📊 Uso del asistente de IA")
    st.caption(
        "Gasto estimado a partir de los tokens reales que devuelve la API de Anthropic en cada "
        "llamada. Es una estimación (tarifa configurada en el código) — para el gasto exacto, "
        "consulta la consola de facturación de Anthropic."
    )
    df = _leer_log_ia_cached()
    if df.empty:
        st.info("Todavía no hay ningún uso registrado del asistente de IA desde que se activó este log.")
        return

    c1, c2 = st.columns(2)
    fecha_min = df["fecha_hora"].min().date()
    fecha_max = df["fecha_hora"].max().date()
    desde = c1.date_input("Desde", value=fecha_min, key="uso_ia_desde")
    hasta = c2.date_input("Hasta", value=fecha_max, key="uso_ia_hasta")

    mask = (df["fecha_hora"].dt.date >= desde) & (df["fecha_hora"].dt.date <= hasta)
    df_f = df[mask]
    if df_f.empty:
        st.warning("No hay uso registrado en ese rango de fechas.")
        return

    mostrar_metricas("Total del periodo", [
        ("Mensajes", f"{len(df_f):,}"),
        ("Tokens totales", f"{int(df_f['tokens_input'].sum() + df_f['tokens_output'].sum()):,}"),
        ("Coste estimado", f"${df_f['coste_estimado_usd'].sum():,.2f}"),
    ])

    st.markdown("### Por usuario")
    resumen = (
        df_f.groupby("usuario")
        .agg(mensajes=("usuario", "count"),
             tokens_input=("tokens_input", "sum"),
             tokens_output=("tokens_output", "sum"),
             coste_estimado_usd=("coste_estimado_usd", "sum"))
        .reset_index()
        .sort_values("coste_estimado_usd", ascending=False)
    )
    resumen["coste_estimado_usd"] = resumen["coste_estimado_usd"].map(lambda v: f"${v:,.4f}")
    st.dataframe(resumen, use_container_width=True, hide_index=True)


def guardar_borrador_nota(tipo_wizard: str, numero_nota: int, datos: dict) -> bool:
    """
    Guarda automáticamente (sin que el usuario tenga que pulsar nada) el resultado de una
    extracción de IA en la hoja BORRADORES_NOTAS, para que sobreviva a reinicios de la app
    (actualizaciones de código, redeploys, etc.) — independiente de CONTROL_NOTAS/CALENDARIO_NOTAS.
    Devuelve True si quedó realmente a salvo en Google Drive (persistente de verdad),
    False si solo se guardó en el disco local (se perderá en el próximo redeploy).
    """
    import json as _json
    try:
        hojas = leer_todas_las_hojas_excel()
        if not hojas:
            st.error("⚠️ No se pudo leer el Excel para guardar el borrador — el resultado de esta extracción NO está a salvo todavía.")
            return False
        fila_nueva = pd.DataFrame([{
            "TIPO": tipo_wizard,
            "NOTA": int(numero_nota),
            "JSON_DATOS": _json.dumps(datos, ensure_ascii=False),
            "FECHA_GUARDADO": pd.Timestamp.now(),
        }])
        if "BORRADORES_NOTAS" in hojas and not hojas["BORRADORES_NOTAS"].empty:
            df_b = hojas["BORRADORES_NOTAS"]
            df_b = df_b[~((df_b["TIPO"] == tipo_wizard) & (pd.to_numeric(df_b["NOTA"], errors="coerce") == numero_nota))]
            hojas["BORRADORES_NOTAS"] = pd.concat([df_b, fila_nueva], ignore_index=True)
        else:
            hojas["BORRADORES_NOTAS"] = fila_nueva
        contenido = excel_hojas_a_bytes(hojas)
        with open(ARCHIVO, "wb") as f:
            f.write(contenido)
        st.cache_data.clear()
        try:
            from postgres_writer import escribir_hojas_postgres
            escribir_hojas_postgres(hojas)
        except Exception:
            pass

        if "gcp_service_account" not in st.secrets:
            st.warning("⚠️ Borrador guardado solo localmente (no hay credenciales de Google configuradas) — se perderá si el servidor se reinicia o actualizas el código. Sube el Excel a Drive manualmente si quieres conservarlo.")
            return False

        exito, mensaje = subir_excel_a_drive(hojas)
        if exito:
            st.success(f"📝 Borrador guardado y a salvo en Google Drive — sobrevivirá aunque cambies el código.")
            return True
        else:
            st.error(f"⚠️ El borrador se guardó solo localmente — la subida a Google Drive falló: {mensaje}. Si actualizas el código o el servidor se reinicia antes de solucionarlo, se PERDERÁ esta extracción.")
            return False
    except Exception as e:
        st.error(f"⚠️ No se pudo guardar el borrador (error: {e}). Esta extracción no está a salvo todavía.")
        return False


def cargar_borrador_nota(tipo_wizard: str, numero_nota: int) -> dict | None:
    """Recupera el borrador guardado (si existe) para esa nota y ese wizard."""
    import json as _json
    try:
        df_b = leer_hoja_excel("BORRADORES_NOTAS")
        if df_b.empty:
            return None
        df_b.columns = [str(c).strip().upper() for c in df_b.columns]
        fila = df_b[(df_b["TIPO"] == tipo_wizard) & (pd.to_numeric(df_b["NOTA"], errors="coerce") == numero_nota)]
        if fila.empty:
            return None
        return _json.loads(fila.iloc[-1]["JSON_DATOS"])
    except Exception:
        return None


def borrar_borrador_nota(tipo_wizard: str, numero_nota: int):
    """Elimina el borrador de una nota (se llama cuando ya se ha guardado de verdad, o al descartarla)."""
    try:
        hojas = leer_todas_las_hojas_excel()
        if not hojas or "BORRADORES_NOTAS" not in hojas or hojas["BORRADORES_NOTAS"].empty:
            return
        df_b = hojas["BORRADORES_NOTAS"]
        hojas["BORRADORES_NOTAS"] = df_b[~((df_b["TIPO"] == tipo_wizard) & (pd.to_numeric(df_b["NOTA"], errors="coerce") == numero_nota))]
        contenido = excel_hojas_a_bytes(hojas)
        with open(ARCHIVO, "wb") as f:
            f.write(contenido)
        st.cache_data.clear()
        try:
            from postgres_writer import escribir_hojas_postgres
            escribir_hojas_postgres(hojas)
        except Exception:
            pass
        if "gcp_service_account" in st.secrets:
            exito, mensaje = subir_excel_a_drive(hojas)
            if not exito:
                st.warning(f"⚠️ El borrador se borró localmente pero no se pudo sincronizar con Drive: {mensaje}")
    except Exception as e:
        st.warning(f"⚠️ No se pudo borrar el borrador correctamente: {e}")


# ══════════════════════════════════════════════════════════════════════════
# BORRADORES DE ALTA DE INVERSIÓN (hoja BORRADORES_INVERSIONES)
# Mismo patrón que BORRADORES_NOTAS: se guarda como JSON bajo una "CLAVE" libre
# que Yuri elige (ej. "PAM-NOTA28"), persiste en Drive, y sobrevive a redeploys.
# ══════════════════════════════════════════════════════════════════════════

def guardar_borrador_inversion(clave: str, datos: dict) -> bool:
    """Igual patrón que guardar_borrador_nota, pero para altas/cierres/reinversiones de INVERSIONES."""
    import json as _json
    try:
        hojas = leer_todas_las_hojas_excel()
        if not hojas:
            st.error("⚠️ No se pudo leer el Excel para guardar el borrador — esta operación NO está a salvo todavía.")
            return False
        fila_nueva = pd.DataFrame([{
            "CLAVE": str(clave),
            "JSON_DATOS": _json.dumps(datos, ensure_ascii=False, default=str),
            "FECHA_GUARDADO": pd.Timestamp.now(),
        }])
        if "BORRADORES_INVERSIONES" in hojas and not hojas["BORRADORES_INVERSIONES"].empty:
            df_b = hojas["BORRADORES_INVERSIONES"]
            df_b = df_b[df_b["CLAVE"].astype(str) != str(clave)]
            hojas["BORRADORES_INVERSIONES"] = pd.concat([df_b, fila_nueva], ignore_index=True)
        else:
            hojas["BORRADORES_INVERSIONES"] = fila_nueva
        contenido = excel_hojas_a_bytes(hojas)
        with open(ARCHIVO, "wb") as f:
            f.write(contenido)
        st.cache_data.clear()
        try:
            from postgres_writer import escribir_hojas_postgres
            escribir_hojas_postgres(hojas)
        except Exception:
            pass

        if "gcp_service_account" not in st.secrets:
            st.warning("⚠️ Borrador guardado solo localmente (no hay credenciales de Google configuradas) — se perderá si el servidor se reinicia o actualizas el código.")
            return False

        exito, mensaje = subir_excel_a_drive(hojas)
        if exito:
            st.success("📝 Borrador guardado y a salvo en Google Drive — sobrevivirá aunque cambies el código.")
            return True
        else:
            st.error(f"⚠️ El borrador se guardó solo localmente — la subida a Google Drive falló: {mensaje}.")
            return False
    except Exception as e:
        st.error(f"⚠️ No se pudo guardar el borrador (error: {e}). Esta operación no está a salvo todavía.")
        return False


def cargar_borrador_inversion(clave: str) -> dict | None:
    """Recupera el borrador guardado (si existe) para esa clave."""
    import json as _json
    try:
        df_b = leer_hoja_excel("BORRADORES_INVERSIONES")
        if df_b.empty:
            return None
        df_b.columns = [str(c).strip().upper() for c in df_b.columns]
        fila = df_b[df_b["CLAVE"].astype(str) == str(clave)]
        if fila.empty:
            return None
        return _json.loads(fila.iloc[-1]["JSON_DATOS"])
    except Exception:
        return None


def listar_claves_borradores_inversion() -> list:
    """Lista las claves de borradores de inversión guardados, para poder retomarlos."""
    try:
        df_b = leer_hoja_excel("BORRADORES_INVERSIONES")
        if df_b.empty:
            return []
        df_b.columns = [str(c).strip().upper() for c in df_b.columns]
        return sorted(df_b["CLAVE"].dropna().astype(str).unique().tolist())
    except Exception:
        return []


def borrar_borrador_inversion(clave: str):
    """Elimina el borrador de una operación (se llama tras guardarla de verdad, o al descartarla)."""
    try:
        hojas = leer_todas_las_hojas_excel()
        if not hojas or "BORRADORES_INVERSIONES" not in hojas or hojas["BORRADORES_INVERSIONES"].empty:
            return
        df_b = hojas["BORRADORES_INVERSIONES"]
        hojas["BORRADORES_INVERSIONES"] = df_b[df_b["CLAVE"].astype(str) != str(clave)]
        contenido = excel_hojas_a_bytes(hojas)
        with open(ARCHIVO, "wb") as f:
            f.write(contenido)
        st.cache_data.clear()
        try:
            from postgres_writer import escribir_hojas_postgres
            escribir_hojas_postgres(hojas)
        except Exception:
            pass
        if "gcp_service_account" in st.secrets:
            exito, mensaje = subir_excel_a_drive(hojas)
            if not exito:
                st.warning(f"⚠️ El borrador se borró localmente pero no se pudo sincronizar con Drive: {mensaje}")
    except Exception as e:
        st.warning(f"⚠️ No se pudo borrar el borrador correctamente: {e}")


def _valores_conocidos_columna(df: pd.DataFrame, columna: str) -> list:
    """Valores distintos ya usados en una columna de INVERSIONES — para ofrecerlos como opciones
    en el formulario de alta en vez de adivinar la convención exacta (metodo_calculo,
    activo_generador_interes, cuenta_cobro, motivo...)."""
    if df is None or df.empty or columna not in df.columns:
        return []
    valores = df[columna].astype(str).str.strip()
    valores = valores[(valores != "") & (valores.str.lower() != "nan")]
    return sorted(valores.unique().tolist())


def _mapa_columnas_reales(hojas: dict, nombre_hoja: str) -> dict:
    """Mapa {nombre_columna_en_minusculas: nombre_columna_real_en_el_excel}, para escribir una fila
    nueva respetando exactamente el nombre/formato de columna que ya tiene la hoja en Drive, sin
    adivinar mayúsculas ni tener que igualar el orden a mano."""
    if not hojas or nombre_hoja not in hojas or hojas[nombre_hoja] is None:
        return {}
    return {str(c).strip().lower(): c for c in hojas[nombre_hoja].columns}


def aplicar_formula_simple(df: pd.DataFrame, operacion: str, columna_a: str, columna_b: str | None, nueva_columna: str) -> pd.DataFrame:
    """Aplica cálculos tipo Excel básicos sobre columnas numéricas."""
    out = df.copy()
    if not nueva_columna or not str(nueva_columna).strip():
        nueva_columna = "columna_calculada"
    nueva_columna = str(nueva_columna).strip()

    a = pd.to_numeric(out[columna_a], errors="coerce") if columna_a in out.columns else pd.Series(0, index=out.index)
    b = pd.to_numeric(out[columna_b], errors="coerce") if columna_b and columna_b in out.columns else pd.Series(0, index=out.index)

    if operacion == "Sumar A + B":
        out[nueva_columna] = a.fillna(0) + b.fillna(0)
    elif operacion == "Restar A - B":
        out[nueva_columna] = a.fillna(0) - b.fillna(0)
    elif operacion == "Multiplicar A x B":
        out[nueva_columna] = a.fillna(0) * b.fillna(0)
    elif operacion == "Dividir A / B":
        out[nueva_columna] = a / b.replace(0, pd.NA)
    elif operacion == "Porcentaje A sobre B":
        out[nueva_columna] = (a / b.replace(0, pd.NA)) * 100
    elif operacion == "Interés mensual: capital x interés / 12":
        if "capital_invertido" in out.columns and "interes_inversor_anual" in out.columns:
            capital = pd.to_numeric(out["capital_invertido"], errors="coerce").fillna(0)
            interes = pd.to_numeric(out["interes_inversor_anual"], errors="coerce").fillna(0)
            out[nueva_columna] = capital * interes / 12
        else:
            st.warning("Para esta fórmula necesitas las columnas capital_invertido e interes_inversor_anual.")
    elif operacion == "Interés nota mensual: capital x interés nota / 12":
        if "capital_invertido" in out.columns and "interes_nota_anual" in out.columns:
            capital = pd.to_numeric(out["capital_invertido"], errors="coerce").fillna(0)
            interes = pd.to_numeric(out["interes_nota_anual"], errors="coerce").fillna(0)
            out[nueva_columna] = capital * interes / 12
        else:
            st.warning("Para esta fórmula necesitas las columnas capital_invertido e interes_nota_anual.")
    return out


def mostrar_sumatorias_excel(df: pd.DataFrame):
    """Muestra sumatorias rápidas de columnas numéricas como apoyo tipo Excel."""
    if df is None or df.empty:
        return
    numericas = []
    for col in df.columns:
        serie = pd.to_numeric(df[col], errors="coerce")
        if serie.notna().sum() > 0:
            numericas.append(col)
    if not numericas:
        st.info("No hay columnas numéricas para sumar.")
        return
    seleccion = st.multiselect("Columnas para calcular sumatorias", numericas, default=numericas[: min(4, len(numericas))])
    if seleccion:
        cols = st.columns(min(4, len(seleccion)))
        for i, col in enumerate(seleccion):
            total = pd.to_numeric(df[col], errors="coerce").sum()
            cols[i % len(cols)].metric(f"Suma {col}", fmt(total))


def seccion_gestion_excel():
    st.markdown("## Gestión de Excel")
    st.caption("Sube, edita, calcula, guarda y descarga la base de datos directamente desde la app.")

    tab_subir, tab_editar, tab_usuarios, tab_descargar = st.tabs(
        ["Subir Excel", "Editar y calcular", "Usuarios y contraseñas", "Descargar copia"]
    )

    with tab_subir:
        st.subheader("Recargar Excel desde Google Drive")
        st.info(
            "El Excel se lee desde Google Drive. "
            "Si has actualizado el archivo en Drive, pulsa el botón para recargar."
        )
        if st.button("🔄 Recargar Excel desde Google Drive", type="primary"):
            st.cache_data.clear()
            for k in list(st.session_state.keys()):
                if str(k).startswith("excel_editor_"):
                    del st.session_state[k]
            ok = descargar_excel_desde_drive()
            if ok:
                st.success("Excel recargado correctamente desde Google Drive.")
            st.rerun()

        st.markdown("---")
        st.subheader("🗄️ Traer cambios de Drive a Postgres ahora")
        st.caption(
            "La app lee de Postgres, que normalmente solo se actualiza desde Drive en cada "
            "despliegue. Si acabas de editar algo a mano en el Excel de Drive y quieres verlo "
            "reflejado ya, sin esperar a un redeploy, pulsa este botón."
        )
        if st.button("🗄️ Traer cambios de Drive a Postgres ahora", type="primary", key="btn_sync_drive_postgres_manual"):
            with st.spinner("Descargando el Excel de Drive y actualizando Postgres..."):
                try:
                    import init_db
                    resultado = init_db.main()
                except Exception as e:
                    resultado = {"ok": False, "detalle": [], "error": str(e)}
            if resultado.get("ok"):
                st.success("Postgres actualizado con lo último de Drive.")
                with st.expander("Ver detalle de filas sincronizadas por hoja"):
                    for linea in resultado.get("detalle", []):
                        st.caption(linea)
                st.cache_data.clear()
                st.rerun()
            else:
                st.error(f"No se pudo sincronizar: {resultado.get('error')}")

        st.markdown("---")
        st.subheader("O sube un Excel manualmente")
        st.caption("Si prefieres subir el archivo directamente, recuerda actualizarlo también en Google Drive para que sea permanente.")
        archivo_subido = st.file_uploader("Sube el archivo actualizado", type=["xlsx"])
        if archivo_subido is not None:
            nombre = archivo_subido.name.strip()
            if nombre != ARCHIVO:
                st.warning(f"El archivo se llama '{nombre}'. El sistema trabaja con '{ARCHIVO}'. Se guardará igualmente como {ARCHIVO}.")
            if st.button("Reemplazar Excel actual", type="primary"):
                with open(ARCHIVO, "wb") as f:
                    f.write(archivo_subido.read())
                st.cache_data.clear()
                for k in list(st.session_state.keys()):
                    if str(k).startswith("excel_editor_"):
                        del st.session_state[k]
                st.success("Excel actualizado. Recuerda actualizarlo también en Google Drive para que sea permanente.")
                st.rerun()

    with tab_editar:
        hojas = leer_todas_las_hojas_excel()
        if not hojas:
            st.error("No se ha podido leer el Excel actual.")
            return

        hoja = st.selectbox("Selecciona la hoja que quieres editar", list(hojas.keys()))
        editor_key = f"excel_editor_{hoja}"

        c1, c2 = st.columns([1, 1])
        if editor_key not in st.session_state:
            st.session_state[editor_key] = hojas[hoja].copy()
        if c1.button("Recargar hoja desde el Excel"):
            st.session_state[editor_key] = hojas[hoja].copy()
            st.rerun()
        if c2.button("Limpiar caché de datos"):
            st.cache_data.clear()
            st.success("Caché limpiada.")

        st.info("Puedes editar celdas, añadir filas nuevas y después guardar los cambios en el Excel.")
        df_editado = st.data_editor(
            st.session_state[editor_key],
            use_container_width=True,
            num_rows="dynamic",
            key=f"data_editor_{hoja}",
        )
        st.session_state[editor_key] = df_editado

        with st.expander("Sumatorias rápidas", expanded=True):
            mostrar_sumatorias_excel(df_editado)

        with st.expander("Añadir columna calculada tipo fórmula", expanded=False):
            columnas = list(df_editado.columns)
            columnas_numericas = [c for c in columnas if pd.to_numeric(df_editado[c], errors="coerce").notna().sum() > 0]
            if not columnas_numericas:
                st.info("No hay columnas numéricas disponibles para crear fórmulas.")
            else:
                operacion = st.selectbox(
                    "Fórmula",
                    [
                        "Sumar A + B",
                        "Restar A - B",
                        "Multiplicar A x B",
                        "Dividir A / B",
                        "Porcentaje A sobre B",
                        "Interés mensual: capital x interés / 12",
                        "Interés nota mensual: capital x interés nota / 12",
                    ],
                )
                c1, c2, c3 = st.columns(3)
                columna_a = c1.selectbox("Columna A", columnas_numericas)
                columna_b = c2.selectbox("Columna B", columnas_numericas) if operacion not in ["Interés mensual: capital x interés / 12", "Interés nota mensual: capital x interés nota / 12"] else None
                nueva_columna = c3.text_input("Nombre nueva columna", value="columna_calculada")
                if st.button("Aplicar fórmula a la hoja"):
                    st.session_state[editor_key] = aplicar_formula_simple(df_editado, operacion, columna_a, columna_b, nueva_columna)
                    st.success("Fórmula aplicada. Revisa la nueva columna en la tabla.")
                    st.rerun()

        c1, c2 = st.columns(2)
        if c1.button("Guardar cambios en inversiones.xlsx", type="primary"):
            hojas_actualizadas = leer_todas_las_hojas_excel()
            hojas_actualizadas[hoja] = st.session_state[editor_key].copy()
            guardar_excel_completo_desde_hojas(hojas_actualizadas)
            st.success("Cambios guardados en el Excel de la app.")
            st.rerun()

        hojas_para_descargar = leer_todas_las_hojas_excel()
        hojas_para_descargar[hoja] = st.session_state[editor_key].copy()
        c2.download_button(
            "Descargar Excel con estos cambios",
            data=excel_hojas_a_bytes(hojas_para_descargar),
            file_name=ARCHIVO,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    with tab_usuarios:
        st.subheader("Usuarios y contraseñas")
        st.caption(
            "Por seguridad, las contraseñas se guardan hasheadas (bcrypt) — ni siquiera "
            "nosotros podemos ver la contraseña real de nadie. Los accesos nuevos y los "
            "reseteos se envían por email como contraseña temporal de un solo uso: el usuario "
            "está obligado a cambiarla en su primer login, y esa contraseña temporal no vuelve "
            "a mostrarse ni guardarse en ningún sitio."
        )

        try:
            _smtp_sender_u = st.secrets["email"]["sender"]
            _smtp_password_u = st.secrets["email"]["password"]
            _display_name_u = st.secrets["email"].get("display_name", "Chaparro Fernández Wealth")
            _email_configurado = True
        except Exception:
            _smtp_sender_u = _smtp_password_u = _display_name_u = ""
            _email_configurado = False

        df_u = _leer_hoja_usuarios()

        # ── Generar acceso nuevo ──
        st.markdown("#### ➕ Generar acceso nuevo")
        if not _email_configurado:
            st.warning("Configura el email de envío (Secrets → [email]) para poder generar accesos nuevos de forma segura.")
        else:
            with st.form("form_generar_acceso_nuevo"):
                c1, c2 = st.columns(2)
                tipo_nuevo = c1.selectbox("Tipo de acceso", ["inversor", "admin"], key="tipo_acceso_nuevo")
                usuario_nuevo = c2.text_input(
                    "Usuario", key="input_usuario_nuevo",
                    help="Para inversores, debe coincidir exactamente con el nombre en la columna 'inversor' del Excel.",
                )
                email_nuevo = st.text_input("Email de destino", key="input_email_nuevo")
                generar = st.form_submit_button("Generar y enviar acceso")
            if generar:
                usuario_nuevo_limpio = (usuario_nuevo or "").strip()
                email_nuevo_limpio = (email_nuevo or "").strip()
                ya_existe = (
                    not df_u.empty and {"usuario", "tipo_usuario"}.issubset(df_u.columns)
                    and not df_u[
                        (df_u["usuario"].astype(str).str.strip().str.lower() == usuario_nuevo_limpio.lower())
                        & (df_u["tipo_usuario"].astype(str).str.strip().str.lower() == tipo_nuevo)
                    ].empty
                )
                if not usuario_nuevo_limpio or not email_nuevo_limpio:
                    st.error("Rellena usuario y email.")
                elif "@" not in email_nuevo_limpio:
                    st.error("El email no parece válido.")
                elif ya_existe:
                    st.error(f"Ya existe un acceso '{tipo_nuevo}' para {usuario_nuevo_limpio}. Usa 'Resetear contraseña' más abajo.")
                else:
                    password_temp = _generar_password_temporal()
                    for col in ["usuario", "tipo_usuario", "password", "debe_cambiar_password", "email"]:
                        if col not in df_u.columns:
                            df_u[col] = "NO" if col == "debe_cambiar_password" else ""
                        df_u[col] = df_u[col].astype(object)
                    fila_nueva = pd.DataFrame([{
                        "usuario": usuario_nuevo_limpio, "tipo_usuario": tipo_nuevo,
                        "password": _hash_password(password_temp), "debe_cambiar_password": "SI",
                        "email": email_nuevo_limpio,
                    }])
                    df_u = pd.concat([df_u, fila_nueva], ignore_index=True)
                    exito, mensaje = _guardar_hoja_usuarios(df_u)
                    if exito:
                        env_ok, env_msg = enviar_email_credenciales_nuevas(
                            email_nuevo_limpio, usuario_nuevo_limpio, password_temp, tipo_nuevo,
                            _smtp_sender_u, _smtp_password_u, _display_name_u,
                        )
                        if env_ok:
                            st.success(f"✅ Acceso creado y contraseña temporal enviada a {email_nuevo_limpio}.")
                        else:
                            st.error(f"El acceso se creó, pero el email falló: {env_msg}. La contraseña temporal no se puede recuperar — resetéala de nuevo.")
                    else:
                        st.warning(f"⚠️ No se pudo guardar en Drive: {mensaje}")

        st.divider()

        if df_u.empty:
            st.info(
                "Todavía no hay nadie en la hoja USUARIOS (nadie ha cambiado su contraseña "
                "desde el portal) — esos usuarios siguen con la contraseña inicial definida en el código."
            )
        else:
            tabla = df_u[["usuario", "tipo_usuario"]].rename(columns={"usuario": "Usuario", "tipo_usuario": "Tipo"})
            tabla["Estado"] = [
                ("🕓 Pendiente de 1er cambio" if str(dcp).strip().upper() == "SI"
                 else ("🔒 Hasheada (segura)" if _es_hash_bcrypt(str(pw)) else "⚠️ Texto plano (pendiente de migrar)"))
                for pw, dcp in zip(df_u["password"], df_u.get("debe_cambiar_password", ["NO"] * len(df_u)))
            ]
            st.dataframe(tabla, use_container_width=True, hide_index=True)

            st.markdown("#### 🔁 Resetear una contraseña")
            usuario_reset = st.selectbox(
                "Usuario", options=df_u["usuario"].astype(str).tolist(), key="select_usuario_reset_pw"
            )
            tipo_reset = df_u.loc[df_u["usuario"].astype(str) == usuario_reset, "tipo_usuario"].iloc[0] if usuario_reset else ""

            tab_reset_email, tab_reset_manual = st.tabs(["📧 Enviar por email (recomendado)", "✍️ Manual (respaldo)"])
            with tab_reset_email:
                st.caption("Genera una contraseña temporal aleatoria, la envía por email y obliga a cambiarla en el próximo login. Nadie ve la contraseña en pantalla.")
                email_reset = st.text_input("Email de destino", key="input_email_reset")
                if st.button("Generar y enviar nueva contraseña", key="btn_resetear_pw_email") and usuario_reset:
                    if not _email_configurado:
                        st.error("Configura el email de envío (Secrets → [email]) antes de usar esta opción.")
                    elif not email_reset or "@" not in email_reset:
                        st.error("Introduce un email válido.")
                    else:
                        password_temp = _generar_password_temporal()
                        idx = df_u[
                            (df_u["usuario"].astype(str) == usuario_reset)
                            & (df_u["tipo_usuario"].astype(str) == tipo_reset)
                        ].index
                        if "debe_cambiar_password" not in df_u.columns:
                            df_u["debe_cambiar_password"] = "NO"
                        if "email" not in df_u.columns:
                            df_u["email"] = ""
                        df_u.loc[idx, "password"] = _hash_password(password_temp)
                        df_u.loc[idx, "debe_cambiar_password"] = "SI"
                        # Aprovechamos el reseteo para dejar registrado el email de contacto —
                        # así, de ahora en adelante, este usuario puede usar 'olvidé mi
                        # contraseña' por su cuenta sin depender de un admin.
                        df_u.loc[idx, "email"] = email_reset
                        exito, mensaje = _guardar_hoja_usuarios(df_u)
                        if exito:
                            env_ok, env_msg = enviar_email_credenciales_nuevas(
                                email_reset, usuario_reset, password_temp, tipo_reset,
                                _smtp_sender_u, _smtp_password_u, _display_name_u,
                            )
                            if env_ok:
                                st.success(f"✅ Contraseña temporal de {usuario_reset} enviada a {email_reset}.")
                            else:
                                st.error(f"Se reseteó, pero el email falló: {env_msg}. Vuelve a resetear cuando el email funcione.")
                        else:
                            st.warning(f"⚠️ Reseteada localmente. {mensaje}")
            with tab_reset_manual:
                st.caption("Úsalo solo si el envío por email no funciona. Tú eliges y ves la contraseña — comunícasela al usuario por un canal seguro (nunca email sin cifrar).")
                nueva_pw_reset = st.text_input(
                    "Nueva contraseña temporal", type="password", key="input_nueva_pw_reset",
                )
                if st.button("Resetear contraseña manualmente", key="btn_resetear_pw") and usuario_reset and nueva_pw_reset:
                    if len(nueva_pw_reset) < 6:
                        st.error("La nueva contraseña debe tener al menos 6 caracteres.")
                    else:
                        idx = df_u[
                            (df_u["usuario"].astype(str) == usuario_reset)
                            & (df_u["tipo_usuario"].astype(str) == tipo_reset)
                        ].index
                        if "debe_cambiar_password" not in df_u.columns:
                            df_u["debe_cambiar_password"] = "NO"
                        df_u.loc[idx, "password"] = _hash_password(nueva_pw_reset)
                        df_u.loc[idx, "debe_cambiar_password"] = "SI"
                        exito, mensaje = _guardar_hoja_usuarios(df_u)
                        if exito:
                            st.success(f"✅ Contraseña de {usuario_reset} reseteada. {mensaje}")
                        else:
                            st.warning(f"⚠️ Reseteada localmente. {mensaje}")

        if st.button("🔄 Recargar desde Drive", key="btn_recargar_usuarios"):
            st.cache_data.clear()
            st.rerun()

    with tab_descargar:
        st.subheader("Descargar copia actual")
        hojas = leer_todas_las_hojas_excel()
        if hojas:
            st.download_button(
                "Descargar inversiones.xlsx",
                data=excel_hojas_a_bytes(hojas),
                file_name=ARCHIVO,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        else:
            st.error("No se ha podido preparar la descarga.")

# =========================
# APP FINAL
# =========================
if __name__ == "__main__":  # carga inicial + hero: solo se ejecuta con `streamlit run`, no al importar
    print("[DIAG] 7. Arranca bloque carga inicial + hero", file=sys.stderr, flush=True)
    tag_sesion = st.session_state.usuario + (" (Portal de inversor)" if st.session_state.get("tipo_usuario") == "inversor" else "")
    mostrar_hero(tag_sesion)
    print("[DIAG] 8. Hero renderizado, empieza carga de datos", file=sys.stderr, flush=True)

    # Límite de tiempo DURO a la carga inicial de datos: si algo se queda colgado (red, Drive,
    # lo que sea), la app NUNCA se queda en blanco indefinidamente — pasados 15s se corta y se
    # muestra un aviso claro con un botón para reintentar, en vez de una pantalla vacía.
    # IMPORTANTE: el executor se crea SIN 'with' — un 'with' espera (bloquea) a que el hilo de
    # fondo termine de verdad antes de dejar continuar el código, aunque ya hayamos detectado el
    # timeout, lo que anulaba por completo el límite de 15s. Así, en cuanto salta el timeout,
    # seguimos adelante inmediatamente y dejamos el hilo colgado morir solo en segundo plano.
    import concurrent.futures as _cf
    _executor_carga = _cf.ThreadPoolExecutor(max_workers=1)
    try:
        with st.spinner("Cargando datos del fondo..."):
            print("[DIAG] 9. Lanzando cargar_excel_completo en hilo aparte", file=sys.stderr, flush=True)
            _future = _executor_carga.submit(cargar_excel_completo)
            df_inv, df_cal, df_control = _future.result(timeout=15)
        print("[DIAG] 10. cargar_excel_completo TERMINÓ BIEN", file=sys.stderr, flush=True)
        _executor_carga.shutdown(wait=False)
    except _cf.TimeoutError:
        print("[DIAG] 10b. TIMEOUT de 15s disparado correctamente", file=sys.stderr, flush=True)
        _executor_carga.shutdown(wait=False)
        st.error(
            "⏱️ La carga de datos está tardando demasiado (más de 15 segundos) y se ha cortado "
            "para no dejarte con la pantalla en blanco. Puede ser un problema temporal de conexión "
            "con Google Drive. Pulsa el botón de abajo para reintentar."
        )
        if st.button("🔄 Reintentar"):
            st.cache_data.clear()
            st.rerun()
        st.stop()
    except Exception as e:
        _executor_carga.shutdown(wait=False)
        st.error("No se ha podido cargar inversiones.xlsx. Revisa que el archivo esté subido a GitHub y que las hojas existan.")
        with st.expander("Ver detalle técnico"):
            st.exception(e)
        st.stop()


def _cobro_notas_jordi_mes(df_inv, df_cal, anio, mes):
    """
    Cobro bruto de la compañía por notas con cuenta_cobro=JORDI en el mes.
    Solo cuentan las NOTAS 1 A 8 (las que se invirtieron con capital personal de Jordi) que
    estén ACTIVAS en ese mes: ni canceladas (tipo_operacion != cancelada) ni ya llamadas
    (motivo call/call final con fecha_final_inversion <= fin de este mes), y que ya existieran
    a esa fecha (fecha_inversion <= fin de este mes). Una nota llamada deja de reducir la deuda
    a partir del mes de su call, aunque haya tenido pagos mientras estuvo activa en meses
    anteriores — por eso el corte se calcula mes a mes, no con 'hoy'.

    IMPORTANTE: toda nota activa ese mes se incluye SIEMPRE en por_nota, aunque ese mes en
    concreto no le tocara pago (p.ej. una nota trimestral en un mes que no es de cobro) —
    en ese caso sale con importe 0, en vez de desaparecer de la lista sin más.

    cobro = capital * interes_nota_anual / 12 * periodicidad, solo si hay PAGO ese mes en el calendario.
    Devuelve (total, {nombre_nota: importe})
    """
    fecha_corte_mes = pd.Timestamp(anio, mes, 1) + pd.offsets.MonthEnd(0)

    cal_tmp = df_cal.copy()
    cal_tmp["_fecha"] = pd.to_datetime(cal_tmp["fecha"], dayfirst=True, errors="coerce")
    cal_tmp["_nota"] = pd.to_numeric(cal_tmp["nota"], errors="coerce")
    pagos_mes = cal_tmp[
        (cal_tmp["tipo_evento"].astype(str).str.upper() == "PAGO") &
        (cal_tmp["_fecha"].dt.year == anio) &
        (cal_tmp["_fecha"].dt.month == mes) &
        (cal_tmp["_nota"] >= 1) & (cal_tmp["_nota"] <= 8)
    ]
    notas_con_pago_este_mes = set(pagos_mes["_nota"].dropna().astype(int).unique())

    # Notas ya llamadas (call o call final) a fecha de cierre de este mes: dejan de contar.
    notas_llamadas = set()
    if "nombre_activo" in df_inv.columns and "motivo" in df_inv.columns:
        motivo_normalizado = df_inv["motivo"].astype(str).str.lower().str.strip()
        llamadas_df = df_inv[motivo_normalizado.isin(["call", "call final"])].copy()
        llamadas_df["fecha_final_inversion"] = pd.to_datetime(llamadas_df["fecha_final_inversion"], errors="coerce")
        llamadas_df["nota_num"] = llamadas_df["nombre_activo"].apply(extraer_numero_nota)
        for nota_num, grupo in llamadas_df.groupby("nota_num"):
            if pd.notna(nota_num) and (grupo["fecha_final_inversion"].notna() & (grupo["fecha_final_inversion"] <= fecha_corte_mes)).all():
                notas_llamadas.add(int(nota_num))

    notas_jordi = df_inv[
        (df_inv["cuenta_cobro"].astype(str).str.strip().str.upper() == "JORDI") &
        (df_inv["activo_generador_interes"].astype(str).str.upper() == "SI") &
        (df_inv["tipo_operacion"].astype(str).str.lower() != "cancelada")
    ].copy()
    notas_jordi["_nota_num"] = pd.to_numeric(notas_jordi["nombre_activo"].apply(extraer_numero_nota), errors="coerce")
    notas_jordi["_fecha_inversion"] = pd.to_datetime(notas_jordi.get("fecha_inversion"), errors="coerce")
    notas_jordi = notas_jordi[
        notas_jordi["_nota_num"].notna() & (notas_jordi["_nota_num"] >= 1) & (notas_jordi["_nota_num"] <= 8) &
        (~notas_jordi["_nota_num"].astype(int).isin(notas_llamadas)) &
        (notas_jordi["_fecha_inversion"].isna() | (notas_jordi["_fecha_inversion"] <= fecha_corte_mes))
    ]

    total = 0.0
    por_nota = {}
    for nota_num in sorted(notas_jordi["_nota_num"].astype(int).unique()):
        nombre = f"NOTA_{nota_num:02d}"
        activas = notas_jordi[notas_jordi["nombre_activo"].astype(str).str.upper() == nombre]
        if activas.empty:
            continue
        sub = 0.0
        if nota_num in notas_con_pago_este_mes:
            for _, r in activas.iterrows():
                cap = float(r.get("capital_invertido", 0))
                tasa_nota = float(r.get("interes_nota_anual", 0))
                periodicidad = int(r.get("periodicidad_meses", 1) or 1)
                cobro = cap * tasa_nota / 12 * periodicidad
                sub += cobro
        # Se incluye SIEMPRE (aunque sub sea 0) para que se vea que la nota está activa
        # pero este mes en concreto no le tocaba cobro.
        por_nota[nombre] = sub
        total += sub
    return total, por_nota


HOJA_REPARTO_DIVIDENDOS = "REPARTO_DIVIDENDOS"


def cargar_reparto_dividendos() -> pd.DataFrame:
    """Lee la hoja REPARTO_DIVIDENDOS del Excel.
    Columnas esperadas: fecha, importe, descripcion.
    Cada fila es un reparto de dividendos que AUMENTA la deuda con Jordi,
    igual que los intereses de JEP.
    """
    try:
        df = pd.read_excel(ARCHIVO, sheet_name=HOJA_REPARTO_DIVIDENDOS)
        df.columns = [str(c).strip().lower() for c in df.columns]
        df["fecha"] = pd.to_datetime(df["fecha"], dayfirst=True, errors="coerce")
        df["importe"] = pd.to_numeric(df["importe"], errors="coerce").fillna(0)
        if "descripcion" not in df.columns:
            df["descripcion"] = ""
        return df.dropna(subset=["fecha"])
    except Exception:
        return pd.DataFrame(columns=["fecha", "importe", "descripcion"])


def cargar_transferencias_jordi() -> pd.DataFrame:
    """Lee la hoja TRANSFERENCIAS_JORDI del Excel.
    Columnas esperadas: fecha, importe, descripcion.
    Cada fila es una transferencia de Jordi HACIA el fondo (capital que aporta) —
    AUMENTA la deuda con Jordi, con la misma lógica que los intereses de JEP y el
    reparto de dividendos: es dinero que la empresa le debe devolver. Si algún día
    hay una transferencia en sentido contrario (el fondo le devuelve a Jordi), basta
    con cargar el importe en negativo en esa fila y esta misma función ya lo resta bien.
    """
    try:
        df = pd.read_excel(ARCHIVO, sheet_name="TRANSFERENCIAS_JORDI")
        df.columns = [str(c).strip().lower() for c in df.columns]
        df["fecha"] = pd.to_datetime(df["fecha"], dayfirst=True, errors="coerce")
        df["importe"] = pd.to_numeric(df["importe"], errors="coerce").fillna(0)
        if "descripcion" not in df.columns:
            df["descripcion"] = ""
        return df.dropna(subset=["fecha"])
    except Exception:
        return pd.DataFrame(columns=["fecha", "importe", "descripcion"])


def _reparto_dividendos_mes(df_rep, anio, mes):
    """
    Total de reparto de dividendos del mes (aumenta la deuda con Jordi).
    Devuelve (total, {descripcion: importe})
    """
    if df_rep.empty:
        return 0.0, {}
    mes_rows = df_rep[
        (df_rep["fecha"].dt.year == anio) &
        (df_rep["fecha"].dt.month == mes)
    ]
    if mes_rows.empty:
        return 0.0, {}
    total = float(mes_rows["importe"].sum())
    por_item = {}
    for _, r in mes_rows.iterrows():
        desc = str(r.get("descripcion", "") or "").strip() or "Reparto de dividendos"
        por_item[desc] = por_item.get(desc, 0) + float(r["importe"])
    return total, por_item


def _transferencias_jordi_mes(df_transf, anio, mes):
    """
    Total de transferencias de Jordi hacia el fondo del mes (aumenta la deuda con Jordi,
    misma lógica que JEP y reparto de dividendos). Devuelve (total, {descripcion: importe})
    """
    if df_transf.empty:
        return 0.0, {}
    mes_rows = df_transf[
        (df_transf["fecha"].dt.year == anio) &
        (df_transf["fecha"].dt.month == mes)
    ]
    if mes_rows.empty:
        return 0.0, {}
    total = float(mes_rows["importe"].sum())
    por_item = {}
    for _, r in mes_rows.iterrows():
        desc = str(r.get("descripcion", "") or "").strip() or "Transferencia de Jordi"
        por_item[desc] = por_item.get(desc, 0) + float(r["importe"])
    return total, por_item


def _intereses_jep_mes(df_inv, anio, mes):
    """
    Intereses devengados a JEP en el mes, agrupados por tipo de activo.
    Lógica pro-rata días igual que el extracto.
    Devuelve (total, {grupo: importe})
    """
    dias_mes = ultimo_dia_mes(anio, mes)
    inicio_mes = pd.Timestamp(anio, mes, 1)
    fin_mes = pd.Timestamp(anio, mes, dias_mes)
    jep = df_inv[df_inv["inversor"].astype(str).str.upper().str.strip() == "JEP"].copy()
    jep["_fi"] = parsear_fecha_robusta(jep["fecha_inversion"])
    jep["_ff"] = parsear_fecha_robusta(jep["fecha_final_inversion"])
    es_nota = jep["tipo_inversion"].astype(str).str.lower().str.contains("nota")
    es_cancelada = jep["tipo_operacion"].astype(str).str.lower() == "cancelada"
    jep = jep[es_nota | (~es_nota & ~es_cancelada)].copy()
    total = 0.0
    por_grupo = {}
    for _, r in jep.iterrows():
        fi, ff = r["_fi"], r["_ff"]
        if pd.isna(fi) or fi > fin_mes: continue
        if pd.notna(ff) and ff < inicio_mes: continue
        ir = max(fi, inicio_mes)
        fr = fin_mes if pd.isna(ff) else min(ff, fin_mes)
        if ir > fr: continue
        dias = (fr - ir).days + 1
        cap = float(r.get("capital_invertido", 0))
        tasa_inv = float(r.get("interes_inversor_anual", 0))
        interes = cap * tasa_inv / 12 * (dias / dias_mes)
        total += interes
        nombre_activo = str(r.get("nombre_activo", "")).upper()
        if "NOTA" in nombre_activo:
            grupo = "Notas"
        else:
            grupo = str(r.get("subtipo_inversion", r.get("nombre_activo", ""))).title()
        por_grupo[grupo] = por_grupo.get(grupo, 0) + interes
    return total, por_grupo


def calcular_deuda_jordi(df_inv, df_cal, df_control, capital_inicial: float, fecha_inicio: pd.Timestamp) -> pd.DataFrame:
    """
    Evolución mes a mes de la deuda con Jordi Chaparro.
    RESTA: beneficio fijos (Paraguay/Motoclick/Futbol/Bolivia) + beneficio notas cuenta JORDI
    SUMA:  intereses devengados a JEP (todos sus activos, pro-rata días)
           + reparto de dividendos (hoja REPARTO_DIVIDENDOS)
           + transferencias de Jordi hacia el fondo (hoja TRANSFERENCIAS_JORDI)
    """
    hoy = pd.Timestamp.today().normalize()
    filas = []
    saldo = capital_inicial
    anio = fecha_inicio.year
    mes = fecha_inicio.month
    df_reparto = cargar_reparto_dividendos()
    df_transferencias = cargar_transferencias_jordi()

    ACTIVOS_FIJOS = [
        ("Paraguay",  "paraguay",  TASA_ANUAL_PARAGUAY),
        ("MotoClick", "motoclick", TASA_ANUAL_MOTOCLICK),
        ("Fútbol",    "futbol",    TASA_ANUAL_FUTBOL),
        ("Bolivia",   "bolivia",   TASA_ANUAL_BOLIVIA),
    ]

    while (anio < hoy.year) or (anio == hoy.year and mes <= hoy.month):
        label = f"{mes:02d}/{anio}"

        # 1) Beneficio por activo fijo
        df_mov_mc = cargar_movimientos_motoclick()
        detalle_fijos = []
        ingreso_fijos_total = 0.0
        for nombre_act, clave, tasa in ACTIVOS_FIJOS:
            det = detalle_activo_mes(df_inv, clave, tasa, anio, mes)
            if not det.empty:
                det["activo"] = clave
                # Para MotoClick usar ingreso real ajustado por movimientos
                if clave == "motoclick":
                    det_aj = ajustar_ingreso_motoclick(det, df_inv, anio, mes)
                    cobro = float(det_aj["ingreso_bruto"].sum())
                else:
                    cobro = float(det["ingreso_bruto"].sum())
                ingreso_fijos_total += cobro
                if cobro > 0:
                    detalle_fijos.append({"activo": nombre_act, "cobro": cobro})

        # 2) Cobro compañía notas cuenta JORDI (cobro bruto = capital * tasa_nota / 12)
        cobro_notas_jordi, por_nota_jordi = _cobro_notas_jordi_mes(df_inv, df_cal, anio, mes)

        # 3) Intereses devengados a JEP (pro-rata días, agrupados por tipo)
        pago_jep, jep_por_grupo = _intereses_jep_mes(df_inv, anio, mes)

        # 4) Reparto de dividendos del mes (aumenta la deuda, igual que JEP)
        reparto_dividendos, reparto_por_item = _reparto_dividendos_mes(df_reparto, anio, mes)

        # 5) Transferencias de Jordi hacia el fondo del mes (aumenta la deuda, igual que JEP)
        transferencias_jordi, transferencias_por_item = _transferencias_jordi_mes(df_transferencias, anio, mes)

        # Total ingresos compañía = ingresos fijos + cobro notas JORDI
        resta = ingreso_fijos_total + cobro_notas_jordi
        suma = pago_jep + reparto_dividendos + transferencias_jordi
        saldo_inicio = saldo
        saldo = saldo + suma - resta

        filas.append({
            "mes": label,
            "saldo_inicio": saldo_inicio,
            "detalle_fijos": detalle_fijos,          # lista con cobro por activo
            "ingreso_fijos": ingreso_fijos_total,
            "por_nota_jordi": por_nota_jordi,         # dict nota->cobro
            "cobro_notas_jordi": cobro_notas_jordi,
            "jep_por_grupo": jep_por_grupo,           # dict grupo->interes
            "pago_jep": pago_jep,
            "reparto_por_item": reparto_por_item,     # dict descripcion->importe
            "reparto_dividendos": reparto_dividendos,
            "transferencias_por_item": transferencias_por_item,  # dict descripcion->importe
            "transferencias_jordi": transferencias_jordi,
            "total_resta": resta,
            "total_suma": suma,
            "variacion_neta": suma - resta,
            "saldo_fin": saldo,
        })

        mes += 1
        if mes == 13:
            mes = 1
            anio += 1

    return pd.DataFrame(filas)


def _posiciones_activas_para_cerrar(df_inv: pd.DataFrame) -> pd.DataFrame:
    """Posiciones NUEVA/REINVERSION sin fecha_final_inversion — candidatas a cerrar (CANCELADA)."""
    if df_inv is None or df_inv.empty:
        return pd.DataFrame()
    d = df_inv.copy()
    if "tipo_operacion" in d.columns:
        d = d[d["tipo_operacion"].astype(str).str.strip().str.upper().isin(["NUEVA", "REINVERSION"])]
    if "fecha_final_inversion" in d.columns:
        d = d[d["fecha_final_inversion"].isna()]
    return d


def _siguiente_id_inversion(df_inv: pd.DataFrame, offset: int = 0) -> str:
    """Sugiere el próximo id_inversion siguiendo la convención OPxxx (ej. si el último usado es
    OP037, sugiere OP038). Busca el número más alto ya usado en toda la hoja INVERSIONES, con
    independencia de cuántos ceros de relleno tenga cada fila existente.
    'offset' permite pedir el 2º, 3º... siguiente número — para cuando se dan de alta varios
    inversores a la vez en la misma operación, cada uno necesita su propio id_inversion distinto,
    y no se pueden sugerir dos veces el mismo número dentro del mismo formulario."""
    try:
        numeros = []
        if df_inv is not None and not df_inv.empty and "id_inversion" in df_inv.columns:
            for valor in df_inv["id_inversion"].dropna().astype(str):
                m = re.match(r"^\s*OP0*(\d+)\s*$", valor.strip(), re.IGNORECASE)
                if m:
                    numeros.append(int(m.group(1)))
        siguiente = (max(numeros) + 1 + offset) if numeros else (1 + offset)
        return f"OP{siguiente:03d}"
    except Exception:
        return ""


def _sugerir_fecha_inicio_nota(numero_nota: int, df_cal: pd.DataFrame):
    """Sugiere la fecha_inversion de una nota: primer PAGO del calendario, menos 1 periodo
    completo (1 mes si el cobro es mensual, 3 meses si es trimestral, etc. — no siempre 1 mes fijo).
    Busca primero en CALENDARIO_NOTAS ya guardado; si la nota todavía no está guardada (se está dando
    de alta a la vez que su inversión), busca en el borrador de extracción de IA de esa misma nota
    en '➕ Nueva inversión' (en sesión o persistido en Drive).
    Devuelve (fecha_sugerida: pd.Timestamp|None, periodicidad_meses: int, fuente: str)."""
    fechas_pago = []
    fuente = ""

    if df_cal is not None and not df_cal.empty and "nota" in df_cal.columns and "tipo_evento" in df_cal.columns:
        cal_n = df_cal[
            (pd.to_numeric(df_cal["nota"], errors="coerce") == numero_nota)
            & (df_cal["tipo_evento"].astype(str).str.upper() == "PAGO")
        ]
        if not cal_n.empty:
            fechas_pago = sorted(pd.to_datetime(cal_n["fecha"], errors="coerce").dropna().tolist())
            if fechas_pago:
                fuente = "calendario ya guardado (CALENDARIO_NOTAS)"

    if not fechas_pago:
        try:
            almacen_notas = st.session_state.get("notas_wizard_datos", {})
            extraido = almacen_notas.get(numero_nota)
            if extraido and extraido.get("calendario"):
                fechas_pago = sorted(pd.to_datetime(
                    [e.get("pago") for e in extraido["calendario"] if e.get("pago") and str(e.get("pago")).strip().upper() != "REVISAR"],
                    errors="coerce",
                ).dropna().tolist())
                if fechas_pago:
                    fuente = "borrador de extracción de IA de la nota (➕ Nueva inversión)"
        except Exception:
            pass

    if not fechas_pago:
        return None, 1, ""

    primer_pago = fechas_pago[0]
    periodicidad_meses = 1
    if len(fechas_pago) >= 2:
        delta_dias = (fechas_pago[1] - fechas_pago[0]).days
        delta_meses_aprox = max(1, round(delta_dias / 30.44))
        periodicidad_meses = min([1, 3, 6], key=lambda p: abs(p - delta_meses_aprox))

    fecha_sugerida = pd.Timestamp(primer_pago) - pd.DateOffset(months=periodicidad_meses)
    return pd.Timestamp(fecha_sugerida), periodicidad_meses, fuente


def seccion_nueva_inversion(df_inv: pd.DataFrame, df_cal: pd.DataFrame, df_control: pd.DataFrame):
    """Alta/cierre/reinversión de posiciones directamente desde la app, escribiendo en Google Drive
    con el mismo patrón que el wizard de 'Notas estructuradas': primero se guarda como BORRADOR
    revisable (hoja BORRADORES_INVERSIONES), y solo se escribe de verdad en INVERSIONES (y en
    REINVERSIONES si aplica) cuando Yuri confirma explícitamente tras revisar la previsualización.
    Solo visible/accesible para el usuario 'Yuri' (se controla también en el menú principal).

    Soporta más de un inversor en la misma inversión (ej. dos inversores en la misma nota): los
    datos del activo/nota (tipo, nombre, fecha, cupón...) se rellenan una sola vez y se comparten;
    cada inversor tiene su propia fila con su capital, su tasa y su id_inversion."""
    st.markdown("## ➕ Nueva inversión")
    st.caption(
        "Da de alta una inversión nueva (nota estructurada O cualquier activo fijo — Paraguay, MotoClick, "
        "Fútbol, Bolivia, Bitcoin — se elige en 'Tipo de inversión' más abajo), cierra una posición existente "
        "(call/vencimiento/retiro/traspaso) o registra una reinversión — sin tocar el Excel a mano. Funciona "
        "igual que el wizard de notas: se guarda como borrador que puedes revisar/retomar, y solo se escribe "
        "en INVERSIONES cuando confirmas explícitamente."
    )

    # ── Referencia del borrador (clave libre para poder guardarlo/retomarlo) ──────────────
    # Tiene un valor por defecto autogenerado para que el formulario aparezca de inmediato sin
    # tener que escribir nada — solo hace falta cambiarla si quieres ponerle un nombre concreto
    # para retomarla más tarde.
    if "inversion_wizard_datos" not in st.session_state:
        st.session_state["inversion_wizard_datos"] = {}
    if "inversion_wizard_clave_actual" not in st.session_state:
        st.session_state["inversion_wizard_clave_actual"] = f"borrador-{pd.Timestamp.now().strftime('%Y%m%d-%H%M%S')}"

    claves_existentes = listar_claves_borradores_inversion()
    col_clave1, col_clave2 = st.columns([2, 1])
    with col_clave1:
        clave = st.text_input(
            "Referencia de esta operación (opcional cambiarla — solo sirve para guardar/retomar el borrador)",
            value=st.session_state["inversion_wizard_clave_actual"],
            key="inversion_wizard_clave_input",
        )
    with col_clave2:
        st.write("")
        st.write("")
        if claves_existentes:
            clave_retomar = st.selectbox("...o retomar borrador guardado", ["—"] + claves_existentes, key="inversion_wizard_retomar_sel", label_visibility="visible")
            if clave_retomar != "—" and st.button("📂 Cargar borrador"):
                borrador = cargar_borrador_inversion(clave_retomar)
                if borrador:
                    st.session_state["inversion_wizard_clave_actual"] = clave_retomar
                    st.session_state["inversion_wizard_datos"] = borrador
                    st.rerun()
                else:
                    st.error("No se pudo cargar ese borrador.")

    clave = clave.strip() or st.session_state["inversion_wizard_clave_actual"]
    st.session_state["inversion_wizard_clave_actual"] = clave
    datos_previos = st.session_state["inversion_wizard_datos"] if st.session_state.get("inversion_wizard_clave_actual") == clave else {}

    st.markdown("---")

    tipo_operacion = st.selectbox(
        "Tipo de operación",
        ["NUEVA", "CANCELADA", "REINVERSION"],
        index=["NUEVA", "CANCELADA", "REINVERSION"].index(datos_previos.get("tipo_operacion", "NUEVA")) if isinstance(datos_previos.get("tipo_operacion"), str) and datos_previos.get("tipo_operacion") in ["NUEVA", "CANCELADA", "REINVERSION"] else 0,
        help=(
            "NUEVA: entra capital fresco de uno o varios inversores a un activo/nota.\n\n"
            "CANCELADA: se cierra una posición existente (call, vencimiento, retiro completo, traspaso a otro inversor).\n\n"
            "REINVERSION: el capital de una posición cerrada se recoloca en otro activo/nota. OJO — según la "
            "lógica actual del fondo, el pago mensual al inversor sigue calculándose sobre la fila NUEVA "
            "original (misma fecha y misma tasa, sin reiniciar el reloj de intereses); esta fila de "
            "REINVERSION solo redirige de qué activo cobra la EMPRESA a partir de ahora. Si tienes dudas en "
            "un caso concreto, revísalo con calma antes de confirmar — es el tipo de operación más delicado."
        ),
    )

    # ═══════════════════════════════════════════════════════════════════
    # RAMA: CANCELADA — cerrar una posición ya existente (o TODAS las de un
    # mismo activo/nota de golpe — útil para un call, que cierra a la vez
    # las posiciones de todos los inversores de esa nota)
    # ═══════════════════════════════════════════════════════════════════
    if tipo_operacion == "CANCELADA":
        activas = _posiciones_activas_para_cerrar(df_inv)
        if activas.empty:
            st.warning("No se han encontrado posiciones activas (NUEVA/REINVERSION sin fecha_final_inversion) para cerrar.")
            return

        modo_cierre = st.radio(
            "¿Qué quieres cerrar?",
            ["Una posición en concreto", "TODAS las posiciones activas de un mismo activo/nota (útil para un call)"],
            key="cancelada_modo",
            help="Cuando llaman a una nota (call), normalmente hay que cerrar de golpe la posición de TODOS los inversores que estaban en esa nota, todos con la misma fecha y motivo — para eso usa la segunda opción.",
        )

        if modo_cierre == "Una posición en concreto":
            activas_et = activas.copy()
            activas_et["_etiqueta"] = activas_et.apply(
                lambda r: f"{r.get('id_inversion','?')} | {r.get('inversor','?')} | {r.get('nombre_activo','?')} | ${float(r.get('capital_invertido',0) or 0):,.2f} | desde {pd.Timestamp(r['fecha_inversion']).strftime('%d/%m/%Y') if pd.notna(r.get('fecha_inversion')) else '?'}",
                axis=1,
            )
            etiqueta_sel = st.selectbox("Posición a cerrar", activas_et["_etiqueta"].tolist(), key="cancelada_posicion_sel")
            filas_a_cerrar = activas_et[activas_et["_etiqueta"] == etiqueta_sel]
            st.dataframe(filas_a_cerrar.drop(columns=["_etiqueta"]), use_container_width=True, hide_index=True)
        else:
            nombres_activo_disp = sorted(activas["nombre_activo"].dropna().astype(str).unique()) if "nombre_activo" in activas.columns else []
            nombre_activo_cerrar = st.selectbox("Activo/nota a cerrar por completo", nombres_activo_disp, key="cancelada_activo_sel")
            filas_a_cerrar = activas[activas["nombre_activo"].astype(str) == nombre_activo_cerrar]
            st.caption(f"Se cerrarán de golpe estas {len(filas_a_cerrar)} posición(es) — una por cada inversor activo en {nombre_activo_cerrar}:")
            st.dataframe(filas_a_cerrar, use_container_width=True, hide_index=True)

        ids_a_cerrar = filas_a_cerrar["id_inversion"].astype(str).tolist()

        fecha_final = st.date_input("Fecha de cierre (fecha_final_inversion)", value=pd.Timestamp.today().date(), key="cancelada_fecha_final")
        motivos_conocidos = _valores_conocidos_columna(df_inv, "motivo") or ["call", "call final", "vencimiento", "retiro", "traspaso"]
        default_motivo_idx = motivos_conocidos.index("call") if "call" in motivos_conocidos else 0
        motivo_sel = st.selectbox("Motivo del cierre", motivos_conocidos + ["Otro (escribir)"], index=default_motivo_idx, key="cancelada_motivo_sel")
        motivo_final = st.text_input("Escribe el motivo", key="cancelada_motivo_libre") if motivo_sel == "Otro (escribir)" else motivo_sel

        # Si es una nota y hay una hoja CALENDARIO_CALLS con fechas de posible call para esa
        # nota, ofrece marcar ese call como ejecutado a la vez que se cierran las posiciones —
        # para no dejar el calendario de calls desactualizado.
        marcar_call_calendario = False
        etiqueta_call_calendario = None
        df_calls_cancelada = leer_hoja_excel("CALENDARIO_CALLS") if motivo_final.strip().lower().startswith("call") else pd.DataFrame()
        if not df_calls_cancelada.empty and "nota" in df_calls_cancelada.columns:
            numero_nota_cerrar = extraer_numero_nota(nombre_activo_cerrar) if modo_cierre != "Una posición en concreto" else extraer_numero_nota(str(filas_a_cerrar.iloc[0].get("nombre_activo", "")))
            if pd.notna(numero_nota_cerrar):
                calls_nota_cancelada = df_calls_cancelada[pd.to_numeric(df_calls_cancelada["nota"], errors="coerce") == numero_nota_cerrar]
                if not calls_nota_cancelada.empty:
                    calls_nota_cancelada = calls_nota_cancelada.copy()
                    calls_nota_cancelada["_etiqueta"] = calls_nota_cancelada.apply(
                        lambda r: f"{pd.Timestamp(r.get('fecha_call')).strftime('%d/%m/%Y') if pd.notna(r.get('fecha_call')) else '?'} | estado actual: {r.get('estado','?')}",
                        axis=1,
                    )
                    etiqueta_call_calendario = st.selectbox(
                        "Marcar también este call como ejecutado en CALENDARIO_CALLS (opcional)",
                        ["(no marcar ninguno)"] + calls_nota_cancelada["_etiqueta"].tolist(),
                        key="cancelada_call_calendario_sel",
                    )
                    if etiqueta_call_calendario != "(no marcar ninguno)":
                        marcar_call_calendario = True

        datos_borrador = {
            "tipo_operacion": "CANCELADA",
            "ids_a_cerrar": ids_a_cerrar,
            "fecha_final_inversion": str(fecha_final),
            "motivo": motivo_final,
        }

        col_a, col_g = st.columns(2)
        with col_a:
            if st.button("📌 Guardar avance del borrador"):
                guardar_borrador_inversion(clave, datos_borrador)
                st.session_state["inversion_wizard_datos"] = datos_borrador

        with col_g:
            if st.button(f"💾 Confirmar cierre de {len(ids_a_cerrar)} posición(es) en el Excel", type="primary"):
                hojas = leer_todas_las_hojas_excel()
                if not hojas or "INVERSIONES" not in hojas:
                    st.error("No se pudo leer INVERSIONES para guardar el cierre.")
                else:
                    df_raw = hojas["INVERSIONES"]
                    mapa = _mapa_columnas_reales(hojas, "INVERSIONES")
                    col_id = mapa.get("id_inversion")
                    col_ff = mapa.get("fecha_final_inversion")
                    col_mot = mapa.get("motivo")
                    if not col_id or not col_ff:
                        st.error("No se encontraron las columnas id_inversion / fecha_final_inversion en la hoja INVERSIONES — revisa los nombres de columna a mano.")
                    else:
                        mascara = df_raw[col_id].astype(str).isin(ids_a_cerrar)
                        if not mascara.any():
                            st.error("No se encontró ninguna de esas filas en la hoja real — puede que el Excel haya cambiado desde que se cargó esta pantalla. Recarga y vuelve a intentarlo.")
                        else:
                            df_raw.loc[mascara, col_ff] = pd.Timestamp(fecha_final)
                            if col_mot:
                                df_raw.loc[mascara, col_mot] = motivo_final
                            hojas["INVERSIONES"] = df_raw

                            if marcar_call_calendario and "CALENDARIO_CALLS" in hojas:
                                mapa_calls = _mapa_columnas_reales(hojas, "CALENDARIO_CALLS")
                                col_nota_calls = mapa_calls.get("nota")
                                col_fecha_call = mapa_calls.get("fecha_call")
                                col_estado_calls = mapa_calls.get("estado")
                                if col_nota_calls and col_fecha_call and col_estado_calls:
                                    fecha_call_marcar = pd.to_datetime(etiqueta_call_calendario.split(" | ")[0], dayfirst=True, errors="coerce")
                                    df_calls_raw = hojas["CALENDARIO_CALLS"]
                                    estados_conocidos = _valores_conocidos_columna(df_calls_raw.rename(columns={col_estado_calls: "estado"}), "estado")
                                    estado_ejecutado = next((e for e in estados_conocidos if "ejecut" in e.lower() or "call" in e.lower() or "confirm" in e.lower()), (estados_conocidos[0] if estados_conocidos else "EJECUTADO"))
                                    mascara_call = (
                                        (pd.to_numeric(df_calls_raw[col_nota_calls], errors="coerce") == numero_nota_cerrar)
                                        & (pd.to_datetime(df_calls_raw[col_fecha_call], errors="coerce", dayfirst=True) == fecha_call_marcar)
                                    )
                                    if mascara_call.any():
                                        df_calls_raw.loc[mascara_call, col_estado_calls] = estado_ejecutado
                                        hojas["CALENDARIO_CALLS"] = df_calls_raw
                                        st.info(f"CALENDARIO_CALLS actualizado: estado = '{estado_ejecutado}' para la fecha {fecha_call_marcar.strftime('%d/%m/%Y')}.")

                            guardar_excel_completo_desde_hojas(hojas)
                            borrar_borrador_inversion(clave)
                            st.session_state["inversion_wizard_datos"] = {}
                            st.session_state["inversion_wizard_clave_actual"] = ""
                            st.success(f"{len(ids_a_cerrar)} posición(es) cerrada(s) el {fecha_final} (motivo: {motivo_final}): {', '.join(ids_a_cerrar)}.")
                            st.cache_data.clear()
        return

    # ═══════════════════════════════════════════════════════════════════
    # RAMA: NUEVA / REINVERSION — datos del activo/nota (compartidos) +
    # uno o varios inversores (cada uno con su propia fila)
    # ═══════════════════════════════════════════════════════════════════
    st.markdown("### 1. Datos del activo/nota (se comparten entre todos los inversores de esta operación)")

    col2, col3 = st.columns(2)
    with col2:
        tipo_inversion_sel = st.selectbox(
            "Tipo de inversión", ["nota", "paraguay", "motoclick", "futbol", "bolivia", "bitcoin", "otro"],
            key="ni_tipo_inv_sel",
            help="Aquí eliges si es una nota estructurada o un activo fijo (Paraguay, MotoClick, Fútbol, Bolivia, Bitcoin) — o cualquier otro tipo nuevo con 'otro'.",
        )
        tipo_inversion_final = st.text_input("Escribe el tipo de inversión", key="ni_tipo_inv_libre") if tipo_inversion_sel == "otro" else tipo_inversion_sel
    es_nota = tipo_inversion_final.strip().lower() == "nota"

    reinversion_en_nota_nueva = False
    if tipo_operacion == "REINVERSION" and es_nota:
        reinversion_en_nota_nueva = st.radio(
            "¿La nota de destino de esta reinversión ya existe en el sistema, o es una nota nueva (con PDF por subir)?",
            ["Nota que ya existe (número ya cargado en CONTROL_NOTAS)", "Nota nueva — subir PDF"],
            key="ni_reinv_nota_existente_o_nueva",
        ) == "Nota nueva — subir PDF"

    if (tipo_operacion == "NUEVA" and es_nota) or reinversion_en_nota_nueva:
        if tipo_operacion == "NUEVA":
            st.info("📄 Nota nueva: sube el PDF y se rellena todo junto — datos de la nota **e** inversor(es), en un único guardado.")
        else:
            st.info("📄 Reinversión en nota nueva: sube el PDF de la nota igual que en una alta normal — se rellena CONTROL_NOTAS/CALENDARIO_NOTAS y el/los inversor(es), marcando la operación como REINVERSION y vinculando la(s) posición(es) de origen.")
        df_calls_para_nota = _leer_calendario_calls_cached()
        _tab_añadir_nota_nueva(df_control, df_cal, df_calls_para_nota, tipo_operacion=tipo_operacion, df_inv=df_inv)
        return

    fecha_sugerida_nota, periodicidad_sugerida_nota, fuente_sugerencia_nota = None, 1, ""

    with col3:
        if es_nota:
            numero_nota_ni = st.number_input("Número de nota (para nombre_activo = NOTA_XX)", min_value=1, max_value=999, step=1, key="ni_numero_nota")
            nombre_activo_final = f"NOTA_{int(numero_nota_ni):02d}"
            st.caption(f"nombre_activo se guardará como: **{nombre_activo_final}**")
            subtipo_inversion_final = "ESTRUCTURADA"
        else:
            nombre_activo_final = st.text_input("nombre_activo", value=tipo_inversion_final, key="ni_nombre_activo_libre")
            subtipo_conocidos = _valores_conocidos_columna(df_inv, "subtipo_inversion")
            subtipo_inversion_final = st.selectbox("subtipo_inversion", ([tipo_inversion_final] if tipo_inversion_final not in subtipo_conocidos else []) + subtipo_conocidos + ["otro"], key="ni_subtipo_sel")
            if subtipo_inversion_final == "otro":
                subtipo_inversion_final = st.text_input("Escribe subtipo_inversion", key="ni_subtipo_libre")

    if es_nota:
        fecha_sugerida_nota, periodicidad_sugerida_nota, fuente_sugerencia_nota = _sugerir_fecha_inicio_nota(int(numero_nota_ni), df_cal)
        if fecha_sugerida_nota is not None:
            periodo_txt = {1: "mensual", 3: "trimestral", 6: "semestral"}.get(periodicidad_sugerida_nota, f"cada {periodicidad_sugerida_nota} meses")
            st.info(
                f"📌 Fecha de inicio sugerida para INVERSIONES: **{fecha_sugerida_nota.strftime('%d/%m/%Y')}** "
                f"— calculada como el primer PAGO menos 1 periodo ({periodo_txt}), a partir de {fuente_sugerencia_nota}."
            )
            if st.button("📅 Usar esta fecha sugerida", key="ni_usar_fecha_sugerida"):
                st.session_state["ni_fecha_inicio"] = fecha_sugerida_nota.date()
                st.rerun()
        else:
            st.caption("No se encontró un calendario de PAGO guardado ni un borrador de extracción para esta nota — pon la fecha de inicio a mano (primer pago menos 1 periodo).")

    col4, col5 = st.columns(2)
    with col4:
        valor_fecha_inicio = fecha_sugerida_nota.date() if (es_nota and fecha_sugerida_nota is not None) else pd.Timestamp.today().date()
        fecha_inversion_final = st.date_input("Fecha de inicio (fecha_inversion)", value=valor_fecha_inicio, key="ni_fecha_inicio",
                                                help="Para notas: es la primera fecha de PAGO del calendario menos 1 periodo completo (1 mes si el cobro es mensual, 3 meses si es trimestral) — no el Initial Valuation Date del PDF.")
    with col5:
        if es_nota:
            interes_nota_anual_pct = st.number_input("Cupón anual de la NOTA (%)", min_value=0.0, max_value=200.0, step=0.5, format="%.2f", key="ni_interes_nota",
                                                        help="El cupón real de la nota — este es el dato que faltaba en el asistente de IA. Debe coincidir con CONTROL_NOTAS/el PDF de la nota.")
        else:
            tasas_fijas = {"paraguay": 15.0, "bolivia": 15.0, "motoclick": 25.0, "futbol": 15.0, "bitcoin": 20.0}
            default_tasa = tasas_fijas.get(tipo_inversion_final.strip().lower(), 0.0)
            interes_nota_anual_pct = st.number_input("Tasa anual que rinde el activo para la empresa (%)", min_value=0.0, max_value=200.0, step=0.5, format="%.2f", value=default_tasa, key="ni_interes_activo")

    with st.expander("Campos adicionales del activo (metodo_calculo, activo_generador_interes, cuenta_cobro)", expanded=False):
        metodo_calculo_ops = _valores_conocidos_columna(df_inv, "metodo_calculo")
        default_metodo = "NOTA" if es_nota and "NOTA" in metodo_calculo_ops else ((metodo_calculo_ops[0] if metodo_calculo_ops else "(dejar en blanco)"))
        opciones_metodo = ["(dejar en blanco)"] + metodo_calculo_ops + ["otro"]
        metodo_calculo_final = st.selectbox("metodo_calculo", opciones_metodo, index=opciones_metodo.index(default_metodo) if default_metodo in opciones_metodo else 0, key="ni_metodo_calculo_sel")
        if metodo_calculo_final == "otro":
            metodo_calculo_final = st.text_input("Escribe metodo_calculo", key="ni_metodo_calculo_libre")
        elif metodo_calculo_final == "(dejar en blanco)":
            metodo_calculo_final = ""

        activo_gen_ops = _valores_conocidos_columna(df_inv, "activo_generador_interes")
        default_activo_gen = "SI" if "SI" in activo_gen_ops else (activo_gen_ops[0] if activo_gen_ops else "(dejar en blanco)")
        opciones_activo_gen = ["(dejar en blanco)"] + activo_gen_ops + ["otro"]
        activo_generador_final = st.selectbox("activo_generador_interes", opciones_activo_gen, index=opciones_activo_gen.index(default_activo_gen) if default_activo_gen in opciones_activo_gen else 0, key="ni_activo_gen_sel")
        if activo_generador_final == "otro":
            activo_generador_final = st.text_input("Escribe activo_generador_interes", key="ni_activo_gen_libre")
        elif activo_generador_final == "(dejar en blanco)":
            activo_generador_final = ""

        cuenta_cobro_ops = _valores_conocidos_columna(df_inv, "cuenta_cobro")
        cuenta_cobro_final = st.selectbox("cuenta_cobro", ["(dejar en blanco)"] + cuenta_cobro_ops + ["otro"], key="ni_cuenta_cobro_sel")
        if cuenta_cobro_final == "otro":
            cuenta_cobro_final = st.text_input("Escribe cuenta_cobro", key="ni_cuenta_cobro_libre")
        elif cuenta_cobro_final == "(dejar en blanco)":
            cuenta_cobro_final = ""

    ids_inversion_origen_reinv = []
    if tipo_operacion == "REINVERSION":
        st.markdown("#### Origen de la reinversión")
        st.caption(
            "Podés elegir una o varias posiciones de origen (por ejemplo, si el capital de esta "
            "reinversión viene de varias posiciones cerradas a la vez). Se guarda en la columna "
            "id_inversion_origen de cada fila (separadas por coma si son varias) — trazabilidad de "
            "dónde viene el capital, no modifica ni cierra las posiciones originales. Si además hace "
            "falta cerrar formalmente alguna posición original, hazlo por separado con una operación CANCELADA."
        )
        activas_origen = _posiciones_activas_para_cerrar(df_inv)
        if not activas_origen.empty:
            activas_origen = activas_origen.copy()
            activas_origen["_etiqueta"] = activas_origen.apply(
                lambda r: f"{r.get('id_inversion','?')} | {r.get('inversor','?')} | {r.get('nombre_activo','?')} | ${float(r.get('capital_invertido',0) or 0):,.2f}",
                axis=1,
            )
            etiquetas_origen = st.multiselect("Posición(es) de origen (opcional, se aplica a todos los inversores de abajo)", activas_origen["_etiqueta"].tolist(), key="ni_origen_reinv_multisel")
            if etiquetas_origen:
                ids_inversion_origen_reinv = [
                    str(activas_origen[activas_origen["_etiqueta"] == et].iloc[0].get("id_inversion", ""))
                    for et in etiquetas_origen
                ]

    # ═══════════════════════════════════════════════════════════════════
    # 2. INVERSOR(ES) — uno o varios en la misma inversión
    # ═══════════════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("### 2. Inversor(es)")
    st.caption("Si en esta misma nota/activo participa más de un inversor, pulsa '➕ Añadir otro inversor' — comparten todo lo de arriba, pero cada uno tiene su propio capital, tasa, id_inversion, etc.")

    try:
        inversores_conocidos = sorted(set(list(USUARIOS_INVERSORES.keys()) + ["CHAPARRO FERNANDEZ"]) | set(df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).str.strip().unique()))
    except NameError:
        inversores_conocidos = sorted(df_inv.get("inversor", pd.Series(dtype=str)).dropna().astype(str).str.strip().unique())

    if "ni_num_inversores" not in st.session_state:
        st.session_state["ni_num_inversores"] = 1

    ids_recientes = df_inv["id_inversion"].dropna().astype(str).tail(5).tolist() if df_inv is not None and "id_inversion" in df_inv.columns else []
    if ids_recientes:
        st.caption(f"Últimos id_inversion usados (para mantener la misma convención): {', '.join(ids_recientes)}")

    filas_inversores = []
    for i in range(st.session_state["ni_num_inversores"]):
        with st.container(border=True):
            st.markdown(f"**Inversor #{i + 1}**")
            colx1, colx2, colx3 = st.columns(3)
            with colx1:
                inv_sel = st.selectbox("Inversor", inversores_conocidos + ["Otro (escribir)"], key=f"ni_inversor_sel_{i}")
                inv_final = st.text_input("Escribe el nombre (exacto, mayúsculas)", key=f"ni_inversor_libre_{i}") if inv_sel == "Otro (escribir)" else inv_sel
            with colx2:
                capital_i = st.number_input("Capital ($)", min_value=0.0, step=1000.0, format="%.2f", key=f"ni_capital_{i}")
            with colx3:
                tasa_inv_i = st.number_input(
                    "Interés anual AL INVERSOR (%)", min_value=0.0, max_value=100.0, step=0.5, format="%.2f", key=f"ni_interes_inversor_{i}",
                    help="OJO: para Roberto Viscafe, Crowe Bolivia y JR Real Estate la tasa es escalonada en el tiempo (5% → 7.5% → 10%) — revisa cuál corresponde a la fecha de inicio.",
                )
            colx4, colx5, colx6 = st.columns(3)
            with colx4:
                pago_i = st.selectbox("pago_intereses", ["reinvierte", "paga"], key=f"ni_pago_intereses_{i}")
            with colx5:
                cap_real_i = st.selectbox("capital_nuevo_real", ["si", "no"], index=0 if tipo_operacion == "NUEVA" else 1, key=f"ni_capital_real_{i}")
            with colx6:
                email_default_i = ""
                if df_inv is not None and not df_inv.empty and "inversor" in df_inv.columns and "email" in df_inv.columns:
                    fe = df_inv[df_inv["inversor"].astype(str).str.strip().str.upper() == inv_final.strip().upper()]
                    em = fe["email"].dropna().astype(str)
                    em = em[em.str.strip() != ""]
                    if not em.empty:
                        email_default_i = em.iloc[-1]
                email_i = st.text_input("email", value=email_default_i, key=f"ni_email_{i}")
            id_sugerido_i = _siguiente_id_inversion(df_inv, offset=i)
            if id_sugerido_i:
                st.caption(f"📌 Número de operación sugerido: **{id_sugerido_i}** (siguiente al último usado) — puedes escribir otro abajo si corresponde.")
            id_inv_i = st.text_input(
                "id_inversion (sigue la convención de arriba, ej. OP005)",
                value=id_sugerido_i, key=f"ni_id_inversion_{i}",
                help="Se sugiere automáticamente el siguiente número tras el último usado — puedes borrarlo y escribir otro a mano.",
            )

            if st.session_state["ni_num_inversores"] > 1:
                if st.button(f"🗑️ Quitar inversor #{i + 1}", key=f"ni_quitar_{i}"):
                    st.session_state["ni_num_inversores"] -= 1
                    st.rerun()

            filas_inversores.append({
                "id_inversion": id_inv_i,
                "inversor": inv_final,
                "tipo_inversion": tipo_inversion_final,
                "subtipo_inversion": subtipo_inversion_final,
                "nombre_activo": nombre_activo_final,
                "metodo_calculo": metodo_calculo_final,
                "cuenta_cobro": cuenta_cobro_final,
                "activo_generador_interes": activo_generador_final,
                "fecha_inversion": str(fecha_inversion_final),
                "fecha_final_inversion": "",
                "motivo": "",
                "capital_invertido": capital_i,
                "interes_nota_anual": round(interes_nota_anual_pct / 100.0, 6),
                "interes_inversor_anual": round(tasa_inv_i / 100.0, 6),
                "tipo_operacion": tipo_operacion,
                "id_inversion_origen": ", ".join(ids_inversion_origen_reinv) if ids_inversion_origen_reinv else "",
                "capital_nuevo_real": cap_real_i,
                "email": email_i,
                "pago_intereses": pago_i,
            })

    if st.button("➕ Añadir otro inversor a esta misma inversión", key="ni_add_inversor"):
        st.session_state["ni_num_inversores"] += 1
        st.rerun()

    # ── Construir preview (una fila por inversor) ────────────────────────
    st.markdown("---")
    st.markdown("#### Previsualización de las filas (revísalas antes de guardar)")
    df_preview = pd.DataFrame(filas_inversores)
    df_preview_editado = st.data_editor(df_preview, use_container_width=True, num_rows="dynamic", key="ni_preview_editor")

    filas_borrador = df_preview_editado.to_dict(orient="records")

    col_a, col_g, col_d = st.columns(3)
    with col_a:
        if st.button("📌 Guardar avance del borrador", key="ni_guardar_avance"):
            guardar_borrador_inversion(clave, {"tipo_operacion": tipo_operacion, "filas": filas_borrador})
            st.session_state["inversion_wizard_datos"] = {"tipo_operacion": tipo_operacion, "filas": filas_borrador}

    with col_g:
        confirmar = st.button("💾 Confirmar y escribir en INVERSIONES", type="primary", key="ni_confirmar")
    with col_d:
        if st.button("🗑️ Descartar borrador", key="ni_descartar"):
            borrar_borrador_inversion(clave)
            st.session_state["inversion_wizard_datos"] = {}
            st.session_state["inversion_wizard_clave_actual"] = ""
            st.session_state["ni_num_inversores"] = 1
            st.rerun()

    if confirmar:
        ids_vacios = [i + 1 for i, f in enumerate(filas_borrador) if not str(f.get("id_inversion", "")).strip()]
        if ids_vacios:
            st.error(f"Falta id_inversion en el/los inversor(es) #{', '.join(map(str, ids_vacios))} — no puede quedar vacío.")
            return
        if not filas_borrador:
            st.error("No hay ningún inversor en la previsualización.")
            return
        hojas = leer_todas_las_hojas_excel()
        if not hojas or "INVERSIONES" not in hojas:
            st.error("No se pudo leer INVERSIONES para guardar las filas nuevas.")
            return
        mapa = _mapa_columnas_reales(hojas, "INVERSIONES")
        if not mapa:
            st.error("No se pudo determinar el nombre real de las columnas de INVERSIONES.")
            return

        filas_reales = []
        columnas_faltantes_global = set()
        for fila_logica in filas_borrador:
            fila_real = {}
            for clave_logica, valor in fila_logica.items():
                if str(clave_logica).startswith("_"):
                    continue
                col_real = mapa.get(clave_logica)
                if col_real:
                    fila_real[col_real] = valor
                else:
                    columnas_faltantes_global.add(clave_logica)
            filas_reales.append(fila_real)
        if columnas_faltantes_global:
            st.warning(f"⚠️ Estas columnas no existen tal cual en la hoja INVERSIONES real y NO se escribieron: {', '.join(sorted(columnas_faltantes_global))}. Añádelas a mano en el Excel si hacen falta.")

        df_raw = hojas["INVERSIONES"]
        hojas["INVERSIONES"] = pd.concat([df_raw, pd.DataFrame(filas_reales)], ignore_index=True)

        # Vínculo adicional en la hoja histórica REINVERSIONES (aparte de la columna
        # id_inversion_origen ya escrita en cada fila de INVERSIONES) — solo si esa hoja existe.
        if tipo_operacion == "REINVERSION" and "REINVERSIONES" in hojas:
            mapa_reinv = _mapa_columnas_reales(hojas, "REINVERSIONES")
            col_origen = mapa_reinv.get("id_inversion_origen")
            col_destino = mapa_reinv.get("id_inversion_destino")
            if col_origen and col_destino:
                filas_reinv = []
                for fila_logica in filas_borrador:
                    ids_origen_f = str(fila_logica.get("id_inversion_origen") or "").strip()
                    if not ids_origen_f:
                        continue
                    for id_origen_uno in [x.strip() for x in ids_origen_f.split(",") if x.strip()]:
                        fila_reinv = {col_origen: id_origen_uno, col_destino: fila_logica.get("id_inversion", "")}
                        if mapa_reinv.get("fecha"):
                            fila_reinv[mapa_reinv["fecha"]] = pd.Timestamp(fecha_inversion_final)
                        if mapa_reinv.get("importe"):
                            fila_reinv[mapa_reinv["importe"]] = fila_logica.get("capital_invertido")
                        filas_reinv.append(fila_reinv)
                if filas_reinv:
                    hojas["REINVERSIONES"] = pd.concat([hojas["REINVERSIONES"], pd.DataFrame(filas_reinv)], ignore_index=True)
            else:
                st.info("No se encontraron las columnas id_inversion_origen/id_inversion_destino en REINVERSIONES — añade el vínculo a mano si hace falta.")

        # Quitar el borrador de esta misma tanda de hojas para no escribir el Excel dos veces
        if "BORRADORES_INVERSIONES" in hojas and not hojas["BORRADORES_INVERSIONES"].empty:
            df_b = hojas["BORRADORES_INVERSIONES"]
            hojas["BORRADORES_INVERSIONES"] = df_b[df_b["CLAVE"].astype(str) != str(clave)]

        guardar_excel_completo_desde_hojas(hojas)
        st.session_state["inversion_wizard_datos"] = {}
        st.session_state["inversion_wizard_clave_actual"] = ""
        st.session_state["ni_num_inversores"] = 1
        ids_guardados = ", ".join(str(f.get("id_inversion", "")) for f in filas_borrador)
        st.success(f"{len(filas_borrador)} fila(s) guardada(s) en INVERSIONES (id_inversion: {ids_guardados}).")
        st.cache_data.clear()


def seccion_deuda_jordi():
    """Pantalla de seguimiento de la deuda con Jordi Chaparro."""
    st.header("🏦 Deuda con Jordi Chaparro")

    df_inv, df_cal, df_control = cargar_excel_completo()

    # Leer hoja DEUDA_JORDI
    try:
        df_deuda = pd.read_excel(ARCHIVO, sheet_name="DEUDA_JORDI")
        df_deuda.columns = [str(c).strip().lower() for c in df_deuda.columns]
        capital_inicial = float(df_deuda["capital_deuda_inicial"].iloc[0])
        fecha_inicio = pd.to_datetime(df_deuda["fecha_inicio_deuda"].iloc[0], dayfirst=True)
    except Exception as e:
        st.error(f"No se puede leer la hoja DEUDA_JORDI del Excel: {e}")
        st.info("Asegúrate de que la hoja DEUDA_JORDI tiene las columnas: fecha_inicio_deuda y capital_deuda_inicial")
        return

    st.caption(f"Deuda inicial: **${capital_inicial:,.2f}** · Fecha inicio: **{fecha_inicio.strftime('%d/%m/%Y')}**")
    st.caption("La deuda **baja** con: beneficio empresa de Paraguay + Motoclick + Fútbol + Bolivia + cobros de notas en cuenta JORDI.  "
               "La deuda **sube** con: intereses mensuales de JEP (todos sus activos) + reparto de dividendos.")

    with st.spinner("Calculando evolución de la deuda..."):
        df_evol = calcular_deuda_jordi(df_inv, df_cal, df_control, capital_inicial, fecha_inicio)

    if df_evol.empty:
        st.warning("No hay datos para calcular la deuda.")
        return

    # Saldo actual = último mes calculado
    ultimo = df_evol.iloc[-1]
    saldo_actual = ultimo["saldo_fin"]

    # ── Mensaje personalizado ──────────────────────────────────────────────
    st.divider()
    # La deuda inicial es lo que la compañía debe a Jordi.
    # saldo > 0 → la compañía aún le debe a Jordi
    # saldo < 0 → Jordi ha devuelto más de lo que debía (la compañía está en positivo)
    if saldo_actual > 0:
        st.markdown(
            f"""
            <div style="background:#d4edda;border-left:6px solid #28a745;padding:18px 24px;border-radius:8px;font-size:1.15rem;">
            🧾 <b>Señor Chaparro, la compañía le debe a usted:</b>
            <span style="font-size:1.6rem;font-weight:bold;color:#155724;margin-left:12px;">${saldo_actual:,.2f}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
    elif saldo_actual < 0:
        st.markdown(
            f"""
            <div style="background:#fff3cd;border-left:6px solid #f5a623;padding:18px 24px;border-radius:8px;font-size:1.15rem;">
            ✅ <b>Señor Chaparro, usted le debe a la compañía:</b>
            <span style="font-size:1.6rem;font-weight:bold;color:#c0392b;margin-left:12px;">${abs(saldo_actual):,.2f}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            """
            <div style="background:#e2e3e5;border-left:6px solid #6c757d;padding:18px 24px;border-radius:8px;font-size:1.15rem;">
            ⚖️ <b>Señor Chaparro, la deuda está completamente saldada.</b>
            </div>
            """,
            unsafe_allow_html=True,
        )
    st.divider()

    col1, col2, col3 = st.columns(3)
    col1.metric("💰 Deuda inicial", f"${capital_inicial:,.2f}")
    col2.metric("📅 Saldo actual (fin mes)", f"${saldo_actual:,.2f}",
                delta=f"${saldo_actual - capital_inicial:,.2f}")
    col3.metric("📆 Meses calculados", len(df_evol))

    st.divider()

    # ── Desglose mes a mes en formato "cuenta" ─────────────────────────────
    st.subheader("📋 Desglose por mes")

    for _, fila in df_evol.iterrows():
        variacion = fila["variacion_neta"]
        saldo_fin = fila["saldo_fin"]
        emoji = "📉" if variacion < 0 else "📈"
        with st.expander(f"{emoji} {fila['mes']}  |  Inicio: ${fila['saldo_inicio']:,.2f}  →  Fin: ${fila['saldo_fin']:,.2f}  (variación: {'+' if variacion>0 else ''}{variacion:,.2f})", expanded=False):

            color_fin = "#c0392b" if saldo_fin > 0 else "#155724"
            color_var = "#c0392b" if variacion > 0 else "#155724"
            signo = "+" if variacion > 0 else ""

            html = f"""<div style="font-family:Arial,sans-serif;font-size:0.93rem;line-height:1.9;">
              <div style="display:flex;justify-content:space-between;border-bottom:2px solid #555;padding-bottom:6px;margin-bottom:10px;">
                <span style="font-weight:bold;font-size:1rem;">💰 Saldo inicio del mes</span>
                <span style="font-weight:bold;font-size:1rem;">${fila['saldo_inicio']:,.2f}</span>
              </div>"""

            # ── JEP: AUMENTA DEUDA ── agrupado por tipo
            jep_grupos = fila.get("jep_por_grupo", {}) or {}
            if jep_grupos:
                html += f"""<div style="background:#fdf0f0;border-left:4px solid #c0392b;padding:8px 12px;margin:8px 0;border-radius:4px;">
                  <div style="color:#c0392b;font-weight:bold;margin-bottom:4px;">➕ AUMENTA LA DEUDA — Intereses devengados a JEP</div>"""
                for grupo, importe in jep_grupos.items():
                    html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#7b241c;">
                      <span>{grupo}</span><span style="font-weight:bold;">+${importe:,.2f}</span></div>"""
                html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#c0392b;font-weight:bold;border-top:1px solid #f5c6cb;margin-top:4px;padding-top:4px;">
                      <span>TOTAL intereses JEP</span><span>+${fila['pago_jep']:,.2f}</span></div></div>"""

            # ── REPARTO DE DIVIDENDOS: AUMENTA DEUDA ──
            reparto_items = fila.get("reparto_por_item", {}) or {}
            if reparto_items:
                html += f"""<div style="background:#fdf0f0;border-left:4px solid #c0392b;padding:8px 12px;margin:8px 0;border-radius:4px;">
                  <div style="color:#c0392b;font-weight:bold;margin-bottom:4px;">➕ AUMENTA LA DEUDA — Reparto de dividendos</div>"""
                for desc, importe in reparto_items.items():
                    html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#7b241c;">
                      <span>{desc}</span><span style="font-weight:bold;">+${importe:,.2f}</span></div>"""
                html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#c0392b;font-weight:bold;border-top:1px solid #f5c6cb;margin-top:4px;padding-top:4px;">
                      <span>TOTAL reparto dividendos</span><span>+${fila['reparto_dividendos']:,.2f}</span></div></div>"""

            # ── TRANSFERENCIAS DE JORDI: AUMENTA DEUDA ──
            transferencias_items = fila.get("transferencias_por_item", {}) or {}
            if transferencias_items:
                html += f"""<div style="background:#fdf0f0;border-left:4px solid #c0392b;padding:8px 12px;margin:8px 0;border-radius:4px;">
                  <div style="color:#c0392b;font-weight:bold;margin-bottom:4px;">➕ AUMENTA LA DEUDA — Transferencias de Jordi al fondo</div>"""
                for desc, importe in transferencias_items.items():
                    html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#7b241c;">
                      <span>{desc}</span><span style="font-weight:bold;">+${importe:,.2f}</span></div>"""
                html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#c0392b;font-weight:bold;border-top:1px solid #f5c6cb;margin-top:4px;padding-top:4px;">
                      <span>TOTAL transferencias Jordi</span><span>+${fila['transferencias_jordi']:,.2f}</span></div></div>"""

            # ── ACTIVOS FIJOS: REDUCE DEUDA ── cobro bruto por activo
            det_fijos = fila.get("detalle_fijos", []) or []
            if det_fijos:
                html += f"""<div style="background:#f0fdf4;border-left:4px solid #28a745;padding:8px 12px;margin:8px 0;border-radius:4px;">
                  <div style="color:#155724;font-weight:bold;margin-bottom:4px;">➖ REDUCE LA DEUDA — Ingresos compañía activos fijos</div>"""
                for d in det_fijos:
                    html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#1a5c30;">
                      <span>{d['activo']}</span><span style="font-weight:bold;">−${d['cobro']:,.2f}</span></div>"""
                html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#155724;font-weight:bold;border-top:1px solid #c3e6cb;margin-top:4px;padding-top:4px;">
                      <span>TOTAL fijos</span><span>−${fila['ingreso_fijos']:,.2f}</span></div></div>"""

            # ── NOTAS JORDI: REDUCE DEUDA ── cobro bruto por nota
            por_nota = fila.get("por_nota_jordi", {}) or {}
            if por_nota:
                html += f"""<div style="background:#f0f8ff;border-left:4px solid #1a73e8;padding:8px 12px;margin:8px 0;border-radius:4px;">
                  <div style="color:#1a3c6e;font-weight:bold;margin-bottom:4px;">➖ REDUCE LA DEUDA — Cobro notas cuenta JORDI</div>"""
                for nombre_nota, cobro in por_nota.items():
                    html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#1a3c6e;">
                      <span>{nombre_nota}</span><span style="font-weight:bold;">−${cobro:,.2f}</span></div>"""
                html += f"""<div style="display:flex;justify-content:space-between;padding-left:12px;color:#1a3c6e;font-weight:bold;border-top:1px solid #b8daff;margin-top:4px;padding-top:4px;">
                      <span>TOTAL notas JORDI</span><span>−${fila['cobro_notas_jordi']:,.2f}</span></div></div>"""

            html += f"""
              <div style="display:flex;justify-content:space-between;border-top:2px solid #555;margin-top:10px;padding-top:8px;">
                <span style="font-weight:bold;">Variación neta del mes</span>
                <span style="color:{color_var};font-weight:bold;">{signo}${variacion:,.2f}</span>
              </div>
              <div style="display:flex;justify-content:space-between;background:#f8f9fa;border-radius:6px;padding:10px 6px;margin-top:6px;">
                <span style="font-weight:bold;font-size:1.05rem;">✅ Saldo fin de mes</span>
                <span style="color:{color_fin};font-size:1.15rem;font-weight:bold;">${saldo_fin:,.2f}</span>
              </div>
            </div>"""
            st.markdown(html, unsafe_allow_html=True)

    st.divider()

    # ── Tabla resumen mes a mes ────────────────────────────────────────────
    st.subheader("📊 Tabla resumen")
    cols_tabla = ["mes","saldo_inicio","ingreso_fijos","cobro_notas_jordi","total_resta","pago_jep","reparto_dividendos","transferencias_jordi","variacion_neta","saldo_fin"]
    tabla = df_evol[cols_tabla].copy()
    tabla = tabla.rename(columns={
        "mes": "Mes",
        "saldo_inicio": "Saldo inicio ($)",
        "ingreso_fijos": "Ingreso fijos ($)",
        "cobro_notas_jordi": "Cobro notas JORDI ($)",
        "total_resta": "Total resta ($)",
        "pago_jep": "Intereses JEP ($)",
        "reparto_dividendos": "Reparto dividendos ($)",
        "transferencias_jordi": "Transferencias Jordi ($)",
        "variacion_neta": "Variación neta ($)",
        "saldo_fin": "Saldo fin mes ($)",
    })

    def color_col(series):
        return [
            ("color: green; font-weight: bold" if v < 0 else ("color: red; font-weight: bold" if v > 0 else ""))
            for v in series
        ]

    styled = (
        tabla.style
        .format({
            "Saldo inicio ($)": "${:,.2f}",
            "Benef. fijos ($)": "${:,.2f}",
            "Benef. notas JORDI ($)": "${:,.2f}",
            "Total resta ($)": "${:,.2f}",
            "Intereses JEP ($)": "${:,.2f}",
            "Reparto dividendos ($)": "${:,.2f}",
            "Variación neta ($)": "${:,.2f}",
            "Saldo fin mes ($)": "${:,.2f}",
        })
        .apply(color_col, subset=["Variación neta ($)"])
        .apply(color_col, subset=["Saldo fin mes ($)"])
    )

    st.dataframe(styled, use_container_width=True, hide_index=True)
    boton_descarga_excel(tabla, f"deuda_jordi_evolucion.xlsx")

    # ── Gráfico evolución saldo ────────────────────────────────────────────
    st.subheader("📈 Gráfico de evolución")
    import plotly.graph_objects as go
    colores = ["#28a745" if v <= 0 else "#e63946" for v in df_evol["saldo_fin"]]
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=df_evol["mes"],
        y=df_evol["saldo_fin"],
        mode="lines+markers",
        name="Saldo deuda",
        line=dict(color="#e63946", width=2),
        marker=dict(size=9, color=colores),
    ))
    fig.add_hline(y=0, line_dash="dash", line_color="green", annotation_text="Deuda saldada")
    fig.update_layout(
        xaxis_title="Mes",
        yaxis_title="Saldo ($)",
        hovermode="x unified",
        height=400,
    )
    st.plotly_chart(fig, use_container_width=True)


# ═══════════════════════════════════════════════════════════════════════════
# CEREBRO DEL ASISTENTE IA — función de nivel de módulo (no depende de Streamlit
# salvo el flag opcional mostrar_debug_ui), para que la reutilice tanto la pestaña
# "Asistente IA" de la app web como el servicio de WhatsApp, con idéntico resultado.
# ═══════════════════════════════════════════════════════════════════════════
def construir_contexto_ia_fondo(pregunta: str, df_inv, df_cal, df_control, fecha_limite=None, mostrar_debug_ui: bool = False) -> str:
    """
    Contexto COMPLETO con fuentes correctas para cada dato:
    - Intereses a inversores: lógica de EXTRACTOS (solo NUEVA+CANCELADA, sin reinversiones)
    - Ingresos compañía / cobros notas: lógica de DASHBOARD (detalle_activo_mes + resumen_notas_mes)
    - Capital activo por inversor: lógica de CENTRO DE CONTROL (capital_activo_en_fecha)
    """
    hoy = pd.Timestamp.today().normalize()
    anio_hoy, mes_hoy = hoy.year, hoy.month
    p = pregunta.lower()

    # Detectar mes mencionado en la pregunta
    meses_map = {"enero":1,"febrero":2,"marzo":3,"abril":4,"mayo":5,"junio":6,
                 "julio":7,"agosto":8,"septiembre":9,"octubre":10,"noviembre":11,"diciembre":12}
    mes_mencionado = next((v for k,v in meses_map.items() if k in p), None)
    mes_pregunta = mes_mencionado if mes_mencionado is not None else mes_hoy

    # Detectar un año explícito de 4 dígitos (ej. "2025") en la pregunta.
    import re as _re
    anio_match = _re.search(r"\b(20\d{2})\b", p)
    anio_pregunta = int(anio_match.group(1)) if anio_match else anio_hoy

    # Si se menciona un año pero NO un mes concreto, se interpreta como petición de TOTAL ANUAL
    # (ej. "cuántos intereses se pagaron a PAM en 2025" → los 12 meses de 2025, no solo el mes actual).
    pedir_total_anual = (anio_match is not None) and (mes_mencionado is None)

    # Calcular resumen dashboard al inicio para usarlo en todos los bloques
    try:
        resumen_dash = obtener_resumen_dashboard(
            df_inv, df_cal, df_control,
            anio=anio_pregunta, mes=mes_pregunta,
            vista_activo="General", incluir_chaparro=True
        )
    except Exception:
        resumen_dash = {}

    lineas = [
        f"Fecha de hoy: {hoy.strftime('%d/%m/%Y')}",
        f"Mes de referencia: {mes_pregunta}/{anio_pregunta}",
        "",
        "FUENTE DE DATOS:",
        "- Intereses pagados a inversores → lógica EXTRACTOS (solo operaciones NUEVA y CANCELADA, sin reinversiones)",
        "- Ingresos y beneficio empresa → lógica DASHBOARD (obtener_resumen_dashboard, incluir_chaparro=True)",
        "- Capital activo → capital_activo_en_fecha + obtener_resumen_dashboard",
    ]

    # ══════════════════════════════════════════════════════════════════════
    # 1. CAPITAL ACTIVO HOY — fuente: capital_activo_en_fecha (Centro de control)
    # ══════════════════════════════════════════════════════════════════════
    try:
        lineas.append(f"\n=== CAPITAL ACTIVO HOY ===")
        # Total con Chaparro Fernández incluido (igual que Dashboard con checkbox activado)
        cap_con_chaparro = resumen_dash.get("capital_total", 0.0)
        # Total sin Chaparro Fernández (igual que Centro de control por defecto)
        df_sin_cf = aplicar_filtro_chaparro_fernandez(df_inv, incluir_chaparro=False)
        cap_sin_chaparro = capital_activo_en_fecha(df_sin_cf, hoy)
        lineas.append(f"TOTAL FONDO (CON Chaparro Fernández, como Dashboard): ${cap_con_chaparro:,.2f}")
        lineas.append(f"TOTAL FONDO (SIN Chaparro Fernández, como Centro de control): ${cap_sin_chaparro:,.2f}")

        # Desglose por activo (con Chaparro incluido)
        df_notas_cap = df_inv[df_inv["tipo_inversion"].astype(str).str.lower().str.strip() == "nota"]
        cap_notas = capital_activo_en_fecha(df_notas_cap, hoy)
        lineas.append(f"  NOTAS: ${cap_notas:,.2f}")
        for activo in ["futbol","paraguay","bolivia","motoclick","bitcoin"]:
            cap = capital_activo_en_fecha(df_inv, hoy, activo)
            lineas.append(f"  {activo.upper()}: ${cap:,.2f}")

        # Capital activo por inversor
        lineas.append("\nCapital activo por inversor (todos, incluyendo Chaparro Fernández):")
        df_cap = df_inv.copy()
        df_cap["fecha_inversion"] = parsear_fecha_robusta(df_cap.get("fecha_inversion"))
        df_cap["fecha_final_inversion"] = parsear_fecha_robusta(df_cap.get("fecha_final_inversion"))
        df_cap["capital_invertido"] = pd.to_numeric(df_cap.get("capital_invertido"), errors="coerce").fillna(0)
        df_cap["tipo_op_n"] = df_cap["tipo_operacion"].astype(str).str.strip().str.upper()
        activas_cc = df_cap[
            (df_cap["fecha_inversion"].notna()) &
            (df_cap["fecha_inversion"] <= hoy) &
            (df_cap["fecha_final_inversion"].isna() | (df_cap["fecha_final_inversion"] >= hoy))
        ].copy()
        es_cancelada = activas_cc["tipo_op_n"] == "CANCELADA"
        activas_cc = activas_cc[~es_cancelada]
        cap_por_inv = activas_cc.groupby("inversor")["capital_invertido"].sum().sort_values(ascending=False)
        for inv, cap in cap_por_inv.items():
            lineas.append(f"  {inv}: ${cap:,.2f}")
    except Exception as e:
        lineas.append(f"[Error capital activo: {e}]")

    # ══════════════════════════════════════════════════════════════════════
    # 2. INTERESES A INVERSORES DEL MES — fuente: lógica EXTRACTOS
    #    Solo NUEVA y CANCELADA. Reinversiones NO generan pago independiente.
    # ══════════════════════════════════════════════════════════════════════
    try:
        df_ext = df_inv.copy()
        for col in ["inversor","tipo_inversion","subtipo_inversion","nombre_activo","tipo_operacion","id_inversion"]:
            if col in df_ext.columns:
                df_ext[col] = df_ext[col].fillna("").astype(str).str.strip()
        df_ext["tipo_op_n"] = df_ext["tipo_operacion"].str.upper()
        df_ext = df_ext[df_ext["tipo_op_n"].isin(["NUEVA","CANCELADA"])].copy()
        df_ext["fecha_inversion"] = parsear_fecha_robusta(df_ext.get("fecha_inversion"))
        df_ext["fecha_final_inversion"] = parsear_fecha_robusta(df_ext.get("fecha_final_inversion"))
        df_ext["capital_invertido"] = pd.to_numeric(df_ext.get("capital_invertido"), errors="coerce").fillna(0)
        df_ext["interes_inversor_anual"] = pd.to_numeric(df_ext.get("interes_inversor_anual"), errors="coerce").fillna(0)

        INVERSORES_TRAMO = {"ROBERTO VISCAFE", "CROWE BOLIVIA", "JR REAL ESTATE"}
        fin_tramo1 = datetime(2026, 1, 31)
        inicio_tramo2 = datetime(2026, 2, 1)
        fin_tramo2 = datetime(2026, 6, 30)
        inicio_tramo3 = datetime(2026, 7, 1)

        def _calcular_intereses_mes(anio_m, mes_m):
            dias_mes_ext = ultimo_dia_mes(anio_m, mes_m)
            fecha_corte_ext = pd.Timestamp(datetime(anio_m, mes_m, dias_mes_ext))
            filas = []
            for _, row in df_ext.iterrows():
                fi = row.get("fecha_inversion")
                if pd.isna(fi):
                    continue
                fi_dt = fi.to_pydatetime()
                tipo_op = row["tipo_op_n"]
                ff = row.get("fecha_final_inversion")

                if tipo_op == "CANCELADA":
                    if pd.isna(ff):
                        continue
                    fecha_fin_dt = min(ff.to_pydatetime(), fecha_corte_ext.to_pydatetime())
                else:
                    fecha_fin_dt = fecha_corte_ext.to_pydatetime()

                inicio_mes_dt = datetime(anio_m, mes_m, 1)
                fin_mes_dt = datetime(anio_m, mes_m, dias_mes_ext)
                inicio_calc = max(fi_dt, inicio_mes_dt)
                fin_calc = min(fecha_fin_dt, fin_mes_dt)
                if inicio_calc > fin_calc:
                    continue

                dias = (fin_calc - inicio_calc).days + 1
                capital = float(row["capital_invertido"])
                tasa = float(row["interes_inversor_anual"])
                inv_upper = str(row.get("inversor","")).strip().upper()

                if inv_upper in INVERSORES_TRAMO:
                    interes_mes = 0.0
                    if inicio_calc <= fin_tramo1:
                        fin_t1 = min(fin_calc, fin_tramo1)
                        dias_t1 = (fin_t1 - inicio_calc).days + 1
                        interes_mes += round((capital * 0.05 / 12) * dias_t1 / dias_mes_ext, 2)
                    ini_t2 = max(inicio_calc, inicio_tramo2)
                    fin_t2 = min(fin_calc, fin_tramo2)
                    if ini_t2 <= fin_t2:
                        dias_t2 = (fin_t2 - ini_t2).days + 1
                        interes_mes += round((capital * 0.075 / 12) * dias_t2 / dias_mes_ext, 2)
                    if fin_calc >= inicio_tramo3:
                        ini_t3 = max(inicio_calc, inicio_tramo3)
                        dias_t3 = (fin_calc - ini_t3).days + 1
                        interes_mes += round((capital * 0.10 / 12) * dias_t3 / dias_mes_ext, 2)
                else:
                    interes_mes = round((capital * tasa / 12) * dias / dias_mes_ext, 2)

                filas.append({
                    "mes": f"{mes_m:02d}/{anio_m}",
                    "inversor": str(row.get("inversor","")),
                    "nombre_activo": str(row.get("nombre_activo","")),
                    "capital": capital,
                    "tasa_anual": tasa,
                    "dias": dias,
                    "interes_mes": interes_mes,
                })
            return filas

        if pedir_total_anual:
            lineas.append(f"\n=== INTERESES A PAGAR A INVERSORES — TOTAL AÑO {anio_pregunta} (fuente: lógica Extractos) ===")
            lineas.append("REGLA: solo operaciones NUEVA y CANCELADA. Reinversiones excluidas. Suma de los 12 meses del año solicitado.")
            filas_int = []
            for m in range(1, 13):
                filas_int.extend(_calcular_intereses_mes(anio_pregunta, m))
        else:
            lineas.append(f"\n=== INTERESES A PAGAR A INVERSORES {mes_pregunta}/{anio_pregunta} (fuente: lógica Extractos) ===")
            lineas.append("REGLA: solo operaciones NUEVA y CANCELADA. Reinversiones excluidas.")
            filas_int = _calcular_intereses_mes(anio_pregunta, mes_pregunta)

        if filas_int:
            df_int = pd.DataFrame(filas_int)
            # Por inversor agrupado (total del periodo solicitado, sea un mes o el año completo)
            por_inv = df_int.groupby("inversor")["interes_mes"].sum().sort_values(ascending=False)
            total_int = por_inv.sum()
            for inv, val in por_inv.items():
                lineas.append(f"  {inv}: ${val:,.2f}")
            lineas.append(f"  >> TOTAL INTERESES A PAGAR: ${total_int:,.2f}")
            if pedir_total_anual:
                lineas.append("\nDesglose mensual por inversor (para el año solicitado):")
                desglose = df_int.groupby(["inversor","mes"])["interes_mes"].sum().reset_index()
                for inv in por_inv.index:
                    sub = desglose[desglose["inversor"] == inv]
                    detalle_meses = ", ".join(f"{r['mes']}: ${r['interes_mes']:,.2f}" for _, r in sub.iterrows())
                    lineas.append(f"  {inv} → {detalle_meses}")
            else:
                lineas.append("\nDetalle por posición:")
                for r in filas_int:
                    lineas.append(f"  {r['inversor']} | {r['nombre_activo']} | Capital: ${r['capital']:,.2f} | Tasa: {r['tasa_anual']*100:.1f}% | Días: {r['dias']} | Interés: ${r['interes_mes']:,.2f}")
        else:
            lineas.append("  Sin datos de intereses para ese periodo.")
    except Exception as e:
        lineas.append(f"[Error intereses extracto: {e}]")

    # ══════════════════════════════════════════════════════════════════════
    # 2b. CAPITAL DE CHAPARRO FERNÁNDEZ Y BENEFICIO REAL DE LA COMPAÑÍA
    #     (precalculado para que la IA no tenga que sumar decenas de filas ella misma)
    # ══════════════════════════════════════════════════════════════════════
    try:
        df_cf = df_ext.copy()  # ya filtrado a NUEVA/CANCELADA, con pago_intereses disponible
        df_cf["inversor_up"] = df_cf["inversor"].str.upper()
        df_cf["pago_intereses"] = df_cf.get("pago_intereses", "").fillna("").astype(str).str.strip().str.lower()

        # Capital hoy a nombre de Chaparro Fernández — se usa la misma función canónica que
        # el Centro de Control (capital_activo_en_fecha), que ya excluye correctamente las
        # posiciones "nueva" que fueron reinvertidas (fecha_final_inversion pasada), evitando
        # duplicar capital entre una fila nueva ya cerrada y su reinversión sucesora.
        df_inv_cf = df_inv[df_inv["inversor"].astype(str).str.strip().str.upper() == "CHAPARRO FERNANDEZ"].copy()
        capital_cf_hoy = capital_activo_en_fecha(df_inv_cf, hoy)

        # Intereses devengados y NO pagados acumulados: todos los inversores con pago_intereses="reinvierte",
        # EXCLUYENDO a JEP (que sí cobra en efectivo) y a Chaparro Fernández (no se debe nada a sí mismo).
        # Se suman TODOS los meses desde que empezó cada posición hasta hoy (deuda acumulada, nunca se resetea).
        df_reinv = df_cf[(df_cf["pago_intereses"] == "reinvierte") & (df_cf["inversor_up"] != "CHAPARRO FERNANDEZ")]
        if not df_reinv.empty:
            fecha_min_reinv = df_reinv["fecha_inversion"].min()
            total_devengado_no_pagado = 0.0
            por_inv_devengado = {}
            f_iter = pd.Timestamp(fecha_min_reinv.year, fecha_min_reinv.month, 1)
            while f_iter <= hoy:
                for _, r in df_reinv.iterrows():
                    v = 0.0
                    fi = r.get("fecha_inversion")
                    if pd.isna(fi):
                        continue
                    dias_m = ultimo_dia_mes(f_iter.year, f_iter.month)
                    inicio_mes_r = datetime(f_iter.year, f_iter.month, 1)
                    fin_mes_r = datetime(f_iter.year, f_iter.month, dias_m)
                    ff = r.get("fecha_final_inversion")
                    if r["tipo_op_n"] == "CANCELADA" and pd.notna(ff):
                        fecha_fin_r = min(ff.to_pydatetime(), fin_mes_r)
                    else:
                        fecha_fin_r = fin_mes_r
                    inicio_calc_r = max(fi.to_pydatetime(), inicio_mes_r)
                    fin_calc_r = min(fecha_fin_r, fin_mes_r)
                    if inicio_calc_r > fin_calc_r:
                        continue
                    dias_r = (fin_calc_r - inicio_calc_r).days + 1
                    v = round(float(r["capital_invertido"]) * float(r["interes_inversor_anual"]) / 12 * dias_r / dias_m, 2)
                    total_devengado_no_pagado += v
                    por_inv_devengado[r["inversor"]] = por_inv_devengado.get(r["inversor"], 0) + v
                f_iter += pd.DateOffset(months=1)

            beneficio_real = capital_cf_hoy - total_devengado_no_pagado
            lineas.append(f"\n=== CAPITAL CHAPARRO FERNÁNDEZ Y BENEFICIO REAL DE LA COMPAÑÍA (precalculado) ===")
            lineas.append(f"Capital hoy a nombre de CHAPARRO FERNANDEZ: ${capital_cf_hoy:,.2f}")
            lineas.append(f"Intereses devengados y NO pagados acumulados (todos los inversores excepto JEP, desde que empezó cada uno hasta hoy): ${total_devengado_no_pagado:,.2f}")
            lineas.append(f"  Desglose: " + ", ".join(f"{k}: ${v:,.2f}" for k, v in sorted(por_inv_devengado.items(), key=lambda x: -x[1])))
            lineas.append(f">> BENEFICIO REAL DE LA COMPAÑÍA = ${capital_cf_hoy:,.2f} - ${total_devengado_no_pagado:,.2f} = ${beneficio_real:,.2f}")
    except Exception as e:
        lineas.append(f"[Error cálculo beneficio real Chaparro Fernández: {e}]")

    # ══════════════════════════════════════════════════════════════════════
    # 3. INGRESOS EMPRESA DEL MES — fuente: obtener_resumen_dashboard (EXACTAMENTE igual que el Dashboard)
    # ══════════════════════════════════════════════════════════════════════
    try:
        cobro_total   = resumen_dash.get("cobro_total_mes", 0.0)
        pago_total    = resumen_dash.get("pago_total_mes", 0.0)
        benef_total   = resumen_dash.get("beneficio_total_mes", 0.0)
        cap_dashboard = resumen_dash.get("capital_total", 0.0)

        lineas.append(f"\n=== INGRESOS COMPAÑÍA {mes_pregunta}/{anio_pregunta} (fuente: Dashboard — obtener_resumen_dashboard) ===")
        lineas.append(f"  COBRO TOTAL COMPAÑÍA: ${cobro_total:,.2f}")
        lineas.append(f"  PAGO TOTAL INVERSORES (dashboard): ${pago_total:,.2f}")
        lineas.append(f"  BENEFICIO EMPRESA: ${benef_total:,.2f}")
        lineas.append(f"  CAPITAL ACTIVO (dashboard): ${cap_dashboard:,.2f}")

        # Desglose por activo si está disponible
        rent_activo = resumen_dash.get("rentabilidad_por_activo", None)
        if rent_activo is not None and not rent_activo.empty:
            lineas.append("  Desglose por activo:")
            for _, r in rent_activo.iterrows():
                lineas.append(f"    {str(r.get('activo','')).upper()}: cobro ${float(r.get('cobro_compania_mes',0)):,.2f} | pago inversores ${float(r.get('pago_inversor_mes',0)):,.2f} | beneficio ${float(r.get('beneficio_empresa_mes',0)):,.2f}")
    except Exception as e:
        lineas.append(f"[Error ingresos empresa: {e}]")

    # ══════════════════════════════════════════════════════════════════════
    # 3b. DESGLOSE DE COBROS DE NOTAS DEL MES SOLICITADO — YA CALCULADO Y SUMADO
    #     Fuente: resumen_notas_mes (idéntica lógica que el Dashboard, calendario real).
    #     Existe para que la IA NO tenga que extraer/sumar filas a mano de la tabla ancha
    #     de "CALENDARIO INTEGRADO PRÓXIMOS 180 DÍAS" de más abajo — eso ya causó un error
    #     real (una nota olvidada y fechas desplazadas al resumir una tabla larga). Para
    #     preguntas tipo "¿cuánto cobraremos de notas en agosto?" o "dame el calendario de
    #     cobros de notas de [mes]", usa ESTE bloque tal cual, no el de 180 días.
    # ══════════════════════════════════════════════════════════════════════
    try:
        _, _, _, detalle_mes_notas, _ = resumen_notas_mes(df_inv, df_cal, df_control, anio_pregunta, mes_pregunta, prorratear=False)
        lineas.append(f"\n=== COBROS DE NOTAS {mes_pregunta}/{anio_pregunta} — DESGLOSE YA SUMADO POR FECHA Y NOTA (usar tal cual, no recalcular) ===")
        if detalle_mes_notas is not None and not detalle_mes_notas.empty:
            agg_notas_mes = detalle_mes_notas.groupby(["fecha_pago", "nota"], as_index=False).agg(
                monto_cobro=("cobro_compania", "sum"),
                resultado_observacion=("resultado_observacion", "first"),
            ).sort_values(["fecha_pago", "nota"])
            for r in agg_notas_mes.itertuples():
                fecha_str = pd.Timestamp(r.fecha_pago).strftime("%d/%m/%Y") if pd.notna(r.fecha_pago) else "sin fecha"
                lineas.append(f"  {fecha_str} | Nota {int(r.nota)} | ${float(r.monto_cobro):,.2f} | Estado: {r.resultado_observacion}")
            lineas.append(f"  >> TOTAL COBROS DE NOTAS {mes_pregunta}/{anio_pregunta}: ${agg_notas_mes['monto_cobro'].sum():,.2f} (debe coincidir con COBRO TOTAL COMPAÑÍA de arriba menos ingresos de activos fijos)")
        else:
            lineas.append(f"  Sin cobros de notas para {mes_pregunta}/{anio_pregunta}.")
    except Exception as e:
        lineas.append(f"[Error desglose cobros notas del mes: {e}]")
    #    Fuente: preparar_calendario_integrado_notas() — EXACTAMENTE igual que Alertas/Calendario
    #    Incluye monto_cobro calculado correctamente por cada nota
    # ══════════════════════════════════════════════════════════════════════
    try:
        hoy_ts = pd.Timestamp.today().normalize()
        limite_180 = hoy_ts + pd.Timedelta(days=180)
        df_calls_ctx = leer_hoja_excel("CALENDARIO_CALLS")
        cal_integrado = preparar_calendario_integrado_notas(
            df_inv=df_inv,
            df_cal=df_cal,
            df_control=df_control,
            df_calls=df_calls_ctx,
            fecha_inicio=hoy_ts,
            fecha_fin=limite_180,
        )

        lineas.append(f"\n=== CALENDARIO INTEGRADO PRÓXIMOS 180 DÍAS (fuente: Alertas/Calendario, misma lógica que pantalla) ===")
        lineas.append("  NOTA: para '¿cuánto cobraremos de notas en [mes]?' usa el bloque 'COBROS DE NOTAS [mes] — DESGLOSE YA SUMADO' de más arriba, no sumes filas de esta tabla a mano — es fácil saltarse una fila en una tabla tan larga.")

        if cal_integrado.empty:
            lineas.append("  Sin eventos en los próximos 180 días.")
        else:
            # Tabla completa
            cols_show = [c for c in ["fecha","tipo_evento","nota","estado","monto_cobro","detalle"] if c in cal_integrado.columns]
            lineas.append(cal_integrado[cols_show].to_string(index=False))

            # ── PRÓXIMO PAGO (cobro de la compañía) ──────────────────────────
            pagos = cal_integrado[
                (cal_integrado["tipo_evento"] == "PAGO") &
                (pd.to_numeric(cal_integrado.get("monto_cobro", pd.Series(dtype=float)), errors="coerce").fillna(0) > 0)
            ].sort_values("fecha") if not cal_integrado.empty else pd.DataFrame()

            lineas.append(f"\n=== PRÓXIMO COBRO DE NOTAS (PAGO con importe > 0) ===")
            if not pagos.empty:
                prox = pagos.iloc[0]
                lineas.append(f"  FECHA: {pd.Timestamp(prox['fecha']).strftime('%d/%m/%Y')}")
                lineas.append(f"  NOTA: {prox.get('nota','')}")
                lineas.append(f"  IMPORTE: ${float(prox.get('monto_cobro', 0)):,.2f}")
                lineas.append(f"  ESTADO: {prox.get('estado','')}")
                if len(pagos) > 1:
                    sig = pagos.iloc[1]
                    lineas.append(f"  SIGUIENTE: {pd.Timestamp(sig['fecha']).strftime('%d/%m/%Y')} | Nota {sig.get('nota','')} | ${float(sig.get('monto_cobro',0)):,.2f}")
                # Total cobros próximos 30 días
                hoy_30 = hoy_ts + pd.Timedelta(days=30)
                pagos_30 = pagos[pd.to_datetime(pagos["fecha"]) <= hoy_30]
                if not pagos_30.empty:
                    tot_30 = pd.to_numeric(pagos_30["monto_cobro"], errors="coerce").fillna(0).sum()
                    lineas.append(f"  >> TOTAL COBROS PRÓXIMOS 30 DÍAS: ${tot_30:,.2f}")
                tot_180 = pd.to_numeric(pagos["monto_cobro"], errors="coerce").fillna(0).sum()
                lineas.append(f"  >> TOTAL COBROS PRÓXIMOS 180 DÍAS: ${tot_180:,.2f}")
            else:
                lineas.append("  Sin pagos con importe > 0 en los próximos 180 días.")

            # ── PRÓXIMA OBSERVACIÓN POR NOTA ─────────────────────────────────
            obs = cal_integrado[cal_integrado["tipo_evento"] == "OBSERVACION"].sort_values("fecha") if not cal_integrado.empty else pd.DataFrame()
            lineas.append(f"\n=== PRÓXIMA OBSERVACIÓN POR NOTA ===")
            if not obs.empty:
                for nota_id, grupo in obs.groupby("nota"):
                    prox_obs = grupo.iloc[0]
                    lineas.append(f"  Nota {nota_id}: {pd.Timestamp(prox_obs['fecha']).strftime('%d/%m/%Y')} | Estado: {prox_obs.get('estado','')}")
            else:
                lineas.append("  Sin observaciones futuras.")

            # ── PRÓXIMOS CALLS ────────────────────────────────────────────────
            calls_cal = cal_integrado[cal_integrado["tipo_evento"] == "CALL"].sort_values("fecha") if not cal_integrado.empty else pd.DataFrame()
            if not calls_cal.empty:
                lineas.append(f"\n=== PRÓXIMOS CALLS / VENCIMIENTOS ===")
                for _, r in calls_cal.head(10).iterrows():
                    lineas.append(f"  Nota {r.get('nota','')}: {pd.Timestamp(r['fecha']).strftime('%d/%m/%Y')}")

    except Exception as e:
        lineas.append(f"[Error calendario notas: {e}]")

    # Estado de riesgo de notas
    # Fuente: RESULTADOS_OBSERVACION (estado real) + CONTROL_NOTAS (precio/barrera) + CALENDARIO_NOTAS (próx obs)
    # Se acumula aparte en 'lineas_riesgo' (no en 'lineas') para poder insertarla SIEMPRE al
    # principio del contexto, antes de truncar. Así, aunque el resto del contexto crezca mucho
    # (más inversores, más meses), la lista de notas en riesgo real nunca puede quedar cortada.
    lineas_riesgo = []

    # ══════════════════════════════════════════════════════════════════════
    # 0. EMISOR Y CUPÓN ANUAL POR NOTA — fuente: CONTROL_NOTAS (emisor) +
    #    INVERSIONES (interes_nota_anual). Se calcula en código y se inserta
    #    en lineas_riesgo (protegido de truncamiento) porque antes NO existía
    #    en ningún bloque del contexto: la IA no tenía forma de saber el cupón
    #    real de una nota confirmada y podía inventarlo. Esta es la ÚNICA
    #    fuente válida para responder "¿qué tasa/cupón paga la nota X?".
    # ══════════════════════════════════════════════════════════════════════
    try:
        lineas_riesgo.append("\n=== NOTAS ACTIVAS: EMISOR Y CUPÓN ANUAL (fuente: CONTROL_NOTAS + INVERSIONES) ===")
        lineas_riesgo.append("(USA SIEMPRE ESTOS DATOS PARA EL CUPÓN/TASA DE UNA NOTA. NO INVENTES NI APROXIMES UN CUPÓN QUE NO ESTÉ AQUÍ.)")
        df_control_activo_cupon = obtener_control_notas_activas(df_inv, df_control)
        notas_cupon = sorted(pd.to_numeric(df_control_activo_cupon.get("nota"), errors="coerce").dropna().unique().astype(int)) if df_control_activo_cupon is not None and not df_control_activo_cupon.empty else []
        if not notas_cupon:
            lineas_riesgo.append("  Sin notas activas en CONTROL_NOTAS.")
        else:
            for nota_id_c in notas_cupon:
                nombre_nota_c = f"NOTA_{nota_id_c:02d}"
                inv_nota_c = df_inv[df_inv.get("nombre_activo", pd.Series(dtype=str)).astype(str).str.upper().str.replace(" ", "_") == nombre_nota_c]
                cupon_c = None
                if not inv_nota_c.empty and "interes_nota_anual" in inv_nota_c.columns and inv_nota_c["interes_nota_anual"].notna().any():
                    cupon_c = float(inv_nota_c["interes_nota_anual"].dropna().iloc[0])
                cupon_s = f"{cupon_c*100:.2f}% anual" if cupon_c is not None else "N/D (revisar columna interes_nota_anual en INVERSIONES)"
                emisor_c = "N/D"
                filas_control_c = df_control_activo_cupon[pd.to_numeric(df_control_activo_cupon.get("nota"), errors="coerce") == nota_id_c]
                if "emisor" in filas_control_c.columns:
                    emisores_c = filas_control_c["emisor"].dropna().astype(str).str.strip()
                    emisores_c = emisores_c[emisores_c != ""]
                    if not emisores_c.empty:
                        emisor_c = emisores_c.iloc[0]
                tickers_c = ", ".join(sorted(filas_control_c["ticker"].dropna().astype(str).unique())) if "ticker" in filas_control_c.columns else ""
                lineas_riesgo.append(f"  NOTA_{nota_id_c:02d} | Emisor: {emisor_c} | Cupón anual: {cupon_s} | Tickers: {tickers_c}")
    except Exception as _e_cup:
        lineas_riesgo.append(f"[Error emisor/cupón por nota: {_e_cup}]")

    try:
        lineas_riesgo.append(f"\n=== ESTADO DE RIESGO DE NOTAS ===")
        hoy_r = pd.Timestamp.today().normalize()

        # NOTA: la hoja RESULTADOS_OBSERVACION NO se usa — no se mantiene actualizada.
        # El estado real de cada nota se calcula en vivo comparando CONTROL_NOTAS
        # (precio_compra/barrera) contra precios actuales de yfinance, más abajo.

        # 2. Precios actuales y variación por ticker — REUTILIZA construir_resumen_actual_notas_alertas(),
        # la MISMA función (idéntico criterio, idéntica fuente de precio: yf.Ticker().history)
        # que alimenta el semáforo consolidado que ve Yuri en la pantalla "Notas estructuradas".
        # Antes esta sección llamaba a obtener_datos_fundamentales() por separado (más pesada:
        # t.info con 4 reintentos + analyst_price_targets + history 1y + calendar por ticker),
        # lo que la exponía a rate-limits de Yahoo con más facilidad; y cuando fallaba, caía
        # silenciosamente en el precio "precio_actual" ya guardado en CONTROL_NOTAS (un valor
        # estático, no en vivo), lo que podía hacer que la IA reportara "sin riesgo" notas que
        # el semáforo sí marcaba en ROJO. Para que la IA NUNCA pueda divergir de lo que se ve
        # en pantalla, se usa aquí exactamente el mismo cálculo, con el mismo filtro de notas
        # activas (obtener_control_notas_activas) que usa el semáforo.
        df_control_riesgo_ia = obtener_control_notas_activas(df_inv, df_control)
        resumen_riesgo_ia = construir_resumen_actual_notas_alertas(df_control_riesgo_ia)

        def _precios_nota(nota_id):
            if resumen_riesgo_ia.empty:
                return ""
            filas = resumen_riesgo_ia[resumen_riesgo_ia["nota"] == nota_id]
            if filas.empty:
                return ""
            partes = []
            for _, r in filas.iterrows():
                ticker = r["ticker"]
                precio_actual = r["precio_actual"]
                variacion = r["variacion_%"]
                precio_contingencia = r["precio_contingencia"]
                margen_a_barrera = r["margen_a_barrera_%"]
                alerta_nivel = r["alerta_riesgo"]
                p_s = f"${precio_actual:,.2f}" if pd.notna(precio_actual) else "N/D"
                var_s = f" ({variacion:+.1f}%)" if pd.notna(variacion) else ""
                b_s = f"${precio_contingencia:,.2f}" if pd.notna(precio_contingencia) else "N/D"
                icono = {"ROJO": " 🔴 EN RIESGO (variación ≤ -30%)", "OK": " ✅ OK"}.get(alerta_nivel, " ⚪ SIN DATO")
                if pd.notna(margen_a_barrera):
                    alerta = f"{icono} | margen a la barrera (dato adicional, no decide el riesgo): {margen_a_barrera:+.1f}%"
                else:
                    alerta = f"{icono} | margen a la barrera: N/D (sin barrera cargada en CONTROL_NOTAS)"
                partes.append(f"{ticker}: precio={p_s}{var_s} | barrera={b_s}{alerta}")
            return " | ".join(partes)

        # Estado por nota = peor ticker de esa nota (mismo criterio que resumen_alertas_por_nota / el semáforo)
        def _estado_nota_precio(nota_id):
            if resumen_riesgo_ia.empty:
                return "PENDIENTE"
            filas = resumen_riesgo_ia[resumen_riesgo_ia["nota"] == nota_id]
            if filas.empty:
                return "PENDIENTE"
            variaciones = pd.to_numeric(filas["variacion_%"], errors="coerce").dropna()
            if variaciones.empty:
                return "PENDIENTE"
            peor_variacion = variaciones.min()
            nivel = clasificar_alerta_riesgo(peor_variacion)
            if nivel == "ROJO": return "ROJA"
            if nivel == "OK": return "OK"
            return "PENDIENTE"

        # 3. Próxima observación futura (CALENDARIO_NOTAS)
        df_c_riesgo = df_cal.copy()
        df_c_riesgo["fecha"] = pd.to_datetime(df_c_riesgo["fecha"], errors="coerce")
        df_c_riesgo["tipo_evento"] = df_c_riesgo["tipo_evento"].fillna("").astype(str).str.upper()
        obs_futuras = df_c_riesgo[(df_c_riesgo["tipo_evento"]=="OBSERVACION") & (df_c_riesgo["fecha"]>=hoy_r)].sort_values("fecha")
        prox_obs = obs_futuras.groupby("nota").first().reset_index()
        prox_obs_dict = {int(r["nota"]): pd.Timestamp(r["fecha"]).strftime("%d/%m/%Y") 
                        for _, r in prox_obs.iterrows() if not pd.isna(r.get("nota"))}

        # 4. Próximo call (CALENDARIO_NOTAS)
        calls_fut = df_c_riesgo[(df_c_riesgo["tipo_evento"]=="CALL") & (df_c_riesgo["fecha"]>=hoy_r)].sort_values("fecha")
        prox_call_dict = {}
        for _, r in calls_fut.iterrows():
            n = pd.to_numeric(r.get("nota"), errors="coerce")
            if pd.isna(n): continue
            ni = int(n)
            if ni not in prox_call_dict:
                prox_call_dict[ni] = pd.Timestamp(r["fecha"]).strftime("%d/%m/%Y")

        # Universo de notas a evaluar: todas las que están en el resumen (mismas notas activas que
        # ve el semáforo) más cualquiera con observación futura programada.
        notas_en_control = set()
        if not resumen_riesgo_ia.empty:
            notas_en_control = set(int(n) for n in pd.to_numeric(resumen_riesgo_ia["nota"], errors="coerce").dropna().unique())
        todas_notas = sorted(notas_en_control | set(prox_obs_dict.keys()))
        negativas, pendientes, positivas = [], [], []
        filas_riesgo_definitivo = []  # (nota_id, ticker, precio_compra, precio_actual, variación, precio_contingencia, margen)

        for nota_id in todas_notas:
            # Estado por variación de precio (igual que pantalla Notas estructuradas / semáforo)
            estado_precio = _estado_nota_precio(nota_id)
            prox_obs_s = prox_obs_dict.get(nota_id, "Sin obs programada")
            prox_call_s = prox_call_dict.get(nota_id, "")
            precios_s = _precios_nota(nota_id)

            linea = f"  NOTA_{nota_id:02d} | Estado: {estado_precio} | Próx obs: {prox_obs_s}"
            if prox_call_s:
                linea += f" | Próx call: {prox_call_s}"
            if precios_s:
                linea += f"\n    Precios: {precios_s}"

            if estado_precio == "ROJA":
                negativas.append(linea)
            elif estado_precio == "OK":
                positivas.append(linea)
            else:
                pendientes.append(linea)  # PENDIENTE = sin precio disponible

            # Detalle por ticker para la lista definitiva (misma fila que ya calculó el semáforo, sin llamadas nuevas)
            if not resumen_riesgo_ia.empty:
                filas_nota = resumen_riesgo_ia[resumen_riesgo_ia["nota"] == nota_id]
                for _, r_t in filas_nota.iterrows():
                    if r_t["alerta_riesgo"] != "ROJO":
                        continue
                    filas_riesgo_definitivo.append((
                        nota_id, r_t["ticker"], r_t["precio_compra"], r_t["precio_actual"],
                        r_t["variacion_%"], r_t["precio_contingencia"], r_t["margen_a_barrera_%"],
                    ))

        if negativas:
            lineas_riesgo.append(f"🔴 ROJAS / EN RIESGO ({len(negativas)}):")
            lineas.extend(negativas)
        else:
            lineas_riesgo.append("🔴 ROJAS / EN RIESGO: Ninguna")

        if pendientes:
            lineas_riesgo.append(f"⚪ SIN PRECIO DISPONIBLE / PENDIENTES DE DATO ({len(pendientes)}):")
            lineas.extend(pendientes)

        lineas_riesgo.append(f"RESUMEN: {len(negativas)} en riesgo (variación ≤ -30%) | {len(pendientes)} sin dato de precio | {len(positivas)} OK (positivas no se muestran)")
        lineas_riesgo.append("(USA SIEMPRE ESTOS DATOS. NO INVENTES NI CALCULES EL ESTADO DE LAS NOTAS.)")

        # ── LISTA DEFINITIVA, PRE-FILTRADA EN PYTHON (no en la IA) ──────────
        # Esta es la ÚNICA fuente válida para responder "¿qué notas están en riesgo?".
        # Calculada con la MISMA función y el mismo filtro de notas activas que el semáforo
        # de la pantalla "Notas estructuradas" — no puede divergir de lo que ve Yuri en pantalla.
        lineas_riesgo.append("\n=== NOTAS EN RIESGO REAL — LISTA DEFINITIVA (calculada en código, no en la IA) ===")
        lineas_riesgo.append(
            "Único criterio válido: variación % desde precio_compra de CADA ticker. "
            "🔴 EN RIESGO = variación ≤ -30%. Cualquier ticker/nota que NO aparezca abajo NO está en "
            "riesgo (variación > -30%) y NO debe presentarse como en riesgo. NO reclasifiques, no "
            "inventes otro criterio, no uses otro umbral. El 'margen a la barrera' que aparece junto a "
            "cada ticker es SOLO información adicional de contexto (puede venir 'N/D' si la nota no "
            "tiene barrera cargada en CONTROL_NOTAS) — NUNCA decide si una nota entra en esta lista, "
            "solo la variación la decide."
        )
        if not filas_riesgo_definitivo:
            lineas_riesgo.append("Ninguna nota está en riesgo (variación ≤ -30%) ahora mismo.")
        else:
            filas_riesgo_definitivo.sort(key=lambda f: f[4])  # peor variación primero
            for nota_id_ia, ticker_ia, compra_ia, actual_ia, variacion_ia, contingencia_ia, margen_ia in filas_riesgo_definitivo:
                margen_s = f"{margen_ia:+.1f}%" if pd.notna(margen_ia) else "N/D"
                contingencia_s = f"${contingencia_ia:,.2f}" if pd.notna(contingencia_ia) else "N/D"
                lineas_riesgo.append(
                    f"  NOTA_{int(nota_id_ia):02d} | {ticker_ia} | precio_compra=${compra_ia:,.2f} | "
                    f"precio_actual=${actual_ia:,.2f} | variación={variacion_ia:+.1f}% | 🔴 EN RIESGO | "
                    f"(dato adicional) precio_contingencia={contingencia_s} | margen a la barrera={margen_s}"
                )

        # ── CALLS POSIBLES — LISTA DEFINITIVA, PRE-FILTRADA EN PYTHON (no en la IA) ──
        # Una nota entra aquí SOLO si, a la vez: (1) tiene fecha de call futura en CALENDARIO_CALLS,
        # (2) TODOS sus cupones pagados hasta hoy fueron POSITIVA (nunca hubo barrera de cupón rota),
        # y (3) TODOS sus tickers están HOY en positivo (variación_% >= 0) frente a precio_compra
        # (mismos precios que el semáforo, vía construir_resumen_actual_notas_alertas). Yuri solo
        # quiere ver aquí las notas que SÍ cumplen — las que no cumplen no aparecen.
        lineas_riesgo.append("\n=== CALLS POSIBLES — LISTA DEFINITIVA (calculada en código, no en la IA) ===")
        lineas_riesgo.append(
            "Una nota es 'CALL POSIBLE' solo si cumple LAS DOS condiciones a la vez: "
            "(a) ha pagado el cupón TODOS los meses hasta hoy sin ninguna observación NEGATIVA "
            "(barrera de cupón rota), y (b) en la fecha de call, TODOS los tickers/acciones de la "
            "nota están en positivo (variación >= 0%) frente a su precio_compra. Si una nota NO "
            "aparece en esta lista es porque NO cumple ambas condiciones — NO la menciones como "
            "posible call. NO inventes ni recalcules este criterio."
        )
        try:
            calls_posibles = evaluar_calls_posibles_notas(df_cal, df_control, df_calls_ctx, resumen_riesgo_ia)
        except Exception as _e_cp:
            calls_posibles = []
            lineas_riesgo.append(f"[Error calculando calls posibles: {_e_cp}]")
        if not calls_posibles:
            lineas_riesgo.append("Ninguna nota cumple ahora mismo ambas condiciones de call posible.")
        else:
            for cp in calls_posibles:
                lineas_riesgo.append(
                    f"  NOTA_{cp['nota']:02d} | 🟢 CALL POSIBLE | fecha de call: "
                    f"{pd.Timestamp(cp['fecha_call']).strftime('%d/%m/%Y')} (en {cp['dias_restantes']} días) | "
                    f"cupones pagados sin fallo: {cp['n_cupones_pagados_sin_fallo']} | "
                    f"peor variación de sus tickers: {cp['peor_variacion_%']:+.1f}% | "
                    f"precios: {cp['detalle_precios']}"
                )
    except Exception as _e_r:
        import traceback as _tb_r
        _tb_texto = _tb_r.format_exc()
        lineas_riesgo.append(f"[Error riesgo notas: {_e_r}]")
        lineas_riesgo.append(f"[Detalle técnico (no mostrar al usuario, solo para diagnóstico): {_tb_texto[-500:]}]")
        if mostrar_debug_ui:
            with st.expander("⚠️ Error interno calculando riesgo de notas (clic para ver detalle técnico)", expanded=True):
                st.error(f"{type(_e_r).__name__}: {_e_r}")
                st.code(_tb_texto, language="python")

    # ══════════════════════════════════════════════════════════════════════
    # 5. TOTALES HISTÓRICOS POR ACTIVO FIJO (acumulado desde el inicio)
    #    Va a su PROPIA lista (lineas_historicos) porque al final de la función se
    #    inserta al PRINCIPIO del contexto, junto con riesgo y borradores. Antes se
    #    añadía al final de `lineas`, después del calendario integrado de 180 días y
    #    de los extractos: en cuanto el contexto crecía, el recorte de caracteres se
    #    lo comía y el asistente contestaba "no tengo el histórico acumulado" a
    #    preguntas tipo "¿cuánto ha cobrado la compañía de Paraguay desde el inicio?".
    #    Es un bloque pequeño y muy consultado — nunca debe poder perderse.
    # ══════════════════════════════════════════════════════════════════════
    lineas_historicos = []
    try:
        lineas_historicos.append("\n=== TOTALES HISTÓRICOS POR ACTIVO FIJO — ACUMULADO DESDE EL INICIO (fuente: Dashboard) ===")
        lineas_historicos.append(
            "  IMPORTANTE: estas cifras YA SON el acumulado histórico completo, calculado mes a mes desde la "
            "primera inversión de cada activo hasta el mes actual incluido. Si te preguntan '¿cuánto ha cobrado "
            "la compañía de [activo] desde el inicio?', la respuesta es el campo 'ingresado' de esa línea: dilo "
            "directamente. NUNCA respondas que no tienes el histórico acumulado, y NUNCA lo estimes con "
            "capital × tasa — el dato exacto está aquí."
        )
        tot_ing_fijos = 0.0
        tot_pag_fijos = 0.0
        for activo, tasa in [("futbol",TASA_ANUAL_FUTBOL),("paraguay",TASA_ANUAL_PARAGUAY),
                              ("bolivia",TASA_ANUAL_BOLIVIA),("motoclick",TASA_ANUAL_MOTOCLICK),
                              ("bitcoin",TASA_ANUAL_BITCOIN)]:
            t_hist = totales_activo_desde_inicio(df_inv, activo, tasa)
            tot_ing_fijos += t_hist["ingresado"]
            tot_pag_fijos += t_hist["pagado"]
            desde_txt = t_hist["desde"].strftime("%m/%Y") if t_hist["desde"] is not None else "N/D"
            lineas_historicos.append(
                f"  {activo.upper()} (primera inversión {desde_txt}): "
                f"ingresado por la compañía ${t_hist['ingresado']:,.2f} | "
                f"pagado a inversores ${t_hist['pagado']:,.2f} | "
                f"beneficio empresa ${t_hist['beneficio']:,.2f}"
            )
            for anio_h in sorted(t_hist["por_anio"]):
                d_h = t_hist["por_anio"][anio_h]
                lineas_historicos.append(
                    f"      {anio_h}: ingresado ${d_h['ingresado']:,.2f} | "
                    f"pagado ${d_h['pagado']:,.2f} | beneficio ${d_h['ingresado'] - d_h['pagado']:,.2f}"
                )
        lineas_historicos.append(
            f"  >> TOTAL TODOS LOS ACTIVOS FIJOS desde el inicio: ingresado ${tot_ing_fijos:,.2f} | "
            f"pagado a inversores ${tot_pag_fijos:,.2f} | beneficio empresa ${tot_ing_fijos - tot_pag_fijos:,.2f}"
        )
        lineas_historicos.append(
            "  (Este bloque cubre SOLO activos fijos. El acumulado histórico de NOTAS ESTRUCTURADAS no está "
            "en este bloque: si te lo piden, dilo explícitamente en vez de sumar cifras de otros bloques.)"
        )
    except Exception as e:
        lineas_historicos.append(f"[Error históricos: {e}]")

    # ══════════════════════════════════════════════════════════════════════
    # 6. EXTRACTO ACUMULADO — PRE-CALCULADO Y FILTRADO POR INVERSOR
    #    Si se detecta un inversor en la pregunta: solo sus datos
    #    Si no: resumen anual de todos (compacto)
    # ══════════════════════════════════════════════════════════════════════
    try:
        df_ext = df_inv.copy()
        for col in ["inversor","tipo_operacion"]:
            if col in df_ext.columns:
                df_ext[col] = df_ext[col].fillna("").astype(str).str.strip()
        df_ext["tipo_op_n"] = df_ext["tipo_operacion"].str.upper()
        df_ext = df_ext[df_ext["tipo_op_n"].isin(["NUEVA","CANCELADA"])].copy()
        df_ext["fecha_inversion"]        = parsear_fecha_robusta(df_ext.get("fecha_inversion"))
        df_ext["fecha_final_inversion"]  = parsear_fecha_robusta(df_ext.get("fecha_final_inversion"))
        df_ext["capital_invertido"]      = pd.to_numeric(df_ext.get("capital_invertido"), errors="coerce").fillna(0)
        df_ext["interes_inversor_anual"] = pd.to_numeric(df_ext.get("interes_inversor_anual"), errors="coerce").fillna(0)

        TRAMO_INV_E = {"ROBERTO VISCAFE", "CROWE BOLIVIA", "JR REAL ESTATE"}
        FIN_T1_E    = datetime(2026, 1, 31)
        INI_T2_E    = datetime(2026, 2, 1)
        FIN_T2_E    = datetime(2026, 6, 30)
        INI_T3_E    = datetime(2026, 7, 1)

        # Detectar inversor en la pregunta
        inversores_todos = df_ext["inversor"].dropna().unique().tolist()
        inv_detectado = None
        p_up = p.upper()
        for inv in inversores_todos:
            if inv.upper() in p_up or any(pt in p_up for pt in inv.upper().split() if len(pt) > 3):
                inv_detectado = inv
                break

        # Detectar año en la pregunta
        import re as _re2
        match_anio2 = _re2.search(r'\b(202[0-9])\b', p)
        anio_filtro = int(match_anio2.group(1)) if match_anio2 else None

        # Si hay inversor detectado, filtrar solo ese
        if inv_detectado:
            df_ext = df_ext[df_ext["inversor"].str.upper() == inv_detectado.upper()].copy()

        # Calcular desde inicio hasta hoy
        fecha_min = df_ext["fecha_inversion"].dropna().min()
        if pd.isna(fecha_min):
            fecha_min = datetime(2025, 9, 1)
        else:
            fecha_min = fecha_min.to_pydatetime()

        filas_e = []
        ai_e, mi_e = fecha_min.year, fecha_min.month
        while (ai_e, mi_e) <= (anio_hoy, mes_hoy):
            # Si hay filtro de año, saltar años que no interesan
            if anio_filtro and ai_e != anio_filtro and not inv_detectado:
                mi_e = mi_e + 1 if mi_e < 12 else 1
                ai_e = ai_e if mi_e > 1 else ai_e + 1
                continue
            dm_e = ultimo_dia_mes(ai_e, mi_e)
            im_e = datetime(ai_e, mi_e, 1)
            fm_e = datetime(ai_e, mi_e, dm_e)
            for _, row in df_ext.iterrows():
                fi = row.get("fecha_inversion")
                if pd.isna(fi): continue
                fi_dt   = fi.to_pydatetime()
                tipo_op = row["tipo_op_n"]
                ff      = row.get("fecha_final_inversion")
                if tipo_op == "CANCELADA":
                    if pd.isna(ff): continue
                    ffd = min(ff.to_pydatetime(), fm_e)
                else:
                    ffd = fm_e
                ic = max(fi_dt, im_e)
                fc = min(ffd, fm_e)
                if ic > fc: continue
                dias    = (fc - ic).days + 1
                capital = float(row["capital_invertido"])
                tasa    = float(row["interes_inversor_anual"])
                inv_up  = str(row.get("inversor","")).strip().upper()
                if inv_up in TRAMO_INV_E:
                    interes = 0.0
                    if ic <= FIN_T1_E:
                        ft1 = min(fc, FIN_T1_E)
                        interes += round((capital*0.05/12)*((ft1-ic).days+1)/dm_e, 2)
                    it2 = max(ic, INI_T2_E)
                    ft2 = min(fc, FIN_T2_E)
                    if it2 <= ft2:
                        interes += round((capital*0.075/12)*((ft2-it2).days+1)/dm_e, 2)
                    if fc >= INI_T3_E:
                        it3 = max(ic, INI_T3_E)
                        interes += round((capital*0.10/12)*((fc-it3).days+1)/dm_e, 2)
                else:
                    interes = round((capital*tasa/12)*dias/dm_e, 2)
                filas_e.append({
                    "inversor": str(row.get("inversor","")),
                    "anio": ai_e,
                    "mes": f"{mi_e:02d}/{ai_e}",
                    "interes_mes": interes,
                })
            mi_e = mi_e + 1 if mi_e < 12 else 1
            ai_e = ai_e if mi_e > 1 else ai_e + 1

        if filas_e:
            df_tot = pd.DataFrame(filas_e)
            lineas.append(f"\n=== EXTRACTO INTERESES PAGADOS A INVERSORES ===")
            lineas.append(f"(LEE ESTOS DATOS DIRECTAMENTE. PROHIBIDO CALCULAR POR TU CUENTA.)")

            if inv_detectado:
                # Detalle completo del inversor detectado
                lineas.append(f"\n-- {inv_detectado} --")
                for anio_k in sorted(df_tot["anio"].unique()):
                    df_a = df_tot[df_tot["anio"] == anio_k]
                    total_anio = float(df_a["interes_mes"].sum())
                    lineas.append(f"  {anio_k}: ${total_anio:,.2f}")
                    for mes_k, int_k in df_a.groupby("mes")["interes_mes"].sum().items():
                        if float(int_k) > 0:
                            lineas.append(f"    {mes_k}: ${float(int_k):,.2f}")
                lineas.append(f"  TOTAL ACUMULADO DESDE INICIO: ${float(df_tot['interes_mes'].sum()):,.2f}")
            else:
                # Resumen compacto de todos — solo totales por inversor y año
                for inv in sorted(df_tot["inversor"].unique()):
                    df_i = df_tot[df_tot["inversor"] == inv]
                    resumen = " | ".join([f"{a}: ${float(df_i[df_i['anio']==a]['interes_mes'].sum()):,.2f}" for a in sorted(df_i["anio"].unique())])
                    lineas.append(f"  {inv}: {resumen} | TOTAL: ${float(df_i['interes_mes'].sum()):,.2f}")
                lineas.append(f"  GRAN TOTAL: ${float(df_tot['interes_mes'].sum()):,.2f}")

    except Exception as e:
        lineas.append(f"[Error extracto acumulado: {e}]")

    # ══════════════════════════════════════════════════════════════════════
    # 7. NOTAS EN BORRADOR (extraídas por IA en el wizard, AÚN NO GUARDADAS)
    # ══════════════════════════════════════════════════════════════════════
    lineas_borradores = []
    try:
        import json as _json_ia
        df_borr = leer_hoja_excel("BORRADORES_NOTAS")
        if not df_borr.empty:
            df_borr.columns = [str(c).strip().upper() for c in df_borr.columns]
            lineas_borradores.append("\n=== NOTAS EN BORRADOR — EXTRAÍDAS POR IA EN EL WIZARD, TODAVÍA NO GUARDADAS EN CONTROL_NOTAS ===")
            lineas_borradores.append(
                "(IMPORTANTE: esto NO son posiciones activas ni confirmadas — es una extracción de un PDF que "
                "Yuri está revisando en '➕ Nueva inversión' o en 'Notas → Análisis completo de nota'. "
                "Si te pregunta por una de estas notas, acláraselo explícitamente ('esto es un borrador, aún no "
                "guardado') y avisa de que puede haber campos marcados como REVISAR pendientes de corregir. "
                "NUNCA mezcles estos datos con las notas oficiales de CONTROL_NOTAS al calcular capital, cobros "
                "o beneficios — son fuentes completamente distintas.)"
            )
            for _, fila_b in df_borr.iterrows():
                tipo_b = str(fila_b.get("TIPO", "")).strip()
                nota_b = fila_b.get("NOTA", "")
                try:
                    datos_b = _json_ia.loads(fila_b.get("JSON_DATOS", "{}"))
                except Exception:
                    continue
                etiqueta_tipo = "nueva nota, aún no creada en CONTROL_NOTAS" if tipo_b == "nueva" else "auditoría en curso de una nota ya existente"
                lineas_borradores.append(f"\n-- Borrador Nota {nota_b} ({etiqueta_tipo}) --")
                lineas_borradores.append(f"  Emisor: {datos_b.get('emisor', 'REVISAR')} | Cupón anual: {datos_b.get('cupon_anual_pct', 'REVISAR')} | Vencimiento: {datos_b.get('fecha_vencimiento', 'REVISAR')}")
                tickers_b = datos_b.get("tickers", [])
                if tickers_b:
                    resumen_tk = ", ".join(
                        f"{t.get('ticker','?')} (barrera cupón {t.get('barrera_cupon_pct','REVISAR')}, barrera capital {t.get('barrera_capital_pct','REVISAR')}, call {t.get('call_level_pct','REVISAR')})"
                        for t in tickers_b
                    )
                    lineas_borradores.append(f"  Tickers: {resumen_tk}")
                calendario_b = datos_b.get("calendario", [])
                if calendario_b:
                    lineas_borradores.append(f"  Calendario: {len(calendario_b)} eventos de observación/pago extraídos")
                fechas_call_b = datos_b.get("fechas_call", [])
                if fechas_call_b:
                    lineas_borradores.append(f"  Fechas de posible call: {', '.join(str(f) for f in fechas_call_b)}")
    except Exception as e:
        lineas_borradores.append(f"[Error notas en borrador: {e}]")

    # Insertamos las secciones de riesgo y de notas en borrador al principio (justo tras la
    # cabecera de fecha/fuente), para que NUNCA puedan perderse por el recorte de caracteres al
    # final del contexto, sin importar cuánto crezca el resto (más inversores, más meses, más notas).
    lineas_final = lineas[:2] + lineas_riesgo + lineas_historicos + lineas_borradores + lineas[2:]
    return "\n".join(lineas_final)


# ═══════════════════════════════════════════════════════════════════════════
# Llamada al modelo (mismo modelo/prompt/criterio para web y WhatsApp)
# ═══════════════════════════════════════════════════════════════════════════
def preguntar_asistente_ia_fondo(pregunta: str, df_inv, df_cal, df_control,
                                   historial_previo=None, pdfs_b64=None, usuario: str = "desconocido") -> str:
    """Llama al asistente IA del fondo: mismo modelo, mismo system prompt, mismo criterio
    de riesgo y misma función de contexto (construir_contexto_ia_fondo) que usa la pestaña
    'Asistente IA' de la app web. La reutiliza también el servicio de WhatsApp, para que
    ambos canales den siempre la misma respuesta ante la misma pregunta.
    historial_previo: lista opcional de {"role":..,"content":str} con turnos previos.
    """
    import re as _re_ia, requests as _req_ia
    if pdfs_b64 is None:
        pdfs_b64 = pdfs_disponibles_notas()
    try:
        # Detectar fecha límite en la pregunta (ej: "hasta el 07/06/2026")
        fecha_limite = None
        m_fecha = _re_ia.search(r"hasta\s+(?:el\s+)?(\d{1,2})[/\-](\d{1,2})(?:[/\-](\d{4}))?", pregunta.lower())
        if m_fecha:
            d, mo = int(m_fecha.group(1)), int(m_fecha.group(2))
            yr = int(m_fecha.group(3)) if m_fecha.group(3) else pd.Timestamp.today().year
            fecha_limite = f"{yr}-{mo:02d}-{d:02d}"

        ctx = construir_contexto_ia_fondo(pregunta, df_inv, df_cal, df_control, fecha_limite=fecha_limite, mostrar_debug_ui=False)

        # Seleccionar PDFs relevantes: SOLO si la pregunta menciona un número de nota concreto
        # que además coincide con un PDF cargado. Antes, si no había coincidencia exacta pero la
        # pregunta contenía palabras como "call"/"cobro"/"próximo", se mandaban los primeros 2 PDFs
        # de la carpeta sin relación real con la pregunta — eso podía hacer que el modelo describiera
        # el cupón/fechas/condiciones de una nota equivocada. Ahora, si no hay coincidencia exacta,
        # no se manda ningún PDF: el modelo responde solo con los datos ya verificados del Excel
        # (CONTROL_NOTAS/CALENDARIO_NOTAS/INVERSIONES), que son más fiables que adivinar el PDF.
        nums = _re_ia.findall(r"nota[_\s]*(\d+)", pregunta.lower())
        pdfs_sel = {}
        if nums:
            def _norm(s): return _re_ia.sub(r"[^a-z0-9]","",s.lower().replace(".pdf",""))
            pdfs_sel = {k:v for k,v in pdfs_b64.items() if any(_norm(k)==f"nota{n}" for n in nums)}
            if not pdfs_sel:
                ctx += (
                    f"\n\n[AVISO: la pregunta menciona la nota {', '.join(nums)} pero no hay un PDF "
                    f"cargado con ese nombre (nota{nums[0]}.pdf) — si necesitas un detalle que solo está "
                    f"en el documento oficial (condición exacta de call, cláusulas legales), dile a Yuri "
                    f"que suba ese PDF en la pestaña de notas; no lo adivines con otro PDF.]"
                )

        # Construir mensaje con PDFs + contexto
        contenido = []
        for nombre_pdf, pdf_b64 in pdfs_sel.items():
            contenido.append({"type":"document","source":{"type":"base64","media_type":"application/pdf","data":pdf_b64},"title":nombre_pdf.replace(".pdf","").upper()})
        # Límite de caracteres del contexto. Antes eran 60.000, y el contexto real ya los
        # superaba (calendario integrado de 180 días + extractos acumulados por inversor),
        # así que las últimas secciones se perdían de forma silenciosa. Se sube el límite y,
        # además, si aun así hay recorte, se le dice EXPLÍCITAMENTE al modelo qué ha pasado
        # para que avise en vez de afirmar que un dato "no existe".
        LIMITE_CTX_IA = 150000
        ctx_enviado = ctx[:LIMITE_CTX_IA]
        aviso_recorte = ""
        if len(ctx) > LIMITE_CTX_IA:
            aviso_recorte = (
                f"\n\n[AVISO INTERNO: el contexto se ha recortado ({len(ctx):,} caracteres reales, "
                f"{LIMITE_CTX_IA:,} enviados). Las secciones de riesgo, totales históricos por activo y "
                f"notas en borrador están al principio y SÍ están completas. Si un dato que esperabas "
                f"encontrar no aparece, di que puede haberse quedado fuera por el recorte y que Yuri "
                f"acote la pregunta — NO afirmes que ese dato no existe en el sistema.]"
            )
        contenido.append({"type":"text","text":f"DATOS DEL FONDO:\n\n{ctx_enviado}{aviso_recorte}\n\n---\nPREGUNTA: {pregunta}"})

        # Solo los últimos 2 turnos del historial (sin datos pesados)
        historial = []
        mensajes_prev = historial_previo or []
        for m in mensajes_prev[-4:]:
            # Solo texto plano del historial, nunca los bloques con PDFs/contexto
            if isinstance(m["content"], str):
                historial.append({"role": m["role"], "content": m["content"]})
        historial.append({"role": "user", "content": contenido})

        try:
            api_key = st.secrets.get("ANTHROPIC_API_KEY","") or st.secrets.get("anthropic",{}).get("api_key","")
        except Exception:
            api_key = ""
        if not api_key:
            import os as _os_key
            api_key = _os_key.environ.get("ANTHROPIC_API_KEY", "")
        resp = _req_ia.post("https://api.anthropic.com/v1/messages",
            headers={"Content-Type":"application/json","x-api-key":api_key,"anthropic-version":"2023-06-01"},
            json={"model":"claude-sonnet-4-5","max_tokens":2000,
                  "tools": [{"type": "web_search_20250305", "name": "web_search", "max_uses": 3}],
                  "system": """Eres el analista financiero senior de Chaparro Fernández Wealth Management, un fondo de inversión privado. Tu trabajo es razonar sobre el negocio como lo haría un socio del fondo, no solo repetir datos sueltos.

== CÓMO FUNCIONA EL NEGOCIO (esto es la base de todo razonamiento) ==

1. CAPTACIÓN: el fondo recibe capital de inversores. Cada inversor tiene una tasa de interés FIJA anual pactada (7.5%, 10%, 15%...) que se le paga siempre, TODOS LOS MESES, sobre su capital, pro-rata por días. Esa tasa es personal del inversor y NO depende de en qué activo esté invertido su dinero.

2. DESPLIEGUE: ese capital se invierte en:
- ACTIVOS FIJOS (Paraguay, MotoClick, Fútbol, Bolivia, Bitcoin): pagan al fondo un % fijo y conocido de antemano.
- NOTAS ESTRUCTURADAS: cada nota tiene su propio cupón anual (ej. 37.5%), pero el cobro es CONDICIONAL: solo se cobra si en la fecha de OBSERVACIÓN el precio de TODAS las acciones subyacentes está por encima de la barrera de cupón (normalmente 40-60% del precio inicial). Si algún activo está por debajo, ese periodo NO se cobra nada de esa nota, pero el fondo SIGUE pagando al inversor su fijo igual — esa pérdida la asume el fondo, nunca el inversor.

3. EL BENEFICIO (spread) es la diferencia entre lo que cobra el fondo del activo y lo que paga al inversor:
- Notas: beneficio = cobro real de la nota (puede ser $0 algún periodo) − pago fijo al inversor (siempre igual).
- Activos fijos: beneficio = (% del activo − % del inversor) × capital.

4. PERIODICIDAD DE COBRO DE NOTAS: la mayoría de notas pagan cupón mensualmente, pero algunas pagan trimestralmente (revisan y cobran cada 3 meses en vez de cada mes). Esto ya está detectado automáticamente a partir del calendario real de cada nota en los datos que se te dan — los importes de cobro que ves en el contexto YA incluyen el ajuste correcto (si es trimestral, el cobro de ese mes ya representa el acumulado de 3 meses). No necesitas ni debes recalcular esto tú mismo.

5. REINVERSIONES: cuando un inversor no retira su capital (por ejemplo tras un call), ese capital se reinvierte en otro activo/nota. Esto SÍ cambia lo que cobra el fondo (nuevo % del activo), pero NO cambia lo que se paga al inversor (sigue siendo su % fijo original, sobre el capital original, desde la fecha ORIGINAL de inversión — la reinversión no reinicia su reloj de intereses).

6. CALLS (llamadas anticipadas de notas): el emisor puede llamar una nota antes de vencimiento, normalmente si TODAS las acciones subyacentes están en positivo respecto al precio INICIAL (¡ojo! este umbral es mucho más exigente que la barrera de cupón — una nota puede estar pagando cupón tranquilamente sin estar ni cerca de que la llamen). Cuando hay call:
- El capital deja de generar cobro para el fondo desde esa fecha hasta que se reinvierta en otro sitio.
- El inversor sigue cobrando su fijo igual aunque el capital esté parado — pérdida pura para el fondo mientras dure.
- La condición EXACTA de call varía nota a nota y está en el documento oficial (PDF) — si tienes el PDF de esa nota cargado, úsalo como fuente principal para analizar la probabilidad de call, comparando el precio actual de cada ticker (sección ESTADO DE RIESGO DE NOTAS) contra su precio INICIAL, no contra la barrera de cupón.

7. CHAPARRO FERNANDEZ es la sociedad gestora, no un inversor externo: no cobra interés (0%), todo lo que "cobra" en su nombre es beneficio íntegro del fondo.

8. DEUDA CON JORDI CHAPARRO: las notas 1 a 8 se invirtieron con capital personal de Jordi, no del fondo. Cuando esas notas cobran, ese dinero reduce la deuda que el fondo tiene con Jordi (por haber usado su capital inicial). Los intereses devengados a JEP (todos sus activos), el reparto de dividendos y las transferencias que Jordi le hace al fondo (hoja TRANSFERENCIAS_JORDI, capital que él aporta) AUMENTAN esa deuda — las tres juegan a favor de Jordi con la misma lógica: son dinero o valor que la empresa le debe devolver.

9. MOTOCLICK — CASO ESPECIAL: a diferencia de Paraguay/Bolivia/Fútbol, el ingreso que MotoClick genera para el fondo NO es simplemente capital_invertido × 25% / 12. El capital que los inversores tienen "en el papel" asignado a MotoClick no siempre coincide con el capital que realmente está desplegado ahí día a día (puede haber devoluciones temporales de capital y reinyecciones posteriores). Por eso el ingreso real se calcula con el capital promedio diario efectivamente activo ese mes. El pago al inversor, en cambio, SIEMPRE se mantiene fijo sobre su capital nominal, sin este ajuste — solo el ingreso de la compañía varía.

10. VENCIMIENTO DE UNA NOTA SIN CALL PREVIO: si una nota llega a su fecha de vencimiento (maturity) sin haber sido llamada antes, se trata exactamente igual que un call — el capital se reinvierte en otro activo/nota, y el inversor sigue cobrando su fijo igual durante todo el proceso. En los datos, esto se marca con motivo "call final" (distinto de "call", que es la llamada anticipada por el emisor), pero a efectos de negocio ambos son equivalentes: la nota se cierra y su capital se reinvierte.

11. LA HOJA RESULTADOS_OBSERVACION NO SE USA NUNCA: no está mantenida y puede contener datos obsoletos. Si el usuario pregunta por ella, dile que esa hoja no se usa como fuente y que el estado real de cada nota se calcula en vivo con precios actuales.

12. FECHA DE INICIO DE UNA NOTA PARA EFECTOS DE PAGO AL INVERSOR: no es la fecha en la que arranca la nota en el mercado (Initial Valuation Date), sino la PRIMERA FECHA DE PAGO (cobro) del calendario de esa nota, menos 1 mes. Ej.: si la nota empieza el 19/09 y el primer pago es el 25/10, la fecha de inicio que cuenta para calcular los intereses del inversor es el 25/09 (un mes antes del primer pago), no el 19/09.

13. LA HOJA REINVERSIONES ES SOLO TRAZABILIDAD HISTÓRICA, no una fuente de cálculo financiero. Vincula cada operación de reinversión (id_inversion_destino) con de dónde venía originalmente ese capital (id_inversion_origen), para poder responder preguntas del tipo "¿de dónde viene el capital de esta nota?" aunque haya pasado por muchas reinversiones. Los importes de esta hoja NO tienen por qué cuadrar exactamente con el capital_invertido actual de la posición destino — no la uses para calcular capital, ingresos ni beneficios, solo para explicar el origen/historial de un capital si te lo preguntan explícitamente.

14. LOS CALLS SIEMPRE SON TOTALES: cuando una nota es llamada, se cancela el 100% de su capital de golpe, nunca una parte. Si un documento menciona la posibilidad de un call parcial por sorteo entre tenedores (algunos prospectos legales lo mencionan como cláusula genérica), ignóralo — en la operativa real del fondo nunca ha pasado y no debe asumirse salvo que el usuario diga lo contrario explícitamente.

15. DIVISA Y REDONDEO: todo el fondo opera en dólares estadounidenses (USD). Los importes siempre se redondean a 2 decimales.

16. TRASPASOS ENTRE INVERSORES: si un inversor cede su posición a otro (ej. Jordi le pasa capital a Eva), se registra como una CANCELACIÓN de la posición del que cede + una inversión NUEVA (o una REINVERSIÓN, si el capital que se traspasa proviene de una posición previa que venció por call o vencimiento) del que recibe. Si ves una cancelación y una nueva/reinversión con fechas coincidentes y capital similar, puede tratarse de un traspaso de este tipo.

17. NO EXISTE EL RETIRO PARCIAL DE CAPITAL: un inversor solo puede retirar los intereses que genera su capital, nunca una parte del principal. Para sacar dinero del principal, la única vía es cancelar la posición COMPLETA (fila entera). Si un inversor tenía $50,000 y ahora aparece con $30,000, eso NO es un retiro parcial — implica que se canceló la posición de $50,000 y se abrió una nueva de $30,000 (dos operaciones distintas, no un ajuste de la misma).

18. NO SE PUEDE ENTRAR CON CAPITAL NUEVO A MITAD DE VIDA DE UNA NOTA YA ACTIVA: la única forma de que un inversor gane exposición a una nota que ya lleva tiempo funcionando es comprando la posición de otro inversor que ya estaba dentro (un traspaso, ver punto 16) — nunca aportando capital fresco directamente a una nota en curso. En un traspaso de este tipo, todos los términos de la nota (precio de compra, barreras, fechas de calendario, calls) se mantienen EXACTAMENTE igual que en la posición original — solo cambia el nombre del inversor en esa fila de INVERSIONES.

19. CUANDO UNA NOTA CON VARIOS INVERSORES ES LLAMADA (call), el capital de cada inversor se reinvierte de forma INDEPENDIENTE, pudiendo acabar en sitios distintos (ej. si la nota tenía 4 inversores, cada uno puede terminar reinvertido en una nota o activo diferente, no necesariamente juntos). No asumas que todos los inversores de una nota llamada siguen juntos después del call.

20. NO HAY COMISIÓN DE GESTIÓN ADICIONAL: el spread (diferencia entre lo que cobra el fondo del activo y lo que paga al inversor) es la ÚNICA fuente de ingreso del fondo. No existe ningún "management fee" u otra comisión extra sobre el capital de los inversores.

21. EL PAGO AL INVERSOR ES UNA GARANTÍA UNIVERSAL, EN CUALQUIER TIPO DE ACTIVO: no es exclusivo de las notas estructuradas — si un activo fijo (Paraguay, MotoClick, Fútbol, Bolivia) no paga su % pactado algún mes por el motivo que sea, el fondo sigue pagando al inversor su tasa fija igualmente, asumiendo la pérdida él mismo. El inversor cobra su fijo pase lo que pase, sea cual sea el activo donde esté su capital.

22. RENTABILIDAD MEDIA — DOS PERSPECTIVAS DISTINTAS, no las confundas:
- "Rentabilidad que PAGA el fondo a los inversores" = media ponderada por capital de interes_inversor_anual sobre el capital activo de TODOS los inversores (excluyendo a Chaparro Fernández, que no es un inversor real). Esta es la perspectiva del COSTE del fondo.
- "Rentabilidad que OBTIENE Chaparro Fernández como inversor" (o "rentabilidad de las inversiones de Chaparro Fernández") = media ponderada por capital de interes_nota_anual / tasa fija del activo, SOLO sobre las filas donde inversor = CHAPARRO FERNANDEZ (incluyendo sus reinversiones). Esta es la perspectiva del RENDIMIENTO del capital propio de la compañía colocado en los activos, y será bastante más alta que la anterior (activos rinden 15-38%, muy por encima del 10-15% que se paga a inversores externos).
Si te preguntan "rentabilidad media de las inversiones de/en Chaparro Fernández" sin más contexto, usa la SEGUNDA interpretación (rendimiento del capital propio como inversor).

23. BENEFICIO REAL DE LA COMPAÑÍA (distinto del "capital a nombre de Chaparro Fernández"): el capital que aparece invertido bajo el nombre "CHAPARRO FERNANDEZ" en INVERSIONES son los beneficios acumulados de la compañía que se han ido reinvirtiendo, PERO no todo ese capital es beneficio realmente disponible — una parte es, en realidad, intereses ya devengados a otros inversores que aún NO se les ha pagado en efectivo (porque su `pago_intereses = "reinvierte"`, es decir, todos los inversores EXCEPTO JEP, que es el único con `pago_intereses = "paga"`). Fórmula:
BENEFICIO REAL = (Capital total a nombre de CHAPARRO FERNANDEZ activo hoy — usando fecha_final_inversion para excluir posiciones ya cerradas/reinvertidas, sin duplicar entre una fila original cerrada y su reinversión sucesora)
         − (Suma de TODOS los intereses devengados desde el inicio de cada posición hasta hoy, de TODOS los inversores con pago_intereses="reinvierte", EXCLUYENDO a JEP y excluyendo al propio Chaparro Fernández)
Este cálculo de "intereses devengados y no pagados" es ACUMULADO desde que cada inversor empezó (no se resetea nunca, es una deuda pendiente continua) — usa la misma lógica de extractos (ignorar fecha_final_inversion en filas NUEVA sin reinversión sucesora, respetarla en CANCELADA).

24. QUÉ ES UNA "NOTA EN RIESGO" — UN SOLO CRITERIO, SIMPLE: la VARIACIÓN % desde el precio de compra de esa acción concreta.
🔴 EN RIESGO si la variación es ≤ -30% (ej. compra $100 → precio actual ≤ $70).
✅ OK si la variación es > -30%.
Se eligió variación fija en vez de margen a la barrera de contingencia (que se usaba antes) por dos motivos: (a) es mucho más fácil de verificar a mano contra el Excel — es una resta simple; y (b) depende solo de precio_compra, que se consulta en vivo con Yahoo Finance al crear la nota y casi nunca falla, a diferencia de barrera_cupon/barrera_capital, que se han visto mal extraídas del PDF varias veces. El "margen a la barrera" TODAVÍA se calcula y se muestra junto a cada ticker en riesgo, como dato adicional de contexto (puede venir "N/D" si la nota no tiene barrera cargada) — pero NUNCA decide si una nota entra en la lista de riesgo ni qué color tiene. Solo la variación decide.

⚠️ IMPORTANTE — LO QUE "RIESGO" **NO** ES: la proximidad de la próxima fecha de observación/cobro NO es un criterio de riesgo. Que una nota tenga observación "mañana" no la hace estar "en riesgo" — puede tener su próxima observación mañana y estar perfectamente por encima del -30%. NUNCA presentes "próxima observación inminente" como sinónimo o señal de riesgo. El riesgo se mide EXCLUSIVAMENTE por la variación % desde compra (la regla de arriba). Las fechas de observación/cobro se añaden DESPUÉS, como información complementaria de una nota que ya identificaste como en riesgo por variación — nunca como el criterio que decide si está en riesgo.

25. CUANDO TE PREGUNTEN POR NOTAS EN RIESGO (o por una nota en riesgo concreta), tu ÚNICA fuente para decidir qué notas están en riesgo es el bloque "NOTAS EN RIESGO REAL — LISTA DEFINITIVA (calculada en código, no en la IA)". Esa lista ya viene filtrada y ordenada (peor variación primero) — NO recalcules, NO reclasifiques, NO añadas notas que no estén ahí, y NO uses ningún otro umbral que no sea el que ya viene aplicado ahí. Si la lista dice que está vacía, contesta que no hay ninguna nota en riesgo ahora mismo — no inventes ninguna. Para CADA nota de esa lista, tu respuesta SIEMPRE debe incluir, en este orden:
- El número de nota y el ticker en riesgo, tal como aparece en la lista.
- Precio de compra, precio actual y variación %, copiados tal cual de la lista (no los recalcules).
- El margen a la barrera de contingencia como dato adicional (si dice "N/D", decilo así — no es un fallo, es que esa nota no tiene barrera cargada en CONTROL_NOTAS).
- Una breve explicación de la caída de cada acción en riesgo (usa la herramienta de búsqueda web SOLO para noticias/contexto cualitativo — resultados recientes, motivo del movimiento del precio, ratings, eventos — NUNCA para re-consultar el precio actual o la variación, que ya vienen calculados con precisión y son más fiables que cualquier precio que encuentres buscando en la web).
- REGLA ESTRICTA DE PRECIOS: cualquier precio, rango de precio o variación % que menciones EN CUALQUIER PARTE de tu respuesta —ya sea en la tabla de datos, en una frase tipo "situación actual de [ticker]", en la conclusión, donde sea— tiene que ser EXACTAMENTE el mismo número que ya viene en el contexto (ESTADO DE RIESGO DE NOTAS / NOTAS EN RIESGO REAL), copiado tal cual, nunca un valor nuevo ni un rango (ej. "$79.80-$79.84") que hayas visto en un resultado de búsqueda web. La búsqueda web es solo para texto cualitativo (noticias, motivo del movimiento, eventos) — jamás como fuente de un número de precio que vas a mostrar. Si en algún momento se te ocurre escribir un precio y no lo tenés ya calculado en el contexto, no lo escribas — usa una frase sin cifra o remití al dato que sí tenés.
- La situación general de esa nota: próxima fecha de observación (bloque PRÓXIMAS OBSERVACIONES) y próxima fecha de cobro/call (bloque CALENDARIO NOTAS / PRÓXIMOS CALLS), siempre como dato complementario al final, nunca como criterio de riesgo ni como parte del título/encabezado de la nota.
- Si para alguna nota falta el precio de compra (aparece "N/D" en CONTROL_NOTAS), dilo así de simple: "no tengo el precio de compra cargado para la Nota X en CONTROL_NOTAS, revísalo". NO intentes compensar buscando datos en la web ni fabriques un análisis con información incompleta.

26. TONO Y FORMATO AL HABLAR DE RIESGO: mantén un tono profesional y calmado, como un analista senior informando a un socio — NUNCA alarmista. Prohibido usar mayúsculas tipo "SITUACIÓN CRÍTICA", "RIESGO EXTREMO", "ESTO ES CRÍTICO", símbolos de alerta grandes o encabezados dramáticos. Presenta los datos con claridad y deja que la gravedad se entienda por las cifras, no por el tono. Si de verdad no puedes completar un análisis (por ejemplo, por límite de búsquedas), dilo en una frase breve y sigue con lo que sí puedas dar — no lo conviertas en el titular de la respuesta.

27. NOTAS EN BORRADOR (sección "NOTAS EN BORRADOR" del contexto): son extracciones de PDF hechas por IA en el wizard "Notas estructuradas", que Yuri puede estar revisando y todavía NO están guardadas en CONTROL_NOTAS/CALENDARIO_NOTAS. Si te pregunta por una nota y solo la encuentras ahí (no en las fuentes oficiales), dile explícitamente que es un borrador sin guardar, con los datos que tengas, y avisa de que puede haber campos "REVISAR" sin corregir todavía. Puedes usarla para responder preguntas sobre esa nota en concreto, pero NUNCA la incluyas en cálculos de capital activo, cobros o beneficios del fondo — para eso solo cuentan las notas ya confirmadas en CONTROL_NOTAS.

28. REGLA GENERAL DE PRECIOS (aplica a CUALQUIER pregunta sobre una nota o ticker, no solo a las de riesgo): cada vez que menciones un precio, un rango de precio o una variación % de un ticker —en un resumen de nota, en una tabla, en una frase de "situación actual", en una conclusión, donde sea—, tiene que ser EXACTAMENTE el número que ya viene calculado en el contexto (ESTADO DE RIESGO DE NOTAS / NOTAS EN RIESGO REAL), copiado tal cual, nunca un valor que hayas visto en un resultado de búsqueda web. Señal de alerta: si estás a punto de escribir un precio como rango (ej. "$79.80-$79.84") en vez de un número único, es casi seguro que lo sacaste de una búsqueda web y no del contexto — pará y usa el número exacto del contexto en su lugar. La búsqueda web sirve únicamente para texto cualitativo (noticias, motivo de un movimiento, eventos de la compañía, resultados trimestrales) — nunca como fuente de un precio, rango o variación % que vayas a mostrarle a Yuri. Si no tenés el precio de un ticker ya calculado en el contexto, no inventes uno ni lo completes con la web — decí que no lo tenés.

29. REGLA GENERAL DE CUPÓN/TASA DE NOTA Y FECHAS DE CALENDARIO (misma lógica que la regla 28, aplicada a cupones y fechas): el cupón/tasa anual y el emisor de una nota SIEMPRE deben salir del bloque "NOTAS ACTIVAS: EMISOR Y CUPÓN ANUAL" del contexto, copiados tal cual — nunca un número parecido de memoria, ni el cupón de otra nota similar, ni un valor deducido del PDF si el bloque ya trae el dato (el PDF solo es fuente principal para condiciones de call, barreras y cláusulas legales, NO para el número de cupón si ya está en ese bloque). Igual con fechas de observación/pago/call: solo son válidas las que aparecen en los bloques CALENDARIO INTEGRADO / PRÓXIMA OBSERVACIÓN POR NOTA / PRÓXIMOS CALLS. Ese calendario solo cubre los próximos 180 días — si te preguntan por una fecha pasada, o más allá de ese rango, dilo explícitamente ("no tengo esa fecha en el rango de datos que manejo, de 180 días vista") en vez de estimarla o inventarla. Si para una nota concreta ves "N/D" en el cupón, dilo así de simple y sugiere revisar la columna interes_nota_anual en INVERSIONES — no lo completes con un valor aproximado.

30. NUNCA USES LA BÚSQUEDA WEB PARA DATOS INTERNOS DEL FONDO. Los capitales, cobros, ingresos, beneficios, pagos a inversores, históricos acumulados, notas, calendarios y nombres de inversores de Chaparro Fernández son datos PRIVADOS que solo existen en el contexto que te paso — no están en internet y buscarlos ahí no puede devolver nada útil. Si un dato interno no lo encuentras en el contexto, la respuesta correcta es decir en una frase que no lo tienes en el contexto y qué haría falta para tenerlo; jamás lanzar una búsqueda web ni comentar que "la búsqueda web no encontró información sobre el fondo". La búsqueda web sirve EXCLUSIVAMENTE para información pública de mercado sobre los tickers subyacentes: noticias, motivo de un movimiento de precio, resultados trimestrales, rating de analistas. Nada más.

31. ANTES DE DECIR "NO TENGO ESE DATO", RELEE EL CONTEXTO. En particular, el acumulado histórico de cada activo fijo (ingresado por la compañía, pagado a inversores y beneficio, total y por año) SIEMPRE viene en el bloque "TOTALES HISTÓRICOS POR ACTIVO FIJO — ACUMULADO DESDE EL INICIO", situado al principio del contexto. Preguntas como "¿cuánto ha cobrado la compañía de Paraguay desde el inicio?" se responden leyendo ese bloque, no estimando ni pidiendo permiso para calcular una aproximación.

== TUS FUENTES DE DATOS EN EL CONTEXTO ==
- TOTALES HISTÓRICOS POR ACTIVO FIJO: acumulado desde el inicio (ingresado, pagado a inversores y beneficio), total y desglosado por año, para Paraguay, MotoClick, Fútbol, Bolivia y Bitcoin. Es tu única fuente para cualquier pregunta de "desde el inicio" o "en total" sobre un activo fijo.
- CALENDARIO NOTAS: fechas y montos de cobro ya calculados con la periodicidad correcta.
- ESTADO DE RIESGO DE NOTAS: precio actual vs barrera, por nota y ticker.
- INTERESES A PAGAR A INVERSORES / EXTRACTO: pagos ya calculados con la lógica de extractos (solo NUEVA/CANCELADA, reinversiones excluidas del pago pero no del capital activo).
- CAPITAL ACTIVO: capital desplegado por activo e inversor.
- NOTAS EN BORRADOR: extracciones de IA aún no guardadas — ver regla 27, úsalas solo para responder sobre esa nota puntual, nunca en cálculos agregados.
- PDFs de notas (si están cargados y son relevantes): documento oficial con condiciones exactas de call, barreras, cupón y fechas — tu fuente más fiable para razonar sobre probabilidad de call o condiciones legales exactas.

== CÓMO RESPONDER ==
- Si la pregunta pide un dato que YA está calculado en el contexto (fecha, importe, capital) → léelo y repítelo tal cual. No lo "mejores" ni lo recalcules — el contexto ya sigue las reglas correctas del negocio explicadas arriba.
- Si la pregunta requiere COMBINAR o RAZONAR sobre varios datos que sí tienes (ej. "¿es probable que llamen la Nota 15?", "¿cuánto perderíamos si no llaman ninguna nota este trimestre?", "compárame el riesgo de estas dos notas") → SÍ debes razonar, aplicando la lógica de negocio de arriba y, si hay PDF, sus condiciones exactas. Explica brevemente tu razonamiento, no solo la conclusión.
- Para preguntas de CALLS concretos: compara precio actual de cada ticker contra su precio INICIAL (no la barrera de cupón) y da un veredicto argumentado — probable / improbable / imposible de saber sin más datos — señalando qué ticker concreto lo impide si aplica.
- Si no tienes el dato ni puedes derivarlo razonando sobre lo que sí tienes, dilo con claridad — no inventes cifras.

== INVERSORES Y TASAS ==
LEO: 10% | JORDI CHAPARRO: 15% | YURI FERNANDEZ: 15%
ROBERTO VISCAFE: 5% hasta 31/01/2026, 7.5% desde 01/02/2026 hasta 30/06/2026, 10% desde 01/07/2026
CROWE BOLIVIA: 5% hasta 31/01/2026, 7.5% desde 01/02/2026 hasta 30/06/2026, 10% desde 01/07/2026
JR REAL ESTATE: 5% hasta 31/01/2026, 7.5% desde 01/02/2026 hasta 30/06/2026, 10% desde 01/07/2026
2012 JACC GROUP: 10% | PEDRO MAGAÑA: 10% | PAM: 10%
CHAPARRO FERNANDEZ: 0% — sociedad gestora, no recibe pago
GOLDEN BRICKS: 10% | TERESA: 10% | JEP: 15%
JORDI ESPECIAL: 10% | EVA CHAPARRO: 15% | PAOLA CHAPARRO: 15% | JAPAN JORDI: 15%

== FORMATO ==
Responde SIEMPRE en español. Sé conciso cuando el dato es directo; desarrolla el razonamiento cuando la pregunta lo requiera. Fechas DD/MM/YYYY, importes con $ y 2 decimales.""",
                  "messages":historial},timeout=60)
        data = resp.json()
        _uso = data.get("usage", {}) or {}
        log_uso_ia(usuario, "fondo", _uso.get("input_tokens", 0), _uso.get("output_tokens", 0))
        respuesta = "".join(b.get("text","") for b in data.get("content",[]) if b.get("type")=="text")
        if not respuesta:
            respuesta = f"Error API: {data.get('error',{}).get('message',str(data))}"
        return respuesta
    except Exception as e:
        return f"Error: {e}"


def pdfs_disponibles_notas() -> dict:
    """Lista y codifica en base64 los PDFs de notas estructuradas (carpeta notas_pdfs/ del repo
    + subidas temporales). Función de nivel de módulo para que la reutilicen tanto la app web
    como el servicio de WhatsApp."""
    import os, re as _re_ia, tempfile, base64 as _b64
    carpeta_tmp = os.path.join(tempfile.gettempdir(), "notas_pdfs_cf")
    carpeta_repo = os.path.join(os.path.dirname(os.path.abspath(__file__)), "notas_pdfs")
    os.makedirs(carpeta_tmp, exist_ok=True)
    pdfs = {}
    for carpeta in [carpeta_repo, carpeta_tmp]:
        if not os.path.exists(carpeta) or not os.path.isdir(carpeta):
            continue
        for fname in sorted(os.listdir(carpeta)):
            if fname.lower().endswith(".pdf"):
                clave = _re_ia.sub(r"[^a-z0-9]", "", fname.lower().replace(".pdf",""))
                if clave not in {_re_ia.sub(r"[^a-z0-9]","",k.lower().replace(".pdf","")) for k in pdfs}:
                    try:
                        with open(os.path.join(carpeta, fname), "rb") as f:
                            pdfs[fname] = _b64.standard_b64encode(f.read()).decode()
                    except Exception:
                        pass
    return pdfs


def seccion_asistente_ia_fondo():
    """IA única: fondo completo + PDFs de notas estructuradas."""
    df_inv, df_cal, df_control = cargar_excel_completo()
    st.header("✨ Asistente IA — Chaparro Fernández Wealth")
    st.caption("Pregúntame cualquier cosa: cobros, calls, barreras, capitales, intereses, inversores, beneficios...")

    import os, re as _re_ia, requests as _req_ia

    # Los PDFs de notas (si existen en el repo) se siguen usando automáticamente
    # como contexto adicional cuando la pregunta menciona una nota concreta,
    # pero ya no se gestionan ni se muestran desde esta pantalla.
    pdfs_b64 = pdfs_disponibles_notas()

    # (El contexto se construye con construir_contexto_ia_fondo(), definida a nivel de módulo)

    # ── Chat ──────────────────────────────────────────────────────────────────
    if "chat_ia_cf" not in st.session_state:
        st.session_state["chat_ia_cf"] = []

    for msg in st.session_state["chat_ia_cf"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    pregunta = st.chat_input("Escribe tu pregunta...", key="chat_input_ia_cf")
    if pregunta:
        st.session_state["chat_ia_cf"].append({"role": "user", "content": pregunta})
        st.rerun()

    if st.session_state["chat_ia_cf"] and st.session_state["chat_ia_cf"][-1]["role"] == "user":
        ultima = st.session_state["chat_ia_cf"][-1]["content"]
        with st.chat_message("assistant"):
            with st.spinner("Analizando..."):
                try:
                    mensajes_prev = st.session_state["chat_ia_cf"][:-1]
                    respuesta = preguntar_asistente_ia_fondo(
                        ultima, df_inv, df_cal, df_control,
                        historial_previo=mensajes_prev, pdfs_b64=pdfs_b64,
                        usuario=str(st.session_state.get("usuario", "desconocido")),
                    )
                except Exception as e:
                    respuesta = f"Error: {e}"
                st.markdown(_md_seguro(respuesta))
        st.session_state["chat_ia_cf"].append({"role":"assistant","content":respuesta})

    if st.session_state["chat_ia_cf"]:
        if st.button("🗑️ Limpiar conversación", key="btn_limpiar_ia_cf"):
            st.session_state["chat_ia_cf"] = []
            st.rerun()


if __name__ == "__main__":  # menu principal / routing: solo se ejecuta con `streamlit run`, no al importar
    # ── Portal de inversor: acceso limitado, se corta aquí antes del menú de administración ──
    if st.session_state.get("tipo_usuario") == "inversor":
        seccion_portal_inversor(st.session_state.usuario)
        st.stop()

    # ── Admin viendo un perfil de inversor adicional (ver ADMIN_VISTAS_INVERSOR_ADICIONALES):
    # se corta igual que el portal de inversor normal, pero sigue siendo su sesión de admin
    # por debajo — con el selector de la barra lateral puede volver al panel cuando quiera. ──
    if st.session_state.get("tipo_usuario") == "admin" and st.session_state.get("vista_admin_como_inversor"):
        seccion_portal_inversor(st.session_state.vista_admin_como_inversor)
        st.stop()

    _es_yuri = str(st.session_state.usuario).strip().lower() == "yuri"
    _es_jordi_o_alan = str(st.session_state.usuario).strip().lower() in ("jordi", "alan")

    if _es_jordi_o_alan:
        # Jordi y Alan ven un menú reducido: sin Extractos (ni, por supuesto, Gestión de
        # Excel / Nueva inversión / Uso IA / Gastos, que ya eran solo para Yuri).
        menu_opciones = [
            "Dashboard financiero", "Centro de control",
            "Notas estructuradas", "Alertas y calendario",
            "📰 Noticias",
            "🏦 Deuda Jordi Chaparro", "✨ Asistente IA",
        ]
    else:
        menu_opciones = [
            "Dashboard financiero", "Centro de control",
            "Notas estructuradas", "Alertas y calendario", "Extractos", "🏦 Contabilidad",
            "📰 Noticias",
            "🏦 Deuda Jordi Chaparro", "✨ Asistente IA",
        ]
        # "Gestión de Excel" y "➕ Nueva inversión" son acceso directo/de escritura al Excel del
        # fondo: solo Yuri debe verlos.
        if _es_yuri:
            menu_opciones.insert(5, "Gestión de Excel")
            menu_opciones.insert(6, "➕ Nueva inversión")
            menu_opciones.insert(7, "📊 Uso IA")
            menu_opciones.insert(8, "💰 Gastos")

    menu = st.sidebar.selectbox("Menú principal", menu_opciones, key="menu_principal_selector")

    if menu == "Dashboard financiero":
        dashboard_financiero()

    elif menu == "Centro de control":
        centro_control_inversiones()
    elif menu == "Notas estructuradas":
        seccion_notas_archivo()
    elif menu == "Alertas y calendario":
        panel_alertas_y_calendario()

    elif menu == "Extractos":
        seccion_extractos()

    elif menu == "🏦 Contabilidad":
        seccion_contabilidad_banco()
    elif menu == "📰 Noticias":
        seccion_noticias()
    elif menu == "Gestión de Excel" and _es_yuri:
        seccion_gestion_excel()
    elif menu == "➕ Nueva inversión" and _es_yuri:
        _df_inv_ni, _df_cal_ni, _df_control_ni = cargar_excel_completo()
        seccion_nueva_inversion(_df_inv_ni, _df_cal_ni, _df_control_ni)
    elif menu == "📊 Uso IA" and _es_yuri:
        seccion_uso_ia()
    elif menu == "💰 Gastos" and _es_yuri:
        seccion_gastos_plataforma()
    elif menu == "🏦 Deuda Jordi Chaparro":
        seccion_deuda_jordi()
    elif menu == "✨ Asistente IA":
        seccion_asistente_ia_fondo()
